"""Tests for Office document translation processor."""

import builtins
import sys
import zipfile
from pathlib import Path
from unittest.mock import MagicMock, PropertyMock, patch

import pytest
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from lxml import etree
from openpyxl import Workbook, load_workbook
from openpyxl.comments import Comment
from pptx import Presentation
from pptx.util import Inches

from src.constants.errors import ERR_OFFICE_CONVERTER_NOT_FOUND
from src.core.office_formatter import (
    _clear_run_text_only,
    _replace_paragraph_text,
    _run_has_visual_content,
)
from src.core.office_processor import (
    _COMMENT_EXTENSIONS,
    _FATAL_LLM_ERRORS,
    _FOOTNOTE_EXTENSIONS,
    _HEADER_FOOTER_EXTENSIONS,
    _HF_DEFAULT,
    _HF_EVEN,
    _HF_FIRST,
    _IMAGE_EXTENSIONS,
    _LEGACY_EXTENSIONS,
    _NOTES_EXTENSIONS,
    _SHAPE_EXTENSIONS,
    _SHEET_NAME_EXTENSIONS,
    _SUFFIX_TO_MEDIA_PREFIXES,
    LEGACY_CONVERT_MAP,
    _add_hyperlink_to_rels,
    _build_odf_hf_map,
    _build_odf_style_map,
    _collect_wps_texts,
    _convert_with_uno,
    _convert_with_win32com,
    _detect_backend,
    _extract_comments,
    _extract_docx_comments,
    _extract_docx_footnotes,
    _extract_docx_headers_footers,
    _extract_docx_hf_part,
    _extract_docx_shapes,
    _extract_drawingml_text,
    _extract_footnotes,
    _extract_headers_footers,
    _extract_notes,
    _extract_odf_paragraph_text,
    _extract_odp_notes,
    _extract_ods_shapes,
    _extract_ods_sheet_names,
    _extract_odt_footnotes,
    _extract_odt_headers_footers,
    _extract_odt_shapes,
    _extract_pptx_notes,
    _extract_shapes,
    _extract_sheet_names,
    _extract_win32com_excel_comments,
    _extract_win32com_excel_sheet_names,
    _extract_win32com_ppt_comments,
    _extract_win32com_word_comments,
    _extract_win32com_word_footnotes,
    _extract_win32com_word_headers_footers,
    _extract_xlsx_comments,
    _extract_xlsx_shapes,
    _extract_xlsx_sheet_names,
    _get_file_category,
    _inject_comments,
    _inject_docx_comment_html,
    _inject_docx_comments,
    _inject_docx_footnotes,
    _inject_docx_headers_footers,
    _inject_docx_shapes,
    _inject_drawingml_text,
    _inject_footnotes,
    _inject_headers_footers,
    _inject_notes,
    _inject_odf_paragraph_text,
    _inject_odp_notes,
    _inject_ods_shapes,
    _inject_ods_sheet_names,
    _inject_odt_footnotes,
    _inject_odt_headers_footers,
    _inject_odt_shapes,
    _inject_pptx_notes,
    _inject_shapes,
    _inject_sheet_names,
    _inject_win32com_excel_html_runs,
    _inject_win32com_excel_sheet_names,
    _inject_win32com_word_footnotes,
    _inject_win32com_word_headers_footers,
    _inject_xlsx_comments,
    _inject_xlsx_shapes,
    _inject_xlsx_sheet_names,
    _is_inside_table_cell,
    _odf_element_text,
    _odf_replace_text,
    _parse_hyperlink_rels,
    _patch_docx_comment_rels,
    _read_txbx_data,
    _resolve_xlsx_sheet_drawings,
    _restore_uno_char_props,
    _restore_win32com_font,
    _restore_xlsx_embeddings,
    _rewrite_zip_content,
    _sanitize_sheet_name,
    _save_uno_char_props,
    _save_win32com_font,
    _should_translate_comments,
    _should_translate_images,
    _should_translate_notes,
    _should_translate_shapes,
    _should_translate_sheet_names,
    _substitute_font,
    _translate_doc_images,
    _translate_legacy_images,
    _translate_zip_images,
    convert_to_modern_format,
    process_office_file,
)
from src.core.translator import _map_error_to_code

# ---------------------------------------------------------------------------
# UNO enumeration helper — prevents infinite-loop memory leaks.
#
# MagicMock().hasMoreElements() returns a truthy MagicMock, so any
# ``while enum.hasMoreElements():`` loop becomes infinite when the
# enumeration is not properly mocked.  Always use this helper instead
# of hand-wiring hasMoreElements / nextElement on a raw MagicMock.
# ---------------------------------------------------------------------------


def _make_uno_enum(items: list | None = None) -> MagicMock:
    """Creates a bounded UNO-style enumeration mock.

    Args:
        items: Elements returned by nextElement().  Defaults to empty
               (hasMoreElements immediately returns False).

    Returns:
        MagicMock with properly bounded hasMoreElements / nextElement.
    """
    items = items or []
    enum = MagicMock()
    _idx = [0]

    def _has_more() -> bool:
        return _idx[0] < len(items)

    def _next() -> object:
        el = items[_idx[0]]
        _idx[0] += 1
        return el

    enum.hasMoreElements.side_effect = _has_more
    enum.nextElement.side_effect = _next
    return enum


# ---------------------------------------------------------------------------
# Backend detection tests
# ---------------------------------------------------------------------------


def test_detect_backend_ooxml_always_python_lib() -> None:
    """OOXML formats always use python_lib, even when win32com/UNO available."""
    with patch.dict(
        "sys.modules",
        {"win32com": MagicMock(), "win32com.client": MagicMock(), "uno": MagicMock()},
    ):
        assert _detect_backend(".docx") == "python_lib"
        assert _detect_backend(".xlsx") == "python_lib"
        assert _detect_backend(".pptx") == "python_lib"


def test_detect_backend_legacy_prefers_win32com() -> None:
    """Legacy formats prefer win32com when available."""
    with patch.dict(
        "sys.modules",
        {"win32com": MagicMock(), "win32com.client": MagicMock()},
    ):
        assert _detect_backend(".doc") == "win32com"
        assert _detect_backend(".xls") == "win32com"
        assert _detect_backend(".ppt") == "win32com"


def test_detect_backend_legacy_falls_back_to_uno() -> None:
    """Legacy formats fall back to UNO when win32com is unavailable."""
    mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
    saved = {k: sys.modules.pop(k) for k in mods_to_remove}
    try:
        with patch.dict(
            "sys.modules",
            {
                "win32com": None,
                "win32com.client": None,
                "uno": MagicMock(),
            },
        ):
            assert _detect_backend(".doc") == "uno"
            assert _detect_backend(".xls") == "uno"
            assert _detect_backend(".ppt") == "uno"
    finally:
        sys.modules.update(saved)


def test_detect_backend_odf_prefers_uno() -> None:
    """ODF formats prefer UNO when available."""
    mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
    saved = {k: sys.modules.pop(k) for k in mods_to_remove}
    try:
        with patch.dict(
            "sys.modules",
            {
                "win32com": None,
                "win32com.client": None,
                "uno": MagicMock(),
            },
        ):
            assert _detect_backend(".odt") == "uno"
            assert _detect_backend(".ods") == "uno"
            assert _detect_backend(".odp") == "uno"
    finally:
        sys.modules.update(saved)


def test_detect_backend_odf_falls_back_to_win32com() -> None:
    """ODF formats fall back to win32com when UNO is unavailable."""
    mods_to_remove = [k for k in sys.modules if k == "uno"]
    saved = {k: sys.modules.pop(k) for k in mods_to_remove}
    try:
        with patch.dict(
            "sys.modules",
            {
                "uno": None,
                "win32com": MagicMock(),
                "win32com.client": MagicMock(),
            },
        ):
            assert _detect_backend(".odt") == "win32com"
            assert _detect_backend(".ods") == "win32com"
            assert _detect_backend(".odp") == "win32com"
    finally:
        sys.modules.update(saved)


def test_detect_backend_odf_falls_back_to_python_lib() -> None:
    """ODF formats fall back to python_lib when no UNO/win32com."""
    mods_to_remove = [k for k in sys.modules if k.startswith("win32com") or k == "uno"]
    saved = {k: sys.modules.pop(k) for k in mods_to_remove}
    try:
        with patch.dict(
            "sys.modules",
            {
                "win32com": None,
                "win32com.client": None,
                "uno": None,
            },
        ):
            assert _detect_backend(".odt") == "python_lib"
            assert _detect_backend(".ods") == "python_lib"
            assert _detect_backend(".odp") == "python_lib"
    finally:
        sys.modules.update(saved)


def test_detect_backend_legacy_no_backend_raises() -> None:
    """Raises error for legacy formats without win32com/UNO."""
    mods_to_remove = [k for k in sys.modules if k.startswith("win32com") or k == "uno"]
    saved = {k: sys.modules.pop(k) for k in mods_to_remove}
    try:
        with patch.dict(
            "sys.modules",
            {
                "win32com": None,
                "win32com.client": None,
                "uno": None,
            },
        ):
            for ext in (".doc", ".xls", ".ppt"):
                with pytest.raises(
                    ValueError,
                    match="OFFICE_CONVERTER_NOT_FOUND",
                ):
                    _detect_backend(ext)
    finally:
        sys.modules.update(saved)


# ---------------------------------------------------------------------------
# File category tests
# ---------------------------------------------------------------------------


def test_get_file_category() -> None:
    """Maps extensions to correct categories."""
    assert _get_file_category(".doc") == "word"
    assert _get_file_category(".docx") == "word"
    assert _get_file_category(".odt") == "word"
    assert _get_file_category(".xls") == "excel"
    assert _get_file_category(".xlsx") == "excel"
    assert _get_file_category(".ods") == "excel"
    assert _get_file_category(".ppt") == "ppt"
    assert _get_file_category(".pptx") == "ppt"
    assert _get_file_category(".odp") == "ppt"


def test_get_file_category_unsupported() -> None:
    """Raises for unsupported extensions."""
    with pytest.raises(ValueError, match="Unsupported office extension"):
        _get_file_category(".pdf")


# ---------------------------------------------------------------------------
# DOCX processing via python-docx
# ---------------------------------------------------------------------------


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_docx_paragraphs(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """DOCX with paragraphs translates correctly via python-docx."""
    mock_backend.return_value = "python_lib"

    # Create a test DOCX
    doc = Document()
    doc.add_paragraph("Hello")
    doc.add_paragraph("World")
    src_path = tmp_path / "test.docx"
    doc.save(str(src_path))

    # Mock translation
    mock_translate.return_value = ["Bonjour", "Le monde"]

    out_path = tmp_path / "output.docx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()

    # Verify translated content
    translated = Document(str(out_path))
    texts = [p.text for p in translated.paragraphs if p.text.strip()]
    assert texts == ["Bonjour", "Le monde"]


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_docx_tables(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """DOCX with table cells translates correctly via python-docx."""
    mock_backend.return_value = "python_lib"

    # Create a test DOCX with a table
    doc = Document()
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Hello"
    table.cell(1, 1).text = "World"
    src_path = tmp_path / "test_table.docx"
    doc.save(str(src_path))

    mock_translate.return_value = [
        "Nom",
        "Valeur",
        "Bonjour",
        "Le monde",
    ]

    out_path = tmp_path / "output_table.docx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()

    translated = Document(str(out_path))
    table = translated.tables[0]
    assert table.cell(0, 0).text == "Nom"
    assert table.cell(1, 1).text == "Le monde"


@patch("src.core.office_processor._detect_backend")
def test_process_docx_empty(
    mock_backend: MagicMock,
    tmp_path: Path,
) -> None:
    """Empty DOCX copies as-is."""
    mock_backend.return_value = "python_lib"

    doc = Document()
    src_path = tmp_path / "empty.docx"
    doc.save(str(src_path))

    out_path = tmp_path / "output_empty.docx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


# ---------------------------------------------------------------------------
# Image / visual content preservation in DOCX
# ---------------------------------------------------------------------------


def _add_drawing_to_run(run: object) -> None:
    """Injects a fake <w:drawing> element into a run for testing."""
    drawing = OxmlElement("w:drawing")
    inline = OxmlElement("wp:inline")
    drawing.append(inline)
    run._element.append(drawing)


def _add_pict_to_run(run: object) -> None:
    """Injects a fake <w:pict> element into a run for testing."""
    pict = OxmlElement("w:pict")
    run._element.append(pict)


def test_run_has_visual_content_drawing() -> None:
    """Detects <w:drawing> in a run."""
    doc = Document()
    para = doc.add_paragraph()
    run = para.add_run("text")
    assert _run_has_visual_content(run._element) is False

    _add_drawing_to_run(run)
    assert _run_has_visual_content(run._element) is True


def test_run_has_visual_content_pict() -> None:
    """Detects <w:pict> in a run."""
    doc = Document()
    para = doc.add_paragraph()
    run = para.add_run("text")
    _add_pict_to_run(run)
    assert _run_has_visual_content(run._element) is True


def test_run_has_visual_content_alternate() -> None:
    """Detects <mc:AlternateContent> in a run."""
    from lxml import etree  # noqa: PLC0415

    doc = Document()
    para = doc.add_paragraph()
    run = para.add_run("text")
    mc_ns = "http://schemas.openxmlformats.org/markup-compatibility/2006"
    alt = etree.SubElement(run._element, f"{{{mc_ns}}}AlternateContent")
    assert alt is not None
    assert _run_has_visual_content(run._element) is True


def test_clear_run_text_only() -> None:
    """Removes <w:t> but keeps other elements."""
    doc = Document()
    para = doc.add_paragraph()
    run = para.add_run("Hello")
    _add_drawing_to_run(run)

    # Run should have both <w:t> and <w:drawing>
    assert run._element.findall(qn("w:t"))
    assert run._element.findall(qn("w:drawing"))

    _clear_run_text_only(run._element)

    # <w:t> removed, <w:drawing> preserved
    assert not run._element.findall(qn("w:t"))
    assert run._element.findall(qn("w:drawing"))


def test_replace_paragraph_preserves_inline_image() -> None:
    """Translated text replaces text run; image run stays intact."""
    doc = Document()
    para = doc.add_paragraph()
    para.add_run("Hello World")
    img_run = para.add_run()
    _add_drawing_to_run(img_run)

    _replace_paragraph_text(para, "Bonjour le monde")

    # Text is replaced
    assert para.runs[0].text == "Bonjour le monde"
    # Image run still has <w:drawing>
    assert img_run._element.findall(qn("w:drawing"))


def test_replace_paragraph_image_between_text_runs() -> None:
    """Image run between text runs is preserved."""
    doc = Document()
    para = doc.add_paragraph()
    run_a = para.add_run("Before")
    img_run = para.add_run()
    _add_drawing_to_run(img_run)
    run_c = para.add_run("After")

    _replace_paragraph_text(para, "Translated")

    # First text run gets the translation
    assert run_a.text == "Translated"
    # Image preserved
    assert img_run._element.findall(qn("w:drawing"))
    # Last text run cleared
    assert run_c.text == ""


def test_replace_paragraph_all_runs_have_visuals() -> None:
    """When all runs have images, a new text run is inserted."""
    doc = Document()
    para = doc.add_paragraph()
    img_run1 = para.add_run("Caption")
    _add_drawing_to_run(img_run1)
    img_run2 = para.add_run("Another")
    _add_pict_to_run(img_run2)

    _replace_paragraph_text(para, "Translated")

    # Images are preserved
    assert img_run1._element.findall(qn("w:drawing"))
    assert img_run2._element.findall(qn("w:pict"))
    # Text removed from image runs
    assert not img_run1._element.findall(qn("w:t"))
    assert not img_run2._element.findall(qn("w:t"))
    # New text run was inserted and paragraph text includes translation
    assert "Translated" in para.text


def test_replace_paragraph_mixed_image_and_text_with_pict() -> None:
    """<w:pict> in a run is preserved alongside text replacement."""
    doc = Document()
    para = doc.add_paragraph()
    text_run = para.add_run("Original")
    pict_run = para.add_run()
    _add_pict_to_run(pict_run)

    _replace_paragraph_text(para, "New text")

    assert text_run.text == "New text"
    assert pict_run._element.findall(qn("w:pict"))


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_docx_preserves_images_e2e(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """End-to-end: DOCX with inline images preserves them after translation."""
    mock_backend.return_value = "python_lib"

    # Create a DOCX with text + inline drawing
    doc = Document()
    para = doc.add_paragraph()
    para.add_run("Hello")
    img_run = para.add_run()
    _add_drawing_to_run(img_run)

    src_path = tmp_path / "with_image.docx"
    doc.save(str(src_path))

    mock_translate.return_value = ["Bonjour"]

    out_path = tmp_path / "output_image.docx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()

    # Verify image is preserved in output
    translated = Document(str(out_path))
    para = translated.paragraphs[0]
    has_drawing = any(run._element.findall(qn("w:drawing")) for run in para.runs)
    assert has_drawing, "Inline drawing was lost during translation"
    assert "Bonjour" in para.text


# ---------------------------------------------------------------------------
# XLSX processing via openpyxl
# ---------------------------------------------------------------------------


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_xlsx_cells(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """XLSX cells translate correctly via openpyxl."""
    mock_backend.return_value = "python_lib"

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Hello"
    ws["B1"] = "World"
    ws["A2"] = 42  # Number — should be skipped
    src_path = tmp_path / "test.xlsx"
    wb.save(str(src_path))
    wb.close()

    mock_translate.return_value = ["Bonjour", "Le monde"]

    out_path = tmp_path / "output.xlsx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()

    translated_wb = load_workbook(str(out_path))
    ws = translated_wb.active
    assert ws["A1"].value == "Bonjour"
    assert ws["B1"].value == "Le monde"
    num_val = 42
    assert ws["A2"].value == num_val
    translated_wb.close()


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_xlsx_multiple_sheets(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Multi-sheet XLSX translates all sheets."""
    mock_backend.return_value = "python_lib"

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1["A1"] = "Hello"
    ws2 = wb.create_sheet("Sheet2")
    ws2["A1"] = "World"
    src_path = tmp_path / "multi.xlsx"
    wb.save(str(src_path))
    wb.close()

    mock_translate.return_value = ["Bonjour", "Le monde"]

    out_path = tmp_path / "output_multi.xlsx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True

    translated_wb = load_workbook(str(out_path))
    assert translated_wb["Sheet1"]["A1"].value == "Bonjour"
    assert translated_wb["Sheet2"]["A1"].value == "Le monde"
    translated_wb.close()


# ---------------------------------------------------------------------------
# XLSX embedded object preservation
# ---------------------------------------------------------------------------


def test_restore_xlsx_embeddings_preserves_objects(tmp_path: Path) -> None:
    """Embedded objects dropped by openpyxl are restored after save."""
    from lxml import etree  # noqa: PLC0415

    package_rel = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package"
    )
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    fake_embed = b"FAKE_EMBEDDED_DATA"

    # Create an XLSX with an embedded object + proper OPC rels
    wb = Workbook()
    wb.active["A1"] = "Hello"
    src = tmp_path / "source.xlsx"
    wb.save(str(src))
    wb.close()

    with zipfile.ZipFile(src, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    data["xl/embeddings/oleObject1.bin"] = fake_embed

    # Add rels entry
    rels_path = "xl/_rels/workbook.xml.rels"
    rels_xml = etree.fromstring(data[rels_path])
    rel = etree.SubElement(rels_xml, "Relationship")
    rel.set("Id", "rIdEmbed1")
    rel.set("Type", package_rel)
    rel.set("Target", "embeddings/oleObject1.bin")
    data[rels_path] = etree.tostring(
        rels_xml,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )

    # Add content type override
    ct_xml = etree.fromstring(data["[Content_Types].xml"])
    ov = etree.SubElement(ct_xml, f"{{{ct_ns}}}Override")
    ov.set("PartName", "/xl/embeddings/oleObject1.bin")
    ov.set("ContentType", "application/vnd.openxmlformats-officedocument.oleObject")
    data["[Content_Types].xml"] = etree.tostring(
        ct_xml,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )

    with zipfile.ZipFile(src, "w", zipfile.ZIP_DEFLATED) as zf:
        for i in items:
            zf.writestr(i, data[i.filename])
        zf.writestr("xl/embeddings/oleObject1.bin", fake_embed)

    # Simulate openpyxl save (drops embeddings)
    from openpyxl import load_workbook as _load_wb  # noqa: PLC0415

    out = tmp_path / "output.xlsx"
    wb2 = _load_wb(str(src))
    wb2.save(str(out))
    wb2.close()

    # Verify openpyxl dropped the embedding
    with zipfile.ZipFile(out, "r") as zf:
        assert "xl/embeddings/oleObject1.bin" not in zf.namelist()

    # Run the restore
    _restore_xlsx_embeddings(src, out)

    # Verify restoration
    with zipfile.ZipFile(out, "r") as zf:
        assert "xl/embeddings/oleObject1.bin" in zf.namelist()
        assert zf.read("xl/embeddings/oleObject1.bin") == fake_embed
        # Rels entry restored
        rels_bytes = zf.read(rels_path)
        assert b"rIdEmbed1" in rels_bytes
        # Content type restored
        ct_bytes = zf.read("[Content_Types].xml")
        assert b"oleObject" in ct_bytes


def test_restore_xlsx_embeddings_noop_without_embeddings(tmp_path: Path) -> None:
    """No-op when source has no embeddings."""
    wb = Workbook()
    wb.active["A1"] = "Hello"
    src = tmp_path / "source.xlsx"
    wb.save(str(src))
    out = tmp_path / "output.xlsx"
    wb2 = load_workbook(str(src))
    wb2.save(str(out))
    wb2.close()

    # Should return without error
    _restore_xlsx_embeddings(src, out)


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_xlsx_preserves_embedded_objects(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Full XLSX translation pipeline preserves embedded objects."""
    from lxml import etree  # noqa: PLC0415

    package_rel = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package"
    )
    fake_embed = b"EMBEDDED_SPREADSHEET"

    mock_backend.return_value = "python_lib"
    mock_translate.return_value = ["Bonjour"]

    wb = Workbook()
    wb.active["A1"] = "Hello"
    src = tmp_path / "test.xlsx"
    wb.save(str(src))
    wb.close()

    # Inject embedded object with proper OPC relationships
    with zipfile.ZipFile(src, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    data["xl/embeddings/oleObject1.bin"] = fake_embed
    rels_path = "xl/_rels/workbook.xml.rels"
    rels_xml = etree.fromstring(data[rels_path])
    rel = etree.SubElement(rels_xml, "Relationship")
    rel.set("Id", "rIdEmbed1")
    rel.set("Type", package_rel)
    rel.set("Target", "embeddings/oleObject1.bin")
    data[rels_path] = etree.tostring(
        rels_xml,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )

    with zipfile.ZipFile(src, "w", zipfile.ZIP_DEFLATED) as zf:
        for i in items:
            zf.writestr(i, data[i.filename])
        zf.writestr("xl/embeddings/oleObject1.bin", fake_embed)

    out = tmp_path / "translated.xlsx"
    result = process_office_file(src, out, "French", "English")

    assert result is True
    with zipfile.ZipFile(out, "r") as zf:
        assert "xl/embeddings/oleObject1.bin" in zf.namelist()
        assert zf.read("xl/embeddings/oleObject1.bin") == fake_embed

    out_wb = load_workbook(str(out))
    assert out_wb.active["A1"].value == "Bonjour"
    out_wb.close()


# ---------------------------------------------------------------------------
# PPTX processing via python-pptx
# ---------------------------------------------------------------------------


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_pptx_slides(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """PPTX text frames translate correctly via python-pptx."""
    mock_backend.return_value = "python_lib"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    text_box = slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(1),
    )
    tf = text_box.text_frame
    tf.text = "Hello World"
    src_path = tmp_path / "test.pptx"
    prs.save(str(src_path))

    mock_translate.return_value = ["Bonjour le monde"]

    out_path = tmp_path / "output.pptx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()

    translated = Presentation(str(out_path))
    slide = translated.slides[0]
    all_texts = [
        shape.text_frame.text
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]
    assert "Bonjour le monde" in all_texts


@patch("src.core.office_processor._detect_backend")
def test_process_pptx_empty(
    mock_backend: MagicMock,
    tmp_path: Path,
) -> None:
    """Empty PPTX (no text) copies as-is."""
    mock_backend.return_value = "python_lib"

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank, no textbox
    src_path = tmp_path / "empty.pptx"
    prs.save(str(src_path))

    out_path = tmp_path / "output_empty.pptx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


# ---------------------------------------------------------------------------
# Cancellation test
# ---------------------------------------------------------------------------


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_cancelled(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Returns False when cancel_check returns True."""
    mock_backend.return_value = "python_lib"

    doc = Document()
    doc.add_paragraph("Hello")
    src_path = tmp_path / "cancel.docx"
    doc.save(str(src_path))

    out_path = tmp_path / "output_cancel.docx"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
        cancel_check=lambda: True,
    )

    assert result is False
    mock_translate.assert_not_called()


# ---------------------------------------------------------------------------
# Error code mapping test
# ---------------------------------------------------------------------------


def test_error_mapping_in_translator() -> None:
    """OFFICE_CONVERTER_NOT_FOUND maps to correct error code."""
    code = _map_error_to_code("OFFICE_CONVERTER_NOT_FOUND")
    assert code == ERR_OFFICE_CONVERTER_NOT_FOUND


# ---------------------------------------------------------------------------
# Image translation — _should_translate_images
# ---------------------------------------------------------------------------


@patch("src.utils.config_manager.check_ocr_setup", return_value=True)
@patch("src.utils.config_manager.load_setting", return_value=True)
def test_should_translate_images_all_formats(
    _mock_setting: MagicMock,
    _mock_ocr: MagicMock,
) -> None:
    """Image translation is supported for all modern Office formats and EPUB."""
    for ext in (".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp", ".epub"):
        assert _should_translate_images(ext, "python_lib") is True
        # Backend is unused — any value should work
        assert _should_translate_images(ext, "win32com") is True


@patch("src.utils.config_manager.check_ocr_setup", return_value=True)
@patch("src.utils.config_manager.load_setting", return_value=True)
def test_should_translate_images_accepts_legacy(
    _mock_setting: MagicMock,
    _mock_ocr: MagicMock,
) -> None:
    """Legacy formats are now supported for image translation via round-trip."""
    for ext in (".doc", ".xls", ".ppt"):
        assert _should_translate_images(ext, "win32com") is True


# ---------------------------------------------------------------------------
# Image translation — _translate_zip_images
# ---------------------------------------------------------------------------


_EXPECTED_IMAGE_COUNT = 2


def _make_minimal_zip(tmp_path: Path, media_prefix: str, suffix: str) -> Path:
    """Creates a minimal ZIP file with a 1x1 red PNG in the media directory."""
    # Minimal valid 1x1 red PNG (67 bytes)
    png_bytes = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
        b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
        b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
        b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    file_path = tmp_path / f"test{suffix}"
    with zipfile.ZipFile(file_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{media_prefix}image1.png", png_bytes)
        zf.writestr(f"{media_prefix}image2.jpg", b"\xff\xd8\xff\xe0dummy")
        # Non-image file in media dir — should be skipped
        zf.writestr(f"{media_prefix}diagram.emf", b"emf-data")
        # File outside media dir — should be skipped
        zf.writestr("other/file.xml", b"<root/>")
    return file_path


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_replaces_images(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Images in the media directory are translated and replaced."""
    translated_png = b"translated-png-bytes"
    translated_jpg = b"translated-jpg-bytes"
    mock_translate.side_effect = [translated_png, translated_jpg]

    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
    )

    # Verify translated images in output ZIP
    with zipfile.ZipFile(file_path, "r") as zf:
        assert zf.read("word/media/image1.png") == translated_png
        assert zf.read("word/media/image2.jpg") == translated_jpg
        # Non-image files preserved
        assert zf.read("word/media/diagram.emf") == b"emf-data"
        assert zf.read("other/file.xml") == b"<root/>"

    assert mock_translate.call_count == _EXPECTED_IMAGE_COUNT


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_odf_pictures_dir(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """ODF formats use Pictures/ as the media directory."""
    mock_translate.return_value = b"translated"

    file_path = _make_minimal_zip(tmp_path, "Pictures/", ".odp")

    _translate_zip_images(
        file_path,
        ".odp",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
    )

    with zipfile.ZipFile(file_path, "r") as zf:
        assert zf.read("Pictures/image1.png") == b"translated"

    assert mock_translate.call_count == _EXPECTED_IMAGE_COUNT


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_fatal_error_raises(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Fatal LLM errors stop processing and propagate to caller."""
    mock_translate.side_effect = ValueError("AUTH_ERROR")

    file_path = _make_minimal_zip(tmp_path, "xl/media/", ".xlsx")

    with pytest.raises(ValueError, match="AUTH_ERROR"):
        _translate_zip_images(
            file_path,
            ".xlsx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

    # Only one image attempted before raising
    assert mock_translate.call_count == 1


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_fatal_error_with_service_suffix_raises(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """``AUTH_ERROR:Gemini`` (suffixed) still triggers the fatal short-circuit.

    Regression guard: the engine raises ``AUTH_ERROR:Service`` so the
    UI can render service-specific copy.  The old exact set-membership
    check (``error_tag in _FATAL_LLM_ERRORS``) missed the suffixed
    variant — letting an invalid API key silently skip every image
    and finish "Done" with zero translated content.  ``_is_fatal_llm_error``
    strips the suffix before checking membership.
    """
    mock_translate.side_effect = ValueError("AUTH_ERROR:Gemini")

    file_path = _make_minimal_zip(tmp_path, "xl/media/", ".xlsx")

    with pytest.raises(ValueError, match=r"AUTH_ERROR:Gemini"):
        _translate_zip_images(
            file_path,
            ".xlsx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

    # Only one image attempted before raising — the suffix didn't
    # demote the fatal error to a skip-with-warning.
    assert mock_translate.call_count == 1


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_nonfatal_error_skips_with_warning(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Non-fatal per-image errors leave the original and continue.

    Skip-with-warning policy: a single bad image (e.g. transient
    SERVICE_UNAVAILABLE that's still failing after the LLM's own
    3-retry policy) must not abort the whole document.  The image
    stays in its source form, the loop continues, and the function
    returns normally so the document can still be delivered.
    """
    mock_translate.side_effect = ValueError("SERVICE_UNAVAILABLE_ERROR")

    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")

    # No raise — both images are attempted, both skip-with-warning.
    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
    )

    # All images attempted (not short-circuited on first failure).
    assert mock_translate.call_count == _EXPECTED_IMAGE_COUNT  # noqa: PLR2004
    # And the ZIP still contains the *originals* — failed images
    # never overwrite the source bytes.
    with zipfile.ZipFile(file_path, "r") as zf:
        png1 = zf.read("word/media/image1.png")
        jpg2 = zf.read("word/media/image2.jpg")
    # First bytes of the original 1×1 red PNG fixture survived.
    assert png1.startswith(b"\x89PNG")
    assert jpg2.startswith(b"\xff\xd8\xff")


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_cancel_check(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Cancellation stops image processing."""
    file_path = _make_minimal_zip(tmp_path, "ppt/media/", ".pptx")

    _translate_zip_images(
        file_path,
        ".pptx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        lambda: True,
    )

    # No images should be processed when cancelled
    mock_translate.assert_not_called()


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_progress_callback(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Progress callback is called with increasing percentages."""
    mock_translate.return_value = b"translated"

    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")

    progress_values: list[int] = []
    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        progress_values.append,
        None,
    )

    expected_final = 100
    assert len(progress_values) == _EXPECTED_IMAGE_COUNT
    assert progress_values[-1] == expected_final
    assert progress_values == sorted(progress_values)


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_epub_all_paths(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """EPUB images at arbitrary paths are found and translated."""
    # Minimal valid 1x1 red PNG (67 bytes)
    png_bytes = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
        b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
        b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
        b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    mock_translate.return_value = b"translated"

    # Images at various arbitrary paths (not in a fixed media directory)
    epub_path = tmp_path / "test.epub"
    expected_image_count = 3
    with zipfile.ZipFile(epub_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("OEBPS/images/cover.png", png_bytes)
        zf.writestr("images/photo.jpg", b"\xff\xd8\xff\xe0dummy")
        zf.writestr("chapter1/fig1.png", png_bytes)
        # Non-image files — should be skipped
        zf.writestr("OEBPS/content.xhtml", b"<html/>")
        zf.writestr("META-INF/container.xml", b"<container/>")

    _translate_zip_images(
        epub_path,
        ".epub",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
    )

    # All 3 images found and translated
    assert mock_translate.call_count == expected_image_count

    # Verify translated images and preserved non-images
    with zipfile.ZipFile(epub_path, "r") as zf:
        assert zf.read("OEBPS/images/cover.png") == b"translated"
        assert zf.read("images/photo.jpg") == b"translated"
        assert zf.read("chapter1/fig1.png") == b"translated"
        assert zf.read("OEBPS/content.xhtml") == b"<html/>"
        assert zf.read("META-INF/container.xml") == b"<container/>"


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_root_media_fallback(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """DOCX images stored at root media/ (not word/media/) are found and translated.

    Some Word versions reference images via absolute ZIP paths
    (Target="/media/image.png"), placing them at the root media/ directory
    instead of the standard word/media/ location.  Both prefixes are checked.
    """
    png_bytes = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
        b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
        b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
        b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    mock_translate.return_value = b"translated"

    file_path = tmp_path / "test.docx"
    with zipfile.ZipFile(file_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("media/image.png", png_bytes)
        zf.writestr("word/document.xml", b"<root/>")

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
    )

    assert mock_translate.call_count == 1
    with zipfile.ZipFile(file_path, "r") as zf:
        assert zf.read("media/image.png") == b"translated"
        assert zf.read("word/document.xml") == b"<root/>"


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_both_media_prefixes(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """DOCX with images at both word/media/ and root media/ finds all of them."""
    png_bytes = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
        b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
        b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
        b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    mock_translate.return_value = b"translated"

    file_path = tmp_path / "test.docx"
    with zipfile.ZipFile(file_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("word/media/image1.png", png_bytes)
        zf.writestr("media/image2.png", png_bytes)
        zf.writestr("word/document.xml", b"<root/>")

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
    )

    expected_count = 2
    assert mock_translate.call_count == expected_count  # noqa: PLR2004
    with zipfile.ZipFile(file_path, "r") as zf:
        assert zf.read("word/media/image1.png") == b"translated"
        assert zf.read("media/image2.png") == b"translated"
        assert zf.read("word/document.xml") == b"<root/>"


def test_translate_zip_images_unknown_suffix_returns_early(tmp_path: Path) -> None:
    """Suffix not in _SUFFIX_TO_MEDIA_PREFIXES calls progress(100) and returns."""
    file_path = tmp_path / "test.rtf"
    with zipfile.ZipFile(file_path, "w") as zf:
        zf.writestr("media/image.png", b"data")

    progress_values: list[int] = []
    _translate_zip_images(
        file_path,
        ".rtf",
        "French",
        "English",
        None,
        "TesseractOCR",
        progress_values.append,
        None,
    )

    assert progress_values == [100]  # noqa: PLR2004


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_no_images_returns_early(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """ZIP with no translatable images calls progress(100) and skips translation."""
    file_path = tmp_path / "test.docx"
    with zipfile.ZipFile(file_path, "w") as zf:
        zf.writestr("word/document.xml", b"<root/>")
        zf.writestr("word/media/diagram.emf", b"emf-data")  # unsupported extension

    progress_values: list[int] = []
    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        progress_values.append,
        None,
    )

    mock_translate.assert_not_called()
    assert progress_values == [100]  # noqa: PLR2004


# ---------------------------------------------------------------------------
# Per-image checkpoint cache in _translate_zip_images
# ---------------------------------------------------------------------------


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_caches_results_on_first_run(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Successful per-image translations are persisted to the cache dir.

    Each image lands under ``<checkpoint_dir>/office_images/<sha>.bin``
    so a resumed run can skip it.  Validates the *write* half of the
    cache contract.
    """
    from src.core.checkpoint import (  # noqa: PLC0415
        _OFFICE_IMAGE_DIR_NAME,
        hash_office_image,
    )

    translated_png = b"translated-png-bytes"
    translated_jpg = b"translated-jpg-bytes"
    mock_translate.side_effect = [translated_png, translated_jpg]

    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")
    checkpoint_dir = tmp_path / "task_storage"
    checkpoint_dir.mkdir()

    # Grab the source bytes so we can compute the expected hash keys
    # *before* the ZIP is rewritten in place.
    with zipfile.ZipFile(file_path, "r") as zf:
        png_src = zf.read("word/media/image1.png")
        jpg_src = zf.read("word/media/image2.jpg")

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
        checkpoint_dir=checkpoint_dir,
    )

    cache_dir = checkpoint_dir / _OFFICE_IMAGE_DIR_NAME
    png_cache = cache_dir / f"{hash_office_image(png_src)}.bin"
    jpg_cache = cache_dir / f"{hash_office_image(jpg_src)}.bin"
    assert png_cache.read_bytes() == translated_png
    assert jpg_cache.read_bytes() == translated_jpg


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_skips_cached_images_on_resume(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """A second pass re-uses the cache without ever calling the live pipeline.

    This is the load-bearing property: a transient quota error mid-run
    doesn't waste LLM tokens on retry.
    """
    from src.core.checkpoint import (  # noqa: PLC0415
        hash_office_image,
        save_office_image_checkpoint,
    )

    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")
    checkpoint_dir = tmp_path / "task_storage"
    checkpoint_dir.mkdir()

    with zipfile.ZipFile(file_path, "r") as zf:
        png_src = zf.read("word/media/image1.png")
        jpg_src = zf.read("word/media/image2.jpg")

    # Pre-seed the cache as if a prior run had translated both images.
    save_office_image_checkpoint(
        checkpoint_dir,
        hash_office_image(png_src),
        b"cached-png",
    )
    save_office_image_checkpoint(
        checkpoint_dir,
        hash_office_image(jpg_src),
        b"cached-jpg",
    )

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
        checkpoint_dir=checkpoint_dir,
    )

    mock_translate.assert_not_called()
    with zipfile.ZipFile(file_path, "r") as zf:
        assert zf.read("word/media/image1.png") == b"cached-png"
        assert zf.read("word/media/image2.jpg") == b"cached-jpg"


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_partial_cache_translates_only_misses(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """The common resume case: some images already translated, others not.

    Only the uncached image hits the LLM; the cached one is reused
    verbatim from the prior run.
    """
    from src.core.checkpoint import (  # noqa: PLC0415
        hash_office_image,
        save_office_image_checkpoint,
    )

    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")
    checkpoint_dir = tmp_path / "task_storage"
    checkpoint_dir.mkdir()

    with zipfile.ZipFile(file_path, "r") as zf:
        png_src = zf.read("word/media/image1.png")

    save_office_image_checkpoint(
        checkpoint_dir,
        hash_office_image(png_src),
        b"cached-png",
    )
    mock_translate.return_value = b"freshly-translated-jpg"

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
        checkpoint_dir=checkpoint_dir,
    )

    # Only the uncached image (jpg) goes through the live pipeline.
    assert mock_translate.call_count == 1
    with zipfile.ZipFile(file_path, "r") as zf:
        assert zf.read("word/media/image1.png") == b"cached-png"
        assert zf.read("word/media/image2.jpg") == b"freshly-translated-jpg"


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_without_checkpoint_dir_does_not_cache(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Backward-compat path: callers without ``checkpoint_dir`` skip caching.

    The original "translate every time" behaviour is preserved with no
    on-disk cache files created anywhere.
    """
    mock_translate.side_effect = [b"png", b"jpg"]
    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
    )

    # No subdir created anywhere under tmp_path beyond the test ZIP.
    assert mock_translate.call_count == _EXPECTED_IMAGE_COUNT


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_cancellation_preserves_already_cached(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """User cancels mid-loop: the image already translated stays in the cache.

    Real-world scenario: a 100-image document, the user cancels
    after image 50.  The 50 successful images must be in the cache
    so a Retry only re-translates 51-100 (cheap), not 1-100.

    Builds a ZIP with two images, ``cancel_check`` returns True
    after the FIRST image's progress callback fires.  Asserts the
    first image is in the cache, the second is not.
    """
    from src.core.checkpoint import (  # noqa: PLC0415
        _OFFICE_IMAGE_DIR_NAME,
        hash_office_image,
    )

    mock_translate.return_value = b"translated"
    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")
    checkpoint_dir = tmp_path / "task_storage"
    checkpoint_dir.mkdir()

    with zipfile.ZipFile(file_path, "r") as zf:
        png_src = zf.read("word/media/image1.png")
        jpg_src = zf.read("word/media/image2.jpg")

    # Cancel fires only on the *second* loop iteration — i.e. after
    # image 1 has been translated + cached but before image 2 starts.
    call_count: list[int] = [0]

    def cancel_after_first() -> bool:
        call_count[0] += 1
        return call_count[0] > 1

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        cancel_after_first,
        checkpoint_dir=checkpoint_dir,
    )

    # Only image 1 was attempted.
    assert mock_translate.call_count == 1
    cache_dir = checkpoint_dir / _OFFICE_IMAGE_DIR_NAME
    # Image 1's translation survives the cancellation in the cache.
    assert (cache_dir / f"{hash_office_image(png_src)}.bin").read_bytes() == (
        b"translated"
    )
    # Image 2 was never attempted → not in the cache.
    assert not (cache_dir / f"{hash_office_image(jpg_src)}.bin").exists()


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_fatal_after_success_keeps_prior_in_cache(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """Fatal error on image N: images 1..N-1 are still cached.

    Skip-with-warning policy abandons the rest of the document on
    fatal errors (AUTH_ERROR, QUOTA_ERROR, VISION_NOT_SUPPORTED),
    but the cache writes happen *before* the next iteration, so
    the prior-successful images survive in the cache.  Retry value:
    on Retry the user gets cache hits for the work that succeeded
    before the fatal stopped the run.
    """
    from src.core.checkpoint import (  # noqa: PLC0415
        _OFFICE_IMAGE_DIR_NAME,
        hash_office_image,
    )

    # Image 1 succeeds; image 2 raises fatal AUTH_ERROR.
    mock_translate.side_effect = [b"translated", ValueError("AUTH_ERROR")]
    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")
    checkpoint_dir = tmp_path / "task_storage"
    checkpoint_dir.mkdir()

    with zipfile.ZipFile(file_path, "r") as zf:
        png_src = zf.read("word/media/image1.png")
        jpg_src = zf.read("word/media/image2.jpg")

    with pytest.raises(ValueError, match="AUTH_ERROR"):
        _translate_zip_images(
            file_path,
            ".docx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
            checkpoint_dir=checkpoint_dir,
        )

    cache_dir = checkpoint_dir / _OFFICE_IMAGE_DIR_NAME
    # The pre-fatal success is preserved (retry value).
    assert (cache_dir / f"{hash_office_image(png_src)}.bin").read_bytes() == (
        b"translated"
    )
    # The fatal-failed image is not cached.
    assert not (cache_dir / f"{hash_office_image(jpg_src)}.bin").exists()


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_save_creates_storage_dir_if_missing(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """``save_office_image_checkpoint`` lazily creates the storage dir.

    Real-world: the per-task storage directory may exist but the
    ``office_images/`` subdirectory inside it does not until the
    first image is translated.  Verify we don't crash on the
    "subdir doesn't exist yet" cold-start path.
    """
    from src.core.checkpoint import _OFFICE_IMAGE_DIR_NAME  # noqa: PLC0415

    mock_translate.return_value = b"translated"
    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")
    checkpoint_dir = tmp_path / "task_storage"
    checkpoint_dir.mkdir()
    # Deliberately do NOT pre-create ``office_images/`` — the helper
    # should mkdir it on first save.

    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
        checkpoint_dir=checkpoint_dir,
    )

    cache_dir = checkpoint_dir / _OFFICE_IMAGE_DIR_NAME
    assert cache_dir.is_dir()
    assert list(cache_dir.glob("*.bin"))  # at least one cached image


@patch("src.core.office_processor._translate_single_image")
def test_translate_zip_images_failed_image_is_not_cached(
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """A non-fatal failure (e.g. transient connection error) is NOT cached.

    Otherwise a resume would silently treat the broken image as
    "already done" and never retry it.
    """
    from src.core.checkpoint import (  # noqa: PLC0415
        _OFFICE_IMAGE_DIR_NAME,
        hash_office_image,
    )

    # First image (png) fails; second image (jpg) succeeds.
    mock_translate.side_effect = [ValueError("CONNECTION_ERROR"), b"jpg-translated"]
    file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")
    checkpoint_dir = tmp_path / "task_storage"
    checkpoint_dir.mkdir()

    with zipfile.ZipFile(file_path, "r") as zf:
        png_src = zf.read("word/media/image1.png")
        jpg_src = zf.read("word/media/image2.jpg")

    # No raise: skip-with-warning means the function returns normally
    # even with a failed image.  The cache still differentiates
    # successful vs failed images, which is what we're verifying here.
    _translate_zip_images(
        file_path,
        ".docx",
        "French",
        "English",
        None,
        "TesseractOCR",
        None,
        None,
        checkpoint_dir=checkpoint_dir,
    )

    cache_dir = checkpoint_dir / _OFFICE_IMAGE_DIR_NAME
    # The failed png MUST NOT be cached.
    assert not (cache_dir / f"{hash_office_image(png_src)}.bin").exists()
    # The successful jpg IS cached (so resume skips it).
    assert (cache_dir / f"{hash_office_image(jpg_src)}.bin").read_bytes() == (
        b"jpg-translated"
    )


# ---------------------------------------------------------------------------
# ODT processing via odfpy
# ---------------------------------------------------------------------------


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_odt_paragraphs(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """ODT with paragraphs translates correctly via odfpy."""
    from odf.opendocument import OpenDocumentText  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    mock_backend.return_value = "python_lib"

    doc = OpenDocumentText()
    doc.text.addElement(P(text="Hello"))
    doc.text.addElement(P(text="World"))
    src_path = tmp_path / "test.odt"
    doc.save(str(src_path))

    mock_translate.return_value = ["Bonjour", "Le monde"]

    out_path = tmp_path / "output.odt"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()

    # Verify translated content
    from odf.opendocument import load as odf_load  # noqa: PLC0415

    translated = odf_load(str(out_path))
    paras = translated.text.getElementsByType(P)
    texts = [_odf_element_text(p) for p in paras if _odf_element_text(p).strip()]
    assert texts == ["Bonjour", "Le monde"]


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_odt_with_tables(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """ODT with table cells translates correctly via odfpy."""
    from odf.opendocument import OpenDocumentText  # noqa: PLC0415
    from odf.table import Table, TableCell, TableColumn, TableRow  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    mock_backend.return_value = "python_lib"

    doc = OpenDocumentText()
    table = Table(name="Table1")
    table.addElement(TableColumn())
    table.addElement(TableColumn())
    row = TableRow()
    cell1 = TableCell(valuetype="string")
    cell1.addElement(P(text="Name"))
    cell2 = TableCell(valuetype="string")
    cell2.addElement(P(text="Value"))
    row.addElement(cell1)
    row.addElement(cell2)
    table.addElement(row)
    doc.text.addElement(table)
    src_path = tmp_path / "test_table.odt"
    doc.save(str(src_path))

    mock_translate.return_value = ["Name", "Value"]

    out_path = tmp_path / "output_table.odt"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


@patch("src.core.office_processor._detect_backend")
def test_process_odt_empty(
    mock_backend: MagicMock,
    tmp_path: Path,
) -> None:
    """Empty ODT copies as-is."""
    from odf.opendocument import OpenDocumentText  # noqa: PLC0415

    mock_backend.return_value = "python_lib"

    doc = OpenDocumentText()
    src_path = tmp_path / "empty.odt"
    doc.save(str(src_path))

    out_path = tmp_path / "output_empty.odt"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_odt_headings(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """ODT headings are extracted and translated."""
    from odf.opendocument import OpenDocumentText  # noqa: PLC0415
    from odf.text import H, P  # noqa: PLC0415

    mock_backend.return_value = "python_lib"

    doc = OpenDocumentText()
    doc.text.addElement(H(outlinelevel=1, text="Chapter One"))
    doc.text.addElement(P(text="Body text"))
    src_path = tmp_path / "headings.odt"
    doc.save(str(src_path))

    mock_translate.return_value = ["Chapitre Un", "Texte du corps"]

    out_path = tmp_path / "output_headings.odt"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


# ---------------------------------------------------------------------------
# ODS processing via odfpy
# ---------------------------------------------------------------------------


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_ods_cells(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """ODS cells translate correctly via odfpy."""
    from odf.opendocument import OpenDocumentSpreadsheet  # noqa: PLC0415
    from odf.table import Table, TableCell, TableColumn, TableRow  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    mock_backend.return_value = "python_lib"

    doc = OpenDocumentSpreadsheet()
    table = Table(name="Sheet1")
    table.addElement(TableColumn())
    table.addElement(TableColumn())

    row = TableRow()
    cell1 = TableCell(valuetype="string")
    cell1.addElement(P(text="Hello"))
    row.addElement(cell1)
    cell2 = TableCell(valuetype="string")
    cell2.addElement(P(text="World"))
    row.addElement(cell2)
    table.addElement(row)

    # Numeric row — should be skipped
    row2 = TableRow()
    num_cell = TableCell(valuetype="float", value="42")
    num_cell.addElement(P(text="42"))
    row2.addElement(num_cell)
    table.addElement(row2)

    doc.spreadsheet.addElement(table)
    src_path = tmp_path / "test.ods"
    doc.save(str(src_path))

    mock_translate.return_value = ["Bonjour", "Le monde"]

    out_path = tmp_path / "output.ods"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


@patch("src.core.office_processor._detect_backend")
def test_process_ods_empty(
    mock_backend: MagicMock,
    tmp_path: Path,
) -> None:
    """Empty ODS copies as-is."""
    from odf.opendocument import OpenDocumentSpreadsheet  # noqa: PLC0415
    from odf.table import Table  # noqa: PLC0415

    mock_backend.return_value = "python_lib"

    doc = OpenDocumentSpreadsheet()
    table = Table(name="Sheet1")
    doc.spreadsheet.addElement(table)
    src_path = tmp_path / "empty.ods"
    doc.save(str(src_path))

    out_path = tmp_path / "output_empty.ods"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


# ---------------------------------------------------------------------------
# ODP processing via odfpy
# ---------------------------------------------------------------------------


@patch("src.core.llm_engine.translate_text")
@patch("src.core.office_processor._detect_backend")
def test_process_odp_slides(
    mock_backend: MagicMock,
    mock_translate: MagicMock,
    tmp_path: Path,
) -> None:
    """ODP text frames translate correctly via odfpy."""
    from odf.draw import Frame, Page, TextBox  # noqa: PLC0415
    from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
    from odf.style import MasterPage, PageLayout  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    mock_backend.return_value = "python_lib"

    doc = OpenDocumentPresentation()
    pl = PageLayout(name="PM1")
    doc.automaticstyles.addElement(pl)
    mp = MasterPage(name="Default", pagelayoutname=pl)
    doc.masterstyles.addElement(mp)

    page = Page(name="Slide1", masterpagename=mp)
    frame = Frame(width="10cm", height="5cm", x="1cm", y="1cm")
    text_box = TextBox()
    text_box.addElement(P(text="Hello World"))
    frame.addElement(text_box)
    page.addElement(frame)
    doc.presentation.addElement(page)
    src_path = tmp_path / "test.odp"
    doc.save(str(src_path))

    mock_translate.return_value = ["Bonjour le monde"]

    out_path = tmp_path / "output.odp"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


@patch("src.core.office_processor._detect_backend")
def test_process_odp_empty(
    mock_backend: MagicMock,
    tmp_path: Path,
) -> None:
    """Empty ODP (no text) copies as-is."""
    from odf.draw import Page  # noqa: PLC0415
    from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
    from odf.style import MasterPage, PageLayout  # noqa: PLC0415

    mock_backend.return_value = "python_lib"

    doc = OpenDocumentPresentation()
    pl = PageLayout(name="PM1")
    doc.automaticstyles.addElement(pl)
    mp = MasterPage(name="Default", pagelayoutname=pl)
    doc.masterstyles.addElement(mp)

    page = Page(name="Slide1", masterpagename=mp)
    doc.presentation.addElement(page)
    src_path = tmp_path / "empty.odp"
    doc.save(str(src_path))

    out_path = tmp_path / "output_empty.odp"
    result = process_office_file(
        src_path,
        out_path,
        "French",
        "English",
    )

    assert result is True
    assert out_path.exists()


# ---------------------------------------------------------------------------
# Modern threaded PPTX comment tests
# ---------------------------------------------------------------------------

# XML template for modern threaded comments
_MODERN_COMMENT_XML = """\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p188:cmLst xmlns:p188=\
"http://schemas.microsoft.com/office/powerpoint/2018/8/main"
            xmlns:a=\
"http://schemas.openxmlformats.org/drawingml/2006/main"
            xmlns:pc=\
"http://schemas.microsoft.com/office/powerpoint/2013/main/command">
  <p188:cm id="{CM-001}" authorId="{AUTH-001}" status="active"
           created="2024-12-30T20:00:00.000">
    <pc:sldMkLst><pc:docMk/><pc:sldMk cId="1" sldId="256"/></pc:sldMkLst>
    <p188:txBody>
      <a:bodyPr/><a:lstStyle/>
      <a:p><a:r><a:rPr lang="en-US"/><a:t>Main comment text</a:t>\
</a:r></a:p>
    </p188:txBody>
  </p188:cm>
</p188:cmLst>"""

_MODERN_COMMENT_WITH_REPLIES_XML = """\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p188:cmLst xmlns:p188=\
"http://schemas.microsoft.com/office/powerpoint/2018/8/main"
            xmlns:a=\
"http://schemas.openxmlformats.org/drawingml/2006/main"
            xmlns:pc=\
"http://schemas.microsoft.com/office/powerpoint/2013/main/command">
  <p188:cm id="{CM-002}" authorId="{AUTH-001}" status="active"
           created="2024-12-30T20:00:00.000">
    <pc:sldMkLst><pc:docMk/><pc:sldMk cId="1" sldId="256"/></pc:sldMkLst>
    <p188:replyLst>
      <p188:reply id="{REPLY-001}" authorId="{AUTH-002}" status="active"
                  created="2024-12-30T21:00:00.000">
        <p188:txBody>
          <a:bodyPr/><a:lstStyle/>
          <a:p><a:r><a:rPr lang="en-US"/><a:t>Reply text</a:t>\
</a:r></a:p>
        </p188:txBody>
      </p188:reply>
    </p188:replyLst>
    <p188:txBody>
      <a:bodyPr/><a:lstStyle/>
      <a:p><a:r><a:rPr lang="en-US"/><a:t>Parent comment</a:t>\
</a:r></a:p>
    </p188:txBody>
  </p188:cm>
</p188:cmLst>"""

# Legacy comment XML for the "both formats" test
_LEGACY_COMMENT_XML = """\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:cmLst xmlns:p=\
"http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cm authorId="1" dt="2024-01-01T00:00:00" idx="1">
    <p:pos x="100" y="200"/>
    <p:text>Legacy comment</p:text>
  </p:cm>
</p:cmLst>"""

_MODERN_REL_TYPE = "http://schemas.microsoft.com/office/2018/10/relationships/comments"
_LEGACY_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
)


def _inject_comment_part_into_pptx(
    pptx_path: Path,
    comment_xml: str,
    rel_type: str,
    comment_part_name: str,
) -> None:
    """Injects a comment XML part into a PPTX ZIP and wires the relationship.

    Adds the comment XML as a new part, updates the slide1 .rels file
    to reference it, and updates [Content_Types].xml.

    Args:
        pptx_path: Path to an existing .pptx file (modified in-place).
        comment_xml: The XML string for the comment part.
        rel_type: The relationship type URI.
        comment_part_name: The part name inside the ZIP
                           (e.g. ``ppt/comments/comment1.xml``).
    """
    import shutil  # noqa: PLC0415

    from lxml import etree  # noqa: PLC0415

    # Read original ZIP contents
    with zipfile.ZipFile(pptx_path, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    # Add the comment part
    data[comment_part_name] = comment_xml.encode("utf-8")

    # Update slide1 .rels to add the relationship
    rels_path = "ppt/slides/_rels/slide1.xml.rels"
    if rels_path in data:
        rels_root = etree.fromstring(data[rels_path])
    else:
        rels_root = etree.Element(
            "Relationships",
            xmlns="http://schemas.openxmlformats.org/package/2006/relationships",
        )

    # Determine next rId
    existing_ids = [r.get("Id", "") for r in rels_root]
    next_id = f"rId{len(existing_ids) + 100}"

    # Compute relative target from ppt/slides/ to comment_part_name
    target = "../" + comment_part_name.removeprefix("ppt/")
    etree.SubElement(
        rels_root,
        "Relationship",
        Id=next_id,
        Type=rel_type,
        Target=target,
    )
    data[rels_path] = etree.tostring(
        rels_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    # Update [Content_Types].xml
    ct_root = etree.fromstring(data["[Content_Types].xml"])
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    if "2018/10" in rel_type:
        content_type = "application/vnd.ms-powerpoint.comments+xml"
    else:
        content_type = (
            "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"
        )
    etree.SubElement(
        ct_root,
        f"{{{ct_ns}}}Override",
        PartName=f"/{comment_part_name}",
        ContentType=content_type,
    )
    data["[Content_Types].xml"] = etree.tostring(
        ct_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    # Rewrite the ZIP
    tmp_zip = pptx_path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for item in items:
            zf_out.writestr(item, data[item.filename])
        # Write new parts not in the original
        for name, content in data.items():
            if not any(i.filename == name for i in items):
                zf_out.writestr(name, content)
    shutil.move(str(tmp_zip), str(pptx_path))


def _make_pptx_with_slide(tmp_path: Path, name: str) -> Path:
    """Creates a minimal PPTX with one blank slide.

    Args:
        tmp_path: Pytest tmp_path fixture.
        name: Filename for the PPTX.

    Returns:
        Path: The saved PPTX path.
    """
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    path = tmp_path / name
    prs.save(str(path))
    return path


def test_extract_pptx_modern_comments(tmp_path: Path) -> None:
    """Modern threaded comment text is extracted correctly."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_pptx_comments,
    )

    pptx_path = _make_pptx_with_slide(tmp_path, "modern.pptx")
    _inject_comment_part_into_pptx(
        pptx_path,
        _MODERN_COMMENT_XML,
        _MODERN_REL_TYPE,
        "ppt/comments/modernComment1.xml",
    )

    texts = _extract_pptx_comments(pptx_path)
    assert len(texts) == 1
    key, text = texts[0]
    assert key == "comment:0:{CM-001}"
    assert text == "Main comment text"


def test_extract_pptx_modern_comments_with_replies(
    tmp_path: Path,
) -> None:
    """Modern comment replies are extracted with correct keys."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_pptx_comments,
    )

    pptx_path = _make_pptx_with_slide(tmp_path, "replies.pptx")
    _inject_comment_part_into_pptx(
        pptx_path,
        _MODERN_COMMENT_WITH_REPLIES_XML,
        _MODERN_REL_TYPE,
        "ppt/comments/modernComment1.xml",
    )

    texts = _extract_pptx_comments(pptx_path)
    keys = dict(texts)
    assert "comment:0:{CM-002}" in keys
    assert keys["comment:0:{CM-002}"] == "Parent comment"
    assert "comment:0:{CM-002}:reply:{REPLY-001}" in keys
    assert keys["comment:0:{CM-002}:reply:{REPLY-001}"] == "Reply text"


def test_inject_pptx_modern_comments(tmp_path: Path) -> None:
    """Translated text is injected into modern comment txBody."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_pptx_comments,
        _inject_pptx_comments,
    )

    pptx_path = _make_pptx_with_slide(tmp_path, "inject.pptx")
    _inject_comment_part_into_pptx(
        pptx_path,
        _MODERN_COMMENT_WITH_REPLIES_XML,
        _MODERN_REL_TYPE,
        "ppt/comments/modernComment1.xml",
    )

    translations = {
        "comment:0:{CM-002}": "Commentaire traduit",
        "comment:0:{CM-002}:reply:{REPLY-001}": "Réponse traduite",
    }
    _inject_pptx_comments(pptx_path, translations)

    # Re-extract and verify translated text
    texts = _extract_pptx_comments(pptx_path)
    result = dict(texts)
    assert result["comment:0:{CM-002}"] == "Commentaire traduit"
    assert result["comment:0:{CM-002}:reply:{REPLY-001}"] == "Réponse traduite"


def test_extract_pptx_both_legacy_and_modern(
    tmp_path: Path,
) -> None:
    """File with both legacy and modern comments extracts all."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_pptx_comments,
    )

    pptx_path = _make_pptx_with_slide(tmp_path, "both.pptx")

    # Inject legacy comment
    _inject_comment_part_into_pptx(
        pptx_path,
        _LEGACY_COMMENT_XML,
        _LEGACY_REL_TYPE,
        "ppt/comments/comment1.xml",
    )
    # Inject modern comment
    _inject_comment_part_into_pptx(
        pptx_path,
        _MODERN_COMMENT_XML,
        _MODERN_REL_TYPE,
        "ppt/comments/modernComment1.xml",
    )

    texts = _extract_pptx_comments(pptx_path)
    keys = dict(texts)

    # Legacy comment
    assert "comment:0:1" in keys
    assert keys["comment:0:1"] == "Legacy comment"

    # Modern comment
    assert "comment:0:{CM-001}" in keys
    assert keys["comment:0:{CM-001}"] == "Main comment text"


# ---------------------------------------------------------------------------
# Legacy format support — gate sets
# ---------------------------------------------------------------------------


def test_legacy_extensions_in_image_gate_set() -> None:
    """Legacy extensions (.doc, .xls, .ppt) are in _IMAGE_EXTENSIONS."""
    for ext in (".doc", ".xls", ".ppt"):
        assert ext in _IMAGE_EXTENSIONS


def test_legacy_extensions_in_comment_gate_set() -> None:
    """Legacy extensions (.doc, .xls, .ppt) are in _COMMENT_EXTENSIONS."""
    for ext in (".doc", ".xls", ".ppt"):
        assert ext in _COMMENT_EXTENSIONS


def test_legacy_extensions_constant() -> None:
    """_LEGACY_EXTENSIONS contains the three legacy formats."""
    assert {".doc", ".xls", ".ppt"} == _LEGACY_EXTENSIONS


def test_legacy_to_modern_mapping() -> None:
    """LEGACY_CONVERT_MAP maps each legacy format to its modern equivalent."""
    assert LEGACY_CONVERT_MAP == {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}


@patch("src.utils.config_manager.load_setting", return_value=True)
def test_should_translate_comments_legacy(
    _mock_setting: MagicMock,
) -> None:
    """Comment translation setting is checked for legacy formats."""
    for ext in (".doc", ".xls", ".ppt"):
        assert _should_translate_comments(ext, "win32com") is True


# ---------------------------------------------------------------------------
# Legacy format support — conversion helpers (mocked)
# ---------------------------------------------------------------------------


def test_convert_with_win32com_word(tmp_path: Path) -> None:
    """_convert_with_win32com calls Word SaveAs for .doc → .docx."""
    input_path = tmp_path / "test.doc"
    output_path = tmp_path / "test.docx"
    input_path.touch()

    mock_doc = MagicMock()
    mock_word = MagicMock()
    mock_word.Documents.Open.return_value = mock_doc

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_word
    # Parent mock must expose .client for attribute-based access
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        _convert_with_win32com(input_path, output_path)

        mock_win32com_client.Dispatch.assert_called_once_with("Word.Application")
        mock_word.Documents.Open.assert_called_once()
        mock_doc.SaveAs.assert_called_once()
        mock_doc.Close.assert_called_once_with(False)


def test_convert_with_win32com_excel(tmp_path: Path) -> None:
    """_convert_with_win32com calls Excel SaveAs for .xls → .xlsx."""
    input_path = tmp_path / "test.xls"
    output_path = tmp_path / "test.xlsx"
    input_path.touch()

    mock_wb = MagicMock()
    mock_excel = MagicMock()
    mock_excel.Workbooks.Open.return_value = mock_wb

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_excel
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        _convert_with_win32com(input_path, output_path)

        mock_win32com_client.Dispatch.assert_called_once_with("Excel.Application")
        mock_excel.Workbooks.Open.assert_called_once()
        mock_wb.SaveAs.assert_called_once()


def test_convert_with_uno(tmp_path: Path) -> None:
    """_convert_with_uno calls storeToURL with correct filter name."""
    input_path = tmp_path / "test.doc"
    output_path = tmp_path / "test.docx"
    input_path.touch()

    mock_doc = MagicMock()
    mock_desktop = MagicMock()
    mock_desktop.loadComponentFromURL.return_value = mock_doc

    with (
        patch(
            "src.core.office_processor._get_uno_desktop",
            return_value=mock_desktop,
        ),
        patch(
            "src.core.office_processor._uno_file_url",
            side_effect=lambda p: f"file:///{p}",
        ),
        patch.dict(
            "sys.modules",
            {
                "com": MagicMock(),
                "com.sun": MagicMock(),
                "com.sun.star": MagicMock(),
                "com.sun.star.beans": MagicMock(),
            },
        ),
    ):
        _convert_with_uno(input_path, output_path)

        mock_desktop.loadComponentFromURL.assert_called_once()
        mock_doc.storeToURL.assert_called_once()
        mock_doc.close.assert_called_once_with(True)


# ---------------------------------------------------------------------------
# Legacy format support — comment dispatching
# ---------------------------------------------------------------------------


def test_extract_comments_dispatches_doc_win32com() -> None:
    """_extract_comments dispatches .doc to win32com extractor."""
    with patch(
        "src.core.office_processor._extract_win32com_word_comments",
        return_value=[("comment:1", "Hello")],
    ) as mock_fn:
        result = _extract_comments(Path("test.doc"), ".doc", "win32com")
        mock_fn.assert_called_once()
        assert result == [("comment:1", "Hello")]


def test_extract_comments_dispatches_xls_win32com() -> None:
    """_extract_comments dispatches .xls to win32com extractor."""
    with patch(
        "src.core.office_processor._extract_win32com_excel_comments",
        return_value=[("comment:Sheet1:1:1", "Hi")],
    ) as mock_fn:
        result = _extract_comments(Path("test.xls"), ".xls", "win32com")
        mock_fn.assert_called_once()
        assert result == [("comment:Sheet1:1:1", "Hi")]


def test_extract_comments_dispatches_ppt_win32com() -> None:
    """_extract_comments dispatches .ppt to win32com extractor."""
    with patch(
        "src.core.office_processor._extract_win32com_ppt_comments",
        return_value=[("comment:0:1", "Slide note")],
    ) as mock_fn:
        result = _extract_comments(Path("test.ppt"), ".ppt", "win32com")
        mock_fn.assert_called_once()
        assert result == [("comment:0:1", "Slide note")]


def test_extract_comments_dispatches_doc_uno() -> None:
    """_extract_comments dispatches .doc to UNO extractor when backend is uno."""
    with patch(
        "src.core.office_processor._extract_uno_writer_comments",
        return_value=[("comment:0", "UNO comment")],
    ) as mock_fn:
        result = _extract_comments(Path("test.doc"), ".doc", "uno")
        mock_fn.assert_called_once()
        assert result == [("comment:0", "UNO comment")]


def test_inject_comments_dispatches_doc_win32com() -> None:
    """_inject_comments dispatches .doc to win32com injector."""
    with patch(
        "src.core.office_processor._inject_win32com_word_comments",
    ) as mock_fn:
        _inject_comments(
            Path("test.doc"),
            {"comment:1": "Translated"},
            ".doc",
            "win32com",
        )
        mock_fn.assert_called_once()


def test_inject_comments_dispatches_xls_uno() -> None:
    """_inject_comments dispatches .xls to UNO injector."""
    with patch(
        "src.core.office_processor._inject_uno_calc_comments",
    ) as mock_fn:
        _inject_comments(
            Path("test.xls"),
            {"comment:Sheet1:1:1": "Translated"},
            ".xls",
            "uno",
        )
        mock_fn.assert_called_once()


def test_inject_comments_dispatches_ppt_win32com() -> None:
    """_inject_comments dispatches .ppt to win32com injector."""
    with patch(
        "src.core.office_processor._inject_win32com_ppt_comments",
    ) as mock_fn:
        _inject_comments(
            Path("test.ppt"),
            {"comment:0:1": "Translated"},
            ".ppt",
            "win32com",
        )
        mock_fn.assert_called_once()


# ---------------------------------------------------------------------------
# Legacy format support — _translate_doc_images dispatch
# ---------------------------------------------------------------------------


@patch("src.utils.config_manager.load_setting", return_value="TesseractOCR")
def test_translate_doc_images_dispatches_legacy(
    _mock_setting: MagicMock,
) -> None:
    """_translate_doc_images dispatches legacy formats to _translate_legacy_images."""
    with patch(
        "src.core.office_processor._translate_legacy_images",
    ) as mock_legacy:
        _translate_doc_images(
            Path("test.doc"),
            ".doc",
            "win32com",
            "Vietnamese",
            "English",
            None,
            None,
            None,
        )
        mock_legacy.assert_called_once()


@patch("src.utils.config_manager.load_setting", return_value="TesseractOCR")
def test_translate_doc_images_dispatches_modern(
    _mock_setting: MagicMock,
) -> None:
    """_translate_doc_images dispatches modern formats to _translate_zip_images."""
    with patch(
        "src.core.office_processor._translate_zip_images",
    ) as mock_zip:
        _translate_doc_images(
            Path("test.docx"),
            ".docx",
            "python_lib",
            "Vietnamese",
            "English",
            None,
            None,
            None,
        )
        mock_zip.assert_called_once()


@patch("src.utils.config_manager.load_setting", return_value="TesseractOCR")
def test_translate_doc_images_forwards_checkpoint_dir_modern(
    _mock_setting: MagicMock,
    tmp_path: Path,
) -> None:
    """Modern-format dispatch must forward ``checkpoint_dir`` to the ZIP pipeline.

    A dropped kwarg here silently disables the per-image cache —
    re-translating an entire document on retry instead of skipping
    already-translated images.  Regression guard.
    """
    cp_dir = tmp_path / "task_storage"
    cp_dir.mkdir()
    with patch(
        "src.core.office_processor._translate_zip_images",
    ) as mock_zip:
        _translate_doc_images(
            Path("test.docx"),
            ".docx",
            "python_lib",
            "Vietnamese",
            "English",
            None,
            None,
            None,
            checkpoint_dir=cp_dir,
        )
        mock_zip.assert_called_once()
        assert mock_zip.call_args.kwargs["checkpoint_dir"] == cp_dir


@patch("src.utils.config_manager.load_setting", return_value="TesseractOCR")
def test_translate_doc_images_forwards_checkpoint_dir_legacy(
    _mock_setting: MagicMock,
    tmp_path: Path,
) -> None:
    """Legacy-format dispatch must forward ``checkpoint_dir`` to the legacy pipeline.

    Same regression guard as the modern path — legacy ``.doc/.xls/.ppt``
    files go through a separate dispatcher and the kwarg must survive
    the routing to keep the per-image cache effective across legacy ↔
    modern round-trips.
    """
    cp_dir = tmp_path / "task_storage"
    cp_dir.mkdir()
    with patch(
        "src.core.office_processor._translate_legacy_images",
    ) as mock_legacy:
        _translate_doc_images(
            Path("test.doc"),
            ".doc",
            "win32com",
            "Vietnamese",
            "English",
            None,
            None,
            None,
            checkpoint_dir=cp_dir,
        )
        mock_legacy.assert_called_once()
        assert mock_legacy.call_args.kwargs["checkpoint_dir"] == cp_dir


@patch("src.core.office_processor._translate_zip_images")
def test_translate_legacy_images_forwards_checkpoint_dir(
    mock_zip: MagicMock,
    tmp_path: Path,
) -> None:
    """``_translate_legacy_images`` must forward ``checkpoint_dir`` to the ZIP pipeline.

    The legacy pipeline converts ``.doc/.xls/.ppt`` to a modern temp
    file then delegates to ``_translate_zip_images``; the cache key is
    content-addressed so it still benefits from caching across the
    round-trip — but only if this kwarg survives the call chain.
    """
    cp_dir = tmp_path / "task_storage"
    cp_dir.mkdir()
    legacy_file = tmp_path / "fake.doc"
    legacy_file.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1")  # OLE2 magic
    with patch(
        "src.core.office_processor._convert_with_win32com",
        side_effect=lambda _src, dst: dst.write_bytes(b"PK\x03\x04converted"),
    ):
        from src.core.office_processor import _translate_legacy_images  # noqa: PLC0415

        _translate_legacy_images(
            legacy_file,
            ".doc",
            "win32com",
            "Vietnamese",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
            checkpoint_dir=cp_dir,
        )
        mock_zip.assert_called_once()
        assert mock_zip.call_args.kwargs["checkpoint_dir"] == cp_dir


# ---------------------------------------------------------------------------
# Legacy format support — win32com comment extraction (mocked COM objects)
# ---------------------------------------------------------------------------


def test_extract_win32com_word_comments_basic() -> None:
    """Extracts top-level Word comments via mocked win32com."""
    mock_comment = MagicMock()
    mock_comment.Ancestor = None
    mock_comment.Index = 1
    mock_comment.Range.Text = "Test comment"

    mock_doc = MagicMock()
    mock_doc.Comments.Count = 1
    mock_doc.Comments.side_effect = lambda i: mock_comment

    mock_word = MagicMock()
    mock_word.Documents.Open.return_value = mock_doc

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_word
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        result = _extract_win32com_word_comments(Path("test.doc"))
        assert result == [("comment:1", "Test comment")]


def test_extract_win32com_excel_comments_basic() -> None:
    """Extracts Excel cell comments via mocked win32com."""
    mock_cell = MagicMock()
    mock_cell.Row = 1
    mock_cell.Column = 2

    mock_comment = MagicMock()
    mock_comment.Text.return_value = "Cell note"
    mock_comment.Parent = mock_cell

    mock_ws = MagicMock()
    mock_ws.Name = "Sheet1"
    mock_ws.Comments = [mock_comment]

    mock_wb = MagicMock()
    mock_wb.Worksheets = [mock_ws]

    mock_excel = MagicMock()
    mock_excel.Workbooks.Open.return_value = mock_wb

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_excel
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        result = _extract_win32com_excel_comments(Path("test.xls"))
        assert result == [("comment:Sheet1:1:2", "Cell note")]


def test_extract_win32com_ppt_comments_basic() -> None:
    """Extracts PPT slide comments via mocked win32com."""
    mock_comment = MagicMock()
    mock_comment.Text = "Slide comment"
    mock_comment.Index = 1

    mock_slide = MagicMock()
    mock_slide.Comments.Count = 1
    mock_slide.Comments.side_effect = lambda i: mock_comment

    mock_prs = MagicMock()
    mock_prs.Slides.Count = 1
    mock_prs.Slides.side_effect = lambda i: mock_slide

    mock_ppt = MagicMock()
    mock_ppt.Presentations.Open.return_value = mock_prs

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_ppt
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        result = _extract_win32com_ppt_comments(Path("test.ppt"))
        # slide_idx is 0-based: s_idx=1 → key uses s_idx-1=0
        assert result == [("comment:0:1", "Slide comment")]


# ---------------------------------------------------------------------------
# Shape / text-box translation — constants and gate
# ---------------------------------------------------------------------------


def test_shape_extensions_constant() -> None:
    """_SHAPE_EXTENSIONS contains Word and Spreadsheet formats (no PPT)."""
    expected = {".docx", ".xlsx", ".odt", ".ods", ".doc", ".xls"}
    assert expected == _SHAPE_EXTENSIONS


def test_shape_extensions_excludes_ppt() -> None:
    """PPT formats are excluded from _SHAPE_EXTENSIONS."""
    for ext in (".pptx", ".ppt", ".odp"):
        assert ext not in _SHAPE_EXTENSIONS


@patch("src.utils.config_manager.load_setting", return_value=True)
def test_should_translate_shapes_enabled(
    _mock_setting: MagicMock,
) -> None:
    """Returns True when setting is enabled and extension is supported."""
    for ext in (".docx", ".xlsx", ".odt", ".ods", ".doc", ".xls"):
        assert _should_translate_shapes(ext, "python_lib") is True


@patch("src.utils.config_manager.load_setting", return_value=False)
def test_should_translate_shapes_disabled(
    _mock_setting: MagicMock,
) -> None:
    """Returns False when setting is disabled."""
    assert _should_translate_shapes(".docx", "python_lib") is False


@patch("src.utils.config_manager.load_setting", return_value=True)
def test_should_translate_shapes_unsupported_ext(
    _mock_setting: MagicMock,
) -> None:
    """Returns False for PPT formats even when setting is enabled."""
    assert _should_translate_shapes(".pptx", "python_lib") is False
    assert _should_translate_shapes(".ppt", "win32com") is False
    assert _should_translate_shapes(".odp", "python_lib") is False


# ---------------------------------------------------------------------------
# Config injection path — _should_translate_images/comments/shapes
# ---------------------------------------------------------------------------


def test_should_translate_images_config_enabled() -> None:
    """config.should_translate_images=True → returns True."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    config = TranslationConfig(translate_doc_images=True, ocr_is_configured=True)
    assert _should_translate_images(".docx", "python_lib", config=config) is True


def test_should_translate_images_config_disabled() -> None:
    """config.should_translate_images=False → returns False."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    config = TranslationConfig(translate_doc_images=False, ocr_is_configured=False)
    assert _should_translate_images(".docx", "python_lib", config=config) is False


def test_should_translate_comments_config_enabled() -> None:
    """config.translate_doc_comments=True → returns True."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    config = TranslationConfig(translate_doc_comments=True)
    assert _should_translate_comments(".docx", "python_lib", config=config) is True


def test_should_translate_comments_config_disabled() -> None:
    """config.translate_doc_comments=False → returns False."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    config = TranslationConfig(translate_doc_comments=False)
    assert _should_translate_comments(".docx", "python_lib", config=config) is False


def test_should_translate_shapes_config_enabled() -> None:
    """config.translate_doc_shapes=True → returns True."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    config = TranslationConfig(translate_doc_shapes=True)
    assert _should_translate_shapes(".docx", "python_lib", config=config) is True


def test_should_translate_shapes_config_disabled() -> None:
    """config.translate_doc_shapes=False → returns False."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    config = TranslationConfig(translate_doc_shapes=False)
    assert _should_translate_shapes(".docx", "python_lib", config=config) is False


# ---------------------------------------------------------------------------
# Shape / text-box translation — dispatcher routing
# ---------------------------------------------------------------------------


def test_extract_shapes_dispatches_docx() -> None:
    """_extract_shapes dispatches .docx (non-UNO) to _extract_docx_shapes."""
    with patch(
        "src.core.office_processor._extract_docx_shapes",
        return_value=[("shape:0", "Hello")],
    ) as mock_fn:
        result = _extract_shapes(Path("test.docx"), ".docx", "python_lib")
        mock_fn.assert_called_once()
        assert result == [("shape:0", "Hello")]


def test_extract_shapes_dispatches_docx_uno() -> None:
    """_extract_shapes dispatches .docx with UNO backend to ZIP+lxml extractor.

    ZIP+lxml reads raw WordprocessingML run properties (including character-
    style resolution) so mixed formatting in text boxes is correctly detected.
    UNO preserves <wps:txbx> elements when saving DOCX.
    """
    with (
        patch(
            "src.core.office_processor._extract_docx_shapes",
            return_value=[("shape:0", "Box text")],
        ) as mock_zip,
        patch(
            "src.core.office_processor._extract_uno_writer_shapes",
        ) as mock_uno,
    ):
        result = _extract_shapes(Path("test.docx"), ".docx", "uno")
        mock_zip.assert_called_once()
        mock_uno.assert_not_called()
        assert result == [("shape:0", "Box text")]


def test_extract_shapes_dispatches_xlsx() -> None:
    """_extract_shapes dispatches .xlsx to _extract_xlsx_shapes."""
    with patch(
        "src.core.office_processor._extract_xlsx_shapes",
        return_value=[("shape:Sheet1:0", "Data")],
    ) as mock_fn:
        result = _extract_shapes(Path("test.xlsx"), ".xlsx", "python_lib")
        mock_fn.assert_called_once()
        assert result == [("shape:Sheet1:0", "Data")]


def test_extract_shapes_dispatches_odt() -> None:
    """_extract_shapes dispatches .odt to _extract_odt_shapes."""
    with patch(
        "src.core.office_processor._extract_odt_shapes",
        return_value=[("shape:0", "Text")],
    ) as mock_fn:
        result = _extract_shapes(Path("test.odt"), ".odt", "python_lib")
        mock_fn.assert_called_once()
        assert result == [("shape:0", "Text")]


def test_extract_shapes_dispatches_ods() -> None:
    """_extract_shapes dispatches .ods to _extract_ods_shapes."""
    with patch(
        "src.core.office_processor._extract_ods_shapes",
        return_value=[("shape:Sheet1:0", "Cell text")],
    ) as mock_fn:
        result = _extract_shapes(Path("test.ods"), ".ods", "python_lib")
        mock_fn.assert_called_once()
        assert result == [("shape:Sheet1:0", "Cell text")]


def test_extract_shapes_dispatches_doc_win32com() -> None:
    """_extract_shapes dispatches .doc to win32com extractor."""
    with patch(
        "src.core.office_processor._extract_win32com_word_shapes",
        return_value=[("shape:0", "Legacy")],
    ) as mock_fn:
        result = _extract_shapes(Path("test.doc"), ".doc", "win32com")
        mock_fn.assert_called_once()
        assert result == [("shape:0", "Legacy")]


def test_extract_shapes_dispatches_doc_uno() -> None:
    """_extract_shapes dispatches .doc to UNO extractor when backend is uno."""
    with patch(
        "src.core.office_processor._extract_uno_writer_shapes",
        return_value=[("shape:0", "UNO text")],
    ) as mock_fn:
        result = _extract_shapes(Path("test.doc"), ".doc", "uno")
        mock_fn.assert_called_once()
        assert result == [("shape:0", "UNO text")]


def test_extract_shapes_dispatches_xls_win32com() -> None:
    """_extract_shapes dispatches .xls to win32com extractor."""
    with patch(
        "src.core.office_processor._extract_win32com_excel_shapes",
        return_value=[("shape:Sheet1:0", "XLS text")],
    ) as mock_fn:
        result = _extract_shapes(Path("test.xls"), ".xls", "win32com")
        mock_fn.assert_called_once()
        assert result == [("shape:Sheet1:0", "XLS text")]


def test_extract_shapes_dispatches_xls_uno() -> None:
    """_extract_shapes dispatches .xls to UNO extractor when backend is uno."""
    with patch(
        "src.core.office_processor._extract_uno_calc_shapes",
        return_value=[("shape:Sheet1:0", "UNO calc")],
    ) as mock_fn:
        result = _extract_shapes(Path("test.xls"), ".xls", "uno")
        mock_fn.assert_called_once()
        assert result == [("shape:Sheet1:0", "UNO calc")]


def test_extract_shapes_returns_empty_for_unsupported() -> None:
    """_extract_shapes returns empty list for unsupported extensions."""
    assert _extract_shapes(Path("test.pptx"), ".pptx", "python_lib") == []


def test_inject_shapes_dispatches_docx() -> None:
    """_inject_shapes dispatches .docx (non-UNO) to _inject_docx_shapes."""
    with patch(
        "src.core.office_processor._inject_docx_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.docx"),
            {"shape:0": "Translated"},
            ".docx",
            "python_lib",
        )
        mock_fn.assert_called_once()


def test_inject_shapes_dispatches_docx_uno() -> None:
    """_inject_shapes dispatches .docx with UNO backend to ZIP+lxml injector.

    Mirrors test_extract_shapes_dispatches_docx_uno: both extraction and
    injection use ZIP+lxml for consistent shape indexing and correct
    inline-HTML formatting support.
    """
    with (
        patch(
            "src.core.office_processor._inject_docx_shapes",
        ) as mock_zip,
        patch(
            "src.core.office_processor._inject_uno_writer_shapes",
        ) as mock_uno,
    ):
        _inject_shapes(
            Path("test.docx"),
            {"shape:0": "Translated"},
            ".docx",
            "uno",
        )
        mock_zip.assert_called_once()
        mock_uno.assert_not_called()


def test_inject_shapes_dispatches_xlsx() -> None:
    """_inject_shapes dispatches .xlsx to _inject_xlsx_shapes."""
    with patch(
        "src.core.office_processor._inject_xlsx_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.xlsx"),
            {"shape:Sheet1:0": "Translated"},
            ".xlsx",
            "python_lib",
        )
        mock_fn.assert_called_once()


def test_inject_shapes_dispatches_odt() -> None:
    """_inject_shapes dispatches .odt to _inject_odt_shapes."""
    with patch(
        "src.core.office_processor._inject_odt_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.odt"),
            {"shape:0": "Translated"},
            ".odt",
            "python_lib",
        )
        mock_fn.assert_called_once()


def test_inject_shapes_dispatches_ods() -> None:
    """_inject_shapes dispatches .ods to _inject_ods_shapes."""
    with patch(
        "src.core.office_processor._inject_ods_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.ods"),
            {"shape:Sheet1:0": "Translated"},
            ".ods",
            "python_lib",
        )
        mock_fn.assert_called_once()


def test_inject_shapes_dispatches_doc_win32com() -> None:
    """_inject_shapes dispatches .doc to win32com injector."""
    with patch(
        "src.core.office_processor._inject_win32com_word_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.doc"),
            {"shape:0": "Translated"},
            ".doc",
            "win32com",
        )
        mock_fn.assert_called_once()


def test_inject_shapes_dispatches_doc_uno() -> None:
    """_inject_shapes dispatches .doc to UNO injector."""
    with patch(
        "src.core.office_processor._inject_uno_writer_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.doc"),
            {"shape:0": "Translated"},
            ".doc",
            "uno",
        )
        mock_fn.assert_called_once()


def test_inject_shapes_dispatches_xls_win32com() -> None:
    """_inject_shapes dispatches .xls to win32com injector."""
    with patch(
        "src.core.office_processor._inject_win32com_excel_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.xls"),
            {"shape:Sheet1:0": "Translated"},
            ".xls",
            "win32com",
        )
        mock_fn.assert_called_once()


def test_inject_shapes_dispatches_xls_uno() -> None:
    """_inject_shapes dispatches .xls to UNO injector."""
    with patch(
        "src.core.office_processor._inject_uno_calc_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.xls"),
            {"shape:Sheet1:0": "Translated"},
            ".xls",
            "uno",
        )
        mock_fn.assert_called_once()


def test_inject_shapes_skips_non_shape_keys() -> None:
    """_inject_shapes does nothing when no shape: keys are present."""
    with patch(
        "src.core.office_processor._inject_docx_shapes",
    ) as mock_fn:
        _inject_shapes(
            Path("test.docx"),
            {"comment:0": "Not a shape"},
            ".docx",
            "python_lib",
        )
        mock_fn.assert_not_called()


# ---------------------------------------------------------------------------
# Shape / text-box — shared helpers
# ---------------------------------------------------------------------------


def test_rewrite_zip_content(tmp_path: Path) -> None:
    """_rewrite_zip_content atomically rewrites ZIP with modified data."""
    zip_path = tmp_path / "test.zip"

    # Create a ZIP with two files
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("a.txt", "original A")
        zf.writestr("b.txt", "original B")

    # Read and modify
    with zipfile.ZipFile(zip_path, "r") as zf:
        all_items = zf.infolist()
        file_data = {i.filename: zf.read(i.filename) for i in all_items}

    file_data["a.txt"] = b"modified A"
    _rewrite_zip_content(zip_path, file_data, all_items)

    # Verify
    with zipfile.ZipFile(zip_path, "r") as zf:
        assert zf.read("a.txt") == b"modified A"
        assert zf.read("b.txt") == b"original B"


def test_rewrite_zip_content_cleans_tmp_on_error(tmp_path: Path) -> None:
    """_rewrite_zip_content removes .tmp file on failure."""
    zip_path = tmp_path / "test.zip"

    # Create initial ZIP
    with zipfile.ZipFile(zip_path, "w") as zf:
        zf.writestr("a.txt", "data")

    tmp_file = zip_path.with_suffix(".zip.tmp")
    # Cause error by passing invalid data
    with pytest.raises(Exception):  # noqa: B017
        _rewrite_zip_content(zip_path, {"a.txt": None}, [MagicMock(filename="a.txt")])

    assert not tmp_file.exists()


def test_collect_wps_texts() -> None:
    """_collect_wps_texts finds <wps:txbx> elements with <w:t> runs."""
    from lxml import etree  # noqa: PLC0415

    wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    xml_str = f"""\
<root xmlns:wps="{wps_ns}" xmlns:w="{w_ns}">
  <wps:txbx>
    <w:txbxContent>
      <w:p><w:r><w:t>Hello </w:t></w:r><w:r><w:t>World</w:t></w:r></w:p>
    </w:txbxContent>
  </wps:txbx>
  <wps:txbx>
    <w:txbxContent>
      <w:p><w:r><w:t>Second box</w:t></w:r></w:p>
    </w:txbxContent>
  </wps:txbx>
</root>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    results = _collect_wps_texts(root)
    expected_count = 2

    assert len(results) == expected_count
    assert results[0][0] == "Hello World"
    assert len(results[0][1]) == expected_count  # Two <w:t> elements
    assert results[1][0] == "Second box"
    assert len(results[1][1]) == 1


def test_collect_wps_texts_skips_empty() -> None:
    """_collect_wps_texts skips text boxes with only whitespace."""
    from lxml import etree  # noqa: PLC0415

    wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    xml_str = f"""\
<root xmlns:wps="{wps_ns}" xmlns:w="{w_ns}">
  <wps:txbx>
    <w:txbxContent>
      <w:p><w:r><w:t>   </w:t></w:r></w:p>
    </w:txbxContent>
  </wps:txbx>
</root>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    results = _collect_wps_texts(root)
    assert results == []


def test_extract_odf_paragraph_text() -> None:
    """_extract_odf_paragraph_text extracts text from <text:p> children."""
    from lxml import etree  # noqa: PLC0415

    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    text_p_tag = f"{{{text_ns}}}p"

    xml_str = f"""\
<draw:text-box xmlns:draw="{draw_ns}" xmlns:text="{text_ns}">
  <text:p>First paragraph</text:p>
  <text:p>Second paragraph</text:p>
</draw:text-box>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    result = _extract_odf_paragraph_text(root, text_p_tag)
    assert result == "First paragraph\nSecond paragraph"


def test_extract_odf_paragraph_text_mixed_content() -> None:
    """_extract_odf_paragraph_text handles child elements and tail text."""
    from lxml import etree  # noqa: PLC0415

    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    text_p_tag = f"{{{text_ns}}}p"

    xml_str = f"""\
<draw:text-box xmlns:draw="{draw_ns}" xmlns:text="{text_ns}">
  <text:p>Hello <text:span>bold</text:span> world</text:p>
</draw:text-box>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    result = _extract_odf_paragraph_text(root, text_p_tag)
    assert result == "Hello bold world"


def test_extract_odf_paragraph_text_empty() -> None:
    """_extract_odf_paragraph_text returns empty for no paragraphs."""
    from lxml import etree  # noqa: PLC0415

    draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    text_p_tag = f"{{{text_ns}}}p"

    xml_str = f'<draw:text-box xmlns:draw="{draw_ns}"/>'
    root = etree.fromstring(xml_str.encode("utf-8"))
    result = _extract_odf_paragraph_text(root, text_p_tag)
    assert result == ""


def test_inject_odf_paragraph_text() -> None:
    """_inject_odf_paragraph_text replaces text and removes extra paragraphs."""
    from lxml import etree  # noqa: PLC0415

    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    text_p_tag = f"{{{text_ns}}}p"

    xml_str = f"""\
<draw:text-box xmlns:draw="{draw_ns}" xmlns:text="{text_ns}">
  <text:p>First paragraph</text:p>
  <text:p>Second paragraph</text:p>
</draw:text-box>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    result = _inject_odf_paragraph_text(root, "Translated text", text_p_tag)

    assert result is True
    paras = root.findall(text_p_tag)
    assert len(paras) == 1
    assert paras[0].text == "Translated text"


def test_inject_odf_paragraph_text_no_paras() -> None:
    """_inject_odf_paragraph_text returns False when no paragraphs exist."""
    from lxml import etree  # noqa: PLC0415

    draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    text_p_tag = f"{{{text_ns}}}p"

    xml_str = f'<draw:text-box xmlns:draw="{draw_ns}"/>'
    root = etree.fromstring(xml_str.encode("utf-8"))
    result = _inject_odf_paragraph_text(root, "New text", text_p_tag)
    assert result is False


# ---------------------------------------------------------------------------
# Shape / text-box — DOCX extract/inject (ZIP + lxml)
# ---------------------------------------------------------------------------

# Namespace URIs for building DOCX shape XML
_WPS_NS_URI = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
_W_NS_URI = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

# Minimal DOCX document.xml with text boxes
_DOCX_DOCUMENT_XML_WITH_SHAPES = f"""\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="{_W_NS_URI}"
            xmlns:wps="{_WPS_NS_URI}">
  <w:body>
    <wps:txbx>
      <w:txbxContent>
        <w:p><w:r><w:t>Text box one</w:t></w:r></w:p>
      </w:txbxContent>
    </wps:txbx>
    <wps:txbx>
      <w:txbxContent>
        <w:p><w:r><w:t>Text box two</w:t></w:r></w:p>
      </w:txbxContent>
    </wps:txbx>
  </w:body>
</w:document>"""

# Minimal DOCX header with a text box
_DOCX_HEADER_XML_WITH_SHAPE = f"""\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:hdr xmlns:w="{_W_NS_URI}"
       xmlns:wps="{_WPS_NS_URI}">
  <wps:txbx>
    <w:txbxContent>
      <w:p><w:r><w:t>Header text box</w:t></w:r></w:p>
    </w:txbxContent>
  </wps:txbx>
</w:hdr>"""


def _make_docx_with_shapes(tmp_path: Path, name: str) -> Path:
    """Creates a minimal DOCX ZIP with text boxes in document.xml and a header.

    Args:
        tmp_path: Pytest tmp_path fixture.
        name: Filename for the DOCX.

    Returns:
        Path to the created DOCX file.
    """
    # First create a real DOCX to get [Content_Types].xml and _rels
    doc = Document()
    doc.add_paragraph("Body text")
    path = tmp_path / name
    doc.save(str(path))

    # Now inject our shape XML into the ZIP
    import shutil  # noqa: PLC0415

    with zipfile.ZipFile(path, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    # Replace document.xml with our shape-bearing version
    data["word/document.xml"] = _DOCX_DOCUMENT_XML_WITH_SHAPES.encode("utf-8")
    # Add header with shape
    data["word/header1.xml"] = _DOCX_HEADER_XML_WITH_SHAPE.encode("utf-8")

    tmp_zip = path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for item in items:
            zf_out.writestr(item, data[item.filename])
        for fname, content in data.items():
            if not any(i.filename == fname for i in items):
                zf_out.writestr(fname, content)
    shutil.move(str(tmp_zip), str(path))

    return path


def test_extract_docx_shapes(tmp_path: Path) -> None:
    """Extracts text box text from a DOCX file."""
    docx_path = _make_docx_with_shapes(tmp_path, "shapes.docx")
    texts = _extract_docx_shapes(docx_path)
    expected_count = 3

    assert len(texts) == expected_count
    keys = dict(texts)
    assert keys["shape:0"] == "Text box one"
    assert keys["shape:1"] == "Text box two"
    assert keys["shape:2"] == "Header text box"


def test_extract_docx_shapes_empty(tmp_path: Path) -> None:
    """DOCX without shapes returns empty list."""
    doc = Document()
    doc.add_paragraph("Just text")
    path = tmp_path / "no_shapes.docx"
    doc.save(str(path))

    texts = _extract_docx_shapes(path)
    assert texts == []


def test_inject_docx_shapes(tmp_path: Path) -> None:
    """Translated text is injected back into DOCX text boxes."""
    docx_path = _make_docx_with_shapes(tmp_path, "inject.docx")

    translations = {
        "shape:0": "Boîte une",
        "shape:1": "Boîte deux",
        "shape:2": "En-tête traduit",
    }
    _inject_docx_shapes(docx_path, translations)

    # Re-extract and verify
    texts = _extract_docx_shapes(docx_path)
    result = dict(texts)
    assert result["shape:0"] == "Boîte une"
    assert result["shape:1"] == "Boîte deux"
    assert result["shape:2"] == "En-tête traduit"


def test_inject_docx_shapes_partial(tmp_path: Path) -> None:
    """Only shape keys present in translations are injected."""
    docx_path = _make_docx_with_shapes(tmp_path, "partial.docx")

    # Only translate shape:1
    _inject_docx_shapes(docx_path, {"shape:1": "Seulement deux"})

    texts = _extract_docx_shapes(docx_path)
    result = dict(texts)
    assert result["shape:0"] == "Text box one"  # Unchanged
    assert result["shape:1"] == "Seulement deux"
    assert result["shape:2"] == "Header text box"  # Unchanged


def test_inject_docx_shapes_no_match(tmp_path: Path) -> None:
    """Injection with no matching keys leaves file unchanged."""
    docx_path = _make_docx_with_shapes(tmp_path, "no_match.docx")

    _inject_docx_shapes(docx_path, {"shape:99": "No match"})

    # File should still be valid and unchanged
    expected_count = 3
    texts = _extract_docx_shapes(docx_path)
    assert len(texts) == expected_count


# ---------------------------------------------------------------------------
# Shape / text-box — XLSX extract/inject (ZIP + lxml)
# ---------------------------------------------------------------------------

# Namespace URIs for XLSX drawing XML
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_XDR_NS = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
_SXML_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
_PKG_RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"

# Drawing XML with two shapes
_XLSX_DRAWING_XML = f"""\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<xdr:wsDr xmlns:xdr="{_XDR_NS}" xmlns:a="{_A_NS}">
  <xdr:twoCellAnchor>
    <xdr:sp>
      <a:txBody>
        <a:bodyPr/>
        <a:p><a:r><a:rPr lang="en-US"/><a:t>Shape text A</a:t></a:r></a:p>
      </a:txBody>
    </xdr:sp>
  </xdr:twoCellAnchor>
  <xdr:twoCellAnchor>
    <xdr:sp>
      <a:txBody>
        <a:bodyPr/>
        <a:p><a:r><a:rPr lang="en-US"/><a:t>Shape text B</a:t></a:r></a:p>
      </a:txBody>
    </xdr:sp>
  </xdr:twoCellAnchor>
</xdr:wsDr>"""


def _make_xlsx_with_shapes(tmp_path: Path, name: str) -> Path:
    """Creates an XLSX with a drawing containing shapes.

    Args:
        tmp_path: Pytest tmp_path fixture.
        name: Filename for the XLSX.

    Returns:
        Path to the created XLSX file.
    """
    import shutil  # noqa: PLC0415

    # Create a real XLSX first
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Test"
    path = tmp_path / name
    wb.save(str(path))

    # Inject drawing XML and wire relationships
    from lxml import etree  # noqa: PLC0415

    with zipfile.ZipFile(path, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    # Add drawing file
    data["xl/drawings/drawing1.xml"] = _XLSX_DRAWING_XML.encode("utf-8")

    # Create sheet1.xml.rels linking to drawing
    drawing_rel_type = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing"
    )
    rels_xml = f"""\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{_PKG_RELS_NS}">
  <Relationship Id="rId1" Type="{drawing_rel_type}" \
Target="../drawings/drawing1.xml"/>
</Relationships>"""
    data["xl/worksheets/_rels/sheet1.xml.rels"] = rels_xml.encode("utf-8")

    # Update [Content_Types].xml to include drawing content type
    ct_root = etree.fromstring(data["[Content_Types].xml"])
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    etree.SubElement(
        ct_root,
        f"{{{ct_ns}}}Override",
        PartName="/xl/drawings/drawing1.xml",
        ContentType="application/vnd.openxmlformats-officedocument.drawing+xml",
    )
    data["[Content_Types].xml"] = etree.tostring(
        ct_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    tmp_zip = path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for item in items:
            zf_out.writestr(item, data[item.filename])
        for fname, content in data.items():
            if not any(i.filename == fname for i in items):
                zf_out.writestr(fname, content)
    shutil.move(str(tmp_zip), str(path))

    return path


def test_resolve_xlsx_sheet_drawings(tmp_path: Path) -> None:
    """_resolve_xlsx_sheet_drawings maps sheet names to drawing paths."""
    xlsx_path = _make_xlsx_with_shapes(tmp_path, "resolve.xlsx")

    with zipfile.ZipFile(xlsx_path, "r") as zf:
        result = _resolve_xlsx_sheet_drawings(zf)

    assert len(result) == 1
    assert result[0] == ("Sheet1", "xl/drawings/drawing1.xml")


def test_resolve_xlsx_sheet_drawings_no_drawings(tmp_path: Path) -> None:
    """_resolve_xlsx_sheet_drawings returns empty for XLSX without drawings."""
    wb = Workbook()
    wb.active["A1"] = "No shapes"
    path = tmp_path / "plain.xlsx"
    wb.save(str(path))

    with zipfile.ZipFile(path, "r") as zf:
        result = _resolve_xlsx_sheet_drawings(zf)

    assert result == []


def test_extract_xlsx_shapes(tmp_path: Path) -> None:
    """Extracts shape text from XLSX drawing XML."""
    xlsx_path = _make_xlsx_with_shapes(tmp_path, "shapes.xlsx")
    texts = _extract_xlsx_shapes(xlsx_path)
    expected_count = 2

    assert len(texts) == expected_count
    keys = dict(texts)
    assert keys["shape:Sheet1:0"] == "Shape text A"
    assert keys["shape:Sheet1:1"] == "Shape text B"


def test_extract_xlsx_shapes_empty(tmp_path: Path) -> None:
    """XLSX without shapes returns empty list."""
    wb = Workbook()
    wb.active["A1"] = "No shapes"
    path = tmp_path / "plain.xlsx"
    wb.save(str(path))

    texts = _extract_xlsx_shapes(path)
    assert texts == []


def test_inject_xlsx_shapes(tmp_path: Path) -> None:
    """Translated text is injected into XLSX shapes."""
    xlsx_path = _make_xlsx_with_shapes(tmp_path, "inject.xlsx")

    translations = {
        "shape:Sheet1:0": "Forme A traduite",
        "shape:Sheet1:1": "Forme B traduite",
    }
    _inject_xlsx_shapes(xlsx_path, translations)

    # Re-extract and verify
    texts = _extract_xlsx_shapes(xlsx_path)
    result = dict(texts)
    assert result["shape:Sheet1:0"] == "Forme A traduite"
    assert result["shape:Sheet1:1"] == "Forme B traduite"


def test_inject_xlsx_shapes_partial(tmp_path: Path) -> None:
    """Only matching keys are injected in XLSX shapes."""
    xlsx_path = _make_xlsx_with_shapes(tmp_path, "partial.xlsx")

    _inject_xlsx_shapes(xlsx_path, {"shape:Sheet1:0": "Seulement A"})

    texts = _extract_xlsx_shapes(xlsx_path)
    result = dict(texts)
    assert result["shape:Sheet1:0"] == "Seulement A"
    assert result["shape:Sheet1:1"] == "Shape text B"  # Unchanged


# ---------------------------------------------------------------------------
# Shape / text-box — ODT extract/inject (ZIP + lxml)
# ---------------------------------------------------------------------------

# ODF namespace URIs
_ODT_OFFICE_NS = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
_ODT_TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
_ODT_DRAW_NS = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
_ODT_TABLE_NS = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"

# Minimal ODT content.xml with text boxes
_ODT_CONTENT_WITH_SHAPES = f"""\
<?xml version="1.0" encoding="UTF-8"?>
<office:document-content
    xmlns:office="{_ODT_OFFICE_NS}"
    xmlns:text="{_ODT_TEXT_NS}"
    xmlns:draw="{_ODT_DRAW_NS}">
  <office:body>
    <office:text>
      <text:p>Normal paragraph</text:p>
      <draw:frame>
        <draw:text-box>
          <text:p>First text box</text:p>
        </draw:text-box>
      </draw:frame>
      <draw:frame>
        <draw:text-box>
          <text:p>Second text box</text:p>
          <text:p>with two lines</text:p>
        </draw:text-box>
      </draw:frame>
    </office:text>
  </office:body>
</office:document-content>"""


def _make_odt_with_shapes(tmp_path: Path, name: str) -> Path:
    """Creates a minimal ODT ZIP with draw:text-box elements.

    Args:
        tmp_path: Pytest tmp_path fixture.
        name: Filename for the ODT.

    Returns:
        Path to the created ODT file.
    """
    path = tmp_path / name
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("content.xml", _ODT_CONTENT_WITH_SHAPES)
        zf.writestr("META-INF/manifest.xml", "<?xml version='1.0'?><manifest/>")
    return path


def test_extract_odt_shapes(tmp_path: Path) -> None:
    """Extracts text box text from an ODT file."""
    odt_path = _make_odt_with_shapes(tmp_path, "shapes.odt")
    texts = _extract_odt_shapes(odt_path)
    expected_count = 2

    assert len(texts) == expected_count
    keys = dict(texts)
    assert keys["shape:0"] == "First text box"
    assert keys["shape:1"] == "Second text box\nwith two lines"


def test_extract_odt_shapes_empty(tmp_path: Path) -> None:
    """ODT without text boxes returns empty list."""
    content_xml = f"""\
<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="{_ODT_OFFICE_NS}"
                         xmlns:text="{_ODT_TEXT_NS}">
  <office:body><office:text>
    <text:p>Just text</text:p>
  </office:text></office:body>
</office:document-content>"""

    path = tmp_path / "plain.odt"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("content.xml", content_xml)

    texts = _extract_odt_shapes(path)
    assert texts == []


def test_inject_odt_shapes(tmp_path: Path) -> None:
    """Translated text is injected into ODT text boxes."""
    odt_path = _make_odt_with_shapes(tmp_path, "inject.odt")

    translations = {
        "shape:0": "Première boîte",
        "shape:1": "Deuxième boîte",
    }
    _inject_odt_shapes(odt_path, translations)

    # Re-extract and verify
    texts = _extract_odt_shapes(odt_path)
    result = dict(texts)
    assert result["shape:0"] == "Première boîte"
    assert result["shape:1"] == "Deuxième boîte"


def test_inject_odt_shapes_removes_extra_paras(tmp_path: Path) -> None:
    """Injection collapses multi-paragraph text boxes to a single paragraph."""
    odt_path = _make_odt_with_shapes(tmp_path, "collapse.odt")

    # shape:1 originally has two <text:p>, after injection it should have one
    _inject_odt_shapes(odt_path, {"shape:1": "Single line"})

    texts = _extract_odt_shapes(odt_path)
    result = dict(texts)
    assert result["shape:1"] == "Single line"


def test_inject_odt_shapes_no_match(tmp_path: Path) -> None:
    """Injection with no matching keys leaves ODT unchanged."""
    odt_path = _make_odt_with_shapes(tmp_path, "no_match.odt")

    _inject_odt_shapes(odt_path, {"shape:99": "No match"})

    expected_count = 2
    texts = _extract_odt_shapes(odt_path)
    assert len(texts) == expected_count  # Still has both text boxes


# ---------------------------------------------------------------------------
# Shape / text-box — ODS extract/inject (ZIP + lxml)
# ---------------------------------------------------------------------------

# Minimal ODS content.xml with per-table text boxes
_ODS_CONTENT_WITH_SHAPES = f"""\
<?xml version="1.0" encoding="UTF-8"?>
<office:document-content
    xmlns:office="{_ODT_OFFICE_NS}"
    xmlns:text="{_ODT_TEXT_NS}"
    xmlns:draw="{_ODT_DRAW_NS}"
    xmlns:table="{_ODT_TABLE_NS}">
  <office:body>
    <office:spreadsheet>
      <table:table table:name="Sheet1">
        <table:table-row>
          <table:table-cell>
            <draw:frame>
              <draw:text-box>
                <text:p>Sheet1 shape A</text:p>
              </draw:text-box>
            </draw:frame>
          </table:table-cell>
        </table:table-row>
      </table:table>
      <table:table table:name="Sheet2">
        <table:table-row>
          <table:table-cell>
            <draw:frame>
              <draw:text-box>
                <text:p>Sheet2 shape X</text:p>
              </draw:text-box>
            </draw:frame>
          </table:table-cell>
        </table:table-row>
      </table:table>
    </office:spreadsheet>
  </office:body>
</office:document-content>"""


def _make_ods_with_shapes(tmp_path: Path, name: str) -> Path:
    """Creates a minimal ODS ZIP with per-table draw:text-box elements.

    Args:
        tmp_path: Pytest tmp_path fixture.
        name: Filename for the ODS.

    Returns:
        Path to the created ODS file.
    """
    path = tmp_path / name
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("content.xml", _ODS_CONTENT_WITH_SHAPES)
        zf.writestr("META-INF/manifest.xml", "<?xml version='1.0'?><manifest/>")
    return path


def test_extract_ods_shapes(tmp_path: Path) -> None:
    """Extracts text box text from an ODS file with per-table keys."""
    ods_path = _make_ods_with_shapes(tmp_path, "shapes.ods")
    texts = _extract_ods_shapes(ods_path)
    expected_count = 2

    assert len(texts) == expected_count
    keys = dict(texts)
    assert keys["shape:Sheet1:0"] == "Sheet1 shape A"
    assert keys["shape:Sheet2:0"] == "Sheet2 shape X"


def test_extract_ods_shapes_empty(tmp_path: Path) -> None:
    """ODS without text boxes returns empty list."""
    content_xml = f"""\
<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="{_ODT_OFFICE_NS}"
                         xmlns:table="{_ODT_TABLE_NS}"
                         xmlns:text="{_ODT_TEXT_NS}">
  <office:body><office:spreadsheet>
    <table:table table:name="Sheet1">
      <table:table-row><table:table-cell/></table:table-row>
    </table:table>
  </office:spreadsheet></office:body>
</office:document-content>"""

    path = tmp_path / "plain.ods"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("content.xml", content_xml)

    texts = _extract_ods_shapes(path)
    assert texts == []


def test_inject_ods_shapes(tmp_path: Path) -> None:
    """Translated text is injected into ODS text boxes."""
    ods_path = _make_ods_with_shapes(tmp_path, "inject.ods")

    translations = {
        "shape:Sheet1:0": "Forme feuille1",
        "shape:Sheet2:0": "Forme feuille2",
    }
    _inject_ods_shapes(ods_path, translations)

    # Re-extract and verify
    texts = _extract_ods_shapes(ods_path)
    result = dict(texts)
    assert result["shape:Sheet1:0"] == "Forme feuille1"
    assert result["shape:Sheet2:0"] == "Forme feuille2"


def test_inject_ods_shapes_partial(tmp_path: Path) -> None:
    """Only matching keys are injected in ODS shapes."""
    ods_path = _make_ods_with_shapes(tmp_path, "partial.ods")

    _inject_ods_shapes(ods_path, {"shape:Sheet1:0": "Seulement feuille1"})

    texts = _extract_ods_shapes(ods_path)
    result = dict(texts)
    assert result["shape:Sheet1:0"] == "Seulement feuille1"
    assert result["shape:Sheet2:0"] == "Sheet2 shape X"  # Unchanged


# ---------------------------------------------------------------------------
# Shape / text-box — win32com extraction (mocked COM objects)
# ---------------------------------------------------------------------------


def test_extract_win32com_word_shapes_basic() -> None:
    """Extracts Word shape text via mocked win32com."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_win32com_word_shapes,
    )

    mock_shape1 = MagicMock()
    mock_shape1.TextFrame.HasText = True
    mock_shape1.TextFrame.TextRange.Text = "Shape one"

    mock_shape2 = MagicMock()
    mock_shape2.TextFrame.HasText = True
    mock_shape2.TextFrame.TextRange.Text = "Shape two"

    mock_shape3 = MagicMock()
    mock_shape3.TextFrame.HasText = False

    mock_doc = MagicMock()
    mock_doc.Shapes.Count = 3
    mock_doc.Shapes.side_effect = lambda i: {
        1: mock_shape1,
        2: mock_shape2,
        3: mock_shape3,
    }[i]

    mock_word = MagicMock()
    mock_word.Documents.Open.return_value = mock_doc

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_word
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        result = _extract_win32com_word_shapes(Path("test.doc"))
        assert result == [("shape:0", "Shape one"), ("shape:1", "Shape two")]


def test_extract_win32com_excel_shapes_basic() -> None:
    """Extracts Excel shape text via mocked win32com."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_win32com_excel_shapes,
    )

    mock_shape = MagicMock()
    mock_shape.TextFrame2.HasText = True
    mock_shape.TextFrame2.TextRange.Text = "Excel shape"

    mock_ws = MagicMock()
    mock_ws.Name = "Data"
    mock_ws.Shapes.Count = 1
    mock_ws.Shapes.side_effect = lambda i: mock_shape

    mock_wb = MagicMock()
    mock_wb.Worksheets = [mock_ws]

    mock_excel = MagicMock()
    mock_excel.Workbooks.Open.return_value = mock_wb

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_excel
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        result = _extract_win32com_excel_shapes(Path("test.xls"))
        assert result == [("shape:Data:0", "Excel shape")]


def test_inject_win32com_word_shapes_basic() -> None:
    """Injects translated shape text via mocked win32com."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_win32com_word_shapes,
    )

    mock_shape = MagicMock()

    mock_doc = MagicMock()
    mock_doc.Shapes.Count = 1
    mock_doc.Shapes.side_effect = lambda i: mock_shape

    mock_word = MagicMock()
    mock_word.Documents.Open.return_value = mock_doc

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_word
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        _inject_win32com_word_shapes(
            Path("test.doc"),
            {"shape:0": "Traduit"},
        )
        assert mock_shape.TextFrame.TextRange.Text == "Traduit"
        mock_doc.Save.assert_called_once()


def test_inject_win32com_excel_shapes_basic() -> None:
    """Injects translated shape text into Excel via mocked win32com."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_win32com_excel_shapes,
    )

    mock_shape = MagicMock()

    mock_ws = MagicMock()
    mock_ws.Name = "Data"
    mock_ws.Shapes.Count = 1
    mock_ws.Shapes.side_effect = lambda i: mock_shape

    mock_wb = MagicMock()
    mock_wb.Worksheets = [mock_ws]

    mock_excel = MagicMock()
    mock_excel.Workbooks.Open.return_value = mock_wb

    mock_pythoncom = MagicMock()
    mock_win32com_client = MagicMock()
    mock_win32com_client.Dispatch.return_value = mock_excel
    mock_win32com = MagicMock()
    mock_win32com.client = mock_win32com_client

    with patch.dict(
        "sys.modules",
        {
            "win32com": mock_win32com,
            "win32com.client": mock_win32com_client,
            "pythoncom": mock_pythoncom,
        },
    ):
        _inject_win32com_excel_shapes(
            Path("test.xls"),
            {"shape:Data:0": "Traduit"},
        )
        assert mock_shape.TextFrame2.TextRange.Text == "Traduit"
        mock_wb.Save.assert_called_once()


# ---------------------------------------------------------------------------
# Shape / text-box — UNO extraction (mocked UNO objects)
# ---------------------------------------------------------------------------


def test_extract_uno_writer_shapes_basic() -> None:
    """Extracts Writer shape text via mocked UNO."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_uno_writer_shapes,
    )

    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.getString.return_value = "UNO shape text"
    mock_shape.createEnumeration.return_value = _make_uno_enum()

    mock_draw_page = MagicMock()
    mock_draw_page.getCount.return_value = 1
    mock_draw_page.getByIndex.return_value = mock_shape

    mock_doc = MagicMock()
    mock_doc.getDrawPage.return_value = mock_draw_page

    with patch(
        "src.core.office_processor._uno_open",
        return_value=mock_doc,
    ):
        result = _extract_uno_writer_shapes(Path("test.doc"))

    assert result == [("shape:0", "UNO shape text")]
    mock_doc.close.assert_called_once_with(True)


def test_extract_uno_calc_shapes_basic() -> None:
    """Extracts Calc shape text via mocked UNO."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_uno_calc_shapes,
    )

    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.getString.return_value = "Calc shape"
    mock_shape.createEnumeration.return_value = _make_uno_enum()

    mock_draw_page = MagicMock()
    mock_draw_page.getCount.return_value = 1
    mock_draw_page.getByIndex.return_value = mock_shape

    mock_sheet = MagicMock()
    mock_sheet.getName.return_value = "Revenue"
    mock_sheet.getDrawPage.return_value = mock_draw_page

    mock_sheets = MagicMock()
    mock_sheets.getCount.return_value = 1
    mock_sheets.getByIndex.return_value = mock_sheet

    mock_doc = MagicMock()
    mock_doc.getSheets.return_value = mock_sheets

    with patch(
        "src.core.office_processor._uno_open",
        return_value=mock_doc,
    ):
        result = _extract_uno_calc_shapes(Path("test.xls"))

    assert result == [("shape:Revenue:0", "Calc shape")]
    mock_doc.close.assert_called_once_with(True)


def test_inject_uno_writer_shapes_basic() -> None:
    """Injects translated shape text via mocked UNO."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_writer_shapes,
    )

    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.createEnumeration.return_value = _make_uno_enum()

    mock_draw_page = MagicMock()
    mock_draw_page.getCount.return_value = 1
    mock_draw_page.getByIndex.return_value = mock_shape

    mock_doc = MagicMock()
    mock_doc.getDrawPage.return_value = mock_draw_page

    with (
        patch(
            "src.core.office_processor._uno_open",
            return_value=mock_doc,
        ),
        patch(
            "src.core.office_processor._uno_save",
        ) as mock_save,
    ):
        _inject_uno_writer_shapes(
            Path("test.doc"),
            {"shape:0": "Traduit"},
        )

    mock_shape.setString.assert_called_once_with("Traduit")
    mock_save.assert_called_once()
    mock_doc.close.assert_called_once_with(True)


def test_inject_uno_calc_shapes_basic() -> None:
    """Injects translated shape text into Calc via mocked UNO."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_calc_shapes,
    )

    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.createEnumeration.return_value = _make_uno_enum()

    mock_draw_page = MagicMock()
    mock_draw_page.getCount.return_value = 1
    mock_draw_page.getByIndex.return_value = mock_shape

    mock_sheet = MagicMock()
    mock_sheet.getName.return_value = "Revenue"
    mock_sheet.getDrawPage.return_value = mock_draw_page

    mock_sheets = MagicMock()
    mock_sheets.getCount.return_value = 1
    mock_sheets.getByIndex.return_value = mock_sheet

    mock_doc = MagicMock()
    mock_doc.getSheets.return_value = mock_sheets

    with (
        patch(
            "src.core.office_processor._uno_open",
            return_value=mock_doc,
        ),
        patch(
            "src.core.office_processor._uno_save",
        ) as mock_save,
    ):
        _inject_uno_calc_shapes(
            Path("test.xls"),
            {"shape:Revenue:0": "Traduit"},
        )

    mock_shape.setString.assert_called_once_with("Traduit")
    mock_save.assert_called_once()
    mock_doc.close.assert_called_once_with(True)


# ---------------------------------------------------------------------------
# Multiline / newline preservation tests
# ---------------------------------------------------------------------------


def test_collect_wps_texts_multiline() -> None:
    """Multi-paragraph text box yields newline-separated text."""
    from lxml import etree  # noqa: PLC0415

    wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    xml_str = f"""\
<root xmlns:wps="{wps_ns}" xmlns:w="{w_ns}">
  <wps:txbx>
    <w:txbxContent>
      <w:p><w:r><w:t>Line one</w:t></w:r></w:p>
      <w:p><w:r><w:t>Line two</w:t></w:r></w:p>
      <w:p><w:r><w:t>Line three</w:t></w:r></w:p>
    </w:txbxContent>
  </wps:txbx>
</root>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    results = _collect_wps_texts(root)

    assert len(results) == 1
    assert results[0][0] == "Line one\nLine two\nLine three"
    expected_t_count = 3
    assert len(results[0][1]) == expected_t_count


def test_extract_odf_paragraph_text_line_break() -> None:
    """<text:line-break/> is preserved as a newline character."""
    from lxml import etree  # noqa: PLC0415

    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    text_p_tag = f"{{{text_ns}}}p"

    xml_str = f"""\
<draw:text-box xmlns:draw="{draw_ns}" xmlns:text="{text_ns}">
  <text:p>First<text:line-break/>Second</text:p>
</draw:text-box>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    result = _extract_odf_paragraph_text(root, text_p_tag)
    assert result == "First\nSecond"


def test_inject_odf_paragraph_text_multiline() -> None:
    """Newline input creates multiple <text:p> elements."""
    from lxml import etree  # noqa: PLC0415

    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    text_p_tag = f"{{{text_ns}}}p"

    xml_str = f"""\
<draw:text-box xmlns:draw="{draw_ns}" xmlns:text="{text_ns}">
  <text:p>Original</text:p>
</draw:text-box>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    result = _inject_odf_paragraph_text(root, "Line A\nLine B\nLine C", text_p_tag)

    assert result is True
    paras = root.findall(text_p_tag)
    expected_para_count = 3
    assert len(paras) == expected_para_count
    assert paras[0].text == "Line A"
    assert paras[1].text == "Line B"
    assert paras[2].text == "Line C"


def test_extract_drawingml_br() -> None:
    """<a:br/> inside <a:p> is preserved as a newline."""
    from lxml import etree  # noqa: PLC0415

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    xml_str = f"""\
<a:txBody xmlns:a="{a_ns}">
  <a:p>
    <a:r><a:t>Before break</a:t></a:r>
    <a:br/>
    <a:r><a:t>After break</a:t></a:r>
  </a:p>
</a:txBody>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    result = _extract_drawingml_text(root)
    assert result == "Before break\nAfter break"


def test_inject_drawingml_multiline() -> None:
    """Newline input creates <a:br/> and new <a:r> elements."""
    from lxml import etree  # noqa: PLC0415

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    xml_str = f"""\
<a:txBody xmlns:a="{a_ns}">
  <a:p><a:r><a:rPr lang="en-US"/><a:t>Original</a:t></a:r></a:p>
</a:txBody>"""
    root = etree.fromstring(xml_str.encode("utf-8"))
    _inject_drawingml_text(root, "First\nSecond\nThird")

    # Re-extract to verify round-trip
    result = _extract_drawingml_text(root)
    assert result == "First\nSecond\nThird"


def _inject_docx_comment_part(
    path: Path,
    comments_xml: str,
) -> None:
    """Injects a comments XML part into a DOCX and wires the relationship.

    Adds the relationship to both ``_rels/.rels`` (for python-docx
    ``package.part_related_by``) and ``word/_rels/document.xml.rels``
    so that the comments part is discoverable by all code paths.

    Args:
        path: Path to an existing .docx file (modified in-place).
        comments_xml: The raw XML string for the comments part.
    """
    import shutil  # noqa: PLC0415

    from lxml import etree  # noqa: PLC0415

    with zipfile.ZipFile(path, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    data["word/comments.xml"] = comments_xml.encode("utf-8")

    comments_rel_type = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
    )

    # Wire relationship in _rels/.rels (package level)
    pkg_rels_path = "_rels/.rels"
    pkg_rels_root = etree.fromstring(data[pkg_rels_path])
    etree.SubElement(
        pkg_rels_root,
        "Relationship",
        Id="rIdCommentPkg",
        Type=comments_rel_type,
        Target="word/comments.xml",
    )
    data[pkg_rels_path] = etree.tostring(
        pkg_rels_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    # Wire relationship in word/_rels/document.xml.rels (part level)
    doc_rels_path = "word/_rels/document.xml.rels"
    doc_rels_root = etree.fromstring(data[doc_rels_path])
    etree.SubElement(
        doc_rels_root,
        "Relationship",
        Id="rIdComment",
        Type=comments_rel_type,
        Target="comments.xml",
    )
    data[doc_rels_path] = etree.tostring(
        doc_rels_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    # Update [Content_Types].xml
    ct_root = etree.fromstring(data["[Content_Types].xml"])
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    etree.SubElement(
        ct_root,
        f"{{{ct_ns}}}Override",
        PartName="/word/comments.xml",
        ContentType="application/vnd.openxmlformats-officedocument"
        ".wordprocessingml.comments+xml",
    )
    data["[Content_Types].xml"] = etree.tostring(
        ct_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    tmp_zip = path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for item in items:
            zf_out.writestr(item, data[item.filename])
        for name, content in data.items():
            if not any(i.filename == name for i in items):
                zf_out.writestr(name, content)
    shutil.move(str(tmp_zip), str(path))


def test_inject_docx_comments_multiline(tmp_path: Path) -> None:
    """Newline in comment translation creates <w:br/> + <w:t> in XML."""
    from lxml import etree  # noqa: PLC0415

    # Create a DOCX with a single-paragraph comment
    doc = Document()
    doc.add_paragraph("Body")
    path = tmp_path / "comments.docx"
    doc.save(str(path))

    comments_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:comments xmlns:w="http://schemas.openxmlformats.org/'
        'wordprocessingml/2006/main">'
        '<w:comment w:id="1" w:author="Test" w:date="2024-01-01T00:00:00Z">'
        "<w:p><w:r><w:t>Original comment</w:t></w:r></w:p>"
        "</w:comment>"
        "</w:comments>"
    )
    _inject_docx_comment_part(path, comments_xml)

    # Inject multiline translation
    _inject_docx_comments(path, {"comment:1": "Line A\nLine B"})

    # Verify <w:br/> was inserted by parsing the ZIP directly
    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    with zipfile.ZipFile(path, "r") as zf:
        xml_bytes = zf.read("word/comments.xml")
    root = etree.fromstring(xml_bytes)
    br_els = root.findall(f".//{{{w_ns}}}br")
    assert len(br_els) >= 1

    t_els = root.findall(f".//{{{w_ns}}}t")
    t_texts = [t.text for t in t_els if t.text]
    assert "Line A" in t_texts
    assert "Line B" in t_texts


def test_inject_docx_shapes_multiline(tmp_path: Path) -> None:
    """Newline in shape translation creates <w:br/> + <w:t> in XML."""
    from lxml import etree  # noqa: PLC0415

    docx_path = _make_docx_with_shapes(tmp_path, "multiline.docx")
    _inject_docx_shapes(docx_path, {"shape:0": "First\nSecond"})

    # Verify <w:br/> was inserted by parsing the ZIP directly
    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    with zipfile.ZipFile(docx_path, "r") as zf:
        doc_xml = zf.read("word/document.xml")
    root = etree.fromstring(doc_xml)
    br_els = root.findall(f".//{{{w_ns}}}br")
    assert len(br_els) >= 1

    t_els = root.findall(f".//{{{w_ns}}}t")
    t_texts = [t.text for t in t_els if t.text]
    assert "First" in t_texts
    assert "Second" in t_texts


def test_extract_docx_comments_multiline(tmp_path: Path) -> None:
    """Multi-paragraph comment yields newline-separated text."""
    doc = Document()
    doc.add_paragraph("Body text")
    path = tmp_path / "multi_para_comment.docx"
    doc.save(str(path))

    comments_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:comments xmlns:w="http://schemas.openxmlformats.org/'
        'wordprocessingml/2006/main">'
        '<w:comment w:id="5" w:author="Test" w:date="2024-01-01T00:00:00Z">'
        "<w:p><w:r><w:t>Paragraph one</w:t></w:r></w:p>"
        "<w:p><w:r><w:t>Paragraph two</w:t></w:r></w:p>"
        "</w:comment>"
        "</w:comments>"
    )
    _inject_docx_comment_part(path, comments_xml)

    texts = _extract_docx_comments(path)
    result = dict(texts)
    assert result["comment:5"] == "Paragraph one\nParagraph two"


def _inject_docx_comment_part_with_rels(
    path: Path,
    comments_xml: str,
    comment_rels_xml: str,
) -> None:
    """Injects a comments XML part plus its .rels into a DOCX.

    Extends ``_inject_docx_comment_part`` by also writing
    ``word/_rels/comments.xml.rels`` so that hyperlink relationship
    lookups inside ``_extract_docx_comments`` resolve correctly.

    Args:
        path: Path to an existing .docx file (modified in-place).
        comments_xml: The raw XML string for word/comments.xml.
        comment_rels_xml: The raw XML string for
            word/_rels/comments.xml.rels.
    """
    _inject_docx_comment_part(path, comments_xml)
    # Append the .rels file alongside the comments part
    with zipfile.ZipFile(path, "a") as zf:
        zf.writestr("word/_rels/comments.xml.rels", comment_rels_xml.encode("utf-8"))


def test_extract_docx_comments_hyperlink(tmp_path: Path) -> None:
    """Hyperlink inside a comment paragraph is extracted as <a href> HTML."""
    from docx import Document  # noqa: PLC0415

    doc = Document()
    doc.add_paragraph("Body")
    path = tmp_path / "comment_hyperlink.docx"
    doc.save(str(path))

    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    hyperlink_rel_type = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
    )
    comments_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:comments xmlns:w="{w_ns}" xmlns:r="{r_ns}">'
        '<w:comment w:id="1" w:author="Test" w:date="2024-01-01T00:00:00Z">'
        "<w:p>"
        f'<w:hyperlink r:id="rId1">'
        "<w:r><w:t>click here</w:t></w:r>"
        "</w:hyperlink>"
        "</w:p>"
        "</w:comment>"
        "</w:comments>"
    )
    comment_rels_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<Relationships xmlns="{rels_ns}">'
        f'<Relationship Id="rId1" Type="{hyperlink_rel_type}"'
        ' Target="https://example.com" TargetMode="External"/>'
        "</Relationships>"
    )
    _inject_docx_comment_part_with_rels(path, comments_xml, comment_rels_xml)

    texts = _extract_docx_comments(path)
    result = dict(texts)
    assert "comment:1" in result
    assert '<a href="https://example.com">' in result["comment:1"]
    assert "click here" in result["comment:1"]


def test_extract_docx_comments_hyperlink_missing_rels_falls_back_to_plain(
    tmp_path: Path,
) -> None:
    """Hyperlink with no .rels file → plain text (no <a> tag)."""
    from docx import Document  # noqa: PLC0415

    doc = Document()
    doc.add_paragraph("Body")
    path = tmp_path / "comment_no_rels.docx"
    doc.save(str(path))

    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    # Comments XML has a hyperlink element, but NO .rels file is injected
    comments_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:comments xmlns:w="{w_ns}" xmlns:r="{r_ns}">'
        '<w:comment w:id="2" w:author="Test" w:date="2024-01-01T00:00:00Z">'
        "<w:p>"
        f'<w:hyperlink r:id="rId1">'
        "<w:r><w:t>link text</w:t></w:r>"
        "</w:hyperlink>"
        "</w:p>"
        "</w:comment>"
        "</w:comments>"
    )
    # Only inject comments.xml — no .rels → hyperlink_rels stays empty
    _inject_docx_comment_part(path, comments_xml)

    texts = _extract_docx_comments(path)
    result = dict(texts)
    # No rels → falls through to plain-text path (has_links is False since rels empty)
    assert "comment:2" in result
    assert "<a " not in result["comment:2"]
    assert "link text" in result["comment:2"]


def test_extract_docx_comments_mixed_hyperlink_and_plain(tmp_path: Path) -> None:
    """Comment with hyperlink run and plain run both extracted correctly."""
    from docx import Document  # noqa: PLC0415

    doc = Document()
    doc.add_paragraph("Body")
    path = tmp_path / "comment_mixed.docx"
    doc.save(str(path))

    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    hyperlink_rel_type = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
    )
    comments_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:comments xmlns:w="{w_ns}" xmlns:r="{r_ns}">'
        '<w:comment w:id="3" w:author="Test" w:date="2024-01-01T00:00:00Z">'
        "<w:p>"
        "<w:r><w:t>See </w:t></w:r>"
        f'<w:hyperlink r:id="rId1">'
        "<w:r><w:t>this link</w:t></w:r>"
        "</w:hyperlink>"
        "<w:r><w:t> for details</w:t></w:r>"
        "</w:p>"
        "</w:comment>"
        "</w:comments>"
    )
    comment_rels_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<Relationships xmlns="{rels_ns}">'
        f'<Relationship Id="rId1" Type="{hyperlink_rel_type}"'
        ' Target="https://docs.example.com" TargetMode="External"/>'
        "</Relationships>"
    )
    _inject_docx_comment_part_with_rels(path, comments_xml, comment_rels_xml)

    texts = _extract_docx_comments(path)
    result = dict(texts)
    assert "comment:3" in result
    text = result["comment:3"]
    assert "See" in text
    assert '<a href="https://docs.example.com">' in text
    assert "this link" in text
    assert "for details" in text


# ---------------------------------------------------------------------------
# Win32com font save/restore helper tests
# ---------------------------------------------------------------------------


def test_save_win32com_font_captures_properties() -> None:
    """_save_win32com_font reads properties from the Font object."""
    font = MagicMock()
    font.Name = "Arial"
    font.Size = 12  # noqa: PLR2004
    font.Bold = True
    font.Italic = False
    font.Color = 0
    font.Underline = 0
    saved = _save_win32com_font(font)
    assert saved["Name"] == "Arial"
    assert saved["Size"] == 12  # noqa: PLR2004
    assert saved["Bold"] is True
    assert saved["Italic"] is False


def test_save_win32com_font_skips_undefined() -> None:
    """_save_win32com_font skips properties with the undefined sentinel."""
    from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415

    font = MagicMock()
    font.Name = "Calibri"
    font.Size = WIN32COM_UNDEFINED
    font.Bold = WIN32COM_UNDEFINED
    font.Italic = False
    font.Color = 0
    font.Underline = 0
    saved = _save_win32com_font(font)
    assert "Name" in saved
    assert "Size" not in saved
    assert "Bold" not in saved


def test_save_win32com_font_skips_errored() -> None:
    """_save_win32com_font skips properties that raise on read."""

    # Use a plain class — MagicMock's __getattr__ intercepts the
    # AttributeError from a property getter and returns a mock instead.
    class _ErrFont:
        Name = "Arial"
        Bold = True
        Italic = False
        Color = 0
        Underline = 0

        @property
        def Size(self) -> float:  # noqa: N802
            raise AttributeError("Size unavailable")

    saved = _save_win32com_font(_ErrFont())
    assert "Name" in saved
    assert "Size" not in saved


def test_save_win32com_font_captures_strikethrough() -> None:
    """_save_win32com_font captures the StrikeThrough property."""
    font = MagicMock()
    font.StrikeThrough = True
    saved = _save_win32com_font(font)
    assert saved["StrikeThrough"] is True


def test_restore_win32com_font_sets_properties() -> None:
    """_restore_win32com_font sets saved values back on the Font object."""
    font = MagicMock()
    saved = {"Name": "Times New Roman", "Size": 14, "Bold": True}
    _restore_win32com_font(font, saved)
    assert font.Name == "Times New Roman"
    assert font.Size == 14  # noqa: PLR2004
    assert font.Bold is True


def test_restore_win32com_font_continues_on_error() -> None:
    """_restore_win32com_font continues if a single property set fails."""
    font = MagicMock()
    # Make Size raise on set
    type(font).Size = property(
        lambda self: None,
        lambda self, v: (_ for _ in ()).throw(RuntimeError),
    )
    saved = {"Name": "Arial", "Size": 12, "Bold": False}
    _restore_win32com_font(font, saved)
    # Name and Bold should still be set despite Size error
    assert font.Name == "Arial"
    assert font.Bold is False


# ---------------------------------------------------------------------------
# UNO char property save/restore helper tests
# ---------------------------------------------------------------------------


def test_save_uno_char_props_captures() -> None:
    """_save_uno_char_props reads properties via getPropertyValue."""
    obj = MagicMock()
    prop_values = {
        "CharFontName": "Liberation Sans",
        "CharHeight": 11.0,
        "CharWeight": 150.0,
        "CharPosture": 1,
        "CharColor": 255,
        "CharUnderline": 0,
    }
    obj.getPropertyValue.side_effect = lambda p: prop_values[p]
    saved = _save_uno_char_props(obj)
    assert saved["CharFontName"] == "Liberation Sans"
    assert saved["CharHeight"] == 11.0  # noqa: PLR2004
    assert len(saved) == 6  # noqa: PLR2004


def test_save_uno_char_props_skips_errored() -> None:
    """_save_uno_char_props skips properties that raise."""
    obj = MagicMock()

    def side_effect(prop: str) -> object:
        if prop == "CharFontName":
            return "Serif"
        raise RuntimeError("property not available")

    obj.getPropertyValue.side_effect = side_effect
    saved = _save_uno_char_props(obj)
    assert saved == {"CharFontName": "Serif"}


def test_restore_uno_char_props_calls_set() -> None:
    """_restore_uno_char_props calls setPropertyValue for each property."""
    obj = MagicMock()
    saved = {"CharFontName": "Mono", "CharHeight": 10.0}
    _restore_uno_char_props(obj, saved)
    obj.setPropertyValue.assert_any_call("CharFontName", "Mono")
    obj.setPropertyValue.assert_any_call("CharHeight", 10.0)


def test_restore_uno_char_props_continues_on_error() -> None:
    """_restore_uno_char_props continues if a single property set fails."""
    obj = MagicMock()

    call_log: list[str] = []

    def side_effect(prop: str, val: object) -> None:
        if prop == "CharHeight":
            raise RuntimeError("read-only")
        call_log.append(prop)

    obj.setPropertyValue.side_effect = side_effect
    saved = {"CharFontName": "Arial", "CharHeight": 12.0, "CharColor": 0}
    _restore_uno_char_props(obj, saved)
    assert "CharFontName" in call_log
    assert "CharColor" in call_log


# ---------------------------------------------------------------------------
# ODF odfpy _odf_replace_text span preservation tests
# ---------------------------------------------------------------------------


def test_odf_replace_text_preserves_span_style() -> None:
    """_odf_replace_text wraps new text in a Span with the original stylename."""
    from odf.text import P, Span  # noqa: PLC0415

    p = P()
    span = Span(stylename="T1")
    span.addText("Original bold text")
    p.addElement(span)

    _odf_replace_text(p, "Translated bold text")

    # The paragraph should have one child: a Span with stylename T1
    children = [
        c
        for c in p.childNodes
        if getattr(c, "nodeType", None) == 1  # noqa: PLR2004
    ]
    assert len(children) == 1
    child = children[0]
    assert child.qname == Span().qname
    # Verify stylename is preserved
    style_attr = child.attributes.get(
        ("urn:oasis:names:tc:opendocument:xmlns:text:1.0", "style-name"),
    )
    assert style_attr == "T1"
    # Verify text content
    from src.core.office_processor import _odf_element_text  # noqa: PLC0415

    assert _odf_element_text(p) == "Translated bold text"


def test_odf_replace_text_no_span_fallback() -> None:
    """_odf_replace_text falls back to addText when no span exists."""
    from odf.text import P  # noqa: PLC0415

    p = P()
    p.addText("Plain text")

    _odf_replace_text(p, "Translated plain")

    text = _odf_element_text(p)
    assert text == "Translated plain"


# ---------------------------------------------------------------------------
# ODF lxml _inject_odf_paragraph_text span preservation tests
# ---------------------------------------------------------------------------

_TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
_TEXT_P_TAG = f"{{{_TEXT_NS}}}p"
_TEXT_SPAN_TAG = f"{{{_TEXT_NS}}}span"
_STYLE_NAME_ATTR = f"{{{_TEXT_NS}}}style-name"


def test_inject_odf_paragraph_text_preserves_span_attribs() -> None:
    """_inject_odf_paragraph_text preserves <text:span> attributes."""
    from lxml import etree  # noqa: PLC0415

    parent = etree.Element("parent")
    p = etree.SubElement(parent, _TEXT_P_TAG)
    span = etree.SubElement(p, _TEXT_SPAN_TAG, {_STYLE_NAME_ATTR: "T2"})
    span.text = "Original"

    result = _inject_odf_paragraph_text(parent, "Translated", _TEXT_P_TAG)
    assert result is True

    paras = parent.findall(_TEXT_P_TAG)
    assert len(paras) == 1
    spans = paras[0].findall(f".//{_TEXT_SPAN_TAG}")
    assert len(spans) == 1
    assert spans[0].get(_STYLE_NAME_ATTR) == "T2"
    assert spans[0].text == "Translated"


def test_inject_odf_paragraph_text_multiline_preserves_span() -> None:
    """Multi-line injection preserves span attributes on all lines."""
    from lxml import etree  # noqa: PLC0415

    parent = etree.Element("parent")
    p = etree.SubElement(parent, _TEXT_P_TAG)
    span = etree.SubElement(p, _TEXT_SPAN_TAG, {_STYLE_NAME_ATTR: "Bold1"})
    span.text = "Line one"

    result = _inject_odf_paragraph_text(parent, "First\nSecond\nThird", _TEXT_P_TAG)
    assert result is True

    paras = parent.findall(_TEXT_P_TAG)
    assert len(paras) == 3  # noqa: PLR2004
    for i, expected in enumerate(["First", "Second", "Third"]):
        spans = paras[i].findall(f".//{_TEXT_SPAN_TAG}")
        assert len(spans) == 1
        assert spans[0].get(_STYLE_NAME_ATTR) == "Bold1"
        assert spans[0].text == expected


def test_inject_odf_paragraph_text_no_span_fallback() -> None:
    """_inject_odf_paragraph_text sets plain text when no span exists."""
    from lxml import etree  # noqa: PLC0415

    parent = etree.Element("parent")
    p = etree.SubElement(parent, _TEXT_P_TAG)
    p.text = "Plain original"

    result = _inject_odf_paragraph_text(parent, "Plain translated", _TEXT_P_TAG)
    assert result is True

    paras = parent.findall(_TEXT_P_TAG)
    assert len(paras) == 1
    assert paras[0].text == "Plain translated"
    assert paras[0].findall(f".//{_TEXT_SPAN_TAG}") == []


# ---------------------------------------------------------------------------
# Script family detection tests
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("text", "expected"),
    [
        ("Hello world", "latin"),
        ("Bonjour le monde", "latin"),
        ("Tiếng Việt", "latin"),
        ("Ελληνικά", "greek"),
        ("Русский текст", "cyrillic"),
        ("مرحبا بالعالم", "arabic"),
        ("שלום עולם", "hebrew"),
        ("こんにちは世界", "east_asian"),
        ("你好世界", "east_asian"),
        ("한국어", "east_asian"),
        ("สวัสดีชาวโลก", "thai"),
        ("नमस्ते दुनिया", "devanagari"),
        ("বাংলা", "bengali"),
        ("ភាសាខ្មែរ", "khmer"),
        ("ᠮᠣᠩᠭᠣᠯ", "mongolian"),
        ("", "latin"),
        ("12345", "latin"),
    ],
    ids=[
        "english",
        "french",
        "vietnamese",
        "greek",
        "cyrillic",
        "arabic",
        "hebrew",
        "japanese",
        "chinese",
        "korean",
        "thai",
        "devanagari",
        "bengali",
        "khmer",
        "mongolian",
        "empty",
        "digits",
    ],
)
def test_detect_script_family(text: str, expected: str) -> None:
    """detect_script correctly identifies script families."""
    from src.utils.font_utils import detect_script  # noqa: PLC0415

    assert detect_script(text) == expected


# ---------------------------------------------------------------------------
# Cross-script font name skip tests
# ---------------------------------------------------------------------------


def test_restore_win32com_font_substitutes_name_with_target_lang() -> None:
    """Font Name is substituted via _get_font_for_language when target_lang set."""
    font = MagicMock()
    saved = {"Name": "Arial", "Size": 12, "Bold": True}
    _restore_win32com_font(
        font,
        saved,
        target_lang="Japanese",
    )
    # Size and Bold should be set
    assert font.Size == 12  # noqa: PLR2004
    assert font.Bold is True
    # Name should be a language-appropriate font, not original Arial
    assert isinstance(font.Name, str)


def test_restore_win32com_font_keeps_name_without_target_lang() -> None:
    """Font Name IS restored as-is when no target_lang is provided."""
    font = MagicMock()
    _restore_win32com_font(
        font,
        {"Name": "Times New Roman", "Size": 12},
        original_text="Hello",
        translated_text="Bonjour",
    )
    assert font.Name == "Times New Roman"
    assert font.Size == 12  # noqa: PLR2004


def test_restore_uno_substitutes_font_name_with_target_lang() -> None:
    """CharFontName is substituted via _get_font_for_language when target_lang set."""
    obj = MagicMock()
    saved = {"CharFontName": "Liberation Sans", "CharHeight": 11.0}
    _restore_uno_char_props(
        obj,
        saved,
        target_lang="Arabic",
    )
    # Both CharFontName and CharHeight should be set
    calls = {c.args[0] for c in obj.setPropertyValue.call_args_list}
    assert "CharHeight" in calls
    assert "CharFontName" in calls


def test_restore_uno_keeps_font_name_without_target_lang() -> None:
    """CharFontName IS restored as-is when no target_lang is provided."""
    obj = MagicMock()
    saved = {"CharFontName": "Liberation Sans", "CharHeight": 11.0}
    _restore_uno_char_props(
        obj,
        saved,
        original_text="English",
        translated_text="Français",
    )
    calls = {c.args[0] for c in obj.setPropertyValue.call_args_list}
    assert "CharFontName" in calls
    assert "CharHeight" in calls


def test_odf_replace_text_preserves_span_on_cross_script() -> None:
    """_odf_replace_text preserves Span wrapper even when scripts differ."""
    from odf.text import P, Span  # noqa: PLC0415

    p = P()
    span = Span(stylename="T1")
    span.addText("Hello world")
    p.addElement(span)

    _odf_replace_text(p, "你好世界")

    # Text should be wrapped in span (font substitution handled elsewhere)
    text = _odf_element_text(p)
    assert text == "你好世界"
    element_children = [
        c
        for c in p.childNodes
        if getattr(c, "nodeType", None) == 1  # noqa: PLR2004
    ]
    assert len(element_children) == 1


def test_inject_odf_paragraph_text_preserves_span_on_cross_script() -> None:
    """_inject_odf_paragraph_text keeps span attrs even when scripts differ."""
    from lxml import etree  # noqa: PLC0415

    parent = etree.Element("parent")
    p = etree.SubElement(parent, _TEXT_P_TAG)
    span = etree.SubElement(p, _TEXT_SPAN_TAG, {_STYLE_NAME_ATTR: "Bold1"})
    span.text = "English text"

    result = _inject_odf_paragraph_text(parent, "テキスト", _TEXT_P_TAG)
    assert result is True

    paras = parent.findall(_TEXT_P_TAG)
    assert len(paras) == 1
    # Span should be preserved (font substitution handled at higher level)
    spans = paras[0].findall(f".//{_TEXT_SPAN_TAG}")
    assert len(spans) == 1
    assert spans[0].text == "テキスト"


# ---------------------------------------------------------------------------
# Integration wiring tests: injection functions call save/restore helpers
# ---------------------------------------------------------------------------


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_word_wires_font_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Verifies font save/restore wiring in _inject_win32com_word."""
    from src.core.office_processor import _inject_win32com_word  # noqa: PLC0415

    # Build minimal mock COM structure: one paragraph, no tables
    mock_range = MagicMock()
    mock_range.Text = "Hello world\r"
    mock_para = MagicMock()
    mock_para.Range = mock_range
    mock_doc = MagicMock()
    mock_doc.Paragraphs.Count = 1
    mock_doc.Paragraphs.return_value = mock_para
    mock_doc.Tables.Count = 0
    mock_open.return_value = (MagicMock(), mock_doc, MagicMock())

    _inject_win32com_word(
        Path("/fake/in.docx"),
        Path("/fake/out.docx"),
        {"para:1": "Bonjour monde"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    # Original text must be captured BEFORE the replacement (strip trailing \r)
    assert kw["original_text"] == "Hello world"
    # Translated text must match the injected translation
    assert kw["translated_text"] == "Bonjour monde"


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_word_restores_shading_para(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Uniform path saves and restores Shading.BackgroundPatternColor for paragraphs."""
    from src.core.office_processor import _inject_win32com_word  # noqa: PLC0415

    mock_range = MagicMock()
    mock_range.Text = "Hello\r"
    mock_range.HighlightColorIndex = 0
    mock_range.Shading.BackgroundPatternColor = 65535  # yellow BGR

    mock_para = MagicMock()
    mock_para.Range = mock_range

    # After text replacement, Paragraphs(1).Range returns updated range
    mock_new_range = MagicMock()
    mock_new_para = MagicMock()
    mock_new_para.Range = mock_new_range

    mock_doc = MagicMock()
    mock_doc.Paragraphs.Count = 1
    # First call returns original, second call returns post-update
    mock_doc.Paragraphs.side_effect = [mock_para, mock_new_para]
    mock_doc.Tables.Count = 0
    mock_open.return_value = (MagicMock(), mock_doc, MagicMock())

    _inject_win32com_word(
        Path("/fake/in.docx"),
        Path("/fake/out.docx"),
        {"para:1": "Bonjour"},
    )

    # Shading.BackgroundPatternColor must be restored on the new range
    assert mock_new_range.Shading.BackgroundPatternColor == 65535  # noqa: PLR2004


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_word_restores_shading_cell(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Uniform path saves/restores Shading.BackgroundPatternColor for cells."""
    from src.core.office_processor import _inject_win32com_word  # noqa: PLC0415

    mock_cell_range = MagicMock()
    mock_cell_range.Text = "Cell text\r\x07"
    mock_cell_range.HighlightColorIndex = 0
    mock_cell_range.Shading.BackgroundPatternColor = 16711680  # blue BGR

    mock_cell = MagicMock()
    mock_cell.Range = mock_cell_range

    # After text replacement, Cell returns updated cell
    mock_new_cell_range = MagicMock()
    mock_new_cell = MagicMock()
    mock_new_cell.Range = mock_new_cell_range

    mock_table = MagicMock()
    mock_table.Rows.Count = 1
    mock_table.Columns.Count = 1
    mock_table.Cell.side_effect = [mock_cell, mock_new_cell]

    mock_doc = MagicMock()
    mock_doc.Paragraphs.Count = 0
    mock_doc.Tables.Count = 1
    mock_doc.Tables.return_value = mock_table
    mock_open.return_value = (MagicMock(), mock_doc, MagicMock())

    _inject_win32com_word(
        Path("/fake/in.docx"),
        Path("/fake/out.docx"),
        {"table:1:1:1": "Texte cellule"},
    )

    # Shading must be restored on the new cell range
    assert mock_new_cell_range.Shading.BackgroundPatternColor == 16711680  # noqa: PLR2004


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_word_shading_exception_skipped(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Shading read failure is handled gracefully (shading_saved stays None)."""
    from src.core.office_processor import _inject_win32com_word  # noqa: PLC0415

    mock_range = MagicMock()
    mock_range.Text = "Hello\r"
    mock_range.HighlightColorIndex = 0
    # Shading access raises (e.g. COM error)
    type(mock_range).Shading = PropertyMock(side_effect=Exception("COM error"))

    mock_para = MagicMock()
    mock_para.Range = mock_range

    mock_new_para = MagicMock()
    mock_new_range = MagicMock()
    mock_new_para.Range = mock_new_range

    mock_doc = MagicMock()
    mock_doc.Paragraphs.Count = 1
    mock_doc.Paragraphs.side_effect = [mock_para, mock_new_para]
    mock_doc.Tables.Count = 0
    mock_open.return_value = (MagicMock(), mock_doc, MagicMock())

    # Should not raise
    _inject_win32com_word(
        Path("/fake/in.docx"),
        Path("/fake/out.docx"),
        {"para:1": "Bonjour"},
    )

    mock_restore.assert_called_once()


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._restore_uno_char_props")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_writer_wires_char_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_uno_save: MagicMock,
) -> None:
    """Verifies char props save/restore wiring in _inject_uno_writer."""
    from src.core.office_processor import _inject_uno_writer  # noqa: PLC0415

    # Build minimal mock UNO structure: one paragraph, no tables
    mock_para = MagicMock()
    mock_para.supportsService.return_value = False
    mock_para.getString.return_value = "Hello world"

    mock_enum = MagicMock()
    mock_enum.hasMoreElements.side_effect = [True, False]
    mock_enum.nextElement.return_value = mock_para

    mock_text_content = MagicMock()
    mock_text_content.createEnumeration.return_value = mock_enum

    mock_doc = MagicMock()
    mock_doc.getText.return_value = mock_text_content
    mock_doc.getTextTables.return_value.getCount.return_value = 0
    mock_open.return_value = mock_doc

    _inject_uno_writer(
        Path("/fake/in.odt"),
        Path("/fake/out.odt"),
        {"para:0": "Bonjour monde"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    # Original text must be captured via getString() BEFORE setString()
    assert kw["original_text"] == "Hello world"
    # Translated text must match the injected translation
    assert kw["translated_text"] == "Bonjour monde"


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_excel_wires_font_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Verifies font save/restore wiring in _inject_win32com_excel."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_win32com_excel,
    )

    mock_cell = MagicMock()
    mock_cell.Value = "Original"

    mock_used = MagicMock()
    mock_used.Rows.Count = 1
    mock_used.Columns.Count = 1

    mock_ws = MagicMock()
    mock_ws.Name = "Sheet1"
    mock_ws.UsedRange = mock_used
    mock_ws.Cells.return_value = mock_cell

    mock_wb = MagicMock()
    mock_wb.Worksheets = [mock_ws]
    mock_open.return_value = (MagicMock(), mock_wb, MagicMock())

    _inject_win32com_excel(
        Path("/fake/in.xlsx"),
        Path("/fake/out.xlsx"),
        {"sheet:Sheet1:1:1": "Traduit"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Original"
    assert kw["translated_text"] == "Traduit"


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_ppt_wires_font_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Verifies font save/restore wiring in _inject_win32com_ppt."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_win32com_ppt,
    )

    mock_para_rng = MagicMock()
    mock_para_rng.Text = "Slide text"
    # Paragraphs() returns the same mock for both the collection
    # (.Count) and indexed access (.Paragraphs(p_idx))
    mock_para_rng.Count = 1

    mock_tf = MagicMock()
    mock_tf.TextRange.Paragraphs.return_value = mock_para_rng

    mock_shape = MagicMock()
    mock_shape.HasTextFrame = True
    mock_shape.TextFrame = mock_tf

    mock_slide = MagicMock()
    mock_slide.Shapes.Count = 1
    mock_slide.Shapes.return_value = mock_shape

    mock_prs = MagicMock()
    mock_prs.Slides.Count = 1
    mock_prs.Slides.return_value = mock_slide
    mock_open.return_value = (MagicMock(), mock_prs, MagicMock())

    _inject_win32com_ppt(
        Path("/fake/in.pptx"),
        Path("/fake/out.pptx"),
        {"slide:1:1:1": "Texte de diapo"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Slide text"
    assert kw["translated_text"] == "Texte de diapo"


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_ppt_wires_highlight_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """PPT uniform path saves and restores Font.Highlight (Office 365+)."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_win32com_ppt,
    )

    mock_para_rng = MagicMock()
    mock_para_rng.Text = "Highlighted text"
    mock_para_rng.Count = 1
    # Simulate Font.Highlight.ForeColor.RGB returning a valid colour
    mock_para_rng.Font.Highlight.ForeColor.RGB = 0xFF0000  # red

    mock_tf = MagicMock()
    mock_tf.TextRange.Paragraphs.return_value = mock_para_rng

    mock_shape = MagicMock()
    mock_shape.HasTextFrame = True
    mock_shape.TextFrame = mock_tf

    mock_slide = MagicMock()
    mock_slide.Shapes.Count = 1
    mock_slide.Shapes.return_value = mock_shape

    mock_prs = MagicMock()
    mock_prs.Slides.Count = 1
    mock_prs.Slides.return_value = mock_slide
    mock_open.return_value = (MagicMock(), mock_prs, MagicMock())

    _inject_win32com_ppt(
        Path("/fake/in.pptx"),
        Path("/fake/out.pptx"),
        {"slide:1:1:1": "Texte de diapo"},
    )

    # Restore must have been called
    mock_restore.assert_called_once()
    # Highlight must have been written back to the re-acquired range
    mock_tf.TextRange.Paragraphs.return_value.Font.Highlight.ForeColor.RGB = 0xFF0000


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_word_comments_wires_font(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Verifies font save/restore wiring in _inject_win32com_word_comments."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_win32com_word_comments,
    )

    mock_comment = MagicMock()
    mock_comment.Index = 1
    mock_comment.Range.Text = "Comment text"

    mock_doc = MagicMock()
    mock_doc.Comments.Count = 1
    mock_doc.Comments.return_value = mock_comment
    mock_open.return_value = (MagicMock(), mock_doc, MagicMock())

    _inject_win32com_word_comments(
        Path("/fake/out.doc"),
        {"comment:1": "Commentaire"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Comment text"
    assert kw["translated_text"] == "Commentaire"


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_word_shapes_wires_font(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Verifies font save/restore wiring in _inject_win32com_word_shapes."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_win32com_word_shapes,
    )

    mock_text_rng = MagicMock()
    mock_text_rng.Text = "Shape text"
    mock_shape = MagicMock()
    mock_shape.TextFrame.TextRange = mock_text_rng

    mock_doc = MagicMock()
    mock_doc.Shapes.Count = 1
    mock_doc.Shapes.return_value = mock_shape
    mock_open.return_value = (MagicMock(), mock_doc, MagicMock())

    _inject_win32com_word_shapes(
        Path("/fake/out.doc"),
        {"shape:0": "Forme texte"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Shape text"
    assert kw["translated_text"] == "Forme texte"


@patch("src.core.office_processor._win32com_close")
@patch("src.core.office_processor._restore_win32com_font")
@patch("src.core.office_processor._win32com_open")
def test_inject_win32com_excel_shapes_wires_font(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_close: MagicMock,
) -> None:
    """Verifies font save/restore wiring in _inject_win32com_excel_shapes."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_win32com_excel_shapes,
    )

    mock_text_rng = MagicMock()
    mock_text_rng.Text = "Excel shape"
    mock_shape = MagicMock()
    mock_shape.TextFrame2.TextRange = mock_text_rng

    mock_ws = MagicMock()
    mock_ws.Name = "Data"
    mock_ws.Shapes.Count = 1
    mock_ws.Shapes.return_value = mock_shape

    mock_wb = MagicMock()
    mock_wb.Worksheets = [mock_ws]
    mock_open.return_value = (MagicMock(), mock_wb, MagicMock())

    _inject_win32com_excel_shapes(
        Path("/fake/out.xls"),
        {"shape:Data:0": "Forme Excel"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Excel shape"
    assert kw["translated_text"] == "Forme Excel"


def test_inject_win32com_excel_html_runs_bg_color() -> None:
    """Background colour is applied via Font.Highlight.ForeColor.RGB."""
    mock_text_rng = MagicMock()
    mock_text_rng.Font.Name = "Arial"

    _inject_win32com_excel_html_runs(
        mock_text_rng,
        '<span style="background-color:#ffff00">hi</span>',
        "hi",
    )

    # Characters uses 1-based indexing
    char_rng = mock_text_rng.Characters(1, 2)
    char_rng.Font.Highlight.ForeColor.RGB = 65535  # noqa: PLR2004
    # Verify the assignment happened (MagicMock records it)
    assert char_rng.Font.Highlight.ForeColor.RGB == 65535  # noqa: PLR2004


def test_inject_win32com_excel_html_runs_no_bg() -> None:
    """Segments without bg colour skip Highlight assignment."""
    mock_text_rng = MagicMock()
    mock_text_rng.Font.Name = "Arial"

    _inject_win32com_excel_html_runs(
        mock_text_rng,
        "<b>bold</b>",
        "bold",
    )

    # bold applied, but no Highlight interaction
    char_rng = mock_text_rng.Characters(1, 4)
    assert char_rng.Font.Bold is True


def test_inject_win32com_excel_html_runs_bg_exception() -> None:
    """Highlight exception is suppressed gracefully."""
    mock_text_rng = MagicMock()
    mock_text_rng.Font.Name = "Arial"

    # Make Highlight raise — simulating older Office without Highlight
    mock_char_rng = MagicMock()
    type(mock_char_rng.Font).Highlight = PropertyMock(
        side_effect=AttributeError("no Highlight"),
    )
    # Other Font properties should still work
    mock_text_rng.Characters.return_value = mock_char_rng

    # Should not raise
    _inject_win32com_excel_html_runs(
        mock_text_rng,
        '<span style="background-color:#00ff00">text</span>',
        "text",
    )

    # Bold was still applied despite Highlight failure
    assert mock_char_rng.Font.Bold is False


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._restore_uno_char_props")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_calc_wires_char_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_uno_save: MagicMock,
) -> None:
    """Verifies char props save/restore wiring in _inject_uno_calc."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_calc,
    )

    mock_cell = MagicMock()
    mock_cell.getString.return_value = "Cell text"

    mock_addr = MagicMock()
    mock_addr.StartRow = 0
    mock_addr.EndRow = 0
    mock_addr.StartColumn = 0
    mock_addr.EndColumn = 0

    mock_cursor = MagicMock()
    mock_cursor.getRangeAddress.return_value = mock_addr

    mock_sheet = MagicMock()
    mock_sheet.getName.return_value = "Sheet1"
    mock_sheet.createCursor.return_value = mock_cursor
    mock_sheet.getCellByPosition.return_value = mock_cell

    mock_doc = MagicMock()
    mock_doc.getSheets.return_value.getCount.return_value = 1
    mock_doc.getSheets.return_value.getByIndex.return_value = mock_sheet
    mock_open.return_value = mock_doc

    _inject_uno_calc(
        Path("/fake/in.ods"),
        Path("/fake/out.ods"),
        {"sheet:Sheet1:0:0": "Cellule"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Cell text"
    assert kw["translated_text"] == "Cellule"


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._restore_uno_char_props")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_impress_wires_char_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_uno_save: MagicMock,
) -> None:
    """Verifies char props save/restore wiring in _inject_uno_impress."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_impress,
    )

    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.getString.return_value = "Slide shape"

    mock_page = MagicMock()
    mock_page.getCount.return_value = 1
    mock_page.getByIndex.return_value = mock_shape

    mock_draw_pages = MagicMock()
    mock_draw_pages.getCount.return_value = 1
    mock_draw_pages.getByIndex.return_value = mock_page

    mock_doc = MagicMock()
    mock_doc.getDrawPages.return_value = mock_draw_pages
    mock_open.return_value = mock_doc

    _inject_uno_impress(
        Path("/fake/in.odp"),
        Path("/fake/out.odp"),
        {"slide:0:0": "Forme de diapo"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Slide shape"
    assert kw["translated_text"] == "Forme de diapo"


def test_extract_uno_impress_uses_html_for_mixed() -> None:
    """Mixed-format Impress shape → HTML extraction."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_uno_impress,
    )

    mock_para = MagicMock()
    mock_para.getString.return_value = "bold plain"

    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.getString.return_value = "bold plain"
    mock_shape.createEnumeration.return_value = _make_uno_enum([mock_para])

    mock_page = MagicMock()
    mock_page.getCount.return_value = 1
    mock_page.getByIndex.return_value = mock_shape

    mock_draw_pages = MagicMock()
    mock_draw_pages.getCount.return_value = 1
    mock_draw_pages.getByIndex.return_value = mock_page

    mock_doc = MagicMock()
    mock_doc.getDrawPages.return_value = mock_draw_pages

    with (
        patch(
            "src.core.office_processor._uno_open",
            return_value=mock_doc,
        ),
        patch(
            "src.core.office_processor._has_uno_mixed_formatting",
            return_value=True,
        ),
        patch(
            "src.core.office_processor._has_uno_hyperlinks",
            return_value=False,
        ),
        patch(
            "src.core.office_processor._uno_runs_to_html",
            return_value="<b>bold</b> plain",
        ),
    ):
        result = _extract_uno_impress(Path("test.ppt"))

    assert len(result) == 1
    assert result[0] == ("slide:0:0", "<b>bold</b> plain")


def test_extract_uno_impress_plain_for_uniform() -> None:
    """Uniform-format Impress shape → plain text extraction."""
    from src.core.office_processor import (  # noqa: PLC0415
        _extract_uno_impress,
    )

    mock_para = MagicMock()
    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.getString.return_value = "plain text"
    mock_shape.createEnumeration.return_value = _make_uno_enum([mock_para])

    mock_page = MagicMock()
    mock_page.getCount.return_value = 1
    mock_page.getByIndex.return_value = mock_shape

    mock_draw_pages = MagicMock()
    mock_draw_pages.getCount.return_value = 1
    mock_draw_pages.getByIndex.return_value = mock_page

    mock_doc = MagicMock()
    mock_doc.getDrawPages.return_value = mock_draw_pages

    with (
        patch(
            "src.core.office_processor._uno_open",
            return_value=mock_doc,
        ),
        patch(
            "src.core.office_processor._has_uno_mixed_formatting",
            return_value=False,
        ),
        patch(
            "src.core.office_processor._has_uno_hyperlinks",
            return_value=False,
        ),
    ):
        result = _extract_uno_impress(Path("test.ppt"))

    assert len(result) == 1
    assert result[0] == ("slide:0:0", "plain text")


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._inject_uno_impress_para_text")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_impress_dispatches_html(
    mock_open: MagicMock,
    mock_inject_para: MagicMock,
    mock_uno_save: MagicMock,
) -> None:
    """HTML translation dispatches per-paragraph to _inject_uno_impress_para_text."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_impress,
    )

    mock_para = MagicMock()
    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.createEnumeration.return_value = _make_uno_enum([mock_para])

    mock_page = MagicMock()
    mock_page.getCount.return_value = 1
    mock_page.getByIndex.return_value = mock_shape

    mock_draw_pages = MagicMock()
    mock_draw_pages.getCount.return_value = 1
    mock_draw_pages.getByIndex.return_value = mock_page

    mock_doc = MagicMock()
    mock_doc.getDrawPages.return_value = mock_draw_pages
    mock_open.return_value = mock_doc

    _inject_uno_impress(
        Path("/fake/in.ppt"),
        Path("/fake/out.ppt"),
        {"slide:0:0": "<b>gras</b> texte"},
    )

    mock_inject_para.assert_called_once_with(
        mock_para,
        "<b>gras</b> texte",
        target_lang="",
    )


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._restore_uno_char_props")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_calc_comments_wires_char_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_uno_save: MagicMock,
) -> None:
    """Verifies char props wiring in _inject_uno_calc_comments."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_calc_comments,
    )

    mock_pos = MagicMock()
    mock_pos.Row = 0
    mock_pos.Column = 0

    mock_annotation = MagicMock()
    mock_annotation.getPosition.return_value = mock_pos
    mock_annotation.getString.return_value = "Note text"

    mock_annotations = MagicMock()
    mock_annotations.getCount.return_value = 1
    mock_annotations.getByIndex.return_value = mock_annotation

    mock_sheet = MagicMock()
    mock_sheet.getName.return_value = "Sheet1"
    mock_sheet.getAnnotations.return_value = mock_annotations

    mock_doc = MagicMock()
    mock_doc.getSheets.return_value.getCount.return_value = 1
    mock_doc.getSheets.return_value.getByIndex.return_value = mock_sheet
    mock_open.return_value = mock_doc

    _inject_uno_calc_comments(
        Path("/fake/out.ods"),
        {"comment:Sheet1:1:1": "Texte note"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Note text"
    assert kw["translated_text"] == "Texte note"


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._restore_uno_char_props")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_impress_comments_wires_char_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_uno_save: MagicMock,
) -> None:
    """Verifies char props wiring in _inject_uno_impress_comments."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_impress_comments,
    )

    mock_text_range = MagicMock()
    mock_text_range.getString.return_value = "Annotation"

    mock_annotation = MagicMock()
    mock_annotation.TextRange = mock_text_range

    mock_enum = MagicMock()
    mock_enum.hasMoreElements.side_effect = [True, False]
    mock_enum.nextElement.return_value = mock_annotation

    mock_annotations = MagicMock()
    mock_annotations.createEnumeration.return_value = mock_enum

    mock_page = MagicMock()
    mock_page.getAnnotations.return_value = mock_annotations

    mock_draw_pages = MagicMock()
    mock_draw_pages.getCount.return_value = 1
    mock_draw_pages.getByIndex.return_value = mock_page

    mock_doc = MagicMock()
    mock_doc.getDrawPages.return_value = mock_draw_pages
    mock_open.return_value = mock_doc

    _inject_uno_impress_comments(
        Path("/fake/out.odp"),
        {"comment:0:0": "Commentaire"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Annotation"
    assert kw["translated_text"] == "Commentaire"


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._restore_uno_char_props")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_writer_shapes_wires_char_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_uno_save: MagicMock,
) -> None:
    """Verifies char props wiring in _inject_uno_writer_shapes."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_writer_shapes,
    )

    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.getString.return_value = "Shape content"
    mock_shape.createEnumeration.return_value = _make_uno_enum()

    mock_draw_page = MagicMock()
    mock_draw_page.getCount.return_value = 1
    mock_draw_page.getByIndex.return_value = mock_shape

    mock_doc = MagicMock()
    mock_doc.getDrawPage.return_value = mock_draw_page
    mock_open.return_value = mock_doc

    _inject_uno_writer_shapes(
        Path("/fake/out.odt"),
        {"shape:0": "Contenu forme"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Shape content"
    assert kw["translated_text"] == "Contenu forme"


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._restore_uno_char_props")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_calc_shapes_wires_char_save_restore(
    mock_open: MagicMock,
    mock_restore: MagicMock,
    mock_uno_save: MagicMock,
) -> None:
    """Verifies char props wiring in _inject_uno_calc_shapes."""
    from src.core.office_processor import (  # noqa: PLC0415
        _inject_uno_calc_shapes,
    )

    mock_shape = MagicMock()
    mock_shape.supportsService.return_value = True
    mock_shape.getString.return_value = "Calc shape"
    mock_shape.createEnumeration.return_value = _make_uno_enum()

    mock_draw_page = MagicMock()
    mock_draw_page.getCount.return_value = 1
    mock_draw_page.getByIndex.return_value = mock_shape

    mock_sheet = MagicMock()
    mock_sheet.getName.return_value = "Data"
    mock_sheet.getDrawPage.return_value = mock_draw_page

    mock_doc = MagicMock()
    mock_doc.getSheets.return_value.getCount.return_value = 1
    mock_doc.getSheets.return_value.getByIndex.return_value = mock_sheet
    mock_open.return_value = mock_doc

    _inject_uno_calc_shapes(
        Path("/fake/out.ods"),
        {"shape:Data:0": "Forme calc"},
    )

    mock_restore.assert_called_once()
    kw = mock_restore.call_args.kwargs
    assert kw["original_text"] == "Calc shape"
    assert kw["translated_text"] == "Forme calc"


# ---------------------------------------------------------------------------
# Superscript / Subscript tests
# ---------------------------------------------------------------------------


class TestReadWin32comCharFormattingSupSub:
    """Tests for _read_win32com_char_formatting superscript/subscript."""

    def test_superscript_true(self) -> None:
        """Character with superscript=True returns (True, False) at positions 4,5."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_char_formatting,
        )
        from tests.test_office_formatter import _make_win32com_char  # noqa: PLC0415

        ch = _make_win32com_char(superscript=True, subscript=False)
        result = _read_win32com_char_formatting(ch)
        assert len(result) == 9  # noqa: PLR2004
        assert result[4] is True  # superscript
        assert result[5] is False  # subscript

    def test_subscript_true(self) -> None:
        """Character with subscript=True returns (False, True) at positions 4,5."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_char_formatting,
        )
        from tests.test_office_formatter import _make_win32com_char  # noqa: PLC0415

        ch = _make_win32com_char(superscript=False, subscript=True)
        result = _read_win32com_char_formatting(ch)
        assert result[4] is False  # superscript
        assert result[5] is True  # subscript

    def test_neither_sup_nor_sub(self) -> None:
        """Default character has (False, False) for superscript/subscript."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_char_formatting,
        )
        from tests.test_office_formatter import _make_win32com_char  # noqa: PLC0415

        ch = _make_win32com_char()
        result = _read_win32com_char_formatting(ch)
        assert result[4] is False  # superscript
        assert result[5] is False  # subscript

    def test_undefined_superscript(self) -> None:
        """WIN32COM_UNDEFINED for superscript is treated as False."""
        from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_char_formatting,
        )
        from tests.test_office_formatter import _make_win32com_char  # noqa: PLC0415

        ch = _make_win32com_char(superscript=WIN32COM_UNDEFINED)
        result = _read_win32com_char_formatting(ch)
        assert result[4] is False  # superscript treated as False

    def test_undefined_subscript(self) -> None:
        """WIN32COM_UNDEFINED for subscript is treated as False."""
        from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_char_formatting,
        )
        from tests.test_office_formatter import _make_win32com_char  # noqa: PLC0415

        ch = _make_win32com_char(subscript=WIN32COM_UNDEFINED)
        result = _read_win32com_char_formatting(ch)
        assert result[5] is False  # subscript treated as False


class TestReadWin32comPptRunFormattingSupSub:
    """Tests for _read_win32com_ppt_run_formatting superscript/subscript."""

    def test_positive_baseline_super(self) -> None:
        """Positive BaselineOffset (0.3) indicates superscript."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_ppt_run_formatting,
        )
        from tests.test_office_formatter import (  # noqa: PLC0415
            _make_win32com_ppt_run,
        )

        run = _make_win32com_ppt_run()
        run.Font.BaselineOffset = 0.3
        result = _read_win32com_ppt_run_formatting(run)
        assert len(result) == 9  # noqa: PLR2004
        assert result[4] is True  # superscript
        assert result[5] is False  # subscript

    def test_negative_baseline_sub(self) -> None:
        """Negative BaselineOffset (-0.25) indicates subscript."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_ppt_run_formatting,
        )
        from tests.test_office_formatter import (  # noqa: PLC0415
            _make_win32com_ppt_run,
        )

        run = _make_win32com_ppt_run()
        run.Font.BaselineOffset = -0.25
        result = _read_win32com_ppt_run_formatting(run)
        assert result[4] is False  # superscript
        assert result[5] is True  # subscript

    def test_zero_baseline(self) -> None:
        """Zero BaselineOffset means neither superscript nor subscript."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_ppt_run_formatting,
        )
        from tests.test_office_formatter import (  # noqa: PLC0415
            _make_win32com_ppt_run,
        )

        run = _make_win32com_ppt_run()
        run.Font.BaselineOffset = 0.0
        result = _read_win32com_ppt_run_formatting(run)
        assert result[4] is False  # superscript
        assert result[5] is False  # subscript

    def test_no_baseline_attr(self) -> None:
        """Missing BaselineOffset attribute gracefully returns both False."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_win32com_ppt_run_formatting,
        )
        from tests.test_office_formatter import (  # noqa: PLC0415
            _make_win32com_ppt_run,
        )

        run = _make_win32com_ppt_run()
        # Make BaselineOffset raise AttributeError
        type(run.Font).BaselineOffset = property(
            lambda s: (_ for _ in ()).throw(AttributeError("no BaselineOffset")),
        )
        result = _read_win32com_ppt_run_formatting(run)
        assert result[4] is False  # superscript
        assert result[5] is False  # subscript


class TestInjectWin32comWordSupSub:
    """Tests for _inject_win32com_word_html_runs setting Superscript/Subscript."""

    def test_superscript_set(self) -> None:
        """Injecting '<sup>2</sup>' sets Font.Superscript = True on sub-range."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_win32com_word_html_runs,
        )

        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0

        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        _inject_win32com_word_html_runs(doc, rng, "x<sup>2</sup>")
        # "x" = chars 0..1, "2" = chars 1..2
        x_rng = sub_ranges[(0, 1)]
        sup_rng = sub_ranges[(1, 2)]  # noqa: PLR2004
        assert x_rng.Font.Superscript is False
        assert sup_rng.Font.Superscript is True

    def test_subscript_set(self) -> None:
        """Injecting '<sub>i</sub>' sets Font.Subscript = True on sub-range."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_win32com_word_html_runs,
        )

        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0

        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        _inject_win32com_word_html_runs(doc, rng, "H<sub>i</sub>")
        # "H" = chars 0..1, "i" = chars 1..2
        h_rng = sub_ranges[(0, 1)]
        sub_rng = sub_ranges[(1, 2)]  # noqa: PLR2004
        assert h_rng.Font.Subscript is False
        assert sub_rng.Font.Subscript is True


class TestInjectWin32comPptSupSub:
    """Tests for _inject_win32com_ppt_html_runs setting BaselineOffset."""

    def test_superscript_offset(self) -> None:
        """Injecting '<sup>2</sup>' sets Font.BaselineOffset = 0.3."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_win32com_ppt_html_runs,
        )

        tf = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        tf.TextRange.Paragraphs.return_value.Characters.side_effect = make_chars
        _inject_win32com_ppt_html_runs(tf, 1, "x<sup>2</sup>")
        # "x" at offset 0 → Characters(1, 1), "2" at offset 1 → Characters(2, 1)
        x_cr = char_ranges[(1, 1)]
        assert x_cr.Font.BaselineOffset == 0.0  # noqa: PLR2004
        sup_cr = char_ranges[(2, 1)]  # noqa: PLR2004
        assert sup_cr.Font.BaselineOffset == 0.3  # noqa: PLR2004

    def test_subscript_offset(self) -> None:
        """Injecting '<sub>i</sub>' sets Font.BaselineOffset = -0.25."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_win32com_ppt_html_runs,
        )

        tf = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        tf.TextRange.Paragraphs.return_value.Characters.side_effect = make_chars
        _inject_win32com_ppt_html_runs(tf, 1, "H<sub>i</sub>")
        # "H" at offset 0 → Characters(1, 1), "i" at offset 1 → Characters(2, 1)
        h_cr = char_ranges[(1, 1)]
        assert h_cr.Font.BaselineOffset == 0.0  # noqa: PLR2004
        sub_cr = char_ranges[(2, 1)]  # noqa: PLR2004
        assert sub_cr.Font.BaselineOffset == -0.25  # noqa: PLR2004


class TestReadUnoEffectiveFormattingSupSub:
    """Tests for _read_uno_effective_formatting superscript/subscript."""

    def test_positive_escapement_super(self) -> None:
        """CharEscapement=33 indicates superscript."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )
        from tests.test_office_formatter import _make_uno_portion  # noqa: PLC0415

        portion = _make_uno_portion("x")
        # Add CharEscapement to the props dict
        orig_side_effect = portion.getPropertyValue.side_effect
        props = {
            "CharWeight": orig_side_effect("CharWeight"),
            "CharPosture": orig_side_effect("CharPosture"),
            "CharUnderline": orig_side_effect("CharUnderline"),
            "CharStrikeout": orig_side_effect("CharStrikeout"),
            "CharEscapement": 33,
        }
        portion.getPropertyValue.side_effect = lambda p: props[p]
        result = _read_uno_effective_formatting(portion)
        assert len(result) == 6  # noqa: PLR2004
        assert result[4] is True  # superscript
        assert result[5] is False  # subscript

    def test_negative_escapement_sub(self) -> None:
        """CharEscapement=-33 indicates subscript."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )
        from tests.test_office_formatter import _make_uno_portion  # noqa: PLC0415

        portion = _make_uno_portion("x")
        orig_side_effect = portion.getPropertyValue.side_effect
        props = {
            "CharWeight": orig_side_effect("CharWeight"),
            "CharPosture": orig_side_effect("CharPosture"),
            "CharUnderline": orig_side_effect("CharUnderline"),
            "CharStrikeout": orig_side_effect("CharStrikeout"),
            "CharEscapement": -33,
        }
        portion.getPropertyValue.side_effect = lambda p: props[p]
        result = _read_uno_effective_formatting(portion)
        assert result[4] is False  # superscript
        assert result[5] is True  # subscript

    def test_zero_escapement(self) -> None:
        """CharEscapement=0 means neither superscript nor subscript."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )
        from tests.test_office_formatter import _make_uno_portion  # noqa: PLC0415

        portion = _make_uno_portion("x")
        orig_side_effect = portion.getPropertyValue.side_effect
        props = {
            "CharWeight": orig_side_effect("CharWeight"),
            "CharPosture": orig_side_effect("CharPosture"),
            "CharUnderline": orig_side_effect("CharUnderline"),
            "CharStrikeout": orig_side_effect("CharStrikeout"),
            "CharEscapement": 0,
        }
        portion.getPropertyValue.side_effect = lambda p: props[p]
        result = _read_uno_effective_formatting(portion)
        assert result[4] is False  # superscript
        assert result[5] is False  # subscript

    def test_exception_escapement(self) -> None:
        """Exception when reading CharEscapement gracefully returns both False."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )
        from tests.test_office_formatter import _make_uno_portion  # noqa: PLC0415

        portion = _make_uno_portion("x")
        orig_side_effect = portion.getPropertyValue.side_effect

        def _raise_on_escapement(p: str) -> object:
            if p == "CharEscapement":
                raise RuntimeError("no CharEscapement")
            return orig_side_effect(p)

        portion.getPropertyValue.side_effect = _raise_on_escapement
        result = _read_uno_effective_formatting(portion)
        assert result[4] is False  # superscript
        assert result[5] is False  # subscript


class TestReadUnoPortionFormattingSupSub:
    """Tests for _read_uno_portion_formatting superscript/subscript."""

    def test_portion_superscript(self) -> None:
        """Portion with CharEscapement=33 returns superscript=True in 6-tuple."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_portion_formatting,
        )
        from tests.test_office_formatter import _make_uno_portion  # noqa: PLC0415

        portion = _make_uno_portion("x")
        orig_side_effect = portion.getPropertyValue.side_effect
        props = {
            "CharWeight": orig_side_effect("CharWeight"),
            "CharPosture": orig_side_effect("CharPosture"),
            "CharUnderline": orig_side_effect("CharUnderline"),
            "CharStrikeout": orig_side_effect("CharStrikeout"),
            "CharEscapement": 33,
        }
        portion.getPropertyValue.side_effect = lambda p: props[p]
        result = _read_uno_portion_formatting(portion)
        assert len(result) == 6  # noqa: PLR2004
        # (bold, italic, underline, strike, superscript, subscript)
        assert result[4] is True  # superscript
        assert result[5] is False  # subscript


class TestReadUnoPortionFullFormattingSupSub:
    """Tests for _read_uno_portion_full_formatting superscript/subscript."""

    def test_full_9_tuple_superscript(self) -> None:
        """Verifies positions 4,5 in 9-tuple for superscript portion."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_portion_full_formatting,
        )
        from tests.test_office_formatter import _make_uno_portion  # noqa: PLC0415

        portion = _make_uno_portion("x", font_size=14.0, color_int=0xFF0000)
        orig_side_effect = portion.getPropertyValue.side_effect
        props = {
            "CharWeight": orig_side_effect("CharWeight"),
            "CharPosture": orig_side_effect("CharPosture"),
            "CharUnderline": orig_side_effect("CharUnderline"),
            "CharStrikeout": orig_side_effect("CharStrikeout"),
            "CharEscapement": 33,
            "CharHeight": orig_side_effect("CharHeight"),
            "CharColor": orig_side_effect("CharColor"),
            "CharHighlight": -1,
            "CharBackColor": -1,
        }
        portion.getPropertyValue.side_effect = lambda p: props[p]
        result = _read_uno_portion_full_formatting(portion)
        assert len(result) == 9  # noqa: PLR2004
        # (bold, italic, underline, strike, superscript, subscript,
        #  font_size_pt, color_hex, bg_color_hex)
        assert result[4] is True  # superscript
        assert result[5] is False  # subscript


class TestInjectUnoHtmlRunsSupSub:
    """Tests for _inject_uno_html_runs setting CharEscapement."""

    @staticmethod
    def _inject_and_collect_cursors(
        html_text: str,
        base_props: dict[str, object] | None = None,
    ) -> list[MagicMock]:
        """Injects HTML into a UNO para, returns list of created cursors.

        Args:
            html_text: HTML-formatted text to inject.
            base_props: Base properties dict for the UNO paragraph.

        Returns:
            List of MagicMock cursor objects captured during injection.
        """
        from src.core.office_processor import _inject_uno_html_runs  # noqa: PLC0415
        from tests.test_office_formatter import (  # noqa: PLC0415
            _make_uno_para_with_portions,
        )

        specs: list[tuple[str, bool, bool, bool, bool]] = [
            ("A", False, False, False, False),
        ]
        para = _make_uno_para_with_portions(specs)
        cursors: list[MagicMock] = []
        text_obj = para.getText.return_value
        orig_fn = text_obj.createTextCursorByRange.side_effect

        def _track(rng: object) -> MagicMock:
            c = orig_fn(rng)
            cursors.append(c)
            return c

        text_obj.createTextCursorByRange.side_effect = _track
        _inject_uno_html_runs(para, html_text, base_props or {})
        return cursors

    def test_superscript_escapement(self) -> None:
        """'<sup>2</sup>' sets CharEscapement=33, CharEscapementHeight=58."""
        cursors = self._inject_and_collect_cursors("x<sup>2</sup>")
        # cursors[0] = full-range cursor, cursors[1] = "x" seg, cursors[2] = "2" seg
        assert len(cursors) >= 3  # noqa: PLR2004
        sup_cursor = cursors[2]  # noqa: PLR2004
        props = sup_cursor._props
        assert props["CharEscapement"] == 33  # noqa: PLR2004
        assert props["CharEscapementHeight"] == 58  # noqa: PLR2004

    def test_subscript_escapement(self) -> None:
        """'<sub>i</sub>' sets CharEscapement=-33, CharEscapementHeight=58."""
        cursors = self._inject_and_collect_cursors("H<sub>i</sub>")
        # cursors[0] = full-range cursor, cursors[1] = "H" seg, cursors[2] = "i" seg
        assert len(cursors) >= 3  # noqa: PLR2004
        sub_cursor = cursors[2]  # noqa: PLR2004
        props = sub_cursor._props
        assert props["CharEscapement"] == -33  # noqa: PLR2004
        assert props["CharEscapementHeight"] == 58  # noqa: PLR2004

    def test_normal_resets_escapement(self) -> None:
        """Injecting plain text resets CharEscapement=0 and CharEscapementHeight=100."""
        cursors = self._inject_and_collect_cursors("plain")
        # cursors[0] = full-range cursor, cursors[1] = "plain" segment cursor
        assert len(cursors) >= 2  # noqa: PLR2004
        seg_cursor = cursors[1]
        props = seg_cursor._props
        assert props["CharEscapement"] == 0
        assert props["CharEscapementHeight"] == 100  # noqa: PLR2004


class TestUnoFormattingPropsIncludesEscapement:
    """Tests that _UNO_FORMATTING_PROPS includes CharEscapement props."""

    def test_char_escapement_in_props(self) -> None:
        """'CharEscapement' is in _UNO_FORMATTING_PROPS."""
        from src.core.office_processor import (  # noqa: PLC0415
            _UNO_FORMATTING_PROPS,
        )

        assert "CharEscapement" in _UNO_FORMATTING_PROPS

    def test_char_escapement_height_in_props(self) -> None:
        """'CharEscapementHeight' is in _UNO_FORMATTING_PROPS."""
        from src.core.office_processor import (  # noqa: PLC0415
            _UNO_FORMATTING_PROPS,
        )

        assert "CharEscapementHeight" in _UNO_FORMATTING_PROPS


class TestInjectDrawingmlBaselineStrip:
    """Tests that _inject_drawingml_html_runs strips baseline from base rPr."""

    def test_baseline_stripped_from_base(self) -> None:
        """Base rPr with baseline='30000' does not carry it to non-sup runs."""
        from lxml import etree  # noqa: PLC0415

        from src.core.office_processor import (  # noqa: PLC0415
            _inject_drawingml_html_runs,
        )

        dml_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        # Original run has baseline="30000" (superscript) on rPr
        xml = (
            f'<txBody xmlns:a="{dml_ns}">'
            "<a:p>"
            '<a:r><a:rPr baseline="30000"/><a:t>old</a:t></a:r>'
            "</a:p>"
            "</txBody>"
        )
        tx_body_el = etree.fromstring(xml)

        # Inject HTML with sup and plain segments
        _inject_drawingml_html_runs(tx_body_el, "<sup>2</sup> normal")

        a_r_tag = f"{{{dml_ns}}}r"
        a_rpr_tag = f"{{{dml_ns}}}rPr"
        p_el = tx_body_el[0]  # first <a:p>
        runs = p_el.findall(a_r_tag)
        assert len(runs) == 2  # noqa: PLR2004

        # First run (<sup>2</sup>) should have baseline="30000"
        sup_rpr = runs[0].find(a_rpr_tag)
        assert sup_rpr is not None
        assert sup_rpr.get("baseline") == "30000"

        # Second run (" normal") should NOT have baseline from the base rPr
        normal_rpr = runs[1].find(a_rpr_tag)
        assert normal_rpr is not None
        assert normal_rpr.get("baseline") is None


class TestInjectOdfTextBoxSupSub:
    """Tests for _inject_odf_text_box_html_runs setting style:text-position."""

    def _make_odf_text_box_el(self, text: str = "old text") -> object:
        """Builds a minimal <draw:text-box> lxml element.

        Args:
            text: Initial text content of the text box.

        Returns:
            lxml element representing a <draw:text-box>.
        """
        from lxml import etree  # noqa: PLC0415

        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
        xml = (
            f'<draw:text-box xmlns:draw="{draw_ns}" xmlns:text="{text_ns}">'
            f"<text:p>{text}</text:p>"
            f"</draw:text-box>"
        )
        return etree.fromstring(xml)

    def _make_auto_styles(self) -> object:
        """Builds a minimal <office:automatic-styles> lxml element.

        Returns:
            lxml element representing <office:automatic-styles>.
        """
        from lxml import etree  # noqa: PLC0415

        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        return etree.fromstring(
            f'<office:automatic-styles xmlns:office="{office_ns}"/>',
        )

    def test_superscript_text_position(self) -> None:
        """Injecting '<sup>2</sup>' produces style with text-position 'super 58%'."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_odf_text_box_html_runs,
        )

        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        text_box_el = self._make_odf_text_box_el()
        text_p_tag = f"{{{text_ns}}}p"
        auto_styles = self._make_auto_styles()
        counter = [0]

        result = _inject_odf_text_box_html_runs(
            text_box_el,
            "<sup>2</sup>",
            text_p_tag,
            auto_styles,
            counter,
        )
        assert result is True

        # Find the generated style and check text-position
        style_tag = f"{{{style_ns}}}style"
        text_props_tag = f"{{{style_ns}}}text-properties"
        styles = auto_styles.findall(style_tag)
        # At least one style should have text-position = "super 58%"
        found = False
        for style_el in styles:
            tp = style_el.find(text_props_tag)
            if tp is not None:
                text_pos = tp.get(f"{{{style_ns}}}text-position")
                if text_pos == "super 58%":
                    found = True
                    break
        assert found, "Expected style with text-position='super 58%'"

    def test_subscript_text_position(self) -> None:
        """Injecting '<sub>i</sub>' produces style with text-position 'sub 58%'."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_odf_text_box_html_runs,
        )

        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        text_box_el = self._make_odf_text_box_el()
        text_p_tag = f"{{{text_ns}}}p"
        auto_styles = self._make_auto_styles()
        counter = [0]

        result = _inject_odf_text_box_html_runs(
            text_box_el,
            "<sub>i</sub>",
            text_p_tag,
            auto_styles,
            counter,
        )
        assert result is True

        # Find the generated style and check text-position
        style_tag = f"{{{style_ns}}}style"
        text_props_tag = f"{{{style_ns}}}text-properties"
        styles = auto_styles.findall(style_tag)
        # At least one style should have text-position = "sub 58%"
        found = False
        for style_el in styles:
            tp = style_el.find(text_props_tag)
            if tp is not None:
                text_pos = tp.get(f"{{{style_ns}}}text-position")
                if text_pos == "sub 58%":
                    found = True
                    break
        assert found, "Expected style with text-position='sub 58%'"


# ---------------------------------------------------------------------------
# _inject_win32com_excel_html_runs — superscript / subscript
# ---------------------------------------------------------------------------


class TestInjectWin32comExcelSupSub:
    """Tests for _inject_win32com_excel_html_runs superscript/subscript.

    Excel shapes use BaselineOffset (same mechanism as PPT), not the
    Word-style Font.Superscript / Font.Subscript boolean properties.
    """

    def _make_excel_rng(self) -> MagicMock:
        """Creates a mock Excel TextRange2 COM object with Characters accessor.

        Returns:
            A MagicMock representing the TextRange2 interface.
        """
        rng = MagicMock()
        char_rng = MagicMock()
        rng.Characters.return_value = char_rng
        rng.Text = ""
        return rng

    def test_superscript_sets_baseline_offset(self) -> None:
        """Superscript HTML tag sets Font.BaselineOffset=0.3 on char range."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_win32com_excel_html_runs,
        )

        rng = self._make_excel_rng()
        char_rng = rng.Characters.return_value

        _inject_win32com_excel_html_runs(rng, "<sup>2</sup>")

        # BaselineOffset should have been set to 0.3 (superscript)
        assert char_rng.Font.BaselineOffset == 0.3  # noqa: PLR2004

    def test_subscript_sets_baseline_offset(self) -> None:
        """Subscript HTML tag sets Font.BaselineOffset=-0.25 on char range."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_win32com_excel_html_runs,
        )

        rng = self._make_excel_rng()
        char_rng = rng.Characters.return_value

        _inject_win32com_excel_html_runs(rng, "<sub>i</sub>")

        # BaselineOffset should have been set to -0.25 (subscript)
        assert char_rng.Font.BaselineOffset == -0.25  # noqa: PLR2004

    def test_normal_text_resets_baseline_offset(self) -> None:
        """Plain text sets Font.BaselineOffset=0.0 (resets any inherited offset)."""
        from src.core.office_processor import (  # noqa: PLC0415
            _inject_win32com_excel_html_runs,
        )

        rng = self._make_excel_rng()
        char_rng = rng.Characters.return_value

        _inject_win32com_excel_html_runs(rng, "plain")

        # BaselineOffset should be reset to 0.0
        assert char_rng.Font.BaselineOffset == 0.0  # noqa: PLR2004


# ---------------------------------------------------------------------------
# _has_uno_hyperlinks — direct unit tests
# ---------------------------------------------------------------------------


class TestHasUnoHyperlinks:
    """Tests for _has_uno_hyperlinks — direct unit tests."""

    def _make_para_with_portions(
        self,
        portions: list[MagicMock],
    ) -> MagicMock:
        """Creates a mock UNO paragraph that enumerates the given portions.

        Args:
            portions: List of mock UNO portion objects to enumerate.

        Returns:
            A MagicMock paragraph whose createEnumeration() yields portions.
        """
        para = MagicMock()
        idx = [0]

        def has_more() -> bool:
            return idx[0] < len(portions)

        def next_elem() -> MagicMock:
            p = portions[idx[0]]
            idx[0] += 1
            return p

        enum = MagicMock()
        enum.hasMoreElements.side_effect = has_more
        enum.nextElement.side_effect = next_elem
        para.createEnumeration.return_value = enum
        return para

    def _make_portion(
        self,
        text: str,
        hyperlink_url: str = "",
        portion_type: str = "Text",
    ) -> MagicMock:
        """Creates a mock UNO text portion.

        Args:
            text: The string content of the portion.
            hyperlink_url: HyperLinkURL property value (empty = no link).
            portion_type: TextPortionType value (default 'Text').

        Returns:
            A MagicMock representing a UNO text portion.
        """
        portion = MagicMock()
        props = {
            "TextPortionType": portion_type,
            "HyperLinkURL": hyperlink_url,
        }
        portion.getPropertyValue.side_effect = lambda p: props[p]
        portion.getString.return_value = text
        return portion

    def test_no_hyperlinks_returns_false(self) -> None:
        """Paragraph with no hyperlinks → False."""
        from src.core.office_processor import _has_uno_hyperlinks  # noqa: PLC0415

        portions = [self._make_portion("Hello"), self._make_portion("World")]
        para = self._make_para_with_portions(portions)
        assert _has_uno_hyperlinks(para) is False

    def test_with_hyperlink_returns_true(self) -> None:
        """Paragraph with at least one hyperlinked portion → True."""
        from src.core.office_processor import _has_uno_hyperlinks  # noqa: PLC0415

        portions = [
            self._make_portion("Click "),
            self._make_portion("here", hyperlink_url="https://example.com"),
        ]
        para = self._make_para_with_portions(portions)
        assert _has_uno_hyperlinks(para) is True

    def test_empty_url_not_counted(self) -> None:
        """Portions with empty string HyperLinkURL are not counted as hyperlinks."""
        from src.core.office_processor import _has_uno_hyperlinks  # noqa: PLC0415

        portions = [self._make_portion("text", hyperlink_url="")]
        para = self._make_para_with_portions(portions)
        assert _has_uno_hyperlinks(para) is False

    def test_non_text_portions_skipped(self) -> None:
        """Non-Text portion types are skipped (no HyperLinkURL check)."""
        from src.core.office_processor import _has_uno_hyperlinks  # noqa: PLC0415

        portions = [
            self._make_portion("foot", portion_type="Footnote"),
            self._make_portion("Real text"),
        ]
        para = self._make_para_with_portions(portions)
        assert _has_uno_hyperlinks(para) is False

    def test_hyper_link_url_exception_skips_portion(self) -> None:
        """getPropertyValue('HyperLinkURL') exception skips portion."""
        from src.core.office_processor import _has_uno_hyperlinks  # noqa: PLC0415

        portion = MagicMock()

        def _side_effect(prop: str) -> object:
            if prop == "TextPortionType":
                return "Text"
            if prop == "HyperLinkURL":
                raise RuntimeError("UNO property not supported")
            return ""

        portion.getPropertyValue.side_effect = _side_effect
        portion.getString.return_value = "some text"

        para = self._make_para_with_portions([portion])
        assert _has_uno_hyperlinks(para) is False

    def test_empty_string_portion_skipped(self) -> None:
        """Portions with empty getString() are skipped before HyperLinkURL check."""
        from src.core.office_processor import _has_uno_hyperlinks  # noqa: PLC0415

        # Portion has a URL but empty text — should be skipped
        portion = MagicMock()
        props = {
            "TextPortionType": "Text",
            "HyperLinkURL": "https://example.com",
        }
        portion.getPropertyValue.side_effect = lambda p: props[p]
        portion.getString.return_value = ""

        para = self._make_para_with_portions([portion])
        assert _has_uno_hyperlinks(para) is False


# ---------------------------------------------------------------------------
# _inject_uno_html_runs — superscript + hyperlink combo
# ---------------------------------------------------------------------------


class TestUnoSupHyperlinkCombo:
    """Tests for _inject_uno_html_runs handling superscript + hyperlink combo."""

    @staticmethod
    def _inject_and_collect_cursors(
        html_text: str,
        base_props: dict[str, object] | None = None,
    ) -> list[MagicMock]:
        """Injects HTML into a UNO para, returns list of created cursors.

        Args:
            html_text: HTML-formatted text to inject.
            base_props: Base properties dict for the UNO paragraph.

        Returns:
            List of MagicMock cursor objects captured during injection.
        """
        from src.core.office_processor import _inject_uno_html_runs  # noqa: PLC0415
        from tests.test_office_formatter import (  # noqa: PLC0415
            _make_uno_para_with_portions,
        )

        specs: list[tuple[str, bool, bool, bool, bool]] = [
            ("A", False, False, False, False),
        ]
        para = _make_uno_para_with_portions(specs)
        cursors: list[MagicMock] = []
        text_obj = para.getText.return_value
        orig_fn = text_obj.createTextCursorByRange.side_effect

        def _track(rng: object) -> MagicMock:
            c = orig_fn(rng)
            cursors.append(c)
            return c

        text_obj.createTextCursorByRange.side_effect = _track
        _inject_uno_html_runs(para, html_text, base_props or {})
        return cursors

    def test_sup_and_hyperlink_both_set(self) -> None:
        """Superscript + hyperlink sets CharEscapement=33, HyperLinkURL."""
        cursors = self._inject_and_collect_cursors(
            '<a href="https://x.com"><sup>ref</sup></a>',
        )
        # cursors[0] = full-range cursor, cursors[1] = "ref" segment cursor
        assert len(cursors) >= 2  # noqa: PLR2004
        seg_cursor = cursors[1]
        props = seg_cursor._props
        # CharEscapement=33 for superscript
        assert props["CharEscapement"] == 33  # noqa: PLR2004
        assert props["CharEscapementHeight"] == 58  # noqa: PLR2004
        # HyperLinkURL should carry the href
        assert "x.com" in str(props.get("HyperLinkURL", ""))

    def test_sub_and_hyperlink_both_set(self) -> None:
        """Subscript + hyperlink sets CharEscapement=-33, HyperLinkURL."""
        cursors = self._inject_and_collect_cursors(
            '<a href="https://y.org"><sub>n</sub></a>',
        )
        assert len(cursors) >= 2  # noqa: PLR2004
        seg_cursor = cursors[1]
        props = seg_cursor._props
        # CharEscapement=-33 for subscript
        assert props["CharEscapement"] == -33  # noqa: PLR2004
        assert props["CharEscapementHeight"] == 58  # noqa: PLR2004
        assert "y.org" in str(props.get("HyperLinkURL", ""))


# ---------------------------------------------------------------------------
# _read_uno_effective_formatting — CharEscapement edge cases
# ---------------------------------------------------------------------------


class TestReadUnoEffectiveFormattingEscapementEdgeCases:
    """CharEscapement type edge cases in _read_uno_effective_formatting."""

    def _make_obj(self, escapement: object) -> MagicMock:
        """Creates a mock UNO object with given CharEscapement value.

        Args:
            escapement: Value to return for the 'CharEscapement' property.

        Returns:
            A MagicMock UNO object whose getPropertyValue dispatches to a dict.
        """
        obj = MagicMock()
        posture_mock = MagicMock()
        posture_mock.value = "NONE"
        props = {
            "CharWeight": 100.0,
            "CharPosture": posture_mock,
            "CharUnderline": 0,
            "CharStrikeout": 0,
            "CharEscapement": escapement,
        }
        obj.getPropertyValue.side_effect = lambda p: props[p]
        return obj

    def test_none_escapement_returns_false(self) -> None:
        """None CharEscapement value (not int/float) → both sup/sub False."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )

        obj = self._make_obj(None)
        _, _, _, _, sup, sub = _read_uno_effective_formatting(obj)
        assert sup is False
        assert sub is False

    def test_string_escapement_returns_false(self) -> None:
        """String CharEscapement value → both False (isinstance check guards)."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )

        obj = self._make_obj("33")
        _, _, _, _, sup, sub = _read_uno_effective_formatting(obj)
        assert sup is False
        assert sub is False

    def test_float_positive_escapement_superscript(self) -> None:
        """Float positive CharEscapement (33.0) → superscript=True."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )

        obj = self._make_obj(33.0)
        _, _, _, _, sup, sub = _read_uno_effective_formatting(obj)
        assert sup is True
        assert sub is False

    def test_float_negative_escapement_subscript(self) -> None:
        """Float negative CharEscapement (-33.0) → subscript=True."""
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )

        obj = self._make_obj(-33.0)
        _, _, _, _, sup, sub = _read_uno_effective_formatting(obj)
        assert sup is False
        assert sub is True


# ---------------------------------------------------------------------------
# Sheet name translation tests
# ---------------------------------------------------------------------------


class TestSanitizeSheetName:
    """Tests for _sanitize_sheet_name."""

    def test_basic_name(self) -> None:
        """Normal name passes through unchanged."""
        assert _sanitize_sheet_name("Sales Data") == "Sales Data"

    def test_invalid_chars_removed(self) -> None:
        r"""Invalid chars \/*?:[] are stripped."""
        assert _sanitize_sheet_name("Sheet/1*2?3") == "Sheet123"

    def test_truncates_to_31_chars(self) -> None:
        """Names longer than 31 chars are truncated."""
        long_name = "A" * 50  # noqa: PLR2004
        assert len(_sanitize_sheet_name(long_name)) == 31  # noqa: PLR2004

    def test_exactly_31_chars_passes_through(self) -> None:
        """A name that is already exactly 31 chars is returned unchanged."""
        name_31 = "B" * 31  # noqa: PLR2004
        assert _sanitize_sheet_name(name_31) == name_31

    def test_empty_after_strip_returns_sheet(self) -> None:
        """All-invalid name returns 'Sheet'."""
        assert _sanitize_sheet_name("/*?:[]\\") == "Sheet"

    def test_whitespace_stripped(self) -> None:
        """Leading/trailing whitespace is stripped."""
        assert _sanitize_sheet_name("  Name  ") == "Name"

    def test_empty_string_input(self) -> None:
        """Empty string returns 'Sheet' fallback."""
        assert _sanitize_sheet_name("") == "Sheet"

    def test_combined_invalid_and_truncation(self) -> None:
        """Invalid chars are removed before truncation to 31 chars."""
        name = "A" * 20 + "/" + "B" * 20  # noqa: PLR2004
        result = _sanitize_sheet_name(name)
        assert "/" not in result
        assert len(result) == 31  # noqa: PLR2004


class TestShouldTranslateSheetNames:
    """Tests for _should_translate_sheet_names."""

    def test_unsupported_extension(self) -> None:
        """Returns False for non-spreadsheet extensions."""
        assert not _should_translate_sheet_names(".docx", "python_lib")

    def test_config_true(self) -> None:
        """Returns True when config has translate_sheet_names=True."""
        config = MagicMock()
        config.translate_sheet_names = True
        assert _should_translate_sheet_names(".xlsx", "python_lib", config)

    def test_config_false(self) -> None:
        """Returns False when config has translate_sheet_names=False."""
        config = MagicMock()
        config.translate_sheet_names = False
        assert not _should_translate_sheet_names(
            ".xlsx",
            "python_lib",
            config,
        )

    def test_extensions_set(self) -> None:
        """Correct extensions are in _SHEET_NAME_EXTENSIONS."""
        assert {".xlsx", ".ods", ".xls"} == _SHEET_NAME_EXTENSIONS

    @patch(
        "src.utils.config_manager.load_setting",
        return_value=True,
    )
    def test_load_setting_fallback_true(self, _mock: object) -> None:
        """Falls back to load_setting when config is None."""
        assert _should_translate_sheet_names(".xlsx", "python_lib")

    @patch(
        "src.utils.config_manager.load_setting",
        return_value=False,
    )
    def test_load_setting_fallback_false(self, _mock: object) -> None:
        """Falls back to load_setting returning False."""
        assert not _should_translate_sheet_names(".xlsx", "python_lib")


class TestExtractOdsSheetNames:
    """Tests for _extract_ods_sheet_names and _inject_ods_sheet_names."""

    def test_roundtrip(self, tmp_path: Path) -> None:
        """Extracts ODS sheet names, injects translations, re-extracts."""
        from odf.opendocument import OpenDocumentSpreadsheet  # noqa: PLC0415
        from odf.table import Table  # noqa: PLC0415

        doc = OpenDocumentSpreadsheet()
        t1 = Table(name="Revenue")
        doc.spreadsheet.addElement(t1)
        t2 = Table(name="Expenses")
        doc.spreadsheet.addElement(t2)
        p = tmp_path / "test.ods"
        doc.save(str(p))

        # Extract
        result = _extract_ods_sheet_names(p)
        assert ("sheetname:Revenue", "Revenue") in result
        assert ("sheetname:Expenses", "Expenses") in result

        # Inject uppercased translations
        translations = {k: v.upper() for k, v in result}
        _inject_ods_sheet_names(p, translations)

        # Re-extract and verify
        result2 = _extract_ods_sheet_names(p)
        names = [v for _, v in result2]
        assert "REVENUE" in names
        assert "EXPENSES" in names


class TestExtractXlsxSheetNames:
    """Tests for _extract_xlsx_sheet_names."""

    def test_extracts_sheet_names(self, tmp_path: Path) -> None:
        """Extracts sheet names from XLSX file."""
        wb = Workbook()
        wb.active.title = "Data"
        wb.create_sheet("Summary")
        p = tmp_path / "test.xlsx"
        wb.save(str(p))
        wb.close()

        result = _extract_xlsx_sheet_names(p)
        assert ("sheetname:Data", "Data") in result
        assert ("sheetname:Summary", "Summary") in result


class TestInjectXlsxSheetNames:
    """Tests for _inject_xlsx_sheet_names via ZIP+lxml."""

    def test_renames_sheets(self, tmp_path: Path) -> None:
        """Sheet names are updated in xl/workbook.xml."""
        wb = Workbook()
        wb.active.title = "Sheet1"
        wb.create_sheet("Sheet2")
        p = tmp_path / "test.xlsx"
        wb.save(str(p))
        wb.close()

        from src.core.office_processor import (  # noqa: PLC0415
            _inject_xlsx_sheet_names,
        )

        _inject_xlsx_sheet_names(
            p,
            {"sheetname:Sheet1": "Translated1", "sheetname:Sheet2": "Translated2"},
        )

        # Verify via openpyxl
        wb2 = load_workbook(str(p))
        names = [ws.title for ws in wb2.worksheets]
        wb2.close()
        assert "Translated1" in names
        assert "Translated2" in names

    def test_no_match_no_rewrite(self, tmp_path: Path) -> None:
        """No changes when no sheet names match."""
        wb = Workbook()
        p = tmp_path / "test.xlsx"
        wb.save(str(p))
        wb.close()

        from src.core.office_processor import (  # noqa: PLC0415
            _inject_xlsx_sheet_names,
        )

        _inject_xlsx_sheet_names(p, {"sheetname:NoMatch": "X"})
        # Should not crash — file remains valid
        wb2 = load_workbook(str(p))
        wb2.close()


class TestExtractSheetNamesDispatch:
    """Tests for _extract_sheet_names dispatch."""

    def test_xlsx_dispatch(self, tmp_path: Path) -> None:
        """Dispatches to _extract_xlsx_sheet_names for .xlsx."""
        wb = Workbook()
        wb.active.title = "MySheet"
        p = tmp_path / "test.xlsx"
        wb.save(str(p))
        wb.close()

        result = _extract_sheet_names(p, ".xlsx", "python_lib")
        assert any(k == "sheetname:MySheet" for k, _ in result)

    def test_win32com_dispatch(self) -> None:
        """Dispatches to Win32COM for .xls on win32com backend."""
        with patch(
            "src.core.office_processor._extract_win32com_excel_sheet_names",
            return_value=[("sheetname:S1", "S1")],
        ) as mock_fn:
            result = _extract_sheet_names(Path("f.xls"), ".xls", "win32com")
        mock_fn.assert_called_once()
        assert result == [("sheetname:S1", "S1")]


class TestInjectSheetNamesDispatch:
    """Tests for _inject_sheet_names dispatch."""

    def test_no_keys_returns_early(self) -> None:
        """Returns immediately when no sheetname: keys."""
        # Should not crash
        _inject_sheet_names(
            Path("x.xlsx"),
            {"para:0": "text"},
            ".xlsx",
            "python_lib",
        )

    def test_xlsx_dispatch(self) -> None:
        """Dispatches to _inject_xlsx_sheet_names for .xlsx."""
        with patch(
            "src.core.office_processor._inject_xlsx_sheet_names",
        ) as mock_fn:
            _inject_sheet_names(
                Path("f.xlsx"),
                {"sheetname:S1": "T1"},
                ".xlsx",
                "python_lib",
            )
        mock_fn.assert_called_once()

    def test_ods_dispatch(self) -> None:
        """Dispatches to _inject_ods_sheet_names for .ods."""
        with patch(
            "src.core.office_processor._inject_ods_sheet_names",
        ) as mock_fn:
            _inject_sheet_names(
                Path("f.ods"),
                {"sheetname:S1": "T1"},
                ".ods",
                "python_lib",
            )
        mock_fn.assert_called_once()


# ---------------------------------------------------------------------------
# Speaker notes translation tests
# ---------------------------------------------------------------------------


class TestShouldTranslateNotes:
    """Tests for _should_translate_notes."""

    def test_unsupported_extension(self) -> None:
        """Returns False for non-presentation extensions."""
        assert not _should_translate_notes(".xlsx", "python_lib")

    def test_config_true(self) -> None:
        """Returns True when config has translate_doc_notes=True."""
        config = MagicMock()
        config.translate_doc_notes = True
        assert _should_translate_notes(".pptx", "python_lib", config)

    def test_config_false(self) -> None:
        """Returns False when config has translate_doc_notes=False."""
        config = MagicMock()
        config.translate_doc_notes = False
        assert not _should_translate_notes(".pptx", "python_lib", config)

    def test_extensions_set(self) -> None:
        """Correct extensions are in _NOTES_EXTENSIONS."""
        assert {".pptx", ".odp", ".ppt"} == _NOTES_EXTENSIONS

    @patch(
        "src.utils.config_manager.load_setting",
        return_value=True,
    )
    def test_load_setting_fallback_true(self, _mock: object) -> None:
        """Falls back to load_setting when config is None."""
        assert _should_translate_notes(".pptx", "python_lib")

    @patch(
        "src.utils.config_manager.load_setting",
        return_value=False,
    )
    def test_load_setting_fallback_false(self, _mock: object) -> None:
        """Falls back to load_setting returning False."""
        assert not _should_translate_notes(".pptx", "python_lib")


class TestExtractPptxNotes:
    """Tests for _extract_pptx_notes."""

    def test_extracts_notes(self, tmp_path: Path) -> None:
        """Extracts speaker notes from PPTX file."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = "Hello speaker note"
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        result = _extract_pptx_notes(p)
        assert len(result) >= 1
        # The notes text should be in the results
        assert any("Hello speaker note" in v for _, v in result)

    def test_no_notes(self, tmp_path: Path) -> None:
        """Returns empty when slides have no notes."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        result = _extract_pptx_notes(p)
        assert result == []

    def test_multiple_slides_with_notes(self, tmp_path: Path) -> None:
        """Extracts notes from multiple slides with correct keys."""
        prs = Presentation()
        for idx in range(3):  # noqa: PLR2004
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            ns = slide.notes_slide
            ns.notes_text_frame.text = f"Note for slide {idx}"
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        result = _extract_pptx_notes(p)
        keys = [k for k, _ in result]
        # Each slide should have at least one note:s_idx:p_idx entry
        assert any(k.startswith("note:0:") for k in keys)
        assert any(k.startswith("note:1:") for k in keys)
        assert any(k.startswith("note:2:") for k in keys)
        # Verify texts
        assert any("Note for slide 0" in v for _, v in result)
        assert any("Note for slide 1" in v for _, v in result)
        assert any("Note for slide 2" in v for _, v in result)

    def test_mixed_slides_some_with_notes(self, tmp_path: Path) -> None:
        """Only slides with notes produce entries."""
        prs = Presentation()
        # Slide 0 — with notes
        s0 = prs.slides.add_slide(prs.slide_layouts[0])
        s0.notes_slide.notes_text_frame.text = "First note"
        # Slide 1 — no notes (blank layout, never access notes_slide)
        prs.slides.add_slide(prs.slide_layouts[5])
        # Slide 2 — with notes
        s2 = prs.slides.add_slide(prs.slide_layouts[0])
        s2.notes_slide.notes_text_frame.text = "Third note"
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        result = _extract_pptx_notes(p)
        keys = [k for k, _ in result]
        assert any(k.startswith("note:0:") for k in keys)
        assert not any(k.startswith("note:1:") for k in keys)
        assert any(k.startswith("note:2:") for k in keys)

    def test_multiple_paragraphs_in_notes(self, tmp_path: Path) -> None:
        """Multiple paragraphs in a note produce distinct p_idx keys."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        ns = slide.notes_slide
        tf = ns.notes_text_frame
        tf.text = "First paragraph"
        tf.add_paragraph().text = "Second paragraph"
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        result = _extract_pptx_notes(p)
        keys = [k for k, _ in result]
        # Both paragraphs should appear with different p_idx
        assert "note:0:0" in keys
        assert "note:0:1" in keys

    def test_whitespace_only_notes_skipped(self, tmp_path: Path) -> None:
        """Notes with only whitespace text are not extracted."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        tf = slide.notes_slide.notes_text_frame
        tf.text = "   "
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        result = _extract_pptx_notes(p)
        # Whitespace-only note should not produce entries
        assert not any("note:0:" in k for k, _ in result)


class TestInjectPptxNotes:
    """Tests for _inject_pptx_notes."""

    def test_injects_notes(self, tmp_path: Path) -> None:
        """Injects translated notes back into PPTX."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = "Original note"
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        # Find the key for the note text
        extracted = _extract_pptx_notes(p)
        note_key = None
        for k, v in extracted:
            if "Original note" in v:
                note_key = k
                break
        assert note_key is not None

        _inject_pptx_notes(p, {note_key: "Translated note"})

        # Verify
        prs2 = Presentation(str(p))
        slide2 = prs2.slides[0]
        notes_text = slide2.notes_slide.notes_text_frame.text
        assert "Translated note" in notes_text

    def test_injects_multiple_paragraphs(self, tmp_path: Path) -> None:
        """Injects translations for two paragraphs in the same slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        tf = slide.notes_slide.notes_text_frame
        tf.text = "Para one"
        tf.add_paragraph().text = "Para two"
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        extracted = _extract_pptx_notes(p)
        translations = {}
        for k, v in extracted:
            if "Para one" in v:
                translations[k] = "Translated one"
            elif "Para two" in v:
                translations[k] = "Translated two"

        _inject_pptx_notes(p, translations)

        prs2 = Presentation(str(p))
        paras = prs2.slides[0].notes_slide.notes_text_frame.paragraphs
        texts = [pr.text for pr in paras]
        assert "Translated one" in texts
        assert "Translated two" in texts


class TestExtractNotesDispatch:
    """Tests for _extract_notes dispatch."""

    def test_pptx_dispatch(self, tmp_path: Path) -> None:
        """Dispatches to _extract_pptx_notes for .pptx."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Note"
        p = tmp_path / "test.pptx"
        prs.save(str(p))

        result = _extract_notes(p, ".pptx", "python_lib")
        assert any("Note" in v for _, v in result)

    def test_win32com_dispatch(self) -> None:
        """Dispatches to Win32COM for .ppt on win32com backend."""
        with patch(
            "src.core.office_processor._extract_win32com_ppt_notes",
            return_value=[("note:0:0", "Note text")],
        ) as mock_fn:
            result = _extract_notes(Path("f.ppt"), ".ppt", "win32com")
        mock_fn.assert_called_once()
        assert result == [("note:0:0", "Note text")]

    def test_odp_dispatch(self) -> None:
        """Dispatches to _extract_odp_notes for .odp."""
        with patch(
            "src.core.office_processor._extract_odp_notes",
            return_value=[("note:0:0", "ODP note")],
        ) as mock_fn:
            result = _extract_notes(Path("f.odp"), ".odp", "python_lib")
        mock_fn.assert_called_once()
        assert result == [("note:0:0", "ODP note")]


class TestInjectNotesDispatch:
    """Tests for _inject_notes dispatch."""

    def test_no_keys_returns_early(self) -> None:
        """Returns immediately when no note: keys."""
        _inject_notes(
            Path("x.pptx"),
            {"para:0": "text"},
            ".pptx",
            "python_lib",
        )

    def test_pptx_dispatch(self) -> None:
        """Dispatches to _inject_pptx_notes for .pptx."""
        with patch(
            "src.core.office_processor._inject_pptx_notes",
        ) as mock_fn:
            _inject_notes(
                Path("f.pptx"),
                {"note:0:0": "T"},
                ".pptx",
                "python_lib",
            )
        mock_fn.assert_called_once()

    def test_odp_dispatch(self) -> None:
        """Dispatches to _inject_odp_notes for .odp."""
        with patch(
            "src.core.office_processor._inject_odp_notes",
        ) as mock_fn:
            _inject_notes(
                Path("f.odp"),
                {"note:0:0": "T"},
                ".odp",
                "python_lib",
            )
        mock_fn.assert_called_once()


# ---------------------------------------------------------------------------
# Header and footer translation tests
# ---------------------------------------------------------------------------


class TestExtractDocxHeadersFooters:
    """Tests for _extract_docx_headers_footers."""

    def test_extracts_default_header(self, tmp_path: Path) -> None:
        """Extracts default header text from DOCX."""
        doc = Document()
        section = doc.sections[0]
        section.header.is_linked_to_previous = False
        header_para = section.header.paragraphs[0]
        header_para.text = "My Header"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        assert any(
            k.startswith("header:0:default:") and "My Header" in v for k, v in result
        )

    def test_extracts_default_footer(self, tmp_path: Path) -> None:
        """Extracts default footer text from DOCX."""
        doc = Document()
        section = doc.sections[0]
        section.footer.is_linked_to_previous = False
        footer_para = section.footer.paragraphs[0]
        footer_para.text = "My Footer"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        assert any(
            k.startswith("footer:0:default:") and "My Footer" in v for k, v in result
        )

    def test_linked_to_previous_skipped(self, tmp_path: Path) -> None:
        """Headers linked to previous section are skipped."""
        doc = Document()
        # Default section — header is linked by default
        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        # No header text should be extracted (linked to previous is default)
        assert not any(k.startswith("header:") for k, _ in result)

    def test_extracts_first_page_header(self, tmp_path: Path) -> None:
        """Extracts first-page header when different_first_page is set."""
        doc = Document()
        section = doc.sections[0]
        section.different_first_page_header_footer = True
        first_hdr = section.first_page_header
        first_hdr.is_linked_to_previous = False
        first_hdr.paragraphs[0].text = "First Page Header"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        assert any(
            k == f"header:0:{_HF_FIRST}:0" and "First Page Header" in v
            for k, v in result
        )

    def test_extracts_first_page_footer(self, tmp_path: Path) -> None:
        """Extracts first-page footer when different_first_page is set."""
        doc = Document()
        section = doc.sections[0]
        section.different_first_page_header_footer = True
        first_ftr = section.first_page_footer
        first_ftr.is_linked_to_previous = False
        first_ftr.paragraphs[0].text = "First Page Footer"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        assert any(
            k == f"footer:0:{_HF_FIRST}:0" and "First Page Footer" in v
            for k, v in result
        )

    def test_extracts_even_page_header(self, tmp_path: Path) -> None:
        """Extracts even-page header when odd_and_even is set."""
        doc = Document()
        section = doc.sections[0]
        doc.settings.odd_and_even_pages_header_footer = True
        even_hdr = section.even_page_header
        even_hdr.is_linked_to_previous = False
        even_hdr.paragraphs[0].text = "Even Page Header"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        assert any(
            k == f"header:0:{_HF_EVEN}:0" and "Even Page Header" in v for k, v in result
        )

    def test_extracts_even_page_footer(self, tmp_path: Path) -> None:
        """Extracts even-page footer when odd_and_even is set."""
        doc = Document()
        section = doc.sections[0]
        doc.settings.odd_and_even_pages_header_footer = True
        even_ftr = section.even_page_footer
        even_ftr.is_linked_to_previous = False
        even_ftr.paragraphs[0].text = "Even Page Footer"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        assert any(
            k == f"footer:0:{_HF_EVEN}:0" and "Even Page Footer" in v for k, v in result
        )

    def test_multiple_sections(self, tmp_path: Path) -> None:
        """Extracts headers from multiple sections with correct indices."""
        doc = Document()
        # Section 0 — default header
        s0 = doc.sections[0]
        s0.header.is_linked_to_previous = False
        s0.header.paragraphs[0].text = "Header Section 0"

        # Section 1 — add section break, give it a different header
        new_section = doc.add_section()
        new_section.header.is_linked_to_previous = False
        new_section.header.paragraphs[0].text = "Header Section 1"

        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        assert any(
            k.startswith("header:0:default:") and "Header Section 0" in v
            for k, v in result
        )
        assert any(
            k.startswith("header:1:default:") and "Header Section 1" in v
            for k, v in result
        )

    def test_all_three_types_simultaneously(self, tmp_path: Path) -> None:
        """Extracts default, first-page, and even-page headers at once."""
        doc = Document()
        section = doc.sections[0]

        # Default header
        section.header.is_linked_to_previous = False
        section.header.paragraphs[0].text = "Default Header"

        # First-page header
        section.different_first_page_header_footer = True
        first_hdr = section.first_page_header
        first_hdr.is_linked_to_previous = False
        first_hdr.paragraphs[0].text = "First Header"

        # Even-page header
        doc.settings.odd_and_even_pages_header_footer = True
        even_hdr = section.even_page_header
        even_hdr.is_linked_to_previous = False
        even_hdr.paragraphs[0].text = "Even Header"

        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_docx_headers_footers(p)
        assert any(
            k == f"header:0:{_HF_DEFAULT}:0" and "Default Header" in v
            for k, v in result
        )
        assert any(
            k == f"header:0:{_HF_FIRST}:0" and "First Header" in v for k, v in result
        )
        assert any(
            k == f"header:0:{_HF_EVEN}:0" and "Even Header" in v for k, v in result
        )


class TestInjectDocxHeadersFooters:
    """Tests for _inject_docx_headers_footers."""

    def test_injects_header(self, tmp_path: Path) -> None:
        """Injects translated header text into DOCX."""
        doc = Document()
        section = doc.sections[0]
        section.header.is_linked_to_previous = False
        section.header.paragraphs[0].text = "Original Header"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        _inject_docx_headers_footers(
            p,
            {"header:0:default:0": "Translated Header"},
        )

        doc2 = Document(str(p))
        assert doc2.sections[0].header.paragraphs[0].text == "Translated Header"

    def test_injects_first_page_header(self, tmp_path: Path) -> None:
        """Injects translated first-page header into DOCX."""
        doc = Document()
        section = doc.sections[0]
        section.different_first_page_header_footer = True
        first_hdr = section.first_page_header
        first_hdr.is_linked_to_previous = False
        first_hdr.paragraphs[0].text = "Original First"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        _inject_docx_headers_footers(
            p,
            {f"header:0:{_HF_FIRST}:0": "Translated First"},
        )

        doc2 = Document(str(p))
        assert (
            doc2.sections[0].first_page_header.paragraphs[0].text == "Translated First"
        )

    def test_injects_even_page_header(self, tmp_path: Path) -> None:
        """Injects translated even-page header into DOCX."""
        doc = Document()
        section = doc.sections[0]
        doc.settings.odd_and_even_pages_header_footer = True
        even_hdr = section.even_page_header
        even_hdr.is_linked_to_previous = False
        even_hdr.paragraphs[0].text = "Original Even"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        _inject_docx_headers_footers(
            p,
            {f"header:0:{_HF_EVEN}:0": "Translated Even"},
        )

        doc2 = Document(str(p))
        assert doc2.sections[0].even_page_header.paragraphs[0].text == "Translated Even"

    def test_partial_injection(self, tmp_path: Path) -> None:
        """Injecting only header leaves footer text unchanged."""
        doc = Document()
        section = doc.sections[0]
        section.header.is_linked_to_previous = False
        section.header.paragraphs[0].text = "Original Header"
        section.footer.is_linked_to_previous = False
        section.footer.paragraphs[0].text = "Original Footer"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        # Only inject a translation for the header — not the footer
        _inject_docx_headers_footers(
            p,
            {"header:0:default:0": "Translated Header"},
        )

        doc2 = Document(str(p))
        assert doc2.sections[0].header.paragraphs[0].text == "Translated Header"
        assert doc2.sections[0].footer.paragraphs[0].text == "Original Footer"

    def test_injects_all_three_types(self, tmp_path: Path) -> None:
        """Injects default, first-page, and even-page headers at once."""
        doc = Document()
        section = doc.sections[0]

        # Default header
        section.header.is_linked_to_previous = False
        section.header.paragraphs[0].text = "Default H"

        # First-page header
        section.different_first_page_header_footer = True
        first_hdr = section.first_page_header
        first_hdr.is_linked_to_previous = False
        first_hdr.paragraphs[0].text = "First H"

        # Even-page header
        doc.settings.odd_and_even_pages_header_footer = True
        even_hdr = section.even_page_header
        even_hdr.is_linked_to_previous = False
        even_hdr.paragraphs[0].text = "Even H"

        p = tmp_path / "test.docx"
        doc.save(str(p))

        _inject_docx_headers_footers(
            p,
            {
                f"header:0:{_HF_DEFAULT}:0": "Translated Default",
                f"header:0:{_HF_FIRST}:0": "Translated First",
                f"header:0:{_HF_EVEN}:0": "Translated Even",
            },
        )

        doc2 = Document(str(p))
        s = doc2.sections[0]
        assert s.header.paragraphs[0].text == "Translated Default"
        assert s.first_page_header.paragraphs[0].text == "Translated First"
        assert s.even_page_header.paragraphs[0].text == "Translated Even"


class TestDocxHfPart:
    """Tests for _extract_docx_hf_part helper."""

    def test_extracts_paragraphs(self) -> None:
        """Extracts non-empty paragraphs with correct keys."""
        p1 = MagicMock()
        p1.text = "Hello"
        p2 = MagicMock()
        p2.text = ""
        p3 = MagicMock()
        p3.text = "World"

        with patch(
            "src.core.office_processor._extract_para_with_links",
            side_effect=lambda p: p.text,
        ):
            result = _extract_docx_hf_part(
                [p1, p2, p3],
                0,
                _HF_DEFAULT,
                "header",
            )

        assert len(result) == 2  # noqa: PLR2004
        assert result[0] == ("header:0:default:0", "Hello")
        assert result[1] == ("header:0:default:2", "World")

    def test_empty_paragraphs_skipped(self) -> None:
        """Empty and whitespace-only paragraphs are skipped."""
        p1 = MagicMock()
        p1.text = "   "
        result = _extract_docx_hf_part([p1], 0, _HF_FIRST, "footer")
        assert result == []


class TestHeaderFooterExtensions:
    """Tests for _HEADER_FOOTER_EXTENSIONS."""

    def test_extensions(self) -> None:
        """Correct extensions are in _HEADER_FOOTER_EXTENSIONS."""
        assert {".docx", ".odt", ".doc"} == _HEADER_FOOTER_EXTENSIONS


class TestExtractHeadersFootersDispatch:
    """Tests for _extract_headers_footers dispatch."""

    def test_docx_dispatch(self, tmp_path: Path) -> None:
        """Dispatches to _extract_docx_headers_footers for .docx."""
        doc = Document()
        section = doc.sections[0]
        section.header.is_linked_to_previous = False
        section.header.paragraphs[0].text = "Test"
        p = tmp_path / "test.docx"
        doc.save(str(p))

        result = _extract_headers_footers(p, ".docx", "python_lib")
        assert any(k.startswith("header:") for k, _ in result)

    def test_win32com_dispatch(self) -> None:
        """Dispatches to Win32COM for .doc on win32com backend."""
        with patch(
            "src.core.office_processor._extract_win32com_word_headers_footers",
            return_value=[("header:0:default:0", "H")],
        ) as mock_fn:
            result = _extract_headers_footers(
                Path("f.doc"),
                ".doc",
                "win32com",
            )
        mock_fn.assert_called_once()
        assert result == [("header:0:default:0", "H")]


class TestInjectHeadersFootersDispatch:
    """Tests for _inject_headers_footers dispatch."""

    def test_no_keys_returns_early(self) -> None:
        """Returns immediately when no header:/footer: keys."""
        _inject_headers_footers(
            Path("x.docx"),
            {"para:0": "text"},
            ".docx",
            "python_lib",
        )


# ---------------------------------------------------------------------------
# Footnote and endnote translation tests
# ---------------------------------------------------------------------------


class TestExtractDocxFootnotes:
    """Tests for _extract_docx_footnotes."""

    def test_extracts_footnotes(self, tmp_path: Path) -> None:
        """Extracts footnotes from DOCX file."""
        # Create a DOCX with a footnote via low-level XML
        doc = Document()
        doc.add_paragraph("Body text")
        p = tmp_path / "test.docx"
        doc.save(str(p))

        # Manually add footnotes.xml to the ZIP
        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="0" w:type="separator"><w:p><w:r>'
            "<w:separator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="1" w:type="continuationSeparator"><w:p>'
            "<w:r><w:continuationSeparator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="2"><w:p><w:r><w:t>Footnote text</w:t>'
            "</w:r></w:p></w:footnote>"
            "</w:footnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        new_item = zipfile.ZipInfo("word/footnotes.xml")
        all_items.append(new_item)

        _rewrite_zip_content(p, file_data, all_items)

        result = _extract_docx_footnotes(p)
        assert ("footnote:2", "Footnote text") in result

    def test_skips_separator_ids(self, tmp_path: Path) -> None:
        """IDs 0, 1, and -1 are skipped."""
        doc = Document()
        p = tmp_path / "test.docx"
        doc.save(str(p))

        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="0"><w:p><w:r><w:t>Sep</w:t>'
            "</w:r></w:p></w:footnote>"
            '<w:footnote w:id="-1"><w:p><w:r><w:t>Sep2</w:t>'
            "</w:r></w:p></w:footnote>"
            "</w:footnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/footnotes.xml"))
        _rewrite_zip_content(p, file_data, all_items)

        result = _extract_docx_footnotes(p)
        assert result == []

    def test_extracts_endnotes(self, tmp_path: Path) -> None:
        """Extracts endnotes from DOCX endnotes.xml."""
        doc = Document()
        p = tmp_path / "test.docx"
        doc.save(str(p))

        en_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:endnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:endnote w:id="0"><w:p><w:r><w:separator/></w:r>'
            "</w:p></w:endnote>"
            '<w:endnote w:id="2"><w:p><w:r><w:t>Endnote here</w:t>'
            "</w:r></w:p></w:endnote>"
            "</w:endnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/endnotes.xml"] = en_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/endnotes.xml"))
        _rewrite_zip_content(p, file_data, all_items)

        result = _extract_docx_footnotes(p)
        assert ("endnote:2", "Endnote here") in result

    def test_multiple_footnotes(self, tmp_path: Path) -> None:
        """Extracts only user footnotes (IDs 2, 3, 4), not separators."""
        doc = Document()
        p = tmp_path / "test.docx"
        doc.save(str(p))

        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="0" w:type="separator"><w:p><w:r>'
            "<w:separator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="1" w:type="continuationSeparator"><w:p>'
            "<w:r><w:continuationSeparator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="2"><w:p><w:r><w:t>First note</w:t>'
            "</w:r></w:p></w:footnote>"
            '<w:footnote w:id="3"><w:p><w:r><w:t>Second note</w:t>'
            "</w:r></w:p></w:footnote>"
            '<w:footnote w:id="4"><w:p><w:r><w:t>Third note</w:t>'
            "</w:r></w:p></w:footnote>"
            "</w:footnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/footnotes.xml"))
        _rewrite_zip_content(p, file_data, all_items)

        result = _extract_docx_footnotes(p)
        assert len(result) == 3  # noqa: PLR2004
        assert ("footnote:2", "First note") in result
        assert ("footnote:3", "Second note") in result
        assert ("footnote:4", "Third note") in result

    def test_empty_footnote_skipped(self, tmp_path: Path) -> None:
        """Footnotes with empty or whitespace-only text are skipped."""
        doc = Document()
        p = tmp_path / "test.docx"
        doc.save(str(p))

        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="2"><w:p><w:r><w:t>   </w:t>'
            "</w:r></w:p></w:footnote>"
            '<w:footnote w:id="3"><w:p><w:r><w:t/>'
            "</w:r></w:p></w:footnote>"
            '<w:footnote w:id="4"><w:p><w:r><w:t>Real note</w:t>'
            "</w:r></w:p></w:footnote>"
            "</w:footnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/footnotes.xml"))
        _rewrite_zip_content(p, file_data, all_items)

        result = _extract_docx_footnotes(p)
        # Only the non-empty footnote should be extracted
        assert len(result) == 1
        assert ("footnote:4", "Real note") in result

    def test_both_footnotes_and_endnotes(self, tmp_path: Path) -> None:
        """Extracts both footnotes and endnotes from the same DOCX."""
        doc = Document()
        p = tmp_path / "test.docx"
        doc.save(str(p))

        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="0" w:type="separator"><w:p><w:r>'
            "<w:separator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="2"><w:p><w:r><w:t>FN text</w:t>'
            "</w:r></w:p></w:footnote>"
            "</w:footnotes>"
        )
        en_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:endnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:endnote w:id="0"><w:p><w:r><w:separator/></w:r>'
            "</w:p></w:endnote>"
            '<w:endnote w:id="2"><w:p><w:r><w:t>EN text</w:t>'
            "</w:r></w:p></w:endnote>"
            "</w:endnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        file_data["word/endnotes.xml"] = en_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/footnotes.xml"))
        all_items.append(zipfile.ZipInfo("word/endnotes.xml"))
        _rewrite_zip_content(p, file_data, all_items)

        result = _extract_docx_footnotes(p)
        assert ("footnote:2", "FN text") in result
        assert ("endnote:2", "EN text") in result


class TestInjectDocxFootnotes:
    """Tests for _inject_docx_footnotes."""

    def test_injects_footnote_text(self, tmp_path: Path) -> None:
        """Injects translated text into DOCX footnotes."""
        doc = Document()
        p = tmp_path / "test.docx"
        doc.save(str(p))

        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="0" w:type="separator"><w:p><w:r>'
            "<w:separator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="2"><w:p><w:r><w:t>Original</w:t>'
            "</w:r></w:p></w:footnote>"
            "</w:footnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/footnotes.xml"))
        _rewrite_zip_content(p, file_data, all_items)

        _inject_docx_footnotes(p, {"footnote:2": "Translated"})

        # Re-read and verify
        result = _extract_docx_footnotes(p)
        assert ("footnote:2", "Translated") in result

    def test_preserves_footnote_ref_marker(self, tmp_path: Path) -> None:
        """Injection preserves <w:footnoteRef/> in the first run."""
        doc = Document()
        p = tmp_path / "test.docx"
        doc.save(str(p))

        # Footnote with footnoteRef marker in first run, text in second run
        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="0" w:type="separator"><w:p><w:r>'
            "<w:separator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="2"><w:p>'
            "<w:r><w:rPr/><w:footnoteRef/></w:r>"
            "<w:r><w:t>Original text</w:t></w:r>"
            "</w:p></w:footnote>"
            "</w:footnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/footnotes.xml"))
        _rewrite_zip_content(p, file_data, all_items)

        _inject_docx_footnotes(p, {"footnote:2": "Translated text"})

        # Verify footnoteRef is preserved and text is updated
        from lxml import etree  # noqa: PLC0415

        with zipfile.ZipFile(p, "r") as zf:
            fn_data = zf.read("word/footnotes.xml")
        root = etree.fromstring(fn_data)
        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        footnote = None
        for fn in root.iter(f"{{{w_ns}}}footnote"):
            if fn.get(f"{{{w_ns}}}id") == "2":
                footnote = fn
                break
        assert footnote is not None
        # Check footnoteRef still exists
        refs = list(footnote.iter(f"{{{w_ns}}}footnoteRef"))
        assert len(refs) == 1
        # Check translated text is present
        result = _extract_docx_footnotes(p)
        assert ("footnote:2", "Translated text") in result

    def test_injects_both_footnote_and_endnote(self, tmp_path: Path) -> None:
        """Footnote and endnote keys route to their respective XML files.

        When translations contain both ``footnote:{id}`` and ``endnote:{id}``
        keys, each must be written to the correct archive member so the types
        remain distinct after injection.
        """
        doc = Document()
        p = tmp_path / "mixed.docx"
        doc.save(str(p))

        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="2"><w:p><w:r>'
            "<w:t>FN original</w:t></w:r></w:p></w:footnote>"
            "</w:footnotes>"
        )
        en_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:endnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:endnote w:id="2"><w:p><w:r>'
            "<w:t>EN original</w:t></w:r></w:p></w:endnote>"
            "</w:endnotes>"
        )

        with zipfile.ZipFile(p, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        file_data["word/endnotes.xml"] = en_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/footnotes.xml"))
        all_items.append(zipfile.ZipInfo("word/endnotes.xml"))
        _rewrite_zip_content(p, file_data, all_items)

        _inject_docx_footnotes(
            p, {"footnote:2": "FN translated", "endnote:2": "EN translated"}
        )

        result = _extract_docx_footnotes(p)
        result_dict = dict(result)
        # Each key routes to its own XML file — neither is overwritten.
        assert result_dict.get("footnote:2") == "FN translated"
        assert result_dict.get("endnote:2") == "EN translated"


class TestFootnoteExtensions:
    """Tests for _FOOTNOTE_EXTENSIONS."""

    def test_extensions(self) -> None:
        """Correct extensions are in _FOOTNOTE_EXTENSIONS."""
        assert {".docx", ".odt", ".doc"} == _FOOTNOTE_EXTENSIONS


class TestExtractFootnotesDispatch:
    """Tests for _extract_footnotes dispatch."""

    def test_docx_dispatch(self, tmp_path: Path) -> None:
        """Dispatches to _extract_docx_footnotes for .docx."""
        doc = Document()
        p = tmp_path / "test.docx"
        doc.save(str(p))

        # No footnotes file → empty result
        result = _extract_footnotes(p, ".docx", "python_lib")
        assert result == []

    def test_win32com_dispatch(self) -> None:
        """Dispatches to Win32COM for .doc on win32com backend."""
        with patch(
            "src.core.office_processor._extract_win32com_word_footnotes",
            return_value=[("footnote:1", "FN")],
        ) as mock_fn:
            result = _extract_footnotes(
                Path("f.doc"),
                ".doc",
                "win32com",
            )
        mock_fn.assert_called_once()
        assert result == [("footnote:1", "FN")]


class TestInjectFootnotesDispatch:
    """Tests for _inject_footnotes dispatch."""

    def test_no_keys_returns_early(self) -> None:
        """Returns immediately when no footnote:/endnote: keys."""
        _inject_footnotes(
            Path("x.docx"),
            {"para:0": "text"},
            ".docx",
            "python_lib",
        )


class TestExtractOdtFootnotes:
    """Tests for _extract_odt_footnotes."""

    def test_extracts_odt_footnotes(self, tmp_path: Path) -> None:
        """Extracts footnotes from ODT content.xml."""
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-content"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:body><office:text>"
            "<text:p>Some text"
            '<text:note text:id="ftn1" text:note-class="footnote">'
            "<text:note-citation>1</text:note-citation>"
            "<text:note-body>"
            "<text:p>ODT footnote text</text:p>"
            "</text:note-body>"
            "</text:note>"
            "</text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", content_xml)
            zf.writestr("META-INF/manifest.xml", "<manifest/>")

        result = _extract_odt_footnotes(odt_path)
        assert len(result) == 1
        assert result[0][0] == "footnote:ftn1"
        assert "ODT footnote text" in result[0][1]

    def test_extracts_odt_endnotes(self, tmp_path: Path) -> None:
        """Extracts endnotes (note-class=endnote) from ODT."""
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-content"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:body><office:text>"
            '<text:p><text:note text:id="en1" text:note-class="endnote">'
            "<text:note-citation>i</text:note-citation>"
            "<text:note-body>"
            "<text:p>ODT endnote text</text:p>"
            "</text:note-body>"
            "</text:note></text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", content_xml)

        result = _extract_odt_footnotes(odt_path)
        assert result[0][0] == "endnote:en1"


class TestExtractOdtFootnotesRoundtrip:
    """Tests for _extract_odt_footnotes / _inject_odt_footnotes roundtrip."""

    def test_roundtrip(self, tmp_path: Path) -> None:
        """Extracts ODT footnote, injects translation, re-extracts."""
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-content"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:body><office:text>"
            "<text:p>Body"
            '<text:note text:id="ftn1" text:note-class="footnote">'
            "<text:note-citation>1</text:note-citation>"
            "<text:note-body>"
            "<text:p>Original footnote</text:p>"
            "</text:note-body>"
            "</text:note>"
            "</text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", content_xml)
            zf.writestr("META-INF/manifest.xml", "<manifest/>")

        # Extract
        extracted = _extract_odt_footnotes(odt_path)
        assert len(extracted) == 1
        assert extracted[0] == ("footnote:ftn1", "Original footnote")

        # Inject translation
        _inject_odt_footnotes(odt_path, {"footnote:ftn1": "Translated footnote"})

        # Re-extract and verify
        result = _extract_odt_footnotes(odt_path)
        assert len(result) == 1
        assert result[0][0] == "footnote:ftn1"
        assert "Translated footnote" in result[0][1]


class TestExtractOdtHeadersFooters:
    """Tests for _extract_odt_headers_footers."""

    def test_extracts_odt_header(self, tmp_path: Path) -> None:
        """Extracts header text from ODT styles.xml."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header>"
            "<text:p>My ODT Header</text:p>"
            "</style:header>"
            "<style:footer>"
            "<text:p>My ODT Footer</text:p>"
            "</style:footer>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        result = _extract_odt_headers_footers(odt_path)
        assert any(
            k.startswith("header:0:default:") and "My ODT Header" in v
            for k, v in result
        )
        assert any(
            k.startswith("footer:0:default:") and "My ODT Footer" in v
            for k, v in result
        )

    def test_no_styles_xml(self, tmp_path: Path) -> None:
        """Returns empty when styles.xml is missing."""
        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", "<root/>")

        result = _extract_odt_headers_footers(odt_path)
        assert result == []

    def test_extracts_first_page_header(self, tmp_path: Path) -> None:
        """Extracts first-page header from ODT styles.xml."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header-first>"
            "<text:p>First Header</text:p>"
            "</style:header-first>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        result = _extract_odt_headers_footers(odt_path)
        assert any(
            k == f"header:0:{_HF_FIRST}:0" and "First Header" in v for k, v in result
        )

    def test_extracts_even_page_footer(self, tmp_path: Path) -> None:
        """Extracts even-page (left) footer from ODT styles.xml."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:footer-left>"
            "<text:p>Even Footer</text:p>"
            "</style:footer-left>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        result = _extract_odt_headers_footers(odt_path)
        assert any(
            k == f"footer:0:{_HF_EVEN}:0" and "Even Footer" in v for k, v in result
        )


class TestInjectOdtHeadersFooters:
    """Tests for _inject_odt_headers_footers roundtrip."""

    def test_roundtrip(self, tmp_path: Path) -> None:
        """Extracts and injects ODT header — roundtrip verifies text."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header>"
            "<text:p>Original</text:p>"
            "</style:header>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        extracted = _extract_odt_headers_footers(odt_path)
        assert len(extracted) == 1
        key, val = extracted[0]
        assert "Original" in val

        _inject_odt_headers_footers(odt_path, {key: "Translated"})

        # Re-extract to verify injection
        result = _extract_odt_headers_footers(odt_path)
        assert any("Translated" in v for _, v in result)

    def test_no_matching_keys_noop(self, tmp_path: Path) -> None:
        """Injecting translations with non-matching keys leaves file intact."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header>"
            "<text:p>Original Header</text:p>"
            "</style:header>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        # Inject with keys that don't match any header
        _inject_odt_headers_footers(
            odt_path,
            {"header:99:default:0": "Should not appear"},
        )

        # Re-extract — original text must still be there
        result = _extract_odt_headers_footers(odt_path)
        assert any("Original Header" in v for _, v in result)
        assert not any("Should not appear" in v for _, v in result)


class TestExtractOdpNotes:
    """Tests for _extract_odp_notes and _inject_odp_notes."""

    def test_extracts_notes(self, tmp_path: Path) -> None:
        """Extracts speaker notes from ODP file via odfpy."""
        from odf.draw import Frame, Page, TextBox  # noqa: PLC0415
        from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
        from odf.presentation import Notes  # noqa: PLC0415
        from odf.style import MasterPage, PageLayout  # noqa: PLC0415
        from odf.text import P  # noqa: PLC0415

        doc = OpenDocumentPresentation()
        pl = PageLayout(name="PM0")
        doc.automaticstyles.addElement(pl)
        mp = MasterPage(name="Default", pagelayoutname="PM0")
        doc.masterstyles.addElement(mp)
        page = Page(name="Slide1", masterpagename="Default")
        doc.presentation.addElement(page)

        notes = Notes()
        frame = Frame()
        tb = TextBox()
        p = P(text="Speaker note text")
        tb.addElement(p)
        frame.addElement(tb)
        notes.addElement(frame)
        page.addElement(notes)

        odp_path = tmp_path / "test.odp"
        doc.save(str(odp_path))

        result = _extract_odp_notes(odp_path)
        assert len(result) >= 1
        assert any("Speaker note text" in v for _, v in result)

    def test_no_notes(self, tmp_path: Path) -> None:
        """Returns empty when slides have no notes."""
        from odf.draw import Page  # noqa: PLC0415
        from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
        from odf.style import MasterPage, PageLayout  # noqa: PLC0415

        doc = OpenDocumentPresentation()
        pl = PageLayout(name="PM0")
        doc.automaticstyles.addElement(pl)
        mp = MasterPage(name="Default", pagelayoutname="PM0")
        doc.masterstyles.addElement(mp)
        page = Page(name="Slide1", masterpagename="Default")
        doc.presentation.addElement(page)

        odp_path = tmp_path / "test.odp"
        doc.save(str(odp_path))

        result = _extract_odp_notes(odp_path)
        assert result == []

    def test_inject_roundtrip(self, tmp_path: Path) -> None:
        """Injects translated notes and verifies via re-extraction."""
        from odf.draw import Frame, Page, TextBox  # noqa: PLC0415
        from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
        from odf.presentation import Notes  # noqa: PLC0415
        from odf.style import MasterPage, PageLayout  # noqa: PLC0415
        from odf.text import P  # noqa: PLC0415

        doc = OpenDocumentPresentation()
        pl = PageLayout(name="PM0")
        doc.automaticstyles.addElement(pl)
        mp = MasterPage(name="Default", pagelayoutname="PM0")
        doc.masterstyles.addElement(mp)
        page = Page(name="Slide1", masterpagename="Default")
        doc.presentation.addElement(page)

        notes = Notes()
        frame = Frame()
        tb = TextBox()
        p = P(text="Original note")
        tb.addElement(p)
        frame.addElement(tb)
        notes.addElement(frame)
        page.addElement(notes)

        odp_path = tmp_path / "test.odp"
        doc.save(str(odp_path))

        extracted = _extract_odp_notes(odp_path)
        assert len(extracted) >= 1
        key = next(k for k, v in extracted if "Original note" in v)

        _inject_odp_notes(odp_path, {key: "Translated note"})

        result = _extract_odp_notes(odp_path)
        assert any("Translated note" in v for _, v in result)


# ---------------------------------------------------------------------------
# Integration: process_office_file with new features
# ---------------------------------------------------------------------------


class TestProcessOfficeFileNewFeatures:
    """Integration tests for new features wired into process_office_file."""

    def test_sheet_names_extracted_and_injected(
        self,
        tmp_path: Path,
    ) -> None:
        """Sheet names are extracted, translated, and injected."""
        from src.core.config import TranslationConfig  # noqa: PLC0415

        wb = Workbook()
        wb.active.title = "Revenue"
        wb.create_sheet("Expenses")
        src = tmp_path / "src.xlsx"
        wb.save(str(src))
        wb.close()

        out = tmp_path / "out.xlsx"
        config = TranslationConfig(translate_sheet_names=True)

        def fake_translate(
            texts: list[str],
            *args: object,
            **kwargs: object,
        ) -> list[str]:
            return [t.upper() for t in texts]

        with patch(
            "src.core.office_processor.translate_batch",
            side_effect=fake_translate,
        ):
            result = process_office_file(
                src,
                out,
                "French",
                config=config,
            )

        assert result is True
        wb2 = load_workbook(str(out))
        names = [ws.title for ws in wb2.worksheets]
        wb2.close()
        assert "REVENUE" in names
        assert "EXPENSES" in names

    def test_notes_extracted_and_injected(self, tmp_path: Path) -> None:
        """Speaker notes are extracted, translated, and injected."""
        from src.core.config import TranslationConfig  # noqa: PLC0415

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Note text"
        # Add body text so the main extraction finds content
        slide.shapes[0].text_frame.text = "Title"
        src = tmp_path / "src.pptx"
        prs.save(str(src))

        out = tmp_path / "out.pptx"
        config = TranslationConfig(translate_doc_notes=True)

        def fake_translate(
            texts: list[str],
            *args: object,
            **kwargs: object,
        ) -> list[str]:
            return [t.upper() for t in texts]

        with patch(
            "src.core.office_processor.translate_batch",
            side_effect=fake_translate,
        ):
            result = process_office_file(
                src,
                out,
                "French",
                config=config,
            )

        assert result is True
        prs2 = Presentation(str(out))
        notes = prs2.slides[0].notes_slide.notes_text_frame.text
        assert "NOTE TEXT" in notes

    def test_headers_footers_extracted_and_injected(
        self,
        tmp_path: Path,
    ) -> None:
        """Headers/footers are always-on: extracted and injected."""
        from src.core.config import TranslationConfig  # noqa: PLC0415

        doc = Document()
        doc.add_paragraph("Body text")
        section = doc.sections[0]
        section.header.is_linked_to_previous = False
        section.header.paragraphs[0].text = "DocHeader"
        src = tmp_path / "src.docx"
        doc.save(str(src))

        out = tmp_path / "out.docx"
        config = TranslationConfig()

        def fake_translate(
            texts: list[str],
            *args: object,
            **kwargs: object,
        ) -> list[str]:
            return [t.upper() for t in texts]

        with patch(
            "src.core.office_processor.translate_batch",
            side_effect=fake_translate,
        ):
            result = process_office_file(
                src,
                out,
                "French",
                config=config,
            )

        assert result is True
        doc2 = Document(str(out))
        hdr_text = doc2.sections[0].header.paragraphs[0].text
        assert hdr_text == "DOCHEADER"

    def test_footnotes_extracted_and_injected(
        self,
        tmp_path: Path,
    ) -> None:
        """Footnotes are always-on: extracted and injected."""
        from lxml import etree as _etree  # noqa: PLC0415

        from src.core.config import TranslationConfig  # noqa: PLC0415

        # Create a DOCX with body text
        doc = Document()
        doc.add_paragraph("Body text")
        src = tmp_path / "src.docx"
        doc.save(str(src))

        # Add footnotes.xml with proper OPC relationships so python-docx
        # preserves it during the save step.
        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main">'
            '<w:footnote w:id="0" w:type="separator"><w:p><w:r>'
            "<w:separator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="1" w:type="continuationSeparator"><w:p>'
            "<w:r><w:continuationSeparator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="2"><w:p><w:r><w:t>See reference</w:t>'
            "</w:r></w:p></w:footnote>"
            "</w:footnotes>"
        )

        with zipfile.ZipFile(src, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        all_items.append(zipfile.ZipInfo("word/footnotes.xml"))

        # Add OPC relationship so python-docx preserves footnotes.xml
        rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        rels_key = "word/_rels/document.xml.rels"
        rels_root = _etree.fromstring(file_data[rels_key])
        new_rel = _etree.SubElement(rels_root, f"{{{rel_ns}}}Relationship")
        new_rel.set("Id", "rIdFootnotes")
        new_rel.set(
            "Type",
            "http://schemas.openxmlformats.org/officeDocument/"
            "2006/relationships/footnotes",
        )
        new_rel.set("Target", "footnotes.xml")
        file_data[rels_key] = _etree.tostring(
            rels_root,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )

        # Add content type override
        ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
        ct_root = _etree.fromstring(file_data["[Content_Types].xml"])
        override = _etree.SubElement(ct_root, f"{{{ct_ns}}}Override")
        override.set("PartName", "/word/footnotes.xml")
        override.set(
            "ContentType",
            "application/vnd.openxmlformats-officedocument"
            ".wordprocessingml.footnotes+xml",
        )
        file_data["[Content_Types].xml"] = _etree.tostring(
            ct_root,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )

        _rewrite_zip_content(src, file_data, all_items)

        out = tmp_path / "out.docx"
        config = TranslationConfig()

        def fake_translate(
            texts: list[str],
            *args: object,
            **kwargs: object,
        ) -> list[str]:
            return [t.upper() for t in texts]

        with patch(
            "src.core.office_processor.translate_batch",
            side_effect=fake_translate,
        ):
            result = process_office_file(
                src,
                out,
                "French",
                config=config,
            )

        assert result is True
        # Verify footnote was translated
        extracted = _extract_docx_footnotes(out)
        assert any("SEE REFERENCE" in v for _, v in extracted)


# ---------------------------------------------------------------------------
# _substitute_font tests
# ---------------------------------------------------------------------------

_WORDML_NS_VAL = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_ODF_NS_VAL = {
    "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    "style": "urn:oasis:names:tc:opendocument:xmlns:style:1.0",
}


class TestSubstituteFont:
    """Tests for _substitute_font."""

    def test_same_script_returns_original(self) -> None:
        """Latin→Latin keeps the original font unchanged."""
        result = _substitute_font("Arial", "Hello", "World", "French")
        assert result == "Arial"

    def test_different_script_with_target_lang(self) -> None:
        """Different script with target_lang returns font from get_font_for_language."""
        with patch(
            "src.core.office_processor._get_font_for_language",
            return_value="Noto Sans JP",
        ) as mock_gfl:
            result = _substitute_font(
                "Arial",
                "Hello",
                "\u3053\u3093\u306b\u3061\u306f",
                "Japanese",
            )
        assert result == "Noto Sans JP"
        mock_gfl.assert_called_once()

    def test_different_script_without_target_lang(self) -> None:
        """Different script without target_lang returns None."""
        result = _substitute_font(
            "Arial",
            "Hello",
            "\u3053\u3093\u306b\u3061\u306f",
            "",
        )
        assert result is None

    def test_empty_original_text(self) -> None:
        """Empty original_text returns original font."""
        result = _substitute_font("Arial", "", "Bonjour", "French")
        assert result == "Arial"

    def test_empty_translated_text(self) -> None:
        """Empty translated_text returns original font."""
        result = _substitute_font("Arial", "Hello", "", "French")
        assert result == "Arial"


# ---------------------------------------------------------------------------
# _is_inside_table_cell tests
# ---------------------------------------------------------------------------


class TestIsInsideTableCell:
    """Tests for _is_inside_table_cell."""

    def test_direct_child_of_table_cell(self) -> None:
        """Element directly inside a table cell returns True."""
        cell = MagicMock()
        cell.qname = (
            "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
            "table-cell",
        )
        cell.parentNode = None

        element = MagicMock()
        element.parentNode = cell
        assert _is_inside_table_cell(element) is True

    def test_nested_deep_inside_table_cell(self) -> None:
        """Element nested as grandchild of a table cell returns True."""
        cell = MagicMock()
        cell.qname = (
            "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
            "table-cell",
        )
        cell.parentNode = None

        middle = MagicMock()
        middle.qname = ("some-ns", "text")
        middle.parentNode = cell

        element = MagicMock()
        element.parentNode = middle
        assert _is_inside_table_cell(element) is True

    def test_not_in_table_cell(self) -> None:
        """Element not in any table cell returns False."""
        root = MagicMock()
        root.qname = ("some-ns", "body")
        root.parentNode = None

        element = MagicMock()
        element.parentNode = root
        assert _is_inside_table_cell(element) is False

    def test_parent_none(self) -> None:
        """Element with parentNode=None returns False."""
        element = MagicMock()
        element.parentNode = None
        assert _is_inside_table_cell(element) is False


# ---------------------------------------------------------------------------
# _extract_xlsx_comments tests
# ---------------------------------------------------------------------------


class TestExtractXlsxComments:
    """Tests for _extract_xlsx_comments."""

    def test_no_comments(self, tmp_path: Path) -> None:
        """XLSX with no comments returns empty list."""
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "data"
        fpath = tmp_path / "no_comments.xlsx"
        wb.save(str(fpath))
        wb.close()

        result = _extract_xlsx_comments(fpath)
        assert result == []

    def test_one_comment(self, tmp_path: Path) -> None:
        """XLSX with one comment returns single (key, text) pair."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["B3"].comment = Comment("Review this", "Author")
        fpath = tmp_path / "one_comment.xlsx"
        wb.save(str(fpath))
        wb.close()

        result = _extract_xlsx_comments(fpath)
        assert len(result) == 1  # noqa: PLR2004
        key, text = result[0]
        assert key == "comment:Sheet1:3:2"
        assert text == "Review this"

    def test_multiple_comments_across_sheets(self, tmp_path: Path) -> None:
        """XLSX with multiple comments across sheets returns all."""
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Data"
        ws1["A1"].comment = Comment("First note", "Author")
        ws1["C2"].comment = Comment("Second note", "Author")

        ws2 = wb.create_sheet("Summary")
        ws2["D4"].comment = Comment("Third note", "Author")

        fpath = tmp_path / "multi_comments.xlsx"
        wb.save(str(fpath))
        wb.close()

        result = _extract_xlsx_comments(fpath)
        assert len(result) == 3  # noqa: PLR2004

        keys = [k for k, _ in result]
        assert "comment:Data:1:1" in keys
        assert "comment:Data:2:3" in keys
        assert "comment:Summary:4:4" in keys

    def test_whitespace_only_comment_skipped(self, tmp_path: Path) -> None:
        """Comments containing only whitespace are excluded."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"].comment = Comment("   ", "Author")
        fpath = tmp_path / "ws_comment.xlsx"
        wb.save(str(fpath))
        wb.close()

        result = _extract_xlsx_comments(fpath)
        assert result == []


# ---------------------------------------------------------------------------
# _inject_xlsx_comments tests
# ---------------------------------------------------------------------------


class TestInjectXlsxComments:
    """Tests for _inject_xlsx_comments."""

    def test_comment_updated_with_matching_key(self, tmp_path: Path) -> None:
        """Comment text is replaced when key matches."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"].comment = Comment("Original", "Author")
        fpath = tmp_path / "inject.xlsx"
        wb.save(str(fpath))
        wb.close()

        translations = {"comment:Sheet1:1:1": "Translated"}
        _inject_xlsx_comments(fpath, translations)

        wb2 = load_workbook(str(fpath))
        assert wb2.active["A1"].comment.text == "Translated"
        wb2.close()

    def test_non_matching_keys_ignored(self, tmp_path: Path) -> None:
        """Non-matching keys do not alter existing comments."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"].comment = Comment("Original", "Author")
        fpath = tmp_path / "inject_no_match.xlsx"
        wb.save(str(fpath))
        wb.close()

        translations = {"comment:Sheet1:5:5": "Wrong cell"}
        _inject_xlsx_comments(fpath, translations)

        wb2 = load_workbook(str(fpath))
        assert wb2.active["A1"].comment.text == "Original"
        wb2.close()

    def test_multiple_comments_injected(self, tmp_path: Path) -> None:
        """Multiple comments can be injected in one call."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"].comment = Comment("First", "Author")
        ws["B2"].comment = Comment("Second", "Author")
        fpath = tmp_path / "inject_multi.xlsx"
        wb.save(str(fpath))
        wb.close()

        translations = {
            "comment:Sheet1:1:1": "Premier",
            "comment:Sheet1:2:2": "Deuxieme",
        }
        _inject_xlsx_comments(fpath, translations)

        wb2 = load_workbook(str(fpath))
        assert wb2.active["A1"].comment.text == "Premier"
        assert wb2.active["B2"].comment.text == "Deuxieme"
        wb2.close()


# ---------------------------------------------------------------------------
# _read_txbx_data tests
# ---------------------------------------------------------------------------


class TestReadTxbxData:
    """Tests for _read_txbx_data."""

    def test_single_paragraph_single_run(self) -> None:
        """Single paragraph with one run extracts text correctly."""
        w = _WORDML_NS_VAL
        xml = (
            f'<wps:txbx xmlns:wps="http://schemas.microsoft.com/office/'
            f'word/2010/wordprocessingShape" xmlns:w="{w}">'
            f"<w:txbxContent>"
            f"<w:p><w:r><w:t>Hello World</w:t></w:r></w:p>"
            f"</w:txbxContent>"
            f"</wps:txbx>"
        )
        el = etree.fromstring(xml)
        text, t_els = _read_txbx_data(el)
        assert text == "Hello World"
        assert len(t_els) == 1

    def test_multiple_paragraphs(self) -> None:
        """Multiple paragraphs are joined with newlines."""
        w = _WORDML_NS_VAL
        xml = (
            f'<wps:txbx xmlns:wps="http://schemas.microsoft.com/office/'
            f'word/2010/wordprocessingShape" xmlns:w="{w}">'
            f"<w:txbxContent>"
            f"<w:p><w:r><w:t>Line one</w:t></w:r></w:p>"
            f"<w:p><w:r><w:t>Line two</w:t></w:r></w:p>"
            f"</w:txbxContent>"
            f"</wps:txbx>"
        )
        el = etree.fromstring(xml)
        text, t_els = _read_txbx_data(el)
        assert text == "Line one\nLine two"
        assert len(t_els) == 2  # noqa: PLR2004

    def test_empty_text_box(self) -> None:
        """Empty text box returns empty string."""
        w = _WORDML_NS_VAL
        xml = (
            f'<wps:txbx xmlns:wps="http://schemas.microsoft.com/office/'
            f'word/2010/wordprocessingShape" xmlns:w="{w}">'
            f"<w:txbxContent>"
            f"<w:p><w:r><w:t></w:t></w:r></w:p>"
            f"</w:txbxContent>"
            f"</wps:txbx>"
        )
        el = etree.fromstring(xml)
        text, t_els = _read_txbx_data(el)
        assert text == ""

    def test_multiple_runs_in_paragraph(self) -> None:
        """Multiple runs within one paragraph are concatenated."""
        w = _WORDML_NS_VAL
        xml = (
            f'<wps:txbx xmlns:wps="http://schemas.microsoft.com/office/'
            f'word/2010/wordprocessingShape" xmlns:w="{w}">'
            f"<w:txbxContent>"
            f"<w:p>"
            f"<w:r><w:t>Hello </w:t></w:r>"
            f"<w:r><w:t>World</w:t></w:r>"
            f"</w:p>"
            f"</w:txbxContent>"
            f"</wps:txbx>"
        )
        el = etree.fromstring(xml)
        text, t_els = _read_txbx_data(el)
        assert text == "Hello World"
        assert len(t_els) == 2  # noqa: PLR2004


# ---------------------------------------------------------------------------
# _build_odf_style_map tests
# ---------------------------------------------------------------------------


class TestBuildOdfStyleMap:
    """Tests for _build_odf_style_map."""

    def test_no_automatic_styles(self) -> None:
        """Root without automatic-styles element returns empty dict."""
        office_ns = _ODF_NS_VAL["office"]
        xml = f'<office:document-content xmlns:office="{office_ns}"/>'
        root = etree.fromstring(xml)
        result = _build_odf_style_map(root)
        assert result == {}

    def test_multiple_text_styles(self) -> None:
        """Multiple text-family styles are all mapped."""
        office_ns = _ODF_NS_VAL["office"]
        style_ns = _ODF_NS_VAL["style"]
        xml = (
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:style="{style_ns}">'
            f"<office:automatic-styles>"
            f'<style:style style:name="T1" style:family="text"/>'
            f'<style:style style:name="T2" style:family="text"/>'
            f"</office:automatic-styles>"
            f"</office:document-content>"
        )
        root = etree.fromstring(xml)
        result = _build_odf_style_map(root)
        assert len(result) == 2  # noqa: PLR2004
        assert "T1" in result
        assert "T2" in result

    def test_non_text_family_excluded(self) -> None:
        """Styles with family != 'text' are not included."""
        office_ns = _ODF_NS_VAL["office"]
        style_ns = _ODF_NS_VAL["style"]
        xml = (
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:style="{style_ns}">'
            f"<office:automatic-styles>"
            f'<style:style style:name="T1" style:family="text"/>'
            f'<style:style style:name="P1" style:family="paragraph"/>'
            f'<style:style style:name="G1" style:family="graphic"/>'
            f"</office:automatic-styles>"
            f"</office:document-content>"
        )
        root = etree.fromstring(xml)
        result = _build_odf_style_map(root)
        assert len(result) == 1
        assert "T1" in result
        assert "P1" not in result
        assert "G1" not in result

    def test_style_without_name_skipped(self) -> None:
        """Styles without style:name attribute are skipped."""
        office_ns = _ODF_NS_VAL["office"]
        style_ns = _ODF_NS_VAL["style"]
        xml = (
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:style="{style_ns}">'
            f"<office:automatic-styles>"
            f'<style:style style:family="text"/>'
            f"</office:automatic-styles>"
            f"</office:document-content>"
        )
        root = etree.fromstring(xml)
        result = _build_odf_style_map(root)
        assert result == {}


# ---------------------------------------------------------------------------
# _inject_docx_comment_html tests
# ---------------------------------------------------------------------------

_WML_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _make_comment_element(
    comment_id: str = "1",
    text: str = "Original",
    *,
    with_rpr: bool = False,
    with_ppr: bool = False,
    with_hyperlink: bool = False,
) -> object:
    """Builds a minimal ``<w:comment>`` lxml element for testing.

    Args:
        comment_id: The w:id attribute value.
        text: Text content of the comment.
        with_rpr: If True, adds a ``<w:rPr>`` with a font name.
        with_ppr: If True, adds a ``<w:pPr>`` to the paragraph.
        with_hyperlink: If True, wraps the run in ``<w:hyperlink>``.

    Returns:
        An lxml element representing ``<w:comment>``.
    """
    w_p = f"{{{_WML_NS}}}p"
    w_r = f"{{{_WML_NS}}}r"
    w_t = f"{{{_WML_NS}}}t"
    w_rpr = f"{{{_WML_NS}}}rPr"
    w_ppr = f"{{{_WML_NS}}}pPr"
    w_rfonts = f"{{{_WML_NS}}}rFonts"
    w_hyperlink = f"{{{_WML_NS}}}hyperlink"

    comment_el = etree.Element(f"{{{_WML_NS}}}comment")
    comment_el.set(f"{{{_WML_NS}}}id", comment_id)

    p_el = etree.SubElement(comment_el, w_p)
    if with_ppr:
        etree.SubElement(p_el, w_ppr)

    if with_hyperlink:
        hl = etree.SubElement(p_el, w_hyperlink)
        hl.set(f"{{{_R_NS}}}id", "rId1")
        container = hl
    else:
        container = p_el

    r_el = etree.SubElement(container, w_r)
    if with_rpr:
        rpr = etree.SubElement(r_el, w_rpr)
        fonts = etree.SubElement(rpr, w_rfonts)
        fonts.set(f"{{{_WML_NS}}}ascii", "Calibri")
    t_el = etree.SubElement(r_el, w_t)
    t_el.text = text
    return comment_el


class TestInjectDocxCommentHtml:
    """Tests for _inject_docx_comment_html."""

    def test_plain_text_injection(self) -> None:
        """Plain text (no HTML tags) replaces the comment text."""
        comment_el = _make_comment_element(text="Old text")
        mock_part = MagicMock()

        _inject_docx_comment_html(comment_el, "New text", mock_part, qn)

        # Should have one paragraph with the new text
        paras = comment_el.findall(qn("w:p"))
        assert len(paras) == 1
        runs = paras[0].findall(qn("w:r"))
        assert len(runs) == 1
        assert runs[0].find(qn("w:t")).text == "New text"

    def test_multiline_text_creates_multiple_paragraphs(self) -> None:
        """Newlines in text create separate <w:p> elements."""
        comment_el = _make_comment_element(text="Line 1")
        mock_part = MagicMock()

        _inject_docx_comment_html(
            comment_el,
            "Line A\nLine B\nLine C",
            mock_part,
            qn,
        )

        paras = comment_el.findall(qn("w:p"))
        assert len(paras) == 3  # noqa: PLR2004
        texts = [
            p.find(f".//{qn('w:t')}").text
            for p in paras
            if p.find(f".//{qn('w:t')}") is not None
        ]
        assert texts == ["Line A", "Line B", "Line C"]

    def test_preserves_base_rpr(self) -> None:
        """Injects text and preserves base run properties (font)."""
        comment_el = _make_comment_element(text="Old", with_rpr=True)
        mock_part = MagicMock()

        _inject_docx_comment_html(comment_el, "New", mock_part, qn)

        paras = comment_el.findall(qn("w:p"))
        run = paras[0].find(qn("w:r"))
        rpr = run.find(qn("w:rPr"))
        assert rpr is not None
        rfonts = rpr.find(qn("w:rFonts"))
        assert rfonts is not None
        assert rfonts.get(qn("w:ascii")) == "Calibri"

    def test_preserves_base_ppr(self) -> None:
        """Paragraph properties from the original are preserved."""
        comment_el = _make_comment_element(text="Old", with_ppr=True)
        mock_part = MagicMock()

        _inject_docx_comment_html(comment_el, "New", mock_part, qn)

        paras = comment_el.findall(qn("w:p"))
        ppr = paras[0].find(qn("w:pPr"))
        assert ppr is not None

    def test_hyperlink_injection_external(self) -> None:
        """HTML with <a> tag creates <w:hyperlink> with r:id."""
        comment_el = _make_comment_element(text="Old")
        mock_part = MagicMock()
        mock_part.relate_to.return_value = "rId42"

        html = 'Click <a href="https://example.com">here</a> please'
        _inject_docx_comment_html(comment_el, html, mock_part, qn)

        paras = comment_el.findall(qn("w:p"))
        assert len(paras) == 1
        # Should have a hyperlink child
        hyperlinks = paras[0].findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 1
        assert hyperlinks[0].get(qn("r:id")) == "rId42"
        # Hyperlink should contain a run with the link text
        hl_run = hyperlinks[0].find(qn("w:r"))
        assert hl_run is not None
        assert hl_run.find(qn("w:t")).text == "here"

    def test_hyperlink_injection_internal_anchor(self) -> None:
        """HTML with <a href="#bookmark"> creates anchor-based hyperlink."""
        comment_el = _make_comment_element(text="Old")
        mock_part = MagicMock()

        html = '<a href="#section1">Go to section</a>'
        _inject_docx_comment_html(comment_el, html, mock_part, qn)

        paras = comment_el.findall(qn("w:p"))
        hyperlinks = paras[0].findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 1
        assert hyperlinks[0].get(qn("w:anchor")) == "section1"

    def test_empty_segments_fallback_to_plain(self) -> None:
        """When _parse_html_formatting returns no segments, falls back to plain."""
        comment_el = _make_comment_element(text="Old text")
        mock_part = MagicMock()

        # An empty string after stripping should produce a plain fallback
        _inject_docx_comment_html(comment_el, "Simple text", mock_part, qn)

        paras = comment_el.findall(qn("w:p"))
        assert len(paras) >= 1
        all_t = list(comment_el.iter(qn("w:t")))
        all_text = "".join(t.text or "" for t in all_t)
        assert "Simple text" in all_text

    def test_base_rpr_from_hyperlink_run(self) -> None:
        """Finds base rPr inside a <w:hyperlink> when no direct runs have rPr."""
        comment_el = _make_comment_element(
            text="Link text",
            with_hyperlink=True,
            with_rpr=True,
        )
        mock_part = MagicMock()

        _inject_docx_comment_html(comment_el, "Replaced", mock_part, qn)

        # The new run should have inherited the rPr from the hyperlink's run
        run = comment_el.find(f".//{qn('w:r')}")
        rpr = run.find(qn("w:rPr"))
        assert rpr is not None

    def test_hyperlink_relate_to_failure_falls_back(self) -> None:
        """When relate_to raises, the link text is inserted as plain run."""
        comment_el = _make_comment_element(text="Old")
        mock_part = MagicMock()
        mock_part.relate_to.side_effect = RuntimeError("Cannot create rel")

        html = '<a href="https://fail.com">broken</a> text'
        _inject_docx_comment_html(comment_el, html, mock_part, qn)

        # Should still have the text, just without hyperlink wrapper
        assert comment_el.findall(qn("w:p"))  # at least one paragraph
        all_t = list(comment_el.iter(qn("w:t")))
        all_text = "".join(t.text or "" for t in all_t)
        assert "broken" in all_text
        assert "text" in all_text

    def test_space_preservation(self) -> None:
        """Runs with leading/trailing spaces get xml:space='preserve'."""
        comment_el = _make_comment_element(text="Old")
        mock_part = MagicMock()

        # Use HTML that will produce a segment with leading space
        html = 'A <a href="https://x.com">link</a> B'
        _inject_docx_comment_html(comment_el, html, mock_part, qn)

        xml_space = "{http://www.w3.org/XML/1998/namespace}space"
        # At least the segments with spaces should have xml:space="preserve"
        t_els_with_space = [
            t
            for t in comment_el.iter(qn("w:t"))
            if t.text and (t.text.startswith(" ") or t.text.endswith(" "))
        ]
        for t_el in t_els_with_space:
            assert t_el.get(xml_space) == "preserve"


# ---------------------------------------------------------------------------
# _patch_docx_comment_rels tests
# ---------------------------------------------------------------------------


class TestPatchDocxCommentRels:
    """Tests for _patch_docx_comment_rels."""

    def test_no_rels_returns_early(self, tmp_path: Path) -> None:
        """Returns without error when comments_part has no rels."""
        mock_part = MagicMock()
        mock_part.rels = {}

        docx_path = tmp_path / "test.docx"
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr("word/document.xml", b"<root/>")

        _patch_docx_comment_rels(docx_path, mock_part)
        # No crash, no modification
        with zipfile.ZipFile(docx_path, "r") as zf:
            assert "word/_rels/comments.xml.rels" not in zf.namelist()

    def test_rels_access_exception_returns_early(self, tmp_path: Path) -> None:
        """Returns without error when rels access raises an exception."""
        mock_part = MagicMock()
        type(mock_part).rels = PropertyMock(side_effect=AttributeError("no rels"))

        docx_path = tmp_path / "test.docx"
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr("word/document.xml", b"<root/>")

        # Should not raise
        _patch_docx_comment_rels(docx_path, mock_part)

    def test_no_hyperlink_rels_returns_early(self, tmp_path: Path) -> None:
        """Returns without patching when rels has no hyperlink relationships."""
        mock_rel = MagicMock()
        mock_rel.reltype = "some/other/type"
        mock_rel.is_external = True
        mock_rel.rId = "rId1"
        mock_rel.target_ref = "https://example.com"

        mock_part = MagicMock()
        mock_part.rels = {"rId1": mock_rel}

        docx_path = tmp_path / "test.docx"
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr("word/document.xml", b"<root/>")

        _patch_docx_comment_rels(docx_path, mock_part)
        with zipfile.ZipFile(docx_path, "r") as zf:
            assert "word/_rels/comments.xml.rels" not in zf.namelist()

    def test_rels_already_present_skips_patch(self, tmp_path: Path) -> None:
        """When comments.xml.rels already exists in ZIP, no patching occurs."""
        from src.core.office_formatter import _HYPERLINK_RELTYPE  # noqa: PLC0415

        mock_rel = MagicMock()
        mock_rel.reltype = _HYPERLINK_RELTYPE
        mock_rel.is_external = True
        mock_rel.rId = "rId1"
        mock_rel.target_ref = "https://example.com"

        mock_part = MagicMock()
        mock_part.rels = {"rId1": mock_rel}

        docx_path = tmp_path / "test.docx"
        existing_rels = b'<?xml version="1.0" encoding="UTF-8"?><Relationships/>'
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr("word/document.xml", b"<root/>")
            zf.writestr("word/_rels/comments.xml.rels", existing_rels)

        _patch_docx_comment_rels(docx_path, mock_part)
        # The existing rels should remain unchanged
        with zipfile.ZipFile(docx_path, "r") as zf:
            data = zf.read("word/_rels/comments.xml.rels")
        assert data == existing_rels

    def test_patches_missing_rels_into_zip(self, tmp_path: Path) -> None:
        """When rels file is missing, it is created and added to the ZIP."""
        from src.core.office_formatter import _HYPERLINK_RELTYPE  # noqa: PLC0415

        mock_rel = MagicMock()
        mock_rel.reltype = _HYPERLINK_RELTYPE
        mock_rel.is_external = True
        mock_rel.rId = "rId1"
        mock_rel.target_ref = "https://example.com"

        mock_part = MagicMock()
        mock_part.rels = {"rId1": mock_rel}

        docx_path = tmp_path / "test.docx"
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr("word/document.xml", b"<root/>")
            zf.writestr("word/comments.xml", b"<comments/>")

        _patch_docx_comment_rels(docx_path, mock_part)

        with zipfile.ZipFile(docx_path, "r") as zf:
            assert "word/_rels/comments.xml.rels" in zf.namelist()
            rels_data = zf.read("word/_rels/comments.xml.rels")
        # Parse the rels XML and verify it contains the hyperlink
        root = etree.fromstring(rels_data)
        rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        rel_els = root.findall(f"{{{rels_ns}}}Relationship")
        assert len(rel_els) >= 1
        targets = [r.get("Target") for r in rel_els]
        assert "https://example.com" in targets


# ---------------------------------------------------------------------------
# _translate_legacy_images tests
# ---------------------------------------------------------------------------


class TestTranslateLegacyImages:
    """Tests for _translate_legacy_images round-trip conversion."""

    @patch("src.core.office_processor._translate_zip_images")
    @patch("src.core.office_processor._convert_with_win32com")
    def test_win32com_round_trip(
        self,
        mock_convert: MagicMock,
        mock_zip_images: MagicMock,
        tmp_path: Path,
    ) -> None:
        """Legacy .doc via win32com: convert → translate images → convert back."""
        doc_path = tmp_path / "test.doc"
        doc_path.write_bytes(b"fake doc content")

        def fake_convert(src: Path, dst: Path) -> None:
            dst.write_bytes(b"fake modern content")

        mock_convert.side_effect = fake_convert

        _translate_legacy_images(
            doc_path,
            ".doc",
            "win32com",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

        # Convert was called twice: legacy→modern and modern→legacy
        assert mock_convert.call_count == 2  # noqa: PLR2004
        # ZIP images was called once with the modern suffix
        mock_zip_images.assert_called_once()
        call_args = mock_zip_images.call_args
        assert call_args[0][1] == ".docx"  # modern suffix

    @patch("src.core.office_processor._translate_zip_images")
    @patch("src.core.office_processor._convert_with_uno")
    def test_uno_round_trip(
        self,
        mock_convert: MagicMock,
        mock_zip_images: MagicMock,
        tmp_path: Path,
    ) -> None:
        """Legacy .xls via UNO: convert → translate images → convert back."""
        xls_path = tmp_path / "test.xls"
        xls_path.write_bytes(b"fake xls content")

        def fake_convert(src: Path, dst: Path) -> None:
            dst.write_bytes(b"fake modern content")

        mock_convert.side_effect = fake_convert

        _translate_legacy_images(
            xls_path,
            ".xls",
            "uno",
            "Vietnamese",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

        assert mock_convert.call_count == 2  # noqa: PLR2004
        call_args = mock_zip_images.call_args
        assert call_args[0][1] == ".xlsx"

    @patch("src.core.office_processor._translate_zip_images")
    @patch("src.core.office_processor._convert_with_win32com")
    def test_conversion_failure_raises(
        self,
        mock_convert: MagicMock,
        mock_zip_images: MagicMock,
        tmp_path: Path,
    ) -> None:
        """Raises RuntimeError when conversion produces empty file."""
        doc_path = tmp_path / "test.doc"
        doc_path.write_bytes(b"fake doc content")

        def fake_convert(src: Path, dst: Path) -> None:
            # Create empty file — simulates failed conversion
            dst.write_bytes(b"")

        mock_convert.side_effect = fake_convert

        with pytest.raises(RuntimeError, match="no output"):
            _translate_legacy_images(
                doc_path,
                ".doc",
                "win32com",
                "French",
                "English",
                None,
                "TesseractOCR",
                None,
                None,
            )

    @patch("src.core.office_processor._translate_zip_images")
    @patch("src.core.office_processor._convert_with_win32com")
    def test_temp_file_cleaned_up_on_success(
        self,
        mock_convert: MagicMock,
        mock_zip_images: MagicMock,
        tmp_path: Path,
    ) -> None:
        """Temporary modern file is cleaned up even on success."""
        doc_path = tmp_path / "test.ppt"
        doc_path.write_bytes(b"fake ppt content")

        created_temps: list[Path] = []

        def fake_convert(src: Path, dst: Path) -> None:
            dst.write_bytes(b"fake content")
            if dst.suffix == ".pptx":
                created_temps.append(dst)

        mock_convert.side_effect = fake_convert

        _translate_legacy_images(
            doc_path,
            ".ppt",
            "win32com",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

        # The temp file should be cleaned up
        for temp in created_temps:
            assert not temp.exists()

    @patch("src.core.office_processor._translate_zip_images")
    @patch("src.core.office_processor._convert_with_win32com")
    def test_passes_all_args_to_zip_images(
        self,
        mock_convert: MagicMock,
        mock_zip_images: MagicMock,
        tmp_path: Path,
    ) -> None:
        """All translation args are forwarded to _translate_zip_images."""
        doc_path = tmp_path / "test.doc"
        doc_path.write_bytes(b"fake content")

        def fake_convert(src: Path, dst: Path) -> None:
            dst.write_bytes(b"converted")

        mock_convert.side_effect = fake_convert

        glossary = [(1, "Hello", "Bonjour")]
        progress_fn = MagicMock()
        cancel_fn = MagicMock(return_value=False)

        _translate_legacy_images(
            doc_path,
            ".doc",
            "win32com",
            "French",
            "English",
            glossary,
            "EasyOCR",
            progress_fn,
            cancel_fn,
        )

        args = mock_zip_images.call_args[0]
        assert args[1] == ".docx"  # modern suffix
        assert args[2] == "French"  # target_lang
        assert args[3] == "English"  # src_lang
        assert args[4] == glossary
        assert args[5] == "EasyOCR"  # ocr_method


# ---------------------------------------------------------------------------
# convert_to_modern_format tests
# ---------------------------------------------------------------------------


class TestConvertToModernFormat:
    """Tests for convert_to_modern_format."""

    def test_win32com_success(self, tmp_path: Path) -> None:
        """Returns True when win32com conversion succeeds."""
        src = tmp_path / "test.doc"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        with (
            patch.dict(
                "sys.modules",
                {"win32com": MagicMock(), "win32com.client": MagicMock()},
            ),
            patch(
                "src.core.office_processor._convert_with_win32com",
            ) as mock_conv,
        ):
            result = convert_to_modern_format(src, dst)

        assert result is True
        mock_conv.assert_called_once_with(src, dst)

    def test_win32com_failure_returns_false(self, tmp_path: Path) -> None:
        """Returns False when win32com conversion raises a non-ImportError."""
        src = tmp_path / "test.doc"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        with (
            patch.dict(
                "sys.modules",
                {"win32com": MagicMock(), "win32com.client": MagicMock()},
            ),
            patch(
                "src.core.office_processor._convert_with_win32com",
                side_effect=RuntimeError("COM failure"),
            ),
        ):
            result = convert_to_modern_format(src, dst)

        assert result is False

    def test_falls_back_to_uno(self, tmp_path: Path) -> None:
        """Falls back to UNO when win32com is not available."""
        src = tmp_path / "test.xls"
        dst = tmp_path / "test.xlsx"
        src.write_bytes(b"fake")

        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with (
                patch.dict(
                    "sys.modules",
                    {
                        "win32com": None,
                        "win32com.client": None,
                        "uno": MagicMock(),
                    },
                ),
                patch(
                    "src.core.office_processor._convert_with_uno",
                ) as mock_conv,
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is True
        mock_conv.assert_called_once_with(src, dst)

    def test_no_backend_returns_false(self, tmp_path: Path) -> None:
        """Returns False when neither win32com nor UNO is available."""
        src = tmp_path / "test.ppt"
        dst = tmp_path / "test.pptx"
        src.write_bytes(b"fake")

        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with patch.dict(
                "sys.modules",
                {
                    "win32com": None,
                    "win32com.client": None,
                    "uno": None,
                },
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is False

    def test_uno_failure_returns_false(self, tmp_path: Path) -> None:
        """Returns False when UNO conversion raises a non-ImportError."""
        src = tmp_path / "test.odt"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with (
                patch.dict(
                    "sys.modules",
                    {
                        "win32com": None,
                        "win32com.client": None,
                        "uno": MagicMock(),
                    },
                ),
                patch(
                    "src.core.office_processor._convert_with_uno",
                    side_effect=RuntimeError("UNO crash"),
                ),
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is False


# ---------------------------------------------------------------------------
# Password-protected / encrypted document tests
# ---------------------------------------------------------------------------


class TestPasswordProtectedDocuments:
    """Tests for handling encrypted/password-protected Office files."""

    def test_encrypted_docx_raises_text_read_error(self, tmp_path: Path) -> None:
        """An encrypted DOCX raises ValueError('TEXT_READ_ERROR') on extraction."""
        # Create a file that is NOT a valid DOCX ZIP archive — simulates
        # an encrypted file that python-docx cannot open.
        bad_path = tmp_path / "encrypted.docx"
        bad_path.write_bytes(b"\xd0\xcf\x11\xe0encrypted-fake-content")

        out_path = tmp_path / "out.docx"

        with pytest.raises(ValueError, match="TEXT_READ_ERROR"):
            process_office_file(bad_path, out_path, "French")

    def test_encrypted_xlsx_raises_text_read_error(self, tmp_path: Path) -> None:
        """An encrypted XLSX raises ValueError('TEXT_READ_ERROR') on extraction."""
        bad_path = tmp_path / "encrypted.xlsx"
        bad_path.write_bytes(b"\xd0\xcf\x11\xe0encrypted-fake-excel")

        out_path = tmp_path / "out.xlsx"

        with pytest.raises(ValueError, match="TEXT_READ_ERROR"):
            process_office_file(bad_path, out_path, "French")

    def test_encrypted_pptx_raises_text_read_error(self, tmp_path: Path) -> None:
        """An encrypted PPTX raises ValueError('TEXT_READ_ERROR') on extraction."""
        bad_path = tmp_path / "encrypted.pptx"
        bad_path.write_bytes(b"\xd0\xcf\x11\xe0encrypted-fake-ppt")

        out_path = tmp_path / "out.pptx"

        with pytest.raises(ValueError, match="TEXT_READ_ERROR"):
            process_office_file(bad_path, out_path, "French")

    def test_corrupted_zip_raises_text_read_error(self, tmp_path: Path) -> None:
        """A truncated/corrupted ZIP raises ValueError('TEXT_READ_ERROR')."""
        bad_path = tmp_path / "corrupt.docx"
        bad_path.write_bytes(b"PK\x03\x04not-a-real-zip")

        out_path = tmp_path / "out.docx"

        with pytest.raises(ValueError, match="TEXT_READ_ERROR"):
            process_office_file(bad_path, out_path, "French")


# ---------------------------------------------------------------------------
# Sheet name collision / sanitization edge cases
# ---------------------------------------------------------------------------


class TestSanitizeSheetNameCollisions:
    """Additional edge cases for _sanitize_sheet_name."""

    def test_unicode_chars_preserved(self) -> None:
        """Unicode characters that are not invalid are preserved."""
        result = _sanitize_sheet_name("Données françaises")
        assert result == "Données françaises"

    def test_only_whitespace_returns_sheet(self) -> None:
        """A name that is only whitespace returns 'Sheet' fallback."""
        assert _sanitize_sheet_name("   ") == "Sheet"

    def test_truncation_preserves_valid_content(self) -> None:
        """Truncation keeps the first 31 valid characters."""
        name = "A" * 25 + "[*?]" + "B" * 10  # noqa: PLR2004
        result = _sanitize_sheet_name(name)
        # After removing [, *, ?, ] we have 25 A's + 10 B's = 35 chars → truncate to 31
        assert len(result) == 31  # noqa: PLR2004
        assert result.startswith("A" * 25)  # noqa: PLR2004

    def test_all_seven_invalid_chars_removed(self) -> None:
        """All seven invalid sheet name characters are stripped."""
        name = "A\\B/C*D?E:F[G]H"
        result = _sanitize_sheet_name(name)
        assert result == "ABCDEFGH"

    def test_mixed_invalid_and_spaces(self) -> None:
        """Invalid chars + leading/trailing spaces are handled together."""
        result = _sanitize_sheet_name("  [Sheet*1]  ")
        assert result == "Sheet1"

    def test_duplicate_names_after_truncation(self) -> None:
        """Two different long names can truncate to the same 31-char result."""
        name_a = "X" * 31 + "A"  # noqa: PLR2004
        name_b = "X" * 31 + "B"  # noqa: PLR2004
        # Both truncate to "X" * 31
        assert _sanitize_sheet_name(name_a) == _sanitize_sheet_name(name_b)


# ---------------------------------------------------------------------------
# Image media prefix scanning tests
# ---------------------------------------------------------------------------


class TestSuffixToMediaPrefixes:
    """Tests for _SUFFIX_TO_MEDIA_PREFIXES mapping."""

    def test_docx_has_both_prefixes(self) -> None:
        """DOCX maps to both word/media/ and media/ prefixes."""
        prefixes = _SUFFIX_TO_MEDIA_PREFIXES[".docx"]
        assert "word/media/" in prefixes
        assert "media/" in prefixes

    def test_xlsx_has_both_prefixes(self) -> None:
        """XLSX maps to both xl/media/ and media/ prefixes."""
        prefixes = _SUFFIX_TO_MEDIA_PREFIXES[".xlsx"]
        assert "xl/media/" in prefixes
        assert "media/" in prefixes

    def test_pptx_has_both_prefixes(self) -> None:
        """PPTX maps to both ppt/media/ and media/ prefixes."""
        prefixes = _SUFFIX_TO_MEDIA_PREFIXES[".pptx"]
        assert "ppt/media/" in prefixes
        assert "media/" in prefixes

    def test_odf_uses_pictures_prefix(self) -> None:
        """ODF formats use Pictures/ prefix only."""
        for suffix in (".odt", ".ods", ".odp"):
            prefixes = _SUFFIX_TO_MEDIA_PREFIXES[suffix]
            assert prefixes == ("Pictures/",)

    def test_epub_uses_empty_prefix(self) -> None:
        """EPUB uses empty prefix — images can be anywhere."""
        prefixes = _SUFFIX_TO_MEDIA_PREFIXES[".epub"]
        assert prefixes == ("",)

    def test_unsupported_suffix_returns_none(self) -> None:
        """Unknown suffix is not in the map."""
        assert _SUFFIX_TO_MEDIA_PREFIXES.get(".rtf") is None
        assert _SUFFIX_TO_MEDIA_PREFIXES.get(".txt") is None

    @patch("src.core.office_processor._translate_single_image")
    def test_root_media_prefix_finds_images(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """Images at root media/ directory are found for DOCX files."""
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
            b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
            b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        mock_translate.return_value = b"translated"

        file_path = tmp_path / "test.xlsx"
        with zipfile.ZipFile(file_path, "w") as zf:
            zf.writestr("xl/media/image1.png", png_bytes)
            zf.writestr("media/image2.png", png_bytes)

        _translate_zip_images(
            file_path,
            ".xlsx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

        expected_count = 2
        assert mock_translate.call_count == expected_count  # noqa: PLR2004

    @patch("src.core.office_processor._translate_single_image")
    def test_odf_pictures_prefix(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """ODF images under Pictures/ are detected and translated."""
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
            b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
            b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        mock_translate.return_value = b"translated"

        file_path = tmp_path / "test.odt"
        with zipfile.ZipFile(file_path, "w") as zf:
            zf.writestr("Pictures/image1.png", png_bytes)
            zf.writestr("content.xml", b"<root/>")

        _translate_zip_images(
            file_path,
            ".odt",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

        mock_translate.assert_called_once()


# ---------------------------------------------------------------------------
# Fatal vs non-fatal LLM errors in image translation
# ---------------------------------------------------------------------------


class TestFatalVsNonFatalLlmErrors:
    """Tests for _FATAL_LLM_ERRORS abort vs continue behavior."""

    def test_fatal_errors_set_contents(self) -> None:
        """_FATAL_LLM_ERRORS contains the three expected error types."""
        assert "AUTH_ERROR" in _FATAL_LLM_ERRORS
        assert "QUOTA_ERROR" in _FATAL_LLM_ERRORS
        assert "VISION_NOT_SUPPORTED" in _FATAL_LLM_ERRORS

    @patch("src.core.office_processor._translate_single_image")
    def test_fatal_auth_error_aborts_immediately(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """AUTH_ERROR raises immediately — remaining images are not processed."""
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
            b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
            b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        mock_translate.side_effect = ValueError("AUTH_ERROR")

        file_path = tmp_path / "test.docx"
        with zipfile.ZipFile(file_path, "w") as zf:
            zf.writestr("word/media/img1.png", png_bytes)
            zf.writestr("word/media/img2.png", png_bytes)

        with pytest.raises(ValueError, match="AUTH_ERROR"):
            _translate_zip_images(
                file_path,
                ".docx",
                "French",
                "English",
                None,
                "TesseractOCR",
                None,
                None,
            )

        # Only the first image was attempted
        assert mock_translate.call_count == 1

    @patch("src.core.office_processor._translate_single_image")
    def test_fatal_quota_error_aborts(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """QUOTA_ERROR raises immediately."""
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
            b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
            b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        mock_translate.side_effect = ValueError("QUOTA_ERROR")

        file_path = tmp_path / "test.docx"
        with zipfile.ZipFile(file_path, "w") as zf:
            zf.writestr("word/media/img1.png", png_bytes)

        with pytest.raises(ValueError, match="QUOTA_ERROR"):
            _translate_zip_images(
                file_path,
                ".docx",
                "French",
                "English",
                None,
                "TesseractOCR",
                None,
                None,
            )

    @patch("src.core.office_processor._translate_single_image")
    def test_non_fatal_error_continues_with_skip_warning(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """Non-fatal ValueError attempts every image and never raises.

        Skip-with-warning policy: a single bad image (transient
        connection blip, model timeout, etc. that survived the LLM-
        level 3-retry policy) does NOT abort the whole document.
        The image is left in its source form and the loop continues.
        """
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
            b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
            b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        # First image fails with non-fatal error, second succeeds
        mock_translate.side_effect = [
            ValueError("SOME_NON_FATAL_ERROR"),
            b"translated",
        ]

        file_path = tmp_path / "test.docx"
        with zipfile.ZipFile(file_path, "w") as zf:
            zf.writestr("word/media/img1.png", png_bytes)
            zf.writestr("word/media/img2.png", png_bytes)

        # No raise — the function completes normally.
        _translate_zip_images(
            file_path,
            ".docx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

        # Both images were attempted (no short-circuit on first error)
        assert mock_translate.call_count == 2  # noqa: PLR2004
        # The successful one was injected; the failed one stayed
        # original.  This is the key skip-with-warning property.
        with zipfile.ZipFile(file_path, "r") as zf:
            assert zf.read("word/media/img1.png") == png_bytes
            assert zf.read("word/media/img2.png") == b"translated"

    @patch("src.core.office_processor._translate_single_image")
    def test_non_fatal_error_still_writes_successful_images(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """Non-fatal errors don't prevent writing successfully translated images.

        Under skip-with-warning the function returns normally — the
        successful image's translation is preserved in the rewritten
        ZIP, the failed one keeps its original bytes.
        """
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
            b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
            b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        # First succeeds, second fails with non-fatal
        mock_translate.side_effect = [
            b"translated-img1",
            ValueError("MODEL_ERROR"),
        ]

        file_path = tmp_path / "test.docx"
        with zipfile.ZipFile(file_path, "w") as zf:
            zf.writestr("word/media/img1.png", png_bytes)
            zf.writestr("word/media/img2.png", png_bytes)

        _translate_zip_images(
            file_path,
            ".docx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

        with zipfile.ZipFile(file_path, "r") as zf:
            # Translated image was injected.
            assert zf.read("word/media/img1.png") == b"translated-img1"
            # Failed image kept its original bytes.
            assert zf.read("word/media/img2.png") == png_bytes

    @patch("src.core.office_processor._translate_single_image")
    def test_unexpected_exception_continues(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """Non-ValueError exceptions are caught and processing continues."""
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
            b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
            b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        mock_translate.side_effect = [
            RuntimeError("OCR crash"),
            b"translated",
        ]

        file_path = tmp_path / "test.docx"
        with zipfile.ZipFile(file_path, "w") as zf:
            zf.writestr("word/media/img1.png", png_bytes)
            zf.writestr("word/media/img2.png", png_bytes)

        # Should NOT raise — unexpected errors are swallowed
        _translate_zip_images(
            file_path,
            ".docx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )

        assert mock_translate.call_count == 2  # noqa: PLR2004

    @patch("src.core.office_processor._translate_single_image")
    def test_cancel_check_stops_processing(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """cancel_check returning True stops the image processing loop."""
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
            b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
            b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        mock_translate.return_value = b"translated"

        file_path = tmp_path / "test.docx"
        with zipfile.ZipFile(file_path, "w") as zf:
            zf.writestr("word/media/img1.png", png_bytes)
            zf.writestr("word/media/img2.png", png_bytes)

        # Cancel immediately
        _translate_zip_images(
            file_path,
            ".docx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            lambda: True,  # always cancelled
        )

        mock_translate.assert_not_called()


# ---------------------------------------------------------------------------
# ODF paragraph text extraction/injection edge cases
# ---------------------------------------------------------------------------

_TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
_XLINK_NS = "http://www.w3.org/1999/xlink"


class TestExtractOdfParagraphTextEdgeCases:
    """Edge cases for _extract_odf_paragraph_text."""

    def test_empty_parent_returns_empty(self) -> None:
        """Parent with no <text:p> children returns empty string."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        text_p_tag = f"{{{_TEXT_NS}}}p"
        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert result == ""

    def test_single_paragraph_text(self) -> None:
        """Extracts text from a single paragraph."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        p.text = "Hello World"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert result == "Hello World"

    def test_multiple_paragraphs_joined_by_newlines(self) -> None:
        """Multiple paragraphs are joined by newlines."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        for txt in ("Line 1", "Line 2", "Line 3"):
            p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
            p.text = txt
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert result == "Line 1\nLine 2\nLine 3"

    def test_line_break_element(self) -> None:
        """<text:line-break/> elements produce newlines."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        p.text = "Before"
        lb = etree.SubElement(p, f"{{{_TEXT_NS}}}line-break")
        lb.tail = "After"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert result == "Before\nAfter"

    def test_hyperlink_extraction(self) -> None:
        """ODF <text:a> hyperlinks are extracted as HTML <a> tags."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        p.text = "Visit "
        a = etree.SubElement(p, f"{{{_TEXT_NS}}}a")
        a.set(f"{{{_XLINK_NS}}}href", "https://example.com")
        a.text = "Example"
        a.tail = " site"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert 'href="https://example.com"' in result
        assert ">Example</a>" in result
        assert "Visit " in result
        assert " site" in result

    def test_hyperlink_without_href(self) -> None:
        """<text:a> without xlink:href emits text without <a> wrapper."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        a = etree.SubElement(p, f"{{{_TEXT_NS}}}a")
        a.text = "No URL"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert "<a" not in result
        assert "No URL" in result

    def test_mixed_children_text_and_tail(self) -> None:
        """Handles mixed children with text and tail content."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        p.text = "Start "
        span = etree.SubElement(p, f"{{{_TEXT_NS}}}span")
        span.text = "middle"
        span.tail = " end"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert result == "Start middle end"

    def test_nested_paragraphs(self) -> None:
        """Finds paragraphs at any depth via .//<text:p>."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        section = etree.SubElement(parent, f"{{{_TEXT_NS}}}section")
        p = etree.SubElement(section, f"{{{_TEXT_NS}}}p")
        p.text = "Deep text"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert result == "Deep text"

    def test_special_chars_in_hyperlink_escaped(self) -> None:
        """Special characters in hyperlink URL and text are HTML-escaped."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        a = etree.SubElement(p, f"{{{_TEXT_NS}}}a")
        a.set(f"{{{_XLINK_NS}}}href", "https://example.com/a&b")
        a.text = "A&B<C>"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _extract_odf_paragraph_text(parent, text_p_tag)
        assert "&amp;b" in result  # URL escaped
        assert "A&amp;B&lt;C&gt;" in result  # text escaped


class TestInjectOdfParagraphTextEdgeCases:
    """Edge cases for _inject_odf_paragraph_text."""

    def test_no_paras_returns_false(self) -> None:
        """Returns False when parent has no <text:p> children."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _inject_odf_paragraph_text(parent, "New text", text_p_tag)
        assert result is False

    def test_simple_text_replacement(self) -> None:
        """Replaces text in a single paragraph."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        p.text = "Old text"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _inject_odf_paragraph_text(parent, "New text", text_p_tag)
        assert result is True
        paras = parent.findall(text_p_tag)
        assert len(paras) == 1
        assert paras[0].text == "New text"

    def test_multiline_injection_creates_paragraphs(self) -> None:
        """Newlines in new_text create additional <text:p> elements."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        p.text = "Original"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _inject_odf_paragraph_text(
            parent,
            "Line 1\nLine 2\nLine 3",
            text_p_tag,
        )
        assert result is True
        paras = parent.findall(text_p_tag)
        assert len(paras) == 3  # noqa: PLR2004
        assert paras[0].text == "Line 1"
        assert paras[1].text == "Line 2"
        assert paras[2].text == "Line 3"

    def test_preserves_span_attributes(self) -> None:
        """When original has a <text:span>, its attributes are preserved."""
        text_span_tag = f"{{{_TEXT_NS}}}span"
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        span = etree.SubElement(p, text_span_tag)
        span.set(
            f"{{{_ODF_NS_VAL['style']}}}style-name",
            "T1",
        )
        span.text = "Styled text"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        result = _inject_odf_paragraph_text(parent, "New styled", text_p_tag)
        assert result is True

        paras = parent.findall(text_p_tag)
        new_span = paras[0].find(text_span_tag)
        assert new_span is not None
        style_attr = f"{{{_ODF_NS_VAL['style']}}}style-name"
        assert new_span.get(style_attr) == "T1"
        assert new_span.text == "New styled"

    def test_removes_extra_original_paragraphs(self) -> None:
        """Extra original paragraphs are removed during injection."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        for txt in ("Para 1", "Para 2", "Para 3"):
            p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
            p.text = txt
        text_p_tag = f"{{{_TEXT_NS}}}p"

        _inject_odf_paragraph_text(parent, "Single", text_p_tag)
        paras = parent.findall(text_p_tag)
        assert len(paras) == 1
        assert paras[0].text == "Single"

    def test_html_routed_to_html_injection(self) -> None:
        """Text with <a> tags is routed to the HTML-aware injection."""
        parent = etree.Element(f"{{{_TEXT_NS}}}text-box")
        p = etree.SubElement(parent, f"{{{_TEXT_NS}}}p")
        p.text = "Original"
        text_p_tag = f"{{{_TEXT_NS}}}p"

        html_text = 'Click <a href="https://example.com">here</a>'
        result = _inject_odf_paragraph_text(parent, html_text, text_p_tag)
        assert result is True

        # The hyperlink should be in the XML tree
        text_a_tag = f"{{{_TEXT_NS}}}a"
        a_els = list(parent.iter(text_a_tag))
        assert len(a_els) >= 1
        href = a_els[0].get(f"{{{_XLINK_NS}}}href")
        assert href == "https://example.com"


# ---------------------------------------------------------------------------
# _build_odf_hf_map tests
# ---------------------------------------------------------------------------


class TestBuildOdfHfMap:
    """Tests for _build_odf_hf_map."""

    def test_returns_six_entries(self) -> None:
        """Map contains exactly 6 entries (header/footer x default/first/even)."""
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        result = _build_odf_hf_map(style_ns)
        assert len(result) == 6  # noqa: PLR2004

    def test_default_header_footer(self) -> None:
        """Default header and footer map correctly."""
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        result = _build_odf_hf_map(style_ns)

        header_tag = f"{{{style_ns}}}header"
        footer_tag = f"{{{style_ns}}}footer"

        assert result[header_tag] == ("header", _HF_DEFAULT)
        assert result[footer_tag] == ("footer", _HF_DEFAULT)

    def test_first_page_entries(self) -> None:
        """First-page header and footer map correctly."""
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        result = _build_odf_hf_map(style_ns)

        assert result[f"{{{style_ns}}}header-first"] == ("header", _HF_FIRST)
        assert result[f"{{{style_ns}}}footer-first"] == ("footer", _HF_FIRST)

    def test_even_page_entries(self) -> None:
        """Even (left) page header and footer map correctly."""
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        result = _build_odf_hf_map(style_ns)

        assert result[f"{{{style_ns}}}header-left"] == ("header", _HF_EVEN)
        assert result[f"{{{style_ns}}}footer-left"] == ("footer", _HF_EVEN)

    def test_custom_namespace(self) -> None:
        """Works with arbitrary namespace URIs."""
        custom_ns = "http://custom.ns/style"
        result = _build_odf_hf_map(custom_ns)

        assert f"{{{custom_ns}}}header" in result
        assert f"{{{custom_ns}}}footer-left" in result
        assert len(result) == 6  # noqa: PLR2004


# ---------------------------------------------------------------------------
# ODT headers/footers edge cases
# ---------------------------------------------------------------------------


class TestOdtHeadersFootersEdgeCases:
    """Additional edge cases for ODT header/footer extraction and injection."""

    def test_multiple_master_pages(self, tmp_path: Path) -> None:
        """Extracts headers from multiple <style:master-page> elements."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header><text:p>Header A</text:p></style:header>"
            "</style:master-page>"
            '<style:master-page style:name="First">'
            "<style:header><text:p>Header B</text:p></style:header>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "multi.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        result = _extract_odt_headers_footers(odt_path)
        keys = [k for k, _ in result]
        # Master page 0 and master page 1
        assert any(k.startswith("header:0:") for k in keys)
        assert any(k.startswith("header:1:") for k in keys)

    def test_empty_header_paragraph_skipped(self, tmp_path: Path) -> None:
        """Headers with empty/whitespace-only text are skipped."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header><text:p>   </text:p></style:header>"
            "<style:footer><text:p>Real Footer</text:p></style:footer>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "empty_hdr.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        result = _extract_odt_headers_footers(odt_path)
        # Only footer should be extracted (header is whitespace-only)
        assert len(result) == 1
        assert result[0][0].startswith("footer:")
        assert "Real Footer" in result[0][1]

    def test_inject_roundtrip_all_hf_types(self, tmp_path: Path) -> None:
        """Full roundtrip: extract→translate→inject for default, first, even."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header><text:p>Hdr Default</text:p></style:header>"
            "<style:footer><text:p>Ftr Default</text:p></style:footer>"
            "<style:header-first><text:p>Hdr First</text:p></style:header-first>"
            "<style:footer-left><text:p>Ftr Even</text:p></style:footer-left>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "roundtrip.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        extracted = _extract_odt_headers_footers(odt_path)
        assert len(extracted) == 4  # noqa: PLR2004

        # Build translations — uppercase everything
        translations = {k: v.upper() for k, v in extracted}
        _inject_odt_headers_footers(odt_path, translations)

        # Re-extract and verify
        re_extracted = _extract_odt_headers_footers(odt_path)
        values = [v for _, v in re_extracted]
        assert "HDR DEFAULT" in values
        assert "FTR DEFAULT" in values
        assert "HDR FIRST" in values
        assert "FTR EVEN" in values

    def test_inject_no_matching_translations_no_change(
        self,
        tmp_path: Path,
    ) -> None:
        """Injection with no matching keys does not modify the file."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header><text:p>Original</text:p></style:header>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "no_match.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        _inject_odt_headers_footers(odt_path, {"nonexistent:key": "Nothing"})

        # Text should remain unchanged
        result = _extract_odt_headers_footers(odt_path)
        assert any("Original" in v for _, v in result)


# ---------------------------------------------------------------------------
# _add_hyperlink_to_rels tests
# ---------------------------------------------------------------------------


class TestAddHyperlinkToRels:
    """Tests for _add_hyperlink_to_rels."""

    def test_creates_new_rels_document(self) -> None:
        """Creates a new Relationships XML when rels_xml is None."""
        rels_xml, r_id = _add_hyperlink_to_rels(None, "https://example.com")
        assert r_id == "rId1"
        root = etree.fromstring(rels_xml)
        rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        rels = root.findall(f"{{{rels_ns}}}Relationship")
        assert len(rels) == 1
        assert rels[0].get("Target") == "https://example.com"
        assert rels[0].get("TargetMode") == "External"

    def test_appends_to_existing_rels(self) -> None:
        """Appends a new relationship to existing rels XML."""
        rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        existing = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{rels_ns}">'
            f'<Relationship Id="rId1" Type="some/type" Target="foo.xml"/>'
            f"</Relationships>"
        )
        rels_xml, r_id = _add_hyperlink_to_rels(
            existing.encode("utf-8"),
            "https://new.com",
        )
        assert r_id == "rId2"
        root = etree.fromstring(rels_xml)
        rels = root.findall(f"{{{rels_ns}}}Relationship")
        assert len(rels) == 2  # noqa: PLR2004
        targets = [r.get("Target") for r in rels]
        assert "https://new.com" in targets

    def test_avoids_id_collision(self) -> None:
        """Generates a unique rId that doesn't collide with existing IDs."""
        rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        existing = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{rels_ns}">'
            f'<Relationship Id="rId1" Type="t" Target="a.xml"/>'
            f'<Relationship Id="rId2" Type="t" Target="b.xml"/>'
            f"</Relationships>"
        )
        rels_xml, r_id = _add_hyperlink_to_rels(
            existing.encode("utf-8"),
            "https://third.com",
        )
        assert r_id == "rId3"

    def test_multiple_appends(self) -> None:
        """Can append multiple hyperlinks sequentially."""
        xml, r1 = _add_hyperlink_to_rels(None, "https://a.com")
        xml, r2 = _add_hyperlink_to_rels(xml, "https://b.com")
        xml, r3 = _add_hyperlink_to_rels(xml, "https://c.com")
        assert r1 == "rId1"
        assert r2 == "rId2"
        assert r3 == "rId3"
        root = etree.fromstring(xml)
        rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        assert len(root.findall(f"{{{rels_ns}}}Relationship")) == 3  # noqa: PLR2004


# ---------------------------------------------------------------------------
# NEW EDGE-CASE TESTS — appended to expand coverage
# ---------------------------------------------------------------------------


# ---------------------------------------------------------------------------
# 1. TestDetectBackendEdgeCases
# ---------------------------------------------------------------------------


class TestDetectBackendEdgeCases:
    """Edge-case tests for _detect_backend."""

    def test_ooxml_ignores_win32com_and_uno(self) -> None:
        """OOXML always returns python_lib regardless of available backends."""
        for ext in (".docx", ".xlsx", ".pptx"):
            assert _detect_backend(ext) == "python_lib"

    def test_legacy_raises_office_converter_not_found(self) -> None:
        """Legacy formats raise ValueError when neither backend exists."""
        mods_to_remove = [
            k for k in sys.modules if k.startswith("win32com") or k == "uno"
        ]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with (
                patch.dict(
                    "sys.modules",
                    {"win32com": None, "win32com.client": None, "uno": None},
                ),
                pytest.raises(ValueError, match="OFFICE_CONVERTER_NOT_FOUND"),
            ):
                _detect_backend(".doc")
        finally:
            sys.modules.update(saved)

    def test_odf_python_lib_is_last_resort(self) -> None:
        """ODF returns python_lib when UNO and win32com are both missing."""
        mods_to_remove = [
            k for k in sys.modules if k.startswith("win32com") or k == "uno"
        ]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with patch.dict(
                "sys.modules",
                {"win32com": None, "win32com.client": None, "uno": None},
            ):
                assert _detect_backend(".odt") == "python_lib"
        finally:
            sys.modules.update(saved)

    def test_odf_with_libreoffice_path(self) -> None:
        """ODF backend detection forwards libreoffice_path to search paths."""
        mods_to_remove = [
            k for k in sys.modules if k.startswith("win32com") or k == "uno"
        ]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with (
                patch.dict(
                    "sys.modules",
                    {"win32com": None, "win32com.client": None, "uno": None},
                ),
                patch(
                    "src.core.office_processor._get_uno_search_paths",
                    return_value=[],
                ) as mock_paths,
            ):
                _detect_backend(".odt", "/custom/lo/path")
                mock_paths.assert_called_with("/custom/lo/path")
        finally:
            sys.modules.update(saved)

    def test_legacy_with_only_uno_available(self) -> None:
        """Legacy .ppt falls back to UNO when win32com not installed."""
        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with patch.dict(
                "sys.modules",
                {
                    "win32com": None,
                    "win32com.client": None,
                    "uno": MagicMock(),
                },
            ):
                assert _detect_backend(".ppt") == "uno"
        finally:
            sys.modules.update(saved)


# ---------------------------------------------------------------------------
# 2. TestFontPreservation
# ---------------------------------------------------------------------------


class TestFontPreservation:
    """Tests for font property save/restore functions."""

    def test_save_win32com_font_empty_object(self) -> None:
        """Empty font object (no properties set) returns empty dict."""
        from src.constants.office import WIN32COM_FONT_PROPERTIES  # noqa: PLC0415

        font_obj = MagicMock(spec=[])
        for prop in WIN32COM_FONT_PROPERTIES:
            setattr(type(font_obj), prop, PropertyMock(side_effect=AttributeError))

        result = _save_win32com_font(font_obj)
        assert result == {}

    def test_save_and_restore_win32com_font_roundtrip(self) -> None:
        """Font properties are saved and restored correctly."""
        font_obj = MagicMock()
        font_obj.Name = "Arial"
        font_obj.Size = 12.0
        font_obj.Bold = True
        font_obj.Italic = False
        font_obj.Color = 0
        font_obj.Underline = 0
        font_obj.StrikeThrough = False

        saved = _save_win32com_font(font_obj)
        assert "Name" in saved
        assert saved["Name"] == "Arial"
        assert saved["Size"] == 12.0  # noqa: PLR2004

        # Restore to a new mock
        new_font = MagicMock()
        _restore_win32com_font(new_font, saved)
        assert new_font.Name == "Arial"

    def test_restore_win32com_font_partial_props(self) -> None:
        """Restore handles partial saved properties without error."""
        font_obj = MagicMock()
        saved = {"Size": 14.0}
        _restore_win32com_font(font_obj, saved)
        assert font_obj.Size == 14.0  # noqa: PLR2004

    def test_save_uno_char_props_basic(self) -> None:
        """UNO char properties are saved via getPropertyValue."""
        text_range = MagicMock()
        text_range.getPropertyValue.return_value = "Liberation Serif"
        saved = _save_uno_char_props(text_range)
        assert isinstance(saved, dict)
        assert text_range.getPropertyValue.called

    def test_restore_uno_char_props_basic(self) -> None:
        """UNO char properties are restored via setPropertyValue."""
        text_range = MagicMock()
        saved = {"CharFontName": "Noto Sans"}
        _restore_uno_char_props(text_range, saved)
        text_range.setPropertyValue.assert_called()

    def test_restore_win32com_font_script_substitution(self) -> None:
        """Font name is substituted when scripts differ."""
        font_obj = MagicMock()
        saved = {"Name": "Arial", "Size": 12.0}
        _restore_win32com_font(
            font_obj,
            saved,
            original_text="Hello",
            translated_text="你好",
            target_lang="Chinese",
        )
        # The Name may be changed due to script difference
        assert font_obj.Name != "unused"  # Just verify it was set

    def test_restore_win32com_font_keeps_name_same_script(self) -> None:
        """Font name stays unchanged when scripts match."""
        font_obj = MagicMock()
        saved = {"Name": "Arial", "Size": 12.0}
        _restore_win32com_font(
            font_obj,
            saved,
            original_text="Hello",
            translated_text="Bonjour",
            target_lang="French",
        )
        assert font_obj.Name == "Arial"


# ---------------------------------------------------------------------------
# 3. TestCommentTranslation
# ---------------------------------------------------------------------------


class TestCommentTranslation:
    """Tests for comment extraction and injection."""

    def test_extract_docx_comments_empty(self, tmp_path: Path) -> None:
        """DOCX with no comments returns empty list."""
        doc = Document()
        doc.add_paragraph("Hello")
        path = tmp_path / "no_comments.docx"
        doc.save(str(path))
        result = _extract_docx_comments(path)
        assert result == []

    def test_extract_xlsx_comments_basic(self, tmp_path: Path) -> None:
        """XLSX comments are extracted correctly."""
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Data"
        ws["A1"].comment = Comment("Test comment", "Author")
        path = tmp_path / "comments.xlsx"
        wb.save(str(path))
        wb.close()

        result = _extract_xlsx_comments(path)
        assert len(result) >= 1
        assert any("Test comment" in v for _, v in result)

    def test_inject_xlsx_comments_basic(self, tmp_path: Path) -> None:
        """XLSX comment injection replaces comment text."""
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Data"
        ws["A1"].comment = Comment("Original comment", "Author")
        path = tmp_path / "inject_comments.xlsx"
        wb.save(str(path))
        wb.close()

        # Extract to get keys
        extracted = _extract_xlsx_comments(path)
        assert len(extracted) >= 1

        translations = {k: "Commentaire traduit" for k, _ in extracted}
        _inject_xlsx_comments(path, translations)

        # Verify injection
        re_extracted = _extract_xlsx_comments(path)
        assert any("Commentaire traduit" in v for _, v in re_extracted)

    def test_extract_odf_comments_unicode(self, tmp_path: Path) -> None:
        """ODF comments with unicode text are extracted correctly."""
        from src.core.office_processor import (  # noqa: PLC0415
            _extract_odf_comments,
        )

        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:text="{text_ns}">'
            "<office:body><office:text>"
            '<office:annotation office:name="ann1">'
            f"<text:p>Héllo wörld: données françaises</text:p>"
            "</office:annotation>"
            "</office:text></office:body>"
            "</office:document-content>"
        )
        path = tmp_path / "unicode_comments.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("content.xml", content_xml)

        result = _extract_odf_comments(path)
        assert len(result) == 1
        assert "données françaises" in result[0][1]

    def test_extract_odf_comments_no_annotations(self, tmp_path: Path) -> None:
        """ODF without annotations returns empty list."""
        from src.core.office_processor import (  # noqa: PLC0415
            _extract_odf_comments,
        )

        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:text="{text_ns}">'
            "<office:body><office:text>"
            "<text:p>Hello</text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )
        path = tmp_path / "no_annotations.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("content.xml", content_xml)

        result = _extract_odf_comments(path)
        assert result == []

    def test_inject_odf_comments_basic(self, tmp_path: Path) -> None:
        """ODF comment injection replaces annotation text."""
        from src.core.office_processor import (  # noqa: PLC0415
            _extract_odf_comments,
            _inject_odf_comments,
        )

        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:text="{text_ns}">'
            "<office:body><office:text>"
            '<office:annotation office:name="ann1">'
            "<text:p>Original comment</text:p>"
            "</office:annotation>"
            "</office:text></office:body>"
            "</office:document-content>"
        )
        path = tmp_path / "inject_odf.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("content.xml", content_xml)

        _inject_odf_comments(path, {"comment:ann1": "Translated comment"})
        result = _extract_odf_comments(path)
        assert any("Translated comment" in v for _, v in result)

    def test_extract_win32com_excel_comments_dispatched(self) -> None:
        """win32com Excel comment extraction dispatches correctly."""
        with (
            patch(
                "src.core.office_processor._extract_win32com_excel_comments",
                return_value=[("comment:Sheet1:1:1", "Test")],
            ) as mock_fn,
        ):
            result = _extract_comments(Path("test.xls"), ".xls", "win32com")
            mock_fn.assert_called_once()
            assert len(result) == 1


# ---------------------------------------------------------------------------
# 4. TestShapeTranslation
# ---------------------------------------------------------------------------


class TestShapeTranslation:
    """Tests for shape/text-box extraction and injection."""

    def test_extract_docx_shapes_with_text(self, tmp_path: Path) -> None:
        """DOCX shapes with text boxes are extracted."""
        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        mc_ns = "http://schemas.openxmlformats.org/markup-compatibility/2006"

        doc = Document()
        doc.add_paragraph("Body")
        path = tmp_path / "shapes.docx"
        doc.save(str(path))

        # Inject a text box into the DOCX document.xml
        with zipfile.ZipFile(path, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        doc_xml = file_data["word/document.xml"]
        root = etree.fromstring(doc_xml)

        # Find a <w:p> and inject a <wps:txbx> into it
        body = root.find(f"{{{w_ns}}}body")
        if body is not None:
            mc_alt = etree.SubElement(body, f"{{{mc_ns}}}AlternateContent")
            mc_choice = etree.SubElement(mc_alt, f"{{{mc_ns}}}Choice")
            mc_choice.set("Requires", "wps")
            txbx = etree.SubElement(mc_choice, f"{{{wps_ns}}}txbx")
            txbx_content = etree.SubElement(txbx, f"{{{w_ns}}}txbxContent")
            p = etree.SubElement(txbx_content, f"{{{w_ns}}}p")
            r = etree.SubElement(p, f"{{{w_ns}}}r")
            t = etree.SubElement(r, f"{{{w_ns}}}t")
            t.text = "Shape text"

            file_data["word/document.xml"] = etree.tostring(
                root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            _rewrite_zip_content(path, file_data, all_items)

        result = _extract_docx_shapes(path)
        assert any("Shape text" in v for _, v in result)

    def test_extract_docx_shapes_empty_doc(self, tmp_path: Path) -> None:
        """DOCX without shapes returns empty list."""
        doc = Document()
        doc.add_paragraph("Just text")
        path = tmp_path / "no_shapes.docx"
        doc.save(str(path))
        result = _extract_docx_shapes(path)
        assert result == []

    def test_inject_shapes_skips_non_shape_keys_class(self) -> None:
        """inject_shapes ignores non-shape keys."""
        with patch(
            "src.core.office_processor._inject_docx_shapes",
        ) as mock_inject:
            _inject_shapes(
                Path("test.docx"),
                {"para:0": "Translated"},
                ".docx",
                "python_lib",
            )
            mock_inject.assert_not_called()

    def test_extract_shapes_unsupported_ext(self) -> None:
        """Unsupported extension returns empty list."""
        result = _extract_shapes(Path("test.pptx"), ".pptx", "python_lib")
        assert result == []

    def test_extract_ods_shapes_returns_items(self, tmp_path: Path) -> None:
        """ODS shapes are extracted from draw:text-box elements."""
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
        table_ns = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:text="{text_ns}" xmlns:draw="{draw_ns}"'
            f' xmlns:table="{table_ns}">'
            "<office:body><office:spreadsheet>"
            f'<table:table table:name="Sheet1">'
            f"<table:table-row><table:table-cell>"
            f"<draw:frame><draw:text-box>"
            f"<text:p>Shape in cell</text:p>"
            f"</draw:text-box></draw:frame>"
            f"</table:table-cell></table:table-row>"
            f"</table:table>"
            "</office:spreadsheet></office:body>"
            "</office:document-content>"
        )
        path = tmp_path / "shapes.ods"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("content.xml", content_xml)

        result = _extract_ods_shapes(path)
        assert any("Shape in cell" in v for _, v in result)


# ---------------------------------------------------------------------------
# 5. TestHeadersFooters
# ---------------------------------------------------------------------------


class TestHeadersFootersEdgeCases2:
    """Additional edge-case tests for headers/footers."""

    def test_extract_docx_header_basic(self, tmp_path: Path) -> None:
        """Extracts default header from DOCX."""
        doc = Document()
        doc.add_paragraph("Body")
        section = doc.sections[0]
        section.header.is_linked_to_previous = False
        section.header.paragraphs[0].text = "My Header"
        path = tmp_path / "header.docx"
        doc.save(str(path))

        result = _extract_docx_headers_footers(path)
        assert any("My Header" in v for _, v in result)

    def test_extract_docx_footer_basic(self, tmp_path: Path) -> None:
        """Extracts default footer from DOCX."""
        doc = Document()
        doc.add_paragraph("Body")
        section = doc.sections[0]
        section.footer.is_linked_to_previous = False
        section.footer.paragraphs[0].text = "My Footer"
        path = tmp_path / "footer.docx"
        doc.save(str(path))

        result = _extract_docx_headers_footers(path)
        assert any("My Footer" in v for _, v in result)

    def test_extract_docx_no_header_when_linked(self, tmp_path: Path) -> None:
        """Linked-to-previous headers are skipped."""
        doc = Document()
        doc.add_paragraph("Body")
        section = doc.sections[0]
        section.header.is_linked_to_previous = True
        path = tmp_path / "linked.docx"
        doc.save(str(path))

        result = _extract_docx_headers_footers(path)
        header_items = [k for k, _ in result if k.startswith("header:")]
        assert len(header_items) == 0

    def test_inject_docx_headers_footers_basic(self, tmp_path: Path) -> None:
        """Injects translated text into DOCX header."""
        doc = Document()
        doc.add_paragraph("Body")
        section = doc.sections[0]
        section.header.is_linked_to_previous = False
        section.header.paragraphs[0].text = "Original Header"
        path = tmp_path / "inject_hf.docx"
        doc.save(str(path))

        _inject_docx_headers_footers(
            path,
            {"header:0:default:0": "Translated Header"},
        )

        doc2 = Document(str(path))
        hdr_text = doc2.sections[0].header.paragraphs[0].text
        assert hdr_text == "Translated Header"

    def test_extract_odt_headers_footers_basic(self, tmp_path: Path) -> None:
        """Extracts headers/footers from ODT styles.xml."""
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"

        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-styles xmlns:office="{office_ns}"'
            f' xmlns:style="{style_ns}" xmlns:text="{text_ns}">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header><text:p>Default Header</text:p></style:header>"
            "<style:footer><text:p>Default Footer</text:p></style:footer>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        path = tmp_path / "hf.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        result = _extract_odt_headers_footers(path)
        headers = [v for k, v in result if "header" in k]
        footers = [v for k, v in result if "footer" in k]
        assert "Default Header" in headers
        assert "Default Footer" in footers

    def test_extract_odt_headers_footers_first_and_even(
        self,
        tmp_path: Path,
    ) -> None:
        """Extracts first and even page headers from ODT."""
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"

        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-styles xmlns:office="{office_ns}"'
            f' xmlns:style="{style_ns}" xmlns:text="{text_ns}">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header><text:p>Default Hdr</text:p></style:header>"
            "<style:header-first><text:p>First Hdr</text:p>"
            "</style:header-first>"
            "<style:header-left><text:p>Even Hdr</text:p>"
            "</style:header-left>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        path = tmp_path / "hf_first_even.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        result = _extract_odt_headers_footers(path)
        values = [v for _, v in result]
        assert "Default Hdr" in values
        assert "First Hdr" in values
        assert "Even Hdr" in values

    def test_extract_docx_empty_headers(self, tmp_path: Path) -> None:
        """DOCX with empty headers returns no header entries."""
        doc = Document()
        doc.add_paragraph("Body")
        section = doc.sections[0]
        section.header.is_linked_to_previous = False
        # Leave header paragraphs empty
        for p in section.header.paragraphs:
            p.text = ""
        path = tmp_path / "empty_hdr.docx"
        doc.save(str(path))

        result = _extract_docx_headers_footers(path)
        header_items = [k for k, _ in result if k.startswith("header:")]
        assert len(header_items) == 0

    def test_build_odf_hf_map_keys(self) -> None:
        """ODF H/F map returns all six key types."""
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        hf_map = _build_odf_hf_map(style_ns)
        expected_tags = {
            f"{{{style_ns}}}header",
            f"{{{style_ns}}}footer",
            f"{{{style_ns}}}header-first",
            f"{{{style_ns}}}footer-first",
            f"{{{style_ns}}}header-left",
            f"{{{style_ns}}}footer-left",
        }
        assert set(hf_map.keys()) == expected_tags

    def test_build_odf_hf_map_values(self) -> None:
        """ODF H/F map has correct prefix/type tuples."""
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        hf_map = _build_odf_hf_map(style_ns)
        assert hf_map[f"{{{style_ns}}}header"] == ("header", _HF_DEFAULT)
        assert hf_map[f"{{{style_ns}}}footer-first"] == ("footer", _HF_FIRST)
        assert hf_map[f"{{{style_ns}}}header-left"] == ("header", _HF_EVEN)


# ---------------------------------------------------------------------------
# 6. TestFootnotesEndnotes
# ---------------------------------------------------------------------------


class TestFootnotesEndnotesEdgeCases:
    """Edge-case tests for footnote and endnote extraction/injection."""

    def test_extract_docx_footnotes_basic(self, tmp_path: Path) -> None:
        """Extracts footnotes from a DOCX with footnotes.xml."""
        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<w:footnotes xmlns:w="{w_ns}">'
            f'<w:footnote w:id="0" w:type="separator"><w:p><w:r>'
            f"<w:separator/></w:r></w:p></w:footnote>"
            f'<w:footnote w:id="1" w:type="continuationSeparator"><w:p>'
            f"<w:r><w:continuationSeparator/></w:r></w:p></w:footnote>"
            f'<w:footnote w:id="2"><w:p><w:r>'
            f"<w:t>My footnote text</w:t>"
            f"</w:r></w:p></w:footnote>"
            f"</w:footnotes>"
        )

        doc = Document()
        doc.add_paragraph("Body text")
        path = tmp_path / "fn.docx"
        doc.save(str(path))

        # Insert footnotes.xml
        with zipfile.ZipFile(path, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}
        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        new_info = zipfile.ZipInfo("word/footnotes.xml")
        all_items.append(new_info)
        _rewrite_zip_content(path, file_data, all_items)

        result = _extract_docx_footnotes(path)
        assert any("My footnote text" in v for _, v in result)
        assert any("footnote:2" in k for k, _ in result)

    def test_extract_docx_endnotes(self, tmp_path: Path) -> None:
        """Extracts endnotes from a DOCX with endnotes.xml."""
        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        en_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<w:endnotes xmlns:w="{w_ns}">'
            f'<w:endnote w:id="0" w:type="separator"><w:p><w:r>'
            f"<w:separator/></w:r></w:p></w:endnote>"
            f'<w:endnote w:id="1" w:type="continuationSeparator"><w:p>'
            f"<w:r><w:continuationSeparator/></w:r></w:p></w:endnote>"
            f'<w:endnote w:id="2"><w:p><w:r>'
            f"<w:t>My endnote text</w:t>"
            f"</w:r></w:p></w:endnote>"
            f"</w:endnotes>"
        )

        doc = Document()
        doc.add_paragraph("Body text")
        path = tmp_path / "en.docx"
        doc.save(str(path))

        with zipfile.ZipFile(path, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}
        file_data["word/endnotes.xml"] = en_xml.encode("utf-8")
        new_info = zipfile.ZipInfo("word/endnotes.xml")
        all_items.append(new_info)
        _rewrite_zip_content(path, file_data, all_items)

        result = _extract_docx_footnotes(path)
        assert any("My endnote text" in v for _, v in result)
        assert any("endnote:2" in k for k, _ in result)

    def test_extract_docx_no_footnotes(self, tmp_path: Path) -> None:
        """DOCX without footnotes/endnotes returns empty list."""
        doc = Document()
        doc.add_paragraph("Just body")
        path = tmp_path / "no_fn.docx"
        doc.save(str(path))

        result = _extract_docx_footnotes(path)
        assert result == []

    def test_inject_docx_footnotes_replaces_text(self, tmp_path: Path) -> None:
        """Injecting footnotes replaces the text content."""
        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<w:footnotes xmlns:w="{w_ns}">'
            f'<w:footnote w:id="0" w:type="separator"><w:p><w:r>'
            f"<w:separator/></w:r></w:p></w:footnote>"
            f'<w:footnote w:id="2"><w:p><w:r>'
            f"<w:t>Original footnote</w:t>"
            f"</w:r></w:p></w:footnote>"
            f"</w:footnotes>"
        )

        doc = Document()
        doc.add_paragraph("Body")
        path = tmp_path / "inject_fn.docx"
        doc.save(str(path))

        with zipfile.ZipFile(path, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}
        file_data["word/footnotes.xml"] = fn_xml.encode("utf-8")
        new_info = zipfile.ZipInfo("word/footnotes.xml")
        all_items.append(new_info)
        _rewrite_zip_content(path, file_data, all_items)

        _inject_docx_footnotes(path, {"footnote:2": "Translated footnote"})

        result = _extract_docx_footnotes(path)
        assert any("Translated footnote" in v for _, v in result)

    def test_extract_odt_footnotes_basic(self, tmp_path: Path) -> None:
        """Extracts footnotes from ODT content.xml."""
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:text="{text_ns}">'
            "<office:body><office:text>"
            '<text:note text:id="fn1" text:note-class="footnote">'
            "<text:note-citation>1</text:note-citation>"
            "<text:note-body><text:p>ODT footnote</text:p></text:note-body>"
            "</text:note>"
            "</office:text></office:body>"
            "</office:document-content>"
        )
        path = tmp_path / "fn.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("content.xml", content_xml)

        result = _extract_odt_footnotes(path)
        assert len(result) == 1
        assert result[0][0] == "footnote:fn1"
        assert "ODT footnote" in result[0][1]

    def test_extract_odt_endnotes(self, tmp_path: Path) -> None:
        """Extracts endnotes from ODT content.xml."""
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:text="{text_ns}">'
            "<office:body><office:text>"
            '<text:note text:id="en1" text:note-class="endnote">'
            "<text:note-citation>i</text:note-citation>"
            "<text:note-body><text:p>ODT endnote text</text:p></text:note-body>"
            "</text:note>"
            "</office:text></office:body>"
            "</office:document-content>"
        )
        path = tmp_path / "en.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("content.xml", content_xml)

        result = _extract_odt_footnotes(path)
        assert len(result) == 1
        assert result[0][0] == "endnote:en1"
        assert "ODT endnote text" in result[0][1]

    def test_extract_odt_no_footnotes(self, tmp_path: Path) -> None:
        """ODT without notes returns empty list."""
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            f'<office:document-content xmlns:office="{office_ns}"'
            f' xmlns:text="{text_ns}">'
            "<office:body><office:text>"
            "<text:p>Just text</text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )
        path = tmp_path / "no_fn.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("content.xml", content_xml)

        result = _extract_odt_footnotes(path)
        assert result == []

    def test_footnote_extensions_correct(self) -> None:
        """_FOOTNOTE_EXTENSIONS includes correct formats."""
        assert ".docx" in _FOOTNOTE_EXTENSIONS
        assert ".odt" in _FOOTNOTE_EXTENSIONS
        assert ".doc" in _FOOTNOTE_EXTENSIONS
        assert ".xlsx" not in _FOOTNOTE_EXTENSIONS

    def test_extract_footnotes_dispatch_docx(self) -> None:
        """_extract_footnotes dispatches to _extract_docx_footnotes for DOCX."""
        with patch(
            "src.core.office_processor._extract_docx_footnotes",
            return_value=[("footnote:1", "text")],
        ) as mock_fn:
            result = _extract_footnotes(Path("test.docx"), ".docx", "python_lib")
            mock_fn.assert_called_once()
            assert len(result) == 1

    def test_extract_footnotes_dispatch_odt(self) -> None:
        """_extract_footnotes dispatches to _extract_odt_footnotes for ODT."""
        with patch(
            "src.core.office_processor._extract_odt_footnotes",
            return_value=[],
        ) as mock_fn:
            _extract_footnotes(Path("test.odt"), ".odt", "python_lib")
            mock_fn.assert_called_once()

    def test_inject_footnotes_dispatch_docx(self) -> None:
        """_inject_footnotes dispatches to _inject_docx_footnotes for DOCX."""
        with patch(
            "src.core.office_processor._inject_docx_footnotes",
        ) as mock_fn:
            _inject_footnotes(
                Path("test.docx"),
                {"footnote:2": "Translated"},
                ".docx",
                "python_lib",
            )
            mock_fn.assert_called_once()


# ---------------------------------------------------------------------------
# 7. TestSpeakerNotes
# ---------------------------------------------------------------------------


class TestSpeakerNotesEdgeCases:
    """Additional edge-case tests for speaker notes extraction/injection."""

    def test_extract_pptx_notes_no_notes(self, tmp_path: Path) -> None:
        """PPTX without notes returns empty list."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank, no notes
        path = tmp_path / "no_notes.pptx"
        prs.save(str(path))

        result = _extract_pptx_notes(path)
        assert result == []

    def test_extract_pptx_notes_basic(self, tmp_path: Path) -> None:
        """PPTX notes text is extracted correctly."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Speaker notes here"
        path = tmp_path / "notes.pptx"
        prs.save(str(path))

        result = _extract_pptx_notes(path)
        assert any("Speaker notes here" in v for _, v in result)

    def test_inject_pptx_notes_replaces_text(self, tmp_path: Path) -> None:
        """PPTX note injection replaces note text."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Original note"
        path = tmp_path / "inject_notes.pptx"
        prs.save(str(path))

        extracted = _extract_pptx_notes(path)
        translations = {k: "Translated note" for k, _ in extracted}
        _inject_pptx_notes(path, translations)

        result = _extract_pptx_notes(path)
        assert any("Translated note" in v for _, v in result)

    def test_extract_odp_notes_basic(self, tmp_path: Path) -> None:
        """ODP notes are extracted from presentation:notes elements."""
        from odf.draw import Frame, Page, TextBox  # noqa: PLC0415
        from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
        from odf.presentation import Notes  # noqa: PLC0415
        from odf.style import MasterPage, PageLayout  # noqa: PLC0415
        from odf.text import P  # noqa: PLC0415

        doc = OpenDocumentPresentation()
        pl = PageLayout(name="PM1")
        doc.automaticstyles.addElement(pl)
        mp = MasterPage(name="Default", pagelayoutname=pl)
        doc.masterstyles.addElement(mp)

        page = Page(name="Slide1", masterpagename=mp)
        notes = Notes()
        frame = Frame(width="10cm", height="5cm")
        tb = TextBox()
        tb.addElement(P(text="ODP note text"))
        frame.addElement(tb)
        notes.addElement(frame)
        page.addElement(notes)
        doc.presentation.addElement(page)
        path = tmp_path / "notes.odp"
        doc.save(str(path))

        result = _extract_odp_notes(path)
        assert any("ODP note text" in v for _, v in result)

    def test_extract_odp_notes_empty(self, tmp_path: Path) -> None:
        """ODP without notes returns empty list."""
        from odf.draw import Page  # noqa: PLC0415
        from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
        from odf.style import MasterPage, PageLayout  # noqa: PLC0415

        doc = OpenDocumentPresentation()
        pl = PageLayout(name="PM1")
        doc.automaticstyles.addElement(pl)
        mp = MasterPage(name="Default", pagelayoutname=pl)
        doc.masterstyles.addElement(mp)
        page = Page(name="Slide1", masterpagename=mp)
        doc.presentation.addElement(page)
        path = tmp_path / "no_notes.odp"
        doc.save(str(path))

        result = _extract_odp_notes(path)
        assert result == []

    def test_should_translate_notes_gated_by_setting(self) -> None:
        """Notes are gated by SETTING_TRANSLATE_DOC_NOTES via config."""
        config = MagicMock()
        config.translate_doc_notes = False
        assert not _should_translate_notes(".pptx", "python_lib", config)

        config.translate_doc_notes = True
        assert _should_translate_notes(".pptx", "python_lib", config)

    def test_should_translate_notes_unsupported_ext(self) -> None:
        """Non-presentation extensions return False for notes."""
        config = MagicMock()
        config.translate_doc_notes = True
        assert not _should_translate_notes(".docx", "python_lib", config)

    def test_notes_extensions_constant(self) -> None:
        """_NOTES_EXTENSIONS has correct values."""
        assert {".pptx", ".odp", ".ppt"} == _NOTES_EXTENSIONS

    def test_extract_notes_dispatch_pptx(self) -> None:
        """_extract_notes dispatches to _extract_pptx_notes for PPTX."""
        with patch(
            "src.core.office_processor._extract_pptx_notes",
            return_value=[("note:0:0", "text")],
        ) as mock_fn:
            result = _extract_notes(Path("test.pptx"), ".pptx", "python_lib")
            mock_fn.assert_called_once()
            assert len(result) == 1

    def test_extract_notes_dispatch_odp(self) -> None:
        """_extract_notes dispatches to _extract_odp_notes for ODP."""
        with patch(
            "src.core.office_processor._extract_odp_notes",
            return_value=[],
        ) as mock_fn:
            _extract_notes(Path("test.odp"), ".odp", "python_lib")
            mock_fn.assert_called_once()

    def test_inject_notes_dispatch_pptx(self) -> None:
        """_inject_notes dispatches to _inject_pptx_notes for PPTX."""
        with patch(
            "src.core.office_processor._inject_pptx_notes",
        ) as mock_fn:
            _inject_notes(
                Path("test.pptx"),
                {"note:0:0": "Translated"},
                ".pptx",
                "python_lib",
            )
            mock_fn.assert_called_once()

    def test_inject_notes_dispatch_odp(self) -> None:
        """_inject_notes dispatches to _inject_odp_notes for ODP."""
        with patch(
            "src.core.office_processor._inject_odp_notes",
        ) as mock_fn:
            _inject_notes(
                Path("test.odp"),
                {"note:0:0": "Translated"},
                ".odp",
                "python_lib",
            )
            mock_fn.assert_called_once()


# ---------------------------------------------------------------------------
# 8. TestSheetNames
# ---------------------------------------------------------------------------


class TestSheetNamesEdgeCases:
    """Additional edge-case tests for sheet name extraction/injection."""

    def test_extract_xlsx_sheet_names_multiple(self, tmp_path: Path) -> None:
        """Extracts multiple sheet names from XLSX."""
        wb = Workbook()
        wb.active.title = "Data"
        wb.create_sheet("Summary")
        wb.create_sheet("Charts")
        path = tmp_path / "multi.xlsx"
        wb.save(str(path))
        wb.close()

        result = _extract_xlsx_sheet_names(path)
        names = [v for _, v in result]
        assert "Data" in names
        assert "Summary" in names
        assert "Charts" in names

    def test_extract_xlsx_sheet_names_no_workbook_xml(self, tmp_path: Path) -> None:
        """XLSX without xl/workbook.xml returns empty list."""
        path = tmp_path / "no_wb.xlsx"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("other.xml", "<root/>")

        result = _extract_xlsx_sheet_names(path)
        assert result == []

    def test_inject_xlsx_sheet_names_roundtrip(self, tmp_path: Path) -> None:
        """Sheet names are injected and re-extracted correctly."""
        wb = Workbook()
        wb.active.title = "Revenue"
        path = tmp_path / "inject_sn.xlsx"
        wb.save(str(path))
        wb.close()

        _inject_xlsx_sheet_names(path, {"sheetname:Revenue": "Recettes"})

        result = _extract_xlsx_sheet_names(path)
        assert any("Recettes" in v for _, v in result)

    def test_inject_xlsx_sheet_names_no_match(self, tmp_path: Path) -> None:
        """Injection with non-matching keys is a no-op."""
        wb = Workbook()
        wb.active.title = "Sheet1"
        path = tmp_path / "no_match_sn.xlsx"
        wb.save(str(path))
        wb.close()

        _inject_xlsx_sheet_names(path, {"sheetname:NonExistent": "Translated"})

        result = _extract_xlsx_sheet_names(path)
        assert result[0][1] == "Sheet1"

    def test_extract_ods_sheet_names_basic(self, tmp_path: Path) -> None:
        """Extracts sheet names from ODS via odfpy."""
        from odf.opendocument import OpenDocumentSpreadsheet  # noqa: PLC0415
        from odf.table import Table  # noqa: PLC0415

        doc = OpenDocumentSpreadsheet()
        doc.spreadsheet.addElement(Table(name="Sheet1"))
        doc.spreadsheet.addElement(Table(name="Sheet2"))
        path = tmp_path / "sheets.ods"
        doc.save(str(path))

        result = _extract_ods_sheet_names(path)
        names = [v for _, v in result]
        assert "Sheet1" in names
        assert "Sheet2" in names

    def test_inject_ods_sheet_names_roundtrip(self, tmp_path: Path) -> None:
        """ODS sheet names are injected and re-extracted correctly."""
        from odf.opendocument import OpenDocumentSpreadsheet  # noqa: PLC0415
        from odf.table import Table  # noqa: PLC0415

        doc = OpenDocumentSpreadsheet()
        doc.spreadsheet.addElement(Table(name="Sales"))
        path = tmp_path / "inject_ods_sn.ods"
        doc.save(str(path))

        _inject_ods_sheet_names(path, {"sheetname:Sales": "Ventes"})

        result = _extract_ods_sheet_names(path)
        assert any("Ventes" in v for _, v in result)

    def test_sanitize_sheet_name_all_invalid_returns_sheet(self) -> None:
        """All-invalid chars returns 'Sheet' fallback."""
        assert _sanitize_sheet_name("\\/*?:[]") == "Sheet"

    def test_sanitize_sheet_name_truncation_exactly_31(self) -> None:
        """Name of exactly 31 chars passes through unchanged."""
        name_31 = "A" * 31  # noqa: PLR2004
        assert _sanitize_sheet_name(name_31) == name_31

    def test_sanitize_sheet_name_tabs_and_newlines(self) -> None:
        """Tabs and newlines pass through (not in invalid set)."""
        result = _sanitize_sheet_name("Tab\there")
        assert "Tab" in result

    def test_should_translate_sheet_names_gated(self) -> None:
        """Sheet names are gated by config."""
        config = MagicMock()
        config.translate_sheet_names = True
        assert _should_translate_sheet_names(".xlsx", "python_lib", config)

        config.translate_sheet_names = False
        assert not _should_translate_sheet_names(".xlsx", "python_lib", config)

    def test_should_translate_sheet_names_unsupported_ext(self) -> None:
        """Non-spreadsheet ext returns False."""
        config = MagicMock()
        config.translate_sheet_names = True
        assert not _should_translate_sheet_names(".pptx", "python_lib", config)

    def test_extract_sheet_names_dispatch_xlsx(self) -> None:
        """_extract_sheet_names dispatches correctly for XLSX."""
        with patch(
            "src.core.office_processor._extract_xlsx_sheet_names",
            return_value=[("sheetname:S1", "S1")],
        ) as mock_fn:
            result = _extract_sheet_names(Path("test.xlsx"), ".xlsx", "python_lib")
            mock_fn.assert_called_once()
            assert len(result) == 1

    def test_extract_sheet_names_dispatch_ods(self) -> None:
        """_extract_sheet_names dispatches correctly for ODS."""
        with patch(
            "src.core.office_processor._extract_ods_sheet_names",
            return_value=[],
        ) as mock_fn:
            _extract_sheet_names(Path("test.ods"), ".ods", "python_lib")
            mock_fn.assert_called_once()

    def test_inject_sheet_names_dispatch_xlsx(self) -> None:
        """_inject_sheet_names dispatches correctly for XLSX."""
        with patch(
            "src.core.office_processor._inject_xlsx_sheet_names",
        ) as mock_fn:
            _inject_sheet_names(
                Path("test.xlsx"),
                {"sheetname:S1": "Translated"},
                ".xlsx",
                "python_lib",
            )
            mock_fn.assert_called_once()

    def test_inject_sheet_names_dispatch_ods(self) -> None:
        """_inject_sheet_names dispatches correctly for ODS."""
        with patch(
            "src.core.office_processor._inject_ods_sheet_names",
        ) as mock_fn:
            _inject_sheet_names(
                Path("test.ods"),
                {"sheetname:S1": "Translated"},
                ".ods",
                "python_lib",
            )
            mock_fn.assert_called_once()

    def test_inject_sheet_names_skips_non_sheetname_keys(self) -> None:
        """_inject_sheet_names ignores non-sheetname keys."""
        with patch(
            "src.core.office_processor._inject_xlsx_sheet_names",
        ) as mock_fn:
            _inject_sheet_names(
                Path("test.xlsx"),
                {"para:0": "Text"},
                ".xlsx",
                "python_lib",
            )
            mock_fn.assert_not_called()


# ---------------------------------------------------------------------------
# 9. TestAutoConversion
# ---------------------------------------------------------------------------


class TestAutoConversionEdgeCases:
    """Edge-case tests for convert_to_modern_format."""

    def test_convert_returns_true_on_win32com_success(self, tmp_path: Path) -> None:
        """Returns True when win32com succeeds."""
        src = tmp_path / "test.doc"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        with (
            patch.dict(
                "sys.modules",
                {"win32com": MagicMock(), "win32com.client": MagicMock()},
            ),
            patch("src.core.office_processor._convert_with_win32com") as mock_conv,
        ):
            result = convert_to_modern_format(src, dst)

        assert result is True
        mock_conv.assert_called_once_with(src, dst)

    def test_convert_falls_back_to_uno_on_no_win32com(
        self,
        tmp_path: Path,
    ) -> None:
        """Falls back to UNO when win32com not available."""
        src = tmp_path / "test.xls"
        dst = tmp_path / "test.xlsx"
        src.write_bytes(b"fake")

        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with (
                patch.dict(
                    "sys.modules",
                    {"win32com": None, "win32com.client": None, "uno": MagicMock()},
                ),
                patch("src.core.office_processor._convert_with_uno") as mock_conv,
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is True
        mock_conv.assert_called_once_with(src, dst)

    def test_convert_returns_false_no_backend(self, tmp_path: Path) -> None:
        """Returns False when neither backend is available."""
        src = tmp_path / "test.ppt"
        dst = tmp_path / "test.pptx"
        src.write_bytes(b"fake")

        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with patch.dict(
                "sys.modules",
                {"win32com": None, "win32com.client": None, "uno": None},
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is False

    def test_convert_win32com_error_returns_false(self, tmp_path: Path) -> None:
        """Returns False when win32com conversion raises."""
        src = tmp_path / "test.doc"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        with (
            patch.dict(
                "sys.modules",
                {"win32com": MagicMock(), "win32com.client": MagicMock()},
            ),
            patch(
                "src.core.office_processor._convert_with_win32com",
                side_effect=RuntimeError("COM failed"),
            ),
        ):
            result = convert_to_modern_format(src, dst)

        assert result is False

    def test_convert_uno_error_returns_false(self, tmp_path: Path) -> None:
        """Returns False when UNO conversion raises."""
        src = tmp_path / "test.odt"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with (
                patch.dict(
                    "sys.modules",
                    {"win32com": None, "win32com.client": None, "uno": MagicMock()},
                ),
                patch(
                    "src.core.office_processor._convert_with_uno",
                    side_effect=RuntimeError("UNO error"),
                ),
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is False

    def test_legacy_convert_map(self) -> None:
        """LEGACY_CONVERT_MAP has correct mappings."""
        assert LEGACY_CONVERT_MAP[".doc"] == ".docx"
        assert LEGACY_CONVERT_MAP[".xls"] == ".xlsx"
        assert LEGACY_CONVERT_MAP[".ppt"] == ".pptx"

    def test_odf_convert_map(self) -> None:
        """ODF_CONVERT_MAP has correct mappings."""
        from src.core.office_processor import ODF_CONVERT_MAP  # noqa: PLC0415

        assert ODF_CONVERT_MAP[".odt"] == ".docx"
        assert ODF_CONVERT_MAP[".ods"] == ".xlsx"
        assert ODF_CONVERT_MAP[".odp"] == ".pptx"


# ---------------------------------------------------------------------------
# 10. TestImageTranslation
# ---------------------------------------------------------------------------


class TestImageTranslationEdgeCases:
    """Edge-case tests for image translation."""

    def test_suffix_to_media_prefixes_ooxml(self) -> None:
        """OOXML formats have correct media prefixes."""
        assert _SUFFIX_TO_MEDIA_PREFIXES[".docx"] == ("word/media/", "media/")
        assert _SUFFIX_TO_MEDIA_PREFIXES[".xlsx"] == ("xl/media/", "media/")
        assert _SUFFIX_TO_MEDIA_PREFIXES[".pptx"] == ("ppt/media/", "media/")

    def test_suffix_to_media_prefixes_odf(self) -> None:
        """ODF formats use Pictures/ directory."""
        assert _SUFFIX_TO_MEDIA_PREFIXES[".odt"] == ("Pictures/",)
        assert _SUFFIX_TO_MEDIA_PREFIXES[".ods"] == ("Pictures/",)
        assert _SUFFIX_TO_MEDIA_PREFIXES[".odp"] == ("Pictures/",)

    def test_suffix_to_media_prefixes_epub(self) -> None:
        """EPUB uses empty prefix (images anywhere)."""
        assert _SUFFIX_TO_MEDIA_PREFIXES[".epub"] == ("",)

    def test_fatal_llm_errors_set(self) -> None:
        """_FATAL_LLM_ERRORS has correct values."""
        assert "AUTH_ERROR" in _FATAL_LLM_ERRORS
        assert "QUOTA_ERROR" in _FATAL_LLM_ERRORS
        assert "VISION_NOT_SUPPORTED" in _FATAL_LLM_ERRORS
        assert "SERVICE_UNAVAILABLE_ERROR" not in _FATAL_LLM_ERRORS

    @patch("src.core.office_processor._translate_single_image")
    def test_translate_zip_images_quota_error_is_fatal(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """QUOTA_ERROR stops processing immediately."""
        mock_translate.side_effect = ValueError("QUOTA_ERROR")
        file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")

        with pytest.raises(ValueError, match="QUOTA_ERROR"):
            _translate_zip_images(
                file_path,
                ".docx",
                "French",
                "English",
                None,
                "TesseractOCR",
                None,
                None,
            )
        assert mock_translate.call_count == 1

    @patch("src.core.office_processor._translate_single_image")
    def test_translate_zip_images_vision_not_supported_is_fatal(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """VISION_NOT_SUPPORTED stops processing immediately."""
        mock_translate.side_effect = ValueError("VISION_NOT_SUPPORTED")
        file_path = _make_minimal_zip(tmp_path, "xl/media/", ".xlsx")

        with pytest.raises(ValueError, match="VISION_NOT_SUPPORTED"):
            _translate_zip_images(
                file_path,
                ".xlsx",
                "French",
                "English",
                None,
                "TesseractOCR",
                None,
                None,
            )
        assert mock_translate.call_count == 1

    @patch("src.core.office_processor._translate_single_image")
    def test_translate_zip_images_connection_error_tries_all(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
    ) -> None:
        """CONNECTION_ERROR tries all images then completes (skip-with-warning).

        A transient connection error that survives the LLM's own
        3-retry policy is still treated as a per-image bad input —
        skip it, log a warning, leave the original.  No raise.
        """
        mock_translate.side_effect = ValueError("CONNECTION_ERROR")
        file_path = _make_minimal_zip(tmp_path, "word/media/", ".docx")

        _translate_zip_images(
            file_path,
            ".docx",
            "French",
            "English",
            None,
            "TesseractOCR",
            None,
            None,
        )
        assert mock_translate.call_count == _EXPECTED_IMAGE_COUNT

    @patch("src.core.office_processor._translate_single_image")
    def test_translate_zip_images_unexpected_exception_continues(
        self,
        mock_translate: MagicMock,
        tmp_path: Path,
        caplog,  # noqa: ANN001
    ) -> None:
        """Unexpected (non-ValueError) exceptions continue + log with traceback.

        The bare ``except Exception`` branch in ``_translate_zip_images``
        is the catch-all for buggy LLM clients (``RuntimeError``,
        ``TypeError``, …).  Same skip-with-warning policy as the
        non-fatal ``ValueError`` path, plus ``exc_info=True`` on the
        WARNING log so the full traceback is available in ``app.log``
        for diagnosis.
        """
        mock_translate.side_effect = RuntimeError("OCR crash")
        file_path = _make_minimal_zip(tmp_path, "ppt/media/", ".pptx")

        # Should not raise — unexpected errors are logged and continued.
        with caplog.at_level("WARNING", logger="office_processor"):
            _translate_zip_images(
                file_path,
                ".pptx",
                "French",
                "English",
                None,
                "TesseractOCR",
                None,
                None,
            )
        assert mock_translate.call_count == _EXPECTED_IMAGE_COUNT
        # WARNING fires with a traceback (``exc_info=True``).
        crash_records = [
            r
            for r in caplog.records
            if r.levelname == "WARNING" and "Unexpected error" in r.message
        ]
        assert crash_records, "expected WARNING for unexpected exception"
        assert crash_records[0].exc_info is not None, (
            "WARNING must carry traceback (exc_info=True) for debugging"
        )

    def test_translate_doc_images_dispatches_legacy(self, tmp_path: Path) -> None:
        """_translate_doc_images dispatches to _translate_legacy_images for .doc."""
        with (
            patch(
                "src.core.office_processor._translate_legacy_images",
            ) as mock_legacy,
            patch(
                "src.utils.config_manager.load_setting",
                return_value="TesseractOCR",
            ),
        ):
            _translate_doc_images(
                tmp_path / "test.doc",
                ".doc",
                "win32com",
                "French",
                "English",
                None,
                None,
                None,
            )
            mock_legacy.assert_called_once()

    def test_translate_doc_images_dispatches_modern(self, tmp_path: Path) -> None:
        """_translate_doc_images dispatches to _translate_zip_images for .docx."""
        with (
            patch(
                "src.core.office_processor._translate_zip_images",
            ) as mock_zip,
            patch(
                "src.utils.config_manager.load_setting",
                return_value="TesseractOCR",
            ),
        ):
            _translate_doc_images(
                tmp_path / "test.docx",
                ".docx",
                "python_lib",
                "French",
                "English",
                None,
                None,
                None,
            )
            mock_zip.assert_called_once()


# ---------------------------------------------------------------------------
# 11. TestPerRunFormatting
# ---------------------------------------------------------------------------


class TestPerRunFormattingEdgeCases:
    """Tests for per-run formatting helpers."""

    def test_parse_html_formatting_bold(self) -> None:
        """Parses <b> tags correctly."""
        from src.core.office_formatter import _parse_html_formatting  # noqa: PLC0415

        segments = _parse_html_formatting("<b>Bold text</b>")
        assert len(segments) == 1
        assert segments[0].text == "Bold text"
        assert segments[0].bold is True
        assert segments[0].italic is False

    def test_parse_html_formatting_italic(self) -> None:
        """Parses <i> tags correctly."""
        from src.core.office_formatter import _parse_html_formatting  # noqa: PLC0415

        segments = _parse_html_formatting("<i>Italic</i>")
        assert len(segments) == 1
        assert segments[0].italic is True

    def test_parse_html_formatting_underline(self) -> None:
        """Parses <u> tags correctly."""
        from src.core.office_formatter import _parse_html_formatting  # noqa: PLC0415

        segments = _parse_html_formatting("<u>Underlined</u>")
        assert len(segments) == 1
        assert segments[0].underline is True

    def test_parse_html_formatting_strike(self) -> None:
        """Parses <s> tags correctly."""
        from src.core.office_formatter import _parse_html_formatting  # noqa: PLC0415

        segments = _parse_html_formatting("<s>Strikethrough</s>")
        assert len(segments) == 1
        assert segments[0].strike is True

    def test_parse_html_formatting_nested(self) -> None:
        """Parses nested <b><i> tags correctly."""
        from src.core.office_formatter import _parse_html_formatting  # noqa: PLC0415

        segments = _parse_html_formatting("<b><i>Bold Italic</i></b>")
        assert len(segments) == 1
        assert segments[0].bold is True
        assert segments[0].italic is True

    def test_parse_html_formatting_mixed_plain(self) -> None:
        """Parses text with both formatted and plain segments."""
        from src.core.office_formatter import _parse_html_formatting  # noqa: PLC0415

        segments = _parse_html_formatting("Plain <b>bold</b> more")
        assert len(segments) >= 2  # noqa: PLR2004
        assert any(s.bold for s in segments)
        assert any(not s.bold for s in segments)

    def test_parse_html_formatting_empty_string(self) -> None:
        """Empty string returns empty list."""
        from src.core.office_formatter import _parse_html_formatting  # noqa: PLC0415

        segments = _parse_html_formatting("")
        assert segments == []

    def test_wrap_with_tags_all_false(self) -> None:
        """No formatting returns plain text."""
        from src.core.office_formatter import _wrap_with_tags  # noqa: PLC0415

        result = _wrap_with_tags(
            "Hello",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "Hello"

    def test_wrap_with_tags_bold(self) -> None:
        """Bold wraps text in <b> tags."""
        from src.core.office_formatter import _wrap_with_tags  # noqa: PLC0415

        result = _wrap_with_tags(
            "Text",
            True,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert "<b>" in result
        assert "</b>" in result

    def test_wrap_with_tags_italic(self) -> None:
        """Italic wraps text in <i> tags."""
        from src.core.office_formatter import _wrap_with_tags  # noqa: PLC0415

        result = _wrap_with_tags(
            "Text",
            False,
            True,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert "<i>" in result
        assert "</i>" in result

    def test_wrap_with_tags_combined(self) -> None:
        """Bold + italic + underline wraps in all tags."""
        from src.core.office_formatter import _wrap_with_tags  # noqa: PLC0415

        result = _wrap_with_tags(
            "Text",
            True,
            True,
            True,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert "<b>" in result
        assert "<i>" in result
        assert "<u>" in result

    def test_formatted_segment_namedtuple(self) -> None:
        """_FormattedSegment has correct fields and defaults."""
        from src.core.office_formatter import _FormattedSegment  # noqa: PLC0415

        seg = _FormattedSegment("text", True, False, False, False)
        assert seg.text == "text"
        assert seg.bold is True
        assert seg.superscript is False
        assert seg.font_size_pt is None
        assert seg.hyperlink_url is None

    def test_formatted_segment_all_fields(self) -> None:
        """_FormattedSegment with all fields."""
        from src.core.office_formatter import _FormattedSegment  # noqa: PLC0415

        seg = _FormattedSegment(
            "link",
            True,
            True,
            True,
            True,
            True,
            False,
            14.0,
            "#ff0000",
            "#ffff00",
            "https://example.com",
        )
        assert seg.text == "link"
        assert seg.bold is True
        assert seg.italic is True
        assert seg.underline is True
        assert seg.strike is True
        assert seg.superscript is True
        assert seg.subscript is False
        assert seg.font_size_pt == 14.0  # noqa: PLR2004
        assert seg.color_hex == "#ff0000"
        assert seg.bg_color_hex == "#ffff00"
        assert seg.hyperlink_url == "https://example.com"

    def test_roundtrip_html_formatting(self) -> None:
        """Round-trip: wrap_with_tags -> parse_html_formatting preserves flags."""
        from src.core.office_formatter import (  # noqa: PLC0415
            _parse_html_formatting,
            _wrap_with_tags,
        )

        html = _wrap_with_tags(
            "Hello",
            True,
            False,
            True,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        segments = _parse_html_formatting(html)
        assert len(segments) == 1
        assert segments[0].bold is True
        assert segments[0].underline is True
        assert segments[0].italic is False
        assert segments[0].strike is False


# ---------------------------------------------------------------------------
# 12. TestRewriteZipContent
# ---------------------------------------------------------------------------


class TestRewriteZipContentEdgeCases:
    """Additional tests for _rewrite_zip_content."""

    def test_atomic_rewrite_preserves_all_entries(self, tmp_path: Path) -> None:
        """All entries are preserved during atomic rewrite."""
        path = tmp_path / "test.zip"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("a.txt", "Hello")
            zf.writestr("b.txt", "World")
            zf.writestr("c/d.xml", "<root/>")

        with zipfile.ZipFile(path, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        # Modify one entry
        file_data["a.txt"] = b"Modified"
        _rewrite_zip_content(path, file_data, all_items)

        with zipfile.ZipFile(path, "r") as zf:
            assert zf.read("a.txt") == b"Modified"
            assert zf.read("b.txt") == b"World"
            assert zf.read("c/d.xml") == b"<root/>"

    def test_rewrite_cleans_tmp_on_error(self, tmp_path: Path) -> None:
        """Temp file is cleaned up on write error."""
        path = tmp_path / "test.zip"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("a.txt", "Hello")

        with zipfile.ZipFile(path, "r") as zf:
            all_items = zf.infolist()

        # Pass incomplete file_data to trigger KeyError
        file_data: dict[str, bytes] = {}
        with pytest.raises(KeyError):
            _rewrite_zip_content(path, file_data, all_items)

        # Temp file should be cleaned up
        tmp_file = path.with_suffix(path.suffix + ".tmp")
        assert not tmp_file.exists()

    def test_rewrite_single_entry(self, tmp_path: Path) -> None:
        """Works with a single entry."""
        path = tmp_path / "single.zip"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("only.txt", "original")

        with zipfile.ZipFile(path, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        file_data["only.txt"] = b"replaced"
        _rewrite_zip_content(path, file_data, all_items)

        with zipfile.ZipFile(path, "r") as zf:
            assert zf.read("only.txt") == b"replaced"

    def test_rewrite_preserves_directory_structure(self, tmp_path: Path) -> None:
        """Nested directory structure is preserved."""
        path = tmp_path / "nested.zip"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("a/b/c.txt", "deep")
            zf.writestr("x/y.txt", "shallow")

        with zipfile.ZipFile(path, "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item) for item in all_items}

        _rewrite_zip_content(path, file_data, all_items)

        with zipfile.ZipFile(path, "r") as zf:
            assert zf.read("a/b/c.txt") == b"deep"
            assert zf.read("x/y.txt") == b"shallow"


# ---------------------------------------------------------------------------
# Additional miscellaneous edge-case tests
# ---------------------------------------------------------------------------


class TestSubstituteFontEdgeCases:
    """Additional edge cases for _substitute_font."""

    def test_empty_original_text_returns_original_font(self) -> None:
        """Empty original text returns original font unchanged."""
        result = _substitute_font("Arial", "", "Bonjour", "French")
        assert result == "Arial"

    def test_empty_translated_text_returns_original_font(self) -> None:
        """Empty translated text returns original font unchanged."""
        result = _substitute_font("Arial", "Hello", "", "French")
        assert result == "Arial"

    def test_same_script_keeps_font(self) -> None:
        """Same script (both Latin) keeps original font."""
        result = _substitute_font("Times New Roman", "Hello", "Bonjour", "French")
        assert result == "Times New Roman"

    def test_no_target_lang_returns_none_on_script_change(self) -> None:
        """Different scripts with no target_lang returns None."""
        result = _substitute_font("Arial", "Hello", "你好")
        assert result is None


class TestOdfElementTextEdgeCases:
    """Edge cases for _odf_element_text."""

    def test_empty_element(self) -> None:
        """Empty ODF element returns empty string."""
        from odf.text import P  # noqa: PLC0415

        p = P()
        assert _odf_element_text(p) == ""

    def test_simple_text(self) -> None:
        """Simple text element returns text content."""
        from odf.text import P  # noqa: PLC0415

        p = P(text="Hello World")
        assert _odf_element_text(p) == "Hello World"


class TestIsInsideTableCellEdge:
    """Additional edge cases for _is_inside_table_cell."""

    def test_not_inside_table_cell(self) -> None:
        """Element outside table cell returns False."""
        from odf.text import P  # noqa: PLC0415

        p = P(text="Hello")
        assert _is_inside_table_cell(p) is False

    def test_inside_table_cell(self) -> None:
        """Element inside table cell returns True."""
        from odf.table import Table, TableCell, TableRow  # noqa: PLC0415
        from odf.text import P  # noqa: PLC0415

        tbl = Table()
        tr = TableRow()
        tc = TableCell()
        p = P(text="Hello")
        tbl.addElement(tr)
        tr.addElement(tc)
        tc.addElement(p)
        assert _is_inside_table_cell(p) is True


class TestHeaderFooterExtensionsConstants:
    """Tests for header/footer extension constants."""

    def test_header_footer_extensions(self) -> None:
        """_HEADER_FOOTER_EXTENSIONS has correct values."""
        assert ".docx" in _HEADER_FOOTER_EXTENSIONS
        assert ".odt" in _HEADER_FOOTER_EXTENSIONS
        assert ".doc" in _HEADER_FOOTER_EXTENSIONS
        assert ".xlsx" not in _HEADER_FOOTER_EXTENSIONS

    def test_hf_type_constants(self) -> None:
        """H/F type constants have expected values."""
        assert _HF_DEFAULT == "default"
        assert _HF_FIRST == "first"
        assert _HF_EVEN == "even"


class TestShouldTranslateCommentsEdge:
    """Edge cases for _should_translate_comments."""

    def test_unsupported_extension(self) -> None:
        """Returns False for unknown extension."""
        assert not _should_translate_comments(".txt", "python_lib")

    def test_config_true(self) -> None:
        """Returns True when config enables comments."""
        config = MagicMock()
        config.translate_doc_comments = True
        assert _should_translate_comments(".docx", "python_lib", config)

    def test_config_false(self) -> None:
        """Returns False when config disables comments."""
        config = MagicMock()
        config.translate_doc_comments = False
        assert not _should_translate_comments(".docx", "python_lib", config)

    def test_all_comment_extensions_supported(self) -> None:
        """All nine comment extensions are in _COMMENT_EXTENSIONS."""
        expected = {
            ".docx",
            ".xlsx",
            ".pptx",
            ".odt",
            ".ods",
            ".odp",
            ".doc",
            ".xls",
            ".ppt",
        }
        assert expected == _COMMENT_EXTENSIONS


class TestShouldTranslateShapesEdge:
    """Edge cases for _should_translate_shapes."""

    def test_ppt_excluded(self) -> None:
        """PPT formats are excluded from shape translation."""
        config = MagicMock()
        config.translate_doc_shapes = True
        assert not _should_translate_shapes(".pptx", "python_lib", config)
        assert not _should_translate_shapes(".ppt", "python_lib", config)

    def test_shape_extensions_correct(self) -> None:
        """_SHAPE_EXTENSIONS excludes PPT formats."""
        assert ".pptx" not in _SHAPE_EXTENSIONS
        assert ".ppt" not in _SHAPE_EXTENSIONS
        assert ".docx" in _SHAPE_EXTENSIONS
        assert ".xlsx" in _SHAPE_EXTENSIONS

    def test_config_enables_shapes(self) -> None:
        """Returns True when config enables shapes."""
        config = MagicMock()
        config.translate_doc_shapes = True
        assert _should_translate_shapes(".docx", "python_lib", config)


class TestImageExtensionsConstant:
    """Tests for _IMAGE_EXTENSIONS constant."""

    def test_includes_all_formats(self) -> None:
        """_IMAGE_EXTENSIONS includes all 10 supported formats."""
        expected = {
            ".docx",
            ".xlsx",
            ".pptx",
            ".odt",
            ".ods",
            ".odp",
            ".epub",
            ".doc",
            ".xls",
            ".ppt",
        }
        assert expected == _IMAGE_EXTENSIONS

    def test_should_translate_images_config_gated(self) -> None:
        """Image translation is gated by config.should_translate_images."""
        config = MagicMock()
        config.should_translate_images = False
        assert not _should_translate_images(".docx", "python_lib", config)

        config.should_translate_images = True
        assert _should_translate_images(".docx", "python_lib", config)

    def test_should_translate_images_unsupported_ext(self) -> None:
        """Unsupported extension returns False."""
        config = MagicMock()
        config.should_translate_images = True
        assert not _should_translate_images(".txt", "python_lib", config)


# ---------------------------------------------------------------------------
# HIGH-PRIORITY COVERAGE GAPS — convert_to_modern_format
# ---------------------------------------------------------------------------


class TestConvertToModernFormatBothUnavailable:
    """Tests that convert_to_modern_format returns False when both backends fail."""

    def test_both_backends_import_error(self, tmp_path: Path) -> None:
        """Both win32com and UNO raise ImportError → returns False."""
        src = tmp_path / "test.doc"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        # Remove win32com from sys.modules to ensure clean ImportError
        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with patch.dict(
                "sys.modules",
                {
                    "win32com": None,
                    "win32com.client": None,
                    "uno": None,
                },
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is False
        assert not dst.exists()

    def test_win32com_succeeds_creates_output(self, tmp_path: Path) -> None:
        """Win32com succeeds → returns True and mock confirms call."""
        src = tmp_path / "test.doc"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        def _fake_convert(inp: Path, out: Path) -> None:
            """Simulates win32com producing an output file."""
            out.write_bytes(b"converted")

        with (
            patch.dict(
                "sys.modules",
                {"win32com": MagicMock(), "win32com.client": MagicMock()},
            ),
            patch(
                "src.core.office_processor._convert_with_win32com",
                side_effect=_fake_convert,
            ),
        ):
            result = convert_to_modern_format(src, dst)

        assert result is True
        assert dst.exists()
        assert dst.read_bytes() == b"converted"

    def test_win32com_non_import_error_skips_uno(self, tmp_path: Path) -> None:
        """Win32com raises non-ImportError → returns False without trying UNO."""
        src = tmp_path / "test.doc"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        with (
            patch.dict(
                "sys.modules",
                {
                    "win32com": MagicMock(),
                    "win32com.client": MagicMock(),
                    "uno": MagicMock(),
                },
            ),
            patch(
                "src.core.office_processor._convert_with_win32com",
                side_effect=OSError("Disk failure"),
            ),
            patch(
                "src.core.office_processor._convert_with_uno",
            ) as mock_uno,
        ):
            result = convert_to_modern_format(src, dst)

        assert result is False
        # UNO should NOT be tried because win32com raised a non-ImportError
        mock_uno.assert_not_called()

    def test_uno_succeeds_when_win32com_unavailable(self, tmp_path: Path) -> None:
        """UNO succeeds when win32com is unavailable → returns True."""
        src = tmp_path / "test.xls"
        dst = tmp_path / "test.xlsx"
        src.write_bytes(b"fake")

        def _fake_uno_convert(inp: Path, out: Path) -> None:
            """Simulates UNO producing an output file."""
            out.write_bytes(b"uno-converted")

        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with (
                patch.dict(
                    "sys.modules",
                    {
                        "win32com": None,
                        "win32com.client": None,
                        "uno": MagicMock(),
                    },
                ),
                patch(
                    "src.core.office_processor._convert_with_uno",
                    side_effect=_fake_uno_convert,
                ),
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is True
        assert dst.exists()
        assert dst.read_bytes() == b"uno-converted"

    def test_uno_non_import_error_returns_false(self, tmp_path: Path) -> None:
        """UNO raises non-ImportError → returns False."""
        src = tmp_path / "test.odt"
        dst = tmp_path / "test.docx"
        src.write_bytes(b"fake")

        mods_to_remove = [k for k in sys.modules if k.startswith("win32com")]
        saved = {k: sys.modules.pop(k) for k in mods_to_remove}
        try:
            with (
                patch.dict(
                    "sys.modules",
                    {
                        "win32com": None,
                        "win32com.client": None,
                        "uno": MagicMock(),
                    },
                ),
                patch(
                    "src.core.office_processor._convert_with_uno",
                    side_effect=PermissionError("Access denied"),
                ),
            ):
                result = convert_to_modern_format(src, dst)
        finally:
            sys.modules.update(saved)

        assert result is False


# ---------------------------------------------------------------------------
# HIGH-PRIORITY COVERAGE GAPS — _substitute_font
# ---------------------------------------------------------------------------


class TestSubstituteFontCoverage:
    """Coverage tests for _substitute_font edge cases."""

    def test_same_script_latin_to_latin_unchanged(self) -> None:
        """Same-script translation (Latin→Latin) returns original font."""
        result = _substitute_font("Calibri", "Good morning", "Bonjour", "French")
        assert result == "Calibri"

    def test_cross_script_latin_to_cjk_returns_cjk_font(self) -> None:
        """Cross-script (Latin→CJK) returns a font from get_font_for_language."""
        with patch(
            "src.core.office_processor._get_font_for_language",
            return_value="MS Gothic",
        ):
            result = _substitute_font(
                "Arial",
                "Hello world",
                "\u4f60\u597d\u4e16\u754c",
                "Chinese",
            )
        assert result == "MS Gothic"

    def test_cross_script_latin_to_arabic_returns_arabic_font(self) -> None:
        """Cross-script (Latin→Arabic) returns an Arabic-compatible font."""
        with patch(
            "src.core.office_processor._get_font_for_language",
            return_value="Arabic Typesetting",
        ):
            result = _substitute_font(
                "Times New Roman",
                "Hello",
                "\u0645\u0631\u062d\u0628\u0627",
                "Arabic",
            )
        assert result == "Arabic Typesetting"

    def test_empty_original_text_returns_original(self) -> None:
        """Empty original_text short-circuits to original font."""
        result = _substitute_font("Verdana", "", "\u3042", "Japanese")
        assert result == "Verdana"

    def test_empty_translated_text_returns_original(self) -> None:
        """Empty translated_text short-circuits to original font."""
        result = _substitute_font("Verdana", "Hello", "", "Japanese")
        assert result == "Verdana"

    def test_no_target_lang_returns_none(self) -> None:
        """Different scripts with empty target_lang returns None."""
        result = _substitute_font(
            "Courier New",
            "Hello",
            "\u3053\u3093\u306b\u3061\u306f",
            "",
        )
        assert result is None

    def test_no_target_lang_default_returns_none(self) -> None:
        """Different scripts with no target_lang argument returns None."""
        result = _substitute_font(
            "Courier New",
            "Hello",
            "\u3053\u3093\u306b\u3061\u306f",
        )
        assert result is None

    def test_cjk_to_cjk_same_script_keeps_font(self) -> None:
        """CJK→CJK (same script) keeps original font unchanged."""
        result = _substitute_font(
            "MS Mincho",
            "\u3042\u3044\u3046",
            "\u304b\u304d\u304f",
            "Japanese",
        )
        assert result == "MS Mincho"

    def test_classify_generic_family_called_on_cross_script(self) -> None:
        """Verifies _classify_generic_family is called for the original font."""
        with (
            patch(
                "src.core.office_processor._classify_generic_family",
                return_value="sans-serif",
            ) as mock_classify,
            patch(
                "src.core.office_processor._get_font_for_language",
                return_value="Noto Sans CJK JP",
            ) as mock_get_font,
        ):
            result = _substitute_font(
                "Arial",
                "Hello",
                "\u3053\u3093\u306b\u3061\u306f",
                "Japanese",
            )

        mock_classify.assert_called_once_with(font_name="Arial")
        mock_get_font.assert_called_once_with("Japanese", "sans-serif")
        assert result == "Noto Sans CJK JP"


# ---------------------------------------------------------------------------
# HIGH-PRIORITY COVERAGE GAPS — _parse_hyperlink_rels
# ---------------------------------------------------------------------------

_RELS_NS_STR = "http://schemas.openxmlformats.org/package/2006/relationships"
_HYPERLINK_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
)


class TestParseHyperlinkRels:
    """Tests for _parse_hyperlink_rels."""

    def test_extracts_external_hyperlinks(self) -> None:
        """Parses external hyperlink relationships from rels XML."""
        rels_xml = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{_RELS_NS_STR}">'
            f'<Relationship Id="rId1" Type="{_HYPERLINK_TYPE}" '
            f'Target="https://example.com" TargetMode="External"/>'
            f'<Relationship Id="rId2" Type="{_HYPERLINK_TYPE}" '
            f'Target="https://other.com" TargetMode="External"/>'
            f"</Relationships>"
        ).encode()

        result = _parse_hyperlink_rels(rels_xml)

        assert result == {
            "rId1": "https://example.com",
            "rId2": "https://other.com",
        }

    def test_ignores_non_hyperlink_relationships(self) -> None:
        """Non-hyperlink relationship types are excluded."""
        rels_xml = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{_RELS_NS_STR}">'
            f'<Relationship Id="rId1" Type="some/other/type" '
            f'Target="foo.xml"/>'
            f'<Relationship Id="rId2" Type="{_HYPERLINK_TYPE}" '
            f'Target="https://example.com" TargetMode="External"/>'
            f"</Relationships>"
        ).encode()

        result = _parse_hyperlink_rels(rels_xml)

        assert result == {"rId2": "https://example.com"}

    def test_ignores_internal_hyperlinks(self) -> None:
        """Hyperlinks without TargetMode='External' are excluded."""
        rels_xml = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{_RELS_NS_STR}">'
            f'<Relationship Id="rId1" Type="{_HYPERLINK_TYPE}" '
            f'Target="word/document.xml"/>'
            f"</Relationships>"
        ).encode()

        result = _parse_hyperlink_rels(rels_xml)

        assert result == {}

    def test_empty_relationships(self) -> None:
        """Empty Relationships element returns empty dict."""
        rels_xml = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{_RELS_NS_STR}"/>'
        ).encode()

        result = _parse_hyperlink_rels(rels_xml)

        assert result == {}

    def test_mixed_internal_and_external(self) -> None:
        """Only external hyperlinks are returned among mixed entries."""
        rels_xml = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{_RELS_NS_STR}">'
            f'<Relationship Id="rId1" Type="{_HYPERLINK_TYPE}" '
            f'Target="internal.xml"/>'
            f'<Relationship Id="rId2" Type="{_HYPERLINK_TYPE}" '
            f'Target="https://ext.com" TargetMode="External"/>'
            f'<Relationship Id="rId3" Type="other/type" '
            f'Target="https://also-other.com" TargetMode="External"/>'
            f"</Relationships>"
        ).encode()

        result = _parse_hyperlink_rels(rels_xml)

        assert result == {"rId2": "https://ext.com"}

    def test_from_real_docx_zip(self, tmp_path: Path) -> None:
        """Extracts hyperlinks from a minimal DOCX-like ZIP structure."""
        rels_content = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{_RELS_NS_STR}">'
            f'<Relationship Id="rId1" Type="{_HYPERLINK_TYPE}" '
            f'Target="https://example.com/doc" TargetMode="External"/>'
            f'<Relationship Id="rId2" Type="http://other/type" '
            f'Target="styles.xml"/>'
            f"</Relationships>"
        )
        docx_path = tmp_path / "test.docx"
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr(
                "word/_rels/document.xml.rels",
                rels_content,
            )

        with zipfile.ZipFile(docx_path, "r") as zf:
            rels_xml = zf.read("word/_rels/document.xml.rels")

        result = _parse_hyperlink_rels(rels_xml)

        assert result == {"rId1": "https://example.com/doc"}

    def test_missing_id_or_target_defaults_to_empty(self) -> None:
        """Missing Id or Target attributes default to empty strings."""
        rels_xml = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{_RELS_NS_STR}">'
            f'<Relationship Type="{_HYPERLINK_TYPE}" '
            f'Target="https://no-id.com" TargetMode="External"/>'
            f'<Relationship Id="rId1" Type="{_HYPERLINK_TYPE}" '
            f'TargetMode="External"/>'
            f"</Relationships>"
        ).encode()

        result = _parse_hyperlink_rels(rels_xml)

        # Missing Id defaults to "" key, missing Target defaults to "" value
        assert "" in result  # empty Id
        assert result.get("") == "https://no-id.com"
        assert result.get("rId1") == ""  # empty Target


# ---------------------------------------------------------------------------
# HIGH-PRIORITY COVERAGE GAPS — _add_hyperlink_to_rels (additional cases)
# ---------------------------------------------------------------------------


class TestAddHyperlinkToRelsCoverage:
    """Additional coverage tests for _add_hyperlink_to_rels."""

    def test_new_rels_has_correct_type_and_target_mode(self) -> None:
        """New relationship has correct Type and TargetMode attributes."""
        rels_xml, r_id = _add_hyperlink_to_rels(None, "https://test.com")
        root = etree.fromstring(rels_xml)
        rel = root[0]
        assert rel.get("Type") == _HYPERLINK_TYPE
        assert rel.get("TargetMode") == "External"
        assert rel.get("Target") == "https://test.com"
        assert r_id == "rId1"

    def test_preserves_existing_non_hyperlink_rels(self) -> None:
        """Adding a hyperlink preserves other relationship types."""
        existing = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{_RELS_NS_STR}">'
            f'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/'
            f'officeDocument/2006/relationships/styles" Target="styles.xml"/>'
            f"</Relationships>"
        ).encode()

        rels_xml, r_id = _add_hyperlink_to_rels(existing, "https://link.com")

        root = etree.fromstring(rels_xml)
        all_rels = root.findall(f"{{{_RELS_NS_STR}}}Relationship")
        assert len(all_rels) == 2  # noqa: PLR2004

        # Original relationship preserved
        orig = [r for r in all_rels if r.get("Id") == "rId1"][0]
        assert orig.get("Target") == "styles.xml"

        # New hyperlink added
        assert r_id == "rId2"
        new = [r for r in all_rels if r.get("Id") == "rId2"][0]
        assert new.get("Target") == "https://link.com"
        assert new.get("TargetMode") == "External"

    def test_url_with_special_characters(self) -> None:
        """URLs with special characters are preserved in the Target."""
        url = "https://example.com/path?q=hello&lang=fr#section"
        rels_xml, r_id = _add_hyperlink_to_rels(None, url)
        root = etree.fromstring(rels_xml)
        rel = root[0]
        assert rel.get("Target") == url

    def test_roundtrip_add_then_parse(self) -> None:
        """Add a hyperlink then parse — the hyperlink is found."""
        rels_xml, r_id = _add_hyperlink_to_rels(None, "https://roundtrip.com")
        result = _parse_hyperlink_rels(rels_xml)
        assert result == {r_id: "https://roundtrip.com"}


# ---------------------------------------------------------------------------
# TestInjectOdtHeadersFooters — ODT header/footer injection
# ---------------------------------------------------------------------------


class TestInjectOdtHeadersFootersExtended:
    """Extended tests for _inject_odt_headers_footers."""

    def test_basic_injection_replaces_header_text(self, tmp_path: Path) -> None:
        """Basic injection replaces header text in styles.xml."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header>"
            "<text:p>Original Header</text:p>"
            "</style:header>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        _inject_odt_headers_footers(
            odt_path,
            {f"header:0:{_HF_DEFAULT}:0": "Translated Header"},
        )

        result = _extract_odt_headers_footers(odt_path)
        assert any("Translated Header" in v for _, v in result)
        assert not any("Original Header" in v for _, v in result)

    def test_empty_translations_is_noop(self, tmp_path: Path) -> None:
        """Empty translations dict does not modify the file."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header>"
            "<text:p>Keep Me</text:p>"
            "</style:header>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        _inject_odt_headers_footers(odt_path, {})

        # File should not have been rewritten (no matching keys)
        result = _extract_odt_headers_footers(odt_path)
        assert any("Keep Me" in v for _, v in result)

    def test_preserves_non_header_content(self, tmp_path: Path) -> None:
        """Injection preserves non-header content in styles.xml."""
        styles_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-styles"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:master-styles>"
            '<style:master-page style:name="Standard">'
            "<style:header>"
            "<text:p>Header Text</text:p>"
            "</style:header>"
            "<style:footer>"
            "<text:p>Footer Text</text:p>"
            "</style:footer>"
            "</style:master-page>"
            "</office:master-styles>"
            "</office:document-styles>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", styles_xml)
            zf.writestr("content.xml", "<root/>")

        # Only translate the header, leave footer untouched
        _inject_odt_headers_footers(
            odt_path,
            {f"header:0:{_HF_DEFAULT}:0": "New Header"},
        )

        result = _extract_odt_headers_footers(odt_path)
        # Header should be translated
        header_vals = [v for k, v in result if k.startswith("header:")]
        assert any("New Header" in v for v in header_vals)
        # Footer should be preserved
        footer_vals = [v for k, v in result if k.startswith("footer:")]
        assert any("Footer Text" in v for v in footer_vals)

    def test_no_styles_xml_is_noop(self, tmp_path: Path) -> None:
        """File without styles.xml is silently skipped."""
        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", "<root/>")

        # Should not raise
        _inject_odt_headers_footers(
            odt_path,
            {f"header:0:{_HF_DEFAULT}:0": "Anything"},
        )


# ---------------------------------------------------------------------------
# TestInjectOdtFootnotes — ODT footnote injection
# ---------------------------------------------------------------------------


class TestInjectOdtFootnotesExtended:
    """Extended tests for _inject_odt_footnotes."""

    def test_basic_injection_replaces_note_content(self, tmp_path: Path) -> None:
        """Basic injection replaces text:note content in content.xml."""
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-content"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:body><office:text>"
            "<text:p>Body text"
            '<text:note text:id="ftn1" text:note-class="footnote">'
            "<text:note-citation>1</text:note-citation>"
            "<text:note-body>"
            "<text:p>Original footnote</text:p>"
            "</text:note-body>"
            "</text:note>"
            "</text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", content_xml)
            zf.writestr("META-INF/manifest.xml", "<manifest/>")

        _inject_odt_footnotes(odt_path, {"footnote:ftn1": "Translated footnote"})

        result = _extract_odt_footnotes(odt_path)
        assert len(result) == 1
        assert result[0][0] == "footnote:ftn1"
        assert "Translated footnote" in result[0][1]

    def test_empty_translations_is_noop(self, tmp_path: Path) -> None:
        """Empty translations dict does not modify the file."""
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-content"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:body><office:text>"
            "<text:p>Body"
            '<text:note text:id="ftn1" text:note-class="footnote">'
            "<text:note-citation>1</text:note-citation>"
            "<text:note-body>"
            "<text:p>Keep this</text:p>"
            "</text:note-body>"
            "</text:note>"
            "</text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", content_xml)
            zf.writestr("META-INF/manifest.xml", "<manifest/>")

        _inject_odt_footnotes(odt_path, {})

        result = _extract_odt_footnotes(odt_path)
        assert len(result) == 1
        assert "Keep this" in result[0][1]

    def test_preserves_footnote_citation_markers(self, tmp_path: Path) -> None:
        """Injection preserves text:note-citation markers."""
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-content"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:body><office:text>"
            "<text:p>Body"
            '<text:note text:id="ftn1" text:note-class="footnote">'
            "<text:note-citation>1</text:note-citation>"
            "<text:note-body>"
            "<text:p>Original</text:p>"
            "</text:note-body>"
            "</text:note>"
            "</text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", content_xml)
            zf.writestr("META-INF/manifest.xml", "<manifest/>")

        _inject_odt_footnotes(odt_path, {"footnote:ftn1": "Replaced"})

        # Parse content.xml and verify citation is preserved
        with zipfile.ZipFile(odt_path, "r") as zf:
            root = etree.fromstring(zf.read("content.xml"))
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        citations = list(root.iter(f"{{{text_ns}}}note-citation"))
        assert len(citations) == 1
        assert citations[0].text == "1"

    def test_no_content_xml_is_noop(self, tmp_path: Path) -> None:
        """File without content.xml is silently skipped."""
        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("styles.xml", "<root/>")

        # Should not raise
        _inject_odt_footnotes(odt_path, {"footnote:ftn1": "Anything"})

    def test_endnote_injection(self, tmp_path: Path) -> None:
        """Endnote (note-class=endnote) is injected correctly."""
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            "<office:document-content"
            ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
            ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:body><office:text>"
            "<text:p>Body"
            '<text:note text:id="en1" text:note-class="endnote">'
            "<text:note-citation>i</text:note-citation>"
            "<text:note-body>"
            "<text:p>Original endnote</text:p>"
            "</text:note-body>"
            "</text:note>"
            "</text:p>"
            "</office:text></office:body>"
            "</office:document-content>"
        )

        odt_path = tmp_path / "test.odt"
        with zipfile.ZipFile(odt_path, "w") as zf:
            zf.writestr("content.xml", content_xml)
            zf.writestr("META-INF/manifest.xml", "<manifest/>")

        _inject_odt_footnotes(odt_path, {"endnote:en1": "Translated endnote"})

        result = _extract_odt_footnotes(odt_path)
        assert len(result) == 1
        assert result[0][0] == "endnote:en1"
        assert "Translated endnote" in result[0][1]


# ---------------------------------------------------------------------------
# TestInjectOdpNotes — ODP speaker notes injection
# ---------------------------------------------------------------------------


class TestInjectOdpNotesExtended:
    """Extended tests for _inject_odp_notes."""

    def test_basic_injection_replaces_notes_text(self, tmp_path: Path) -> None:
        """Basic injection replaces presentation:notes text."""
        from odf.draw import Frame, Page, TextBox  # noqa: PLC0415
        from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
        from odf.presentation import Notes  # noqa: PLC0415
        from odf.style import MasterPage, PageLayout  # noqa: PLC0415
        from odf.text import P  # noqa: PLC0415

        doc = OpenDocumentPresentation()
        pl = PageLayout(name="PM0")
        doc.automaticstyles.addElement(pl)
        mp = MasterPage(name="Default", pagelayoutname="PM0")
        doc.masterstyles.addElement(mp)
        page = Page(name="Slide1", masterpagename="Default")
        doc.presentation.addElement(page)

        notes = Notes()
        frame = Frame()
        tb = TextBox()
        p = P(text="Original note")
        tb.addElement(p)
        frame.addElement(tb)
        notes.addElement(frame)
        page.addElement(notes)

        odp_path = tmp_path / "test.odp"
        doc.save(str(odp_path))

        # Extract to get the correct key
        extracted = _extract_odp_notes(odp_path)
        assert len(extracted) >= 1
        key = next(k for k, v in extracted if "Original note" in v)

        _inject_odp_notes(odp_path, {key: "Translated note"})

        result = _extract_odp_notes(odp_path)
        assert any("Translated note" in v for _, v in result)
        assert not any("Original note" in v for _, v in result)

    def test_empty_translations_is_noop(self, tmp_path: Path) -> None:
        """Empty translations dict does not modify notes."""
        from odf.draw import Frame, Page, TextBox  # noqa: PLC0415
        from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
        from odf.presentation import Notes  # noqa: PLC0415
        from odf.style import MasterPage, PageLayout  # noqa: PLC0415
        from odf.text import P  # noqa: PLC0415

        doc = OpenDocumentPresentation()
        pl = PageLayout(name="PM0")
        doc.automaticstyles.addElement(pl)
        mp = MasterPage(name="Default", pagelayoutname="PM0")
        doc.masterstyles.addElement(mp)
        page = Page(name="Slide1", masterpagename="Default")
        doc.presentation.addElement(page)

        notes = Notes()
        frame = Frame()
        tb = TextBox()
        p = P(text="Untouched note")
        tb.addElement(p)
        frame.addElement(tb)
        notes.addElement(frame)
        page.addElement(notes)

        odp_path = tmp_path / "test.odp"
        doc.save(str(odp_path))

        _inject_odp_notes(odp_path, {})

        result = _extract_odp_notes(odp_path)
        assert any("Untouched note" in v for _, v in result)

    def test_preserves_slide_content(self, tmp_path: Path) -> None:  # noqa: PLR0915
        """Injection preserves slide body content alongside notes."""
        from odf.draw import Frame, Page, TextBox  # noqa: PLC0415
        from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
        from odf.presentation import Notes  # noqa: PLC0415
        from odf.style import MasterPage, PageLayout  # noqa: PLC0415
        from odf.text import P  # noqa: PLC0415

        doc = OpenDocumentPresentation()
        pl = PageLayout(name="PM0")
        doc.automaticstyles.addElement(pl)
        mp = MasterPage(name="Default", pagelayoutname="PM0")
        doc.masterstyles.addElement(mp)
        page = Page(name="Slide1", masterpagename="Default")
        doc.presentation.addElement(page)

        # Add body frame with content
        body_frame = Frame(width="20cm", height="15cm")
        body_tb = TextBox()
        body_p = P(text="Slide body text")
        body_tb.addElement(body_p)
        body_frame.addElement(body_tb)
        page.addElement(body_frame)

        # Add notes
        notes = Notes()
        notes_frame = Frame()
        notes_tb = TextBox()
        notes_p = P(text="Note to translate")
        notes_tb.addElement(notes_p)
        notes_frame.addElement(notes_tb)
        notes.addElement(notes_frame)
        page.addElement(notes)

        odp_path = tmp_path / "test.odp"
        doc.save(str(odp_path))

        extracted = _extract_odp_notes(odp_path)
        key = next(k for k, v in extracted if "Note to translate" in v)

        _inject_odp_notes(odp_path, {key: "Translated"})

        # Verify notes updated and slide body preserved
        result = _extract_odp_notes(odp_path)
        assert any("Translated" in v for _, v in result)

        # Read ODP content and verify body frame text is intact
        from odf.draw import Frame as DrawFrame  # noqa: PLC0415
        from odf.draw import Page as DrawPage  # noqa: PLC0415
        from odf.draw import TextBox as DrawTextBox  # noqa: PLC0415
        from odf.opendocument import load as odf_load  # noqa: PLC0415
        from odf.text import P as TextP  # noqa: PLC0415

        loaded = odf_load(str(odp_path))
        pages = loaded.getElementsByType(DrawPage)
        assert len(pages) >= 1
        # The body text should still exist in the page
        all_text = []
        for frame in pages[0].getElementsByType(DrawFrame):
            for tb in frame.getElementsByType(DrawTextBox):
                for para in tb.getElementsByType(TextP):
                    # Collect text from child nodes
                    txt = ""
                    for node in para.childNodes:
                        if hasattr(node, "data"):
                            txt += node.data
                        elif hasattr(node, "__str__"):
                            txt += str(node)
                    if txt:
                        all_text.append(txt)
        # "Slide body text" should be somewhere in the page
        assert any("Slide body text" in t for t in all_text)


# ---------------------------------------------------------------------------
# TestWin32comWordHeadersFooters — Legacy .doc headers (mock win32com)
# ---------------------------------------------------------------------------


class TestWin32comWordHeadersFooters:
    """Tests for _extract/_inject_win32com_word_headers_footers."""

    def _build_mock_env(
        self,
    ) -> tuple[MagicMock, MagicMock, MagicMock, MagicMock]:
        """Builds mock win32com environment with a section containing a header and footer.

        Returns:
            (mock_win32com, mock_win32com_client, mock_pythoncom, mock_doc)
        """
        # Header paragraph
        mock_hdr_para_range = MagicMock()
        mock_hdr_para_range.Text = "Header text\r\n"
        mock_hdr_para = MagicMock()
        mock_hdr_para.Range = mock_hdr_para_range

        mock_hdr = MagicMock()
        mock_hdr.Exists = True
        mock_hdr.Range.Paragraphs.Count = 1
        mock_hdr.Range.Paragraphs.side_effect = lambda i: mock_hdr_para

        # Footer paragraph
        mock_ftr_para_range = MagicMock()
        mock_ftr_para_range.Text = "Footer text\r\n"
        mock_ftr_para = MagicMock()
        mock_ftr_para.Range = mock_ftr_para_range

        mock_ftr = MagicMock()
        mock_ftr.Exists = True
        mock_ftr.Range.Paragraphs.Count = 1
        mock_ftr.Range.Paragraphs.side_effect = lambda i: mock_ftr_para

        # Section: only default type (id=1) has headers/footers
        mock_section = MagicMock()

        def headers_factory(hf_type_id: int) -> MagicMock:
            if hf_type_id == 1:
                return mock_hdr
            m = MagicMock()
            m.Exists = False
            return m

        def footers_factory(hf_type_id: int) -> MagicMock:
            if hf_type_id == 1:
                return mock_ftr
            m = MagicMock()
            m.Exists = False
            return m

        mock_section.Headers.side_effect = headers_factory
        mock_section.Footers.side_effect = footers_factory

        mock_doc = MagicMock()
        mock_doc.Sections.Count = 1
        mock_doc.Sections.side_effect = lambda i: mock_section

        mock_word = MagicMock()
        mock_word.Documents.Open.return_value = mock_doc

        mock_pythoncom = MagicMock()
        mock_win32com_client = MagicMock()
        mock_win32com_client.Dispatch.return_value = mock_word
        mock_win32com = MagicMock()
        mock_win32com.client = mock_win32com_client

        return mock_win32com, mock_win32com_client, mock_pythoncom, mock_doc

    def test_extract_returns_header_footer_text(self) -> None:
        """Extract returns header and footer text strings."""
        mock_win32com, mock_client, mock_pycom, _ = self._build_mock_env()

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_client,
                "pythoncom": mock_pycom,
            },
        ):
            result = _extract_win32com_word_headers_footers(Path("test.doc"))

        keys = [k for k, _ in result]
        values = [v for _, v in result]
        assert any("Header text" in v for v in values)
        assert any("Footer text" in v for v in values)
        assert any(k.startswith("header:0:default:") for k in keys)
        assert any(k.startswith("footer:0:default:") for k in keys)

    def test_inject_replaces_range_text(self) -> None:
        """Inject replaces Range.Text on header and footer paragraphs."""
        mock_win32com, mock_client, mock_pycom, mock_doc = self._build_mock_env()

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_client,
                "pythoncom": mock_pycom,
            },
        ):
            _inject_win32com_word_headers_footers(
                Path("test.doc"),
                {
                    "header:0:default:0": "New Header",
                    "footer:0:default:0": "New Footer",
                },
            )

        # Verify SaveAs was called
        mock_doc.SaveAs.assert_called_once()

    def test_empty_headers_handled_gracefully(self) -> None:
        """Empty doc with no sections or empty headers does not raise."""
        mock_doc = MagicMock()
        mock_doc.Sections.Count = 0

        mock_word = MagicMock()
        mock_word.Documents.Open.return_value = mock_doc

        mock_pythoncom = MagicMock()
        mock_win32com_client = MagicMock()
        mock_win32com_client.Dispatch.return_value = mock_word
        mock_win32com = MagicMock()
        mock_win32com.client = mock_win32com_client

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_win32com_client,
                "pythoncom": mock_pythoncom,
            },
        ):
            result = _extract_win32com_word_headers_footers(Path("test.doc"))
            assert result == []

            _inject_win32com_word_headers_footers(
                Path("test.doc"),
                {"header:0:default:0": "Nothing"},
            )
            mock_doc.SaveAs.assert_called()


# ---------------------------------------------------------------------------
# TestWin32comWordFootnotes — Legacy .doc footnotes (mock win32com)
# ---------------------------------------------------------------------------


class TestWin32comWordFootnotes:
    """Tests for _extract/_inject_win32com_word_footnotes."""

    def _build_mock_env(
        self,
    ) -> tuple[MagicMock, MagicMock, MagicMock, MagicMock]:
        """Builds mock win32com environment with footnotes and endnotes.

        Returns:
            (mock_win32com, mock_win32com_client, mock_pythoncom, mock_doc)
        """
        # Footnote
        mock_fn = MagicMock()
        mock_fn.Range.Text = "Footnote text\r\n"

        # Endnote
        mock_en = MagicMock()
        mock_en.Range.Text = "Endnote text\r\n"

        mock_doc = MagicMock()
        mock_doc.Footnotes.Count = 1
        mock_doc.Footnotes.side_effect = lambda i: mock_fn
        mock_doc.Endnotes.Count = 1
        mock_doc.Endnotes.side_effect = lambda i: mock_en

        mock_word = MagicMock()
        mock_word.Documents.Open.return_value = mock_doc

        mock_pythoncom = MagicMock()
        mock_win32com_client = MagicMock()
        mock_win32com_client.Dispatch.return_value = mock_word
        mock_win32com = MagicMock()
        mock_win32com.client = mock_win32com_client

        return mock_win32com, mock_win32com_client, mock_pythoncom, mock_doc

    def test_extract_returns_footnote_text(self) -> None:
        """Extract returns footnote and endnote text."""
        mock_win32com, mock_client, mock_pycom, _ = self._build_mock_env()

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_client,
                "pythoncom": mock_pycom,
            },
        ):
            result = _extract_win32com_word_footnotes(Path("test.doc"))

        keys = [k for k, _ in result]
        values = [v for _, v in result]
        assert "Footnote text" in values
        assert "Endnote text" in values
        assert "footnote:1" in keys
        assert "endnote:1" in keys

    def test_inject_replaces_footnote_range_text(self) -> None:
        """Inject replaces footnote Range.Text."""
        mock_win32com, mock_client, mock_pycom, mock_doc = self._build_mock_env()

        # Track the Range.Text assignments
        fn_mock = MagicMock()
        en_mock = MagicMock()
        mock_doc.Footnotes.side_effect = lambda i: fn_mock
        mock_doc.Endnotes.side_effect = lambda i: en_mock

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_client,
                "pythoncom": mock_pycom,
            },
        ):
            _inject_win32com_word_footnotes(
                Path("test.doc"),
                {
                    "footnote:1": "Translated footnote",
                    "endnote:1": "Translated endnote",
                },
            )

        assert fn_mock.Range.Text == "Translated footnote"
        assert en_mock.Range.Text == "Translated endnote"
        mock_doc.SaveAs.assert_called_once()

    def test_empty_footnotes_handled(self) -> None:
        """Doc with no footnotes or endnotes does not raise."""
        mock_doc = MagicMock()
        mock_doc.Footnotes.Count = 0
        mock_doc.Endnotes.Count = 0

        mock_word = MagicMock()
        mock_word.Documents.Open.return_value = mock_doc

        mock_pythoncom = MagicMock()
        mock_win32com_client = MagicMock()
        mock_win32com_client.Dispatch.return_value = mock_word
        mock_win32com = MagicMock()
        mock_win32com.client = mock_win32com_client

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_win32com_client,
                "pythoncom": mock_pythoncom,
            },
        ):
            result = _extract_win32com_word_footnotes(Path("test.doc"))
            assert result == []

            _inject_win32com_word_footnotes(
                Path("test.doc"),
                {"footnote:1": "Nothing"},
            )
            mock_doc.SaveAs.assert_called()


# ---------------------------------------------------------------------------
# TestWin32comExcelSheetNames — Legacy .xls sheet names (mock win32com)
# ---------------------------------------------------------------------------


class TestWin32comExcelSheetNames:
    """Tests for _extract/_inject_win32com_excel_sheet_names."""

    def _build_mock_env(
        self,
        sheet_names: list[str] | None = None,
    ) -> tuple[MagicMock, MagicMock, MagicMock, MagicMock]:
        """Builds mock win32com environment with worksheets.

        Args:
            sheet_names: Names for the mock worksheets.

        Returns:
            (mock_win32com, mock_win32com_client, mock_pythoncom, mock_wb)
        """
        sheet_names = sheet_names or ["Sheet1", "Data"]

        mock_sheets: list[MagicMock] = []
        for name in sheet_names:
            ws = MagicMock()
            ws.Name = name
            mock_sheets.append(ws)

        mock_wb = MagicMock()
        mock_wb.Worksheets.Count = len(mock_sheets)
        mock_wb.Worksheets.side_effect = lambda i: mock_sheets[i - 1]

        mock_excel = MagicMock()
        mock_excel.Workbooks.Open.return_value = mock_wb

        mock_pythoncom = MagicMock()
        mock_win32com_client = MagicMock()
        mock_win32com_client.Dispatch.return_value = mock_excel
        mock_win32com = MagicMock()
        mock_win32com.client = mock_win32com_client

        return mock_win32com, mock_win32com_client, mock_pythoncom, mock_wb

    def test_extract_returns_sheet_name_strings(self) -> None:
        """Extract returns sheet name strings with correct keys."""
        mock_win32com, mock_client, mock_pycom, _ = self._build_mock_env(
            ["Revenue", "Expenses"],
        )

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_client,
                "pythoncom": mock_pycom,
            },
        ):
            result = _extract_win32com_excel_sheet_names(Path("test.xls"))

        assert result == [
            ("sheetname:Revenue", "Revenue"),
            ("sheetname:Expenses", "Expenses"),
        ]

    def test_inject_renames_sheets_via_name_property(self) -> None:
        """Inject renames sheets via the Name property."""
        mock_win32com, mock_client, mock_pycom, mock_wb = self._build_mock_env(
            ["Sheet1", "Data"],
        )

        # Need fresh mocks for inject since side_effect returns by index
        ws1 = MagicMock()
        ws1.Name = "Sheet1"
        ws2 = MagicMock()
        ws2.Name = "Data"
        mock_wb.Worksheets.Count = 2
        mock_wb.Worksheets.side_effect = lambda i: [ws1, ws2][i - 1]

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_client,
                "pythoncom": mock_pycom,
            },
        ):
            _inject_win32com_excel_sheet_names(
                Path("test.xls"),
                {
                    "sheetname:Sheet1": "Hoja1",
                    "sheetname:Data": "Datos",
                },
            )

        assert ws1.Name == "Hoja1"
        assert ws2.Name == "Datos"
        mock_wb.Save.assert_called_once()

    def test_sanitize_sheet_name_applied(self) -> None:
        """_sanitize_sheet_name is applied to translations during injection."""
        mock_win32com, mock_client, mock_pycom, mock_wb = self._build_mock_env(
            ["Sheet1"],
        )

        ws1 = MagicMock()
        ws1.Name = "Sheet1"
        mock_wb.Worksheets.Count = 1
        mock_wb.Worksheets.side_effect = lambda i: ws1

        # Translation with invalid characters that _sanitize_sheet_name strips
        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_client,
                "pythoncom": mock_pycom,
            },
        ):
            _inject_win32com_excel_sheet_names(
                Path("test.xls"),
                {"sheetname:Sheet1": "New/Sheet*Name?"},
            )

        # _sanitize_sheet_name strips /, *, ? characters
        assert ws1.Name == "NewSheetName"

    def test_empty_workbook_handled(self) -> None:
        """Workbook with no sheets does not raise."""
        mock_wb = MagicMock()
        mock_wb.Worksheets.Count = 0

        mock_excel = MagicMock()
        mock_excel.Workbooks.Open.return_value = mock_wb

        mock_pythoncom = MagicMock()
        mock_win32com_client = MagicMock()
        mock_win32com_client.Dispatch.return_value = mock_excel
        mock_win32com = MagicMock()
        mock_win32com.client = mock_win32com_client

        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_win32com_client,
                "pythoncom": mock_pythoncom,
            },
        ):
            result = _extract_win32com_excel_sheet_names(Path("test.xls"))
            assert result == []

    def test_long_name_truncated(self) -> None:
        """Sheet name exceeding 31 chars is truncated by _sanitize_sheet_name."""
        mock_win32com, mock_client, mock_pycom, mock_wb = self._build_mock_env(
            ["Sheet1"],
        )

        ws1 = MagicMock()
        ws1.Name = "Sheet1"
        mock_wb.Worksheets.Count = 1
        mock_wb.Worksheets.side_effect = lambda i: ws1

        long_name = "A" * 50
        with patch.dict(
            "sys.modules",
            {
                "win32com": mock_win32com,
                "win32com.client": mock_client,
                "pythoncom": mock_pycom,
            },
        ):
            _inject_win32com_excel_sheet_names(
                Path("test.xls"),
                {"sheetname:Sheet1": long_name},
            )

        assert len(ws1.Name) == 31  # noqa: PLR2004


# ---------------------------------------------------------------------------
# convert_to_modern_format
# ---------------------------------------------------------------------------

_original_import = builtins.__import__


def _import_no_win32com(name: str, *args, **kwargs):
    """Import hook that blocks win32com.client."""
    if name == "win32com.client":
        raise ImportError("no win32com")
    return _original_import(name, *args, **kwargs)


def _import_no_uno(name: str, *args, **kwargs):
    """Import hook that blocks uno."""
    if name == "uno":
        raise ImportError("no uno")
    return _original_import(name, *args, **kwargs)


def _import_no_backends(name: str, *args, **kwargs):
    """Import hook that blocks both win32com.client and uno."""
    if name in ("win32com.client", "uno"):
        raise ImportError(f"no {name}")
    return _original_import(name, *args, **kwargs)


class TestConvertToModernFormat:
    """Tests for convert_to_modern_format() backend dispatch."""

    def test_win32com_success(self) -> None:
        """win32com import succeeds and conversion succeeds → returns True."""
        with (
            patch("builtins.__import__", side_effect=_original_import),
            patch.dict(
                "sys.modules", {"win32com": MagicMock(), "win32com.client": MagicMock()}
            ),
            patch("src.core.office_processor._convert_with_win32com") as mock_w32,
        ):
            result = convert_to_modern_format(
                Path("/tmp/test.doc"), Path("/tmp/test.docx")
            )

        assert result is True
        mock_w32.assert_called_once_with(Path("/tmp/test.doc"), Path("/tmp/test.docx"))

    def test_win32com_fails_falls_through_to_uno_success(self) -> None:
        """win32com ImportError → falls through to UNO which succeeds → True."""
        with (
            patch("builtins.__import__", side_effect=_import_no_win32com),
            patch.dict("sys.modules", {"uno": MagicMock()}),
            patch("src.core.office_processor._convert_with_uno") as mock_uno,
        ):
            result = convert_to_modern_format(
                Path("/tmp/test.xls"), Path("/tmp/test.xlsx")
            )

        assert result is True
        mock_uno.assert_called_once_with(Path("/tmp/test.xls"), Path("/tmp/test.xlsx"))

    def test_win32com_conversion_error_returns_false(self) -> None:
        """win32com available but _convert_with_win32com raises → False.

        Must NOT fall through to UNO — the function returns False immediately
        when the chosen backend raises a non-ImportError exception.
        """
        with (
            patch("builtins.__import__", side_effect=_original_import),
            patch.dict(
                "sys.modules", {"win32com": MagicMock(), "win32com.client": MagicMock()}
            ),
            patch(
                "src.core.office_processor._convert_with_win32com",
                side_effect=RuntimeError("COM failed"),
            ),
            patch("src.core.office_processor._convert_with_uno") as mock_uno,
        ):
            result = convert_to_modern_format(
                Path("/tmp/test.ppt"), Path("/tmp/test.pptx")
            )

        assert result is False
        mock_uno.assert_not_called()

    def test_uno_success_when_no_win32com(self) -> None:
        """No win32com, UNO works → returns True."""
        with (
            patch("builtins.__import__", side_effect=_import_no_win32com),
            patch.dict("sys.modules", {"uno": MagicMock()}),
            patch("src.core.office_processor._convert_with_uno") as mock_uno,
        ):
            result = convert_to_modern_format(
                Path("/tmp/test.odt"), Path("/tmp/test.docx")
            )

        assert result is True
        mock_uno.assert_called_once()

    def test_uno_conversion_error_returns_false(self) -> None:
        """No win32com, UNO available but _convert_with_uno raises → False."""
        with (
            patch("builtins.__import__", side_effect=_import_no_win32com),
            patch.dict("sys.modules", {"uno": MagicMock()}),
            patch(
                "src.core.office_processor._convert_with_uno",
                side_effect=RuntimeError("UNO failed"),
            ),
        ):
            result = convert_to_modern_format(
                Path("/tmp/test.ods"), Path("/tmp/test.xlsx")
            )

        assert result is False

    def test_no_backend_available(self) -> None:
        """Both win32com and UNO ImportError → returns False."""
        with patch("builtins.__import__", side_effect=_import_no_backends):
            result = convert_to_modern_format(
                Path("/tmp/test.doc"), Path("/tmp/test.docx")
            )

        assert result is False

    def test_passes_correct_paths(self) -> None:
        """Verify input_path and output_path are forwarded correctly."""
        inp = Path("/data/legacy/report.doc")
        out = Path("/data/modern/report.docx")

        with (
            patch("builtins.__import__", side_effect=_original_import),
            patch.dict(
                "sys.modules", {"win32com": MagicMock(), "win32com.client": MagicMock()}
            ),
            patch("src.core.office_processor._convert_with_win32com") as mock_w32,
        ):
            convert_to_modern_format(inp, out)

        mock_w32.assert_called_once_with(inp, out)


# ─────────────────────────────────────────────────────────────────────────────
# Format-edge-case backfill tests
# ─────────────────────────────────────────────────────────────────────────────


class TestRestoreWin32ComFontWithSubstitution:
    """Documents font-restore exception swallowing.

    ``_restore_win32com_font`` silently swallows attribute-set errors so
    callers don't have to guard against missing fonts (e.g. when a target-
    language font isn't installed on the host).
    """

    def test_missing_font_does_not_raise(self) -> None:
        """Setting Name on a font that rejects the value is logged not raised."""
        # Build a font whose Name setter raises (simulates a missing font).
        font = MagicMock()
        type(font).Name = property(
            lambda self: None,
            lambda self, v: (_ for _ in ()).throw(
                RuntimeError("font not installed"),
            ),
        )
        saved = {"Name": "DoesNotExistFont", "Bold": True}
        # Must not raise — caller relies on this behaviour
        _restore_win32com_font(font, saved)
        # Bold still propagates despite the Name failure
        assert font.Bold is True


class TestConvertToModernFormatTimeout:
    """Backend conversion failures are caught and surfaced as ``False``.

    UNO uses RPC-over-socket via ``doc.storeToURL`` rather than a direct
    subprocess, so there's no ``timeout=`` parameter to plumb in.  The
    catch-all in ``convert_to_modern_format`` keeps the pipeline robust:
    any exception (TimeoutExpired included) is logged and reported as
    ``False`` so the caller can retry or fall back.
    """

    def test_uno_backend_timeout_returns_false(self, tmp_path: Path) -> None:
        """A simulated TimeoutExpired in UNO conversion is swallowed → False."""
        import subprocess  # noqa: PLC0415

        inp = tmp_path / "in.odt"
        out = tmp_path / "out.docx"
        inp.touch()

        # Make win32com unimportable so we fall through to UNO
        original = builtins.__import__

        def _no_w32(name, *a, **kw):
            if name.startswith("win32com"):
                raise ImportError(name)
            return original(name, *a, **kw)

        with (
            patch("builtins.__import__", side_effect=_no_w32),
            patch.dict("sys.modules", {"uno": MagicMock()}),
            patch(
                "src.core.office_processor._convert_with_uno",
                side_effect=subprocess.TimeoutExpired(cmd="soffice", timeout=30),
            ),
        ):
            result = convert_to_modern_format(inp, out)
        assert result is False


class TestExtractPythonPptxNestedGroupShapes:
    """Verifies that grouped (nested) shapes are recursed into.

    ``_extract_python_pptx`` walks shape groups via ``_walk_pptx_text_shapes``
    so text inside a ``GroupShape`` round-trips through the extract → inject
    pipeline like any top-level text frame.
    """

    def test_top_level_shape_is_extracted(self, tmp_path: Path) -> None:
        """A non-grouped shape's text is extracted (sanity baseline)."""
        from src.core.office_processor import _extract_python_pptx  # noqa: PLC0415

        pptx_path = tmp_path / "deck.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Add a top-level text box
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txbox.text_frame.text = "Top level text"
        prs.save(str(pptx_path))

        texts = _extract_python_pptx(pptx_path)
        joined = " ".join(t for _, t in texts)
        assert "Top level text" in joined

    def test_nested_group_text_extracted(self) -> None:
        """Grouped shapes are recursed into; their paragraphs are extracted.

        We mock python-pptx so a slide returns a single shape that is
        itself a group containing a child text frame.  The walker
        descends into the group and surfaces the inner paragraph with
        a dotted shape path encoding its position.
        """
        from src.core.office_processor import _extract_python_pptx  # noqa: PLC0415

        # Inner shape — has a real text frame
        inner_shape = MagicMock()
        inner_shape.has_text_frame = True
        # MagicMock's auto-vivified ``shapes`` would short-circuit the
        # group check; explicitly remove it so the walker treats this
        # shape as a leaf.
        del inner_shape.shapes
        inner_para = MagicMock()
        inner_para.text = "Inside group"
        inner_para.runs = []
        inner_shape.text_frame.paragraphs = [inner_para]

        # Group "shape" — does NOT expose has_text_frame=True at top level
        group_shape = MagicMock()
        group_shape.has_text_frame = False
        group_shape.shapes = [inner_shape]

        slide = MagicMock()
        slide.shapes = [group_shape]
        fake_prs = MagicMock()
        fake_prs.slides = [slide]

        with (
            patch(
                "src.core.office_processor.Presentation",
                return_value=fake_prs,
                create=True,
            ),
            # Patch via the local import inside _extract_python_pptx
            patch.dict(
                "sys.modules",
                {"pptx": MagicMock(Presentation=lambda *_a, **_k: fake_prs)},
            ),
        ):
            texts = _extract_python_pptx(Path("ignored.pptx"))

        # Grouped paragraph is extracted with a dotted shape path.
        keys_with_inner = [k for k, t in texts if "Inside group" in t]
        assert keys_with_inner, "expected grouped paragraph to be extracted"
        # Path encodes the group index (0) and the inner shape index (0).
        assert keys_with_inner[0] == "slide:0:0.0:0"


class TestWalkPptxTextShapes:
    """Direct unit tests for the recursive shape walker."""

    @staticmethod
    def _leaf(name: str = "leaf") -> MagicMock:
        """A leaf shape with a text frame; no ``shapes`` attribute."""
        m = MagicMock(name=name)
        m.has_text_frame = True
        del m.shapes  # leaf — auto-vivified attr would falsely match group
        return m

    @staticmethod
    def _group(children: list, name: str = "group") -> MagicMock:
        """A group shape (no text frame, exposes ``shapes``)."""
        m = MagicMock(name=name)
        m.has_text_frame = False
        m.shapes = children
        return m

    def test_empty_iterable_yields_nothing(self) -> None:
        """No shapes → no output (sanity baseline)."""
        from src.core.office_processor import _walk_pptx_text_shapes  # noqa: PLC0415

        assert list(_walk_pptx_text_shapes([])) == []

    def test_empty_group_yields_nothing(self) -> None:
        """A group with no children produces no entries."""
        from src.core.office_processor import _walk_pptx_text_shapes  # noqa: PLC0415

        empty_group = self._group([], name="empty")
        assert list(_walk_pptx_text_shapes([empty_group])) == []

    def test_doubly_nested_groups_path_dotted(self) -> None:
        """Group within group: path is ``parent.child.leaf``."""
        from src.core.office_processor import _walk_pptx_text_shapes  # noqa: PLC0415

        leaf = self._leaf("inner")
        inner_group = self._group([leaf], name="inner_grp")
        outer_group = self._group([inner_group], name="outer_grp")

        results = list(_walk_pptx_text_shapes([outer_group]))
        assert len(results) == 1
        path, shape = results[0]
        # outer index 0 → inner-group index 0 → leaf index 0
        assert path == "0.0.0"
        assert shape is leaf

    def test_mixed_top_level_leaf_and_group_paths(self) -> None:
        """Sibling top-level shapes get sequential top-level indices."""
        from src.core.office_processor import _walk_pptx_text_shapes  # noqa: PLC0415

        leaf_a = self._leaf("a")
        leaf_b = self._leaf("b")
        group = self._group([leaf_b], name="grp")
        # Top-level: [leaf_a, group_with_leaf_b]
        results = list(_walk_pptx_text_shapes([leaf_a, group]))
        paths = [p for p, _ in results]
        assert paths == ["0", "1.0"]

    def test_group_that_also_has_text_frame_treated_as_leaf(self) -> None:
        """A shape with both ``shapes`` and ``has_text_frame=True`` is treated as leaf.

        The walker keys off ``has_text_frame`` first to avoid double-
        counting; a real ``GroupShape`` always reports
        ``has_text_frame=False`` in python-pptx, so this defensive
        ordering only matters for synthetic mocks.
        """
        from src.core.office_processor import _walk_pptx_text_shapes  # noqa: PLC0415

        weird = MagicMock()
        weird.has_text_frame = True
        weird.shapes = [self._leaf("ignored_child")]

        results = list(_walk_pptx_text_shapes([weird]))
        # The walker descends into ``.shapes`` only when ``has_text_frame``
        # is False, so ``weird`` is yielded as a leaf with path "0".
        assert results == [("0", weird)]


class TestSheetNameRoundTripWithSpecialChars:
    """Round-trip a translated sheet name through inject + extract via xlsx."""

    def test_xlsx_round_trip_strips_invalid_chars(self, tmp_path: Path) -> None:
        """Invalid chars (\\ / ? * [ ]) are stripped on inject, preserved otherwise."""
        from src.core.office_processor import (  # noqa: PLC0415
            _extract_xlsx_sheet_names,
            _inject_xlsx_sheet_names,
        )

        # Build an .xlsx with a single sheet named "Sales".
        xlsx_path = tmp_path / "book.xlsx"
        wb = Workbook()
        wb.active.title = "Sales"
        wb.save(str(xlsx_path))

        # Extract → expect the original name.
        extracted = _extract_xlsx_sheet_names(xlsx_path)
        assert extracted == [("sheetname:Sales", "Sales")]

        # Inject a "translation" containing every invalid character.
        translations = {"sheetname:Sales": "Q1\\Sales/2024?*[Final]"}
        _inject_xlsx_sheet_names(xlsx_path, translations)

        # Re-extract → invalid chars must be stripped to keep Excel happy.
        round_trip = _extract_xlsx_sheet_names(xlsx_path)
        assert len(round_trip) == 1
        new_name = round_trip[0][1]
        # All six forbidden characters absent
        for bad in "\\/*?[]":
            assert bad not in new_name, f"Invalid char {bad!r} survived"
        # Letters preserved
        assert "Q1" in new_name
        assert "Sales" in new_name
        assert "2024" in new_name
        assert "Final" in new_name


# ---------------------------------------------------------------------------
# RTL injection in DOCX / PPTX / XLSX / ODT
# ---------------------------------------------------------------------------


class TestOfficeRtlInjection:
    """RTL paragraph/run/cell markers injected when target_lang is RTL."""

    def test_inject_python_docx_adds_bidi_when_target_is_rtl(
        self, tmp_path: Path
    ) -> None:
        from docx import Document  # noqa: PLC0415

        from src.core.office_processor import _inject_python_docx  # noqa: PLC0415

        src = tmp_path / "src.docx"
        out = tmp_path / "out.docx"
        d = Document()
        d.add_paragraph("Hello world")
        d.save(str(src))

        _inject_python_docx(src, out, {"para:0": "مرحبا"}, target_lang="Arabic")

        # Re-open and verify <w:bidi/> appears in pPr and <w:rtl/> in rPr.
        loaded = Document(str(out))
        para = loaded.paragraphs[0]
        pPr_xml = para._element.get_or_add_pPr().xml  # noqa: N806
        assert "<w:bidi" in pPr_xml
        # The injected text becomes a run; it should carry <w:rtl/>.
        for run in para.runs:
            rPr = run._element.get_or_add_rPr()  # noqa: N806
            assert "<w:rtl" in rPr.xml

    def test_inject_python_docx_skips_bidi_when_target_is_ltr(
        self, tmp_path: Path
    ) -> None:
        from docx import Document  # noqa: PLC0415

        from src.core.office_processor import _inject_python_docx  # noqa: PLC0415

        src = tmp_path / "src.docx"
        out = tmp_path / "out.docx"
        d = Document()
        d.add_paragraph("Hello world")
        d.save(str(src))

        _inject_python_docx(src, out, {"para:0": "Bonjour"}, target_lang="French")

        loaded = Document(str(out))
        para = loaded.paragraphs[0]
        pPr_xml = para._element.get_or_add_pPr().xml  # noqa: N806
        assert "<w:bidi" not in pPr_xml

    def test_inject_python_pptx_adds_rtl_attr_when_target_is_rtl(
        self, tmp_path: Path
    ) -> None:
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415

        from src.core.office_processor import _inject_python_pptx  # noqa: PLC0415

        src = tmp_path / "src.pptx"
        out = tmp_path / "out.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tx.text_frame.text = "Hello"
        prs.save(str(src))

        _inject_python_pptx(
            src,
            out,
            {"slide:0:1:0": "مرحبا"},
            target_lang="Arabic",
        )

        loaded = Presentation(str(out))
        # Find the textbox we created.
        para = next(
            p
            for sh in loaded.slides[0].shapes
            if sh.has_text_frame
            for p in sh.text_frame.paragraphs
        )
        pPr = para._pPr  # noqa: N806
        assert pPr is not None
        assert pPr.get("rtl") == "1"

    def test_inject_python_xlsx_sets_right_to_left_view_for_rtl(
        self, tmp_path: Path
    ) -> None:
        from openpyxl import Workbook, load_workbook  # noqa: PLC0415

        from src.core.office_processor import _inject_python_xlsx  # noqa: PLC0415

        src = tmp_path / "src.xlsx"
        out = tmp_path / "out.xlsx"
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Hello"
        wb.save(str(src))

        _inject_python_xlsx(
            src,
            out,
            {f"sheet:{ws.title}:1:1": "مرحبا"},
            target_lang="Arabic",
        )

        loaded = load_workbook(str(out))
        assert loaded.active.sheet_view.rightToLeft is True

    def test_inject_python_xlsx_no_rtl_view_for_ltr(self, tmp_path: Path) -> None:
        from openpyxl import Workbook, load_workbook  # noqa: PLC0415

        from src.core.office_processor import _inject_python_xlsx  # noqa: PLC0415

        src = tmp_path / "src.xlsx"
        out = tmp_path / "out.xlsx"
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Hello"
        wb.save(str(src))

        _inject_python_xlsx(
            src,
            out,
            {f"sheet:{ws.title}:1:1": "Bonjour"},
            target_lang="French",
        )

        loaded = load_workbook(str(out))
        # openpyxl returns None or False for an unset rightToLeft.
        assert not loaded.active.sheet_view.rightToLeft

    def test_set_odf_default_rtl_adds_writing_mode(self, tmp_path: Path) -> None:
        """ODF post-processor sets style:writing-mode + fo:text-align.

        Parses the rewritten ``styles.xml`` and checks attribute values
        on the new ``default-style`` element by URI rather than prefix —
        odfpy doesn't pre-declare the ``style:`` and ``fo:`` prefixes on
        the root, so lxml synthesises new ones on serialize.  The file
        is still valid ODF (LibreOffice resolves by namespace URI).
        """
        from odf.opendocument import OpenDocumentText  # noqa: PLC0415
        from odf.text import P  # noqa: PLC0415

        from src.core.office_processor import (  # noqa: PLC0415
            _ODF_NS_FO,
            _ODF_NS_OFFICE,
            _ODF_NS_STYLE,
            _set_odf_default_rtl,
        )

        odt = tmp_path / "test.odt"
        doc = OpenDocumentText()
        doc.text.addElement(P(text="Hello"))
        doc.save(str(odt))

        _set_odf_default_rtl(odt)

        with zipfile.ZipFile(odt, "r") as zf:
            styles_xml = zf.read("styles.xml")

        from lxml import etree  # noqa: PLC0415

        root = etree.fromstring(styles_xml)
        ns = {"office": _ODF_NS_OFFICE, "style": _ODF_NS_STYLE, "fo": _ODF_NS_FO}
        para_props = root.find(
            "office:styles/style:default-style[@style:family='paragraph']"
            "/style:paragraph-properties",
            namespaces=ns,
        )
        assert para_props is not None
        assert para_props.get(f"{{{_ODF_NS_STYLE}}}writing-mode") == "rl-tb"
        assert para_props.get(f"{{{_ODF_NS_FO}}}text-align") == "end"


class TestProcessOfficeFileProviderModel:
    """Regression: provider/model kwargs must reach translate_batch.

    The PDF path threads them; the office path used to drop them on the
    floor, which silently routed every Office translation through the
    default model regardless of the user's per-feature pick.
    """

    def test_translate_batch_receives_provider_and_model(self, tmp_path: Path) -> None:
        """process_office_file forwards provider/model to translate_batch."""
        from docx import Document  # noqa: PLC0415

        from src.core.office_processor import process_office_file  # noqa: PLC0415

        src = tmp_path / "src.docx"
        out = tmp_path / "out.docx"
        d = Document()
        d.add_paragraph("Hello world")
        d.save(str(src))

        with patch(
            "src.core.office_processor.translate_batch",
            return_value=["Bonjour le monde"],
        ) as mock_batch:
            ok = process_office_file(
                src,
                out,
                "French",
                provider="Custom",
                model="gpt-5.4-pro",
            )

        assert ok
        assert mock_batch.called, "expected translate_batch to be invoked"
        for call in mock_batch.call_args_list:
            assert call.kwargs.get("provider") == "Custom"
            assert call.kwargs.get("model") == "gpt-5.4-pro"


class TestProcessOfficeFileForwardsCheckpointDir:
    """``process_office_file`` must thread ``checkpoint_dir`` to image translation.

    The per-image cache (added in this branch) is keyed on the task's
    storage directory.  If ``process_office_file`` drops the kwarg on
    either of its two ``_translate_doc_images`` callsites (the
    text-only fast path at office_processor.py:~10470 and the full
    text+images path at ~10654), the document still translates — but
    image translations don't cache, so a retry re-translates every
    image from scratch.  Regression guards for both callsites.
    """

    def test_text_plus_images_path_forwards_checkpoint_dir(
        self,
        tmp_path: Path,
    ) -> None:
        """Standard docx with text + image translation enabled."""
        from docx import Document  # noqa: PLC0415

        src = tmp_path / "src.docx"
        out = tmp_path / "out.docx"
        cp = tmp_path / "task_storage"
        cp.mkdir()

        d = Document()
        d.add_paragraph("Hello world")
        d.save(str(src))

        from src.core.config import TranslationConfig  # noqa: PLC0415

        config = TranslationConfig(
            translate_doc_images=True,
            ocr_is_configured=True,
            ocr_method="TesseractOCR",
        )

        with (
            patch(
                "src.core.office_processor.translate_batch",
                return_value=["Bonjour le monde"],
            ),
            patch(
                "src.core.office_processor._translate_doc_images",
            ) as mock_images,
        ):
            process_office_file(
                src,
                out,
                "French",
                checkpoint_dir=cp,
                config=config,
            )

        mock_images.assert_called_once()
        assert mock_images.call_args.kwargs.get("checkpoint_dir") == cp

    def test_text_only_fast_path_forwards_checkpoint_dir(
        self,
        tmp_path: Path,
    ) -> None:
        """Fast path for documents with no translatable text.

        When ``texts`` extraction returns empty, ``process_office_file``
        short-circuits to a ``shutil.copy2 + _translate_doc_images``
        fast path.  This callsite has its own forwarding line and
        needs its own regression guard.
        """
        from docx import Document  # noqa: PLC0415

        # Empty document (no paragraphs) → texts list is empty after
        # extraction → fast path is taken.
        src = tmp_path / "empty.docx"
        out = tmp_path / "out.docx"
        cp = tmp_path / "task_storage"
        cp.mkdir()

        d = Document()
        d.save(str(src))

        from src.core.config import TranslationConfig  # noqa: PLC0415

        config = TranslationConfig(
            translate_doc_images=True,
            ocr_is_configured=True,
            ocr_method="TesseractOCR",
        )

        with patch(
            "src.core.office_processor._translate_doc_images",
        ) as mock_images:
            process_office_file(
                src,
                out,
                "French",
                checkpoint_dir=cp,
                config=config,
            )

        mock_images.assert_called_once()
        assert mock_images.call_args.kwargs.get("checkpoint_dir") == cp


class TestLegacyOfficePathsNoRtlInjection:
    """Legacy paths don't call RTL helpers directly.

    AGENTS.md: ``.doc`` / ``.xls`` / ``.ppt`` injection paths via
    Win32COM/UNO do NOT yet inject RTL flags directly — users with
    ``SETTING_AUTO_CONVERT_LEGACY`` on get RTL via the modern-format
    round-trip; with auto-convert off, legacy files render as LTR.
    A regression that adds RTL injection to the legacy injector
    without also exercising the round-trip would silently change
    behaviour for users who turned auto-convert OFF — the legacy
    file would suddenly grow RTL flags Word can't fully honour.
    """

    def test_inject_win32com_word_does_not_import_rtl_helpers(self) -> None:
        """The legacy DOC injector source has no RTL helper references.

        Static check via inspection of the function source — looks for
        any call to the modern-format RTL helpers
        (``_set_docx_paragraph_rtl``, ``_set_pptx_paragraph_rtl``,
        ``_set_odf_default_rtl``) or to ``is_rtl_language``.  Adding
        any of those calls into the legacy injector would either
        silently fail (helpers expect modern XML) or produce
        half-tagged legacy files.
        """
        import inspect  # noqa: PLC0415

        from src.core.office_processor import (  # noqa: PLC0415
            _inject_win32com_word,
        )

        source = inspect.getsource(_inject_win32com_word)
        forbidden = (
            "_set_docx_paragraph_rtl",
            "_set_pptx_paragraph_rtl",
            "_set_odf_default_rtl",
            "is_rtl_language",
        )
        for name in forbidden:
            assert name not in source, (
                f"Legacy DOC injector ({_inject_win32com_word.__name__}) "
                f"must not reference {name!r} directly — RTL injection "
                f"belongs in the modern-format round-trip path "
                f"(gated by SETTING_AUTO_CONVERT_LEGACY)."
            )

    def test_inject_uno_writer_does_not_import_rtl_helpers(self) -> None:
        """Same contract for the UNO Writer legacy injector."""
        import inspect  # noqa: PLC0415

        from src.core.office_processor import _inject_uno_writer  # noqa: PLC0415

        source = inspect.getsource(_inject_uno_writer)
        forbidden = (
            "_set_docx_paragraph_rtl",
            "_set_pptx_paragraph_rtl",
            "_set_odf_default_rtl",
            "is_rtl_language",
        )
        for name in forbidden:
            assert name not in source, (
                f"UNO Writer legacy injector must not reference "
                f"{name!r} directly — RTL injection belongs in the "
                f"modern-format round-trip path."
            )

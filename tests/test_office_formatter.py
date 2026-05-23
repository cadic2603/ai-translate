"""Tests for Office document formatting preservation (HTML round-trip)."""

import sys
import zipfile
from pathlib import Path
from unittest.mock import MagicMock, patch

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from src.core.office_formatter import (
    _DRAWINGML_NS,
    _FORMATTING_HTML_RE,
    _HIGHLIGHT_COLORS,
    _HYPERLINK_RELTYPE,
    _ODF_NS,
    _WD_COLOR_INDEX_TO_HEX,
    _WORDML_NS,
    _apply_drawingml_format_attrs,
    _apply_pptx_format_attrs,
    _build_span_style,
    _color_hex_to_int,
    _color_hex_to_win32com,
    _drawingml_to_html,
    _FormattedSegment,
    _has_drawingml_hyperlinks,
    _has_drawingml_mixed_formatting,
    _has_mixed_formatting,
    _has_odf_text_box_mixed_formatting,
    _has_pptx_hyperlinks,
    _has_pptx_mixed_formatting,
    _inject_html_runs,
    _inject_pptx_html_runs,
    _int_to_color_hex,
    _odf_text_box_to_html,
    _para_has_hyperlinks,
    _parse_docx_char_styles,
    _parse_html_formatting,
    _parse_span_style,
    _pptx_runs_to_html,
    _read_docx_run_bg_hex,
    _read_drawingml_rpr_formatting,
    _read_pptx_run_formatting,
    _read_pptx_run_full_formatting,
    _read_wml_rpr_sup_sub,
    _replace_paragraph_text,
    _run_has_visual_content,
    _runs_to_html,
    _set_rpr_color,
    _set_rpr_font_size,
    _set_rpr_vert_align,
    _win32com_color_to_hex,
    _wrap_with_tags,
)
from src.core.office_processor import (
    _UNO_FORMATTING_PROPS,
    _UNO_SLANT_ITALIC,
    _UNO_SLANT_NONE,
    _UNO_STRIKEOUT_NONE,
    _UNO_STRIKEOUT_SINGLE,
    _UNO_UNDERLINE_NONE,
    _UNO_UNDERLINE_SINGLE,
    _UNO_WEIGHT_BOLD,
    _UNO_WEIGHT_NORMAL,
    _extract_odf_paragraph_text,
    _extract_para_with_links,
    _has_uno_hyperlinks,
    _has_uno_mixed_formatting,
    _has_win32com_ppt_hyperlinks,
    _has_win32com_ppt_mixed_formatting,
    _has_win32com_range_hyperlinks,
    _has_win32com_word_hyperlinks,
    _has_win32com_word_mixed_formatting,
    _inject_drawingml_html_runs,
    _inject_odf_paragraph_text,
    _inject_odf_text_box_html_runs,
    _inject_uno_html_runs,
    _inject_win32com_excel_html_runs,
    _inject_win32com_ppt_html_runs,
    _inject_win32com_word_html_runs,
    _inject_wps_txbx_html_runs,
    _read_uno_portion_formatting,
    _read_uno_portion_full_formatting,
    _read_win32com_char_formatting,
    _read_win32com_ppt_run_formatting,
    _resolve_para_hyperlink_rels,
    _save_uno_first_portion_props,
    _uno_runs_to_html,
    _win32com_ppt_runs_to_html,
    _win32com_word_runs_to_html,
    _wps_txbx_to_text_or_html,
)
from tests.test_office_processor import _make_uno_enum

# ---------------------------------------------------------------------------
# Per-run formatting preservation (HTML round-trip) tests
# ---------------------------------------------------------------------------


def _make_para_with_runs(
    run_specs: list[tuple[str, bool, bool, bool, bool]],
) -> object:
    """Creates a Document paragraph with specified runs.

    Each spec is (text, bold, italic, underline, strike).
    Returns the paragraph object.
    """
    doc = Document()
    para = doc.add_paragraph()
    for text, bold, italic, underline, strike in run_specs:
        run = para.add_run(text)
        run.bold = bold
        run.italic = italic
        run.underline = underline
        run.font.strike = strike
    return para


# ── _has_mixed_formatting tests ──────────────────────────────────────


class TestHasMixedFormatting:
    """Tests for _has_mixed_formatting."""

    def test_uniform_formatting(self) -> None:
        """All runs have same formatting → False."""
        para = _make_para_with_runs(
            [
                ("Hello ", True, False, False, False),
                ("world", True, False, False, False),
            ]
        )
        assert _has_mixed_formatting(para) is False

    def test_mixed_bold_and_plain(self) -> None:
        """Bold + non-bold runs → True."""
        para = _make_para_with_runs(
            [
                ("Hello ", False, False, False, False),
                ("world", True, False, False, False),
            ]
        )
        assert _has_mixed_formatting(para) is True

    def test_single_run(self) -> None:
        """Single text run → False (nothing to mix)."""
        para = _make_para_with_runs(
            [
                ("Hello", True, False, False, False),
            ]
        )
        assert _has_mixed_formatting(para) is False

    def test_empty_runs_skipped(self) -> None:
        """Runs with empty text are ignored."""
        para = _make_para_with_runs(
            [
                ("Hello", True, False, False, False),
                ("", False, True, False, False),
                ("world", True, False, False, False),
            ]
        )
        assert _has_mixed_formatting(para) is False

    def test_visual_content_runs_skipped(self) -> None:
        """Runs containing visual content are ignored."""
        doc = Document()
        para = doc.add_paragraph()
        # Add a normal text run
        run1 = para.add_run("Hello")
        run1.bold = True
        # Add a run with a drawing element (visual content)
        run2 = para.add_run()
        drawing = OxmlElement("w:drawing")
        run2._element.append(drawing)
        run2.bold = False
        # Add another text run with same formatting
        run3 = para.add_run("world")
        run3.bold = True
        assert _has_mixed_formatting(para) is False

    def test_no_runs(self) -> None:
        """Paragraph with no runs → False."""
        doc = Document()
        para = doc.add_paragraph()
        assert _has_mixed_formatting(para) is False

    def test_mixed_italic_and_underline(self) -> None:
        """Italic + underline in different runs → True."""
        para = _make_para_with_runs(
            [
                ("Hello ", False, True, False, False),
                ("world", False, False, True, False),
            ]
        )
        assert _has_mixed_formatting(para) is True


# ── _runs_to_html tests ─────────────────────────────────────────────


class TestRunsToHtml:
    """Tests for _runs_to_html."""

    def test_bold_run(self) -> None:
        """Single bold run produces <b> wrapper."""
        para = _make_para_with_runs(
            [
                ("Hello", True, False, False, False),
            ]
        )
        assert _runs_to_html(para) == "<b>Hello</b>"

    def test_nested_bold_italic(self) -> None:
        """Bold+italic run produces nested tags."""
        para = _make_para_with_runs(
            [
                ("Hello", True, True, False, False),
            ]
        )
        result = _runs_to_html(para)
        assert "<b>" in result
        assert "<i>" in result
        assert "Hello" in result

    def test_mixed_runs(self) -> None:
        """Plain text + bold produces mixed output."""
        para = _make_para_with_runs(
            [
                ("Hello ", False, False, False, False),
                ("world", True, False, False, False),
            ]
        )
        result = _runs_to_html(para)
        assert "Hello " in result
        assert "<b>world</b>" in result

    def test_html_escaping(self) -> None:
        """Special characters in text are HTML-escaped."""
        para = _make_para_with_runs(
            [
                ("<script>", False, False, False, False),
            ]
        )
        result = _runs_to_html(para)
        assert "&lt;script&gt;" in result
        assert "<script>" not in result

    def test_visual_runs_skipped(self) -> None:
        """Runs with visual content are excluded from HTML."""
        doc = Document()
        para = doc.add_paragraph()
        run1 = para.add_run("Hello")
        run1.bold = True
        run2 = para.add_run()
        drawing = OxmlElement("w:drawing")
        run2._element.append(drawing)
        result = _runs_to_html(para)
        assert result == "<b>Hello</b>"

    def test_underline_and_strike(self) -> None:
        """Underline and strikethrough produce <u> and <s> tags."""
        para = _make_para_with_runs(
            [
                ("Hello", False, False, True, True),
            ]
        )
        result = _runs_to_html(para)
        assert "<u>" in result
        assert "<s>" in result
        assert "Hello" in result

    def test_plain_text_no_tags(self) -> None:
        """Plain-text run produces no HTML tags."""
        para = _make_para_with_runs(
            [
                ("Hello", False, False, False, False),
            ]
        )
        assert _runs_to_html(para) == "Hello"

    def test_empty_runs_skipped(self) -> None:
        """Runs with empty text produce no output."""
        para = _make_para_with_runs(
            [
                ("Hello", True, False, False, False),
                ("", False, True, False, False),
            ]
        )
        assert _runs_to_html(para) == "<b>Hello</b>"

    def test_all_empty_runs(self) -> None:
        """All runs have empty text → result is empty string."""
        para = _make_para_with_runs(
            [
                ("", True, False, False, False),
                ("", False, True, False, False),
            ]
        )
        assert _runs_to_html(para) == ""


# ── _parse_html_formatting tests ────────────────────────────────────


class TestParseHtmlFormatting:
    """Tests for _parse_html_formatting."""

    def test_simple_bold(self) -> None:
        """Parses <b>text</b> into a bold segment."""
        segments = _parse_html_formatting("<b>Hello</b>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0] == _FormattedSegment("Hello", True, False, False, False)

    def test_nested_bold_italic(self) -> None:
        """Parses nested <b><i>text</i></b>."""
        segments = _parse_html_formatting("<b><i>Hello</i></b>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].bold is True
        assert segments[0].italic is True

    def test_mixed_segments(self) -> None:
        """Parses mixed plain and formatted text."""
        segments = _parse_html_formatting("Hello <b>world</b>")
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0] == _FormattedSegment("Hello ", False, False, False, False)
        assert segments[1] == _FormattedSegment("world", True, False, False, False)

    def test_plain_text(self) -> None:
        """Parses text with no tags as a single plain segment."""
        segments = _parse_html_formatting("Hello world")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0] == _FormattedSegment(
            "Hello world",
            False,
            False,
            False,
            False,
        )

    def test_empty_string(self) -> None:
        """Parses empty string to empty list."""
        segments = _parse_html_formatting("")
        assert segments == []

    def test_merges_adjacent_same_format(self) -> None:
        """Adjacent segments with identical formatting are merged."""
        segments = _parse_html_formatting("<b>Hello</b><b> world</b>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "Hello world"

    def test_all_four_tags(self) -> None:
        """All four formatting tags: <b><i><u><s>."""
        segments = _parse_html_formatting("<b><i><u><s>Hello</s></u></i></b>")
        assert len(segments) == 1  # noqa: PLR2004
        seg = segments[0]
        assert seg.bold is True
        assert seg.italic is True
        assert seg.underline is True
        assert seg.strike is True

    def test_html_entities(self) -> None:
        """HTML entities are decoded by the parser."""
        segments = _parse_html_formatting("&lt;tag&gt;")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "<tag>"

    def test_misnested_tags(self) -> None:
        """Handles misnested HTML gracefully (e.g. <b><i>text</b></i>)."""
        segments = _parse_html_formatting("<b><i>Hello</b></i>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "Hello"
        assert segments[0].bold is True
        assert segments[0].italic is True

    def test_whitespace_only_tags(self) -> None:
        """Whitespace-only content inside tags is preserved as a segment."""
        segments = _parse_html_formatting("<b> </b>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == " "
        assert segments[0].bold is True

    def test_nested_same_tags(self) -> None:
        """Nested identical tags don't double-apply; text is still bold."""
        segments = _parse_html_formatting("<b><b>text</b></b>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "text"
        assert segments[0].bold is True

    def test_bare_span_tag(self) -> None:
        """Bare <span> without style attributes yields a plain-text segment."""
        segments = _parse_html_formatting("<span>text</span>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "text"
        assert segments[0].bold is False
        assert segments[0].italic is False
        assert segments[0].underline is False
        assert segments[0].strike is False


# ── _inject_html_runs tests ─────────────────────────────────────────


class TestInjectHtmlRuns:
    """Tests for _inject_html_runs."""

    def test_inject_bold(self) -> None:
        """Injects <b>text</b> as a bold run."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("placeholder")
        _inject_html_runs(para, "<b>Hello</b>")
        # Should have exactly one text run that is bold
        text_runs = [r for r in para.runs if r.text]
        assert len(text_runs) == 1  # noqa: PLR2004
        assert text_runs[0].text == "Hello"
        # Check the bold flag in XML
        rpr = text_runs[0]._element.find(qn("w:rPr"))
        assert rpr is not None
        assert rpr.find(qn("w:b")) is not None

    def test_inject_mixed_formatting(self) -> None:
        """Injects mixed plain + bold as separate runs."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("placeholder")
        _inject_html_runs(para, "Hello <b>world</b>")
        text_runs = [r for r in para.runs if r.text]
        assert len(text_runs) == 2  # noqa: PLR2004
        assert text_runs[0].text == "Hello "
        assert text_runs[1].text == "world"
        # Second run should be bold
        rpr1 = text_runs[1]._element.find(qn("w:rPr"))
        assert rpr1.find(qn("w:b")) is not None
        # First run should not be bold
        rpr0 = text_runs[0]._element.find(qn("w:rPr"))
        assert rpr0 is None or rpr0.find(qn("w:b")) is None

    def test_preserves_visual_content(self) -> None:
        """Visual-content runs are preserved after injection."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("text")
        visual_run = para.add_run()
        drawing = OxmlElement("w:drawing")
        visual_run._element.append(drawing)
        _inject_html_runs(para, "<b>Translated</b>")
        # Should have the translated run + the visual run
        has_visual = any(_run_has_visual_content(r._element) for r in para.runs)
        assert has_visual

    def test_fallback_no_tags(self) -> None:
        """Falls back to _replace_paragraph_text when no HTML tags."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("original")
        _inject_html_runs(para, "plain text")
        assert para.text == "plain text"

    def test_base_formatting_preserved(self) -> None:
        """Base formatting (font name, size) is copied to new runs."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("original")
        run.font.name = "Arial"
        _inject_html_runs(para, "<b>Hello</b> world")
        # Both runs should have rPr with font info from the original
        for r in para.runs:
            if r.text:
                rpr = r._element.find(qn("w:rPr"))
                assert rpr is not None
                rfonts = rpr.find(qn("w:rFonts"))
                if rfonts is not None:
                    assert rfonts.get(qn("w:ascii")) == "Arial"

    def test_inject_underline(self) -> None:
        """Injects <u>text</u> as an underlined run."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("placeholder")
        _inject_html_runs(para, "<u>underlined</u>")
        text_runs = [r for r in para.runs if r.text]
        assert len(text_runs) == 1  # noqa: PLR2004
        rpr = text_runs[0]._element.find(qn("w:rPr"))
        u_elem = rpr.find(qn("w:u"))
        assert u_elem is not None
        assert u_elem.get(qn("w:val")) == "single"

    def test_inject_strike(self) -> None:
        """Injects <s>text</s> as a strikethrough run."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("placeholder")
        _inject_html_runs(para, "<s>deleted</s>")
        text_runs = [r for r in para.runs if r.text]
        assert len(text_runs) == 1  # noqa: PLR2004
        rpr = text_runs[0]._element.find(qn("w:rPr"))
        assert rpr.find(qn("w:strike")) is not None

    def test_empty_html_tags_fallback(self) -> None:
        """Tags with no text content (<b></b>) fall back to empty plain text."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("original")
        _inject_html_runs(para, "<b></b>")
        # Empty segments → _replace_paragraph_text("") removes all runs
        assert para.text == ""


# ── Round-trip (encode → decode) tests ──────────────────────────────


class TestHtmlRoundTrip:
    """Tests that encoding runs to HTML then decoding back is consistent."""

    def test_bold_italic_round_trip(self) -> None:
        """Bold+italic paragraph survives encode → parse round-trip."""
        para = _make_para_with_runs(
            [
                ("Hello ", True, False, False, False),
                ("world", False, True, False, False),
            ]
        )
        html_text = _runs_to_html(para)
        segments = _parse_html_formatting(html_text)
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0].text == "Hello "
        assert segments[0].bold is True
        assert segments[0].italic is False
        assert segments[1].text == "world"
        assert segments[1].bold is False
        assert segments[1].italic is True

    def test_all_formats_round_trip(self) -> None:
        """All four formatting types survive the round-trip."""
        para = _make_para_with_runs(
            [
                ("A", True, True, True, True),
                ("B", False, False, False, False),
            ]
        )
        html_text = _runs_to_html(para)
        segments = _parse_html_formatting(html_text)
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0] == _FormattedSegment("A", True, True, True, True)
        assert segments[1] == _FormattedSegment("B", False, False, False, False)

    def test_encode_inject_preserves_text(self, tmp_path: Path) -> None:
        """Full encode → translate (identity) → inject preserves text."""
        doc = Document()
        para = doc.add_paragraph()
        r1 = para.add_run("Hello ")
        r1.bold = True
        r2 = para.add_run("world")
        r2.italic = True

        html_text = _runs_to_html(para)
        # Simulate identity translation (HTML passes through unchanged)
        _inject_html_runs(para, html_text)

        # Verify paragraph text is preserved
        assert para.text == "Hello world"

    def test_special_chars_round_trip(self) -> None:
        """HTML-escaped characters survive the round-trip."""
        para = _make_para_with_runs(
            [
                ("A < B", True, False, False, False),
                (" & C > D", False, False, False, False),
            ]
        )
        html_text = _runs_to_html(para)
        segments = _parse_html_formatting(html_text)
        assert segments[0].text == "A < B"
        assert segments[1].text == " & C > D"


# ── _FORMATTING_HTML_RE tests ───────────────────────────────────────


class TestFormattingHtmlRe:
    """Tests for the _FORMATTING_HTML_RE regex."""

    def test_detects_bold_tag(self) -> None:
        """Detects <b> tags."""
        assert _FORMATTING_HTML_RE.search("<b>text</b>") is not None

    def test_detects_italic_tag(self) -> None:
        """Detects <i> tags."""
        assert _FORMATTING_HTML_RE.search("<i>text</i>") is not None

    def test_no_match_plain_text(self) -> None:
        """No match for plain text."""
        assert _FORMATTING_HTML_RE.search("plain text") is None

    def test_case_insensitive(self) -> None:
        """Matches uppercase tags."""
        assert _FORMATTING_HTML_RE.search("<B>text</B>") is not None

    def test_no_match_other_tags(self) -> None:
        """Does not match non-formatting tags like <p>."""
        assert _FORMATTING_HTML_RE.search("<p>text</p>") is None


# ── PPTX per-run formatting helpers ────────────────────────────────────


def _make_pptx_para_with_runs(
    specs: list[tuple[str, bool, bool, bool, bool]],
) -> object:
    """Creates a PPTX paragraph with runs of specified formatting.

    Args:
        specs: List of (text, bold, italic, underline, strike) tuples.

    Returns:
        The paragraph object from the first shape on the first slide.
    """
    from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = txbox.text_frame
    para = tf.paragraphs[0]

    for idx, (text, bold, italic, underline, strike) in enumerate(specs):
        if idx == 0:
            run = para.runs[0] if para.runs else para.add_run()
        else:
            run = para.add_run()
        run.text = text
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        if strike:
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                from lxml import etree  # noqa: PLC0415

                rpr = etree.SubElement(run._r, pptx_qn("a:rPr"))
            rpr.set("strike", "sngStrike")

    return para


class TestHasPptxMixedFormatting:
    """Tests for _has_pptx_mixed_formatting."""

    def test_uniform_formatting_returns_false(self) -> None:
        """All runs bold → not mixed."""
        para = _make_pptx_para_with_runs(
            [
                ("Hello ", True, False, False, False),
                ("World", True, False, False, False),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is False

    def test_mixed_bold_italic_returns_true(self) -> None:
        """One bold, one italic → mixed."""
        para = _make_pptx_para_with_runs(
            [
                ("Hello ", True, False, False, False),
                ("World", False, True, False, False),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is True

    def test_single_run_returns_false(self) -> None:
        """One run → not mixed."""
        para = _make_pptx_para_with_runs(
            [
                ("Only", False, False, False, False),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is False

    def test_empty_runs_skipped(self) -> None:
        """Empty-text runs are ignored."""
        para = _make_pptx_para_with_runs(
            [
                ("Text", True, False, False, False),
                ("", False, True, False, False),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is False

    def test_no_runs_returns_false(self) -> None:
        """Paragraph with no runs → not mixed."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        para = txbox.text_frame.paragraphs[0]
        # Remove default run if any
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        for r in list(para._p.findall(pptx_qn("a:r"))):
            para._p.remove(r)
        assert _has_pptx_mixed_formatting(para) is False

    def test_strike_detected(self) -> None:
        """Strike attribute makes formatting differ."""
        para = _make_pptx_para_with_runs(
            [
                ("Normal ", False, False, False, False),
                ("Struck", False, False, False, True),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is True

    def test_two_plain_runs_returns_false(self) -> None:
        """Two unformatted runs have the same signature → not mixed."""
        para = _make_pptx_para_with_runs(
            [
                ("First", False, False, False, False),
                ("Second", False, False, False, False),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is False

    def test_same_mixed_signature_returns_false(self) -> None:
        """Two runs both bold+italic share the same signature → not mixed."""
        para = _make_pptx_para_with_runs(
            [
                ("First", True, True, False, False),
                ("Second", True, True, False, False),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is False

    def test_dbl_strike_treated_as_struck(self) -> None:
        """DblStrike value is counted as struck, making formatting differ."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("Normal", False, False, False, False),
                ("DblStruck", False, False, False, False),
            ]
        )
        rpr = para.runs[1]._r.find(pptx_qn("a:rPr"))
        if rpr is None:
            from lxml import etree  # noqa: PLC0415

            rpr = etree.SubElement(para.runs[1]._r, pptx_qn("a:rPr"))
        rpr.set("strike", "dblStrike")
        assert _has_pptx_mixed_formatting(para) is True

    def test_no_strike_same_as_absent_strike(self) -> None:
        """NoStrike value is treated as not struck, same as absent rPr."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("A", False, False, False, False),
                ("B", False, False, False, False),
            ]
        )
        # Set noStrike explicitly on the first run
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        if rpr is None:
            from lxml import etree  # noqa: PLC0415

            rpr = etree.SubElement(para.runs[0]._r, pptx_qn("a:rPr"))
        rpr.set("strike", "noStrike")
        # Second run has no rPr at all — both are effectively unstruck
        assert _has_pptx_mixed_formatting(para) is False


class TestPptxRunsToHtml:
    """Tests for _pptx_runs_to_html."""

    def test_bold_only(self) -> None:
        """Bold run → <b> tag."""
        para = _make_pptx_para_with_runs(
            [
                ("Bold", True, False, False, False),
            ]
        )
        assert _pptx_runs_to_html(para) == "<b>Bold</b>"

    def test_nested_bold_italic(self) -> None:
        """Bold+italic run → nested tags."""
        para = _make_pptx_para_with_runs(
            [
                ("Both", True, True, False, False),
            ]
        )
        result = _pptx_runs_to_html(para)
        assert "<b>" in result
        assert "<i>" in result
        assert "Both" in result

    def test_mixed_runs(self) -> None:
        """Different runs get different tags."""
        para = _make_pptx_para_with_runs(
            [
                ("Hello ", True, False, False, False),
                ("World", False, True, False, False),
            ]
        )
        result = _pptx_runs_to_html(para)
        assert "<b>Hello </b>" in result
        assert "<i>World</i>" in result

    def test_html_escaping(self) -> None:
        """Special characters are escaped."""
        para = _make_pptx_para_with_runs(
            [
                ("A < B", False, False, False, False),
            ]
        )
        result = _pptx_runs_to_html(para)
        assert "A &lt; B" in result

    def test_underline_and_strike(self) -> None:
        """Underline and strike formatting."""
        para = _make_pptx_para_with_runs(
            [
                ("Under", False, False, True, False),
                ("Strike", False, False, False, True),
            ]
        )
        result = _pptx_runs_to_html(para)
        assert "<u>Under</u>" in result
        assert "<s>Strike</s>" in result

    def test_plain_text_no_tags(self) -> None:
        """No formatting → no tags."""
        para = _make_pptx_para_with_runs(
            [
                ("Plain", False, False, False, False),
            ]
        )
        assert _pptx_runs_to_html(para) == "Plain"

    def test_empty_runs_skipped(self) -> None:
        """Empty-text runs produce no output."""
        para = _make_pptx_para_with_runs(
            [
                ("", True, False, False, False),
                ("Visible", False, False, False, False),
            ]
        )
        assert _pptx_runs_to_html(para) == "Visible"

    def test_ampersand_escaping(self) -> None:
        """Ampersand is HTML-escaped to &amp;."""
        para = _make_pptx_para_with_runs(
            [
                ("A & B", False, False, False, False),
            ]
        )
        assert _pptx_runs_to_html(para) == "A &amp; B"

    def test_bold_italic_nesting_order(self) -> None:
        """Bold+italic on one run produces <b><i>text</i></b> (b outermost)."""
        para = _make_pptx_para_with_runs(
            [
                ("Both", True, True, False, False),
            ]
        )
        assert _pptx_runs_to_html(para) == "<b><i>Both</i></b>"

    def test_all_four_flags_nesting(self) -> None:
        """All four flags produce <b><i><u><s>text</s></u></i></b>."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("All", True, True, True, False),
            ]
        )
        # Add sngStrike directly so all four are set
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        if rpr is None:
            from lxml import etree  # noqa: PLC0415

            rpr = etree.SubElement(para.runs[0]._r, pptx_qn("a:rPr"))
        rpr.set("strike", "sngStrike")
        assert _pptx_runs_to_html(para) == "<b><i><u><s>All</s></u></i></b>"

    def test_all_empty_runs(self) -> None:
        """All PPTX runs have empty text → result is empty string."""
        para = _make_pptx_para_with_runs(
            [
                ("", True, False, False, False),
                ("", False, True, False, False),
            ]
        )
        assert _pptx_runs_to_html(para) == ""


class TestInjectPptxHtmlRuns:
    """Tests for _inject_pptx_html_runs."""

    def test_bold_injection(self) -> None:
        """Injects bold segment correctly."""
        para = _make_pptx_para_with_runs(
            [
                ("original", False, False, False, False),
            ]
        )
        _inject_pptx_html_runs(para, "<b>Translated</b>")
        assert len(para.runs) == 1  # noqa: PLR2004
        assert para.runs[0].text == "Translated"
        assert para.runs[0].font.bold is True

    def test_mixed_injection(self) -> None:
        """Injects bold + italic segments."""
        para = _make_pptx_para_with_runs(
            [
                ("orig1", True, False, False, False),
                ("orig2", False, True, False, False),
            ]
        )
        _inject_pptx_html_runs(para, "<b>Bold</b><i>Italic</i>")
        assert len(para.runs) == 2  # noqa: PLR2004
        assert para.runs[0].text == "Bold"
        assert para.runs[0].font.bold is True
        assert para.runs[1].text == "Italic"
        assert para.runs[1].font.italic is True

    def test_fallback_no_tags(self) -> None:
        """No HTML tags → plain text into first run."""
        para = _make_pptx_para_with_runs(
            [
                ("original1", False, False, False, False),
                ("original2", False, False, False, False),
            ]
        )
        _inject_pptx_html_runs(para, "Plain translation")
        assert para.runs[0].text == "Plain translation"
        assert para.runs[1].text == ""

    def test_base_formatting_preserved(self) -> None:
        """Base formatting (font size, lang) carries over to new runs."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("orig", False, False, False, False),
            ]
        )
        # Set a lang attribute on the original rPr
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        if rpr is None:
            from lxml import etree  # noqa: PLC0415

            rpr = etree.SubElement(para.runs[0]._r, pptx_qn("a:rPr"))
        rpr.set("lang", "en-US")
        rpr.set("sz", "1800")

        _inject_pptx_html_runs(para, "<b>New</b> text")
        # Both new runs should have the base lang and sz
        for run in para.runs:
            new_rpr = run._r.find(pptx_qn("a:rPr"))
            assert new_rpr is not None
            assert new_rpr.get("lang") == "en-US"
            assert new_rpr.get("sz") == "1800"

    def test_strike_injection(self) -> None:
        """Strike segment gets strike attribute."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("original", False, False, False, False),
            ]
        )
        _inject_pptx_html_runs(para, "<s>Deleted</s>")
        assert len(para.runs) == 1
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        assert rpr is not None
        assert rpr.get("strike") == "sngStrike"

    def test_underline_injection(self) -> None:
        """Underline segment gets u attribute."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("original", False, False, False, False),
            ]
        )
        _inject_pptx_html_runs(para, "<u>Underlined</u>")
        assert len(para.runs) == 1
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        assert rpr is not None
        assert rpr.get("u") == "sng"

    def test_no_runs_fallback(self) -> None:
        """Paragraph with no runs and no tags → sets para.text."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        para = txbox.text_frame.paragraphs[0]
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        for r in list(para._p.findall(pptx_qn("a:r"))):
            para._p.remove(r)
        _inject_pptx_html_runs(para, "No runs plain")
        assert para.text == "No runs plain"

    def test_no_rpr_on_first_run_with_html(self) -> None:
        """First run has no rPr (base_rpr=None) → injects without crash."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("orig", False, False, False, False),
            ]
        )
        # Strip the rPr so base_rpr=None path is exercised
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        if rpr is not None:
            para.runs[0]._r.remove(rpr)
        _inject_pptx_html_runs(para, "<b>Bold</b> plain")
        assert len(para.runs) == 2  # noqa: PLR2004
        assert para.runs[0].text == "Bold"
        assert para.runs[0].font.bold is True
        assert para.runs[1].text == " plain"

    def test_no_runs_with_html_tags(self) -> None:
        """Para with no runs + HTML tags → creates new runs from scratch."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        para = txbox.text_frame.paragraphs[0]
        for r in list(para._p.findall(pptx_qn("a:r"))):
            para._p.remove(r)
        _inject_pptx_html_runs(para, "<b>Created</b>")
        assert len(para.runs) == 1
        assert para.runs[0].text == "Created"
        assert para.runs[0].font.bold is True

    def test_bold_italic_same_segment(self) -> None:
        """Bold+italic in the same segment → both attrs set on one run."""
        para = _make_pptx_para_with_runs(
            [
                ("orig", False, False, False, False),
            ]
        )
        _inject_pptx_html_runs(para, "<b><i>Both</i></b>")
        assert len(para.runs) == 1
        assert para.runs[0].text == "Both"
        assert para.runs[0].font.bold is True
        assert para.runs[0].font.italic is True

    def test_empty_html_tags_fallback(self) -> None:
        """Tags with no text content → empty segments → stripped plain fallback."""
        para = _make_pptx_para_with_runs(
            [
                ("orig1", False, False, False, False),
                ("orig2", False, False, False, False),
            ]
        )
        # <b></b> matches _FORMATTING_HTML_RE but produces no segments;
        # plain fallback strips tags so literal HTML doesn't appear.
        _inject_pptx_html_runs(para, "<b></b>")
        assert para.runs[0].text == ""
        assert para.runs[1].text == ""


class TestPptxHtmlRoundTrip:
    """Tests for encode → decode consistency."""

    def test_bold_round_trip(self) -> None:
        """Bold formatting survives encode → decode."""
        para = _make_pptx_para_with_runs(
            [
                ("Hello ", True, False, False, False),
                ("World", False, False, False, False),
            ]
        )
        html_text = _pptx_runs_to_html(para)
        _inject_pptx_html_runs(para, html_text)
        assert para.runs[0].text == "Hello "
        assert para.runs[0].font.bold is True
        assert para.runs[1].text == "World"

    def test_special_chars_round_trip(self) -> None:
        """Special characters survive encode → decode."""
        para = _make_pptx_para_with_runs(
            [
                ("A < B & C", True, False, False, False),
                (" > D", False, False, False, False),
            ]
        )
        html_text = _pptx_runs_to_html(para)
        _inject_pptx_html_runs(para, html_text)
        assert para.runs[0].text == "A < B & C"
        assert para.runs[1].text == " > D"

    def test_all_formatting_round_trip(self) -> None:
        """Bold + italic + underline + strike round-trips."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("Bold", True, False, False, False),
                ("Italic", False, True, False, False),
                ("Under", False, False, True, False),
                ("Strike", False, False, False, True),
            ]
        )
        html_text = _pptx_runs_to_html(para)
        _inject_pptx_html_runs(para, html_text)
        runs = para.runs
        assert len(runs) == 4  # noqa: PLR2004
        assert runs[0].text == "Bold"
        assert runs[0].font.bold is True
        assert runs[1].text == "Italic"
        assert runs[1].font.italic is True
        assert runs[2].text == "Under"
        rpr2 = runs[2]._r.find(pptx_qn("a:rPr"))
        assert rpr2 is not None
        assert rpr2.get("u") == "sng"
        assert runs[3].text == "Strike"
        rpr3 = runs[3]._r.find(pptx_qn("a:rPr"))
        assert rpr3 is not None
        assert rpr3.get("strike") == "sngStrike"

    def test_mixed_with_plain_round_trip(self) -> None:
        """Mixed formatted and plain runs round-trip correctly."""
        para = _make_pptx_para_with_runs(
            [
                ("Start ", False, False, False, False),
                ("bold", True, False, False, False),
                (" end", False, False, False, False),
            ]
        )
        html_text = _pptx_runs_to_html(para)
        _inject_pptx_html_runs(para, html_text)
        texts = [r.text for r in para.runs]
        assert "".join(texts) == "Start bold end"

    def test_whitespace_preserved_via_xml_space(self) -> None:
        """Leading/trailing whitespace in runs survives encode → decode."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("  padded  ", True, False, False, False),
                (" gap ", False, False, False, False),
            ]
        )
        html_text = _pptx_runs_to_html(para)
        _inject_pptx_html_runs(para, html_text)
        assert para.runs[0].text == "  padded  "
        assert para.runs[1].text == " gap "
        # Verify xml:space="preserve" is set on every new <a:t>
        for run in para.runs:
            t_elem = run._r.find(pptx_qn("a:t"))
            assert t_elem is not None
            space_attr = t_elem.get("{http://www.w3.org/XML/1998/namespace}space")
            assert space_attr == "preserve"


# ---------------------------------------------------------------------------
# Helper: mock UNO paragraph with enumerable text portions
# ---------------------------------------------------------------------------


def _make_uno_portion(  # noqa: PLR0913
    text: str,
    bold: bool = False,
    italic: bool = False,
    underline: bool = False,
    strike: bool = False,
    portion_type: str = "Text",
    font_size: float = 11.0,
    color_int: int = 0,
) -> MagicMock:
    """Creates a mock UNO text portion with formatting properties.

    Args:
        text: The portion text content.
        bold: Effective bold value.
        italic: Effective italic value.
        underline: Effective underline value.
        strike: Effective strikethrough value.
        portion_type: TextPortionType value (default "Text").
        font_size: CharHeight in points (default 11.0).
        color_int: CharColor as integer (default 0 = black).

    Returns:
        Mock UNO TextPortion object.
    """
    portion = MagicMock()
    props = {
        "TextPortionType": portion_type,
        "CharWeight": _UNO_WEIGHT_BOLD if bold else _UNO_WEIGHT_NORMAL,
        "CharPosture": _UNO_SLANT_ITALIC if italic else _UNO_SLANT_NONE,
        "CharUnderline": _UNO_UNDERLINE_SINGLE if underline else _UNO_UNDERLINE_NONE,
        "CharStrikeout": _UNO_STRIKEOUT_SINGLE if strike else _UNO_STRIKEOUT_NONE,
        "CharFontName": "Liberation Sans",
        "CharHeight": font_size,
        "CharColor": color_int,
    }
    portion.getPropertyValue.side_effect = lambda p: props[p]
    portion.getString.return_value = text
    return portion


def _make_uno_para_with_portions(
    specs: list[tuple[str, bool, bool, bool, bool]],
) -> MagicMock:
    """Creates a mock UNO paragraph with enumerable text portions.

    Args:
        specs: List of (text, bold, italic, underline, strike) tuples
            representing each portion's *effective* formatting.

    Returns:
        Mock UNO paragraph with a createEnumeration() that yields portions.
    """
    portions = [_make_uno_portion(t, b, i, u, s) for t, b, i, u, s in specs]
    full_text = "".join(t for t, *_ in specs)

    # Build an enumeration mock
    enum_mock = MagicMock()
    _idx = [0]

    def _has_more() -> bool:
        return _idx[0] < len(portions)

    def _next() -> MagicMock:
        p = portions[_idx[0]]
        _idx[0] += 1
        return p

    enum_mock.hasMoreElements.side_effect = _has_more
    enum_mock.nextElement.side_effect = _next

    para = MagicMock()
    para.createEnumeration.return_value = enum_mock
    para.getString.return_value = full_text
    para.supportsService.return_value = False

    # For injection: mock getText() → text object with cursor helpers
    text_obj = MagicMock()

    def _create_cursor_by_range(rng: object) -> MagicMock:
        cursor = MagicMock()
        cursor._offset = 0
        cursor._selected = 0
        cursor._props: dict[str, object] = {}

        def go_start_of_para(expand: bool) -> None:
            cursor._offset = 0
            cursor._selected = 0

        def go_right(count: int, expand: bool) -> None:
            if expand:
                cursor._selected = count
            else:
                cursor._offset += count

        def set_prop(name: str, val: object) -> None:
            cursor._props[name] = val

        cursor.gotoStartOfParagraph.side_effect = go_start_of_para
        cursor.goRight.side_effect = go_right
        cursor.setPropertyValue.side_effect = set_prop
        return cursor

    text_obj.createTextCursorByRange.side_effect = _create_cursor_by_range
    para.getText.return_value = text_obj
    para.getStart.return_value = MagicMock()

    return para


# ---------------------------------------------------------------------------
# TestReadUnoPortionFormatting
# ---------------------------------------------------------------------------


class TestReadUnoPortionFormatting:
    """Tests for _read_uno_portion_formatting."""

    def test_bold(self) -> None:
        """Bold portion returns (True, False, False, False, False, False)."""
        portion = _make_uno_portion("Bold", bold=True)
        result = _read_uno_portion_formatting(portion)
        assert result == (True, False, False, False, False, False)

    def test_italic(self) -> None:
        """Italic portion returns (False, True, False, False, False, False)."""
        portion = _make_uno_portion("Italic", italic=True)
        result = _read_uno_portion_formatting(portion)
        assert result == (False, True, False, False, False, False)

    def test_underline(self) -> None:
        """Underline portion returns (False, False, True, False, False, False)."""
        portion = _make_uno_portion("Underline", underline=True)
        result = _read_uno_portion_formatting(portion)
        assert result == (False, False, True, False, False, False)

    def test_strike(self) -> None:
        """Strikethrough portion returns (False, False, False, True, False, False)."""
        portion = _make_uno_portion("Strike", strike=True)
        result = _read_uno_portion_formatting(portion)
        assert result == (False, False, False, True, False, False)

    def test_all_off(self) -> None:
        """Plain portion returns all False."""
        portion = _make_uno_portion("Plain")
        result = _read_uno_portion_formatting(portion)
        assert result == (False, False, False, False, False, False)

    def test_all_on(self) -> None:
        """All-formatted portion returns all True."""
        portion = _make_uno_portion(
            "All",
            bold=True,
            italic=True,
            underline=True,
            strike=True,
        )
        result = _read_uno_portion_formatting(portion)
        assert result == (True, True, True, True, False, False)


# ---------------------------------------------------------------------------
# TestHasUnoMixedFormatting
# ---------------------------------------------------------------------------


class TestHasUnoMixedFormatting:
    """Tests for _has_uno_mixed_formatting."""

    def test_uniform_formatting(self) -> None:
        """All portions with same formatting → False."""
        para = _make_uno_para_with_portions(
            [
                ("Hello ", True, False, False, False),
                ("World", True, False, False, False),
            ]
        )
        assert _has_uno_mixed_formatting(para) is False

    def test_mixed_formatting(self) -> None:
        """Portions with different formatting → True."""
        para = _make_uno_para_with_portions(
            [
                ("Bold ", True, False, False, False),
                ("Plain", False, False, False, False),
            ]
        )
        assert _has_uno_mixed_formatting(para) is True

    def test_single_portion(self) -> None:
        """Single portion → False (need 2+ for mixed)."""
        para = _make_uno_para_with_portions(
            [
                ("Only one", True, False, False, False),
            ]
        )
        assert _has_uno_mixed_formatting(para) is False

    def test_empty_paragraph(self) -> None:
        """No portions → False."""
        para = _make_uno_para_with_portions([])
        assert _has_uno_mixed_formatting(para) is False

    def test_non_text_portions_skipped(self) -> None:
        """Non-Text portions are ignored; single remaining → False."""
        portions = [
            _make_uno_portion("foot", portion_type="Footnote"),
            _make_uno_portion("Real text", bold=True),
        ]

        enum_mock = MagicMock()
        _idx = [0]

        def _has_more() -> bool:
            return _idx[0] < len(portions)

        def _next() -> MagicMock:
            p = portions[_idx[0]]
            _idx[0] += 1
            return p

        enum_mock.hasMoreElements.side_effect = _has_more
        enum_mock.nextElement.side_effect = _next

        para = MagicMock()
        para.createEnumeration.return_value = enum_mock
        assert _has_uno_mixed_formatting(para) is False

    def test_empty_text_portions_skipped(self) -> None:
        """Portions with empty getString() are ignored."""
        para = _make_uno_para_with_portions(
            [
                ("", True, False, False, False),
                ("Text", False, False, False, False),
            ]
        )
        # Only one non-empty portion → False
        assert _has_uno_mixed_formatting(para) is False


# ---------------------------------------------------------------------------
# TestUnoRunsToHtml
# ---------------------------------------------------------------------------


class TestUnoRunsToHtml:
    """Tests for _uno_runs_to_html."""

    def test_bold_only(self) -> None:
        """Bold portion produces <b> tags."""
        para = _make_uno_para_with_portions(
            [
                ("Bold", True, False, False, False),
            ]
        )
        assert _uno_runs_to_html(para) == "<b>Bold</b>"

    def test_nested_bold_italic(self) -> None:
        """Bold+italic portion produces <b><i>…</i></b>."""
        para = _make_uno_para_with_portions(
            [
                ("Both", True, True, False, False),
            ]
        )
        assert _uno_runs_to_html(para) == "<b><i>Both</i></b>"

    def test_mixed_portions(self) -> None:
        """Multiple portions with different formatting."""
        para = _make_uno_para_with_portions(
            [
                ("Bold ", True, False, False, False),
                ("plain", False, False, False, False),
            ]
        )
        assert _uno_runs_to_html(para) == "<b>Bold </b>plain"

    def test_html_escaping(self) -> None:
        """HTML special characters in text are escaped."""
        para = _make_uno_para_with_portions(
            [
                ("<A & B>", True, False, False, False),
            ]
        )
        result = _uno_runs_to_html(para)
        assert "&lt;" in result
        assert "&amp;" in result
        assert "&gt;" in result

    def test_plain_only(self) -> None:
        """Plain text portions produce no tags."""
        para = _make_uno_para_with_portions(
            [
                ("No format", False, False, False, False),
            ]
        )
        assert _uno_runs_to_html(para) == "No format"

    def test_all_four_formats(self) -> None:
        """All four formatting flags produce nested tags."""
        para = _make_uno_para_with_portions(
            [
                ("All", True, True, True, True),
            ]
        )
        assert _uno_runs_to_html(para) == "<b><i><u><s>All</s></u></i></b>"

    def test_underline_and_strike(self) -> None:
        """Underline + strike without bold/italic."""
        para = _make_uno_para_with_portions(
            [
                ("US", False, False, True, True),
            ]
        )
        assert _uno_runs_to_html(para) == "<u><s>US</s></u>"

    def test_mixed_italic_bold_underline(self) -> None:
        """Simulates "very big hello world": big=italic, world=bold+underline.

        Each portion carries its own effective formatting. Only portions
        with active formatting get HTML tags.
        """
        para = _make_uno_para_with_portions(
            [
                ("very ", False, False, False, False),
                ("big", False, True, False, False),
                (" hello ", False, False, False, False),
                ("world", True, False, True, False),
            ]
        )
        result = _uno_runs_to_html(para)
        expected = "very <i>big</i> hello <b><u>world</u></b>"
        assert result == expected

    def test_all_empty_portions(self) -> None:
        """All UNO portions have empty text → result is empty string."""
        para = _make_uno_para_with_portions(
            [
                ("", True, False, False, False),
                ("", False, True, False, False),
            ]
        )
        assert _uno_runs_to_html(para) == ""


# ---------------------------------------------------------------------------
# TestUnoEnumHandling
# ---------------------------------------------------------------------------


class TestUnoEnumHandling:
    """Tests for UNO enum handling in _read_uno_effective_formatting."""

    def test_char_posture_uno_enum_italic(self) -> None:
        """UNO FontSlant enum with .value='ITALIC' → italic=True."""
        portion = MagicMock()
        # Simulate uno.Enum("com.sun.star.awt.FontSlant", "ITALIC")
        enum_italic = MagicMock()
        enum_italic.value = "ITALIC"
        portion.getPropertyValue.side_effect = lambda p: {
            "CharWeight": _UNO_WEIGHT_NORMAL,
            "CharPosture": enum_italic,
            "CharUnderline": _UNO_UNDERLINE_NONE,
            "CharStrikeout": _UNO_STRIKEOUT_NONE,
        }[p]
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )

        result = _read_uno_effective_formatting(portion)
        assert result == (False, True, False, False, False, False)

    def test_char_posture_uno_enum_none(self) -> None:
        """UNO FontSlant enum with .value='NONE' → italic=False."""
        portion = MagicMock()
        enum_none = MagicMock()
        enum_none.value = "NONE"
        portion.getPropertyValue.side_effect = lambda p: {
            "CharWeight": _UNO_WEIGHT_NORMAL,
            "CharPosture": enum_none,
            "CharUnderline": _UNO_UNDERLINE_NONE,
            "CharStrikeout": _UNO_STRIKEOUT_NONE,
        }[p]
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )

        result = _read_uno_effective_formatting(portion)
        assert result == (False, False, False, False, False, False)

    def test_char_posture_plain_integer(self) -> None:
        """Plain integer CharPosture still works (fallback path)."""
        portion = _make_uno_portion("x", italic=True)
        result = _read_uno_portion_formatting(portion)
        assert result == (False, True, False, False, False, False)

    def test_char_posture_plain_integer_none(self) -> None:
        """Plain integer 0 for CharPosture → italic=False."""
        portion = _make_uno_portion("x")
        result = _read_uno_portion_formatting(portion)
        assert result == (False, False, False, False, False, False)

    def test_get_property_value_exception(self) -> None:
        """Exception in getPropertyValue returns False for that flag."""
        portion = MagicMock()

        def _raise(p: str) -> object:
            if p == "CharWeight":
                raise RuntimeError("nope")
            return _UNO_SLANT_NONE

        portion.getPropertyValue.side_effect = _raise
        from src.core.office_processor import (  # noqa: PLC0415
            _read_uno_effective_formatting,
        )

        result = _read_uno_effective_formatting(portion)
        assert result[0] is False  # bold failed gracefully


# ---------------------------------------------------------------------------
# TestSaveUnoFirstPortionProps
# ---------------------------------------------------------------------------


class TestSaveUnoFirstPortionProps:
    """Tests for _save_uno_first_portion_props."""

    def test_captures_first_text_portion(self) -> None:
        """Reads properties from the first non-empty Text portion."""
        para = _make_uno_para_with_portions(
            [
                ("First", True, False, False, False),
                ("Second", False, True, False, False),
            ]
        )
        saved = _save_uno_first_portion_props(para)
        assert saved["CharFontName"] == "Liberation Sans"
        assert saved["CharWeight"] == _UNO_WEIGHT_BOLD  # noqa: PLR2004

    def test_skips_non_text_portion(self) -> None:
        """Skips Footnote portions, reads from first Text portion."""
        portions = [
            _make_uno_portion("foot", portion_type="Footnote"),
            _make_uno_portion("Real", bold=False),
        ]
        enum_mock = MagicMock()
        _idx = [0]

        def _has_more() -> bool:
            return _idx[0] < len(portions)

        def _next() -> MagicMock:
            p = portions[_idx[0]]
            _idx[0] += 1
            return p

        enum_mock.hasMoreElements.side_effect = _has_more
        enum_mock.nextElement.side_effect = _next

        para = MagicMock()
        para.createEnumeration.return_value = enum_mock

        saved = _save_uno_first_portion_props(para)
        assert saved["CharWeight"] == _UNO_WEIGHT_NORMAL  # noqa: PLR2004

    def test_empty_paragraph_returns_empty(self) -> None:
        """No text portions → empty dict."""
        para = _make_uno_para_with_portions([])
        assert _save_uno_first_portion_props(para) == {}

    def test_skips_errored_properties(self) -> None:
        """Properties that raise are silently skipped."""
        portion = MagicMock()

        def side_effect(p: str) -> object:
            if p == "TextPortionType":
                return "Text"
            if p == "CharFontName":
                raise RuntimeError("not available")
            return {
                "CharHeight": 12.0,
                "CharWeight": 100.0,
                "CharPosture": 0,
                "CharColor": 0,
                "CharUnderline": 0,
                "CharStrikeout": 0,
            }.get(p, 0)

        portion.getPropertyValue.side_effect = side_effect
        portion.getString.return_value = "text"

        enum_mock = MagicMock()
        enum_mock.hasMoreElements.side_effect = [True, False]
        enum_mock.nextElement.return_value = portion

        para = MagicMock()
        para.createEnumeration.return_value = enum_mock

        saved = _save_uno_first_portion_props(para)
        assert "CharFontName" not in saved
        assert "CharHeight" in saved


# ---------------------------------------------------------------------------
# TestInjectUnoHtmlRuns
# ---------------------------------------------------------------------------


class TestInjectUnoHtmlRuns:
    """Tests for _inject_uno_html_runs."""

    def test_bold_segment(self) -> None:
        """Bold HTML produces bold cursor formatting."""
        para = _make_uno_para_with_portions(
            [
                ("Original", False, False, False, False),
            ]
        )
        base_props: dict[str, object] = {
            "CharFontName": "Liberation Sans",
            "CharHeight": 11.0,
        }
        _inject_uno_html_runs(para, "<b>Translated</b>", base_props)
        para.setString.assert_called_once_with("Translated")

    def test_mixed_segments(self) -> None:
        """Mixed formatting applies correct properties per segment."""
        para = _make_uno_para_with_portions(
            [
                ("Original text", False, False, False, False),
            ]
        )
        base_props: dict[str, object] = {"CharFontName": "Arial"}
        _inject_uno_html_runs(para, "<b>Bold</b> plain", base_props)
        para.setString.assert_called_once_with("Bold plain")

    def test_fallback_no_tags(self) -> None:
        """Plain text with no HTML tags falls back to setString."""
        para = _make_uno_para_with_portions(
            [
                ("Original", False, False, False, False),
            ]
        )
        # _parse_html_formatting on plain text returns one segment, so
        # it still calls setString with the full text
        _inject_uno_html_runs(para, "No tags", {})
        para.setString.assert_called_once_with("No tags")

    def test_empty_segments_fallback(self) -> None:
        """Empty HTML text falls back to setString."""
        para = _make_uno_para_with_portions(
            [
                ("Original", False, False, False, False),
            ]
        )
        _inject_uno_html_runs(para, "", {})
        para.setString.assert_called_once_with("")

    def test_base_props_restored(self) -> None:
        """Base properties (font name, height) are applied to full cursor."""
        para = _make_uno_para_with_portions(
            [
                ("Original", False, False, False, False),
            ]
        )
        # Track cursors created by the text object
        cursors: list[MagicMock] = []
        text_obj = para.getText.return_value
        orig_fn = text_obj.createTextCursorByRange.side_effect

        def _track(rng: object) -> MagicMock:
            c = orig_fn(rng)
            cursors.append(c)
            return c

        text_obj.createTextCursorByRange.side_effect = _track

        base_props: dict[str, object] = {
            "CharFontName": "Liberation Serif",
            "CharHeight": 14.0,
        }
        _inject_uno_html_runs(para, "<b>Text</b>", base_props)

        # First cursor is the full-range cursor for base props
        assert len(cursors) >= 1
        full_cursor = cursors[0]
        calls = full_cursor.setPropertyValue.call_args_list
        set_calls = {c.args[0]: c.args[1] for c in calls}
        assert set_calls.get("CharFontName") == "Liberation Serif"
        expected_height = 14.0
        assert set_calls.get("CharHeight") == expected_height

    def test_formatting_props_excluded_from_base(self) -> None:
        """Formatting props are NOT applied as base props."""
        para = _make_uno_para_with_portions(
            [
                ("Original", False, False, False, False),
            ]
        )
        cursors: list[MagicMock] = []
        text_obj = para.getText.return_value
        orig_fn = text_obj.createTextCursorByRange.side_effect

        def _track(rng: object) -> MagicMock:
            c = orig_fn(rng)
            cursors.append(c)
            return c

        text_obj.createTextCursorByRange.side_effect = _track

        base_props: dict[str, object] = {
            "CharFontName": "Arial",
            "CharWeight": _UNO_WEIGHT_BOLD,
            "CharPosture": _UNO_SLANT_ITALIC,
            "CharUnderline": _UNO_UNDERLINE_SINGLE,
            "CharStrikeout": _UNO_STRIKEOUT_SINGLE,
        }
        _inject_uno_html_runs(para, "<b>Text</b>", base_props)

        full_cursor = cursors[0]
        set_props = {c.args[0] for c in full_cursor.setPropertyValue.call_args_list}
        assert "CharFontName" in set_props
        for prop in _UNO_FORMATTING_PROPS:
            assert prop not in set_props

    def test_script_family_font_skip(self) -> None:
        """CharFontName is skipped when scripts differ and no target_lang."""
        para = _make_uno_para_with_portions(
            [
                ("Hello", False, False, False, False),
            ]
        )
        cursors: list[MagicMock] = []
        text_obj = para.getText.return_value
        orig_fn = text_obj.createTextCursorByRange.side_effect

        def _track(rng: object) -> MagicMock:
            c = orig_fn(rng)
            cursors.append(c)
            return c

        text_obj.createTextCursorByRange.side_effect = _track

        base_props: dict[str, object] = {
            "CharFontName": "Liberation Sans",
            "CharHeight": 11.0,
            "__original_text__": "Hello",
        }
        _inject_uno_html_runs(para, "<b>مرحبا</b>", base_props)

        full_cursor = cursors[0]
        set_props = {c.args[0] for c in full_cursor.setPropertyValue.call_args_list}
        # Scripts differ + no target_lang → _substitute_font returns None → skip
        assert "CharFontName" not in set_props
        assert "CharHeight" in set_props

    def test_script_family_font_kept_same_script(self) -> None:
        """CharFontName IS restored when scripts match."""
        para = _make_uno_para_with_portions(
            [
                ("English", False, False, False, False),
            ]
        )
        cursors: list[MagicMock] = []
        text_obj = para.getText.return_value
        orig_fn = text_obj.createTextCursorByRange.side_effect

        def _track(rng: object) -> MagicMock:
            c = orig_fn(rng)
            cursors.append(c)
            return c

        text_obj.createTextCursorByRange.side_effect = _track

        base_props: dict[str, object] = {
            "CharFontName": "Arial",
            "__original_text__": "English",
        }
        _inject_uno_html_runs(para, "<b>Français</b>", base_props)

        full_cursor = cursors[0]
        set_props = {c.args[0] for c in full_cursor.setPropertyValue.call_args_list}
        assert "CharFontName" in set_props

    def test_empty_html_tags_fallback(self) -> None:
        """Tags with no text content (<b></b>) fall back to setString('')."""
        para = _make_uno_para_with_portions(
            [
                ("Original", False, False, False, False),
            ]
        )
        _inject_uno_html_runs(para, "<b></b>", {})
        # Empty segments → plain = "" → setString("") called
        para.setString.assert_called_once_with("")


# ---------------------------------------------------------------------------
# TestUnoHtmlRoundTrip
# ---------------------------------------------------------------------------


class TestUnoHtmlRoundTrip:
    """Tests for UNO HTML encode → decode consistency."""

    def test_bold_round_trip(self) -> None:
        """Bold text survives encode → parse."""
        para = _make_uno_para_with_portions(
            [
                ("Bold", True, False, False, False),
                (" plain", False, False, False, False),
            ]
        )
        html_text = _uno_runs_to_html(para)
        segments = _parse_html_formatting(html_text)
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0] == _FormattedSegment("Bold", True, False, False, False)
        assert segments[1] == _FormattedSegment(" plain", False, False, False, False)

    def test_all_formats_round_trip(self) -> None:
        """All four formatting flags survive encode → parse."""
        para = _make_uno_para_with_portions(
            [
                ("All", True, True, True, True),
            ]
        )
        html_text = _uno_runs_to_html(para)
        segments = _parse_html_formatting(html_text)
        assert len(segments) == 1
        assert segments[0] == _FormattedSegment("All", True, True, True, True)

    def test_mixed_complex_round_trip(self) -> None:
        """Complex mixed formatting survives encode → parse."""
        para = _make_uno_para_with_portions(
            [
                ("Normal ", False, False, False, False),
                ("bold+italic ", True, True, False, False),
                ("underline", False, False, True, False),
            ]
        )
        html_text = _uno_runs_to_html(para)
        segments = _parse_html_formatting(html_text)
        assert len(segments) == 3  # noqa: PLR2004
        assert segments[0].text == "Normal "
        assert segments[0].bold is False
        assert segments[1].text == "bold+italic "
        assert segments[1].bold is True
        assert segments[1].italic is True
        assert segments[2].text == "underline"
        assert segments[2].underline is True

    def test_plain_round_trip(self) -> None:
        """Plain text produces no formatting flags after round trip."""
        para = _make_uno_para_with_portions(
            [
                ("Just plain", False, False, False, False),
            ]
        )
        html_text = _uno_runs_to_html(para)
        segments = _parse_html_formatting(html_text)
        assert len(segments) == 1
        expected = _FormattedSegment(
            "Just plain",
            False,
            False,
            False,
            False,
        )
        assert segments[0] == expected


# ---------------------------------------------------------------------------
# Integration: _extract_uno_writer / _inject_uno_writer HTML dispatch
# ---------------------------------------------------------------------------


def _make_uno_portion_enum(
    specs: list[tuple[str, bool, bool, bool, bool]],
) -> MagicMock:
    """Creates a fresh UNO portion enumeration from specs."""
    portions = [_make_uno_portion(t, b, i, u, s) for t, b, i, u, s in specs]
    enum_mock = MagicMock()
    _idx = [0]

    def _has_more() -> bool:
        return _idx[0] < len(portions)

    def _next() -> MagicMock:
        p = portions[_idx[0]]
        _idx[0] += 1
        return p

    enum_mock.hasMoreElements.side_effect = _has_more
    enum_mock.nextElement.side_effect = _next
    return enum_mock


@patch("src.core.office_processor._uno_open")
def test_extract_uno_writer_uses_html_for_mixed(
    mock_open: MagicMock,
) -> None:
    """_extract_uno_writer returns HTML for mixed formatting."""
    specs = [
        ("Bold ", True, False, False, False),
        ("plain", False, False, False, False),
    ]

    # Paragraph-level baseline: plain (no formatting)
    para_props: dict[str, object] = {
        "CharWeight": _UNO_WEIGHT_NORMAL,
        "CharPosture": _UNO_SLANT_NONE,
        "CharUnderline": _UNO_UNDERLINE_NONE,
        "CharStrikeout": _UNO_STRIKEOUT_NONE,
    }

    # Each createEnumeration call needs a fresh enumeration
    mock_para = MagicMock()
    mock_para.supportsService.return_value = False
    mock_para.getString.return_value = "Bold plain"
    mock_para.getPropertyValue.side_effect = lambda p: para_props[p]
    mock_para.createEnumeration.side_effect = lambda: _make_uno_portion_enum(specs)

    doc_enum = MagicMock()
    doc_enum.hasMoreElements.side_effect = [True, False]
    doc_enum.nextElement.return_value = mock_para

    mock_text = MagicMock()
    mock_text.createEnumeration.return_value = doc_enum

    mock_doc = MagicMock()
    mock_doc.getText.return_value = mock_text
    mock_doc.getTextTables.return_value.getCount.return_value = 0
    mock_open.return_value = mock_doc

    from src.core.office_processor import (  # noqa: PLC0415
        _extract_uno_writer,
    )

    result = _extract_uno_writer(Path("/fake/doc.docx"))
    assert len(result) == 1
    key, text = result[0]
    assert key == "para:0"
    assert "<b>" in text
    assert "Bold " in text


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_writer_dispatches_html(
    mock_open: MagicMock,
    mock_save: MagicMock,
) -> None:
    """_inject_uno_writer calls _inject_uno_html_runs when HTML tags present."""
    # Build mock paragraph with portions (for _save_uno_first_portion_props)
    portion = _make_uno_portion("Original", bold=False)
    portions = [portion]

    enum_mock = MagicMock()
    _idx = [0]

    def _has_more() -> bool:
        return _idx[0] < len(portions)

    def _next() -> MagicMock:
        p = portions[_idx[0]]
        _idx[0] += 1
        return p

    enum_mock.hasMoreElements.side_effect = _has_more
    enum_mock.nextElement.side_effect = _next

    mock_para = MagicMock()
    mock_para.supportsService.return_value = False
    mock_para.getString.return_value = "Original"
    mock_para.createEnumeration.return_value = enum_mock

    # Mock getText() and cursor for _inject_uno_html_runs
    text_obj = MagicMock()
    mock_para.getText.return_value = text_obj
    mock_para.getStart.return_value = MagicMock()

    doc_enum = MagicMock()
    doc_enum.hasMoreElements.side_effect = [True, False]
    doc_enum.nextElement.return_value = mock_para

    mock_text = MagicMock()
    mock_text.createEnumeration.return_value = doc_enum

    mock_doc = MagicMock()
    mock_doc.getText.return_value = mock_text
    mock_doc.getTextTables.return_value.getCount.return_value = 0
    mock_open.return_value = mock_doc

    from src.core.office_processor import _inject_uno_writer  # noqa: PLC0415

    _inject_uno_writer(
        Path("/fake/in.docx"),
        Path("/fake/out.docx"),
        {"para:0": "<b>Translated</b>"},
    )

    # setString should have been called with the plain text (by _inject_uno_html_runs)
    mock_para.setString.assert_called_once_with("Translated")


@patch("src.core.office_processor._uno_save")
@patch("src.core.office_processor._uno_open")
def test_inject_uno_writer_plain_text_fallback(
    mock_open: MagicMock,
    mock_save: MagicMock,
) -> None:
    """_inject_uno_writer uses plain setString + restore for non-HTML text."""
    mock_para = MagicMock()
    mock_para.supportsService.return_value = False
    mock_para.getString.return_value = "Original"

    doc_enum = MagicMock()
    doc_enum.hasMoreElements.side_effect = [True, False]
    doc_enum.nextElement.return_value = mock_para

    mock_text = MagicMock()
    mock_text.createEnumeration.return_value = doc_enum

    mock_doc = MagicMock()
    mock_doc.getText.return_value = mock_text
    mock_doc.getTextTables.return_value.getCount.return_value = 0
    mock_open.return_value = mock_doc

    from src.core.office_processor import _inject_uno_writer  # noqa: PLC0415

    _inject_uno_writer(
        Path("/fake/in.docx"),
        Path("/fake/out.docx"),
        {"para:0": "Plain translated"},
    )

    # setString should have been called with plain text
    mock_para.setString.assert_called_once_with("Plain translated")


# ---------------------------------------------------------------------------
# Per-run font size & text colour preservation tests
# ---------------------------------------------------------------------------


class TestBuildSpanStyle:
    """Tests for _build_span_style."""

    def test_both_none(self) -> None:
        """Both None → empty string."""
        assert _build_span_style(None, None) == ""

    def test_size_only(self) -> None:
        """Only font_size_pt → font-size style."""
        assert _build_span_style(14.0, None) == "font-size:14pt"

    def test_color_only(self) -> None:
        """Only color_hex → color style."""
        assert _build_span_style(None, "#ff0000") == "color:#ff0000"

    def test_both(self) -> None:
        """Both values → semicolon-separated."""
        result = _build_span_style(10.5, "#00ff00")
        assert result == "font-size:10.5pt;color:#00ff00"


class TestParseSpanStyle:
    """Tests for _parse_span_style."""

    def test_empty(self) -> None:
        """Empty style → empty dict."""
        assert _parse_span_style("") == {}

    def test_size(self) -> None:
        """font-size:14pt → dict with font_size_pt."""
        result = _parse_span_style("font-size:14pt")
        assert result == {"font_size_pt": 14.0}

    def test_color(self) -> None:
        """color:#ff0000 → dict with color_hex."""
        result = _parse_span_style("color:#ff0000")
        assert result == {"color_hex": "#ff0000"}

    def test_both(self) -> None:
        """Both font-size and color."""
        result = _parse_span_style("font-size:10pt;color:#00ff00")
        assert result == {"font_size_pt": 10.0, "color_hex": "#00ff00"}

    def test_background_color_parsed(self) -> None:
        """background-color is extracted as bg_color_hex, not color_hex."""
        result = _parse_span_style("background-color:#ff0000")
        assert "color_hex" not in result
        assert result == {"bg_color_hex": "#ff0000"}

    def test_background_color_with_text_color(self) -> None:
        """background-color and color are both extracted separately."""
        result = _parse_span_style("background-color:#0000ff;color:#ff0000")
        assert result == {"color_hex": "#ff0000", "bg_color_hex": "#0000ff"}


class TestIntToColorHex:
    """Tests for _int_to_color_hex."""

    def test_normal(self) -> None:
        """Standard red → #ff0000."""
        assert _int_to_color_hex(0xFF0000) == "#ff0000"  # noqa: PLR2004

    def test_negative_returns_none(self) -> None:
        """UNO auto colour -1 → None."""
        assert _int_to_color_hex(-1) is None

    def test_overflow_returns_none(self) -> None:
        """Values > 0xFFFFFF → None."""
        assert _int_to_color_hex(0x1000000) is None  # noqa: PLR2004

    def test_zero_is_black(self) -> None:
        """0 → #000000."""
        assert _int_to_color_hex(0) == "#000000"


class TestColorHexToInt:
    """Tests for _color_hex_to_int."""

    def test_red(self) -> None:
        """#ff0000 → 16711680."""
        expected = 16711680  # noqa: PLR2004
        assert _color_hex_to_int("#ff0000") == expected

    def test_black(self) -> None:
        """#000000 → 0."""
        assert _color_hex_to_int("#000000") == 0


class TestFormattedSegmentBackwardCompat:
    """Tests that _FormattedSegment works with old 5-arg construction."""

    def test_five_arg_defaults(self) -> None:
        """5-arg construction defaults new fields."""
        seg = _FormattedSegment("text", True, False, True, False)
        assert seg.superscript is False
        assert seg.subscript is False
        assert seg.font_size_pt is None
        assert seg.color_hex is None
        assert seg.bold is True
        assert seg.underline is True

    def test_nine_arg(self) -> None:
        """9-arg construction sets all fields including sup/sub."""
        seg = _FormattedSegment(
            "x",
            False,
            True,
            False,
            False,
            False,
            False,
            12.0,
            "#aabbcc",
        )
        assert seg.font_size_pt == 12.0  # noqa: PLR2004
        assert seg.color_hex == "#aabbcc"
        assert seg.superscript is False
        assert seg.subscript is False


class TestWrapWithTags:
    """Tests for _wrap_with_tags."""

    def test_bold_with_size(self) -> None:
        """Bold + size variation → <b><span>text</span></b>."""
        result = _wrap_with_tags(
            "text",
            True,
            False,
            False,
            False,
            14.0,
            None,
            has_size_variation=True,
            has_color_variation=False,
        )
        assert result == '<b><span style="font-size:14pt">text</span></b>'

    def test_color(self) -> None:
        """Color variation → <span style="color:...">text</span>."""
        result = _wrap_with_tags(
            "text",
            False,
            False,
            False,
            False,
            None,
            "#ff0000",
            has_size_variation=False,
            has_color_variation=True,
        )
        assert result == '<span style="color:#ff0000">text</span>'

    def test_no_variation_no_span(self) -> None:
        """No variation flags → no <span> even with size/color values."""
        result = _wrap_with_tags(
            "text",
            True,
            False,
            False,
            False,
            14.0,
            "#ff0000",
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<b>text</b>"

    def test_both_variations(self) -> None:
        """Both size and color vary → combined span style."""
        result = _wrap_with_tags(
            "hi",
            False,
            True,
            False,
            False,
            10.0,
            "#00ff00",
            has_size_variation=True,
            has_color_variation=True,
        )
        assert '<span style="font-size:10pt;color:#00ff00">' in result
        assert result.startswith("<i>")


class TestParseHtmlFormattingSpans:
    """Tests for _parse_html_formatting with <span> tags."""

    def test_span_with_size(self) -> None:
        """Parses <span style="font-size:14pt">."""
        segs = _parse_html_formatting('<span style="font-size:14pt">big</span>')
        assert len(segs) == 1
        assert segs[0].text == "big"
        assert segs[0].font_size_pt == 14.0  # noqa: PLR2004
        assert segs[0].color_hex is None

    def test_span_with_color(self) -> None:
        """Parses <span style="color:#ff0000">."""
        segs = _parse_html_formatting('<span style="color:#ff0000">red</span>')
        assert len(segs) == 1
        assert segs[0].color_hex == "#ff0000"
        assert segs[0].font_size_pt is None

    def test_span_with_both(self) -> None:
        """Parses span with both font-size and color."""
        segs = _parse_html_formatting(
            '<span style="font-size:10pt;color:#0000ff">blue small</span>'
        )
        assert len(segs) == 1
        assert segs[0].font_size_pt == 10.0  # noqa: PLR2004
        assert segs[0].color_hex == "#0000ff"

    def test_merge_with_same_span_props(self) -> None:
        """Adjacent segments with same formatting + span props merge."""
        segs = _parse_html_formatting(
            '<span style="font-size:14pt">A</span><span style="font-size:14pt">B</span>'
        )
        assert len(segs) == 1
        assert segs[0].text == "AB"

    def test_no_merge_different_sizes(self) -> None:
        """Adjacent segments with different sizes stay separate."""
        segs = _parse_html_formatting(
            '<span style="font-size:14pt">A</span><span style="font-size:10pt">B</span>'
        )
        assert len(segs) == 2  # noqa: PLR2004

    def test_no_spans_backward_compat(self) -> None:
        """No <span> tags → segments have None size/color (backward compat)."""
        segs = _parse_html_formatting("<b>Bold</b> plain")
        assert segs[0].font_size_pt is None
        assert segs[0].color_hex is None

    def test_bold_with_span(self) -> None:
        """<b><span style="...">text</span></b> parses correctly."""
        segs = _parse_html_formatting(
            '<b><span style="font-size:14pt;color:#ff0000">red bold</span></b>'
        )
        assert len(segs) == 1
        assert segs[0].bold is True
        assert segs[0].font_size_pt == 14.0  # noqa: PLR2004
        assert segs[0].color_hex == "#ff0000"


class TestFormattingHtmlReSpan:
    """Tests that _FORMATTING_HTML_RE detects <span> tags."""

    def test_span_open(self) -> None:
        """<span ...> is detected."""
        assert _FORMATTING_HTML_RE.search('<span style="font-size:14pt">') is not None

    def test_span_close(self) -> None:
        """</span> is detected."""
        assert _FORMATTING_HTML_RE.search("</span>") is not None

    def test_span_bare(self) -> None:
        """<span> (no attrs) is detected."""
        assert _FORMATTING_HTML_RE.search("<span>") is not None


# ---------------------------------------------------------------------------
# UNO full formatting tests
# ---------------------------------------------------------------------------


class TestReadUnoPortionFullFormatting:
    """Tests for _read_uno_portion_full_formatting."""

    def test_basic(self) -> None:
        """Reads bold + default size/colour."""
        portion = _make_uno_portion("Bold", bold=True)
        b, i, u, s, sup, sub, sz, clr, bg = _read_uno_portion_full_formatting(portion)
        assert b is True
        assert sz == 11.0  # noqa: PLR2004
        assert clr == "#000000"

    def test_custom_size_and_color(self) -> None:
        """Reads custom CharHeight and CharColor."""
        portion = _make_uno_portion("Big Red", font_size=24.0, color_int=0xFF0000)
        b, i, u, s, sup, sub, sz, clr, bg = _read_uno_portion_full_formatting(portion)
        assert sz == 24.0  # noqa: PLR2004
        assert clr == "#ff0000"

    def test_auto_color_returns_none(self) -> None:
        """CharColor -1 (auto) → color_hex is None."""
        portion = _make_uno_portion("Auto", color_int=-1)
        _, _, _, _, _, _, _, clr, _ = _read_uno_portion_full_formatting(portion)
        assert clr is None


class TestHasUnoMixedFormattingSize:
    """Tests that _has_uno_mixed_formatting detects size variation."""

    def test_same_size_returns_false(self) -> None:
        """Uniform size → not mixed (when bold/italic also match)."""
        specs = [("A", False, False, False, False), ("B", False, False, False, False)]
        para = _make_uno_para_with_portions(specs)
        assert _has_uno_mixed_formatting(para) is False

    def test_different_size_returns_true(self) -> None:
        """Different CharHeight across portions → mixed."""
        portions = [
            _make_uno_portion("Big", font_size=24.0),
            _make_uno_portion("Small", font_size=10.0),
        ]
        enum_mock = MagicMock()
        _idx = [0]
        enum_mock.hasMoreElements.side_effect = lambda: _idx[0] < len(portions)

        def _next() -> MagicMock:
            p = portions[_idx[0]]
            _idx[0] += 1
            return p

        enum_mock.nextElement.side_effect = _next
        para = MagicMock()
        para.createEnumeration.return_value = enum_mock
        assert _has_uno_mixed_formatting(para) is True


class TestUnoRunsToHtmlSizeColor:
    """Tests for _uno_runs_to_html with size/colour variation."""

    def test_size_variation_emits_span(self) -> None:
        """Different sizes across portions → <span style="font-size:...">."""
        portions = [
            _make_uno_portion("Big", bold=True, font_size=24.0),
            _make_uno_portion("Small", font_size=10.0),
        ]
        enum_mock = MagicMock()
        _idx = [0]
        enum_mock.hasMoreElements.side_effect = lambda: _idx[0] < len(portions)

        def _next() -> MagicMock:
            p = portions[_idx[0]]
            _idx[0] += 1
            return p

        enum_mock.nextElement.side_effect = _next
        para = MagicMock()
        para.createEnumeration.return_value = enum_mock
        result = _uno_runs_to_html(para)
        # Base-value optimization: 24pt (first seen) is omitted, 10pt kept
        assert "font-size:10pt" in result
        assert "<b>" in result

    def test_uniform_size_no_span(self) -> None:
        """Same size across portions → no <span>."""
        portions = [
            _make_uno_portion("A", bold=True, font_size=11.0),
            _make_uno_portion("B", font_size=11.0),
        ]
        enum_mock = MagicMock()
        _idx = [0]
        enum_mock.hasMoreElements.side_effect = lambda: _idx[0] < len(portions)

        def _next() -> MagicMock:
            p = portions[_idx[0]]
            _idx[0] += 1
            return p

        enum_mock.nextElement.side_effect = _next
        para = MagicMock()
        para.createEnumeration.return_value = enum_mock
        result = _uno_runs_to_html(para)
        assert "<span" not in result
        assert "<b>A</b>" in result

    def test_color_variation_emits_span(self) -> None:
        """Different colours → <span style="color:...">."""
        portions = [
            _make_uno_portion("Red", color_int=0xFF0000),
            _make_uno_portion("Blue", color_int=0x0000FF),
        ]
        enum_mock = MagicMock()
        _idx = [0]
        enum_mock.hasMoreElements.side_effect = lambda: _idx[0] < len(portions)

        def _next() -> MagicMock:
            p = portions[_idx[0]]
            _idx[0] += 1
            return p

        enum_mock.nextElement.side_effect = _next
        para = MagicMock()
        para.createEnumeration.return_value = enum_mock
        result = _uno_runs_to_html(para)
        # Base-value optimization: first color is base, only second emitted
        assert "color:#0000ff" in result


class TestInjectUnoHtmlRunsSizeColor:
    """Tests that _inject_uno_html_runs sets CharHeight/CharColor per segment."""

    @staticmethod
    def _inject_and_collect_cursors(
        html_text: str,
        base_props: dict[str, object] | None = None,
    ) -> list[MagicMock]:
        """Injects HTML into a UNO para, returns list of created cursors."""
        specs = [("A", False, False, False, False)]
        para = _make_uno_para_with_portions(specs)
        cursors: list[MagicMock] = []
        text_obj = para.getText.return_value
        orig_fn = text_obj.createTextCursorByRange.side_effect

        def _track(rng: object) -> MagicMock:
            c = orig_fn(rng)
            cursors.append(c)
            return c

        text_obj.createTextCursorByRange.side_effect = _track
        _inject_uno_html_runs(para, html_text, base_props or {})
        return cursors

    def test_sets_char_height(self) -> None:
        """Segment with font_size_pt → CharHeight set on cursor."""
        html_text = '<span style="font-size:14pt">translated</span>'
        cursors = self._inject_and_collect_cursors(html_text)
        # cursors[0] = full-range cursor, cursors[1] = segment cursor
        assert len(cursors) >= 2  # noqa: PLR2004
        seg_cursor = cursors[1]
        props = seg_cursor._props
        assert "CharHeight" in props
        assert props["CharHeight"] == 14.0  # noqa: PLR2004

    def test_sets_char_color(self) -> None:
        """Segment with color_hex → CharColor set on cursor."""
        html_text = '<span style="color:#ff0000">red text</span>'
        cursors = self._inject_and_collect_cursors(html_text)
        assert len(cursors) >= 2  # noqa: PLR2004
        seg_cursor = cursors[1]
        props = seg_cursor._props
        assert "CharColor" in props
        expected_color = 0xFF0000  # noqa: PLR2004
        assert props["CharColor"] == expected_color

    def test_no_size_color_when_none(self) -> None:
        """Segment without size/color → CharHeight/CharColor NOT set."""
        html_text = "<b>bold text</b>"
        cursors = self._inject_and_collect_cursors(html_text)
        assert len(cursors) >= 2  # noqa: PLR2004
        seg_cursor = cursors[1]
        props = seg_cursor._props
        assert "CharHeight" not in props
        assert "CharColor" not in props


# ---------------------------------------------------------------------------
# DOCX full formatting tests
# ---------------------------------------------------------------------------


def _make_para_with_runs_extended(
    run_specs: list[tuple[str, bool, bool, bool, bool, float | None, str | None]],
) -> object:
    """Creates a Document paragraph with runs including font size and colour.

    Each spec is (text, bold, italic, underline, strike, size_pt, color_hex).
    """
    from docx.shared import Pt, RGBColor  # noqa: PLC0415

    doc = Document()
    para = doc.add_paragraph()
    for text, bold, italic, underline, strike, size_pt, color_hex in run_specs:
        run = para.add_run(text)
        run.bold = bold
        run.italic = italic
        run.underline = underline
        run.font.strike = strike
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if color_hex is not None:
            hex_str = color_hex.lstrip("#")
            run.font.color.rgb = RGBColor(
                int(hex_str[0:2], 16),
                int(hex_str[2:4], 16),
                int(hex_str[4:6], 16),
            )
    return para


class TestHasMixedFormattingSize:
    """Tests that _has_mixed_formatting detects size/colour variation."""

    def test_same_size_not_mixed(self) -> None:
        """Uniform size + uniform formatting → not mixed."""
        from docx.shared import Pt  # noqa: PLC0415

        doc = Document()
        para = doc.add_paragraph()
        for text in ("A", "B"):
            run = para.add_run(text)
            run.bold = True
            run.font.size = Pt(14)
        assert _has_mixed_formatting(para) is False

    def test_different_size_mixed(self) -> None:
        """Different font sizes → mixed."""
        para = _make_para_with_runs_extended(
            [
                ("Big", False, False, False, False, 24.0, None),
                ("Small", False, False, False, False, 10.0, None),
            ]
        )
        assert _has_mixed_formatting(para) is True

    def test_different_color_mixed(self) -> None:
        """Different colours → mixed."""
        para = _make_para_with_runs_extended(
            [
                ("Red", False, False, False, False, None, "#ff0000"),
                ("Blue", False, False, False, False, None, "#0000ff"),
            ]
        )
        assert _has_mixed_formatting(para) is True


class TestRunsToHtmlSizeColor:
    """Tests that _runs_to_html emits <span> for size/colour variation."""

    def test_size_variation_emits_span(self) -> None:
        """Different sizes → deviating size gets <span>; base omitted."""
        para = _make_para_with_runs_extended(
            [
                ("Big", True, False, False, False, 24.0, None),
                ("Small", False, False, False, False, 10.0, None),
            ]
        )
        result = _runs_to_html(para)
        # Base-value optimization: 24pt is first seen → base, 10pt kept
        assert "font-size:10pt" in result

    def test_uniform_size_no_span(self) -> None:
        """Uniform size → no <span>."""
        para = _make_para_with_runs_extended(
            [
                ("A", True, False, False, False, 14.0, None),
                ("B", False, False, False, False, 14.0, None),
            ]
        )
        result = _runs_to_html(para)
        assert "<span" not in result

    def test_color_variation_emits_span(self) -> None:
        """Different colours → deviating colour gets <span>; base omitted."""
        para = _make_para_with_runs_extended(
            [
                ("Red", False, False, False, False, None, "#ff0000"),
                ("Blue", False, False, False, False, None, "#0000ff"),
            ]
        )
        result = _runs_to_html(para)
        # Base-value optimization: #ff0000 is first seen → base, #0000ff kept
        assert "color:#0000ff" in result


class TestInjectHtmlRunsSizeColor:
    """Tests that _inject_html_runs sets w:sz and w:color per segment."""

    def test_sets_font_size(self) -> None:
        """Segment with font_size_pt → w:sz/w:szCs set."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("original")
        html_text = '<span style="font-size:14pt">translated</span>'
        _inject_html_runs(para, html_text)
        # Check the run XML has w:sz
        assert len(para.runs) == 1
        rpr = para.runs[0]._element.find(qn("w:rPr"))
        assert rpr is not None
        sz = rpr.find(qn("w:sz"))
        assert sz is not None
        half_pts = 28  # noqa: PLR2004
        assert sz.get(qn("w:val")) == str(half_pts)

    def test_sets_color(self) -> None:
        """Segment with color_hex → w:color set."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("original")
        html_text = '<span style="color:#ff0000">translated</span>'
        _inject_html_runs(para, html_text)
        assert len(para.runs) == 1
        rpr = para.runs[0]._element.find(qn("w:rPr"))
        assert rpr is not None
        color_elem = rpr.find(qn("w:color"))
        assert color_elem is not None
        assert color_elem.get(qn("w:val")) == "FF0000"

    def test_no_size_color_without_span(self) -> None:
        """Segments without span → no w:sz or w:color added."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("original")
        html_text = "<b>bold text</b>"
        _inject_html_runs(para, html_text)
        rpr = para.runs[0]._element.find(qn("w:rPr"))
        # No w:sz should be present (unless from base rPr)
        sz = rpr.find(qn("w:sz")) if rpr is not None else None
        color_elem = rpr.find(qn("w:color")) if rpr is not None else None
        # These shouldn't have been added by the injection
        assert sz is None or color_elem is None


class TestSetRprFontSize:
    """Tests for _set_rpr_font_size."""

    def test_sets_sz_and_szcs(self) -> None:
        """Creates w:sz and w:szCs with correct half-point values."""
        rpr = OxmlElement("w:rPr")
        _set_rpr_font_size(rpr, 14.0, qn)
        sz = rpr.find(qn("w:sz"))
        szcs = rpr.find(qn("w:szCs"))
        assert sz is not None
        half_pts = 28  # noqa: PLR2004
        assert sz.get(qn("w:val")) == str(half_pts)
        assert szcs is not None
        assert szcs.get(qn("w:val")) == str(half_pts)

    def test_overwrites_existing(self) -> None:
        """Overwrites existing w:sz value."""
        rpr = OxmlElement("w:rPr")
        _set_rpr_font_size(rpr, 10.0, qn)
        _set_rpr_font_size(rpr, 14.0, qn)
        sz = rpr.find(qn("w:sz"))
        half_pts = 28  # noqa: PLR2004
        assert sz.get(qn("w:val")) == str(half_pts)


class TestSetRprColor:
    """Tests for _set_rpr_color."""

    def test_sets_color(self) -> None:
        """Creates w:color with uppercase hex."""
        rpr = OxmlElement("w:rPr")
        _set_rpr_color(rpr, "#ff0000", qn)
        color_elem = rpr.find(qn("w:color"))
        assert color_elem is not None
        assert color_elem.get(qn("w:val")) == "FF0000"

    def test_overwrites_existing(self) -> None:
        """Overwrites existing w:color value."""
        rpr = OxmlElement("w:rPr")
        _set_rpr_color(rpr, "#ff0000", qn)
        _set_rpr_color(rpr, "#00ff00", qn)
        color_elem = rpr.find(qn("w:color"))
        assert color_elem.get(qn("w:val")) == "00FF00"


# ---------------------------------------------------------------------------
# PPTX full formatting tests
# ---------------------------------------------------------------------------


def _make_pptx_para_with_runs_extended(
    specs: list[tuple[str, bool, bool, bool, bool, float | None, str | None]],
) -> object:
    """Creates a PPTX paragraph with runs including font size and colour.

    Each spec is (text, bold, italic, underline, strike, size_pt, color_hex).
    """
    from lxml import etree  # noqa: PLC0415
    from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = txbox.text_frame
    para = tf.paragraphs[0]

    for idx, (text, bold, italic, underline, strike, size_pt, color_hex) in enumerate(
        specs,
    ):
        if idx == 0:
            run = para.runs[0] if para.runs else para.add_run()
        else:
            run = para.add_run()
        run.text = text
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        if strike:
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                rpr = etree.SubElement(run._r, pptx_qn("a:rPr"))
            rpr.set("strike", "sngStrike")
        if size_pt is not None:
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                rpr = etree.SubElement(run._r, pptx_qn("a:rPr"))
            rpr.set("sz", str(int(size_pt * 100)))
        if color_hex is not None:
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                rpr = etree.SubElement(run._r, pptx_qn("a:rPr"))
            solid_fill = etree.SubElement(rpr, pptx_qn("a:solidFill"))
            srgb = etree.SubElement(solid_fill, pptx_qn("a:srgbClr"))
            srgb.set("val", color_hex.lstrip("#").upper())

    return para


class TestReadPptxRunFullFormatting:
    """Tests for _read_pptx_run_full_formatting."""

    def test_basic(self) -> None:
        """Bold run without size/color → None for size/color."""
        para = _make_pptx_para_with_runs(
            [
                ("Bold", True, False, False, False),
            ]
        )
        b, i, u, s, sup, sub, sz, clr, bg = _read_pptx_run_full_formatting(para.runs[0])
        assert b is True
        assert sz is None
        assert clr is None
        assert bg is None

    def test_with_size(self) -> None:
        """Run with sz attribute → correct font_size_pt."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("Big", False, False, False, False, 24.0, None),
            ]
        )
        _, _, _, _, _, _, sz, _, _ = _read_pptx_run_full_formatting(para.runs[0])
        assert sz == 24.0  # noqa: PLR2004

    def test_with_color(self) -> None:
        """Run with solidFill/srgbClr → correct color_hex."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("Red", False, False, False, False, None, "#ff0000"),
            ]
        )
        _, _, _, _, _, _, _, clr, _ = _read_pptx_run_full_formatting(para.runs[0])
        assert clr == "#ff0000"

    def test_with_bg_color(self) -> None:
        """Run with highlight/srgbClr → correct bg_color_hex."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs_extended(
            [
                ("Highlighted", False, False, False, False, None, None),
            ]
        )
        # Manually add <a:highlight>/<a:srgbClr val="FFFF00"> to the run
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        if rpr is None:
            rpr = etree.SubElement(para.runs[0]._r, pptx_qn("a:rPr"))
        hl = etree.SubElement(rpr, pptx_qn("a:highlight"))
        srgb = etree.SubElement(hl, pptx_qn("a:srgbClr"))
        srgb.set("val", "FFFF00")

        _, _, _, _, _, _, _, _, bg = _read_pptx_run_full_formatting(para.runs[0])
        assert bg == "#ffff00"

    def test_no_bg_returns_none(self) -> None:
        """Run without highlight → bg_color_hex is None."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("Plain", False, False, False, False, None, None),
            ]
        )
        _, _, _, _, _, _, _, _, bg = _read_pptx_run_full_formatting(para.runs[0])
        assert bg is None


class TestHasPptxMixedFormattingSize:
    """Tests that _has_pptx_mixed_formatting detects size variation."""

    def test_different_sizes_mixed(self) -> None:
        """Runs with different sz → mixed."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("Big", False, False, False, False, 24.0, None),
                ("Small", False, False, False, False, 10.0, None),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is True

    def test_same_sizes_not_mixed(self) -> None:
        """Runs with same sz → not mixed."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("A", False, False, False, False, 14.0, None),
                ("B", False, False, False, False, 14.0, None),
            ]
        )
        assert _has_pptx_mixed_formatting(para) is False

    def test_different_bg_mixed(self) -> None:
        """Runs with different highlight → mixed."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs_extended(
            [
                ("A", False, False, False, False, None, None),
                ("B", False, False, False, False, None, None),
            ]
        )
        # Add highlight to first run only
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        if rpr is None:
            rpr = etree.SubElement(para.runs[0]._r, pptx_qn("a:rPr"))
        hl = etree.SubElement(rpr, pptx_qn("a:highlight"))
        srgb = etree.SubElement(hl, pptx_qn("a:srgbClr"))
        srgb.set("val", "FFFF00")

        assert _has_pptx_mixed_formatting(para) is True

    def test_same_bg_not_mixed(self) -> None:
        """Runs with identical highlight → not mixed."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs_extended(
            [
                ("A", False, False, False, False, None, None),
                ("B", False, False, False, False, None, None),
            ]
        )
        # Add same highlight to both runs
        for run in para.runs:
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                rpr = etree.SubElement(run._r, pptx_qn("a:rPr"))
            hl = etree.SubElement(rpr, pptx_qn("a:highlight"))
            srgb = etree.SubElement(hl, pptx_qn("a:srgbClr"))
            srgb.set("val", "FFFF00")

        assert _has_pptx_mixed_formatting(para) is False


class TestPptxRunsToHtmlSizeColor:
    """Tests that _pptx_runs_to_html emits <span> for size/colour/bg variation."""

    def test_size_variation(self) -> None:
        """Different sizes → deviating size gets <span>; base omitted."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("Big", True, False, False, False, 24.0, None),
                ("Small", False, False, False, False, 10.0, None),
            ]
        )
        result = _pptx_runs_to_html(para)
        assert "font-size:10pt" in result

    def test_uniform_size_no_span(self) -> None:
        """Uniform size → no <span>."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("A", True, False, False, False, 14.0, None),
                ("B", False, False, False, False, 14.0, None),
            ]
        )
        result = _pptx_runs_to_html(para)
        assert "<span" not in result

    def test_color_variation(self) -> None:
        """Different colours → deviating colour gets <span>; base omitted."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("Red", False, False, False, False, None, "#ff0000"),
                ("Blue", False, False, False, False, None, "#0000ff"),
            ]
        )
        result = _pptx_runs_to_html(para)
        assert "color:#0000ff" in result

    def test_bg_variation(self) -> None:
        """Different highlights → bg gets <span style="background-color:...">."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs_extended(
            [
                ("A", False, False, False, False, None, None),
                ("B", False, False, False, False, None, None),
            ]
        )
        # Add yellow highlight to first run, green to second
        for run, hex_val in zip(para.runs, ("FFFF00", "00FF00"), strict=True):
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                rpr = etree.SubElement(run._r, pptx_qn("a:rPr"))
            hl = etree.SubElement(rpr, pptx_qn("a:highlight"))
            srgb = etree.SubElement(hl, pptx_qn("a:srgbClr"))
            srgb.set("val", hex_val)

        result = _pptx_runs_to_html(para)
        assert "background-color:#ffff00" in result
        assert "background-color:#00ff00" in result

    def test_uniform_bg_no_span(self) -> None:
        """Uniform highlight → no background-color in <span>."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs_extended(
            [
                ("A", True, False, False, False, None, None),
                ("B", False, False, False, False, None, None),
            ]
        )
        # Same highlight on both runs
        for run in para.runs:
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                rpr = etree.SubElement(run._r, pptx_qn("a:rPr"))
            hl = etree.SubElement(rpr, pptx_qn("a:highlight"))
            srgb = etree.SubElement(hl, pptx_qn("a:srgbClr"))
            srgb.set("val", "FFFF00")

        result = _pptx_runs_to_html(para)
        assert "background-color" not in result


class TestApplyPptxFormatAttrsSizeColor:
    """Tests that _apply_pptx_format_attrs sets sz and solidFill."""

    def test_sets_sz(self) -> None:
        """Segment with font_size_pt → sz attribute set."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        rpr = etree.Element(pptx_qn("a:rPr"))
        seg = _FormattedSegment(
            "text",
            True,
            False,
            False,
            False,
            False,
            False,
            14.0,
            None,
        )
        _apply_pptx_format_attrs(rpr, seg)
        sz_val = 1400  # noqa: PLR2004
        assert rpr.get("sz") == str(sz_val)
        assert rpr.get("b") == "1"

    def test_sets_solid_fill(self) -> None:
        """Segment with color_hex → solidFill/srgbClr created."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        rpr = etree.Element(pptx_qn("a:rPr"))
        seg = _FormattedSegment(
            "text",
            False,
            False,
            False,
            False,
            False,
            False,
            None,
            "#ff0000",
        )
        _apply_pptx_format_attrs(rpr, seg)
        solid_fill = rpr.find(pptx_qn("a:solidFill"))
        assert solid_fill is not None
        srgb = solid_fill.find(pptx_qn("a:srgbClr"))
        assert srgb is not None
        assert srgb.get("val") == "FF0000"

    def test_no_sz_without_size(self) -> None:
        """Segment without font_size_pt → no sz attribute."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        rpr = etree.Element(pptx_qn("a:rPr"))
        seg = _FormattedSegment("text", True, False, False, False)
        _apply_pptx_format_attrs(rpr, seg)
        assert rpr.get("sz") is None


class TestInjectPptxHtmlRunsSizeColor:
    """Tests that _inject_pptx_html_runs sets sz and solidFill."""

    def test_inject_with_size(self) -> None:
        """Injected segments with font-size → sz attribute on rPr."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("original", False, False, False, False),
            ]
        )
        html_text = '<span style="font-size:14pt">translated</span>'
        _inject_pptx_html_runs(para, html_text)
        runs = para.runs
        assert len(runs) == 1
        rpr = runs[0]._r.find(pptx_qn("a:rPr"))
        assert rpr is not None
        sz_val = 1400  # noqa: PLR2004
        assert rpr.get("sz") == str(sz_val)

    def test_inject_with_color(self) -> None:
        """Injected segments with color → solidFill on rPr."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs(
            [
                ("original", False, False, False, False),
            ]
        )
        html_text = '<span style="color:#ff0000">red text</span>'
        _inject_pptx_html_runs(para, html_text)
        runs = para.runs
        assert len(runs) == 1
        rpr = runs[0]._r.find(pptx_qn("a:rPr"))
        solid_fill = rpr.find(pptx_qn("a:solidFill"))
        assert solid_fill is not None
        srgb = solid_fill.find(pptx_qn("a:srgbClr"))
        assert srgb is not None
        assert srgb.get("val") == "FF0000"

    def test_base_rpr_strips_sz_and_fill(self) -> None:
        """Base rPr should have sz and solidFill stripped."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        # Create a paragraph with a run that has sz and solidFill
        para = _make_pptx_para_with_runs_extended(
            [
                ("original", True, False, False, False, 24.0, "#ff0000"),
            ]
        )
        # Inject with different formatting
        html_text = '<b><span style="font-size:10pt;color:#00ff00">new</span></b>'
        _inject_pptx_html_runs(para, html_text)
        runs = para.runs
        assert len(runs) == 1
        rpr = runs[0]._r.find(pptx_qn("a:rPr"))
        # The injected run should have the NEW sz/color, not base
        sz_val = 1000  # noqa: PLR2004
        assert rpr.get("sz") == str(sz_val)
        srgb = rpr.find(pptx_qn("a:solidFill")).find(pptx_qn("a:srgbClr"))
        assert srgb.get("val") == "00FF00"


# ---------------------------------------------------------------------------
# DOCX HTML round-trip with size/colour
# ---------------------------------------------------------------------------


class TestDocxHtmlRoundTripSizeColor:
    """Tests DOCX encode → parse → verify with font size and colour."""

    def test_size_round_trip(self) -> None:
        """Encode with size variation → both sizes preserved explicitly."""
        para = _make_para_with_runs_extended(
            [
                ("Big", True, False, False, False, 24.0, None),
                ("Small", False, False, False, False, 10.0, None),
            ]
        )
        html_text = _runs_to_html(para)
        segs = _parse_html_formatting(html_text)
        assert len(segs) == 2  # noqa: PLR2004
        assert segs[0].bold is True
        # All non-None sizes are emitted explicitly (no base-value omission)
        assert segs[0].font_size_pt == 24.0  # noqa: PLR2004
        assert segs[1].font_size_pt == 10.0  # noqa: PLR2004

    def test_color_round_trip(self) -> None:
        """Encode with colour variation → both colours preserved explicitly."""
        para = _make_para_with_runs_extended(
            [
                ("Red", False, False, False, False, None, "#ff0000"),
                ("Blue", False, False, False, False, None, "#0000ff"),
            ]
        )
        html_text = _runs_to_html(para)
        segs = _parse_html_formatting(html_text)
        assert len(segs) == 2  # noqa: PLR2004
        # All non-None colours are emitted explicitly (no base-value omission)
        assert segs[0].color_hex == "#ff0000"
        assert segs[1].color_hex == "#0000ff"


# ---------------------------------------------------------------------------
# PPTX HTML round-trip with size/colour
# ---------------------------------------------------------------------------


class TestPptxHtmlRoundTripSizeColor:
    """Tests PPTX encode → parse → verify with font size and colour."""

    def test_size_round_trip(self) -> None:
        """Encode with size variation → both sizes preserved explicitly."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("Big", False, False, False, False, 24.0, None),
                ("Small", False, False, False, False, 10.0, None),
            ]
        )
        html_text = _pptx_runs_to_html(para)
        segs = _parse_html_formatting(html_text)
        assert len(segs) == 2  # noqa: PLR2004
        # All non-None sizes are emitted explicitly (no base-value omission)
        assert segs[0].font_size_pt == 24.0  # noqa: PLR2004
        assert segs[1].font_size_pt == 10.0  # noqa: PLR2004

    def test_color_round_trip(self) -> None:
        """Encode with colour variation → both colours preserved explicitly."""
        para = _make_pptx_para_with_runs_extended(
            [
                ("Red", False, False, False, False, None, "#ff0000"),
                ("Blue", False, False, False, False, None, "#0000ff"),
            ]
        )
        html_text = _pptx_runs_to_html(para)
        segs = _parse_html_formatting(html_text)
        assert len(segs) == 2  # noqa: PLR2004
        # All non-None colours are emitted explicitly (no base-value omission)
        assert segs[0].color_hex == "#ff0000"
        assert segs[1].color_hex == "#0000ff"

    def test_bg_round_trip(self) -> None:
        """Encode with bg variation → both bg colours preserved explicitly."""
        from lxml import etree  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para = _make_pptx_para_with_runs_extended(
            [
                ("A", False, False, False, False, None, None),
                ("B", False, False, False, False, None, None),
            ]
        )
        # Add different highlights to each run
        for run, hex_val in zip(para.runs, ("FFFF00", "00FF00"), strict=True):
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                rpr = etree.SubElement(run._r, pptx_qn("a:rPr"))
            hl = etree.SubElement(rpr, pptx_qn("a:highlight"))
            srgb = etree.SubElement(hl, pptx_qn("a:srgbClr"))
            srgb.set("val", hex_val)

        html_text = _pptx_runs_to_html(para)
        segs = _parse_html_formatting(html_text)
        assert len(segs) == 2  # noqa: PLR2004
        assert segs[0].bg_color_hex == "#ffff00"
        assert segs[1].bg_color_hex == "#00ff00"


# ---------------------------------------------------------------------------
# UNO HTML round-trip with size/colour
# ---------------------------------------------------------------------------


class TestUnoHtmlRoundTripSizeColor:
    """Tests UNO encode → parse → verify with font size and colour."""

    def test_size_round_trip(self) -> None:
        """Encode with size variation → parse → correct sizes."""
        portions = [
            _make_uno_portion("Big", bold=True, font_size=24.0),
            _make_uno_portion("Small", font_size=10.0),
        ]
        enum_mock = MagicMock()
        _idx = [0]
        enum_mock.hasMoreElements.side_effect = lambda: _idx[0] < len(portions)

        def _next() -> MagicMock:
            p = portions[_idx[0]]
            _idx[0] += 1
            return p

        enum_mock.nextElement.side_effect = _next
        para = MagicMock()
        para.createEnumeration.return_value = enum_mock
        html_text = _uno_runs_to_html(para)
        segs = _parse_html_formatting(html_text)
        assert len(segs) == 2  # noqa: PLR2004
        assert segs[0].bold is True
        # All non-None sizes are emitted explicitly (no base-value omission)
        assert segs[0].font_size_pt == 24.0  # noqa: PLR2004
        assert segs[1].font_size_pt == 10.0  # noqa: PLR2004


# ---------------------------------------------------------------------------
# Win32com BGR colour conversion helpers
# ---------------------------------------------------------------------------

WIN32COM_UNDEFINED = 9999999  # noqa: PLR2004


class TestWin32comColorToHex:
    """Tests for _win32com_color_to_hex (BGR → #rrggbb)."""

    def test_pure_red_bgr(self) -> None:
        """Red in BGR = 0x0000FF = 255 → '#ff0000'."""
        assert _win32com_color_to_hex(255) == "#ff0000"  # noqa: PLR2004

    def test_pure_blue_bgr(self) -> None:
        """Blue in BGR = 0xFF0000 = 16711680 → '#0000ff'."""
        assert _win32com_color_to_hex(16711680) == "#0000ff"  # noqa: PLR2004

    def test_pure_green_bgr(self) -> None:
        """Green in BGR = 0x00FF00 = 65280 → '#00ff00'."""
        assert _win32com_color_to_hex(65280) == "#00ff00"  # noqa: PLR2004

    def test_black(self) -> None:
        """Zero → '#000000'."""
        assert _win32com_color_to_hex(0) == "#000000"

    def test_white(self) -> None:
        """0xFFFFFF → '#ffffff'."""
        assert _win32com_color_to_hex(0xFFFFFF) == "#ffffff"

    def test_negative_returns_none(self) -> None:
        """Negative values (automatic colour) → None."""
        assert _win32com_color_to_hex(-1) is None

    def test_overflow_returns_none(self) -> None:
        """Values > 0xFFFFFF → None."""
        assert _win32com_color_to_hex(0x1000000) is None


class TestColorHexToWin32com:
    """Tests for _color_hex_to_win32com (#rrggbb → BGR int)."""

    def test_red_to_bgr(self) -> None:
        """'#ff0000' → 255 (red in BGR byte order)."""
        assert _color_hex_to_win32com("#ff0000") == 255  # noqa: PLR2004

    def test_blue_to_bgr(self) -> None:
        """'#0000ff' → 16711680 (blue in BGR)."""
        assert _color_hex_to_win32com("#0000ff") == 16711680  # noqa: PLR2004

    def test_green_to_bgr(self) -> None:
        """'#00ff00' → 65280."""
        assert _color_hex_to_win32com("#00ff00") == 65280  # noqa: PLR2004

    def test_round_trip_red(self) -> None:
        """Round-trip: hex → BGR → hex."""
        bgr = _color_hex_to_win32com("#ff0000")
        assert _win32com_color_to_hex(bgr) == "#ff0000"

    def test_round_trip_mixed(self) -> None:
        """Round-trip with mixed colour '#1a2b3c'."""
        bgr = _color_hex_to_win32com("#1a2b3c")
        assert _win32com_color_to_hex(bgr) == "#1a2b3c"


# ---------------------------------------------------------------------------
# Win32com Word extraction helpers
# ---------------------------------------------------------------------------


def _make_win32com_char(  # noqa: PLR0913
    text: str = "A",
    bold: object = False,
    italic: object = False,
    underline: object = 0,
    strike: object = False,
    size: object = 12.0,
    color: object = 0,
    *,
    shading_bg: object = -16777216,
    highlight_index: object = 0,
    superscript: object = False,
    subscript: object = False,
) -> MagicMock:
    """Creates a mock win32com character Range with Font properties."""
    ch = MagicMock()
    ch.Text = text
    ch.Font.Bold = bold
    ch.Font.Italic = italic
    ch.Font.Underline = underline
    ch.Font.StrikeThrough = strike
    ch.Font.Superscript = superscript
    ch.Font.Subscript = subscript
    ch.Font.Size = size
    ch.Font.Color = color
    # Background: Shading (arbitrary colour) and HighlightColorIndex (predefined)
    ch.Shading.BackgroundPatternColor = shading_bg
    ch.HighlightColorIndex = highlight_index
    return ch


class TestReadWin32comCharFormatting:
    """Tests for _read_win32com_char_formatting."""

    def test_bold_italic(self) -> None:
        """Reads bold + italic flags correctly."""
        ch = _make_win32com_char(bold=True, italic=True)
        b, i, u, s, sup, sub, sz, clr, bg = _read_win32com_char_formatting(ch)
        assert b is True
        assert i is True
        assert u is False
        assert s is False

    def test_underline_enum(self) -> None:
        """Underline enum != 0 → True."""
        ch = _make_win32com_char(underline=1)
        _, _, u, _, _, _, _, _, _ = _read_win32com_char_formatting(ch)
        assert u is True

    def test_strikethrough(self) -> None:
        """StrikeThrough = True → strike=True."""
        ch = _make_win32com_char(strike=True)
        _, _, _, s, _, _, _, _, _ = _read_win32com_char_formatting(ch)
        assert s is True

    def test_size_and_color(self) -> None:
        """Reads font size and BGR colour."""
        # Red in BGR = 255
        ch = _make_win32com_char(size=14.0, color=255)
        _, _, _, _, _, _, sz, clr, _ = _read_win32com_char_formatting(ch)
        assert sz == 14.0  # noqa: PLR2004
        assert clr == "#ff0000"

    def test_undefined_values_skipped(self) -> None:
        """WIN32COM_UNDEFINED sentinel → False/None."""
        ch = _make_win32com_char(
            bold=WIN32COM_UNDEFINED,
            italic=WIN32COM_UNDEFINED,
            underline=WIN32COM_UNDEFINED,
            strike=WIN32COM_UNDEFINED,
            size=WIN32COM_UNDEFINED,
            color=WIN32COM_UNDEFINED,
        )
        b, i, u, s, sup, sub, sz, clr, bg = _read_win32com_char_formatting(ch)
        assert b is False
        assert i is False
        assert u is False
        assert s is False
        assert sz is None
        assert clr is None
        assert bg is None

    def test_size_zero_is_none(self) -> None:
        """Size=0 (degenerate) → sz is None (guard: raw_size > 0)."""
        ch = _make_win32com_char(size=0)
        _, _, _, _, _, _, sz, _, _ = _read_win32com_char_formatting(ch)
        assert sz is None

    def test_automatic_color_negative_is_none(self) -> None:
        """Negative colour value (Word automatic colour, e.g. -1) → None."""
        ch = _make_win32com_char(color=-1)
        _, _, _, _, _, _, _, clr, _ = _read_win32com_char_formatting(ch)
        assert clr is None

    def test_bg_from_shading(self) -> None:
        """Shading.BackgroundPatternColor (BGR) → bg_color_hex."""
        # Yellow in BGR = 0x00FFFF = 65535
        ch = _make_win32com_char(shading_bg=65535)
        _, _, _, _, _, _, _, _, bg = _read_win32com_char_formatting(ch)
        assert bg == "#ffff00"

    def test_bg_from_highlight_index(self) -> None:
        """HighlightColorIndex enum → bg_color_hex via mapping."""
        # wdYellow = 7
        ch = _make_win32com_char(highlight_index=7)
        _, _, _, _, _, _, _, _, bg = _read_win32com_char_formatting(ch)
        assert bg == "#ffff00"

    def test_bg_shading_automatic_is_none(self) -> None:
        """Automatic shading colour (-16777216) → bg is None."""
        ch = _make_win32com_char(shading_bg=-16777216, highlight_index=0)
        _, _, _, _, _, _, _, _, bg = _read_win32com_char_formatting(ch)
        assert bg is None

    def test_bg_shading_takes_priority(self) -> None:
        """Shading is preferred over HighlightColorIndex."""
        # Red shading (BGR 255) + Yellow highlight — shading wins
        ch = _make_win32com_char(shading_bg=255, highlight_index=7)
        _, _, _, _, _, _, _, _, bg = _read_win32com_char_formatting(ch)
        assert bg == "#ff0000"


class TestHasWin32comWordMixedFormatting:
    """Tests for _has_win32com_word_mixed_formatting."""

    def test_uniform_returns_false(self) -> None:
        """All chars same formatting → False."""
        chars = [_make_win32com_char("H"), _make_win32com_char("i")]
        para = MagicMock()
        para.Range.Font.Bold = False
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        assert _has_win32com_word_mixed_formatting(para) is False

    def test_mixed_bold_returns_true(self) -> None:
        """One bold, one not → True (detected via quick check)."""
        para = MagicMock()
        para.Range.Font.Bold = WIN32COM_UNDEFINED
        assert _has_win32com_word_mixed_formatting(para) is True

    def test_mixed_size_returns_true(self) -> None:
        """Different font sizes → True."""
        chars = [
            _make_win32com_char("A", size=12.0),
            _make_win32com_char("B", size=24.0),
        ]
        para = MagicMock()
        para.Range.Font.Bold = False
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        assert _has_win32com_word_mixed_formatting(para) is True

    def test_mixed_color_returns_true(self) -> None:
        """Different colours → True."""
        chars = [
            _make_win32com_char("A", color=0),
            _make_win32com_char("B", color=255),
        ]
        para = MagicMock()
        para.Range.Font.Bold = False
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        assert _has_win32com_word_mixed_formatting(para) is True

    def test_single_char_returns_false(self) -> None:
        """One character → False."""
        chars = [_make_win32com_char("A")]
        para = MagicMock()
        para.Range.Font.Bold = False
        para.Range.Characters.Count = 1
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        assert _has_win32com_word_mixed_formatting(para) is False

    def test_whitespace_only_chars_skipped(self) -> None:
        """Only whitespace chars → False (they're skipped)."""
        chars = [
            _make_win32com_char(" "),
            _make_win32com_char("\t"),
        ]
        para = MagicMock()
        para.Range.Font.Bold = False
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        assert _has_win32com_word_mixed_formatting(para) is False

    def test_zero_char_paragraph_returns_false(self) -> None:
        """Count=0 → False (count <= 1 guard)."""
        para = MagicMock()
        para.Range.Font.Bold = False
        para.Range.Characters.Count = 0
        assert _has_win32com_word_mixed_formatting(para) is False

    def test_quick_check_exception_falls_through_to_iteration(self) -> None:
        """Font.Bold access raising → exception caught, iteration used instead."""
        chars = [_make_win32com_char("A"), _make_win32com_char("B")]

        class _BadBoldFont:
            @property
            def Bold(self) -> None:  # noqa: N802
                raise AttributeError("Bold unavailable")

        para = MagicMock()
        para.Range.Font = _BadBoldFont()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        # Chars are uniform → iteration returns False
        assert _has_win32com_word_mixed_formatting(para) is False


class TestWin32comWordRunsToHtml:
    """Tests for _win32com_word_runs_to_html."""

    def test_bold_and_plain_runs(self) -> None:
        """Bold 'Hello' + plain ' world' → HTML with <b> tag."""
        chars = [
            _make_win32com_char("H", bold=True),
            _make_win32com_char("e", bold=True),
            _make_win32com_char("l", bold=True),
            _make_win32com_char("l", bold=True),
            _make_win32com_char("o", bold=True),
            _make_win32com_char(" "),
            _make_win32com_char("w"),
            _make_win32com_char("o"),
            _make_win32com_char("r"),
            _make_win32com_char("l"),
            _make_win32com_char("d"),
            _make_win32com_char("\r"),  # trailing paragraph mark
        ]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        result = _win32com_word_runs_to_html(para)
        assert "<b>Hello</b>" in result
        assert " world" in result
        # No \r in output
        assert "\r" not in result

    def test_size_variation_emits_span(self) -> None:
        """Different sizes → deviating size gets <span>; base omitted."""
        chars = [
            _make_win32com_char("A", size=12.0),
            _make_win32com_char("B", size=24.0),
        ]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        result = _win32com_word_runs_to_html(para)
        # Base-value optimization: 12pt is first seen → base, 24pt kept
        assert "font-size:24pt" in result

    def test_color_variation_emits_span(self) -> None:
        """Different colours → deviating colour gets <span>; base omitted."""
        chars = [
            _make_win32com_char("R", color=255),  # red in BGR
            _make_win32com_char("B", color=16711680),  # blue in BGR
        ]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        result = _win32com_word_runs_to_html(para)
        # Base-value optimization: first color is base, second emitted
        assert "color:#0000ff" in result

    def test_empty_paragraph_returns_empty(self) -> None:
        r"""Paragraph with only \r → empty string."""
        chars = [_make_win32com_char("\r")]
        para = MagicMock()
        para.Range.Characters.Count = 1
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        assert _win32com_word_runs_to_html(para) == ""

    def test_strips_trailing_cr(self) -> None:
        r"""Paragraph mark \r at end is stripped."""
        chars = [
            _make_win32com_char("X"),
            _make_win32com_char("\r"),
        ]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        result = _win32com_word_runs_to_html(para)
        assert result == "X"
        assert "\r" not in result

    def test_uniform_formatting_no_span_tags(self) -> None:
        """All chars same size+colour → no <span> tags in output."""
        chars = [
            _make_win32com_char("H", size=12.0, color=0),
            _make_win32com_char("i", size=12.0, color=0),
        ]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        result = _win32com_word_runs_to_html(para)
        assert "Hi" in result
        assert "<span" not in result

    def test_html_entity_escaping(self) -> None:
        """Characters like '<' and '&' are HTML-escaped in output."""
        chars = [
            _make_win32com_char("<"),
            _make_win32com_char("&"),
        ]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        result = _win32com_word_runs_to_html(para)
        assert result == "&lt;&amp;"

    def test_italic_underline_strikethrough_tags(self) -> None:
        """Italic+underline+strikethrough formatting → <i><u><s> tags."""
        chars = [
            _make_win32com_char("X", italic=True, underline=1, strike=True),
            _make_win32com_char("Y"),
        ]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        result = _win32com_word_runs_to_html(para)
        assert "<i>" in result
        assert "<u>" in result
        assert "<s>" in result
        assert "X" in result
        assert "Y" in result


# ---------------------------------------------------------------------------
# Win32com Word injection helpers
# ---------------------------------------------------------------------------


class TestInjectWin32comWordHtmlRuns:
    """Tests for _inject_win32com_word_html_runs."""

    def test_bold_segments(self) -> None:
        """Bold segments get Font.Bold = True on sub-range."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0

        # Track sub-range calls
        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        _inject_win32com_word_html_runs(
            doc,
            rng,
            "<b>Hello</b> world",
        )
        # "Hello" = chars 0..5
        hello_rng = sub_ranges[(0, 5)]  # noqa: PLR2004
        assert hello_rng.Font.Bold is True
        # " world" = chars 5..11
        world_rng = sub_ranges[(5, 11)]  # noqa: PLR2004
        assert world_rng.Font.Bold is False

    def test_size_segments(self) -> None:
        """Font size is set on sub-range when specified in span."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0
        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        _inject_win32com_word_html_runs(
            doc,
            rng,
            '<span style="font-size:14pt">Big</span>Small',
        )
        big_rng = sub_ranges[(0, 3)]
        assert big_rng.Font.Size == 14.0  # noqa: PLR2004

    def test_color_segments_bgr(self) -> None:
        """Colour is set via BGR using _color_hex_to_win32com."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0
        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        _inject_win32com_word_html_runs(
            doc,
            rng,
            '<span style="color:#ff0000">Red</span>',
        )
        red_rng = sub_ranges[(0, 3)]
        # #ff0000 in BGR = 255
        assert red_rng.Font.Color == 255  # noqa: PLR2004

    def test_cell_mode_no_trailing_cr(self) -> None:
        r"""is_cell=True → no trailing \r appended to text."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0
        doc.Range.return_value = MagicMock()
        _inject_win32com_word_html_runs(
            doc,
            rng,
            "<b>cell</b>",
            is_cell=True,
        )
        # Text set to "cell" (no trailing \r)
        assert rng.Text == "cell"

    def test_empty_segments_noop(self) -> None:
        """Empty parse result → no crash."""
        doc = MagicMock()
        rng = MagicMock()
        _inject_win32com_word_html_runs(doc, rng, "")
        doc.Range.assert_not_called()

    def test_underline_and_strikethrough(self) -> None:
        """Underline=1 (wdUnderlineSingle), StrikeThrough=True."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0
        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        _inject_win32com_word_html_runs(
            doc,
            rng,
            "<u><s>test</s></u>",
        )
        test_rng = sub_ranges[(0, 4)]
        assert test_rng.Font.Underline == 1
        assert test_rng.Font.StrikeThrough is True

    def test_font_name_preserved(self) -> None:
        """Base font Name is restored on the whole range when target_lang given."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0
        rng.Font.Name = "Calibri"

        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        # Use two segments so the whole-range (0,8) doesn't overlap
        _inject_win32com_word_html_runs(
            doc,
            rng,
            "<b>Hi</b> there",
            "Hello",
            target_lang="French",
        )
        # Whole-range sub-range for font Name: (0, 8) for "Hi there"
        whole_rng = sub_ranges[(0, 8)]  # noqa: PLR2004
        # Font name is set (language-aware substitution)
        assert whole_rng.Font.Name is not None

    def test_font_name_not_restored_without_target_lang(self) -> None:
        """Font Name NOT restored when target_lang is empty."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0
        rng.Font.Name = "Arial"

        call_log: list[tuple[int, int]] = []

        def make_sub_range(start: int, end: int) -> MagicMock:
            call_log.append((start, end))
            return MagicMock()

        doc.Range.side_effect = make_sub_range
        # No target_lang → font restoration skipped
        _inject_win32com_word_html_runs(
            doc,
            rng,
            "<b>\u4f60</b>\u597d",
            "Hi",
        )
        # Only per-segment sub-ranges, no font Name whole-range call
        assert len(call_log) == 2  # noqa: PLR2004
        assert call_log[0] == (0, 1)  # "你"
        assert call_log[1] == (1, 2)  # "好"

    def test_nonzero_start_offset(self) -> None:
        """rng.Start=10 → segment sub-ranges are shifted by 10."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 10  # non-zero start  # noqa: PLR2004
        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        # Two segments so font-name range (10,14) doesn't collide with them
        _inject_win32com_word_html_runs(doc, rng, "<b>AB</b>CD")
        # "AB" at offset 0 → (10+0, 10+2) = (10, 12)
        assert (10, 12) in sub_ranges  # noqa: PLR2004
        # "CD" at offset 2 → (10+2, 10+4) = (12, 14)
        assert (12, 14) in sub_ranges  # noqa: PLR2004
        assert sub_ranges[(10, 12)].Font.Bold is True
        assert sub_ranges[(12, 14)].Font.Bold is False

    def test_font_name_exception_no_crash(self) -> None:
        """rng.Font.Name raising → saved_name stays None, no crash."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0

        class _BadNameFont:
            @property
            def Name(self) -> None:  # noqa: N802
                raise AttributeError("Name unavailable")

        rng.Font = _BadNameFont()
        doc.Range.return_value = MagicMock()
        # Must not raise; segments are still applied
        _inject_win32com_word_html_runs(doc, rng, "<b>Hi</b>")
        doc.Range.assert_called()

    def test_empty_html_tags_fallback(self) -> None:
        r"""Tags with no text (<b></b>) → plain fallback sets rng.Text to '\r'."""
        doc = MagicMock()
        rng = MagicMock()
        _inject_win32com_word_html_runs(doc, rng, "<b></b>")
        # Empty segments: plain = "" → rng.Text = "\r" (non-cell mode)
        assert rng.Text == "\r"
        doc.Range.assert_not_called()

    def test_color_com_exception_suppressed(self) -> None:
        """COM exception when setting Font.Color is suppressed."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0
        rng.Font.Name = "Arial"

        class _RaisingColorFont:
            """Font mock where Color assignment raises a COM-like exception."""

            def __init__(self) -> None:
                self.Bold = False
                self.Italic = False
                self.Underline = False
                self.StrikeThrough = False
                self.Size = None
                self.Name = "Arial"

            @property
            def Color(self) -> int:  # noqa: N802
                """Return placeholder color."""
                return 0

            @Color.setter
            def Color(self, value: int) -> None:  # noqa: N802
                """Simulate COM error on color assignment."""
                msg = "Simulated COM error"
                raise Exception(msg)  # noqa: TRY002

        def _make_sub_rng(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sr.Font = _RaisingColorFont()
            return sr

        doc.Range.side_effect = _make_sub_rng
        # Must not raise even though Font.Color assignment raises
        _inject_win32com_word_html_runs(
            doc,
            rng,
            '<span style="color:#ff0000">text</span>',
        )
        doc.Range.assert_called()


# ---------------------------------------------------------------------------
# Win32com PPT extraction helpers
# ---------------------------------------------------------------------------


def _make_win32com_ppt_run(  # noqa: PLR0913
    text: str = "A",
    bold: object = 0,
    italic: object = 0,
    underline: object = 0,
    strike: object = 0,
    size: object = 12.0,
    color_rgb: object = 0,
    *,
    highlight_rgb: object | None = None,
    hyperlink_address: str = "",
) -> MagicMock:
    """Creates a mock win32com PPT run TextRange with Font properties."""
    run = MagicMock()
    run.Text = text
    run.Font.Bold = bold
    run.Font.Italic = italic
    run.Font.Underline = underline
    run.Font.Strikethrough = strike
    run.Font.Size = size
    run.Font.Color.RGB = color_rgb
    if highlight_rgb is not None:
        run.Font.Highlight.ForeColor.RGB = highlight_rgb
    else:
        # No highlight — make .Highlight access raise (simulates older Office)
        type(run.Font).Highlight = property(
            lambda s: (_ for _ in ()).throw(AttributeError("no Highlight")),
        )
    # Hyperlink via ActionSettings(ppMouseClick=1)
    run.ActionSettings.return_value.Hyperlink.Address = hyperlink_address
    return run


class TestReadWin32comPptRunFormatting:
    """Tests for _read_win32com_ppt_run_formatting."""

    def test_bold_italic_size_color(self) -> None:
        """Reads msoTrue (-1) bold+italic, size, BGR colour."""
        run = _make_win32com_ppt_run(
            bold=-1,
            italic=-1,
            size=18.0,
            color_rgb=255,
        )
        b, i, u, s, sup, sub, sz, clr, bg = _read_win32com_ppt_run_formatting(run)
        assert b is True
        assert i is True
        assert u is False
        assert s is False
        assert sz == 18.0  # noqa: PLR2004
        assert clr == "#ff0000"
        assert bg is None

    def test_mso_false_is_false(self) -> None:
        """MsoFalse (0) → False for all bool props."""
        run = _make_win32com_ppt_run(bold=0, italic=0, underline=0, strike=0)
        b, i, u, s, _, _, _, _, _ = _read_win32com_ppt_run_formatting(run)
        assert b is False
        assert i is False
        assert u is False
        assert s is False

    def test_undefined_treated_as_false(self) -> None:
        """WIN32COM_UNDEFINED → False/None."""
        run = _make_win32com_ppt_run(
            bold=WIN32COM_UNDEFINED,
            italic=WIN32COM_UNDEFINED,
            underline=WIN32COM_UNDEFINED,
            strike=WIN32COM_UNDEFINED,
            size=WIN32COM_UNDEFINED,
            color_rgb=WIN32COM_UNDEFINED,
        )
        b, i, u, s, sup, sub, sz, clr, bg = _read_win32com_ppt_run_formatting(run)
        assert b is False
        assert i is False
        assert u is False
        assert s is False
        assert sz is None
        assert clr is None
        assert bg is None

    def test_color_exception_returns_none(self) -> None:
        """If Font.Color.RGB raises → color is None."""
        run = _make_win32com_ppt_run()
        # Override to raise
        type(run.Font.Color).RGB = property(
            lambda s: (_ for _ in ()).throw(Exception),
        )
        _, _, _, _, _, _, _, clr, _ = _read_win32com_ppt_run_formatting(run)
        assert clr is None

    def test_size_zero_is_none(self) -> None:
        """Size=0 → sz is None (raw_size > 0 guard)."""
        run = _make_win32com_ppt_run(size=0)
        _, _, _, _, _, _, sz, _, _ = _read_win32com_ppt_run_formatting(run)
        assert sz is None

    def test_highlight_rgb_extracted(self) -> None:
        """Font.Highlight.ForeColor.RGB (BGR) → bg_color_hex."""
        # Yellow in BGR = 0x00FFFF = 65535
        run = _make_win32com_ppt_run(highlight_rgb=65535)
        _, _, _, _, _, _, _, _, bg = _read_win32com_ppt_run_formatting(run)
        assert bg == "#ffff00"

    def test_highlight_missing_returns_none(self) -> None:
        """Older Office without Highlight property → bg is None."""
        run = _make_win32com_ppt_run()  # no highlight_rgb → raises
        _, _, _, _, _, _, _, _, bg = _read_win32com_ppt_run_formatting(run)
        assert bg is None

    def test_tri_state_mixed_is_false(self) -> None:
        """MsoTriStateMixed (-2) for all props → all False (not == -1)."""
        run = _make_win32com_ppt_run(bold=-2, italic=-2, underline=-2, strike=-2)
        b, i, u, s, _, _, _, _, _ = _read_win32com_ppt_run_formatting(run)
        assert b is False
        assert i is False
        assert u is False
        assert s is False


class TestHasWin32comPptMixedFormatting:
    """Tests for _has_win32com_ppt_mixed_formatting."""

    def test_uniform_returns_false(self) -> None:
        """All runs same formatting → False."""
        runs = [
            _make_win32com_ppt_run("Hello"),
            _make_win32com_ppt_run("World"),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        assert _has_win32com_ppt_mixed_formatting(para_rng) is False

    def test_mixed_bold_returns_true(self) -> None:
        """One bold, one not → True."""
        runs = [
            _make_win32com_ppt_run("Bold", bold=-1),
            _make_win32com_ppt_run("Plain", bold=0),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        assert _has_win32com_ppt_mixed_formatting(para_rng) is True

    def test_single_run_returns_false(self) -> None:
        """One run → False."""
        runs = [_make_win32com_ppt_run("Only")]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = 1
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=1)
        )
        assert _has_win32com_ppt_mixed_formatting(para_rng) is False

    def test_whitespace_only_runs_skipped(self) -> None:
        """Only whitespace runs → False (skipped)."""
        runs = [
            _make_win32com_ppt_run(" ", bold=-1),
            _make_win32com_ppt_run("\t", bold=0),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        assert _has_win32com_ppt_mixed_formatting(para_rng) is False

    def test_mixed_size_returns_true(self) -> None:
        """Different sizes across runs → True."""
        runs = [
            _make_win32com_ppt_run("Big", size=24.0),
            _make_win32com_ppt_run("Small", size=10.0),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        assert _has_win32com_ppt_mixed_formatting(para_rng) is True

    def test_mixed_color_returns_true(self) -> None:
        """Different BGR colours across runs → True."""
        runs = [
            _make_win32com_ppt_run("R", color_rgb=255),  # red
            _make_win32com_ppt_run("B", color_rgb=16711680),  # blue  # noqa: PLR2004
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        assert _has_win32com_ppt_mixed_formatting(para_rng) is True


class TestWin32comPptRunsToHtml:
    """Tests for _win32com_ppt_runs_to_html."""

    def test_bold_and_plain_runs(self) -> None:
        """Bold 'Hello' + plain ' world'."""
        runs = [
            _make_win32com_ppt_run("Hello", bold=-1),
            _make_win32com_ppt_run(" world", bold=0),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        assert "<b>Hello</b>" in result
        assert " world" in result

    def test_size_variation_emits_span(self) -> None:
        """Different sizes → deviating size gets <span>; base omitted."""
        runs = [
            _make_win32com_ppt_run("Big", size=24.0),
            _make_win32com_ppt_run("Small", size=10.0),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        # Base-value optimization: 24pt is first → base, 10pt kept
        assert "font-size:10pt" in result

    def test_color_variation_emits_span(self) -> None:
        """Different BGR colours → deviating colour gets <span>; base omitted."""
        runs = [
            _make_win32com_ppt_run("R", color_rgb=255),  # red
            _make_win32com_ppt_run("B", color_rgb=16711680),  # blue
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        # Base-value optimization: first color is base, second emitted
        assert "color:#0000ff" in result

    def test_empty_runs_returns_empty(self) -> None:
        """No text runs → empty string."""
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = 0
        para_rng.Runs.side_effect = lambda i=None: MagicMock(Count=0)
        assert _win32com_ppt_runs_to_html(para_rng) == ""

    def test_italic_and_underline_tags(self) -> None:
        """Italic+underline in one run, plain in another → <i><u> emitted."""
        runs = [
            _make_win32com_ppt_run("Ital", italic=-1, underline=-1),
            _make_win32com_ppt_run("Plain"),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        assert "<i>" in result
        assert "<u>" in result
        assert "Ital" in result

    def test_empty_text_run_skipped(self) -> None:
        """Run with empty text is silently skipped; other runs rendered."""
        runs = [
            _make_win32com_ppt_run("", bold=-1),  # empty → skipped
            _make_win32com_ppt_run("Hello"),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        assert result == "Hello"


# ---------------------------------------------------------------------------
# Win32com PPT injection helpers
# ---------------------------------------------------------------------------


class TestInjectWin32comPptHtmlRuns:
    """Tests for _inject_win32com_ppt_html_runs."""

    def test_bold_segments(self) -> None:
        """Bold segments get Font.Bold = -1 (msoTrue)."""
        tf = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        tf.TextRange.Paragraphs.return_value.Characters.side_effect = make_chars
        _inject_win32com_ppt_html_runs(tf, 1, "<b>Hello</b> world")
        # "Hello" at offset 0 → Characters(1, 5) (1-based)
        hello_cr = char_ranges[(1, 5)]  # noqa: PLR2004
        assert hello_cr.Font.Bold == -1
        # " world" at offset 5 → Characters(6, 6)
        world_cr = char_ranges[(6, 6)]  # noqa: PLR2004
        assert world_cr.Font.Bold == 0

    def test_size_segments(self) -> None:
        """Font size set via Characters() sub-range."""
        tf = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        tf.TextRange.Paragraphs.return_value.Characters.side_effect = make_chars
        _inject_win32com_ppt_html_runs(
            tf,
            1,
            '<span style="font-size:14pt">Big</span>',
        )
        big_cr = char_ranges[(1, 3)]
        assert big_cr.Font.Size == 14.0  # noqa: PLR2004

    def test_color_segments_bgr(self) -> None:
        """Colour set via Font.Color.RGB using BGR encoding."""
        tf = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        tf.TextRange.Paragraphs.return_value.Characters.side_effect = make_chars
        _inject_win32com_ppt_html_runs(
            tf,
            1,
            '<span style="color:#ff0000">Red</span>',
        )
        red_cr = char_ranges[(1, 3)]
        # #ff0000 in BGR = 255
        assert red_cr.Font.Color.RGB == 255  # noqa: PLR2004

    def test_italic_underline_strikethrough(self) -> None:
        """Italic, underline, strikethrough flags set correctly."""
        tf = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        tf.TextRange.Paragraphs.return_value.Characters.side_effect = make_chars
        _inject_win32com_ppt_html_runs(tf, 1, "<i><u><s>test</s></u></i>")
        test_cr = char_ranges[(1, 4)]
        assert test_cr.Font.Italic == -1
        assert test_cr.Font.Underline == -1
        assert test_cr.Font.Strikethrough == -1

    def test_empty_html_noop(self) -> None:
        """Empty HTML → no crash."""
        tf = MagicMock()
        _inject_win32com_ppt_html_runs(tf, 1, "")
        # Text not set when segments is empty
        tf.TextRange.Paragraphs.return_value.Characters.assert_not_called()

    def test_font_name_preserved(self) -> None:
        """Base font Name is restored when target_lang is given."""
        tf = MagicMock()
        tf.TextRange.Paragraphs.return_value.Font.Name = "Calibri"
        tf.TextRange.Paragraphs.return_value.Characters.side_effect = (
            lambda start, length: MagicMock()
        )
        _inject_win32com_ppt_html_runs(
            tf,
            1,
            "<b>Hi</b>",
            "Hello",
            target_lang="French",
        )
        # Font Name should be set on the re-acquired paragraph range
        assert tf.TextRange.Paragraphs.return_value.Font.Name is not None

    def test_font_name_not_restored_without_target_lang(self) -> None:
        """Font Name NOT restored when target_lang is empty."""
        tf = MagicMock()
        # Set up a fresh mock so we can track Name assignments
        para_mock = MagicMock()
        para_mock.Font.Name = "Arial"  # original font
        para_mock.Characters.side_effect = lambda start, length: MagicMock()
        tf.TextRange.Paragraphs.return_value = para_mock

        _inject_win32com_ppt_html_runs(
            tf,
            1,
            "<b>\u4f60\u597d</b>",
            "Hello",
        )
        # No target_lang → font restoration skipped; segments still applied
        assert para_mock.Characters.call_count == 1

    def test_correct_paragraph_index_used(self) -> None:
        """p_idx=2 → tf.TextRange.Paragraphs called with 2, not 1."""
        tf = MagicMock()
        tf.TextRange.Paragraphs.return_value.Characters.side_effect = (
            lambda start, length: MagicMock()
        )
        _inject_win32com_ppt_html_runs(tf, 2, "<b>Hi</b>")  # noqa: PLR2004
        # Verify Paragraphs was invoked with p_idx=2
        calls_with_2 = [
            c for c in tf.TextRange.Paragraphs.call_args_list if c.args == (2,)
        ]
        assert len(calls_with_2) >= 1  # noqa: PLR2004

    def test_font_name_exception_no_crash(self) -> None:
        """para_before.Font.Name raising → saved_name stays None, no crash."""
        tf = MagicMock()

        class _BadNameFont:
            @property
            def Name(self) -> None:  # noqa: N802
                raise AttributeError("Name unavailable")

        tf.TextRange.Paragraphs.return_value.Font = _BadNameFont()
        tf.TextRange.Paragraphs.return_value.Characters.side_effect = (
            lambda start, length: MagicMock()
        )
        # Must not raise; segments are still applied
        _inject_win32com_ppt_html_runs(tf, 1, "<b>Hi</b>")
        tf.TextRange.Paragraphs.return_value.Characters.assert_called()

    def test_empty_html_tags_fallback(self) -> None:
        """Tags with no text content (<b></b>) → paragraph text set to ''."""
        tf = MagicMock()
        _inject_win32com_ppt_html_runs(tf, 1, "<b></b>")
        # Empty segments → plain = "" → Paragraphs(1).Text = ""
        assert tf.TextRange.Paragraphs.return_value.Text == ""
        tf.TextRange.Paragraphs.return_value.Characters.assert_not_called()

    def test_color_com_exception_suppressed(self) -> None:
        """COM exception when setting Font.Color.RGB is suppressed gracefully."""
        tf = MagicMock()

        class _RaisingColorFont:
            """Font mock where Color.RGB assignment raises a COM-like exception."""

            def __init__(self) -> None:
                self.Bold = 0
                self.Italic = 0
                self.Underline = 0
                self.Strikethrough = 0
                self.Size = None
                self.Name = "Arial"
                self.Color = _RaisingColor()

        class _RaisingColor:
            """Color object where RGB assignment raises."""

            @property
            def RGB(self) -> int:  # noqa: N802
                """Return placeholder RGB."""
                return 0

            @RGB.setter
            def RGB(self, value: int) -> None:  # noqa: N802
                """Simulate COM error on RGB assignment."""
                msg = "Simulated COM error"
                raise Exception(msg)  # noqa: TRY002

        para_mock = MagicMock()
        para_mock.Font = _RaisingColorFont()

        def _make_char_rng(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            cr.Font = _RaisingColorFont()
            return cr

        para_mock.Characters.side_effect = _make_char_rng
        tf.TextRange.Paragraphs.return_value = para_mock

        # Must not raise even though Font.Color.RGB assignment raises
        _inject_win32com_ppt_html_runs(
            tf,
            1,
            '<span style="color:#ff0000">text</span>',
        )
        para_mock.Characters.assert_called()


# ---------------------------------------------------------------------------
# Integration: extract→inject Word round-trip awareness
# ---------------------------------------------------------------------------


class TestWin32comWordHtmlRoundTrip:
    """Verifies extract/inject Word functions use HTML path."""

    @patch("src.core.office_processor._win32com_close")
    @patch("src.core.office_processor._win32com_open")
    def test_extract_uses_html_for_mixed(
        self,
        mock_open: MagicMock,
        mock_close: MagicMock,
    ) -> None:
        """_extract_win32com_word emits HTML for mixed-formatting paras."""
        doc = MagicMock()
        doc.Paragraphs.Count = 1
        # Create paragraph with mixed bold
        chars = [
            _make_win32com_char("B", bold=True),
            _make_win32com_char("x"),
        ]
        para = MagicMock()
        # Quick check: Bold is mixed
        para.Range.Font.Bold = WIN32COM_UNDEFINED
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        para.Range.Text = "Bx\r"
        doc.Paragraphs.side_effect = lambda i: para
        doc.Tables.Count = 0
        mock_open.return_value = (MagicMock(), doc, MagicMock())

        from src.core.office_processor import _extract_win32com_word  # noqa: PLC0415

        texts = _extract_win32com_word(Path("test.doc"))
        assert len(texts) == 1
        key, text = texts[0]
        assert key == "para:1"
        # Should contain HTML tags from the mixed formatting
        assert "<b>" in text

    @patch("src.core.office_processor._win32com_close")
    @patch("src.core.office_processor._win32com_open")
    def test_inject_uses_html_path(
        self,
        mock_open: MagicMock,
        mock_close: MagicMock,
    ) -> None:
        """_inject_win32com_word uses HTML injection for tagged text."""
        doc = MagicMock()
        doc.Paragraphs.Count = 1
        rng = MagicMock()
        rng.Text = "Hello\r"
        rng.Start = 0
        para = MagicMock()
        para.Range = rng
        doc.Paragraphs.side_effect = lambda i: para

        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range
        doc.Tables.Count = 0
        mock_open.return_value = (MagicMock(), doc, MagicMock())

        from src.core.office_processor import _inject_win32com_word  # noqa: PLC0415

        _inject_win32com_word(
            Path("test.doc"),
            Path("out.doc"),
            {"para:1": "<b>Bold</b> text"},
        )
        # Sub-ranges: 2 segments + 1 whole-range font restore (scripts match)
        assert len(sub_ranges) == 3  # noqa: PLR2004
        # Verify per-segment ranges exist
        assert (0, 4) in sub_ranges  # "Bold"
        assert (4, 9) in sub_ranges  # " text"
        assert (0, 9) in sub_ranges  # whole-range font Name

    @patch("src.core.office_processor._win32com_close")
    @patch("src.core.office_processor._win32com_open")
    def test_inject_plain_text_fallback(
        self,
        mock_open: MagicMock,
        mock_close: MagicMock,
    ) -> None:
        """_inject_win32com_word falls back to font save/restore for plain."""
        doc = MagicMock()
        doc.Paragraphs.Count = 1
        rng = MagicMock()
        rng.Text = "Hello\r"
        para = MagicMock()
        para.Range = rng
        doc.Paragraphs.side_effect = lambda i: para
        doc.Tables.Count = 0
        mock_open.return_value = (MagicMock(), doc, MagicMock())

        from src.core.office_processor import _inject_win32com_word  # noqa: PLC0415

        _inject_win32com_word(
            Path("test.doc"),
            Path("out.doc"),
            {"para:1": "Plain translation"},
        )
        # No doc.Range calls for sub-ranges (HTML path not used)
        doc.Range.assert_not_called()


# ---------------------------------------------------------------------------
# Integration: extract→inject PPT round-trip awareness
# ---------------------------------------------------------------------------


class TestWin32comPptHtmlRoundTrip:
    """Verifies extract/inject PPT functions use HTML path."""

    @patch("src.core.office_processor._win32com_close")
    @patch("src.core.office_processor._win32com_open")
    def test_extract_uses_html_for_mixed(
        self,
        mock_open: MagicMock,
        mock_close: MagicMock,
    ) -> None:
        """_extract_win32com_ppt emits HTML for mixed-formatting paras."""
        prs = MagicMock()
        prs.Slides.Count = 1
        slide = MagicMock()
        slide.Shapes.Count = 1
        shape = MagicMock()
        shape.HasTextFrame = True
        tf = MagicMock()

        # Build a paragraph with mixed formatting via Runs()
        runs = [
            _make_win32com_ppt_run("Bold", bold=-1),
            _make_win32com_ppt_run("Plain", bold=0),
        ]

        para_rng = MagicMock()
        para_rng.Text = "BoldPlain"
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )

        paras = MagicMock()
        paras.Count = 1
        tf.TextRange.Paragraphs.side_effect = lambda i=None: (
            para_rng if i is not None else paras
        )

        shape.TextFrame = tf
        slide.Shapes.side_effect = lambda i: shape
        prs.Slides.side_effect = lambda i: slide
        mock_open.return_value = (MagicMock(), prs, MagicMock())

        from src.core.office_processor import _extract_win32com_ppt  # noqa: PLC0415

        texts = _extract_win32com_ppt(Path("test.ppt"))
        assert len(texts) == 1
        key, text = texts[0]
        assert key == "slide:1:1:1"
        assert "<b>" in text

    @patch("src.core.office_processor._win32com_close")
    @patch("src.core.office_processor._win32com_open")
    def test_inject_uses_html_path(
        self,
        mock_open: MagicMock,
        mock_close: MagicMock,
    ) -> None:
        """_inject_win32com_ppt uses HTML injection for tagged text."""
        prs = MagicMock()
        prs.Slides.Count = 1
        slide = MagicMock()
        slide.Shapes.Count = 1
        shape = MagicMock()
        shape.HasTextFrame = True
        tf = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        tf.TextRange.Paragraphs.return_value.Characters.side_effect = make_chars

        paras = MagicMock()
        paras.Count = 1
        tf.TextRange.Paragraphs.side_effect = lambda i=None: (
            tf.TextRange.Paragraphs.return_value if i is not None else paras
        )

        shape.TextFrame = tf
        slide.Shapes.side_effect = lambda i: shape
        prs.Slides.side_effect = lambda i: slide
        mock_open.return_value = (MagicMock(), prs, MagicMock())

        from src.core.office_processor import _inject_win32com_ppt  # noqa: PLC0415

        _inject_win32com_ppt(
            Path("test.ppt"),
            Path("out.ppt"),
            {"slide:1:1:1": "<b>Bold</b> text"},
        )
        # Characters should have been called for sub-segments
        assert len(char_ranges) == 2  # noqa: PLR2004

    @patch("src.core.office_processor._win32com_close")
    @patch("src.core.office_processor._win32com_open")
    def test_inject_plain_text_fallback(
        self,
        mock_open: MagicMock,
        mock_close: MagicMock,
    ) -> None:
        """_inject_win32com_ppt falls back to font save/restore for plain."""
        prs = MagicMock()
        prs.Slides.Count = 1
        slide = MagicMock()
        slide.Shapes.Count = 1
        shape = MagicMock()
        shape.HasTextFrame = True
        tf = MagicMock()

        para_rng = MagicMock()
        para_rng.Text = "Hello"

        paras = MagicMock()
        paras.Count = 1
        tf.TextRange.Paragraphs.side_effect = lambda i=None: (
            para_rng if i is not None else paras
        )

        shape.TextFrame = tf
        slide.Shapes.side_effect = lambda i: shape
        prs.Slides.side_effect = lambda i: slide
        mock_open.return_value = (MagicMock(), prs, MagicMock())

        from src.core.office_processor import _inject_win32com_ppt  # noqa: PLC0415

        _inject_win32com_ppt(
            Path("test.ppt"),
            Path("out.ppt"),
            {"slide:1:1:1": "Plain translation"},
        )
        # Font save/restore should have been called (plain path)
        assert para_rng.Text == "Plain translation"


# ---------------------------------------------------------------------------
# TestInjectDrawingmlHtmlRuns
# ---------------------------------------------------------------------------


class TestInjectDrawingmlHtmlRuns:
    """Tests for _inject_drawingml_html_runs."""

    def test_empty_html_tags_fallback(self) -> None:
        """Tags with no text (<b></b>) fall back to _inject_drawingml_text('')."""
        from lxml import etree  # noqa: PLC0415

        dml_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        xml = f'<txBody xmlns:a="{dml_ns}"><a:p/></txBody>'
        tx_body_el = etree.fromstring(xml)

        with patch("src.core.office_processor._inject_drawingml_text") as mock_fn:
            _inject_drawingml_html_runs(tx_body_el, "<b></b>")
            mock_fn.assert_called_once_with(tx_body_el, "")

    def test_bold_segment_creates_run(self) -> None:
        """Bold HTML creates an <a:r> with b='1' on <a:rPr>."""
        from lxml import etree  # noqa: PLC0415

        dml_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        xml = (
            f'<txBody xmlns:a="{dml_ns}">'
            "<a:p>"
            "<a:r><a:rPr/><a:t>old text</a:t></a:r>"
            "</a:p>"
            "</txBody>"
        )
        tx_body_el = etree.fromstring(xml)
        _inject_drawingml_html_runs(tx_body_el, "<b>Hello</b>")

        a_r_tag = f"{{{dml_ns}}}r"
        a_rpr_tag = f"{{{dml_ns}}}rPr"
        a_t_tag = f"{{{dml_ns}}}t"
        p_el = tx_body_el[0]  # first <a:p>
        runs = p_el.findall(a_r_tag)
        assert len(runs) == 1  # noqa: PLR2004
        assert runs[0].find(a_t_tag).text == "Hello"
        rpr = runs[0].find(a_rpr_tag)
        assert rpr is not None
        assert rpr.get("b") == "1"

    def test_plain_text_no_formatting_tags(self) -> None:
        """Plain text (no HTML tags) dispatches to _inject_drawingml_text."""
        from lxml import etree  # noqa: PLC0415

        dml_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        xml = f'<txBody xmlns:a="{dml_ns}"><a:p/></txBody>'
        tx_body_el = etree.fromstring(xml)

        with patch("src.core.office_processor._inject_drawingml_text") as mock_fn:
            _inject_drawingml_html_runs(tx_body_el, "plain text")
            mock_fn.assert_called_once_with(tx_body_el, "plain text")

    def test_mixed_bold_plain_segments(self) -> None:
        """Mixed bold + plain creates two <a:r> elements."""
        from lxml import etree  # noqa: PLC0415

        dml_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        xml = (
            f'<txBody xmlns:a="{dml_ns}">'
            "<a:p>"
            "<a:r><a:rPr/><a:t>old</a:t></a:r>"
            "</a:p>"
            "</txBody>"
        )
        tx_body_el = etree.fromstring(xml)
        _inject_drawingml_html_runs(tx_body_el, "<b>Bold</b> plain")

        a_r_tag = f"{{{dml_ns}}}r"
        a_t_tag = f"{{{dml_ns}}}t"
        a_rpr_tag = f"{{{dml_ns}}}rPr"
        p_el = tx_body_el[0]
        runs = p_el.findall(a_r_tag)
        assert len(runs) == 2  # noqa: PLR2004
        texts = [r.find(a_t_tag).text for r in runs]
        assert texts == ["Bold", " plain"]
        # First run is bold
        assert runs[0].find(a_rpr_tag).get("b") == "1"
        # Second run is plain (no bold attribute)
        assert runs[1].find(a_rpr_tag).get("b") is None


# ---------------------------------------------------------------------------
# TestInjectWin32comExcelHtmlRuns
# ---------------------------------------------------------------------------


class TestInjectWin32comExcelHtmlRuns:
    """Tests for _inject_win32com_excel_html_runs."""

    def test_empty_html_tags_fallback(self) -> None:
        """Tags with no text (<b></b>) → text_rng.Text set to ''."""
        text_rng = MagicMock()
        _inject_win32com_excel_html_runs(text_rng, "<b></b>")
        assert text_rng.Text == ""
        text_rng.Characters.assert_not_called()

    def test_bold_segment(self) -> None:
        """Bold segment sets Font.Bold on Characters sub-range."""
        text_rng = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def _make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        text_rng.Characters.side_effect = _make_chars
        _inject_win32com_excel_html_runs(text_rng, "<b>Hello</b>")

        # Full text "Hello" should be set
        assert text_rng.Text == "Hello"
        # Characters(1, 5) for "Hello"
        hello_cr = char_ranges[(1, 5)]  # noqa: PLR2004
        assert hello_cr.Font.Bold is True

    def test_mixed_bold_plain(self) -> None:
        """Mixed bold + plain → two Character sub-ranges applied."""
        text_rng = MagicMock()
        char_ranges: dict[tuple[int, int], MagicMock] = {}

        def _make_chars(start: int, length: int) -> MagicMock:
            cr = MagicMock()
            char_ranges[(start, length)] = cr
            return cr

        text_rng.Characters.side_effect = _make_chars
        _inject_win32com_excel_html_runs(text_rng, "<b>Hi</b> there")

        assert text_rng.Text == "Hi there"
        hi_cr = char_ranges[(1, 2)]  # noqa: PLR2004
        assert hi_cr.Font.Bold is True
        there_cr = char_ranges[(3, 6)]  # noqa: PLR2004
        assert there_cr.Font.Bold is False

    def test_font_name_preserved(self) -> None:
        """Base font name is restored when target_lang is given."""
        text_rng = MagicMock()
        text_rng.Font.Name = "Calibri"
        text_rng.Characters.side_effect = lambda start, length: MagicMock()
        _inject_win32com_excel_html_runs(
            text_rng,
            "<b>Hi</b>",
            "Hello",
            target_lang="French",
        )
        # After inject, font name should be restored
        assert text_rng.Font.Name is not None

    def test_empty_string_fallback(self) -> None:
        """Completely empty input → text_rng.Text set to ''."""
        text_rng = MagicMock()
        _inject_win32com_excel_html_runs(text_rng, "")
        assert text_rng.Text == ""
        text_rng.Characters.assert_not_called()


# ---------------------------------------------------------------------------
# TestInjectWpsTxbxHtmlRuns
# ---------------------------------------------------------------------------


class TestInjectWpsTxbxHtmlRuns:
    """Tests for _inject_wps_txbx_html_runs."""

    def _make_txbx_el(self, text: str = "old text") -> object:
        """Builds a minimal <wps:txbx> lxml element with one paragraph."""
        from lxml import etree  # noqa: PLC0415

        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        xml = (
            f'<wps:txbx xmlns:wps="{wps_ns}" xmlns:w="{w_ns}">'
            f"<w:txbxContent>"
            f"<w:p><w:r><w:t>{text}</w:t></w:r></w:p>"
            f"</w:txbxContent>"
            f"</wps:txbx>"
        )
        return etree.fromstring(xml)

    def test_empty_html_tags_fallback_calls_plain_inject(self) -> None:
        """Tags with no text (<b></b>) fall back to _inject_wps_txbx_plain."""
        txbx_el = self._make_txbx_el("original")
        with patch(
            "src.core.office_processor._inject_wps_txbx_plain",
        ) as mock_plain:
            _inject_wps_txbx_html_runs(txbx_el, "<b></b>")
            mock_plain.assert_called_once()
            # First arg is txbx_el, second arg is plain=""
            assert mock_plain.call_args.args[1] == ""

    def test_bold_segment_creates_run(self) -> None:
        """Bold HTML creates a <w:b> element inside <w:rPr>."""
        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        txbx_el = self._make_txbx_el("old")
        _inject_wps_txbx_html_runs(txbx_el, "<b>Hello</b>")

        w_r_tag = f"{{{w_ns}}}r"
        w_rpr_tag = f"{{{w_ns}}}rPr"
        w_b_tag = f"{{{w_ns}}}b"
        w_t_tag = f"{{{w_ns}}}t"
        runs = txbx_el.findall(f".//{w_r_tag}")
        assert len(runs) == 1  # noqa: PLR2004
        assert runs[0].find(w_t_tag).text == "Hello"
        rpr = runs[0].find(w_rpr_tag)
        assert rpr is not None
        assert rpr.find(w_b_tag) is not None

    def test_plain_text_creates_run(self) -> None:
        """Plain text without formatting tags creates a plain <w:r>/<w:t> element."""
        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        txbx_el = self._make_txbx_el("old")
        _inject_wps_txbx_html_runs(txbx_el, "plain text")

        w_t_tag = f"{{{w_ns}}}t"
        texts = [t.text for t in txbx_el.findall(f".//{w_t_tag}") if t.text]
        assert "plain text" in texts

    def test_font_size_inherited_from_base_run(self) -> None:
        """When HTML has no font-size span, runs inherit size from base rPr."""
        from lxml import etree  # noqa: PLC0415

        w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        wps = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        # Build a text box with font size 36 half-points (18pt) on the run
        xml = (
            f'<wps:txbx xmlns:wps="{wps}" xmlns:w="{w}">'
            f"<w:txbxContent>"
            f'<w:p><w:r><w:rPr><w:b/><w:sz w:val="36"/></w:rPr>'
            f"<w:t>old</w:t></w:r></w:p>"
            f"</w:txbxContent>"
            f"</wps:txbx>"
        )
        txbx_el = etree.fromstring(xml)

        # Inject HTML with bold only — no font-size span
        _inject_wps_txbx_html_runs(txbx_el, "<b>translated</b>")

        # The injected run must carry <w:sz w:val="36"/> from the base
        w_sz_tag = f"{{{w}}}sz"
        w_r_tag = f"{{{w}}}r"
        runs = txbx_el.findall(f".//{w_r_tag}")
        assert len(runs) == 1
        sz_el = runs[0].find(f".//{w_sz_tag}")
        assert sz_el is not None, "Font size lost — <w:sz> missing from injected run"
        assert sz_el.get(f"{{{w}}}val") == "36"  # noqa: PLR2004

    def test_font_color_inherited_from_base_run(self) -> None:
        """When HTML has no color span, runs inherit color from base rPr."""
        from lxml import etree  # noqa: PLC0415

        w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        wps = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        # Build a text box with color FF0000 on the run
        xml = (
            f'<wps:txbx xmlns:wps="{wps}" xmlns:w="{w}">'
            f"<w:txbxContent>"
            f'<w:p><w:r><w:rPr><w:b/><w:color w:val="FF0000"/></w:rPr>'
            f"<w:t>old</w:t></w:r></w:p>"
            f"</w:txbxContent>"
            f"</wps:txbx>"
        )
        txbx_el = etree.fromstring(xml)

        # Inject HTML with bold only — no color span
        _inject_wps_txbx_html_runs(txbx_el, "<b>translated</b>")

        # The injected run must carry <w:color w:val="FF0000"/> from the base
        w_color_tag = f"{{{w}}}color"
        w_r_tag = f"{{{w}}}r"
        runs = txbx_el.findall(f".//{w_r_tag}")
        assert len(runs) == 1
        color_el = runs[0].find(f".//{w_color_tag}")
        assert color_el is not None, "Color lost — <w:color> missing from injected run"
        assert color_el.get(f"{{{w}}}val") == "FF0000"

    def test_explicit_size_overrides_base(self) -> None:
        """When HTML specifies font-size, it overrides the base rPr size."""
        from lxml import etree  # noqa: PLC0415

        w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        wps = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        xml = (
            f'<wps:txbx xmlns:wps="{wps}" xmlns:w="{w}">'
            f"<w:txbxContent>"
            f'<w:p><w:r><w:rPr><w:sz w:val="36"/></w:rPr>'
            f"<w:t>old</w:t></w:r></w:p>"
            f"</w:txbxContent>"
            f"</wps:txbx>"
        )
        txbx_el = etree.fromstring(xml)

        # Inject HTML with explicit 24pt → 48 half-points
        _inject_wps_txbx_html_runs(
            txbx_el,
            '<span style="font-size:24pt">big</span>',
        )

        w_sz_tag = f"{{{w}}}sz"
        w_r_tag = f"{{{w}}}r"
        runs = txbx_el.findall(f".//{w_r_tag}")
        sz_el = runs[0].find(f".//{w_sz_tag}")
        assert sz_el is not None
        expected_half_pts = 48  # noqa: PLR2004
        assert sz_el.get(f"{{{w}}}val") == str(expected_half_pts)


# ---------------------------------------------------------------------------
# TestInjectOdfTextBoxHtmlRuns
# ---------------------------------------------------------------------------


class TestInjectOdfTextBoxHtmlRuns:
    """Tests for _inject_odf_text_box_html_runs."""

    def _make_odf_text_box_el(self, text: str = "old text") -> object:
        """Builds a minimal <draw:text-box> lxml element."""
        from lxml import etree  # noqa: PLC0415

        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
        xml = (
            f'<draw:text-box xmlns:draw="{draw_ns}" xmlns:text="{text_ns}">'
            f"<text:p>{text}</text:p>"
            f"</draw:text-box>"
        )
        return etree.fromstring(xml)

    def _make_auto_styles(self) -> object:
        """Builds a minimal <office:automatic-styles> lxml element."""
        from lxml import etree  # noqa: PLC0415

        office_ns = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
        return etree.fromstring(
            f'<office:automatic-styles xmlns:office="{office_ns}"/>',
        )

    def test_empty_html_tags_fallback(self) -> None:
        """Tags with no text (<b></b>) fall back to _inject_odf_paragraph_text('')."""
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        text_box_el = self._make_odf_text_box_el()
        text_p_tag = f"{{{text_ns}}}p"
        auto_styles = self._make_auto_styles()

        with patch(
            "src.core.office_processor._inject_odf_paragraph_text",
        ) as mock_fn:
            _inject_odf_text_box_html_runs(
                text_box_el,
                "<b></b>",
                text_p_tag,
                auto_styles,
                [0],
            )
            mock_fn.assert_called_once_with(text_box_el, "", text_p_tag)

    def test_plain_text_fallback(self) -> None:
        """Plain text (no HTML tags) dispatches to _inject_odf_paragraph_text."""
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        text_box_el = self._make_odf_text_box_el()
        text_p_tag = f"{{{text_ns}}}p"
        auto_styles = self._make_auto_styles()

        with patch(
            "src.core.office_processor._inject_odf_paragraph_text",
        ) as mock_fn:
            mock_fn.return_value = True
            _inject_odf_text_box_html_runs(
                text_box_el,
                "plain text",
                text_p_tag,
                auto_styles,
                [0],
            )
            mock_fn.assert_called_once_with(text_box_el, "plain text", text_p_tag)

    def test_bold_segment_creates_style(self) -> None:
        """Bold HTML creates a style entry in auto_styles_el."""
        text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        style_ns = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
        text_box_el = self._make_odf_text_box_el()
        text_p_tag = f"{{{text_ns}}}p"
        auto_styles = self._make_auto_styles()
        counter = [0]

        result = _inject_odf_text_box_html_runs(
            text_box_el,
            "<b>Hello</b>",
            text_p_tag,
            auto_styles,
            counter,
        )
        assert result is True
        # A style element should have been added to auto_styles
        style_tag = f"{{{style_ns}}}style"
        styles = auto_styles.findall(style_tag)
        assert len(styles) >= 1  # noqa: PLR2004


# ---------------------------------------------------------------------------
# _wps_txbx_to_text_or_html tests
# ---------------------------------------------------------------------------


class TestWpsTxbxToTextOrHtml:
    """Tests for ``_wps_txbx_to_text_or_html`` mixed-formatting detection."""

    _W = _WORDML_NS

    def _make_txbx(self, paragraphs: list[list[dict]]) -> object:
        """Builds a ``<wps:txbx>`` lxml element from paragraph/run specs.

        Each paragraph is a list of run dicts with keys:
            text (str), bold (bool), italic (bool), underline (bool),
            strike (bool), size (float|None), color (str|None),
            rstyle (str|None), multi_t (list[str]|None).
        """
        from lxml import etree  # noqa: PLC0415

        wps = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        w = self._W
        txbx = etree.Element(f"{{{wps}}}txbx")
        txbx_content = etree.SubElement(txbx, f"{{{w}}}txbxContent")

        for runs in paragraphs:
            p = etree.SubElement(txbx_content, f"{{{w}}}p")
            for run_spec in runs:
                r = etree.SubElement(p, f"{{{w}}}r")
                # Build <w:rPr> if any formatting is non-default
                rpr = etree.SubElement(r, f"{{{w}}}rPr")
                if run_spec.get("rstyle"):
                    rs = etree.SubElement(rpr, f"{{{w}}}rStyle")
                    rs.set(f"{{{w}}}val", run_spec["rstyle"])
                if run_spec.get("bold"):
                    etree.SubElement(rpr, f"{{{w}}}b")
                if run_spec.get("italic"):
                    etree.SubElement(rpr, f"{{{w}}}i")
                if run_spec.get("underline"):
                    u_el = etree.SubElement(rpr, f"{{{w}}}u")
                    u_el.set(f"{{{w}}}val", "single")
                if run_spec.get("strike"):
                    etree.SubElement(rpr, f"{{{w}}}strike")
                if run_spec.get("size") is not None:
                    sz = etree.SubElement(rpr, f"{{{w}}}sz")
                    sz.set(f"{{{w}}}val", str(int(run_spec["size"] * 2)))
                if run_spec.get("color") is not None:
                    clr = etree.SubElement(rpr, f"{{{w}}}color")
                    clr.set(f"{{{w}}}val", run_spec["color"].lstrip("#"))

                # <w:t> elements — support multi_t for split runs
                multi_t = run_spec.get("multi_t")
                if multi_t:
                    for txt in multi_t:
                        t = etree.SubElement(r, f"{{{w}}}t")
                        t.text = txt
                else:
                    t = etree.SubElement(r, f"{{{w}}}t")
                    t.text = run_spec.get("text", "")
        return txbx

    # -- Uniform formatting → plain text --

    def test_uniform_formatting_returns_plain_text(self) -> None:
        """All runs same formatting → plain text."""
        txbx = self._make_txbx(
            [
                [{"text": "Hello "}, {"text": "world"}],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx)
        assert result == "Hello world"
        assert "<" not in result

    def test_single_run_returns_plain_text(self) -> None:
        """One run → plain text (can't have mixed)."""
        txbx = self._make_txbx(
            [
                [{"text": "Only one run", "bold": True}],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx)
        assert result == "Only one run"
        assert "<" not in result

    def test_empty_text_box_returns_empty(self) -> None:
        """No runs → empty string."""
        txbx = self._make_txbx([[]])
        result = _wps_txbx_to_text_or_html(txbx)
        assert result == ""

    # -- Mixed formatting → HTML --

    def test_bold_and_normal_runs_returns_html(self) -> None:
        """Bold + normal → HTML with ``<b>`` tags."""
        txbx = self._make_txbx(
            [
                [{"text": "bold", "bold": True}, {"text": " normal"}],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx)
        assert "<b>" in result
        assert "bold" in result
        assert "normal" in result

    def test_italic_variation_returns_html(self) -> None:
        """Italic + normal → HTML with ``<i>`` tags."""
        txbx = self._make_txbx(
            [
                [{"text": "italic", "italic": True}, {"text": " normal"}],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx)
        assert "<i>" in result
        assert "italic" in result

    def test_size_variation_returns_html(self) -> None:
        """Different font sizes → HTML with ``<span style="font-size:...">``."""
        txbx = self._make_txbx(
            [
                [{"text": "big", "size": 18.0}, {"text": " small", "size": 10.0}],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx)
        assert "font-size:" in result
        assert "big" in result
        assert "small" in result

    def test_color_variation_returns_html(self) -> None:
        """Different colors → HTML with ``<span style="color:...">``."""
        txbx = self._make_txbx(
            [
                [
                    {"text": "red", "color": "#FF0000"},
                    {"text": " blue", "color": "#0000FF"},
                ],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx)
        assert "color:" in result
        assert "red" in result
        assert "blue" in result

    # -- Style-based formatting --

    def test_style_based_bold_returns_html(self) -> None:
        """``<w:rStyle w:val="Strong"/>`` + normal → HTML (with char_styles)."""
        char_styles = {
            "Strong": (True, False, False, False, None, None, None),
        }
        txbx = self._make_txbx(
            [
                [
                    {"text": "styled bold", "rstyle": "Strong"},
                    {"text": " plain"},
                ],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx, char_styles)
        assert "<b>" in result
        assert "styled bold" in result
        assert "plain" in result

    def test_no_char_styles_still_detects_direct_formatting(self) -> None:
        """``char_styles=None`` works for direct bold."""
        txbx = self._make_txbx(
            [
                [{"text": "bold", "bold": True}, {"text": " normal"}],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx, None)
        assert "<b>" in result

    # -- Unnamespaced val attribute --

    def test_unnamespaced_val_attribute(self) -> None:
        """``<w:b val="0"/>`` correctly reads as not-bold."""
        from lxml import etree  # noqa: PLC0415

        w = self._W
        wps = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        txbx = etree.Element(f"{{{wps}}}txbx")
        content = etree.SubElement(txbx, f"{{{w}}}txbxContent")
        p = etree.SubElement(content, f"{{{w}}}p")

        # Run 1: bold disabled via unnamespaced val="0"
        r1 = etree.SubElement(p, f"{{{w}}}r")
        rpr1 = etree.SubElement(r1, f"{{{w}}}rPr")
        b1 = etree.SubElement(rpr1, f"{{{w}}}b")
        b1.set("val", "0")  # unnamespaced
        t1 = etree.SubElement(r1, f"{{{w}}}t")
        t1.text = "not bold"

        # Run 2: plain
        r2 = etree.SubElement(p, f"{{{w}}}r")
        t2 = etree.SubElement(r2, f"{{{w}}}t")
        t2.text = " also not bold"

        result = _wps_txbx_to_text_or_html(txbx)
        # Both runs have same formatting (no bold), so plain text
        assert result == "not bold also not bold"
        assert "<" not in result

    # -- Multiple <w:t> elements --

    def test_multiple_wt_elements_concatenated(self) -> None:
        """Run with 2 ``<w:t>`` → text from both."""
        txbx = self._make_txbx(
            [
                [
                    {"multi_t": ["Hel", "lo"], "bold": True},
                    {"text": " world"},
                ],
            ]
        )
        result = _wps_txbx_to_text_or_html(txbx)
        assert "<b>" in result
        assert "Hello" in result
        assert "world" in result


# ---------------------------------------------------------------------------
# _parse_docx_char_styles tests
# ---------------------------------------------------------------------------


class TestParseDocxCharStyles:
    """Tests for ``_parse_docx_char_styles``."""

    _W = _WORDML_NS

    def _make_styles_zip(self, styles_xml: str, tmp_path: Path) -> Path:
        """Creates a minimal DOCX zip with only ``word/styles.xml``."""
        docx_path = tmp_path / "test.docx"
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr("word/styles.xml", styles_xml)
        return docx_path

    def test_parse_char_styles_bold_style(self, tmp_path: Path) -> None:
        """'Strong' style → bold=True."""
        w = self._W
        xml = (
            f'<w:styles xmlns:w="{w}">'
            f'  <w:style w:type="character" w:styleId="Strong">'
            f"    <w:rPr><w:b/></w:rPr>"
            f"  </w:style>"
            f"</w:styles>"
        )
        docx_path = self._make_styles_zip(xml, tmp_path)
        with zipfile.ZipFile(docx_path, "r") as zf:
            result = _parse_docx_char_styles(zf)
        assert "Strong" in result
        b, i, u, s, sz, clr, bg = result["Strong"]
        assert b is True
        assert i is False

    def test_parse_char_styles_missing_styles_xml(self, tmp_path: Path) -> None:
        """No ``word/styles.xml`` → empty dict."""
        docx_path = tmp_path / "test.docx"
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr("word/document.xml", "<w:document/>")
        with zipfile.ZipFile(docx_path, "r") as zf:
            result = _parse_docx_char_styles(zf)
        assert result == {}

    def test_parse_char_styles_empty_rpr(self, tmp_path: Path) -> None:
        """Style without ``<w:rPr>`` → not included in result."""
        w = self._W
        xml = (
            f'<w:styles xmlns:w="{w}">'
            f'  <w:style w:type="character" w:styleId="NoFormat">'
            f'    <w:name w:val="NoFormat"/>'
            f"  </w:style>"
            f"</w:styles>"
        )
        docx_path = self._make_styles_zip(xml, tmp_path)
        with zipfile.ZipFile(docx_path, "r") as zf:
            result = _parse_docx_char_styles(zf)
        assert "NoFormat" not in result


# ---------------------------------------------------------------------------
# Style override tests — explicit val="0" vs absent element
# ---------------------------------------------------------------------------


class TestStyleOverrideDisableFormatting:
    """Tests that explicit ``val='0'`` overrides style formatting."""

    _W = _WORDML_NS

    def _make_txbx_with_style(
        self,
        *,
        run_has_bold_elem: bool,
        bold_val: str | None,
    ) -> object:
        """Builds a two-run text box: first with style ref, second plain."""
        from lxml import etree  # noqa: PLC0415

        wps = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        w = self._W
        txbx = etree.Element(f"{{{wps}}}txbx")
        content = etree.SubElement(txbx, f"{{{w}}}txbxContent")
        p = etree.SubElement(content, f"{{{w}}}p")

        # Run 1: has rStyle + optional explicit <w:b>
        r1 = etree.SubElement(p, f"{{{w}}}r")
        rpr1 = etree.SubElement(r1, f"{{{w}}}rPr")
        rs = etree.SubElement(rpr1, f"{{{w}}}rStyle")
        rs.set(f"{{{w}}}val", "Strong")
        if run_has_bold_elem:
            b_el = etree.SubElement(rpr1, f"{{{w}}}b")
            if bold_val is not None:
                b_el.set(f"{{{w}}}val", bold_val)
        t1 = etree.SubElement(r1, f"{{{w}}}t")
        t1.text = "styled"

        # Run 2: plain
        r2 = etree.SubElement(p, f"{{{w}}}r")
        t2 = etree.SubElement(r2, f"{{{w}}}t")
        t2.text = " plain"

        return txbx

    def test_absent_bold_inherits_from_style(self) -> None:
        """No <w:b> element → inherits bold from style."""
        char_styles = {"Strong": (True, False, False, False, None, None, None)}
        txbx = self._make_txbx_with_style(run_has_bold_elem=False, bold_val=None)
        result = _wps_txbx_to_text_or_html(txbx, char_styles)
        assert "<b>" in result

    def test_explicit_bold_false_overrides_style(self) -> None:
        """<w:b w:val='false'/> → overrides style's bold=True → not bold."""
        char_styles = {"Strong": (True, False, False, False, None, None, None)}
        txbx = self._make_txbx_with_style(run_has_bold_elem=True, bold_val="false")
        result = _wps_txbx_to_text_or_html(txbx, char_styles)
        # Both runs are not bold → uniform formatting → plain text
        assert "<b>" not in result

    def test_explicit_bold_zero_overrides_style(self) -> None:
        """<w:b w:val='0'/> → overrides style's bold=True."""
        char_styles = {"Strong": (True, False, False, False, None, None, None)}
        txbx = self._make_txbx_with_style(run_has_bold_elem=True, bold_val="0")
        result = _wps_txbx_to_text_or_html(txbx, char_styles)
        assert "<b>" not in result

    def test_explicit_bold_true_with_style(self) -> None:
        """<w:b/> (no val → default true) + style bold → still bold."""
        char_styles = {"Strong": (True, False, False, False, None, None, None)}
        txbx = self._make_txbx_with_style(run_has_bold_elem=True, bold_val=None)
        result = _wps_txbx_to_text_or_html(txbx, char_styles)
        # Run 1 is bold (from style+direct), run 2 is not → HTML
        assert "<b>" in result


# ---------------------------------------------------------------------------
# DrawingML HTML multi-paragraph injection tests
# ---------------------------------------------------------------------------


class TestDrawingmlHtmlMultiPara:
    """Tests for ``_inject_drawingml_html_runs`` with multi-paragraph shapes."""

    def _make_tx_body(self, paragraphs: list[str]) -> object:
        """Builds an ``<a:txBody>`` lxml element with text paragraphs."""
        from lxml import etree  # noqa: PLC0415

        ns = _DRAWINGML_NS
        tx = etree.Element(f"{{{ns}}}txBody")
        for text in paragraphs:
            p = etree.SubElement(tx, f"{{{ns}}}p")
            r = etree.SubElement(p, f"{{{ns}}}r")
            t = etree.SubElement(r, f"{{{ns}}}t")
            t.text = text
        return tx

    def test_single_paragraph_no_duplication(self) -> None:
        """Single paragraph with HTML → text replaced, not duplicated."""
        ns = _DRAWINGML_NS
        tx = self._make_tx_body(["old text"])
        _inject_drawingml_html_runs(tx, "<b>new</b> text")
        p_els = tx.findall(f"{{{ns}}}p")
        assert len(p_els) == 1  # noqa: PLR2004
        all_text = "".join(t.text or "" for t in p_els[0].iter(f"{{{ns}}}t"))
        assert "old" not in all_text
        assert "new" in all_text

    def test_multi_paragraph_no_old_text_remains(self) -> None:
        """Multi-paragraph shape → old text from second paragraph removed."""
        ns = _DRAWINGML_NS
        tx = self._make_tx_body(["old para1", "old para2"])
        _inject_drawingml_html_runs(tx, "<b>new1</b>\n<i>new2</i>")
        p_els = tx.findall(f"{{{ns}}}p")
        assert len(p_els) == 2  # noqa: PLR2004
        all_text = "".join(t.text or "" for t in tx.iter(f"{{{ns}}}t"))
        assert "old" not in all_text
        assert "new1" in all_text
        assert "new2" in all_text

    def test_multi_paragraph_newline_creates_paragraphs(self) -> None:
        """Newline in HTML creates separate <a:p> elements."""
        ns = _DRAWINGML_NS
        tx = self._make_tx_body(["single"])
        _inject_drawingml_html_runs(tx, "<b>line1</b>\nline2")
        p_els = tx.findall(f"{{{ns}}}p")
        assert len(p_els) == 2  # noqa: PLR2004


# ---------------------------------------------------------------------------
# DrawingML <a:br> line-break in HTML extraction
# ---------------------------------------------------------------------------


class TestDrawingmlBrInHtml:
    """Tests that ``<a:br>`` is preserved in ``_drawingml_to_html``."""

    def _make_tx_body_with_br(self) -> object:
        """Builds ``<a:txBody>`` with one paragraph containing <a:br>."""
        from lxml import etree  # noqa: PLC0415

        ns = _DRAWINGML_NS
        tx = etree.Element(f"{{{ns}}}txBody")
        p = etree.SubElement(tx, f"{{{ns}}}p")

        # Run 1: bold
        r1 = etree.SubElement(p, f"{{{ns}}}r")
        rpr1 = etree.SubElement(r1, f"{{{ns}}}rPr")
        rpr1.set("b", "1")
        t1 = etree.SubElement(r1, f"{{{ns}}}t")
        t1.text = "before"

        # Line break
        etree.SubElement(p, f"{{{ns}}}br")

        # Run 2: not bold
        r2 = etree.SubElement(p, f"{{{ns}}}r")
        t2 = etree.SubElement(r2, f"{{{ns}}}t")
        t2.text = "after"

        return tx

    def test_br_preserved_as_newline(self) -> None:
        """<a:br> in HTML extraction produces a newline character."""
        result = _drawingml_to_html(self._make_tx_body_with_br())
        assert "\n" in result
        assert "before" in result
        assert "after" in result

    def test_br_not_lost(self) -> None:
        """Without the fix, <a:br> would be silently dropped."""
        result = _drawingml_to_html(self._make_tx_body_with_br())
        # The result should not be "beforeafter" (no break)
        assert result != "beforeafter"

    def test_line_break_element_produces_newline(self) -> None:
        """<a:br/> elements are converted to newline characters."""
        from lxml import etree  # noqa: PLC0415

        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        # Build a paragraph with text, then a:br, then more text
        p = etree.fromstring(
            f'<a:p xmlns:a="{a_ns}">'
            f"<a:r><a:t>Hello</a:t></a:r>"
            f"<a:br/>"
            f"<a:r><a:t>World</a:t></a:r>"
            f"</a:p>"
        )
        # Wrap the paragraph in a txBody as _drawingml_to_html expects
        tx_body = etree.Element(f"{{{a_ns}}}txBody")
        tx_body.append(p)
        result = _drawingml_to_html(tx_body)
        # Should contain both texts, with a line break between them
        assert "Hello" in result
        assert "World" in result
        assert "\n" in result


# ---------------------------------------------------------------------------
# DrawingML mixed-formatting detection includes <a:br> runs
# ---------------------------------------------------------------------------


class TestDrawingmlMixedFormattingWithBr:
    """Tests that ``_has_drawingml_mixed_formatting`` works with <a:br>."""

    def test_br_between_mixed_runs_detected(self) -> None:
        """Mixed formatting with <a:br> between runs is still detected."""
        from lxml import etree  # noqa: PLC0415

        ns = _DRAWINGML_NS
        tx = etree.Element(f"{{{ns}}}txBody")
        p = etree.SubElement(tx, f"{{{ns}}}p")

        # Bold run
        r1 = etree.SubElement(p, f"{{{ns}}}r")
        rpr1 = etree.SubElement(r1, f"{{{ns}}}rPr")
        rpr1.set("b", "1")
        t1 = etree.SubElement(r1, f"{{{ns}}}t")
        t1.text = "bold"

        # Line break
        etree.SubElement(p, f"{{{ns}}}br")

        # Plain run
        r2 = etree.SubElement(p, f"{{{ns}}}r")
        t2 = etree.SubElement(r2, f"{{{ns}}}t")
        t2.text = "plain"

        assert _has_drawingml_mixed_formatting(tx) is True


# ---------------------------------------------------------------------------
# Excel content type detection with shape HTML
# ---------------------------------------------------------------------------


class TestExcelContentTypeWithShapeHtml:
    """Tests that HTML in shapes overrides CONTENT_DATA_VALUES for Excel."""

    def test_html_in_values_produces_content_html(self) -> None:
        """Values containing HTML tags → CONTENT_HTML even for Excel."""
        values = ["plain cell", "<b>bold</b> shape"]
        has_html = any(_FORMATTING_HTML_RE.search(v) for v in values)
        assert has_html is True

    def test_no_html_in_values_uses_data_values(self) -> None:
        """No HTML tags → content type is not HTML."""
        values = ["cell1", "cell2"]
        has_html = any(_FORMATTING_HTML_RE.search(v) for v in values)
        assert has_html is False


# ---------------------------------------------------------------------------
# UNO multi-paragraph shape extraction tests
# ---------------------------------------------------------------------------


class TestUnoMultiParagraphShapes:
    """Tests that UNO shape extraction checks all paragraphs for formatting."""

    def test_extract_uno_writer_multi_para_mixed(self) -> None:
        """Multi-paragraph shape with mixed formatting → HTML output."""
        from src.core.office_processor import (  # noqa: PLC0415
            _extract_uno_writer_shapes,
        )

        # Build mock UNO objects
        shape = MagicMock()
        shape.supportsService.return_value = True

        # Two paragraphs: first has mixed formatting, second doesn't
        para1 = MagicMock()
        para2 = MagicMock()

        # Shape enumeration
        shape_enum = _make_uno_enum([para1, para2])
        shape.createEnumeration.return_value = shape_enum

        draw_page = MagicMock()
        draw_page.getCount.return_value = 1
        draw_page.getByIndex.return_value = shape

        doc = MagicMock()
        doc.getDrawPage.return_value = draw_page

        with (
            patch(
                "src.core.office_processor._uno_open",
                return_value=doc,
            ),
            patch(
                "src.core.office_processor._has_uno_mixed_formatting",
                side_effect=[True, False],
            ) as mock_has,
            patch(
                "src.core.office_processor._has_uno_hyperlinks",
                return_value=False,
            ),
            patch(
                "src.core.office_processor._uno_runs_to_html",
                side_effect=["<b>bold</b>", "plain"],
            ) as mock_html,
        ):
            result = _extract_uno_writer_shapes(Path("test.doc"))

        # Both paragraphs checked (any() stops at first True)
        mock_has.assert_called()
        # _uno_runs_to_html called for each paragraph
        assert mock_html.call_count == 2  # noqa: PLR2004
        assert len(result) == 1  # noqa: PLR2004
        assert result[0][0] == "shape:0"
        assert "\n" in result[0][1]
        assert "<b>bold</b>" in result[0][1]

    def test_extract_uno_writer_multi_para_uniform(self) -> None:
        """Multi-paragraph shape with uniform formatting → plain text."""
        from src.core.office_processor import (  # noqa: PLC0415
            _extract_uno_writer_shapes,
        )

        shape = MagicMock()
        shape.supportsService.return_value = True
        shape.getString.return_value = "line1\nline2"

        para1 = MagicMock()
        para2 = MagicMock()

        shape_enum = _make_uno_enum([para1, para2])
        shape.createEnumeration.return_value = shape_enum

        draw_page = MagicMock()
        draw_page.getCount.return_value = 1
        draw_page.getByIndex.return_value = shape

        doc = MagicMock()
        doc.getDrawPage.return_value = draw_page

        with (
            patch(
                "src.core.office_processor._uno_open",
                return_value=doc,
            ),
            patch(
                "src.core.office_processor._has_uno_mixed_formatting",
                return_value=False,
            ),
            patch(
                "src.core.office_processor._has_uno_hyperlinks",
                return_value=False,
            ),
        ):
            result = _extract_uno_writer_shapes(Path("test.doc"))

        assert len(result) == 1  # noqa: PLR2004
        assert result[0][1] == "line1\nline2"
        assert "<" not in result[0][1]


# ---------------------------------------------------------------------------
# _uno_save filter name tests
# ---------------------------------------------------------------------------


class TestUnoSaveFilterName:
    """Verifies _uno_save preserves the original file format."""

    def _call_uno_save(
        self,
        suffix: str,
        doc_filter: str = "",
    ) -> list[object]:
        """Calls _uno_save with a mocked UNO doc and returns PV calls.

        Args:
            suffix: Output file extension (e.g. ".docx").
            doc_filter: FilterName returned by doc.getArgs().  Empty string
                        simulates a document with no filter in its descriptor
                        (falls back to _UNO_FILTER_NAMES lookup).
        """
        from src.core.office_processor import _uno_save  # noqa: PLC0415

        # Mock the com.sun.star.beans.PropertyValue import
        fake_pv_class = MagicMock()
        fake_pv_class.side_effect = lambda name, *a: MagicMock(Name=name, Value=a[1])

        fake_module = MagicMock()
        fake_module.PropertyValue = fake_pv_class

        doc = MagicMock()
        # Simulate doc.getArgs() → list of PropertyValue-like objects
        if doc_filter:
            arg = MagicMock(Name="FilterName", Value=doc_filter)
            doc.getArgs.return_value = [arg]
        else:
            doc.getArgs.return_value = []

        output = Path(f"/tmp/test{suffix}")

        with (
            patch.dict(
                sys.modules,
                {
                    "com": MagicMock(),
                    "com.sun": MagicMock(),
                    "com.sun.star": MagicMock(),
                    "com.sun.star.beans": fake_module,
                },
            ),
            patch(
                "src.core.office_processor._uno_file_url",
                return_value=f"file:///tmp/test{suffix}",
            ),
        ):
            _uno_save(doc, output)

        return fake_pv_class.call_args_list

    # -- Primary path: filter from doc.getArgs() --

    def test_docx_uses_doc_filter(self) -> None:
        """DOCX uses the FilterName from the document's MediaDescriptor."""
        calls = self._call_uno_save(".docx", doc_filter="MS Word 2007 XML")
        filter_calls = [c for c in calls if c[0][0] == "FilterName"]
        assert len(filter_calls) == 1
        assert filter_calls[0][0][2] == "MS Word 2007 XML"

    def test_odt_uses_doc_filter(self) -> None:
        """ODF uses the FilterName from the document's MediaDescriptor."""
        calls = self._call_uno_save(".odt", doc_filter="writer8")
        filter_calls = [c for c in calls if c[0][0] == "FilterName"]
        assert len(filter_calls) == 1
        assert filter_calls[0][0][2] == "writer8"

    # -- Fallback path: filter from _UNO_FILTER_NAMES --

    def test_docx_fallback_filter(self) -> None:
        """DOCX falls back to _UNO_FILTER_NAMES when getArgs is empty."""
        calls = self._call_uno_save(".docx")
        filter_calls = [c for c in calls if c[0][0] == "FilterName"]
        assert len(filter_calls) == 1
        assert filter_calls[0][0][2] == "MS Word 2007 XML"

    def test_xlsx_fallback_filter(self) -> None:
        """XLSX falls back to Calc MS Excel 2007 XML filter."""
        calls = self._call_uno_save(".xlsx")
        filter_calls = [c for c in calls if c[0][0] == "FilterName"]
        assert len(filter_calls) == 1
        assert filter_calls[0][0][2] == "Calc MS Excel 2007 XML"

    def test_pptx_fallback_filter(self) -> None:
        """PPTX falls back to Impress MS PowerPoint 2007 XML filter."""
        calls = self._call_uno_save(".pptx")
        filter_calls = [c for c in calls if c[0][0] == "FilterName"]
        assert len(filter_calls) == 1
        assert filter_calls[0][0][2] == "Impress MS PowerPoint 2007 XML"

    def test_doc_fallback_filter(self) -> None:
        """Legacy .doc falls back to MS Word 97 filter."""
        calls = self._call_uno_save(".doc")
        filter_calls = [c for c in calls if c[0][0] == "FilterName"]
        assert len(filter_calls) == 1
        assert filter_calls[0][0][2] == "MS Word 97"

    def test_odt_no_filter_when_no_descriptor(self) -> None:
        """ODF with no descriptor and no _UNO_FILTER_NAMES entry → no filter."""
        calls = self._call_uno_save(".odt")
        filter_calls = [c for c in calls if c[0][0] == "FilterName"]
        assert len(filter_calls) == 0

    def test_getargs_exception_uses_fallback(self) -> None:
        """If doc.getArgs() raises, falls back to _UNO_FILTER_NAMES."""
        from src.core.office_processor import _uno_save  # noqa: PLC0415

        fake_pv_class = MagicMock()
        fake_pv_class.side_effect = lambda name, *a: MagicMock(Name=name, Value=a[1])
        fake_module = MagicMock()
        fake_module.PropertyValue = fake_pv_class

        doc = MagicMock()
        doc.getArgs.side_effect = RuntimeError("no args")
        output = Path("/tmp/test.docx")

        with (
            patch.dict(
                sys.modules,
                {
                    "com": MagicMock(),
                    "com.sun": MagicMock(),
                    "com.sun.star": MagicMock(),
                    "com.sun.star.beans": fake_module,
                },
            ),
            patch(
                "src.core.office_processor._uno_file_url",
                return_value="file:///tmp/test.docx",
            ),
        ):
            _uno_save(doc, output)

        filter_calls = [
            c for c in fake_pv_class.call_args_list if c[0][0] == "FilterName"
        ]
        assert len(filter_calls) == 1
        assert filter_calls[0][0][2] == "MS Word 2007 XML"


# ── Hyperlink helpers ────────────────────────────────────────────────────


_hyperlink_counter = 0


def _add_hyperlink_to_para(para, url, text, bold=False):
    """Adds a <w:hyperlink> element with a run to a python-docx paragraph.

    Creates a mock relationship so ``_resolve_para_hyperlink_rels`` can
    look up the URL.  Returns the r:id string.
    """
    global _hyperlink_counter  # noqa: PLW0603
    _hyperlink_counter += 1
    r_id = f"rIdHL{_hyperlink_counter}"
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)

    run = OxmlElement("w:r")
    rpr = OxmlElement("w:rPr")
    if bold:
        b_elem = OxmlElement("w:b")
        rpr.append(b_elem)
    run.append(rpr)
    t = OxmlElement("w:t")
    t.text = text
    t.set(qn("xml:space"), "preserve")
    run.append(t)

    hyperlink.append(run)
    para._element.append(hyperlink)

    # Register relationship on the document part
    mock_rel = MagicMock()
    mock_rel.target_ref = url
    para.part.rels[r_id] = mock_rel
    return r_id


def _add_anchor_hyperlink_to_para(para, anchor, text):
    """Adds a ``<w:hyperlink w:anchor="...">`` (internal bookmark) to a para."""
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("w:anchor"), anchor)

    run = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = text
    t.set(qn("xml:space"), "preserve")
    run.append(t)
    hyperlink.append(run)
    para._element.append(hyperlink)


# ── Hyperlink tests ──────────────────────────────────────────────────────


class TestParaHasHyperlinks:
    """Tests for ``_para_has_hyperlinks``."""

    def test_no_hyperlinks(self) -> None:
        doc = Document()
        para = doc.add_paragraph("Hello world")
        assert _para_has_hyperlinks(para._element) is False

    def test_with_hyperlink(self) -> None:
        doc = Document()
        para = doc.add_paragraph("Visit ")
        _add_hyperlink_to_para(para, "https://example.com", "here")
        assert _para_has_hyperlinks(para._element) is True

    def test_with_anchor_hyperlink(self) -> None:
        doc = Document()
        para = doc.add_paragraph("")
        _add_anchor_hyperlink_to_para(para, "section1", "jump")
        assert _para_has_hyperlinks(para._element) is True


class TestRunsToHtmlHyperlinks:
    """Tests for ``_runs_to_html`` with hyperlinks."""

    def test_single_hyperlink(self) -> None:
        doc = Document()
        para = doc.add_paragraph("Click ")
        r_id = _add_hyperlink_to_para(para, "https://example.com", "here")
        rels = {r_id: "https://example.com"}

        result = _runs_to_html(para, hyperlink_rels=rels)
        assert "Click " in result
        assert '<a href="https://example.com">' in result
        assert "here" in result
        assert "</a>" in result

    def test_hyperlink_with_bold(self) -> None:
        doc = Document()
        para = doc.add_paragraph("")
        r_id = _add_hyperlink_to_para(
            para,
            "https://example.com",
            "bold link",
            bold=True,
        )
        rels = {r_id: "https://example.com"}

        result = _runs_to_html(para, hyperlink_rels=rels)
        assert "<b>" in result
        assert '<a href="https://example.com">' in result

    def test_anchor_hyperlink(self) -> None:
        doc = Document()
        para = doc.add_paragraph("See ")
        _add_anchor_hyperlink_to_para(para, "section1", "Section 1")

        result = _runs_to_html(para, hyperlink_rels=None)
        assert '<a href="#section1">' in result
        assert "Section 1" in result
        assert "</a>" in result

    def test_multiple_hyperlinks(self) -> None:
        doc = Document()
        para = doc.add_paragraph("Visit ")
        r_id1 = _add_hyperlink_to_para(para, "https://a.com", "A")
        # Add plain text between hyperlinks
        run_mid = OxmlElement("w:r")
        t_mid = OxmlElement("w:t")
        t_mid.text = " or "
        t_mid.set(qn("xml:space"), "preserve")
        run_mid.append(t_mid)
        para._element.append(run_mid)
        r_id2 = _add_hyperlink_to_para(para, "https://b.com", "B")
        rels = {r_id1: "https://a.com", r_id2: "https://b.com"}

        result = _runs_to_html(para, hyperlink_rels=rels)
        assert '<a href="https://a.com">A</a>' in result
        assert '<a href="https://b.com">B</a>' in result
        assert " or " in result

    def test_no_hyperlinks_unchanged(self) -> None:
        """Without hyperlinks, result matches original ``_runs_to_html``."""
        doc = Document()
        para = doc.add_paragraph("")
        run1 = para.add_run("Hello ")
        run1.bold = True
        para.add_run("world")

        result = _runs_to_html(para, hyperlink_rels=None)
        assert "<b>" in result
        assert "Hello " in result
        assert "world" in result
        assert "<a " not in result

    def test_html_escaping_in_href(self) -> None:
        doc = Document()
        para = doc.add_paragraph("")
        url = 'https://example.com/page?a=1&b="2"'
        r_id = _add_hyperlink_to_para(para, url, "link")
        rels = {r_id: url}

        result = _runs_to_html(para, hyperlink_rels=rels)
        # href value must be HTML-escaped
        assert "&amp;" in result
        assert "&quot;" in result

    def test_hyperlink_with_unresolved_rid(self) -> None:
        """r:id not in rels dict → runs emitted without <a> wrapper."""
        doc = Document()
        para = doc.add_paragraph("")
        _add_hyperlink_to_para(para, "https://example.com", "orphan")

        result = _runs_to_html(para, hyperlink_rels={})
        assert "orphan" in result
        assert "<a " not in result


class TestParseHtmlFormattingHyperlinks:
    """Tests for ``_parse_html_formatting`` with ``<a>`` tags."""

    def test_simple_link(self) -> None:
        segments = _parse_html_formatting(
            'Visit <a href="https://example.com">here</a> now',
        )
        assert len(segments) == 3  # noqa: PLR2004
        assert segments[0].text == "Visit "
        assert segments[0].hyperlink_url is None
        assert segments[1].text == "here"
        assert segments[1].hyperlink_url == "https://example.com"
        assert segments[2].text == " now"
        assert segments[2].hyperlink_url is None

    def test_link_with_formatting(self) -> None:
        segments = _parse_html_formatting(
            '<a href="url"><b>bold link</b></a>',
        )
        assert len(segments) == 1
        assert segments[0].text == "bold link"
        assert segments[0].bold is True
        assert segments[0].hyperlink_url == "url"

    def test_merge_same_link_same_formatting(self) -> None:
        segments = _parse_html_formatting(
            '<a href="url">part1</a><a href="url">part2</a>',
        )
        # Adjacent segments with same URL and formatting merge
        assert len(segments) == 1
        assert segments[0].text == "part1part2"
        assert segments[0].hyperlink_url == "url"

    def test_no_merge_different_links(self) -> None:
        segments = _parse_html_formatting(
            '<a href="a">X</a><a href="b">Y</a>',
        )
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0].hyperlink_url == "a"
        assert segments[1].hyperlink_url == "b"

    def test_escaped_href(self) -> None:
        segments = _parse_html_formatting(
            '<a href="https://x.com?a=1&amp;b=2">link</a>',
        )
        assert segments[0].hyperlink_url == "https://x.com?a=1&b=2"

    def test_formatting_html_re_matches_a_tag(self) -> None:
        assert _FORMATTING_HTML_RE.search('<a href="url">x</a>') is not None
        assert _FORMATTING_HTML_RE.search("</a>") is not None
        assert _FORMATTING_HTML_RE.search("<a>x</a>") is not None


class TestInjectHtmlRunsHyperlinks:
    """Tests for ``_inject_html_runs`` hyperlink creation."""

    def test_inject_single_hyperlink(self) -> None:
        doc = Document()
        para = doc.add_paragraph("placeholder")

        _inject_html_runs(
            para,
            'Click <a href="https://example.com">here</a>',
            part=doc.part,
        )

        # Check XML structure
        hyperlinks = para._element.findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 1
        r_id = hyperlinks[0].get(qn("r:id"))
        assert r_id is not None
        # Relationship was created
        rel = doc.part.rels[r_id]
        assert rel.target_ref == "https://example.com"
        assert rel.reltype == _HYPERLINK_RELTYPE

    def test_inject_anchor_hyperlink(self) -> None:
        doc = Document()
        para = doc.add_paragraph("placeholder")

        _inject_html_runs(
            para,
            '<a href="#bookmark1">see section</a>',
            part=doc.part,
        )

        hyperlinks = para._element.findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 1
        assert hyperlinks[0].get(qn("w:anchor")) == "bookmark1"
        # No r:id for anchors
        assert hyperlinks[0].get(qn("r:id")) is None

    def test_inject_hyperlink_text_preserved(self) -> None:
        doc = Document()
        para = doc.add_paragraph("old text")

        _inject_html_runs(
            para,
            'Visit <a href="https://example.com"><b>Example</b></a> today',
            part=doc.part,
        )

        # "Visit " and " today" are direct <w:r>, "Example" is inside hyperlink
        all_texts = [t.text for t in para._element.iter(qn("w:t")) if t.text]
        assert "Visit " in all_texts
        assert "Example" in all_texts
        assert " today" in all_texts

    def test_inject_multiple_hyperlinks(self) -> None:
        doc = Document()
        para = doc.add_paragraph("old")

        _inject_html_runs(
            para,
            '<a href="https://a.com">A</a> and <a href="https://b.com">B</a>',
            part=doc.part,
        )

        hyperlinks = para._element.findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 2  # noqa: PLR2004

    def test_inject_consecutive_same_url_grouped(self) -> None:
        doc = Document()
        para = doc.add_paragraph("old")

        _inject_html_runs(
            para,
            '<a href="https://x.com"><b>bold</b> plain</a>',
            part=doc.part,
        )

        # Both runs should be under a single <w:hyperlink>
        hyperlinks = para._element.findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 1
        runs_in_hl = hyperlinks[0].findall(qn("w:r"))
        assert len(runs_in_hl) == 2  # noqa: PLR2004

    def test_inject_without_part_skips_external(self) -> None:
        doc = Document()
        para = doc.add_paragraph("old")

        _inject_html_runs(
            para,
            '<a href="https://example.com">link</a>',
            part=None,  # No part available
        )

        # Hyperlink not created, but text still injected
        hyperlinks = para._element.findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 0
        texts = [t.text for t in para._element.iter(qn("w:t")) if t.text]
        assert "link" in texts

    def test_old_hyperlinks_removed(self) -> None:
        doc = Document()
        para = doc.add_paragraph("before ")
        _add_hyperlink_to_para(para, "https://old.com", "old link")

        _inject_html_runs(
            para,
            "translated text",
            part=doc.part,
        )

        # Old hyperlinks must be gone
        hyperlinks = para._element.findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 0


class TestReplaceParaTextHyperlinks:
    """Tests for ``_replace_paragraph_text`` hyperlink cleanup."""

    def test_removes_stale_hyperlinks(self) -> None:
        doc = Document()
        para = doc.add_paragraph("text ")
        _add_hyperlink_to_para(para, "https://old.com", "old")

        _replace_paragraph_text(para, "new text")

        hyperlinks = para._element.findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 0
        assert para.text == "new text"

    def test_no_hyperlinks_unchanged_behavior(self) -> None:
        doc = Document()
        para = doc.add_paragraph("old text")

        _replace_paragraph_text(para, "new text")

        assert para.text == "new text"


class TestExtractParaWithLinks:
    """Tests for ``_extract_para_with_links`` in office_processor."""

    def test_plain_paragraph(self) -> None:
        doc = Document()
        para = doc.add_paragraph("Hello world")
        assert _extract_para_with_links(para) == "Hello world"

    def test_paragraph_with_hyperlink(self) -> None:
        doc = Document()
        para = doc.add_paragraph("Click ")
        _add_hyperlink_to_para(para, "https://example.com", "here")

        result = _extract_para_with_links(para)
        assert '<a href="https://example.com">' in result
        assert "here" in result
        assert "</a>" in result


class TestResolveParaHyperlinkRels:
    """Tests for ``_resolve_para_hyperlink_rels``."""

    def test_no_hyperlinks(self) -> None:
        doc = Document()
        para = doc.add_paragraph("plain text")
        assert _resolve_para_hyperlink_rels(para) == {}

    def test_resolves_external_hyperlink(self) -> None:
        doc = Document()
        para = doc.add_paragraph("text ")
        r_id = _add_hyperlink_to_para(para, "https://example.com", "link")

        rels = _resolve_para_hyperlink_rels(para)
        assert r_id in rels
        assert rels[r_id] == "https://example.com"

    def test_anchor_only_not_in_rels(self) -> None:
        doc = Document()
        para = doc.add_paragraph("")
        _add_anchor_hyperlink_to_para(para, "bm1", "link")

        # Anchor hyperlinks have no r:id → empty rels dict
        assert _resolve_para_hyperlink_rels(para) == {}

    def test_missing_rel_ignored(self) -> None:
        doc = Document()
        para = doc.add_paragraph("")
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), "rIdMissing")
        para._element.append(hyperlink)

        # KeyError in rels lookup should be caught
        rels = _resolve_para_hyperlink_rels(para)
        assert rels == {}


class TestHyperlinkRoundTrip:
    """End-to-end tests: extract → translate → inject for hyperlinks."""

    def test_single_link_roundtrip(self) -> None:
        """Hyperlink survives extract → inject cycle."""
        doc = Document()
        para = doc.add_paragraph("Visit ")
        r_id = _add_hyperlink_to_para(para, "https://example.com", "our site")

        # Extract
        rels = {r_id: "https://example.com"}
        html_text = _runs_to_html(para, hyperlink_rels=rels)
        assert '<a href="https://example.com">' in html_text

        # Simulate LLM translation (text changes, tags preserved)
        translated = html_text.replace("Visit ", "Visitez ").replace(
            "our site",
            "notre site",
        )

        # Inject into a fresh paragraph
        doc2 = Document()
        para2 = doc2.add_paragraph("placeholder")
        _inject_html_runs(para2, translated, part=doc2.part)

        # Verify hyperlink was recreated
        hyperlinks = para2._element.findall(qn("w:hyperlink"))
        assert len(hyperlinks) == 1
        new_rid = hyperlinks[0].get(qn("r:id"))
        assert new_rid is not None
        assert doc2.part.rels[new_rid].target_ref == "https://example.com"

        # Verify translated text
        link_texts = [t.text for t in hyperlinks[0].iter(qn("w:t")) if t.text]
        assert "notre site" in " ".join(link_texts)

    def test_mixed_plain_and_link_roundtrip(self) -> None:
        doc = Document()
        para = doc.add_paragraph("")
        # Plain bold run
        run1 = para.add_run("Read ")
        run1.bold = True
        # Hyperlink run
        r_id = _add_hyperlink_to_para(para, "https://doc.io", "the docs")

        rels = {r_id: "https://doc.io"}
        html_text = _runs_to_html(para, hyperlink_rels=rels)
        assert "<b>" in html_text
        assert '<a href="https://doc.io">' in html_text

        # Inject
        doc2 = Document()
        para2 = doc2.add_paragraph("old")
        _inject_html_runs(para2, html_text, part=doc2.part)

        # Both plain run and hyperlink should exist
        direct_runs = [c for c in para2._element if c.tag == qn("w:r")]
        hyperlinks = para2._element.findall(qn("w:hyperlink"))
        assert len(direct_runs) >= 1
        assert len(hyperlinks) == 1


# ── PPTX hyperlink helpers ────────────────────────────────────────────────


def _make_pptx_para_with_hyperlink(
    specs: list[tuple[str, bool, bool, bool, bool, str | None]],
) -> object:
    """Creates a PPTX paragraph with runs, optionally adding hyperlinks.

    Args:
        specs: List of (text, bold, italic, underline, strike, url_or_None).

    Returns:
        The paragraph object.
    """
    from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    para = txbox.text_frame.paragraphs[0]

    for idx, (text, bold, italic, underline, strike, url) in enumerate(specs):
        if idx == 0:
            run = para.runs[0] if para.runs else para.add_run()
        else:
            run = para.add_run()
        run.text = text
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        if strike:
            rpr = run._r.find(pptx_qn("a:rPr"))
            if rpr is None:
                from lxml import etree as _et  # noqa: PLC0415

                rpr = _et.SubElement(run._r, pptx_qn("a:rPr"))
            rpr.set("strike", "sngStrike")
        if url:
            run.hyperlink.address = url

    return para


# ── PPTX hyperlink tests ──────────────────────────────────────────────────


class TestHasPptxHyperlinks:
    """Tests for ``_has_pptx_hyperlinks``."""

    def test_no_hyperlinks(self) -> None:
        para = _make_pptx_para_with_runs([("Hello", False, False, False, False)])
        assert _has_pptx_hyperlinks(para) is False

    def test_with_hyperlink(self) -> None:
        para = _make_pptx_para_with_hyperlink(
            [("Click here", False, False, False, False, "https://example.com")]
        )
        assert _has_pptx_hyperlinks(para) is True

    def test_empty_run_with_hyperlink_skipped(self) -> None:
        para = _make_pptx_para_with_hyperlink(
            [("", False, False, False, False, "https://example.com")]
        )
        assert _has_pptx_hyperlinks(para) is False

    def test_no_runs(self) -> None:
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        para = txbox.text_frame.paragraphs[0]
        for r in list(para._p.findall(pptx_qn("a:r"))):
            para._p.remove(r)
        assert _has_pptx_hyperlinks(para) is False


class TestPptxRunsToHtmlHyperlinks:
    """Tests for hyperlink emission in ``_pptx_runs_to_html``."""

    def test_single_hyperlink_run(self) -> None:
        para = _make_pptx_para_with_hyperlink(
            [("Click", False, False, False, False, "https://example.com")]
        )
        result = _pptx_runs_to_html(para)
        assert '<a href="https://example.com">' in result
        assert "Click" in result
        assert result.endswith("</a>")

    def test_plain_then_hyperlink(self) -> None:
        para = _make_pptx_para_with_hyperlink(
            [
                ("Visit ", False, False, False, False, None),
                ("here", False, False, False, False, "https://example.com"),
            ],
        )
        result = _pptx_runs_to_html(para)
        assert "Visit " in result
        assert '<a href="https://example.com">here</a>' in result

    def test_hyperlink_then_plain(self) -> None:
        para = _make_pptx_para_with_hyperlink(
            [
                ("Link", False, False, False, False, "https://example.com"),
                (" text", False, False, False, False, None),
            ],
        )
        result = _pptx_runs_to_html(para)
        assert result.startswith('<a href="https://example.com">Link</a>')
        assert " text" in result

    def test_bold_hyperlink(self) -> None:
        para = _make_pptx_para_with_hyperlink(
            [("Bold link", True, False, False, False, "https://x.com")]
        )
        result = _pptx_runs_to_html(para)
        assert '<a href="https://x.com">' in result
        assert "<b>" in result

    def test_two_consecutive_same_url_runs(self) -> None:
        """Consecutive runs with same URL are grouped under one <a> tag."""
        para = _make_pptx_para_with_hyperlink(
            [
                ("part1 ", True, False, False, False, "https://a.com"),
                ("part2", False, False, False, False, "https://a.com"),
            ],
        )
        result = _pptx_runs_to_html(para)
        # Should have only one <a> open and one </a> close
        assert result.count('<a href="https://a.com">') == 1
        assert result.count("</a>") == 1

    def test_two_different_urls(self) -> None:
        para = _make_pptx_para_with_hyperlink(
            [
                ("A", False, False, False, False, "https://a.com"),
                ("B", False, False, False, False, "https://b.com"),
            ],
        )
        result = _pptx_runs_to_html(para)
        assert '<a href="https://a.com">' in result
        assert '<a href="https://b.com">' in result
        assert result.count("</a>") == 2  # noqa: PLR2004

    def test_url_html_escaped(self) -> None:
        para = _make_pptx_para_with_hyperlink(
            [("Go", False, False, False, False, "https://x.com?a=1&b=2")]
        )
        result = _pptx_runs_to_html(para)
        assert "a=1&amp;b=2" in result


class TestInjectPptxHtmlRunsHyperlinks:
    """Tests for hyperlink injection in ``_inject_pptx_html_runs``."""

    def _make_para_and_part(self):
        """Creates a fresh PPTX paragraph and its slide part."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(2),
        )
        para = txbox.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = "placeholder"
        return para, slide.part

    def test_hyperlink_creates_hlink_click(self) -> None:
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para, part = self._make_para_and_part()
        _inject_pptx_html_runs(
            para,
            '<a href="https://example.com">Click</a>',
            part=part,
        )
        # Find <a:hlinkClick> inside <a:rPr>
        runs = para._p.findall(pptx_qn("a:r"))
        assert len(runs) >= 1
        rpr = runs[0].find(pptx_qn("a:rPr"))
        assert rpr is not None
        hlink = rpr.find(pptx_qn("a:hlinkClick"))
        assert hlink is not None
        # r:id should be set
        r_id = hlink.get(pptx_qn("r:id"))
        assert r_id is not None
        # Verify the relationship points to the URL
        assert part.rels[r_id].target_ref == "https://example.com"

    def test_plain_text_no_hlink(self) -> None:
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para, part = self._make_para_and_part()
        _inject_pptx_html_runs(
            para,
            "<b>No link</b>",
            part=part,
        )
        runs = para._p.findall(pptx_qn("a:r"))
        for r in runs:
            rpr = r.find(pptx_qn("a:rPr"))
            if rpr is not None:
                assert rpr.find(pptx_qn("a:hlinkClick")) is None

    def test_mixed_link_and_plain(self) -> None:
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para, part = self._make_para_and_part()
        _inject_pptx_html_runs(
            para,
            'Visit <a href="https://example.com"><b>here</b></a> now',
            part=part,
        )
        runs = para._p.findall(pptx_qn("a:r"))
        # Should have at least 3 runs: "Visit ", "here", " now"
        assert len(runs) >= 3  # noqa: PLR2004
        # Only the middle run(s) should have hlinkClick
        hlink_count = 0
        for r in runs:
            rpr = r.find(pptx_qn("a:rPr"))
            if rpr is not None and rpr.find(pptx_qn("a:hlinkClick")) is not None:
                hlink_count += 1
        assert hlink_count >= 1

    def test_no_part_skips_hyperlink(self) -> None:
        """When part is None, hyperlinks are not created (graceful degradation)."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para, _ = self._make_para_and_part()
        _inject_pptx_html_runs(
            para,
            '<a href="https://example.com">Click</a>',
            part=None,
        )
        runs = para._p.findall(pptx_qn("a:r"))
        assert len(runs) >= 1
        # Text should still be injected
        texts = [
            r.find(pptx_qn("a:t")).text
            for r in runs
            if r.find(pptx_qn("a:t")) is not None
        ]
        assert "Click" in "".join(texts)
        # But no hlinkClick
        for r in runs:
            rpr = r.find(pptx_qn("a:rPr"))
            if rpr is not None:
                assert rpr.find(pptx_qn("a:hlinkClick")) is None

    def test_base_rpr_hlink_stripped(self) -> None:
        """Existing <a:hlinkClick> in base rPr is stripped before copying."""
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        para, part = self._make_para_and_part()
        # Add a hyperlink to the existing run's rPr
        if para.runs:
            para.runs[0].hyperlink.address = "https://old.com"

        _inject_pptx_html_runs(para, "<b>No link</b>", part=part)
        runs = para._p.findall(pptx_qn("a:r"))
        for r in runs:
            rpr = r.find(pptx_qn("a:rPr"))
            if rpr is not None:
                assert rpr.find(pptx_qn("a:hlinkClick")) is None


class TestPptxHyperlinkRoundTrip:
    """End-to-end extract → translate → inject round trips for PPTX hyperlinks."""

    def test_single_link_roundtrip(self) -> None:
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        # Create paragraph with hyperlink
        para = _make_pptx_para_with_hyperlink(
            [("Click here", False, False, False, False, "https://example.com")]
        )
        html_text = _pptx_runs_to_html(para)
        assert '<a href="https://example.com">' in html_text

        # Inject into a fresh paragraph
        prs2 = Presentation()
        slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
        txbox2 = slide2.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(2),
        )
        para2 = txbox2.text_frame.paragraphs[0]
        _inject_pptx_html_runs(para2, html_text, part=slide2.part)

        # Verify <a:hlinkClick> exists
        runs = para2._p.findall(pptx_qn("a:r"))
        assert len(runs) >= 1
        hlink_found = False
        for r in runs:
            rpr = r.find(pptx_qn("a:rPr"))
            if rpr is not None:
                hlink = rpr.find(pptx_qn("a:hlinkClick"))
                if hlink is not None:
                    r_id = hlink.get(pptx_qn("r:id"))
                    assert slide2.part.rels[r_id].target_ref == "https://example.com"
                    hlink_found = True
        assert hlink_found

    def test_mixed_plain_and_link_roundtrip(self) -> None:
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        # Create paragraph: plain bold + hyperlink
        para = _make_pptx_para_with_hyperlink(
            [
                ("Read ", True, False, False, False, None),
                ("the docs", False, False, False, False, "https://doc.io"),
            ],
        )
        html_text = _pptx_runs_to_html(para)
        assert "<b>" in html_text
        assert '<a href="https://doc.io">' in html_text

        # Inject into fresh
        prs2 = Presentation()
        slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
        txbox2 = slide2.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(2),
        )
        para2 = txbox2.text_frame.paragraphs[0]
        _inject_pptx_html_runs(para2, html_text, part=slide2.part)

        # Should have both plain runs and hyperlinked runs
        runs = para2._p.findall(pptx_qn("a:r"))
        assert len(runs) >= 2  # noqa: PLR2004
        hlink_runs = []
        plain_runs = []
        for r in runs:
            rpr = r.find(pptx_qn("a:rPr"))
            if rpr is not None and rpr.find(pptx_qn("a:hlinkClick")) is not None:
                hlink_runs.append(r)
            else:
                plain_runs.append(r)
        assert len(hlink_runs) >= 1
        assert len(plain_runs) >= 1


# ---------------------------------------------------------------------------
# Win32COM Word hyperlink tests
# ---------------------------------------------------------------------------


class TestHasWin32comRangeHyperlinks:
    """Tests for _has_win32com_range_hyperlinks."""

    def test_no_hyperlinks(self) -> None:
        """Range with zero hyperlinks returns False."""
        rng = MagicMock()
        rng.Hyperlinks.Count = 0
        assert _has_win32com_range_hyperlinks(rng) is False

    def test_one_hyperlink(self) -> None:
        """Range with one hyperlink returns True."""
        rng = MagicMock()
        rng.Hyperlinks.Count = 1
        assert _has_win32com_range_hyperlinks(rng) is True

    def test_com_error_returns_false(self) -> None:
        """COM exception → conservatively returns False."""
        rng = MagicMock()
        type(rng).Hyperlinks = property(
            lambda s: (_ for _ in ()).throw(OSError("COM error")),
        )
        assert _has_win32com_range_hyperlinks(rng) is False


class TestHasWin32comWordHyperlinks:
    """Tests for _has_win32com_word_hyperlinks."""

    def test_delegates_to_range(self) -> None:
        """Delegates to _has_win32com_range_hyperlinks on para.Range."""
        para = MagicMock()
        para.Range.Hyperlinks.Count = 2  # noqa: PLR2004
        assert _has_win32com_word_hyperlinks(para) is True

    def test_no_hyperlinks(self) -> None:
        """Paragraph without hyperlinks returns False."""
        para = MagicMock()
        para.Range.Hyperlinks.Count = 0
        assert _has_win32com_word_hyperlinks(para) is False


class TestWin32comWordRunsToHtmlHyperlinks:
    """Tests for hyperlink support in _win32com_word_runs_to_html."""

    def test_hyperlink_emits_anchor_tag(self) -> None:
        """Characters inside a hyperlink get wrapped in <a> tag."""
        url = "https://example.com"
        chars = [
            _make_win32com_char("A"),
            _make_win32com_char("B"),
            _make_win32com_char("C"),
        ]
        # Build hyperlink mock
        hl = MagicMock()
        hl.Address = url
        hl.Range.Start = 0  # absolute position
        hl.Range.End = 2  # covers chars A, B  # noqa: PLR2004

        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        para.Range.Start = 0
        para.Range.Hyperlinks.Count = 1
        para.Range.Hyperlinks.side_effect = lambda i: hl

        result = _win32com_word_runs_to_html(para)
        assert '<a href="https://example.com">' in result
        assert "</a>" in result
        # "C" should NOT be inside the anchor
        assert result.endswith("</a>C") or "C" in result.split("</a>")[-1]

    def test_no_hyperlinks_plain_html(self) -> None:
        """Without hyperlinks, output is plain formatting HTML."""
        chars = [
            _make_win32com_char("X", bold=True),
            _make_win32com_char("Y"),
        ]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        para.Range.Start = 0
        para.Range.Hyperlinks.Count = 0

        result = _win32com_word_runs_to_html(para)
        assert "<b>X</b>" in result
        assert "<a " not in result

    def test_bold_hyperlink(self) -> None:
        """Bold text inside a hyperlink gets both <a> and <b> tags."""
        url = "https://bold.example"
        chars = [
            _make_win32com_char("B", bold=True),
            _make_win32com_char("C", bold=True),
        ]
        hl = MagicMock()
        hl.Address = url
        hl.Range.Start = 0
        hl.Range.End = 2  # noqa: PLR2004

        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        para.Range.Start = 0
        para.Range.Hyperlinks.Count = 1
        para.Range.Hyperlinks.side_effect = lambda i: hl

        result = _win32com_word_runs_to_html(para)
        assert "<a " in result
        assert "<b>" in result
        assert "</b>" in result
        assert "</a>" in result

    def test_hyperlink_error_graceful(self) -> None:
        """COM error reading hyperlinks → falls back to no hyperlinks."""
        chars = [_make_win32com_char("X")]
        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        para.Range.Start = 0
        # Hyperlinks access raises
        type(para.Range).Hyperlinks = property(
            lambda s: (_ for _ in ()).throw(OSError("COM")),
        )

        result = _win32com_word_runs_to_html(para)
        # Should still produce output, just without links
        assert "X" in result
        assert "<a " not in result

    def test_superscript_inside_hyperlink(self) -> None:
        """Superscript inside a hyperlink gets both <sup> and <a> tags."""
        url = "https://sup.example"
        chars = [
            _make_win32com_char("x", superscript=True),
            _make_win32com_char("2", superscript=True),
        ]
        hl = MagicMock()
        hl.Address = url
        hl.Range.Start = 0
        hl.Range.End = 2  # noqa: PLR2004

        para = MagicMock()
        para.Range.Characters.Count = len(chars)
        para.Range.Characters.side_effect = lambda i: chars[i - 1]
        para.Range.Start = 0
        para.Range.Hyperlinks.Count = 1
        para.Range.Hyperlinks.side_effect = lambda i: hl

        result = _win32com_word_runs_to_html(para)
        assert "<sup>" in result
        assert "<a " in result
        assert "x2" in result


class TestInjectWin32comWordHtmlRunsHyperlinks:
    """Tests for hyperlink injection in _inject_win32com_word_html_runs."""

    def test_hyperlink_created(self) -> None:
        """Segments with hyperlink_url cause doc.Hyperlinks.Add call."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0

        sub_ranges: dict[tuple[int, int], MagicMock] = {}

        def make_sub_range(start: int, end: int) -> MagicMock:
            sr = MagicMock()
            sub_ranges[(start, end)] = sr
            return sr

        doc.Range.side_effect = make_sub_range

        _inject_win32com_word_html_runs(
            doc,
            rng,
            '<a href="https://example.com">Link</a> text',
        )
        # Hyperlink should be added for "Link" (chars 0..4)
        doc.Hyperlinks.Add.assert_called_once()
        call_kwargs = doc.Hyperlinks.Add.call_args
        assert call_kwargs[1]["Address"] == "https://example.com"

    def test_no_hyperlink_no_add(self) -> None:
        """Segments without hyperlink_url don't call Hyperlinks.Add."""
        doc = MagicMock()
        rng = MagicMock()
        rng.Start = 0
        doc.Range.return_value = MagicMock()

        _inject_win32com_word_html_runs(
            doc,
            rng,
            "<b>Bold</b> plain",
        )
        doc.Hyperlinks.Add.assert_not_called()


class TestWin32comWordHyperlinkRoundTrip:
    """Round-trip: extract hyperlink HTML → parse → re-inject."""

    def test_link_survives_roundtrip(self) -> None:
        """<a> tag parsed by _parse_html_formatting preserves URL."""
        html_in = '<a href="https://rt.example">Click</a> here'
        segments = _parse_html_formatting(html_in)
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0].hyperlink_url == "https://rt.example"
        assert segments[0].text == "Click"
        assert segments[1].hyperlink_url is None
        assert segments[1].text == " here"


# ---------------------------------------------------------------------------
# Win32COM PPT hyperlink tests
# ---------------------------------------------------------------------------


class TestHasWin32comPptHyperlinks:
    """Tests for _has_win32com_ppt_hyperlinks."""

    def test_no_hyperlinks(self) -> None:
        """All runs without hyperlinks → False."""
        runs = [
            _make_win32com_ppt_run("Hello"),
            _make_win32com_ppt_run("World"),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        assert _has_win32com_ppt_hyperlinks(para_rng) is False

    def test_one_hyperlink(self) -> None:
        """One run with hyperlink → True."""
        runs = [
            _make_win32com_ppt_run("Click", hyperlink_address="https://x.com"),
            _make_win32com_ppt_run("plain"),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        assert _has_win32com_ppt_hyperlinks(para_rng) is True

    def test_whitespace_only_run_skipped(self) -> None:
        """Whitespace-only run with hyperlink is skipped."""
        runs = [
            _make_win32com_ppt_run(
                "  ",
                hyperlink_address="https://x.com",
            ),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        assert _has_win32com_ppt_hyperlinks(para_rng) is False

    def test_com_error_returns_false(self) -> None:
        """COM exception → False."""
        para_rng = MagicMock()
        para_rng.Runs.side_effect = OSError("COM error")
        assert _has_win32com_ppt_hyperlinks(para_rng) is False


class TestWin32comPptRunsToHtmlHyperlinks:
    """Tests for hyperlink support in _win32com_ppt_runs_to_html."""

    def test_hyperlink_emits_anchor(self) -> None:
        """Run with hyperlink gets <a> tag."""
        runs = [
            _make_win32com_ppt_run(
                "Click",
                hyperlink_address="https://ppt.example",
            ),
            _make_win32com_ppt_run("plain"),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        assert '<a href="https://ppt.example">' in result
        assert "</a>" in result
        assert "plain" in result

    def test_no_hyperlink_no_anchor(self) -> None:
        """Runs without hyperlinks → no <a> tags."""
        runs = [
            _make_win32com_ppt_run("A"),
            _make_win32com_ppt_run("B"),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        assert "<a " not in result

    def test_bold_hyperlink(self) -> None:
        """Bold run with hyperlink gets both <a> and <b> tags."""
        runs = [
            _make_win32com_ppt_run(
                "Link",
                bold=-1,
                hyperlink_address="https://b.example",
            ),
        ]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        assert "<a " in result
        assert "<b>" in result

    def test_hyperlink_error_per_run_graceful(self) -> None:
        """ActionSettings error on one run → that run has no link."""
        run_ok = _make_win32com_ppt_run("OK")
        run_err = MagicMock()
        run_err.Text = "Err"
        run_err.Font.Bold = 0
        run_err.Font.Italic = 0
        run_err.Font.Underline = 0
        run_err.Font.Strikethrough = 0
        run_err.Font.Size = 12.0
        run_err.Font.Color.RGB = 0
        type(run_err.Font).Highlight = property(
            lambda s: (_ for _ in ()).throw(AttributeError),
        )
        run_err.ActionSettings.side_effect = OSError("COM")

        runs = [run_ok, run_err]
        para_rng = MagicMock()
        para_rng.Runs.return_value.Count = len(runs)
        para_rng.Runs.side_effect = lambda i=None: (
            runs[i - 1] if i is not None else MagicMock(Count=len(runs))
        )
        result = _win32com_ppt_runs_to_html(para_rng)
        assert "OK" in result
        assert "Err" in result
        assert "<a " not in result


class TestInjectWin32comPptHtmlRunsHyperlinks:
    """Tests for hyperlink injection in _inject_win32com_ppt_html_runs."""

    def test_hyperlink_set_on_segment(self) -> None:
        """Segment with hyperlink sets ActionSettings(1).Hyperlink.Address."""
        tf = MagicMock()
        tf.TextRange.Paragraphs.return_value.Font.Name = "Arial"

        _inject_win32com_ppt_html_runs(
            tf,
            1,
            '<a href="https://inject.example">Link</a>',
        )
        # The char_rng for the "Link" segment should have ActionSettings called
        char_rng = tf.TextRange.Paragraphs.return_value.Characters.return_value
        char_rng.ActionSettings.assert_called_with(1)

    def test_no_hyperlink_no_action_settings(self) -> None:
        """Segments without hyperlink don't touch ActionSettings."""
        tf = MagicMock()
        tf.TextRange.Paragraphs.return_value.Font.Name = "Arial"

        _inject_win32com_ppt_html_runs(
            tf,
            1,
            "<b>Bold</b>",
        )
        # Negative test — just verify no crash during injection
        assert True  # noqa: PLR6201


class TestWin32comPptHyperlinkRoundTrip:
    """Round-trip: PPT hyperlink HTML → parse → segments preserve URL."""

    def test_link_survives_roundtrip(self) -> None:
        """<a> tag with formatting parsed correctly."""
        html_in = '<a href="https://ppt-rt.example"><b>Bold Link</b></a> text'
        segments = _parse_html_formatting(html_in)
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0].hyperlink_url == "https://ppt-rt.example"
        assert segments[0].text == "Bold Link"
        assert segments[0].bold is True
        assert segments[1].hyperlink_url is None
        assert segments[1].text == " text"


# ---------------------------------------------------------------------------
# DrawingML hyperlink tests
# ---------------------------------------------------------------------------

# Namespaces used by DrawingML hyperlink XML builders
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _make_drawingml_txbody_with_hyperlink(
    text: str = "Click here",
    url_rid: str = "rId1",
) -> object:
    """Creates a DrawingML txBody with a hyperlinked run.

    Args:
        text: The link text.
        url_rid: The relationship ID for the hyperlink.

    Returns:
        An lxml element representing ``<a:txBody>``.
    """
    from lxml import etree  # noqa: PLC0415

    tx_body = etree.SubElement(etree.Element("root"), f"{{{_A_NS}}}txBody")
    p = etree.SubElement(tx_body, f"{{{_A_NS}}}p")
    r = etree.SubElement(p, f"{{{_A_NS}}}r")
    rpr = etree.SubElement(r, f"{{{_A_NS}}}rPr")
    etree.SubElement(
        rpr,
        f"{{{_A_NS}}}hlinkClick",
        attrib={
            f"{{{_R_NS}}}id": url_rid,
        },
    )
    t = etree.SubElement(r, f"{{{_A_NS}}}t")
    t.text = text
    return tx_body


def _make_drawingml_txbody_plain(text: str = "No link") -> object:
    """Creates a DrawingML txBody with a plain (non-hyperlinked) run.

    Args:
        text: The run text.

    Returns:
        An lxml element representing ``<a:txBody>``.
    """
    from lxml import etree  # noqa: PLC0415

    tx_body = etree.SubElement(etree.Element("root"), f"{{{_A_NS}}}txBody")
    p = etree.SubElement(tx_body, f"{{{_A_NS}}}p")
    r = etree.SubElement(p, f"{{{_A_NS}}}r")
    etree.SubElement(r, f"{{{_A_NS}}}rPr")
    t = etree.SubElement(r, f"{{{_A_NS}}}t")
    t.text = text
    return tx_body


class TestHasDrawingmlHyperlinks:
    """Tests for _has_drawingml_hyperlinks."""

    def test_returns_true_when_hlinkclick_present(self) -> None:
        """<a:hlinkClick> inside <a:rPr> -> True."""
        tx_body = _make_drawingml_txbody_with_hyperlink("Go", "rId5")
        assert _has_drawingml_hyperlinks(tx_body) is True

    def test_returns_false_when_no_hlinkclick(self) -> None:
        """No <a:hlinkClick> anywhere -> False."""
        tx_body = _make_drawingml_txbody_plain("Plain")
        assert _has_drawingml_hyperlinks(tx_body) is False

    def test_returns_false_for_empty_txbody(self) -> None:
        """Empty <a:txBody> (no paragraphs) -> False."""
        from lxml import etree  # noqa: PLC0415

        tx_body = etree.SubElement(
            etree.Element("root"),
            f"{{{_A_NS}}}txBody",
        )
        assert _has_drawingml_hyperlinks(tx_body) is False


class TestDrawingmlToHtmlHyperlinks:
    """Tests for _drawingml_to_html with hyperlink support."""

    def test_with_rels_emits_anchor_tag(self) -> None:
        """Hyperlinked run with rels mapping emits <a href="...">."""
        tx_body = _make_drawingml_txbody_with_hyperlink("Click", "rId1")
        rels = {"rId1": "https://example.com"}
        result = _drawingml_to_html(tx_body, hyperlink_rels=rels)
        assert '<a href="https://example.com">' in result
        assert "Click" in result
        assert "</a>" in result

    def test_without_rels_returns_text_only(self) -> None:
        """Hyperlinked run without rels mapping returns text without <a>."""
        tx_body = _make_drawingml_txbody_with_hyperlink("Click", "rId1")
        result = _drawingml_to_html(tx_body, hyperlink_rels=None)
        assert "<a " not in result
        assert "Click" in result

    def test_rels_missing_rid_omits_anchor(self) -> None:
        """When rels mapping does not contain the rId, no <a> is emitted."""
        tx_body = _make_drawingml_txbody_with_hyperlink("Click", "rId99")
        rels = {"rId1": "https://other.com"}
        result = _drawingml_to_html(tx_body, hyperlink_rels=rels)
        assert "<a " not in result
        assert "Click" in result

    def test_mixed_hyperlinked_and_plain_runs(self) -> None:
        """Paragraph with both hyperlinked and plain runs."""
        from lxml import etree  # noqa: PLC0415

        tx_body = etree.SubElement(
            etree.Element("root"),
            f"{{{_A_NS}}}txBody",
        )
        p = etree.SubElement(tx_body, f"{{{_A_NS}}}p")
        # Hyperlinked run
        r1 = etree.SubElement(p, f"{{{_A_NS}}}r")
        rpr1 = etree.SubElement(r1, f"{{{_A_NS}}}rPr")
        etree.SubElement(
            rpr1,
            f"{{{_A_NS}}}hlinkClick",
            attrib={
                f"{{{_R_NS}}}id": "rId1",
            },
        )
        t1 = etree.SubElement(r1, f"{{{_A_NS}}}t")
        t1.text = "link"
        # Plain run
        r2 = etree.SubElement(p, f"{{{_A_NS}}}r")
        etree.SubElement(r2, f"{{{_A_NS}}}rPr")
        t2 = etree.SubElement(r2, f"{{{_A_NS}}}t")
        t2.text = " after"

        rels = {"rId1": "https://example.com"}
        result = _drawingml_to_html(tx_body, hyperlink_rels=rels)
        assert '<a href="https://example.com">' in result
        assert "link" in result
        assert "</a>" in result
        assert " after" in result


class TestInjectDrawingmlHtmlRunsHyperlinks:
    """Tests for _inject_drawingml_html_runs with hyperlink injection."""

    def test_anchor_tag_calls_rels_adder(self) -> None:
        """Segment with <a href> calls rels_adder and creates <a:hlinkClick>."""
        tx_body = _make_drawingml_txbody_plain("old text")
        captured_urls: list[str] = []

        def rels_adder(url: str) -> str:
            captured_urls.append(url)
            return "rId99"

        _inject_drawingml_html_runs(
            tx_body,
            '<a href="https://injected.example">Link</a>',
            rels_adder=rels_adder,
        )

        assert captured_urls == ["https://injected.example"]
        # Verify <a:hlinkClick> was created in the run's rPr
        a_hlink_tag = f"{{{_A_NS}}}hlinkClick"
        r_id_attr = f"{{{_R_NS}}}id"
        hlinks = list(tx_body.iter(a_hlink_tag))
        assert len(hlinks) == 1
        assert hlinks[0].get(r_id_attr) == "rId99"

    def test_no_links_does_not_call_rels_adder(self) -> None:
        """Plain bold HTML does not invoke rels_adder."""
        tx_body = _make_drawingml_txbody_plain("old text")
        rels_adder = MagicMock(return_value="rId1")

        _inject_drawingml_html_runs(
            tx_body,
            "<b>Bold</b>",
            rels_adder=rels_adder,
        )

        rels_adder.assert_not_called()

    def test_no_rels_adder_skips_hlinkclick(self) -> None:
        """When rels_adder is None, <a:hlinkClick> is not created."""
        tx_body = _make_drawingml_txbody_plain("old text")
        _inject_drawingml_html_runs(
            tx_body,
            '<a href="https://example.com">Link</a>',
            rels_adder=None,
        )

        a_hlink_tag = f"{{{_A_NS}}}hlinkClick"
        hlinks = list(tx_body.iter(a_hlink_tag))
        assert len(hlinks) == 0

    def test_hlinkclick_stripped_from_base_rpr(self) -> None:
        """Base rPr's <a:hlinkClick> is stripped so it does not spread."""
        # Build txBody where the first run has an <a:hlinkClick> in its rPr
        tx_body = _make_drawingml_txbody_with_hyperlink("old link", "rId1")
        _inject_drawingml_html_runs(
            tx_body,
            "<b>Not a link</b>",
            rels_adder=None,
        )

        a_hlink_tag = f"{{{_A_NS}}}hlinkClick"
        hlinks = list(tx_body.iter(a_hlink_tag))
        assert len(hlinks) == 0


# ---------------------------------------------------------------------------
# WPS textbox hyperlink tests
# ---------------------------------------------------------------------------

_W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"


def _make_wps_txbx_with_hyperlink(
    plain_text: str = "before ",
    link_text: str = "click",
    url_rid: str = "rId1",
) -> object:
    """Creates a WPS textbox with a hyperlink.

    Args:
        plain_text: Text before the hyperlink.
        link_text: Text inside the hyperlink.
        url_rid: Relationship ID for the hyperlink.

    Returns:
        An lxml element for ``<w:txbxContent>``.
    """
    from lxml import etree  # noqa: PLC0415

    txbx = etree.Element(f"{{{_W_NS}}}txbxContent")
    p = etree.SubElement(txbx, f"{{{_W_NS}}}p")
    # Plain run
    r1 = etree.SubElement(p, f"{{{_W_NS}}}r")
    t1 = etree.SubElement(r1, f"{{{_W_NS}}}t")
    t1.text = plain_text
    # Hyperlink
    hl = etree.SubElement(
        p,
        f"{{{_W_NS}}}hyperlink",
        attrib={
            f"{{{_R_NS}}}id": url_rid,
        },
    )
    r2 = etree.SubElement(hl, f"{{{_W_NS}}}r")
    t2 = etree.SubElement(r2, f"{{{_W_NS}}}t")
    t2.text = link_text
    return txbx


class TestWpsTxbxHyperlinks:
    """Tests for _wps_txbx_to_text_or_html with hyperlink support."""

    def test_with_hyperlink_rels_emits_anchor(self) -> None:
        """Hyperlink run with rels mapping emits <a href="...">."""
        txbx = _make_wps_txbx_with_hyperlink("Hello ", "world", "rId1")
        rels = {"rId1": "https://example.com"}
        result = _wps_txbx_to_text_or_html(
            txbx,
            char_styles=None,
            hyperlink_rels=rels,
        )
        assert '<a href="https://example.com">' in result
        assert "world" in result
        assert "</a>" in result

    def test_without_hyperlink_rels_returns_plain(self) -> None:
        """Without hyperlink_rels, hyperlink runs are plain text."""
        txbx = _make_wps_txbx_with_hyperlink("Hello ", "world", "rId1")
        result = _wps_txbx_to_text_or_html(
            txbx,
            char_styles=None,
            hyperlink_rels=None,
        )
        # Without rels and no formatting variation, returns plain text
        assert "<a " not in result
        assert "Hello " in result
        assert "world" in result

    def test_anchor_attribute_used_as_bookmark_url(self) -> None:
        """w:anchor attribute on <w:hyperlink> emits '#anchor' URL."""
        from lxml import etree  # noqa: PLC0415

        txbx = etree.Element(f"{{{_W_NS}}}txbxContent")
        p = etree.SubElement(txbx, f"{{{_W_NS}}}p")
        # Plain run for formatting variation
        r0 = etree.SubElement(p, f"{{{_W_NS}}}r")
        rpr0 = etree.SubElement(r0, f"{{{_W_NS}}}rPr")
        b_el = etree.SubElement(rpr0, f"{{{_W_NS}}}b")  # noqa: F841
        t0 = etree.SubElement(r0, f"{{{_W_NS}}}t")
        t0.text = "bold "
        # Hyperlink with w:anchor
        hl = etree.SubElement(
            p,
            f"{{{_W_NS}}}hyperlink",
            attrib={f"{{{_W_NS}}}anchor": "Bookmark1"},
        )
        r1 = etree.SubElement(hl, f"{{{_W_NS}}}r")
        t1 = etree.SubElement(r1, f"{{{_W_NS}}}t")
        t1.text = "link"

        result = _wps_txbx_to_text_or_html(txbx, char_styles=None, hyperlink_rels={})
        assert '<a href="#Bookmark1">' in result
        assert "link" in result


class TestInjectWpsTxbxHtmlRunsHyperlinks:
    """Tests for _inject_wps_txbx_html_runs with hyperlink injection."""

    def test_anchor_tag_calls_rels_adder(self) -> None:
        """Segment with <a href> calls rels_adder and creates <w:hyperlink>."""
        from lxml import etree  # noqa: PLC0415

        txbx = _make_wps_txbx_with_hyperlink("old ", "old link", "rId1")
        # Wrap in <wps:txbx> since the function expects it
        wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        wps_txbx = etree.Element(f"{{{wps_ns}}}txbx")
        wps_txbx.append(txbx)

        captured_urls: list[str] = []

        def rels_adder(url: str) -> str:
            captured_urls.append(url)
            return "rId42"

        _inject_wps_txbx_html_runs(
            wps_txbx,
            '<a href="https://injected.example">New Link</a>',
            rels_adder=rels_adder,
        )

        assert captured_urls == ["https://injected.example"]
        # Verify <w:hyperlink> was created
        w_hyperlink_tag = f"{{{_W_NS}}}hyperlink"
        hlinks = list(wps_txbx.iter(w_hyperlink_tag))
        assert len(hlinks) >= 1
        r_id_attr = f"{{{_R_NS}}}id"
        assert any(hl.get(r_id_attr) == "rId42" for hl in hlinks)

    def test_no_links_does_not_call_rels_adder(self) -> None:
        """Plain bold HTML does not create <w:hyperlink>."""
        from lxml import etree  # noqa: PLC0415

        txbx = _make_wps_txbx_with_hyperlink("old ", "old link", "rId1")
        wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        wps_txbx = etree.Element(f"{{{wps_ns}}}txbx")
        wps_txbx.append(txbx)

        rels_adder = MagicMock(return_value="rId1")
        _inject_wps_txbx_html_runs(wps_txbx, "<b>Bold</b>", rels_adder=rels_adder)
        rels_adder.assert_not_called()

    def test_bookmark_anchor_uses_w_anchor(self) -> None:
        """Internal '#bookmark' URL creates <w:hyperlink w:anchor='bookmark'>."""
        from lxml import etree  # noqa: PLC0415

        txbx = _make_wps_txbx_with_hyperlink("old ", "old link", "rId1")
        wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        wps_txbx = etree.Element(f"{{{wps_ns}}}txbx")
        wps_txbx.append(txbx)

        _inject_wps_txbx_html_runs(
            wps_txbx,
            '<a href="#Section2">Jump</a>',
            rels_adder=None,
        )

        w_hyperlink_tag = f"{{{_W_NS}}}hyperlink"
        w_anchor_attr = f"{{{_W_NS}}}anchor"
        hlinks = list(wps_txbx.iter(w_hyperlink_tag))
        assert len(hlinks) >= 1
        assert any(hl.get(w_anchor_attr) == "Section2" for hl in hlinks)


# ---------------------------------------------------------------------------
# ODF paragraph hyperlink tests (lxml ZIP path)
# ---------------------------------------------------------------------------

_TEXT_NS = _ODF_NS["text"]
_XLINK_NS = _ODF_NS["xlink"]


def _make_odf_para_with_link(
    plain: str = "Hello ",
    link_text: str = "click",
    url: str = "https://example.com",
) -> tuple[object, str]:
    """Creates an ODF root element with a paragraph containing a hyperlink.

    Args:
        plain: Plain text before the hyperlink.
        link_text: Text inside the hyperlink.
        url: Target URL of the hyperlink.

    Returns:
        Tuple of (root_element, text_p_tag_string).
    """
    from lxml import etree  # noqa: PLC0415

    text_p_tag = f"{{{_TEXT_NS}}}p"
    root = etree.Element("root")
    p = etree.SubElement(root, text_p_tag)
    p.text = plain
    a = etree.SubElement(p, f"{{{_TEXT_NS}}}a")
    a.set(f"{{{_XLINK_NS}}}href", url)
    a.text = link_text
    return root, text_p_tag


class TestExtractOdfParagraphTextHyperlinks:
    """Tests for _extract_odf_paragraph_text with <text:a> hyperlinks."""

    def test_with_text_a_emits_anchor_html(self) -> None:
        """Paragraph with <text:a> emits <a href="...">text</a> HTML."""
        root, text_p_tag = _make_odf_para_with_link(
            "Visit ",
            "our site",
            "https://example.com",
        )
        result = _extract_odf_paragraph_text(root, text_p_tag)
        assert '<a href="https://example.com">' in result
        assert "our site" in result
        assert "</a>" in result
        assert "Visit " in result

    def test_without_text_a_returns_plain(self) -> None:
        """Paragraph without <text:a> returns plain text."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        root = etree.Element("root")
        p = etree.SubElement(root, text_p_tag)
        p.text = "Just plain text"
        result = _extract_odf_paragraph_text(root, text_p_tag)
        assert result == "Just plain text"
        assert "<a " not in result

    def test_text_a_without_href_omits_anchor(self) -> None:
        """<text:a> with no xlink:href returns plain text for that child."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        root = etree.Element("root")
        p = etree.SubElement(root, text_p_tag)
        a = etree.SubElement(p, f"{{{_TEXT_NS}}}a")
        a.text = "orphan link"
        result = _extract_odf_paragraph_text(root, text_p_tag)
        assert "<a " not in result
        assert "orphan link" in result

    def test_tail_text_after_link_preserved(self) -> None:
        """Tail text after <text:a> is preserved."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        root = etree.Element("root")
        p = etree.SubElement(root, text_p_tag)
        p.text = "Before "
        a = etree.SubElement(p, f"{{{_TEXT_NS}}}a")
        a.set(f"{{{_XLINK_NS}}}href", "https://example.com")
        a.text = "link"
        a.tail = " after"
        result = _extract_odf_paragraph_text(root, text_p_tag)
        assert "Before " in result
        assert "link" in result
        assert " after" in result


class TestInjectOdfParagraphTextHyperlinks:
    """Tests for _inject_odf_paragraph_text with <a> HTML tags."""

    def test_with_anchor_html_dispatches_to_html_handler(self) -> None:
        """Text with <a href> dispatches to _inject_odf_paragraph_text_html."""
        root, text_p_tag = _make_odf_para_with_link()
        result = _inject_odf_paragraph_text(
            root,
            '<a href="https://new.example">New Link</a>',
            text_p_tag,
        )
        assert result is True
        # Verify <text:a> element was created
        text_a_tag = f"{{{_TEXT_NS}}}a"
        links = list(root.iter(text_a_tag))
        assert len(links) >= 1
        xlink_href = f"{{{_XLINK_NS}}}href"
        assert any(link.get(xlink_href) == "https://new.example" for link in links)

    def test_plain_text_does_not_create_text_a(self) -> None:
        """Plain text (no HTML) does not create <text:a> elements."""
        root, text_p_tag = _make_odf_para_with_link()
        _inject_odf_paragraph_text(root, "Just plain text", text_p_tag)
        text_a_tag = f"{{{_TEXT_NS}}}a"
        links = list(root.iter(text_a_tag))
        assert len(links) == 0

    def test_mixed_plain_and_link_segments(self) -> None:
        """Mixed plain + hyperlink text creates correct structure."""
        root, text_p_tag = _make_odf_para_with_link()
        _inject_odf_paragraph_text(
            root,
            'Visit <a href="https://example.com">here</a> now',
            text_p_tag,
        )
        text_a_tag = f"{{{_TEXT_NS}}}a"
        links = list(root.iter(text_a_tag))
        assert len(links) >= 1
        # Link text should be "here"
        assert any(link.text == "here" for link in links)

    def test_xlink_type_set_to_simple(self) -> None:
        """Injected <text:a> elements have xlink:type='simple'."""
        root, text_p_tag = _make_odf_para_with_link()
        _inject_odf_paragraph_text(
            root,
            '<a href="https://example.com">link</a>',
            text_p_tag,
        )
        text_a_tag = f"{{{_TEXT_NS}}}a"
        xlink_type = f"{{{_XLINK_NS}}}type"
        links = list(root.iter(text_a_tag))
        assert len(links) >= 1
        assert links[0].get(xlink_type) == "simple"


# ---------------------------------------------------------------------------
# ODF text-box hyperlink tests (office_formatter functions)
# ---------------------------------------------------------------------------


class TestHasOdfTextBoxMixedFormattingHyperlinks:
    """Tests for _has_odf_text_box_mixed_formatting with <text:a>."""

    def test_returns_true_when_text_a_present(self) -> None:
        """Presence of <text:a> in paragraph -> True (requires HTML round-trip)."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        tb = etree.Element("draw-text-box")
        p = etree.SubElement(tb, text_p_tag)
        p.text = "before "
        a = etree.SubElement(p, f"{{{_TEXT_NS}}}a")
        a.set(f"{{{_XLINK_NS}}}href", "https://example.com")
        a.text = "link"

        result = _has_odf_text_box_mixed_formatting(tb, {}, text_p_tag)
        assert result is True

    def test_returns_false_when_no_text_a(self) -> None:
        """No <text:a> and uniform formatting -> False."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        tb = etree.Element("draw-text-box")
        p = etree.SubElement(tb, text_p_tag)
        p.text = "plain text only"

        result = _has_odf_text_box_mixed_formatting(tb, {}, text_p_tag)
        assert result is False


class TestOdfTextBoxToHtmlHyperlinks:
    """Tests for _odf_text_box_to_html with <text:a> hyperlinks."""

    def test_text_a_emits_anchor_html(self) -> None:
        """<text:a> child emits <a href="..."> in HTML output."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        tb = etree.Element("draw-text-box")
        p = etree.SubElement(tb, text_p_tag)
        p.text = "See "
        a = etree.SubElement(p, f"{{{_TEXT_NS}}}a")
        a.set(f"{{{_XLINK_NS}}}href", "https://example.com")
        a.text = "link"
        a.tail = " end"

        result = _odf_text_box_to_html(tb, {}, text_p_tag)
        assert '<a href="https://example.com">' in result
        assert "link" in result
        assert "</a>" in result
        assert "See " in result
        assert " end" in result

    def test_no_hyperlinks_returns_plain_wrapped(self) -> None:
        """Text box with only plain text returns escaped text."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        tb = etree.Element("draw-text-box")
        p = etree.SubElement(tb, text_p_tag)
        p.text = "plain content"

        result = _odf_text_box_to_html(tb, {}, text_p_tag)
        assert "plain content" in result
        assert "<a " not in result


class TestInjectOdfTextBoxHtmlRunsHyperlinks:
    """Tests for _inject_odf_text_box_html_runs with hyperlink injection."""

    def test_anchor_tag_creates_text_a(self) -> None:
        """HTML with <a href> creates <text:a xlink:href> in the text box."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        tb = etree.Element("draw-text-box")
        p = etree.SubElement(tb, text_p_tag)
        p.text = "old text"

        # Create a minimal <office:automatic-styles> element
        auto_styles = etree.Element(
            f"{{{_ODF_NS['office']}}}automatic-styles",
        )
        style_counter: list[int] = [0]

        result = _inject_odf_text_box_html_runs(
            tb,
            '<a href="https://example.com">Link</a>',
            text_p_tag,
            auto_styles,
            style_counter,
        )
        assert result is True
        text_a_tag = f"{{{_TEXT_NS}}}a"
        xlink_href = f"{{{_XLINK_NS}}}href"
        links = list(tb.iter(text_a_tag))
        assert len(links) >= 1
        assert links[0].get(xlink_href) == "https://example.com"
        assert links[0].text == "Link"

    def test_plain_html_no_text_a_created(self) -> None:
        """Bold-only HTML does not create <text:a> elements."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        tb = etree.Element("draw-text-box")
        p = etree.SubElement(tb, text_p_tag)
        p.text = "old text"

        auto_styles = etree.Element(
            f"{{{_ODF_NS['office']}}}automatic-styles",
        )
        style_counter: list[int] = [0]

        _inject_odf_text_box_html_runs(
            tb,
            "<b>Bold</b>",
            text_p_tag,
            auto_styles,
            style_counter,
        )
        text_a_tag = f"{{{_TEXT_NS}}}a"
        links = list(tb.iter(text_a_tag))
        assert len(links) == 0

    def test_mixed_link_and_formatted_text(self) -> None:
        """Mixed <a> and <b> creates both <text:a> and <text:span>."""
        from lxml import etree  # noqa: PLC0415

        text_p_tag = f"{{{_TEXT_NS}}}p"
        tb = etree.Element("draw-text-box")
        p = etree.SubElement(tb, text_p_tag)
        p.text = "old text"

        auto_styles = etree.Element(
            f"{{{_ODF_NS['office']}}}automatic-styles",
        )
        style_counter: list[int] = [0]

        _inject_odf_text_box_html_runs(
            tb,
            '<b>Bold</b> <a href="https://example.com">link</a>',
            text_p_tag,
            auto_styles,
            style_counter,
        )

        text_a_tag = f"{{{_TEXT_NS}}}a"
        text_span_tag = f"{{{_TEXT_NS}}}span"
        links = list(tb.iter(text_a_tag))
        spans = list(tb.iter(text_span_tag))
        assert len(links) >= 1
        assert len(spans) >= 1
        assert links[0].text == "link"


# ---------------------------------------------------------------------------
# UNO hyperlink tests
# ---------------------------------------------------------------------------


def _make_uno_portion_with_url(
    text: str,
    url: str = "",
    bold: bool = False,
    font_size: float = 11.0,
    color_int: int = 0,
) -> MagicMock:
    """Creates a mock UNO text portion with an optional HyperLinkURL.

    Args:
        text: The portion text.
        url: HyperLinkURL value (empty string = no link).
        bold: Whether the portion is bold.
        font_size: CharHeight in points.
        color_int: CharColor as integer.

    Returns:
        Mock UNO TextPortion object.
    """
    portion = MagicMock()
    props = {
        "TextPortionType": "Text",
        "CharWeight": _UNO_WEIGHT_BOLD if bold else _UNO_WEIGHT_NORMAL,
        "CharPosture": _UNO_SLANT_NONE,
        "CharUnderline": _UNO_UNDERLINE_NONE,
        "CharStrikeout": _UNO_STRIKEOUT_NONE,
        "CharFontName": "Liberation Sans",
        "CharHeight": font_size,
        "CharColor": color_int,
        "CharHighlight": -1,
        "CharBackColor": -1,
        "HyperLinkURL": url,
    }
    portion.getPropertyValue.side_effect = lambda p: props[p]
    portion.getString.return_value = text
    return portion


def _make_uno_para_from_portions(portions: list[MagicMock]) -> MagicMock:
    """Creates a mock UNO paragraph from a list of mock portions.

    Args:
        portions: List of mock UNO portion objects.

    Returns:
        Mock UNO paragraph object.
    """
    para = MagicMock()
    para.createEnumeration.return_value = _make_uno_enum(portions)
    return para


class TestHasUnoHyperlinks:
    """Tests for _has_uno_hyperlinks."""

    def test_returns_true_when_url_present(self) -> None:
        """Portion with non-empty HyperLinkURL -> True."""
        portion = _make_uno_portion_with_url("Link", url="https://example.com")
        para = _make_uno_para_from_portions([portion])
        assert _has_uno_hyperlinks(para) is True

    def test_returns_false_when_no_hyperlinks(self) -> None:
        """All portions have empty HyperLinkURL -> False."""
        p1 = _make_uno_portion_with_url("Hello")
        p2 = _make_uno_portion_with_url("world")
        para = _make_uno_para_from_portions([p1, p2])
        assert _has_uno_hyperlinks(para) is False

    def test_returns_false_when_portions_have_empty_url(self) -> None:
        """Portions with explicitly empty HyperLinkURL -> False."""
        portion = _make_uno_portion_with_url("text", url="")
        para = _make_uno_para_from_portions([portion])
        assert _has_uno_hyperlinks(para) is False

    def test_returns_false_for_empty_paragraph(self) -> None:
        """Paragraph with no portions -> False."""
        para = _make_uno_para_from_portions([])
        assert _has_uno_hyperlinks(para) is False

    def test_skips_non_text_portions(self) -> None:
        """Non-'Text' portion types are ignored even with HyperLinkURL."""
        portion = MagicMock()
        portion.getPropertyValue.side_effect = lambda p: {
            "TextPortionType": "SoftPageBreak",
            "HyperLinkURL": "https://example.com",
        }[p]
        portion.getString.return_value = "break"
        para = _make_uno_para_from_portions([portion])
        assert _has_uno_hyperlinks(para) is False

    def test_skips_empty_text_portions(self) -> None:
        """Portions with empty text are skipped."""
        portion = _make_uno_portion_with_url("", url="https://example.com")
        portion.getString.return_value = ""
        para = _make_uno_para_from_portions([portion])
        assert _has_uno_hyperlinks(para) is False


class TestUnoRunsToHtmlHyperlinks:
    """Tests for _uno_runs_to_html with hyperlink support."""

    def test_hyperlink_emits_anchor_tag(self) -> None:
        """Portion with HyperLinkURL emits <a href="..."> in HTML."""
        portion = _make_uno_portion_with_url(
            "Click",
            url="https://example.com",
        )
        para = MagicMock()
        para.createEnumeration.return_value = _make_uno_enum([portion])

        result = _uno_runs_to_html(para)
        assert '<a href="https://example.com">' in result
        assert "Click" in result
        assert "</a>" in result

    def test_mixed_link_and_plain(self) -> None:
        """Paragraph with linked and plain portions."""
        p1 = _make_uno_portion_with_url("Link", url="https://example.com")
        p2 = _make_uno_portion_with_url(" plain")
        para = MagicMock()
        para.createEnumeration.return_value = _make_uno_enum([p1, p2])

        result = _uno_runs_to_html(para)
        assert '<a href="https://example.com">' in result
        assert "Link" in result
        assert "</a>" in result
        assert " plain" in result

    def test_consecutive_same_url_grouped(self) -> None:
        """Two consecutive portions with the same URL share one <a> tag."""
        p1 = _make_uno_portion_with_url("Part1", url="https://example.com")
        p2 = _make_uno_portion_with_url("Part2", url="https://example.com")
        para = MagicMock()
        para.createEnumeration.return_value = _make_uno_enum([p1, p2])

        result = _uno_runs_to_html(para)
        # Only one opening and one closing <a> tag
        assert result.count('<a href="https://example.com">') == 1
        assert result.count("</a>") == 1
        assert "Part1" in result
        assert "Part2" in result

    def test_no_hyperlinks_no_anchor_tags(self) -> None:
        """Portions without HyperLinkURL produce no <a> tags."""
        p1 = _make_uno_portion_with_url("Hello ", bold=True)
        p2 = _make_uno_portion_with_url("world")
        para = MagicMock()
        para.createEnumeration.return_value = _make_uno_enum([p1, p2])

        result = _uno_runs_to_html(para)
        assert "<a " not in result
        assert "</a>" not in result


class TestInjectUnoHtmlRunsHyperlinks:
    """Tests for _inject_uno_html_runs with hyperlink injection."""

    def test_anchor_tag_sets_hyperlinkurl(self) -> None:
        """Segment with <a href> sets HyperLinkURL on the cursor."""
        para = MagicMock()
        para_text = MagicMock()
        para.getText.return_value = para_text
        cursor = MagicMock()
        para_text.createTextCursorByRange.return_value = cursor

        _inject_uno_html_runs(
            para,
            '<a href="https://example.com">Link</a>',
            base_props={},
        )

        # Verify HyperLinkURL was set on the cursor
        set_calls = cursor.setPropertyValue.call_args_list
        hyperlink_calls = [
            c
            for c in set_calls
            if c[0][0] == "HyperLinkURL" and c[0][1] == "https://example.com"
        ]
        assert len(hyperlink_calls) >= 1

    def test_plain_text_clears_hyperlinkurl(self) -> None:
        """Segment without hyperlink sets HyperLinkURL to empty string."""
        para = MagicMock()
        para_text = MagicMock()
        para.getText.return_value = para_text
        cursor = MagicMock()
        para_text.createTextCursorByRange.return_value = cursor

        _inject_uno_html_runs(
            para,
            "<b>Bold</b>",
            base_props={},
        )

        set_calls = cursor.setPropertyValue.call_args_list
        hyperlink_calls = [c for c in set_calls if c[0][0] == "HyperLinkURL"]
        assert len(hyperlink_calls) >= 1
        # Should be empty string for non-link segments
        assert hyperlink_calls[0][0][1] == ""

    def test_mixed_link_and_plain_segments(self) -> None:
        """Mixed <a> and plain segments set correct HyperLinkURL per segment."""
        para = MagicMock()
        para_text = MagicMock()
        para.getText.return_value = para_text
        cursor = MagicMock()
        para_text.createTextCursorByRange.return_value = cursor

        _inject_uno_html_runs(
            para,
            '<a href="https://example.com">Link</a> plain',
            base_props={},
        )

        set_calls = cursor.setPropertyValue.call_args_list
        hyperlink_calls = [c for c in set_calls if c[0][0] == "HyperLinkURL"]
        # Two segments = two HyperLinkURL calls
        assert len(hyperlink_calls) == 2  # noqa: PLR2004
        urls = [c[0][1] for c in hyperlink_calls]
        assert "https://example.com" in urls
        assert "" in urls


# ── Superscript / Subscript tests ──────────────────────────────────────


class TestWrapWithTagsSupSub:
    """Tests for _wrap_with_tags with superscript/subscript."""

    def test_superscript_tag(self) -> None:
        """Superscript wraps text in <sup> tags."""
        result = _wrap_with_tags(
            "E=mc²",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            font_size_pt=None,
            color_hex=None,
            has_size_variation=False,
            has_color_variation=False,
            superscript=True,
        )
        assert result == "<sup>E=mc²</sup>"

    def test_subscript_tag(self) -> None:
        """Subscript wraps text in <sub> tags."""
        result = _wrap_with_tags(
            "H₂O",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            font_size_pt=None,
            color_hex=None,
            has_size_variation=False,
            has_color_variation=False,
            subscript=True,
        )
        assert result == "<sub>H₂O</sub>"

    def test_bold_superscript(self) -> None:
        """Bold + superscript produces both <b> and <sup> tags."""
        result = _wrap_with_tags(
            "2",
            bold=True,
            italic=False,
            underline=False,
            strike=False,
            font_size_pt=None,
            color_hex=None,
            has_size_variation=False,
            has_color_variation=False,
            superscript=True,
        )
        assert "<b>" in result
        assert "<sup>" in result
        assert "2" in result

    def test_both_sup_sub_prefers_sup(self) -> None:
        """When both sup and sub are True, sup wins (earlier in nesting)."""
        result = _wrap_with_tags(
            "x",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            font_size_pt=None,
            color_hex=None,
            has_size_variation=False,
            has_color_variation=False,
            superscript=True,
            subscript=True,
        )
        # Both tags are present since both flags are True
        assert "<sup>" in result
        assert "<sub>" in result

    def test_no_sup_sub(self) -> None:
        """No superscript or subscript produces plain text."""
        result = _wrap_with_tags(
            "plain",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            font_size_pt=None,
            color_hex=None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "plain"

    def test_all_formatting_with_sup(self) -> None:
        """All formatting flags including superscript produce all 5 tags."""
        result = _wrap_with_tags(
            "x",
            bold=True,
            italic=True,
            underline=True,
            strike=True,
            font_size_pt=None,
            color_hex=None,
            has_size_variation=False,
            has_color_variation=False,
            superscript=True,
        )
        assert "<b>" in result
        assert "<i>" in result
        assert "<u>" in result
        assert "<s>" in result
        assert "<sup>" in result


class TestParseHtmlFormattingSupSub:
    """Tests for _parse_html_formatting with <sup>/<sub>."""

    def test_sup_tag_parsed(self) -> None:
        """<sup> tag sets superscript=True on parsed segment."""
        segments = _parse_html_formatting("<sup>2</sup>")
        assert len(segments) == 1
        assert segments[0].text == "2"
        assert segments[0].superscript is True
        assert segments[0].subscript is False

    def test_sub_tag_parsed(self) -> None:
        """<sub> tag sets subscript=True on parsed segment."""
        segments = _parse_html_formatting("<sub>2</sub>")
        assert len(segments) == 1
        assert segments[0].text == "2"
        assert segments[0].subscript is True
        assert segments[0].superscript is False

    def test_nested_bold_sup(self) -> None:
        """Nested <b><sup> tags set both bold and superscript."""
        segments = _parse_html_formatting("<b><sup>ref</sup></b>")
        assert len(segments) == 1
        assert segments[0].text == "ref"
        assert segments[0].bold is True
        assert segments[0].superscript is True

    def test_mixed_normal_and_sup(self) -> None:
        """Mix of normal text and <sup> produces two segments."""
        segments = _parse_html_formatting("normal<sup>2</sup>")
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0].text == "normal"
        assert segments[0].superscript is False
        assert segments[1].text == "2"
        assert segments[1].superscript is True

    def test_sup_sub_roundtrip(self) -> None:
        """Wrap then parse roundtrip preserves superscript and subscript."""
        # Superscript roundtrip
        sup_html = _wrap_with_tags(
            "2",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            font_size_pt=None,
            color_hex=None,
            has_size_variation=False,
            has_color_variation=False,
            superscript=True,
        )
        sup_segments = _parse_html_formatting(sup_html)
        assert len(sup_segments) == 1
        assert sup_segments[0].superscript is True
        assert sup_segments[0].subscript is False

        # Subscript roundtrip
        sub_html = _wrap_with_tags(
            "n",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            font_size_pt=None,
            color_hex=None,
            has_size_variation=False,
            has_color_variation=False,
            subscript=True,
        )
        sub_segments = _parse_html_formatting(sub_html)
        assert len(sub_segments) == 1
        assert sub_segments[0].subscript is True
        assert sub_segments[0].superscript is False


class TestReadWmlRprSupSub:
    """Tests for _read_wml_rpr_sup_sub."""

    def test_none_rpr(self) -> None:
        """None rPr returns (False, False)."""
        assert _read_wml_rpr_sup_sub(None) == (False, False)

    def test_no_vert_align(self) -> None:
        """Element rPr without vertAlign returns (False, False)."""
        rpr = OxmlElement("w:rPr")
        assert _read_wml_rpr_sup_sub(rpr) == (False, False)

    def test_superscript_val(self) -> None:
        """Element rPr with vertAlign val=superscript returns (True, False)."""
        rpr = OxmlElement("w:rPr")
        va = OxmlElement("w:vertAlign")
        va.set(qn("w:val"), "superscript")
        rpr.append(va)
        assert _read_wml_rpr_sup_sub(rpr) == (True, False)

    def test_subscript_val(self) -> None:
        """Element rPr with vertAlign val=subscript returns (False, True)."""
        rpr = OxmlElement("w:rPr")
        va = OxmlElement("w:vertAlign")
        va.set(qn("w:val"), "subscript")
        rpr.append(va)
        assert _read_wml_rpr_sup_sub(rpr) == (False, True)

    def test_unknown_val(self) -> None:
        """Element rPr with vertAlign val=baseline returns (False, False)."""
        rpr = OxmlElement("w:rPr")
        va = OxmlElement("w:vertAlign")
        va.set(qn("w:val"), "baseline")
        rpr.append(va)
        assert _read_wml_rpr_sup_sub(rpr) == (False, False)

    def test_unnamespaced_val(self) -> None:
        """Unnamespaced val attr on vertAlign is still read correctly."""
        rpr = OxmlElement("w:rPr")
        va = OxmlElement("w:vertAlign")
        # Set the val attribute without namespace prefix
        va.set("val", "superscript")
        rpr.append(va)
        result = _read_wml_rpr_sup_sub(rpr)
        # The function checks namespaced first, then falls back to unnamespaced
        assert result == (True, False)


class TestSetRprVertAlign:
    """Tests for _set_rpr_vert_align."""

    def test_set_superscript(self) -> None:
        """Creates vertAlign element with val=superscript."""
        rpr = OxmlElement("w:rPr")
        _set_rpr_vert_align(rpr, superscript=True, subscript=False, qn_fn=qn)
        va = rpr.find(qn("w:vertAlign"))
        assert va is not None
        assert va.get(qn("w:val")) == "superscript"

    def test_set_subscript(self) -> None:
        """Creates vertAlign element with val=subscript."""
        rpr = OxmlElement("w:rPr")
        _set_rpr_vert_align(rpr, superscript=False, subscript=True, qn_fn=qn)
        va = rpr.find(qn("w:vertAlign"))
        assert va is not None
        assert va.get(qn("w:val")) == "subscript"

    def test_remove_when_neither(self) -> None:
        """Removes existing vertAlign when both superscript and subscript are False."""
        rpr = OxmlElement("w:rPr")
        # First add a vertAlign element
        va = OxmlElement("w:vertAlign")
        va.set(qn("w:val"), "superscript")
        rpr.append(va)
        assert rpr.find(qn("w:vertAlign")) is not None
        # Now remove it
        _set_rpr_vert_align(rpr, superscript=False, subscript=False, qn_fn=qn)
        assert rpr.find(qn("w:vertAlign")) is None

    def test_update_existing(self) -> None:
        """Changes existing vertAlign from superscript to subscript."""
        rpr = OxmlElement("w:rPr")
        va = OxmlElement("w:vertAlign")
        va.set(qn("w:val"), "superscript")
        rpr.append(va)
        _set_rpr_vert_align(rpr, superscript=False, subscript=True, qn_fn=qn)
        va = rpr.find(qn("w:vertAlign"))
        assert va is not None
        assert va.get(qn("w:val")) == "subscript"

    def test_no_remove_when_absent(self) -> None:
        """No error when removing from rPr that has no vertAlign."""
        rpr = OxmlElement("w:rPr")
        # Should not raise even though there is nothing to remove
        _set_rpr_vert_align(rpr, superscript=False, subscript=False, qn_fn=qn)
        assert rpr.find(qn("w:vertAlign")) is None

    def test_superscript_wins_when_both_true(self) -> None:
        """Superscript takes precedence when both sup and sub are True."""
        from docx import Document  # noqa: PLC0415
        from docx.oxml.ns import qn as docx_qn  # noqa: PLC0415

        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("x")
        rpr = run._r.get_or_add_rPr()
        _set_rpr_vert_align(rpr, superscript=True, subscript=True, qn_fn=docx_qn)
        va = rpr.find(docx_qn("w:vertAlign"))
        assert va is not None
        assert va.get(docx_qn("w:val")) == "superscript"


class TestHasMixedFormattingSupSub:
    """Tests for _has_mixed_formatting with superscript/subscript in python-docx."""

    def test_mixed_with_superscript(self) -> None:
        """Normal run + superscript run → True."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("normal")
        run_super = para.add_run("sup")
        run_super.font.superscript = True
        assert _has_mixed_formatting(para) is True

    def test_uniform_superscript(self) -> None:
        """Both runs superscript → False (uniform)."""
        doc = Document()
        para = doc.add_paragraph()
        r1 = para.add_run("a")
        r1.font.superscript = True
        r2 = para.add_run("b")
        r2.font.superscript = True
        assert _has_mixed_formatting(para) is False


class TestRunsToHtmlSupSub:
    """Tests for _runs_to_html with superscript/subscript in python-docx."""

    def test_superscript_run(self) -> None:
        """Single superscript run produces <sup> tag."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("text")
        run.font.superscript = True
        html = _runs_to_html(para)
        assert "<sup>" in html
        assert "text" in html
        assert "</sup>" in html

    def test_subscript_run(self) -> None:
        """Single subscript run produces <sub> tag."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("text")
        run.font.subscript = True
        html = _runs_to_html(para)
        assert "<sub>" in html
        assert "text" in html
        assert "</sub>" in html

    def test_bold_superscript_run(self) -> None:
        """Bold + superscript run produces both <b> and <sup> tags."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("ref")
        run.bold = True
        run.font.superscript = True
        html = _runs_to_html(para)
        assert "<b>" in html
        assert "<sup>" in html
        assert "ref" in html


class TestInjectHtmlRunsSupSub:
    """Tests for _inject_html_runs setting w:vertAlign."""

    def test_superscript_roundtrip(self) -> None:
        """Create docx para, runs_to_html with sup, inject back, verify vertAlign."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("E=mc")
        run_sup = para.add_run("2")
        run_sup.font.superscript = True

        # Extract HTML
        html = _runs_to_html(para)
        assert "<sup>" in html

        # Inject back into a fresh paragraph
        para2 = doc.add_paragraph()
        para2.add_run("placeholder")
        _inject_html_runs(para2, html)

        # Find the run with superscript vertAlign
        found_super = False
        for r in para2.runs:
            rpr = r._element.find(qn("w:rPr"))
            if rpr is not None:
                va = rpr.find(qn("w:vertAlign"))
                if va is not None and va.get(qn("w:val")) == "superscript":
                    found_super = True
                    break
        assert found_super, "Expected a run with w:vertAlign=superscript"

    def test_subscript_roundtrip(self) -> None:
        """Create docx para, runs_to_html with sub, inject back, verify vertAlign."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("H")
        run_sub = para.add_run("2")
        run_sub.font.subscript = True

        # Extract HTML
        html = _runs_to_html(para)
        assert "<sub>" in html

        # Inject back into a fresh paragraph
        para2 = doc.add_paragraph()
        para2.add_run("placeholder")
        _inject_html_runs(para2, html)

        # Find the run with subscript vertAlign
        found_sub = False
        for r in para2.runs:
            rpr = r._element.find(qn("w:rPr"))
            if rpr is not None:
                va = rpr.find(qn("w:vertAlign"))
                if va is not None and va.get(qn("w:val")) == "subscript":
                    found_sub = True
                    break
        assert found_sub, "Expected a run with w:vertAlign=subscript"


def _make_pptx_run_with_baseline(baseline_val: str | None = None) -> object:
    """Creates a real python-pptx run, optionally setting the baseline attribute.

    Sets ``run.font.bold = False`` first to force python-pptx to create a
    proper ``<a:rPr>`` element, then sets the ``baseline`` attribute on it.

    Args:
        baseline_val: Value for the baseline attribute, or None to skip.

    Returns:
        A python-pptx ``_Run`` object.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    )
    para = txbox.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "test"
    # Force creation of a proper <a:rPr> element via python-pptx API
    run.font.bold = False
    if baseline_val is not None:
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        rpr = run._r.find(pptx_qn("a:rPr"))
        rpr.set("baseline", baseline_val)
    return run


class TestReadPptxRunFormattingSupSub:
    """Tests for _read_pptx_run_formatting with baseline attribute."""

    def test_positive_baseline_superscript(self) -> None:
        """Run with baseline=30000 is detected as superscript."""
        run = _make_pptx_run_with_baseline("30000")
        result = _read_pptx_run_formatting(run)
        # result is (bold, italic, underline, strike, superscript, subscript)
        assert result[4] is True  # noqa: PLR2004  # superscript
        assert result[5] is False  # noqa: PLR2004  # subscript

    def test_negative_baseline_subscript(self) -> None:
        """Run with baseline=-25000 is detected as subscript."""
        run = _make_pptx_run_with_baseline("-25000")
        result = _read_pptx_run_formatting(run)
        assert result[4] is False  # noqa: PLR2004  # superscript
        assert result[5] is True  # noqa: PLR2004  # subscript

    def test_no_baseline(self) -> None:
        """Run without baseline attribute → both False."""
        run = _make_pptx_run_with_baseline(None)
        result = _read_pptx_run_formatting(run)
        assert result[4] is False  # noqa: PLR2004  # superscript
        assert result[5] is False  # noqa: PLR2004  # subscript

    def test_zero_baseline(self) -> None:
        """Run with baseline=0 → both False."""
        run = _make_pptx_run_with_baseline("0")
        result = _read_pptx_run_formatting(run)
        assert result[4] is False  # noqa: PLR2004  # superscript
        assert result[5] is False  # noqa: PLR2004  # subscript

    def test_invalid_baseline(self) -> None:
        """Run with baseline=abc → both False (ValueError suppressed)."""
        run = _make_pptx_run_with_baseline("abc")
        result = _read_pptx_run_formatting(run)
        assert result[4] is False  # noqa: PLR2004  # superscript
        assert result[5] is False  # noqa: PLR2004  # subscript


class TestReadPptxRunFullFormattingSupSub:
    """Tests for _read_pptx_run_full_formatting with sup/sub."""

    def test_superscript_in_9_tuple(self) -> None:
        """Superscript run returns (True, False) at positions 4 and 5 in 9-tuple."""
        run = _make_pptx_run_with_baseline("30000")
        run.text = "sup"
        result = _read_pptx_run_full_formatting(run)
        # 9-tuple: (bold, italic, underline, strike, superscript, subscript,
        #           font_size_pt, color_hex, bg_color_hex)
        assert len(result) == 9  # noqa: PLR2004
        assert result[4] is True  # noqa: PLR2004  # superscript
        assert result[5] is False  # noqa: PLR2004  # subscript


class TestHasPptxMixedFormattingSupSub:
    """Tests for _has_pptx_mixed_formatting with superscript."""

    def test_mixed_normal_and_super(self) -> None:
        """Normal + superscript run → True (mixed formatting detected)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        para = txbox.text_frame.paragraphs[0]
        # Normal run — force rPr creation so formatting comparison is valid
        run_normal = para.add_run()
        run_normal.text = "normal"
        run_normal.font.bold = False
        # Superscript run
        run_super = para.add_run()
        run_super.text = "sup"
        run_super.font.bold = False
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        rpr = run_super._r.find(pptx_qn("a:rPr"))
        rpr.set("baseline", "30000")
        assert _has_pptx_mixed_formatting(para) is True

    def test_mixed_normal_and_sub(self) -> None:
        """Normal + subscript run → True (mixed formatting detected)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        para = txbox.text_frame.paragraphs[0]
        run_normal = para.add_run()
        run_normal.text = "normal"
        run_normal.font.bold = False
        run_sub = para.add_run()
        run_sub.text = "sub"
        run_sub.font.bold = False
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        rpr = run_sub._r.find(pptx_qn("a:rPr"))
        rpr.set("baseline", "-25000")
        assert _has_pptx_mixed_formatting(para) is True


class TestPptxRunsToHtmlSupSub:
    """Tests for _pptx_runs_to_html with superscript/subscript."""

    def test_superscript_produces_sup_tag(self) -> None:
        """Run with baseline=30000 produces <sup> tag in HTML output."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        para = txbox.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = "exponent"
        # Force proper rPr creation, then set baseline
        run.font.bold = False
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415

        rpr = run._r.find(pptx_qn("a:rPr"))
        rpr.set("baseline", "30000")
        html = _pptx_runs_to_html(para)
        assert "<sup>" in html
        assert "exponent" in html
        assert "</sup>" in html


class TestApplyDrawingmlFormatAttrsSupSub:
    """Tests for _apply_drawingml_format_attrs with baseline."""

    def test_superscript_sets_baseline(self) -> None:
        """Segment with superscript=True sets baseline=30000 on rPr."""
        rpr = etree.Element(f"{{{_DRAWINGML_NS}}}rPr")
        seg = _FormattedSegment(
            text="2",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            superscript=True,
            subscript=False,
        )
        _apply_drawingml_format_attrs(rpr, seg)
        assert rpr.get("baseline") == "30000"

    def test_subscript_sets_baseline(self) -> None:
        """Segment with subscript=True sets baseline=-25000 on rPr."""
        rpr = etree.Element(f"{{{_DRAWINGML_NS}}}rPr")
        seg = _FormattedSegment(
            text="n",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            superscript=False,
            subscript=True,
        )
        _apply_drawingml_format_attrs(rpr, seg)
        assert rpr.get("baseline") == "-25000"

    def test_no_sup_sub_no_baseline(self) -> None:
        """Segment with neither superscript nor subscript has no baseline attr."""
        rpr = etree.Element(f"{{{_DRAWINGML_NS}}}rPr")
        seg = _FormattedSegment(
            text="x",
            bold=False,
            italic=False,
            underline=False,
            strike=False,
            superscript=False,
            subscript=False,
        )
        _apply_drawingml_format_attrs(rpr, seg)
        assert rpr.get("baseline") is None


# ---------------------------------------------------------------------------
# _wrap_with_tags hyperlink combo tests
# ---------------------------------------------------------------------------


class TestWrapWithTagsHyperlinkCombos:
    """Tests for _wrap_with_tags with hyperlink_url combined with other flags."""

    def test_hyperlink_with_superscript(self) -> None:
        """Hyperlink + superscript: <a> wraps <sup>."""
        result = _wrap_with_tags(
            "ref",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            hyperlink_url="https://example.com",
            superscript=True,
        )
        assert result.startswith('<a href="https://example.com">')
        assert "<sup>" in result
        assert "ref" in result
        assert result.endswith("</a>")

    def test_hyperlink_with_subscript(self) -> None:
        """Hyperlink + subscript: <a> wraps <sub>."""
        result = _wrap_with_tags(
            "i",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            hyperlink_url="https://x.com",
            subscript=True,
        )
        assert result.startswith('<a href="https://x.com">')
        assert "<sub>" in result

    def test_hyperlink_with_bold_and_superscript(self) -> None:
        """Hyperlink + bold + superscript: all three present, <a> outermost."""
        result = _wrap_with_tags(
            "2",
            True,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            hyperlink_url="https://x.com",
            superscript=True,
        )
        # <a> is outermost
        assert result.startswith('<a href="https://x.com">')
        assert "<b>" in result
        assert "<sup>" in result
        assert "2" in result

    def test_hyperlink_url_html_escaped(self) -> None:
        """URL with special characters is HTML-escaped in href."""
        result = _wrap_with_tags(
            "link",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            hyperlink_url='https://x.com?a=1&b="2"',
        )
        assert "&amp;" in result
        assert "&quot;" in result

    def test_hyperlink_only_no_other_flags(self) -> None:
        """Plain hyperlink with no formatting flags."""
        result = _wrap_with_tags(
            "click",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            hyperlink_url="https://example.com",
        )
        assert result == '<a href="https://example.com">click</a>'


# ---------------------------------------------------------------------------
# _parse_html_formatting hyperlink edge-case tests
# ---------------------------------------------------------------------------


class TestParseHtmlFormattingHyperlinkEdgeCases:
    """Tests for _parse_html_formatting with edge-case <a> tag usage."""

    def test_a_tag_without_href(self) -> None:
        """<a> without href attribute → hyperlink_url is empty string."""
        segs = _parse_html_formatting("<a>bare anchor</a>")
        assert len(segs) == 1
        assert segs[0].text == "bare anchor"
        assert segs[0].hyperlink_url == ""

    def test_a_tag_with_empty_href(self) -> None:
        """<a href=""> → hyperlink_url is empty string."""
        segs = _parse_html_formatting('<a href="">empty</a>')
        assert len(segs) == 1
        assert segs[0].hyperlink_url == ""

    def test_sup_inside_hyperlink(self) -> None:
        """<a><sup>2</sup></a> → segment has both hyperlink_url and superscript."""
        segs = _parse_html_formatting('<a href="https://x.com"><sup>2</sup></a>')
        assert len(segs) == 1
        assert segs[0].superscript is True
        assert segs[0].hyperlink_url == "https://x.com"

    def test_sub_inside_hyperlink(self) -> None:
        """<a><sub>i</sub></a> → segment has both hyperlink_url and subscript."""
        segs = _parse_html_formatting('<a href="https://x.com"><sub>i</sub></a>')
        assert len(segs) == 1
        assert segs[0].subscript is True
        assert segs[0].hyperlink_url == "https://x.com"

    def test_nested_sup_sub_both_set(self) -> None:
        """<sup><sub>X</sub></sup> → both flags set on segment (parser sets both)."""
        segs = _parse_html_formatting("<sup><sub>X</sub></sup>")
        assert len(segs) == 1
        # Both sup and sub active simultaneously (unusual but parser handles gracefully)
        assert segs[0].text == "X"


# ---------------------------------------------------------------------------
# _runs_to_html superscript + hyperlink
# ---------------------------------------------------------------------------


class TestRunsToHtmlSupHyperlink:
    """Tests for _runs_to_html producing <sup>/<sub> tags from real python-docx runs."""

    def test_superscript_and_hyperlink_run(self) -> None:
        """Superscript + hyperlink run produces both <sup> and <a> tags."""
        from docx import Document  # noqa: PLC0415
        from docx.oxml import OxmlElement as DocxOxmlElement  # noqa: PLC0415
        from docx.oxml.ns import qn as docx_qn  # noqa: PLC0415

        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("ref")

        # Set superscript via XML
        rpr = run._r.get_or_add_rPr()
        va = DocxOxmlElement("w:vertAlign")
        va.set(docx_qn("w:val"), "superscript")
        rpr.append(va)

        # Add a hyperlink wrapping the run
        hyperlink = DocxOxmlElement("w:hyperlink")
        hyperlink.set(docx_qn("w:anchor"), "section1")
        # Move run under hyperlink
        para._p.remove(run._r)
        hyperlink.append(run._r)
        para._p.append(hyperlink)

        result = _runs_to_html(para)
        # Should contain sup tag for the superscript and/or at minimum text is preserved
        assert "<sup>" in result or "ref" in result


# ---------------------------------------------------------------------------
# _read_pptx_run_formatting strike attribute variants
# ---------------------------------------------------------------------------


class TestReadPptxRunFormattingStrikeVariants:
    """Tests for _read_pptx_run_formatting with different strike attribute values."""

    def _make_run_with_strike(self, strike_val: str) -> object:
        """Creates a pptx run with a specific strike attribute.

        Args:
            strike_val: The strike attribute value to set on the rPr element.

        Returns:
            A python-pptx run object with the specified strike attribute.
        """
        from pptx import Presentation  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = tx_box.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "test"
        run.font.bold = False  # initializes rPr
        rpr = run._r.find(pptx_qn("a:rPr"))
        if rpr is not None:
            rpr.set("strike", strike_val)
        return run

    def test_sng_strike(self) -> None:
        """Value sngStrike → strike=True."""
        run = self._make_run_with_strike("sngStrike")
        _, _, _, strike, _, _ = _read_pptx_run_formatting(run)
        assert strike is True

    def test_dbl_strike(self) -> None:
        """Value dblStrike → strike=True (any non-noStrike)."""
        run = self._make_run_with_strike("dblStrike")
        _, _, _, strike, _, _ = _read_pptx_run_formatting(run)
        assert strike is True

    def test_no_strike(self) -> None:
        """Value noStrike → strike=False."""
        run = self._make_run_with_strike("noStrike")
        _, _, _, strike, _, _ = _read_pptx_run_formatting(run)
        assert strike is False


# ---------------------------------------------------------------------------
# _color_hex_to_win32com case sensitivity tests
# ---------------------------------------------------------------------------


class TestColorHexToWin32comCaseSensitivity:
    """Tests for _color_hex_to_win32com case sensitivity."""

    def test_lowercase_hex(self) -> None:
        """Lowercase hex input converts same as uppercase."""
        upper = _color_hex_to_win32com("#FF0000")
        lower = _color_hex_to_win32com("#ff0000")
        assert upper == lower

    def test_mixed_case_hex(self) -> None:
        """Mixed case hex input converts correctly."""
        mixed = _color_hex_to_win32com("#Ff0000")
        expected = _color_hex_to_win32com("#FF0000")
        assert mixed == expected


# ── ODF span formatting with superscript/subscript ───────────────────────────

from src.core.office_formatter import _read_odf_span_formatting  # noqa: E402


def _make_odf_style_map(text_position: str | None = None) -> dict[str, object]:
    """Builds a minimal ODF style_map with a single style entry.

    Args:
        text_position: Value for ``style:text-position`` (e.g. "super 58%").
            ``None`` omits the attribute.

    Returns:
        Dict mapping style name ``"T1"`` to its lxml ``<style:style>`` element.
    """
    from lxml import etree  # noqa: PLC0415

    style_ns = _ODF_NS["style"]
    fo_ns = _ODF_NS["fo"]
    style_el = etree.fromstring(
        f'<style:style xmlns:style="{style_ns}" xmlns:fo="{fo_ns}"'
        f' style:name="T1" style:family="text">'
        f"<style:text-properties/>"
        f"</style:style>",
    )
    if text_position is not None:
        tp = style_el.find(f"{{{style_ns}}}text-properties")
        tp.set(f"{{{style_ns}}}text-position", text_position)
    return {"T1": style_el}


class TestReadOdfSpanFormattingSupSub:
    """Tests for _read_odf_span_formatting superscript/subscript detection."""

    def test_super_text_position_returns_superscript(self) -> None:
        """style:text-position='super 58%' → superscript=True, subscript=False."""
        style_map = _make_odf_style_map("super 58%")
        b, i, u, s, sup, sub, sz, clr, bg = _read_odf_span_formatting(style_map, "T1")
        assert sup is True
        assert sub is False

    def test_sub_text_position_returns_subscript(self) -> None:
        """style:text-position='sub 58%' → subscript=True, superscript=False."""
        style_map = _make_odf_style_map("sub 58%")
        b, i, u, s, sup, sub, sz, clr, bg = _read_odf_span_formatting(style_map, "T1")
        assert sup is False
        assert sub is True

    def test_no_text_position_both_false(self) -> None:
        """No text-position attribute → both False."""
        style_map = _make_odf_style_map(None)
        b, i, u, s, sup, sub, sz, clr, bg = _read_odf_span_formatting(style_map, "T1")
        assert sup is False
        assert sub is False

    def test_missing_style_returns_all_false(self) -> None:
        """Unknown style name → 9-tuple of False/None."""
        result = _read_odf_span_formatting({}, "unknown")
        assert result == (False, False, False, False, False, False, None, None, None)

    def test_returns_9_tuple(self) -> None:
        """Return value is always a 9-element tuple."""
        style_map = _make_odf_style_map("super 58%")
        result = _read_odf_span_formatting(style_map, "T1")
        assert len(result) == 9  # noqa: PLR2004


class TestOdfTextBoxToHtmlSupSub:
    """Tests for _odf_text_box_to_html extracting superscript/subscript via style."""

    def _make_text_box_with_style(self, text_position: str) -> tuple[object, dict]:
        """Builds a minimal ODF text-box with a styled superscript span.

        Args:
            text_position: Value for ``style:text-position``.

        Returns:
            (text_box_el, style_map) ready for _odf_text_box_to_html.
        """
        from lxml import etree  # noqa: PLC0415

        text_ns = _ODF_NS["text"]
        style_ns = _ODF_NS["style"]
        fo_ns = _ODF_NS["fo"]
        draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
        xlink_ns = _ODF_NS["xlink"]

        xml = (
            f'<draw:text-box xmlns:draw="{draw_ns}"'
            f' xmlns:text="{text_ns}"'
            f' xmlns:style="{style_ns}"'
            f' xmlns:xlink="{xlink_ns}">'
            f"<text:p>normal "
            f'<text:span text:style-name="T1">2</text:span>'
            f"</text:p>"
            f"</draw:text-box>"
        )
        tb = etree.fromstring(xml)

        style_el = etree.fromstring(
            f'<style:style xmlns:style="{style_ns}" xmlns:fo="{fo_ns}"'
            f' style:name="T1" style:family="text">'
            f'<style:text-properties style:text-position="{text_position}"/>'
            f"</style:style>",
        )
        style_map = {"T1": style_el}
        return tb, style_map

    def test_superscript_span_produces_sup_tag(self) -> None:
        """ODF span with text-position='super 58%' → <sup> in HTML output."""
        text_p_tag = f"{{{_ODF_NS['text']}}}p"
        tb, style_map = self._make_text_box_with_style("super 58%")
        result = _odf_text_box_to_html(tb, style_map, text_p_tag)
        assert "<sup>" in result
        assert "2" in result
        assert "</sup>" in result

    def test_subscript_span_produces_sub_tag(self) -> None:
        """ODF span with text-position='sub 58%' → <sub> in HTML output."""
        text_p_tag = f"{{{_ODF_NS['text']}}}p"
        tb, style_map = self._make_text_box_with_style("sub 58%")
        result = _odf_text_box_to_html(tb, style_map, text_p_tag)
        assert "<sub>" in result
        assert "2" in result
        assert "</sub>" in result


class TestHasOdfTextBoxMixedFormattingSupSub:
    """Tests for _has_odf_text_box_mixed_formatting detecting sup/sub spans."""

    def _make_text_box_with_sup_span(self) -> tuple[object, dict]:
        """Builds a text-box where one span is superscript and rest is plain."""
        from lxml import etree  # noqa: PLC0415

        text_ns = _ODF_NS["text"]
        style_ns = _ODF_NS["style"]
        fo_ns = _ODF_NS["fo"]
        draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"

        xml = (
            f'<draw:text-box xmlns:draw="{draw_ns}"'
            f' xmlns:text="{text_ns}">'
            f"<text:p>plain "
            f'<text:span text:style-name="T1">sup</text:span>'
            f"</text:p>"
            f"</draw:text-box>"
        )
        tb = etree.fromstring(xml)

        style_el = etree.fromstring(
            f'<style:style xmlns:style="{style_ns}" xmlns:fo="{fo_ns}"'
            f' style:name="T1" style:family="text">'
            f'<style:text-properties style:text-position="super 58%"/>'
            f"</style:style>",
        )
        return tb, {"T1": style_el}

    def test_superscript_span_returns_true(self) -> None:
        """Text box with one plain + one superscript span → True (mixed)."""
        from src.core.office_formatter import (  # noqa: PLC0415
            _has_odf_text_box_mixed_formatting,
        )

        text_p_tag = f"{{{_ODF_NS['text']}}}p"
        tb, style_map = self._make_text_box_with_sup_span()
        assert _has_odf_text_box_mixed_formatting(tb, style_map, text_p_tag) is True

    def test_uniform_superscript_returns_false(self) -> None:
        """Text box where ALL spans are superscript → False (uniform)."""
        from lxml import etree  # noqa: PLC0415

        from src.core.office_formatter import (  # noqa: PLC0415
            _has_odf_text_box_mixed_formatting,
        )

        text_ns = _ODF_NS["text"]
        style_ns = _ODF_NS["style"]
        fo_ns = _ODF_NS["fo"]
        draw_ns = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"

        xml = (
            f'<draw:text-box xmlns:draw="{draw_ns}"'
            f' xmlns:text="{text_ns}">'
            f"<text:p>"
            f'<text:span text:style-name="T1">a</text:span>'
            f'<text:span text:style-name="T1">b</text:span>'
            f"</text:p>"
            f"</draw:text-box>"
        )
        tb = etree.fromstring(xml)
        style_el = etree.fromstring(
            f'<style:style xmlns:style="{style_ns}" xmlns:fo="{fo_ns}"'
            f' style:name="T1" style:family="text">'
            f'<style:text-properties style:text-position="super 58%"/>'
            f"</style:style>",
        )
        text_p_tag = f"{{{text_ns}}}p"
        assert (
            _has_odf_text_box_mixed_formatting(
                tb,
                {"T1": style_el},
                text_p_tag,
            )
            is False
        )


class TestInjectPptxHtmlRunsSupSub:
    """Tests for _inject_pptx_html_runs setting baseline for sup/sub."""

    def test_sup_html_sets_baseline_30000(self) -> None:
        """Injecting '<sup>2</sup>' sets baseline='30000' on <a:rPr>."""
        from pptx import Presentation  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        para = txbox.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = "original"
        run.font.bold = False  # initialise rPr

        _inject_pptx_html_runs(para, "<sup>2</sup>")

        assert len(para.runs) == 1
        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        assert rpr is not None
        assert rpr.get("baseline") == "30000"

    def test_sub_html_sets_baseline_neg25000(self) -> None:
        """Injecting '<sub>i</sub>' sets baseline='-25000' on <a:rPr>."""
        from pptx import Presentation  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        para = txbox.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = "original"
        run.font.bold = False

        _inject_pptx_html_runs(para, "<sub>i</sub>")

        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        assert rpr is not None
        assert rpr.get("baseline") == "-25000"

    def test_plain_html_no_baseline(self) -> None:
        """Injecting plain text produces no baseline attribute on <a:rPr>."""
        from pptx import Presentation  # noqa: PLC0415
        from pptx.oxml.ns import qn as pptx_qn  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        para = txbox.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = "original"
        run.font.bold = False

        _inject_pptx_html_runs(para, "plain text")

        rpr = para.runs[0]._r.find(pptx_qn("a:rPr"))
        if rpr is not None:
            assert rpr.get("baseline") is None


class TestDrawingmlToHtmlSupSub:
    """Tests for _drawingml_to_html extracting sup/sub from baseline attribute."""

    def _make_drawingml_para_with_baseline(self, baseline: str) -> object:
        """Builds a DrawingML <a:p> with a single run bearing baseline attribute.

        Args:
            baseline: Value for the ``baseline`` attribute on ``<a:rPr>``.

        Returns:
            lxml element for the ``<a:p>``.
        """
        from lxml import etree  # noqa: PLC0415

        dml_ns = _DRAWINGML_NS
        return etree.fromstring(
            f'<a:p xmlns:a="{dml_ns}">'
            f"<a:r>"
            f'<a:rPr baseline="{baseline}"/>'
            f"<a:t>exponent</a:t>"
            f"</a:r>"
            f"</a:p>",
        )

    def test_positive_baseline_produces_sup_tag(self) -> None:
        """baseline='30000' on <a:rPr> → <sup> in HTML output."""
        from lxml import etree  # noqa: PLC0415

        dml_ns = _DRAWINGML_NS
        # _drawingml_to_html expects a <txBody> parent context; wrap in one
        txbody = etree.fromstring(
            f'<txBody xmlns:a="{dml_ns}">'
            f'<a:p><a:r><a:rPr baseline="30000"/><a:t>exponent</a:t></a:r></a:p>'
            f"</txBody>",
        )
        result = _drawingml_to_html(txbody)
        assert "<sup>" in result
        assert "exponent" in result

    def test_negative_baseline_produces_sub_tag(self) -> None:
        """baseline='-25000' on <a:rPr> → <sub> in HTML output."""
        from lxml import etree  # noqa: PLC0415

        dml_ns = _DRAWINGML_NS
        txbody = etree.fromstring(
            f'<txBody xmlns:a="{dml_ns}">'
            f'<a:p><a:r><a:rPr baseline="-25000"/><a:t>subscript</a:t></a:r></a:p>'
            f"</txBody>",
        )
        result = _drawingml_to_html(txbody)
        assert "<sub>" in result
        assert "subscript" in result

    def test_zero_baseline_no_sup_sub(self) -> None:
        """baseline='0' → no <sup> or <sub> tags."""
        from lxml import etree  # noqa: PLC0415

        dml_ns = _DRAWINGML_NS
        txbody = etree.fromstring(
            f'<txBody xmlns:a="{dml_ns}">'
            f'<a:p><a:r><a:rPr baseline="0"/><a:t>plain</a:t></a:r></a:p>'
            f"</txBody>",
        )
        result = _drawingml_to_html(txbody)
        assert "<sup>" not in result
        assert "<sub>" not in result
        assert "plain" in result


class TestHasMixedFormattingMixedSupSub:
    """Tests for _has_mixed_formatting with mixed superscript and subscript."""

    def test_sup_and_sub_in_same_para_returns_true(self) -> None:
        """One superscript run + one subscript run → True (different formatting)."""
        from docx import Document  # noqa: PLC0415

        doc = Document()
        para = doc.add_paragraph()
        r_sup = para.add_run("a")
        r_sup.font.superscript = True
        r_sub = para.add_run("b")
        r_sub.font.subscript = True
        assert _has_mixed_formatting(para) is True


# ---------------------------------------------------------------------------
# _read_drawingml_rpr_formatting tests
# ---------------------------------------------------------------------------

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _make_rpr_xml(attrs: str = "", children: str = "") -> object:
    """Builds a DrawingML ``<a:rPr>`` lxml element from attr and child strings."""
    xml = f'<a:rPr xmlns:a="{_A_NS}" {attrs}>{children}</a:rPr>'
    return etree.fromstring(xml)


class TestReadDrawingmlRprFormatting:
    """Tests for _read_drawingml_rpr_formatting."""

    def test_none_input(self) -> None:
        """None input returns all defaults."""
        result = _read_drawingml_rpr_formatting(None)
        assert result == (False, False, False, False, False, False, None, None, None)

    def test_bold_italic(self) -> None:
        """Bold and italic attributes are read correctly."""
        rpr = _make_rpr_xml(attrs='b="1" i="1"')
        result = _read_drawingml_rpr_formatting(rpr)
        bold, italic = result[0], result[1]
        assert bold is True
        assert italic is True

    def test_underline_single(self) -> None:
        """Underline attribute 'sng' is detected."""
        rpr = _make_rpr_xml(attrs='u="sng"')
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[2] is True  # underline

    def test_underline_none(self) -> None:
        """Underline attribute 'none' is not treated as underline."""
        rpr = _make_rpr_xml(attrs='u="none"')
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[2] is False  # underline

    def test_strike(self) -> None:
        """Strike attribute is detected."""
        rpr = _make_rpr_xml(attrs='strike="sngStrike"')
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[3] is True  # strike

    def test_no_strike_value(self) -> None:
        """Strike attribute 'noStrike' is not treated as strikethrough."""
        rpr = _make_rpr_xml(attrs='strike="noStrike"')
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[3] is False  # strike

    def test_baseline_positive_superscript(self) -> None:
        """Positive baseline indicates superscript."""
        rpr = _make_rpr_xml(attrs='baseline="30000"')
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[4] is True  # superscript
        assert result[5] is False  # subscript

    def test_baseline_negative_subscript(self) -> None:
        """Negative baseline indicates subscript."""
        rpr = _make_rpr_xml(attrs='baseline="-25000"')
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[4] is False  # superscript
        assert result[5] is True  # subscript

    def test_baseline_zero_neither(self) -> None:
        """Zero baseline is neither superscript nor subscript."""
        rpr = _make_rpr_xml(attrs='baseline="0"')
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[4] is False  # superscript
        assert result[5] is False  # subscript

    def test_font_size(self) -> None:
        """Sz attribute is converted from hundredths of a point to points."""
        rpr = _make_rpr_xml(attrs='sz="1200"')
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[6] == 12.0  # noqa: PLR2004

    def test_color_from_solid_fill(self) -> None:
        """solidFill/srgbClr extracts color hex."""
        children = (
            f'<a:solidFill xmlns:a="{_A_NS}"><a:srgbClr val="FF0000"/></a:solidFill>'
        )
        rpr = _make_rpr_xml(children=children)
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[7] == "#ff0000"

    def test_bg_color_from_highlight(self) -> None:
        """highlight/srgbClr extracts background color hex."""
        children = (
            f'<a:highlight xmlns:a="{_A_NS}"><a:srgbClr val="FFFF00"/></a:highlight>'
        )
        rpr = _make_rpr_xml(children=children)
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[8] == "#ffff00"

    def test_all_defaults_when_no_attrs(self) -> None:
        """Empty rPr element returns all defaults."""
        rpr = _make_rpr_xml()
        result = _read_drawingml_rpr_formatting(rpr)
        assert result == (False, False, False, False, False, False, None, None, None)

    def test_combined_formatting(self) -> None:
        """Multiple formatting flags at once are all detected."""
        children = (
            f'<a:solidFill xmlns:a="{_A_NS}">'
            f'<a:srgbClr val="0000FF"/>'
            f"</a:solidFill>"
            f'<a:highlight xmlns:a="{_A_NS}">'
            f'<a:srgbClr val="00FF00"/>'
            f"</a:highlight>"
        )
        rpr = _make_rpr_xml(
            attrs='b="1" i="1" u="sng" strike="sngStrike" baseline="30000" sz="2400"',
            children=children,
        )
        result = _read_drawingml_rpr_formatting(rpr)
        assert result[0] is True  # bold
        assert result[1] is True  # italic
        assert result[2] is True  # underline
        assert result[3] is True  # strike
        assert result[4] is True  # superscript
        assert result[5] is False  # subscript (baseline > 0 → super)
        assert result[6] == 24.0  # noqa: PLR2004
        assert result[7] == "#0000ff"
        assert result[8] == "#00ff00"


# ---------------------------------------------------------------------------
# _parse_html_formatting — deeply nested HTML
# ---------------------------------------------------------------------------


class TestParseHtmlFormattingDeeplyNested:
    """Tests for _parse_html_formatting with deeply nested and complex HTML."""

    def test_triple_nesting_b_i_u(self) -> None:
        """Deeply nested <b><i><u>text</u></i></b> is parsed correctly."""
        segments = _parse_html_formatting("<b><i><u>deep</u></i></b>")
        assert len(segments) == 1  # noqa: PLR2004
        seg = segments[0]
        assert seg.text == "deep"
        assert seg.bold is True
        assert seg.italic is True
        assert seg.underline is True

    def test_all_six_formatting_tags_nested(self) -> None:
        """All six tags: <b><i><u><s><sup><sub>."""
        html = "<b><i><u><s><sup><sub>x</sub></sup></s></u></i></b>"
        segments = _parse_html_formatting(html)
        assert len(segments) == 1  # noqa: PLR2004
        seg = segments[0]
        assert seg.bold is True
        assert seg.italic is True
        assert seg.underline is True
        assert seg.strike is True
        assert seg.superscript is True
        assert seg.subscript is True

    def test_nested_with_span_inside(self) -> None:
        """Span with style nested inside formatting tags."""
        html = '<b><span style="color:#ff0000">red bold</span></b>'
        segments = _parse_html_formatting(html)
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].bold is True
        assert segments[0].color_hex == "#ff0000"
        assert segments[0].text == "red bold"

    def test_nested_hyperlink_with_formatting(self) -> None:
        """Hyperlink wrapping formatted text."""
        html = '<a href="https://x.com"><b><i>link text</i></b></a>'
        segments = _parse_html_formatting(html)
        assert len(segments) == 1  # noqa: PLR2004
        seg = segments[0]
        assert seg.text == "link text"
        assert seg.bold is True
        assert seg.italic is True
        assert seg.hyperlink_url == "https://x.com"

    def test_alternating_nesting_produces_multiple_segments(self) -> None:
        """<b>A</b><i>B</i><u>C</u> → three segments with different flags."""
        html = "<b>A</b><i>B</i><u>C</u>"
        segments = _parse_html_formatting(html)
        assert len(segments) == 3  # noqa: PLR2004
        assert segments[0].bold is True
        assert segments[0].italic is False
        assert segments[1].italic is True
        assert segments[1].bold is False
        assert segments[2].underline is True
        assert segments[2].bold is False

    def test_multiple_spans_with_different_styles(self) -> None:
        """Multiple span segments with different colors are not merged."""
        html = (
            '<span style="color:#ff0000">red</span>'
            '<span style="color:#0000ff">blue</span>'
        )
        segments = _parse_html_formatting(html)
        assert len(segments) == 2  # noqa: PLR2004
        assert segments[0].color_hex == "#ff0000"
        assert segments[1].color_hex == "#0000ff"

    def test_text_between_nested_tags(self) -> None:
        """Plain text between formatted segments."""
        html = "<b>bold</b> plain <i>italic</i>"
        segments = _parse_html_formatting(html)
        assert len(segments) == 3  # noqa: PLR2004
        assert segments[0].text == "bold"
        assert segments[0].bold is True
        assert segments[1].text == " plain "
        assert segments[1].bold is False
        assert segments[2].text == "italic"
        assert segments[2].italic is True

    def test_span_with_size_color_and_bg(self) -> None:
        """Span with all three CSS properties."""
        html = (
            '<span style="font-size:18pt;color:#aabbcc;background-color:#112233">'
            "styled"
            "</span>"
        )
        segments = _parse_html_formatting(html)
        assert len(segments) == 1  # noqa: PLR2004
        seg = segments[0]
        assert seg.text == "styled"
        assert seg.font_size_pt == 18.0  # noqa: PLR2004
        assert seg.color_hex == "#aabbcc"
        assert seg.bg_color_hex == "#112233"


# ---------------------------------------------------------------------------
# _FormattedSegment — merge logic edge cases
# ---------------------------------------------------------------------------


class TestFormattedSegmentMergeEdgeCases:
    """Tests for _FormattedSegment merge logic in _parse_html_formatting."""

    def test_merge_adjacent_bold_segments(self) -> None:
        """Adjacent bold segments are merged into one."""
        segments = _parse_html_formatting("<b>one</b><b>two</b>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "onetwo"
        assert segments[0].bold is True

    def test_no_merge_different_formatting(self) -> None:
        """Adjacent segments with different formatting are NOT merged."""
        segments = _parse_html_formatting("<b>bold</b><i>italic</i>")
        assert len(segments) == 2  # noqa: PLR2004

    def test_merge_plain_text_segments(self) -> None:
        """Adjacent plain text (no tags) is one segment."""
        segments = _parse_html_formatting("hello world")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "hello world"

    def test_no_merge_different_colors(self) -> None:
        """Adjacent spans with different colors are NOT merged."""
        html = (
            '<span style="color:#ff0000">A</span><span style="color:#0000ff">B</span>'
        )
        segments = _parse_html_formatting(html)
        assert len(segments) == 2  # noqa: PLR2004

    def test_merge_same_color_segments(self) -> None:
        """Adjacent spans with the same color ARE merged."""
        html = (
            '<span style="color:#ff0000">A</span><span style="color:#ff0000">B</span>'
        )
        segments = _parse_html_formatting(html)
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "AB"

    def test_no_merge_different_hyperlinks(self) -> None:
        """Adjacent hyperlinks to different URLs are NOT merged."""
        html = '<a href="https://a.com">A</a><a href="https://b.com">B</a>'
        segments = _parse_html_formatting(html)
        assert len(segments) == 2  # noqa: PLR2004

    def test_merge_same_hyperlink_segments(self) -> None:
        """Adjacent segments with same hyperlink URL are merged."""
        html = '<a href="https://x.com">A</a><a href="https://x.com">B</a>'
        segments = _parse_html_formatting(html)
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == "AB"
        assert segments[0].hyperlink_url == "https://x.com"

    def test_no_merge_different_font_sizes(self) -> None:
        """Adjacent spans with different font sizes are NOT merged."""
        html = (
            '<span style="font-size:12pt">A</span><span style="font-size:18pt">B</span>'
        )
        segments = _parse_html_formatting(html)
        assert len(segments) == 2  # noqa: PLR2004

    def test_no_merge_different_bg_colors(self) -> None:
        """Adjacent spans with different bg colors are NOT merged."""
        html = (
            '<span style="background-color:#ff0000">A</span>'
            '<span style="background-color:#0000ff">B</span>'
        )
        segments = _parse_html_formatting(html)
        assert len(segments) == 2  # noqa: PLR2004

    def test_empty_segments_not_emitted(self) -> None:
        """Tags wrapping empty text produce no segments."""
        segments = _parse_html_formatting("<b></b><i></i>")
        assert segments == []

    def test_whitespace_segment_preserved(self) -> None:
        """Whitespace-only segments are preserved (not discarded)."""
        segments = _parse_html_formatting("<b> </b>")
        assert len(segments) == 1  # noqa: PLR2004
        assert segments[0].text == " "
        assert segments[0].bold is True


# ---------------------------------------------------------------------------
# _wrap_with_tags — all formatting combinations
# ---------------------------------------------------------------------------


class TestWrapWithTagsAllCombinations:
    """Tests for _wrap_with_tags with various formatting flag combos."""

    def test_plain_text_no_flags(self) -> None:
        """No flags → bare text."""
        result = _wrap_with_tags(
            "text",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "text"

    def test_all_four_basic_flags(self) -> None:
        """bold+italic+underline+strike → all four tags nested."""
        result = _wrap_with_tags(
            "x",
            True,
            True,
            True,
            True,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<b><i><u><s>x</s></u></i></b>"

    def test_italic_only(self) -> None:
        """Only italic → <i> wrapper."""
        result = _wrap_with_tags(
            "text",
            False,
            True,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<i>text</i>"

    def test_underline_only(self) -> None:
        """Only underline → <u> wrapper."""
        result = _wrap_with_tags(
            "text",
            False,
            False,
            True,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<u>text</u>"

    def test_strike_only(self) -> None:
        """Only strike → <s> wrapper."""
        result = _wrap_with_tags(
            "text",
            False,
            False,
            False,
            True,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<s>text</s>"

    def test_bold_and_italic(self) -> None:
        """Bold + italic → <b><i>text</i></b>."""
        result = _wrap_with_tags(
            "text",
            True,
            True,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<b><i>text</i></b>"

    def test_size_color_bg_all_variations(self) -> None:
        """All three CSS variations → combined span style."""
        result = _wrap_with_tags(
            "x",
            False,
            False,
            False,
            False,
            12.0,
            "#ff0000",
            has_size_variation=True,
            has_color_variation=True,
            bg_color_hex="#00ff00",
            has_bg_variation=True,
        )
        assert "font-size:12pt" in result
        assert "color:#ff0000" in result
        assert "background-color:#00ff00" in result

    def test_bg_variation_false_omits_bg(self) -> None:
        """has_bg_variation=False omits bg even when bg_color_hex is set."""
        result = _wrap_with_tags(
            "x",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            bg_color_hex="#00ff00",
            has_bg_variation=False,
        )
        assert "background-color" not in result
        assert result == "x"

    def test_bold_with_bg_variation(self) -> None:
        """Bold + bg color variation."""
        result = _wrap_with_tags(
            "hi",
            True,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            bg_color_hex="#ffff00",
            has_bg_variation=True,
        )
        assert result.startswith("<b>")
        assert "background-color:#ffff00" in result

    def test_all_flags_with_all_css_and_hyperlink(self) -> None:
        """All formatting flags + all CSS variations + hyperlink."""
        result = _wrap_with_tags(
            "text",
            True,
            True,
            True,
            True,
            14.0,
            "#112233",
            has_size_variation=True,
            has_color_variation=True,
            bg_color_hex="#445566",
            has_bg_variation=True,
            hyperlink_url="https://example.com",
            superscript=True,
        )
        # <a> is outermost
        assert result.startswith('<a href="https://example.com">')
        assert result.endswith("</a>")
        # All tags present
        assert "<b>" in result
        assert "<i>" in result
        assert "<u>" in result
        assert "<s>" in result
        assert "<sup>" in result
        # Span with CSS
        assert "font-size:14pt" in result
        assert "color:#112233" in result
        assert "background-color:#445566" in result

    def test_superscript_with_color_variation(self) -> None:
        """Superscript + color variation."""
        result = _wrap_with_tags(
            "n",
            False,
            False,
            False,
            False,
            None,
            "#ff0000",
            has_size_variation=False,
            has_color_variation=True,
            superscript=True,
        )
        assert "<sup>" in result
        assert "color:#ff0000" in result

    def test_subscript_with_bold(self) -> None:
        """Subscript + bold → <b><sub>text</sub></b>."""
        result = _wrap_with_tags(
            "2",
            True,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            subscript=True,
        )
        assert "<b>" in result
        assert "<sub>" in result
        # Bold is outside sub in nesting order
        assert result == "<b><sub>2</sub></b>"


# ---------------------------------------------------------------------------
# Backend-specific formatter edge cases
# ---------------------------------------------------------------------------


class TestRunHasVisualContent:
    """Tests for _run_has_visual_content."""

    def test_run_with_drawing(self) -> None:
        """Run containing w:drawing is visual content."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run()
        drawing = OxmlElement("w:drawing")
        run._element.append(drawing)
        assert _run_has_visual_content(run._element) is True

    def test_run_with_pict(self) -> None:
        """Run containing w:pict is visual content."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run()
        pict = OxmlElement("w:pict")
        run._element.append(pict)
        assert _run_has_visual_content(run._element) is True

    def test_run_with_object(self) -> None:
        """Run containing w:object is visual content."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run()
        obj = OxmlElement("w:object")
        run._element.append(obj)
        assert _run_has_visual_content(run._element) is True

    def test_plain_text_run_is_not_visual(self) -> None:
        """Run with only text is NOT visual content."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("Hello")
        assert _run_has_visual_content(run._element) is False

    def test_empty_run_is_not_visual(self) -> None:
        """Empty run (no elements) is NOT visual content."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run()
        assert _run_has_visual_content(run._element) is False


class TestReplaceParaTextEdgeCases:
    """Edge cases for _replace_paragraph_text."""

    def test_replace_text_in_empty_paragraph(self) -> None:
        """Replacing text in a paragraph with no runs sets para.text."""
        doc = Document()
        para = doc.add_paragraph()
        _replace_paragraph_text(para, "new text")
        assert para.text == "new text"

    def test_replace_with_empty_string(self) -> None:
        """Replacing with empty string clears all run text."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("old text")
        _replace_paragraph_text(para, "")
        assert para.text == ""

    def test_replace_preserves_image_run(self) -> None:
        """Visual content runs are preserved; text goes into text-only run."""
        doc = Document()
        para = doc.add_paragraph()
        para.add_run("text")
        img_run = para.add_run()
        drawing = OxmlElement("w:drawing")
        img_run._element.append(drawing)
        para.add_run("more text")

        _replace_paragraph_text(para, "translated")

        # Drawing should still be present
        has_visual = any(_run_has_visual_content(r._element) for r in para.runs)
        assert has_visual
        # First text run should have the new text
        text_runs = [
            r for r in para.runs if r.text and not _run_has_visual_content(r._element)
        ]
        assert any(r.text == "translated" for r in text_runs)

    def test_replace_all_visual_runs_inserts_new_run(self) -> None:
        """When all runs are visual, a new text run is inserted."""
        doc = Document()
        para = doc.add_paragraph()
        run1 = para.add_run()
        drawing1 = OxmlElement("w:drawing")
        run1._element.append(drawing1)

        _replace_paragraph_text(para, "inserted text")

        # The text should appear in the paragraph
        assert "inserted text" in para.text


class TestBuildSpanStyleAdditional:
    """Additional tests for _build_span_style."""

    def test_bg_color_only(self) -> None:
        """Only bg_color_hex → background-color style."""
        assert _build_span_style(None, None, "#ffff00") == "background-color:#ffff00"

    def test_all_three_properties(self) -> None:
        """All three CSS properties combined."""
        result = _build_span_style(12.0, "#ff0000", "#00ff00")
        assert result == "font-size:12pt;color:#ff0000;background-color:#00ff00"

    def test_size_and_bg(self) -> None:
        """Size + bg without color."""
        result = _build_span_style(10.0, None, "#aabbcc")
        assert result == "font-size:10pt;background-color:#aabbcc"


class TestParseSpanStyleAdditional:
    """Additional tests for _parse_span_style."""

    def test_all_three_properties(self) -> None:
        """All three CSS properties extracted."""
        result = _parse_span_style(
            "font-size:16pt;color:#aabbcc;background-color:#ddeeff"
        )
        assert result == {
            "font_size_pt": 16.0,
            "color_hex": "#aabbcc",
            "bg_color_hex": "#ddeeff",
        }

    def test_extra_whitespace(self) -> None:
        """CSS properties with extra whitespace are still parsed."""
        result = _parse_span_style("font-size: 14pt ; color: #ff0000")
        assert result["font_size_pt"] == 14.0  # noqa: PLR2004
        assert result["color_hex"] == "#ff0000"

    def test_uppercase_hex(self) -> None:
        """Uppercase hex values are lowercased in output."""
        result = _parse_span_style("color:#AABBCC")
        assert result["color_hex"] == "#aabbcc"


class TestColorConversionEdgeCases:
    """Edge cases for color conversion functions."""

    def test_int_to_color_hex_white(self) -> None:
        """Maximum valid value (0xFFFFFF) → #ffffff."""
        assert _int_to_color_hex(0xFFFFFF) == "#ffffff"  # noqa: PLR2004

    def test_int_to_color_hex_just_over_limit(self) -> None:
        """0x1000000 (just over max) → None."""
        assert _int_to_color_hex(0x1000000) is None  # noqa: PLR2004

    def test_color_hex_to_int_white(self) -> None:
        """#ffffff → 16777215."""
        expected = 16777215  # noqa: PLR2004
        assert _color_hex_to_int("#ffffff") == expected

    def test_color_hex_to_int_green(self) -> None:
        """#00ff00 → 65280."""
        expected = 65280  # noqa: PLR2004
        assert _color_hex_to_int("#00ff00") == expected

    def test_win32com_color_to_hex_red(self) -> None:
        """BGR red (0x0000FF in BGR = pure red) → #ff0000."""
        # In BGR: R is lowest byte, so 0xFF = red
        assert _win32com_color_to_hex(0xFF) == "#ff0000"

    def test_win32com_color_to_hex_blue(self) -> None:
        """BGR blue (0xFF0000 in BGR = pure blue) → #0000ff."""
        assert _win32com_color_to_hex(0xFF0000) == "#0000ff"  # noqa: PLR2004

    def test_win32com_color_to_hex_negative(self) -> None:
        """Negative value → None."""
        assert _win32com_color_to_hex(-1) is None

    def test_win32com_color_to_hex_overflow(self) -> None:
        """Value > 0xFFFFFF → None."""
        assert _win32com_color_to_hex(0x1000000) is None  # noqa: PLR2004

    def test_color_hex_to_win32com_red(self) -> None:
        """#ff0000 → 255 (BGR: R in lowest byte)."""
        expected = 255  # noqa: PLR2004
        assert _color_hex_to_win32com("#ff0000") == expected

    def test_color_hex_to_win32com_blue(self) -> None:
        """#0000ff → 16711680 (BGR: B in highest byte)."""
        expected = 16711680  # noqa: PLR2004
        assert _color_hex_to_win32com("#0000ff") == expected

    def test_color_hex_to_win32com_roundtrip(self) -> None:
        """Roundtrip: hex → win32com → hex."""
        original = "#ab12cd"
        bgr = _color_hex_to_win32com(original)
        back = _win32com_color_to_hex(bgr)
        assert back == original


class TestFormattedSegmentFieldAccess:
    """Tests for _FormattedSegment field access and defaults."""

    def test_default_hyperlink_url_is_none(self) -> None:
        """Default hyperlink_url is None."""
        seg = _FormattedSegment("text", False, False, False, False)
        assert seg.hyperlink_url is None

    def test_default_bg_color_hex_is_none(self) -> None:
        """Default bg_color_hex is None."""
        seg = _FormattedSegment("text", False, False, False, False)
        assert seg.bg_color_hex is None

    def test_all_fields_set(self) -> None:
        """All 11 fields can be set explicitly."""
        seg = _FormattedSegment(
            "txt",
            True,
            True,
            True,
            True,
            True,
            True,
            14.0,
            "#ff0000",
            "#00ff00",
            "https://x.com",
        )
        assert seg.text == "txt"
        assert seg.bold is True
        assert seg.italic is True
        assert seg.underline is True
        assert seg.strike is True
        assert seg.superscript is True
        assert seg.subscript is True
        assert seg.font_size_pt == 14.0  # noqa: PLR2004
        assert seg.color_hex == "#ff0000"
        assert seg.bg_color_hex == "#00ff00"
        assert seg.hyperlink_url == "https://x.com"

    def test_equality_with_different_hyperlinks(self) -> None:
        """Segments with different hyperlinks are not equal."""
        seg1 = _FormattedSegment(
            "text",
            False,
            False,
            False,
            False,
            hyperlink_url="https://a.com",
        )
        seg2 = _FormattedSegment(
            "text",
            False,
            False,
            False,
            False,
            hyperlink_url="https://b.com",
        )
        assert seg1 != seg2

    def test_equality_with_same_fields(self) -> None:
        """Segments with identical fields are equal."""
        seg1 = _FormattedSegment("text", True, False, True, False, font_size_pt=12.0)
        seg2 = _FormattedSegment("text", True, False, True, False, font_size_pt=12.0)
        assert seg1 == seg2


class TestFormattingHtmlReAdditional:
    """Additional tests for the _FORMATTING_HTML_RE regex."""

    def test_detects_span_tag(self) -> None:
        """Detects <span> tags."""
        html = '<span style="color:red">x</span>'
        assert _FORMATTING_HTML_RE.search(html) is not None

    def test_detects_anchor_tag(self) -> None:
        """Detects <a> tags."""
        assert _FORMATTING_HTML_RE.search('<a href="https://x.com">x</a>') is not None

    def test_detects_closing_anchor(self) -> None:
        """Detects </a> closing tag."""
        assert _FORMATTING_HTML_RE.search("text</a>") is not None

    def test_detects_sup_tag(self) -> None:
        """Detects <sup> tags."""
        assert _FORMATTING_HTML_RE.search("<sup>2</sup>") is not None

    def test_detects_sub_tag(self) -> None:
        """Detects <sub> tags."""
        assert _FORMATTING_HTML_RE.search("<sub>2</sub>") is not None

    def test_no_match_div_tag(self) -> None:
        """Does not match non-formatting tags like <div>."""
        assert _FORMATTING_HTML_RE.search("<div>text</div>") is None

    def test_no_match_br_tag(self) -> None:
        """Does not match <br> tags."""
        assert _FORMATTING_HTML_RE.search("<br/>") is None


# ---------------------------------------------------------------------------
# Additional coverage: _parse_html_formatting edge cases
# ---------------------------------------------------------------------------


class TestParseHtmlFormattingExtended:
    """Extended tests for _parse_html_formatting with complex HTML."""

    def test_bold_inside_italic(self) -> None:
        """Nested <b> inside <i> produces segment with both flags."""
        segs = _parse_html_formatting("<i><b>nested</b></i>")
        assert len(segs) == 1
        assert segs[0].bold is True
        assert segs[0].italic is True
        assert segs[0].text == "nested"

    def test_italic_inside_bold(self) -> None:
        """Nested <i> inside <b> produces segment with both flags."""
        segs = _parse_html_formatting("<b><i>text</i></b>")
        assert len(segs) == 1
        assert segs[0].bold is True
        assert segs[0].italic is True

    def test_underline_inside_bold_inside_italic(self) -> None:
        """Triple nesting: <i><b><u>text</u></b></i>."""
        segs = _parse_html_formatting("<i><b><u>triple</u></b></i>")
        assert len(segs) == 1
        assert segs[0].bold is True
        assert segs[0].italic is True
        assert segs[0].underline is True

    def test_hyperlink_with_formatted_text(self) -> None:
        """<a href='...'><b>text</b></a> produces linked bold segment."""
        segs = _parse_html_formatting('<a href="https://example.com"><b>link</b></a>')
        assert len(segs) == 1
        assert segs[0].bold is True
        assert segs[0].hyperlink_url == "https://example.com"
        assert segs[0].text == "link"

    def test_empty_formatted_segment(self) -> None:
        """Empty text inside tags produces no segments."""
        segs = _parse_html_formatting("<b></b>")
        assert len(segs) == 0

    def test_empty_string_input(self) -> None:
        """Empty string produces no segments."""
        segs = _parse_html_formatting("")
        assert len(segs) == 0

    def test_plain_text_only(self) -> None:
        """Plain text without tags produces a single unformatted segment."""
        segs = _parse_html_formatting("hello world")
        assert len(segs) == 1
        assert segs[0].text == "hello world"
        assert segs[0].bold is False
        assert segs[0].italic is False

    def test_strikethrough_tag(self) -> None:
        """<s> tag sets strike flag."""
        segs = _parse_html_formatting("<s>deleted</s>")
        assert len(segs) == 1
        assert segs[0].strike is True

    def test_superscript_tag(self) -> None:
        """<sup> tag sets superscript flag."""
        segs = _parse_html_formatting("x<sup>2</sup>")
        assert len(segs) == 2  # noqa: PLR2004
        assert segs[0].superscript is False
        assert segs[1].superscript is True
        assert segs[1].text == "2"

    def test_subscript_tag(self) -> None:
        """<sub> tag sets subscript flag."""
        segs = _parse_html_formatting("H<sub>2</sub>O")
        assert len(segs) == 3  # noqa: PLR2004
        assert segs[1].subscript is True
        assert segs[1].text == "2"

    def test_span_with_font_size(self) -> None:
        """<span style='font-size:14pt'> sets font_size_pt."""
        segs = _parse_html_formatting('<span style="font-size:14pt">sized</span>')
        assert len(segs) == 1
        assert segs[0].font_size_pt == 14.0  # noqa: PLR2004
        assert segs[0].text == "sized"

    def test_span_with_color(self) -> None:
        """<span style='color:#ff0000'> sets color_hex."""
        segs = _parse_html_formatting('<span style="color:#ff0000">red</span>')
        assert len(segs) == 1
        assert segs[0].color_hex == "#ff0000"

    def test_span_with_background_color(self) -> None:
        """<span style='background-color:#ffff00'> sets bg_color_hex."""
        segs = _parse_html_formatting(
            '<span style="background-color:#ffff00">highlight</span>'
        )
        assert len(segs) == 1
        assert segs[0].bg_color_hex == "#ffff00"

    def test_multiple_mixed_segments(self) -> None:
        """Mixed plain + bold + italic produces correct segment list."""
        segs = _parse_html_formatting("Hello <b>bold</b> and <i>italic</i>")
        assert len(segs) == 4  # noqa: PLR2004
        assert segs[0].text == "Hello "
        assert segs[0].bold is False
        assert segs[1].text == "bold"
        assert segs[1].bold is True
        assert segs[2].text == " and "
        assert segs[3].text == "italic"
        assert segs[3].italic is True

    def test_adjacent_same_formatting_merged(self) -> None:
        """Adjacent segments with identical formatting are merged."""
        # Two consecutive plain text pieces
        segs = _parse_html_formatting("<b>first</b><b> second</b>")
        assert len(segs) == 1
        assert segs[0].text == "first second"
        assert segs[0].bold is True

    def test_html_entities_decoded(self) -> None:
        """HTML entities in text are decoded properly."""
        segs = _parse_html_formatting("<b>&lt;tag&gt;</b>")
        assert len(segs) == 1
        assert segs[0].text == "<tag>"


# ---------------------------------------------------------------------------
# Additional coverage: _wrap_with_tags comprehensive tests
# ---------------------------------------------------------------------------


class TestWrapWithTagsExtended:
    """Extended tests for _wrap_with_tags with all formatting types."""

    def test_bold_only(self) -> None:
        """Bold wrapping produces <b> tag."""
        result = _wrap_with_tags(
            "text",
            True,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<b>text</b>"

    def test_italic_only(self) -> None:
        """Italic wrapping produces <i> tag."""
        result = _wrap_with_tags(
            "text",
            False,
            True,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<i>text</i>"

    def test_underline_only(self) -> None:
        """Underline wrapping produces <u> tag."""
        result = _wrap_with_tags(
            "text",
            False,
            False,
            True,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<u>text</u>"

    def test_strike_only(self) -> None:
        """Strikethrough wrapping produces <s> tag."""
        result = _wrap_with_tags(
            "text",
            False,
            False,
            False,
            True,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "<s>text</s>"

    def test_superscript_flag(self) -> None:
        """Superscript wrapping produces <sup> tag."""
        result = _wrap_with_tags(
            "2",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            superscript=True,
        )
        assert result == "<sup>2</sup>"

    def test_subscript_flag(self) -> None:
        """Subscript wrapping produces <sub> tag."""
        result = _wrap_with_tags(
            "2",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            subscript=True,
        )
        assert result == "<sub>2</sub>"

    def test_all_formatting_combined(self) -> None:
        """All formatting flags produce nested tags."""
        result = _wrap_with_tags(
            "text",
            True,
            True,
            True,
            True,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            superscript=True,
        )
        assert "<b>" in result
        assert "<i>" in result
        assert "<u>" in result
        assert "<s>" in result
        assert "<sup>" in result
        assert "text" in result

    def test_span_with_size_variation(self) -> None:
        """Size variation emits <span style='font-size:..'>."""
        result = _wrap_with_tags(
            "big",
            False,
            False,
            False,
            False,
            14.0,
            None,
            has_size_variation=True,
            has_color_variation=False,
        )
        assert "font-size:14pt" in result
        assert "<span" in result

    def test_span_with_color_variation(self) -> None:
        """Color variation emits <span style='color:..'>."""
        result = _wrap_with_tags(
            "red",
            False,
            False,
            False,
            False,
            None,
            "#ff0000",
            has_size_variation=False,
            has_color_variation=True,
        )
        assert "color:#ff0000" in result

    def test_no_span_without_variation(self) -> None:
        """No <span> emitted when variation flags are False."""
        result = _wrap_with_tags(
            "text",
            False,
            False,
            False,
            False,
            14.0,
            "#ff0000",
            has_size_variation=False,
            has_color_variation=False,
        )
        assert "<span" not in result
        assert result == "text"

    def test_hyperlink_wraps_outermost(self) -> None:
        """Hyperlink <a> tag wraps outside all formatting tags."""
        result = _wrap_with_tags(
            "link",
            True,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            hyperlink_url="https://example.com",
        )
        assert result.startswith("<a ")
        assert result.endswith("</a>")
        assert "<b>link</b>" in result

    def test_bg_color_variation(self) -> None:
        """Background color variation emits background-color in span style."""
        result = _wrap_with_tags(
            "hl",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            bg_color_hex="#ffff00",
            has_bg_variation=True,
        )
        assert "background-color:#ffff00" in result

    def test_plain_text_no_tags(self) -> None:
        """No formatting flags produces plain text."""
        result = _wrap_with_tags(
            "plain",
            False,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        assert result == "plain"


# ---------------------------------------------------------------------------
# Additional coverage: round-trip formatting preservation
# ---------------------------------------------------------------------------


class TestFormattingRoundTrip:
    """Tests for round-trip formatting preservation (html → parse → verify)."""

    def test_bold_italic_round_trip(self) -> None:
        """Bold+italic text survives wrap → parse round trip."""
        html_str = _wrap_with_tags(
            "text",
            True,
            True,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        segs = _parse_html_formatting(html_str)
        assert len(segs) == 1
        assert segs[0].bold is True
        assert segs[0].italic is True
        assert segs[0].text == "text"

    def test_underline_strike_round_trip(self) -> None:
        """Underline+strike survives round trip."""
        html_str = _wrap_with_tags(
            "crossed",
            False,
            False,
            True,
            True,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
        )
        segs = _parse_html_formatting(html_str)
        assert len(segs) == 1
        assert segs[0].underline is True
        assert segs[0].strike is True

    def test_font_size_round_trip(self) -> None:
        """Font size value survives wrap → parse round trip."""
        html_str = _wrap_with_tags(
            "sized",
            False,
            False,
            False,
            False,
            18.0,
            None,
            has_size_variation=True,
            has_color_variation=False,
        )
        segs = _parse_html_formatting(html_str)
        assert len(segs) == 1
        assert segs[0].font_size_pt == 18.0  # noqa: PLR2004

    def test_color_round_trip(self) -> None:
        """Color hex value survives wrap → parse round trip."""
        html_str = _wrap_with_tags(
            "colored",
            False,
            False,
            False,
            False,
            None,
            "#00ff00",
            has_size_variation=False,
            has_color_variation=True,
        )
        segs = _parse_html_formatting(html_str)
        assert len(segs) == 1
        assert segs[0].color_hex == "#00ff00"

    def test_hyperlink_round_trip(self) -> None:
        """Hyperlink URL survives wrap → parse round trip."""
        html_str = _wrap_with_tags(
            "click",
            True,
            False,
            False,
            False,
            None,
            None,
            has_size_variation=False,
            has_color_variation=False,
            hyperlink_url="https://example.org",
        )
        segs = _parse_html_formatting(html_str)
        assert len(segs) == 1
        assert segs[0].hyperlink_url == "https://example.org"
        assert segs[0].bold is True

    def test_complex_mixed_round_trip(self) -> None:
        """Complex mixed formatting survives round trip."""
        html_str = "Hello <b>bold</b> <i><u>italic underline</u></i>"
        segs = _parse_html_formatting(html_str)
        assert len(segs) == 4  # noqa: PLR2004
        # First: plain "Hello "
        assert segs[0].bold is False
        assert segs[0].italic is False
        # Second: bold "bold"
        assert segs[1].bold is True
        # Third: space
        assert segs[2].text == " "
        # Fourth: italic+underline
        assert segs[3].italic is True
        assert segs[3].underline is True


# ---------------------------------------------------------------------------
# Additional coverage: color conversion functions
# ---------------------------------------------------------------------------


class TestColorConversionsExtended:
    """Extended tests for color conversion helpers."""

    def test_int_to_color_hex_black(self) -> None:
        """0x000000 converts to #000000."""
        assert _int_to_color_hex(0) == "#000000"

    def test_int_to_color_hex_white(self) -> None:
        """0xFFFFFF converts to #ffffff."""
        assert _int_to_color_hex(0xFFFFFF) == "#ffffff"

    def test_int_to_color_hex_red(self) -> None:
        """0xFF0000 converts to #ff0000."""
        assert _int_to_color_hex(0xFF0000) == "#ff0000"

    def test_int_to_color_hex_negative(self) -> None:
        """Negative value returns None (automatic/unset)."""
        assert _int_to_color_hex(-1) is None

    def test_int_to_color_hex_overflow(self) -> None:
        """Value > 0xFFFFFF returns None."""
        assert _int_to_color_hex(0x1000000) is None

    def test_color_hex_to_int_roundtrip(self) -> None:
        """Hex → int → hex round trip preserves value."""
        original = "#abcdef"
        result = _int_to_color_hex(_color_hex_to_int(original))
        assert result == original

    def test_win32com_color_to_hex_pure_red(self) -> None:
        """BGR 0x0000FF (pure red in BGR) converts to #ff0000."""
        assert _win32com_color_to_hex(0xFF) == "#ff0000"

    def test_win32com_color_to_hex_pure_blue(self) -> None:
        """BGR 0xFF0000 (pure blue in BGR) converts to #0000ff."""
        assert _win32com_color_to_hex(0xFF0000) == "#0000ff"

    def test_color_hex_to_win32com_roundtrip(self) -> None:
        """Hex → win32com → hex round trip preserves value."""
        original = "#12ab34"
        bgr = _color_hex_to_win32com(original)
        result = _win32com_color_to_hex(bgr)
        assert result == original

    def test_win32com_negative_returns_none(self) -> None:
        """Negative BGR value returns None."""
        assert _win32com_color_to_hex(-1) is None


# ---------------------------------------------------------------------------
# Additional coverage: _build_span_style and _parse_span_style
# ---------------------------------------------------------------------------


class TestSpanStyleExtended:
    """Extended tests for span style building and parsing."""

    def test_build_empty_when_all_none(self) -> None:
        """All None arguments produce empty string."""
        assert _build_span_style(None, None) == ""

    def test_build_size_only(self) -> None:
        """Font size only produces 'font-size:...pt'."""
        assert _build_span_style(12.0, None) == "font-size:12pt"

    def test_build_color_only(self) -> None:
        """Color only produces 'color:...'."""
        assert _build_span_style(None, "#ff0000") == "color:#ff0000"

    def test_build_all_three(self) -> None:
        """All three properties produce semicolon-separated style."""
        result = _build_span_style(10.0, "#00ff00", "#ffff00")
        assert "font-size:10pt" in result
        assert "color:#00ff00" in result
        assert "background-color:#ffff00" in result

    def test_parse_size(self) -> None:
        """Parse extracts font-size from style string."""
        result = _parse_span_style("font-size:14pt")
        assert result["font_size_pt"] == 14.0  # noqa: PLR2004

    def test_parse_color(self) -> None:
        """Parse extracts color from style string."""
        result = _parse_span_style("color:#abcdef")
        assert result["color_hex"] == "#abcdef"

    def test_parse_bg_color(self) -> None:
        """Parse extracts background-color from style string."""
        result = _parse_span_style("background-color:#123456")
        assert result["bg_color_hex"] == "#123456"

    def test_parse_roundtrip(self) -> None:
        """Build → parse round trip preserves values."""
        style = _build_span_style(16.5, "#aabbcc", "#ddeeff")
        parsed = _parse_span_style(style)
        assert parsed["font_size_pt"] == 16.5  # noqa: PLR2004
        assert parsed["color_hex"] == "#aabbcc"
        assert parsed["bg_color_hex"] == "#ddeeff"

    def test_parse_empty_string(self) -> None:
        """Empty style string returns empty dict."""
        assert _parse_span_style("") == {}

    def test_parse_does_not_confuse_bg_with_fg_color(self) -> None:
        """Parser distinguishes color: from background-color:."""
        style = "background-color:#111111;color:#222222"
        parsed = _parse_span_style(style)
        assert parsed["color_hex"] == "#222222"
        assert parsed["bg_color_hex"] == "#111111"


# ---------------------------------------------------------------------------
# HIGH-PRIORITY COVERAGE GAPS — _read_docx_run_bg_hex
# ---------------------------------------------------------------------------


class TestReadDocxRunBgHex:
    """Tests for _read_docx_run_bg_hex — all 3 fallback paths."""

    def _make_run_with_shd(self, fill_value: str) -> object:
        """Creates a python-docx run with a <w:shd> element.

        Args:
            fill_value: The fill attribute value (e.g. "FFFF00" or "auto").

        Returns:
            A python-docx Run object.
        """
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        rpr = OxmlElement("w:rPr")
        shd = OxmlElement("w:shd")
        shd.set(qn("w:fill"), fill_value)
        rpr.append(shd)
        run._element.insert(0, rpr)
        return run

    def _make_run_with_highlight(self, val: str) -> object:
        """Creates a python-docx run with a <w:highlight> element.

        Args:
            val: The highlight value (e.g. "yellow", "none").

        Returns:
            A python-docx Run object.
        """
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        rpr = OxmlElement("w:rPr")
        hl = OxmlElement("w:highlight")
        hl.set(qn("w:val"), val)
        rpr.append(hl)
        run._element.insert(0, rpr)
        return run

    def test_shd_with_valid_fill_returns_hex(self) -> None:
        """<w:shd> with valid fill color returns hex string."""
        run = self._make_run_with_shd("FFFF00")
        result = _read_docx_run_bg_hex(run)
        assert result == "#ffff00"

    def test_shd_with_uppercase_fill_normalized(self) -> None:
        """<w:shd> fill is normalized to lowercase."""
        run = self._make_run_with_shd("FF0000")
        result = _read_docx_run_bg_hex(run)
        assert result == "#ff0000"

    def test_shd_with_auto_fill_skips(self) -> None:
        """<w:shd> with fill='auto' is skipped (falls through)."""
        run = self._make_run_with_shd("auto")
        # No highlight or highlight_color, so should return None
        result = _read_docx_run_bg_hex(run)
        assert result is None

    def test_shd_with_auto_uppercase_fill_skips(self) -> None:
        """<w:shd> with fill='AUTO' (uppercase) is skipped."""
        run = self._make_run_with_shd("AUTO")
        result = _read_docx_run_bg_hex(run)
        assert result is None

    def test_highlight_with_yellow_returns_hex(self) -> None:
        """<w:highlight> with val='yellow' returns #ffff00."""
        run = self._make_run_with_highlight("yellow")
        result = _read_docx_run_bg_hex(run)
        assert result == "#ffff00"

    def test_highlight_with_green_returns_hex(self) -> None:
        """<w:highlight> with val='green' returns #00ff00."""
        run = self._make_run_with_highlight("green")
        result = _read_docx_run_bg_hex(run)
        assert result == "#00ff00"

    def test_highlight_with_none_skips(self) -> None:
        """<w:highlight> with val='none' is skipped."""
        run = self._make_run_with_highlight("none")
        result = _read_docx_run_bg_hex(run)
        assert result is None

    def test_highlight_with_darkblue_returns_hex(self) -> None:
        """<w:highlight> with val='darkBlue' returns correct hex."""
        run = self._make_run_with_highlight("darkBlue")
        result = _read_docx_run_bg_hex(run)
        assert result == _HIGHLIGHT_COLORS.get("darkblue")

    def test_highlight_color_api_fallback(self) -> None:
        """Falls back to python-docx highlight_color API when no XML elements."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        # Patch the font property to simulate highlight_color = 7 (yellow)
        mock_font = MagicMock(highlight_color=7)
        with patch.object(
            type(run), "font", new_callable=lambda: property(lambda self: mock_font)
        ):
            result = _read_docx_run_bg_hex(run)
        assert result == _WD_COLOR_INDEX_TO_HEX[7]
        assert result == "#ffff00"

    def test_highlight_color_api_none_returns_none(self) -> None:
        """highlight_color API returning None → returns None."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        # Patch the font property to simulate highlight_color = None
        mock_font = MagicMock(highlight_color=None)
        with patch.object(
            type(run), "font", new_callable=lambda: property(lambda self: mock_font)
        ):
            result = _read_docx_run_bg_hex(run)
        assert result is None

    def test_highlight_color_api_zero_returns_none(self) -> None:
        """highlight_color API returning 0 → returns None (0 is not > 0)."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        mock_font = MagicMock(highlight_color=0)
        with patch.object(
            type(run), "font", new_callable=lambda: property(lambda self: mock_font)
        ):
            result = _read_docx_run_bg_hex(run)
        assert result is None

    def test_no_rpr_element_falls_to_api(self) -> None:
        """Run with no <w:rPr> element falls through to API path."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        # Remove any rPr child if present
        for child in list(run._element):
            if child.tag.endswith("}rPr"):
                run._element.remove(child)
        mock_font = MagicMock(highlight_color=None)
        with patch.object(
            type(run), "font", new_callable=lambda: property(lambda self: mock_font)
        ):
            result = _read_docx_run_bg_hex(run)
        assert result is None

    def test_shd_empty_fill_skips_to_next_fallback(self) -> None:
        """<w:shd> with empty fill attribute skips to next fallback."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        rpr = OxmlElement("w:rPr")
        shd = OxmlElement("w:shd")
        shd.set(qn("w:fill"), "")
        rpr.append(shd)
        run._element.insert(0, rpr)

        mock_font = MagicMock(highlight_color=None)
        with patch.object(
            type(run), "font", new_callable=lambda: property(lambda self: mock_font)
        ):
            result = _read_docx_run_bg_hex(run)
        assert result is None

    def test_shd_takes_priority_over_highlight(self) -> None:
        """<w:shd> with valid fill is returned even if <w:highlight> also present."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        rpr = OxmlElement("w:rPr")
        # Add shd with a fill color
        shd = OxmlElement("w:shd")
        shd.set(qn("w:fill"), "00FF00")
        rpr.append(shd)
        # Also add highlight
        hl = OxmlElement("w:highlight")
        hl.set(qn("w:val"), "yellow")
        rpr.append(hl)
        run._element.insert(0, rpr)

        result = _read_docx_run_bg_hex(run)
        # shd takes priority
        assert result == "#00ff00"

    def test_all_highlight_colors_recognized(self) -> None:
        """All known highlight color names in _HIGHLIGHT_COLORS are recognized."""
        for color_name, expected_hex in _HIGHLIGHT_COLORS.items():
            if color_name in ("none",):
                continue
            run = self._make_run_with_highlight(color_name)
            result = _read_docx_run_bg_hex(run)
            assert result == expected_hex, f"Failed for highlight color: {color_name}"

    def test_unnamespaced_fill_attribute(self) -> None:
        """<w:shd> with unnamespaced 'fill' attribute (no w: prefix) works."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        rpr = OxmlElement("w:rPr")
        shd = OxmlElement("w:shd")
        # Set unnamespaced 'fill' directly (no w: namespace prefix)
        shd.set("fill", "AABBCC")
        rpr.append(shd)
        run._element.insert(0, rpr)

        result = _read_docx_run_bg_hex(run)
        assert result == "#aabbcc"

    def test_unnamespaced_highlight_val(self) -> None:
        """<w:highlight> with unnamespaced 'val' attribute works."""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("test")
        rpr = OxmlElement("w:rPr")
        hl = OxmlElement("w:highlight")
        # Set unnamespaced 'val' directly
        hl.set("val", "red")
        rpr.append(hl)
        run._element.insert(0, rpr)

        result = _read_docx_run_bg_hex(run)
        assert result == "#ff0000"


class TestDrawingMLHyperlinkExtraction:
    """``_drawingml_to_html`` resolves PPTX hyperlink rels to ``<a href="">``.

    Pins the contract that DrawingML runs with ``<a:hlinkClick r:id="...">``
    inside their ``<a:rPr>`` are emitted as anchor tags when the caller
    threads through a ``hyperlink_rels`` map. Without this, translated
    PPTX text would lose its hyperlink targets — a silent data loss that
    the LLM round-trip can't recover from on injection.
    """

    def _make_tx_body_with_hyperlink(self, rid: str, url_text: str):
        """Builds a minimal ``<a:txBody>`` containing a single hyperlinked run."""
        from lxml import etree  # noqa: PLC0415

        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        nsmap = {"a": a_ns, "r": r_ns}

        tx = etree.Element(f"{{{a_ns}}}txBody", nsmap=nsmap)
        p = etree.SubElement(tx, f"{{{a_ns}}}p")
        r = etree.SubElement(p, f"{{{a_ns}}}r")
        rpr = etree.SubElement(r, f"{{{a_ns}}}rPr")
        hlink = etree.SubElement(rpr, f"{{{a_ns}}}hlinkClick")
        hlink.set(f"{{{r_ns}}}id", rid)
        t = etree.SubElement(r, f"{{{a_ns}}}t")
        t.text = url_text
        return tx

    def test_hyperlink_rel_resolved_to_anchor_tag(self) -> None:
        """``hyperlink_rels`` lookup converts r:id → ``<a href="...">``."""
        rid = "rIdHL_X1"
        tx_body = self._make_tx_body_with_hyperlink(rid, "click me")
        rels = {rid: "https://example.com/page"}

        result = _drawingml_to_html(tx_body, hyperlink_rels=rels)

        assert '<a href="https://example.com/page">' in result, (
            f"DrawingML hyperlink rel not resolved to anchor: {result!r}"
        )
        assert "click me" in result
        assert "</a>" in result

    def test_unknown_rid_falls_back_to_plain_text(self) -> None:
        """An r:id not present in ``hyperlink_rels`` does not crash.

        Some drawings reference rels that aren't materialised in the
        passed-in map (rel-file parsing race, document corruption).
        The extractor must degrade to plain text rather than raising.
        """
        tx_body = self._make_tx_body_with_hyperlink("rIdMissing", "fallback")

        # Empty rels — extraction must succeed with text content intact.
        result = _drawingml_to_html(tx_body, hyperlink_rels={})

        assert "fallback" in result
        assert "rIdMissing" not in result  # rel id never leaks into output

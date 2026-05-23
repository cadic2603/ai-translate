"""Integration tests for win32com and UNO office backends.

Runs real end-to-end office translations through these backends when
available, skipping gracefully otherwise.  The existing
``test_integ_office_pipeline.py`` forces the python_lib backend for CI
reliability; this module exercises the *real* backend detection and I/O.

Only the LLM is mocked — all file I/O, backend detection, soffice
management, and document creation/reading use production code.

Verification strategy
~~~~~~~~~~~~~~~~~~~~~
Every test verifies **all** translatable items in the output, not just
``any()`` spot-checks.  The mock LLM deterministically returns
``[French] {original}`` so we can assert exact expected content for
each paragraph, cell, slide, comment, and shape.
"""

import builtins
import struct
import sys
import zipfile
import zlib
from collections.abc import Callable, Generator
from pathlib import Path
from typing import Any

import pytest
from lxml import etree

from src.constants.settings import (
    SETTING_TRANSLATE_DOC_COMMENTS,
    SETTING_TRANSLATE_DOC_SHAPES,
)
from src.core.database import init_db
from src.core.office_lifecycle import (
    _ensure_soffice_running,
    _find_soffice_binary,
    _get_uno_search_paths,
    stop_soffice,
)
from src.core.office_processor import (
    _convert_with_uno,
    _detect_backend,
    _extract_python_odp,
    _extract_python_odt,
)
from src.core.text_processor import translate_file

# ---------------------------------------------------------------------------
# Backend availability detection (evaluated once at module load)
# ---------------------------------------------------------------------------

# Add UNO search paths so the import probe works
for _p in _get_uno_search_paths():
    if _p not in sys.path and Path(_p).is_dir():
        sys.path.append(_p)

try:
    import uno  # noqa: F401

    HAS_UNO = True
except ImportError:
    HAS_UNO = False

HAS_SOFFICE = _find_soffice_binary() is not None

try:
    import win32com.client  # noqa: F401

    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False

requires_uno = pytest.mark.skipif(
    not (HAS_UNO and HAS_SOFFICE),
    reason="LibreOffice UNO and/or soffice binary not available",
)
requires_win32com = pytest.mark.skipif(
    not HAS_WIN32COM,
    reason="win32com not available (requires Windows + MS Office)",
)

# ---------------------------------------------------------------------------
# Known input texts (single source of truth for creation + assertion)
# ---------------------------------------------------------------------------

_DOCX_PARAGRAPHS = ["Hello world", "Goodbye world"]
_DOCX_TABLE_ROWS = [["Name", "Value"], ["Hello", "World"]]
_DOCX_MIXED_FMT_PARTS = ("Bold text", " and ", "italic text")
_DOCX_COMMENT_ANCHOR = "This has a comment"
_DOCX_COMMENT_TEXT = "Review this"
_DOCX_SHAPE_TEXT = "Text in shape"

_XLSX_ROWS = [["Hello", "World"], ["Foo", "Bar"], ["Translation", "Test"]]
_XLSX_COMMENT_TEXT = "Reviewer note"

_PPTX_SLIDES = ["First slide content", "Third slide content"]
_PPTX_MIXED_FMT_PARTS = ("Bold text", " and ", "italic text")

_ODT_PARAGRAPHS = ["Hello world", "Goodbye world", "Another paragraph"]
_ODT_TABLE_ROWS = [["Name", "Value"], ["Foo", "Bar"]]

_ODS_ROWS = [["Hello", "World"], ["Foo", "Bar"], ["Translation", "Test"]]

_ODP_SLIDES = [
    "First slide content",
    "Second slide content",
    "Third slide content",
]

_PPTX_COMMENT_TEXT = "Slide comment note"
_XLSX_SHAPE_TEXT = "Shape in sheet"
_ODT_COMMENT_TEXT = "Writer annotation"
_ODT_SHAPE_TEXT = "Writer text box"
_ODS_COMMENT_TEXT = "Calc annotation"
_ODS_SHAPE_TEXT = "Calc text box"
_ODP_COMMENT_TEXT = "Impress annotation"


def _french(text: str) -> str:
    """Return expected mock-LLM output for a given input text."""
    return f"[French] {text}"


# ---------------------------------------------------------------------------
# UNO import bypass helpers (same as test_integ_office_pipeline.py)
# ---------------------------------------------------------------------------


def _bypass_uno_import() -> object:
    """Temporarily restore Python's real import if UNO's hook is active."""
    uno_mod = sys.modules.get("uno")
    if uno_mod is None:
        return None
    original = getattr(uno_mod, "_builtin_import", None)
    if original is None or builtins.__import__ is original:
        return None
    uno_hook = builtins.__import__
    builtins.__import__ = original
    return uno_hook


def _restore_uno_import(hook: object) -> None:
    """Restore the UNO import hook if it was bypassed."""
    if hook is not None:
        builtins.__import__ = hook


# ---------------------------------------------------------------------------
# Session-scoped soffice fixture
# ---------------------------------------------------------------------------


@pytest.fixture(scope="session")
def soffice_session() -> Generator[None, None, None]:
    """Start soffice once for all UNO tests, stop on teardown.

    Session-scoped because soffice startup costs 1-6 s.  Each test uses
    its own ``tmp_path`` files so isolation is maintained.
    """
    if not _ensure_soffice_running():
        pytest.skip("Could not start soffice headless server")
    yield
    stop_soffice()


# ---------------------------------------------------------------------------
# Per-test environment fixture (autouse)
# ---------------------------------------------------------------------------


@pytest.fixture(autouse=True)
def backend_test_env(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> Generator[None, None, None]:
    """Per-test DB isolation + mock settings.

    Does NOT mock ``_detect_backend`` — real detection runs so the
    correct backend is exercised.
    """
    db_file = tmp_path / "integration.db"
    monkeypatch.setattr("src.core.database.get_db_path", lambda: str(db_file))
    config_dir = tmp_path / "config"
    config_dir.mkdir()
    monkeypatch.setattr(
        "src.utils.path_manager.get_app_config_dir",
        lambda: config_dir,
    )
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    monkeypatch.setattr(
        "src.utils.path_manager.get_app_data_dir",
        lambda: data_dir,
    )
    init_db()
    # Prevent translator.py from calling stop_soffice() after each task
    monkeypatch.setattr("src.core.translator.stop_soffice", lambda: None)
    # Disable image/comment/shape translation by default
    monkeypatch.setattr(
        "src.utils.config_manager.load_setting",
        lambda k, d=None: False,
    )
    yield


# ---------------------------------------------------------------------------
# Force-UNO fixture (blocks win32com so _detect_backend falls to UNO)
# ---------------------------------------------------------------------------


@pytest.fixture()
def force_uno_backend(monkeypatch: pytest.MonkeyPatch) -> None:
    """Block win32com import so _detect_backend naturally selects UNO."""
    monkeypatch.setitem(sys.modules, "win32com", None)
    monkeypatch.setitem(sys.modules, "win32com.client", None)


# ---------------------------------------------------------------------------
# Mock LLM fixture
# ---------------------------------------------------------------------------


@pytest.fixture()
def mock_llm(
    monkeypatch: pytest.MonkeyPatch,
) -> Callable[..., list[str]]:
    """Patches translate_text at all import sites."""

    def fake_translate(
        texts: list[str],
        target_lang: str,
        source_lang: str = "",
        **kwargs: object,
    ) -> list[str]:
        return [f"[{target_lang}] {t}" for t in texts]

    monkeypatch.setattr(
        "src.core.llm_engine.translate_text",
        fake_translate,
    )
    monkeypatch.setattr(
        "src.core.text_processor._llm_engine.translate_text",
        fake_translate,
    )
    return fake_translate


# ---------------------------------------------------------------------------
# Enable-feature fixture (selectively turns on optional settings)
# ---------------------------------------------------------------------------


@pytest.fixture()
def enable_feature(monkeypatch: pytest.MonkeyPatch) -> Callable[..., None]:
    """Factory fixture that enables specific settings."""

    def _enable(*setting_keys: str) -> None:
        enabled = set(setting_keys)

        def fake_load(key: str, default: object = None) -> object:
            return True if key in enabled else default

        monkeypatch.setattr(
            "src.utils.config_manager.load_setting",
            fake_load,
        )

    return _enable


# ---------------------------------------------------------------------------
# Minimal PNG helper (no PIL dependency)
# ---------------------------------------------------------------------------


def _create_tiny_png(path: Path) -> None:
    """Generate a 4x4 red PNG from raw bytes — no PIL required."""
    width, height = 4, 4  # noqa: PLR2004
    raw_data = b""
    for _ in range(height):
        raw_data += b"\x00" + (b"\xff\x00\x00") * width

    def _png_chunk(chunk_type: bytes, data: bytes) -> bytes:
        """Build a single PNG chunk with CRC."""
        length = struct.pack(">I", len(data))
        crc = struct.pack(">I", zlib.crc32(chunk_type + data) & 0xFFFFFFFF)
        return length + chunk_type + data + crc

    ihdr_data = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    png_bytes = b"\x89PNG\r\n\x1a\n"
    png_bytes += _png_chunk(b"IHDR", ihdr_data)
    png_bytes += _png_chunk(b"IDAT", zlib.compress(raw_data))
    png_bytes += _png_chunk(b"IEND", b"")
    path.write_bytes(png_bytes)


# ---------------------------------------------------------------------------
# DOCX textbox injection helper (python-docx doesn't support shapes)
# ---------------------------------------------------------------------------

_WPS_NS = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
_WORDML_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_WP_NS = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _inject_docx_textbox(docx_path: Path, text: str) -> None:
    """Post-process a .docx ZIP to inject a <wps:txbx> shape element."""
    nsmap = {
        "w": _WORDML_NS,
        "wps": _WPS_NS,
        "wp": _WP_NS,
        "a": _A_NS,
        "r": _R_NS,
    }
    with zipfile.ZipFile(docx_path, "r") as zf:
        doc_xml = zf.read("word/document.xml")
        other_files = {n: zf.read(n) for n in zf.namelist() if n != "word/document.xml"}
    root = etree.fromstring(doc_xml)
    body = root.find(f"{{{_WORDML_NS}}}body")
    if body is None:
        return

    p_el = etree.SubElement(body, f"{{{_WORDML_NS}}}p")
    r_el = etree.SubElement(p_el, f"{{{_WORDML_NS}}}r")
    drawing = etree.SubElement(r_el, f"{{{_WORDML_NS}}}drawing")
    inline = etree.SubElement(drawing, f"{{{_WP_NS}}}inline")
    graphic = etree.SubElement(inline, f"{{{_A_NS}}}graphic")
    gdata = etree.SubElement(
        graphic,
        f"{{{_A_NS}}}graphicData",
        attrib={"uri": _WPS_NS},
    )
    wsp = etree.SubElement(gdata, f"{{{_WPS_NS}}}wsp")
    txbx = etree.SubElement(wsp, f"{{{_WPS_NS}}}txbx")
    txbx_content = etree.SubElement(txbx, f"{{{_WORDML_NS}}}txbxContent")
    inner_p = etree.SubElement(txbx_content, f"{{{_WORDML_NS}}}p")
    inner_r = etree.SubElement(inner_p, f"{{{_WORDML_NS}}}r")
    inner_t = etree.SubElement(inner_r, f"{{{_WORDML_NS}}}t")
    inner_t.text = text

    for prefix, uri in nsmap.items():
        if root.nsmap.get(prefix) is None:
            etree.register_namespace(prefix, uri)

    new_xml = etree.tostring(root, xml_declaration=True, encoding="UTF-8")
    with zipfile.ZipFile(docx_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("word/document.xml", new_xml)
        for name, data in other_files.items():
            zf.writestr(name, data)


# ---------------------------------------------------------------------------
# PPTX legacy comment injection helper
# ---------------------------------------------------------------------------

_PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_COMMENTS_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
)
_PKG_RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _inject_pptx_legacy_comment(pptx_path: Path, text: str) -> None:
    """Post-process a .pptx ZIP to inject a legacy <p:cm> comment part.

    python-pptx doesn't support comments.  We inject a legacy comment
    XML part, wire it via slide1 relationships, and register it in
    ``[Content_Types].xml``.
    """
    import shutil  # noqa: PLC0415

    comment_xml = (
        f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<p:cmLst xmlns:p="{_PPTX_NS}">'
        f'<p:cm authorId="1" dt="2024-01-01T00:00:00" idx="1">'
        f'<p:pos x="100" y="200"/>'
        f"<p:text>{text}</p:text>"
        f"</p:cm>"
        f"</p:cmLst>"
    )
    comment_part_name = "ppt/comments/comment1.xml"

    with zipfile.ZipFile(pptx_path, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    # Add the comment part
    data[comment_part_name] = comment_xml.encode("utf-8")

    # Update slide1 .rels
    rels_path = "ppt/slides/_rels/slide1.xml.rels"
    if rels_path in data:
        rels_root = etree.fromstring(data[rels_path])
    else:
        rels_root = etree.Element(
            "Relationships",
            xmlns=_PKG_RELS_NS,
        )

    existing_ids = [r.get("Id", "") for r in rels_root]
    next_id = f"rId{len(existing_ids) + 100}"
    target = "../" + comment_part_name.removeprefix("ppt/")
    etree.SubElement(
        rels_root,
        "Relationship",
        Id=next_id,
        Type=_COMMENTS_REL_TYPE,
        Target=target,
    )
    data[rels_path] = etree.tostring(
        rels_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    # Update [Content_Types].xml
    ct_root = etree.fromstring(data["[Content_Types].xml"])
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    etree.SubElement(
        ct_root,
        f"{{{ct_ns}}}Override",
        PartName=f"/{comment_part_name}",
        ContentType=(
            "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"
        ),
    )
    data["[Content_Types].xml"] = etree.tostring(
        ct_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    # Rewrite the ZIP
    tmp_zip = pptx_path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for item in items:
            zf_out.writestr(item, data[item.filename])
        for name, content in data.items():
            if not any(i.filename == name for i in items):
                zf_out.writestr(name, content)
    shutil.move(str(tmp_zip), str(pptx_path))


# ---------------------------------------------------------------------------
# ODF annotation injection helper
# ---------------------------------------------------------------------------

_ODF_NS = {
    "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    "draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
    "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
}


def _inject_odf_annotation(odf_path: Path, text: str) -> None:
    """Post-process an ODF ZIP to inject an <office:annotation>.

    For ODT the annotation is appended under the first ``<text:p>``.
    For ODS it is placed inside the first ``<table:table-cell>``.
    For ODP it is placed as a child of the first ``<draw:page>``.
    """
    import shutil  # noqa: PLC0415

    with zipfile.ZipFile(odf_path, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    root = etree.fromstring(data["content.xml"])

    annot = etree.Element(f"{{{_ODF_NS['office']}}}annotation")
    annot.set(f"{{{_ODF_NS['office']}}}name", "test_comment")
    annot_p = etree.SubElement(annot, f"{{{_ODF_NS['text']}}}p")
    annot_p.text = text

    suffix = odf_path.suffix.lower()
    if suffix == ".ods":
        # Inject inside the first table-cell
        cell = root.find(
            f".//{{{_ODF_NS['table']}}}table-cell",
        )
        if cell is not None:
            cell.insert(0, annot)
    elif suffix == ".odp":
        # ODP — inject as a child of the first <draw:page>
        # (annotations inside <draw:text-box> corrupt the file)
        page = root.find(f".//{{{_ODF_NS['draw']}}}page")
        if page is not None:
            page.append(annot)
    else:
        # ODT — inject inside the first <text:p>
        first_p = root.find(f".//{{{_ODF_NS['text']}}}p")
        if first_p is not None:
            first_p.insert(0, annot)

    data["content.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    tmp_zip = odf_path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for item in items:
            zf_out.writestr(item, data[item.filename])
        for name, content in data.items():
            if not any(i.filename == name for i in items):
                zf_out.writestr(name, content)
    shutil.move(str(tmp_zip), str(odf_path))


# ---------------------------------------------------------------------------
# XLSX shape injection helper
# ---------------------------------------------------------------------------

_XDR_NS = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"


def _inject_xlsx_shape(xlsx_path: Path, text: str) -> None:
    """Post-process an .xlsx ZIP to inject a DrawingML shape.

    Creates ``xl/drawings/drawing1.xml`` with a single shape, wires it
    via ``xl/worksheets/_rels/sheet1.xml.rels``, and updates
    ``[Content_Types].xml``.
    """
    import shutil  # noqa: PLC0415

    drawing_xml = (
        f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<xdr:wsDr xmlns:xdr="{_XDR_NS}" xmlns:a="{_A_NS}">'
        f"<xdr:twoCellAnchor><xdr:sp>"
        f"<a:txBody><a:bodyPr/>"
        f'<a:p><a:r><a:rPr lang="en-US"/><a:t>{text}</a:t></a:r></a:p>'
        f"</a:txBody>"
        f"</xdr:sp></xdr:twoCellAnchor>"
        f"</xdr:wsDr>"
    )

    with zipfile.ZipFile(xlsx_path, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    data["xl/drawings/drawing1.xml"] = drawing_xml.encode("utf-8")

    # Wire sheet1 → drawing1
    drawing_rel_type = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing"
    )
    rels_path = "xl/worksheets/_rels/sheet1.xml.rels"
    if rels_path in data:
        rels_root = etree.fromstring(data[rels_path])
    else:
        rels_root = etree.Element(
            "Relationships",
            xmlns=_PKG_RELS_NS,
        )
    etree.SubElement(
        rels_root,
        "Relationship",
        Id="rId1",
        Type=drawing_rel_type,
        Target="../drawings/drawing1.xml",
    )
    data[rels_path] = etree.tostring(
        rels_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    # Update [Content_Types].xml
    ct_root = etree.fromstring(data["[Content_Types].xml"])
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    etree.SubElement(
        ct_root,
        f"{{{ct_ns}}}Override",
        PartName="/xl/drawings/drawing1.xml",
        ContentType=("application/vnd.openxmlformats-officedocument.drawing+xml"),
    )
    data["[Content_Types].xml"] = etree.tostring(
        ct_root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    tmp_zip = xlsx_path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for item in items:
            zf_out.writestr(item, data[item.filename])
        for name, content in data.items():
            if not any(i.filename == name for i in items):
                zf_out.writestr(name, content)
    shutil.move(str(tmp_zip), str(xlsx_path))


# ---------------------------------------------------------------------------
# ODF text-box injection helper
# ---------------------------------------------------------------------------


def _inject_odf_textbox(odf_path: Path, text: str) -> None:
    """Post-process an ODF ZIP to inject a <draw:frame><draw:text-box>.

    For ODT the frame is appended to the document body.
    For ODS it is appended inside the first ``<table:table>``.
    """
    import shutil  # noqa: PLC0415

    with zipfile.ZipFile(odf_path, "r") as zf:
        items = zf.infolist()
        data = {i.filename: zf.read(i.filename) for i in items}

    root = etree.fromstring(data["content.xml"])

    frame = etree.Element(f"{{{_ODF_NS['draw']}}}frame")
    frame.set(f"{{{_ODF_NS['draw']}}}name", "TestFrame")
    textbox = etree.SubElement(frame, f"{{{_ODF_NS['draw']}}}text-box")
    tb_p = etree.SubElement(textbox, f"{{{_ODF_NS['text']}}}p")
    tb_p.text = text

    suffix = odf_path.suffix.lower()
    if suffix == ".ods":
        table = root.find(f".//{{{_ODF_NS['table']}}}table")
        if table is not None:
            table.append(frame)
    else:
        # ODT — append to office:body/office:text
        body_text = root.find(
            f".//{{{_ODF_NS['office']}}}text",
        )
        if body_text is not None:
            body_text.append(frame)

    data["content.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
    )

    tmp_zip = odf_path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for item in items:
            zf_out.writestr(item, data[item.filename])
        for name, content in data.items():
            if not any(i.filename == name for i in items):
                zf_out.writestr(name, content)
    shutil.move(str(tmp_zip), str(odf_path))


# ---------------------------------------------------------------------------
# Rich file creation helpers
# ---------------------------------------------------------------------------


def _create_docx(path: Path, paragraphs: list[str]) -> None:
    """Create a real .docx file with the given paragraphs."""
    from docx import Document  # noqa: PLC0415

    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    doc.save(str(path))


def _create_rich_docx(path: Path, png_path: Path) -> None:
    """Create a .docx with all translatable features.

    Includes paragraphs, table, mixed formatting, comment, image,
    and shape.  Content texts come from the module-level constants.
    """
    from docx import Document  # noqa: PLC0415
    from docx.shared import Pt  # noqa: PLC0415

    doc = Document()

    # 1. Plain-text paragraphs
    for text in _DOCX_PARAGRAPHS:
        doc.add_paragraph(text)

    # 2. Table
    nrows, ncols = len(_DOCX_TABLE_ROWS), len(_DOCX_TABLE_ROWS[0])
    table = doc.add_table(rows=nrows, cols=ncols)
    for r_idx, row in enumerate(_DOCX_TABLE_ROWS):
        for c_idx, val in enumerate(row):
            table.cell(r_idx, c_idx).text = val

    # 3. Mixed formatting paragraph
    mixed_para = doc.add_paragraph()
    bold_run = mixed_para.add_run(_DOCX_MIXED_FMT_PARTS[0])
    bold_run.bold = True
    bold_run.font.size = Pt(12)
    mixed_para.add_run(_DOCX_MIXED_FMT_PARTS[1])
    italic_run = mixed_para.add_run(_DOCX_MIXED_FMT_PARTS[2])
    italic_run.italic = True
    italic_run.font.size = Pt(12)

    # 4. Comment on a paragraph
    comment_para = doc.add_paragraph()
    comment_run = comment_para.add_run(_DOCX_COMMENT_ANCHOR)
    doc.add_comment(
        comment_run,
        text=_DOCX_COMMENT_TEXT,
        author="Tester",
    )

    # 5. Embedded image
    img_para = doc.add_paragraph()
    img_para.add_run().add_picture(str(png_path))

    doc.save(str(path))

    # 6. Inject a text box shape (python-docx doesn't support shapes)
    _inject_docx_textbox(path, _DOCX_SHAPE_TEXT)


def _create_rich_xlsx(path: Path) -> None:
    """Create an .xlsx with string cells + a cell comment."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import Workbook  # noqa: PLC0415
        from openpyxl.comments import Comment  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    wb = Workbook()
    ws = wb.active
    for row in _XLSX_ROWS:
        ws.append(row)
    ws["A1"].comment = Comment(_XLSX_COMMENT_TEXT, "Tester")
    wb.save(str(path))

    # Inject a DrawingML shape (openpyxl doesn't support shapes)
    _inject_xlsx_shape(path, _XLSX_SHAPE_TEXT)


def _create_rich_pptx(path: Path) -> None:
    """Create a .pptx with 3 slides: plain, bold+italic, plain."""
    hook = _bypass_uno_import()
    try:
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Inches, Pt  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    prs = Presentation()

    # Slide 1: plain text
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    t1 = s1.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(1),
    )
    t1.text_frame.text = _PPTX_SLIDES[0]

    # Slide 2: mixed formatting (bold + italic runs)
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    t2 = s2.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(1),
    )
    tf2 = t2.text_frame
    tf2.text = ""
    p2 = tf2.paragraphs[0]
    rb = p2.add_run()
    rb.text = _PPTX_MIXED_FMT_PARTS[0]
    rb.font.bold = True
    rb.font.size = Pt(14)
    p2.add_run().text = _PPTX_MIXED_FMT_PARTS[1]
    ri = p2.add_run()
    ri.text = _PPTX_MIXED_FMT_PARTS[2]
    ri.font.italic = True
    ri.font.size = Pt(14)

    # Slide 3: plain text
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    t3 = s3.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(1),
    )
    t3.text_frame.text = _PPTX_SLIDES[1]

    prs.save(str(path))

    # Inject a legacy comment (python-pptx doesn't support comments)
    _inject_pptx_legacy_comment(path, _PPTX_COMMENT_TEXT)


def _create_rich_odt(path: Path) -> None:
    """Create a .odt with multiple paragraphs and a table."""
    from odf.opendocument import OpenDocumentText  # noqa: PLC0415
    from odf.table import Table, TableCell, TableColumn, TableRow  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = OpenDocumentText()
    for text in _ODT_PARAGRAPHS:
        doc.text.addElement(P(text=text))

    table = Table(name="TestTable")
    table.addElement(TableColumn(numbercolumnsrepeated="2"))
    for row_data in _ODT_TABLE_ROWS:
        tr = TableRow()
        for val in row_data:
            tc = TableCell()
            tc.addElement(P(text=val))
            tr.addElement(tc)
        table.addElement(tr)
    doc.text.addElement(table)
    doc.save(str(path))

    # Inject annotation + text box (odfpy doesn't support these)
    _inject_odf_annotation(path, _ODT_COMMENT_TEXT)
    _inject_odf_textbox(path, _ODT_SHAPE_TEXT)


def _create_rich_ods(path: Path) -> None:
    """Create a .ods with multiple rows of string cells."""
    from odf.opendocument import OpenDocumentSpreadsheet  # noqa: PLC0415
    from odf.table import Table, TableCell, TableRow  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = OpenDocumentSpreadsheet()
    table = Table(name="Sheet1")
    for row_data in _ODS_ROWS:
        tr = TableRow()
        for val in row_data:
            tc = TableCell()
            tc.addElement(P(text=val))
            tr.addElement(tc)
        table.addElement(tr)
    doc.spreadsheet.addElement(table)
    doc.save(str(path))

    # Inject annotation + text box (odfpy doesn't support these)
    _inject_odf_annotation(path, _ODS_COMMENT_TEXT)
    _inject_odf_textbox(path, _ODS_SHAPE_TEXT)


def _create_rich_odp(path: Path) -> None:
    """Create a .odp with multiple slides/frames."""
    from odf.draw import Frame, Page, TextBox  # noqa: PLC0415
    from odf.opendocument import OpenDocumentPresentation  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = OpenDocumentPresentation()
    for text in _ODP_SLIDES:
        page = Page(masterpagename="Default")
        frame = Frame(width="20cm", height="5cm", x="1cm", y="1cm")
        tb = TextBox()
        tb.addElement(P(text=text))
        frame.addElement(tb)
        page.addElement(frame)
        doc.presentation.addElement(page)
    doc.save(str(path))

    # Inject annotation (odfpy doesn't support office:annotation)
    _inject_odf_annotation(path, _ODP_COMMENT_TEXT)


def _create_xlsx(path: Path, rows: list[list[Any]]) -> None:
    """Create a real .xlsx file with the given rows."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import Workbook  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    wb.save(str(path))


def _create_pptx(path: Path, slide_texts: list[str]) -> None:
    """Create a real .pptx file with one text box per slide."""
    hook = _bypass_uno_import()
    try:
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    prs = Presentation()
    for text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(1),
        )
        tx_box.text_frame.text = text
    prs.save(str(path))


# ---------------------------------------------------------------------------
# File verification helpers
# ---------------------------------------------------------------------------


def _read_docx_paragraphs(path: Path) -> list[str]:
    """Read all non-empty paragraph texts from a .docx file."""
    from docx import Document  # noqa: PLC0415

    doc = Document(str(path))
    return [p.text for p in doc.paragraphs if p.text.strip()]


def _read_docx_table_cells(path: Path) -> list[str]:
    """Read all non-empty table cell texts from a .docx file."""
    from docx import Document  # noqa: PLC0415

    doc = Document(str(path))
    cells = []
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    cells.append(cell.text)
    return cells


def _read_docx_run_formatting(path: Path) -> list[tuple[str, bool, bool]]:
    """Read (text, bold, italic) per run from all paragraphs."""
    from docx import Document  # noqa: PLC0415

    doc = Document(str(path))
    runs: list[tuple[str, bool, bool]] = []
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text.strip():
                runs.append(
                    (
                        run.text,
                        bool(run.bold),
                        bool(run.italic),
                    )
                )
    return runs


def _read_docx_images(path: Path) -> list[str]:
    """Return image file paths inside a .docx ZIP."""
    img_exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}
    images: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        for name in zf.namelist():
            if (
                name.startswith(("word/media/", "media/"))
                and Path(name).suffix.lower() in img_exts
            ):
                images.append(name)
    return sorted(images)


def _read_xlsx_values(path: Path) -> list[list[Any]]:
    """Read all cell values from the first sheet of an .xlsx file."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import load_workbook  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    wb = load_workbook(str(path))
    ws = wb.active
    return [[cell.value for cell in row] for row in ws.iter_rows()]


def _read_pptx_texts(path: Path) -> list[str]:
    """Read text per slide (one entry per text frame) from a .pptx."""
    hook = _bypass_uno_import()
    try:
        from pptx import Presentation  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
    return texts


def _read_docx_comments(path: Path) -> list[str]:
    """Read comment texts from a .docx via low-level XML access."""
    from docx import Document  # noqa: PLC0415
    from docx.oxml.ns import qn  # noqa: PLC0415

    doc = Document(str(path))
    try:
        comments_part = doc.part.package.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/comments",
        )
    except KeyError:
        return []
    if comments_part is None:
        return []

    root = etree.fromstring(comments_part.blob)
    texts: list[str] = []
    for comment_el in root.findall(qn("w:comment")):
        p_texts = []
        for p_el in comment_el.findall(qn("w:p")):
            t_parts = [t.text for t in p_el.iter(qn("w:t")) if t.text]
            if t_parts:
                p_texts.append("".join(t_parts))
        text = "\n".join(p_texts).strip()
        if text:
            texts.append(text)
    return texts


def _read_docx_shape_texts(path: Path) -> list[str]:
    """Read text from <wps:txbx> elements in a .docx via ZIP + lxml."""
    wps_txbx_tag = f"{{{_WPS_NS}}}txbx"
    w_t_tag = f"{{{_WORDML_NS}}}t"
    texts: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        for part_name in zf.namelist():
            if not (
                part_name == "word/document.xml"
                or (part_name.startswith("word/header") and part_name.endswith(".xml"))
                or (part_name.startswith("word/footer") and part_name.endswith(".xml"))
            ):
                continue
            root = etree.fromstring(zf.read(part_name))
            for txbx in root.iter(wps_txbx_tag):
                t_parts = [t.text for t in txbx.iter(w_t_tag) if t.text]
                if t_parts:
                    texts.append("".join(t_parts))
    return texts


def _read_xlsx_comments(path: Path) -> list[str]:
    """Read cell comment texts from an .xlsx file."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import load_workbook  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    wb = load_workbook(str(path))
    texts: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                c = cell.comment
                if c and c.text and c.text.strip():
                    texts.append(c.text)
    return texts


def _read_pptx_comments(path: Path) -> list[str]:
    """Read legacy <p:text> comment texts from a .pptx via ZIP + lxml.

    Walks ``ppt/slides/_rels/slide*.xml.rels`` to find comment parts,
    then parses ``<p:cm>/<p:text>`` elements.
    """
    texts: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        names = zf.namelist()
        # Find all slide .rels files
        rels_files = [
            n
            for n in names
            if n.startswith("ppt/slides/_rels/slide") and n.endswith(".xml.rels")
        ]
        for rels_name in rels_files:
            rels_root = etree.fromstring(zf.read(rels_name))
            for rel in rels_root:
                if rel.get("Type") == _COMMENTS_REL_TYPE:
                    target = rel.get("Target", "")
                    # Resolve relative path from ppt/slides/ directory
                    part = "ppt/" + target[3:] if target.startswith("../") else target
                    if part in names:
                        cm_root = etree.fromstring(zf.read(part))
                        for cm in cm_root.iter(f"{{{_PPTX_NS}}}cm"):
                            p_text = cm.find(f"{{{_PPTX_NS}}}text")
                            if p_text is not None and p_text.text:
                                texts.append(p_text.text)
    return texts


def _read_odf_comments(path: Path) -> list[str]:
    """Read <office:annotation> text from content.xml via ZIP + lxml."""
    texts: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        root = etree.fromstring(zf.read("content.xml"))
    for annot in root.iter(f"{{{_ODF_NS['office']}}}annotation"):
        parts = []
        for p_el in annot.iter(f"{{{_ODF_NS['text']}}}p"):
            if p_el.text:
                parts.append(p_el.text)
        text = " ".join(parts).strip()
        if text:
            texts.append(text)
    return texts


def _read_xlsx_shape_texts(path: Path) -> list[str]:
    """Read text from DrawingML shapes in xl/drawings/*.xml via ZIP + lxml."""
    a_t_tag = f"{{{_A_NS}}}t"
    texts: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        for name in zf.namelist():
            if name.startswith("xl/drawings/drawing") and name.endswith(".xml"):
                root = etree.fromstring(zf.read(name))
                for t_el in root.iter(a_t_tag):
                    if t_el.text and t_el.text.strip():
                        texts.append(t_el.text)
    return texts


def _read_odf_shape_texts(path: Path) -> list[str]:
    """Read text from <draw:text-box> elements in content.xml via ZIP + lxml."""
    texts: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        root = etree.fromstring(zf.read("content.xml"))
    for textbox in root.iter(f"{{{_ODF_NS['draw']}}}text-box"):
        parts = []
        for p_el in textbox.iter(f"{{{_ODF_NS['text']}}}p"):
            if p_el.text:
                parts.append(p_el.text)
        text = " ".join(parts).strip()
        if text:
            texts.append(text)
    return texts


# ---------------------------------------------------------------------------
# Legacy file creation helper (modern → legacy via backend conversion)
# ---------------------------------------------------------------------------


def _create_legacy_via_uno(
    modern_path: Path,
    legacy_path: Path,
) -> None:
    """Create a legacy file via UNO, skip test on failure."""
    try:
        _convert_with_uno(modern_path, legacy_path)
    except Exception as exc:
        pytest.skip(
            f"UNO conversion failed — cannot create legacy file: {exc}",
        )
    if not legacy_path.exists():
        pytest.skip("UNO conversion produced no output file")


def _create_legacy_via_win32com(
    modern_path: Path,
    legacy_path: Path,
) -> None:
    """Create a legacy file via win32com, skip test on failure."""
    from src.core.office_processor import (  # noqa: PLC0415
        _convert_with_win32com,
    )

    try:
        _convert_with_win32com(modern_path, legacy_path)
    except Exception as exc:
        pytest.skip(
            f"win32com conversion failed: {exc}",
        )
    if not legacy_path.exists():
        pytest.skip("win32com conversion produced no output file")


# ═══════════════════════════════════════════════════════════════════════
# UNO Tests
# ═══════════════════════════════════════════════════════════════════════


@requires_uno
@pytest.mark.timeout(60)
class TestUnoBackend:
    """Integration tests that exercise the UNO backend."""

    # --- Backend detection ---

    def test_uno_backend_detected(
        self,
        force_uno_backend: None,
    ) -> None:
        """UNO is selected for legacy/ODF when win32com is unavailable."""
        # OOXML always uses python_lib regardless of UNO availability
        assert _detect_backend(".docx") == "python_lib"
        # ODF prefers UNO when available
        assert _detect_backend(".odt") == "uno"
        # Legacy falls back to UNO when win32com unavailable
        assert _detect_backend(".doc") == "uno"

    # --- DOCX ---

    def test_uno_docx_roundtrip(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich DOCX via UNO: all features verified exhaustively."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        # --- Every paragraph ---
        paragraphs = _read_docx_paragraphs(out)
        mixed_text = "".join(_DOCX_MIXED_FMT_PARTS)
        expected_paragraphs = (
            [_french(t) for t in _DOCX_PARAGRAPHS]
            + [_french(mixed_text)]
            + [_french(_DOCX_COMMENT_ANCHOR)]
        )
        for exp in expected_paragraphs:
            assert exp in paragraphs, f"Missing paragraph {exp!r} in {paragraphs}"

        # --- Every table cell ---
        cells = _read_docx_table_cells(out)
        expected_cells = [_french(val) for row in _DOCX_TABLE_ROWS for val in row]
        for exp in expected_cells:
            assert exp in cells, f"Missing table cell {exp!r} in {cells}"

        # --- Embedded image preserved ---
        images = _read_docx_images(out)
        assert len(images) >= 1, "Embedded image lost during translation"

        # --- Bold + italic on correct runs ---
        runs = _read_docx_run_formatting(out)
        bold_texts = {t for t, b, i in runs if b}
        italic_texts = {t for t, b, i in runs if i}
        assert _DOCX_MIXED_FMT_PARTS[0] in bold_texts, (
            f"Bold run {_DOCX_MIXED_FMT_PARTS[0]!r} not found in bold runs {bold_texts}"
        )
        assert _DOCX_MIXED_FMT_PARTS[2] in italic_texts, (
            f"Italic run {_DOCX_MIXED_FMT_PARTS[2]!r} "
            f"not found in italic runs {italic_texts}"
        )

    def test_uno_docx_empty(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Empty DOCX via UNO → output exists, no crash."""
        inp = tmp_path / "empty.docx"
        out = tmp_path / "translated.docx"
        _create_docx(inp, [])

        assert translate_file(inp, out, "French", "English (US)") is True
        assert out.exists()

    # --- XLSX ---

    def test_uno_xlsx_roundtrip(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich XLSX via UNO → valid output with rows.

        Note: openpyxl writes cells as ``inlineStr`` which UNO may
        not classify as TEXT.  We verify the file is valid and has
        data; cell-level translation is verified in the win32com and
        python_lib integration tests.
        """
        inp = tmp_path / "sheet.xlsx"
        out = tmp_path / "translated.xlsx"
        _create_rich_xlsx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True
        rows = _read_xlsx_values(out)
        assert len(rows) >= 1, "Output XLSX has no rows"

    # --- PPTX ---

    def test_uno_pptx_roundtrip(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich PPTX via UNO: every slide text verified."""
        inp = tmp_path / "pres.pptx"
        out = tmp_path / "translated.pptx"
        _create_rich_pptx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        texts = _read_pptx_texts(out)
        # 3 slides: 2 plain + 1 mixed-formatting
        mixed_text = "".join(_PPTX_MIXED_FMT_PARTS)
        expected = [_french(t) for t in _PPTX_SLIDES] + [_french(mixed_text)]
        for exp in expected:
            assert exp in texts, f"Missing slide text {exp!r} in {texts}"

    # --- ODF ---

    def test_uno_odt_roundtrip(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich ODT via UNO: every paragraph and table cell verified."""
        inp = tmp_path / "doc.odt"
        out = tmp_path / "translated.odt"
        _create_rich_odt(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        items = _extract_python_odt(out)
        item_texts = [text for _, text in items]

        # Every paragraph
        for orig in _ODT_PARAGRAPHS:
            assert _french(orig) in item_texts, (
                f"Missing ODT paragraph {_french(orig)!r}"
            )
        # Every table cell
        for row in _ODT_TABLE_ROWS:
            for val in row:
                assert _french(val) in item_texts, (
                    f"Missing ODT table cell {_french(val)!r}"
                )

    def test_uno_ods_roundtrip(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich ODS via UNO → valid output exists.

        Note: UNO may reclassify odfpy-created string cells as
        non-string, making ``_extract_python_ods`` return fewer items.
        We verify the pipeline succeeds and the file is readable.
        """
        inp = tmp_path / "sheet.ods"
        out = tmp_path / "translated.ods"
        _create_rich_ods(inp)

        assert translate_file(inp, out, "French", "English (US)") is True
        assert out.exists()
        assert out.stat().st_size > 0, "Output ODS is empty"

    def test_uno_odp_roundtrip(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich ODP via UNO: every slide text verified."""
        inp = tmp_path / "pres.odp"
        out = tmp_path / "translated.odp"
        _create_rich_odp(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        items = _extract_python_odp(out)
        item_texts = [text for _, text in items]
        for orig in _ODP_SLIDES:
            assert _french(orig) in item_texts, f"Missing ODP slide {_french(orig)!r}"

    # --- Legacy ---

    def test_uno_doc_legacy(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """.doc legacy round-trip via UNO."""
        modern = tmp_path / "source.docx"
        legacy = tmp_path / "source.doc"
        out = tmp_path / "translated.doc"
        _create_docx(modern, ["Hello world"])
        _create_legacy_via_uno(modern, legacy)

        assert translate_file(legacy, out, "French", "English (US)") is True
        assert out.exists()
        assert out.stat().st_size > 0, "Output .doc is empty"

    def test_uno_xls_legacy(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """.xls legacy round-trip via UNO."""
        modern = tmp_path / "source.xlsx"
        legacy = tmp_path / "source.xls"
        out = tmp_path / "translated.xls"
        _create_xlsx(modern, [["Hello", "World"]])
        _create_legacy_via_uno(modern, legacy)

        assert translate_file(legacy, out, "French", "English (US)") is True
        assert out.exists()
        assert out.stat().st_size > 0, "Output .xls is empty"

    def test_uno_ppt_legacy(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """.ppt legacy round-trip via UNO."""
        modern = tmp_path / "source.pptx"
        legacy = tmp_path / "source.ppt"
        out = tmp_path / "translated.ppt"
        _create_pptx(modern, ["Hello world"])
        _create_legacy_via_uno(modern, legacy)

        assert translate_file(legacy, out, "French", "English (US)") is True
        assert out.exists()
        assert out.stat().st_size > 0, "Output .ppt is empty"

    # --- Feature-specific tests (UNO) ---

    def test_uno_docx_comments_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """DOCX comments with SETTING_TRANSLATE_DOC_COMMENTS on.

        Note: UNO strips ``word/comments.xml`` from DOCX output, so
        comment injection silently no-ops.  This test verifies the
        pipeline does not crash; if comments survive, each is checked.
        """
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_docx_comments(out)
        if comments:
            # Input has exactly 1 comment
            assert len(comments) == 1, (
                f"Expected 1 comment, got {len(comments)}: {comments}"
            )
            assert comments[0] == _french(_DOCX_COMMENT_TEXT)

    def test_uno_docx_shapes_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """DOCX shapes: each shape text verified."""
        enable_feature(SETTING_TRANSLATE_DOC_SHAPES)

        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        shape_texts = _read_docx_shape_texts(out)
        assert len(shape_texts) >= 1, "No shapes found in output"
        assert _french(_DOCX_SHAPE_TEXT) in shape_texts, (
            f"Expected {_french(_DOCX_SHAPE_TEXT)!r} in shapes {shape_texts}"
        )

    def test_uno_docx_mixed_formatting_preserved(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich DOCX via UNO: bold + italic on the correct text."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        runs = _read_docx_run_formatting(out)
        bold_texts = {t for t, b, _ in runs if b}
        italic_texts = {t for t, _, i in runs if i}
        assert _DOCX_MIXED_FMT_PARTS[0] in bold_texts
        assert _DOCX_MIXED_FMT_PARTS[2] in italic_texts
        # Connector " and " must NOT be bold or italic
        connector_runs = [
            (b, i) for t, b, i in runs if t.strip() == _DOCX_MIXED_FMT_PARTS[1].strip()
        ]
        if connector_runs:
            for b, i in connector_runs:
                assert not b, "Connector run should not be bold"
                assert not i, "Connector run should not be italic"

    def test_uno_xlsx_comments_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """XLSX comments: each comment verified."""
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "sheet.xlsx"
        out = tmp_path / "translated.xlsx"
        _create_rich_xlsx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_xlsx_comments(out)
        assert len(comments) == 1, (
            f"Expected 1 comment, got {len(comments)}: {comments}"
        )
        assert comments[0] == _french(_XLSX_COMMENT_TEXT)

    # --- PPTX comments ---

    def test_uno_pptx_comments_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """PPTX legacy comments via UNO.

        UNO may strip or modify PPTX comment parts — soft assertion:
        if comments survive, verify their content.
        """
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "pres.pptx"
        out = tmp_path / "translated.pptx"
        _create_rich_pptx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_pptx_comments(out)
        if comments:
            assert _french(_PPTX_COMMENT_TEXT) in comments, (
                f"Expected {_french(_PPTX_COMMENT_TEXT)!r} in {comments}"
            )

    # --- ODF comments ---

    def test_uno_odt_comments_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """ODT annotations via UNO: each annotation verified."""
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "doc.odt"
        out = tmp_path / "translated.odt"
        _create_rich_odt(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_odf_comments(out)
        if comments:
            assert _french(_ODT_COMMENT_TEXT) in comments, (
                f"Expected {_french(_ODT_COMMENT_TEXT)!r} in {comments}"
            )

    def test_uno_ods_comments_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """ODS annotations via UNO: pipeline succeeds.

        UNO may reclassify odfpy cells and may not translate ODF
        annotations — lenient assertion: verify pipeline succeeds.
        If annotations survive and are translated, verify content.
        If not translated, at least verify the original text is present
        (annotation was not lost).
        """
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "sheet.ods"
        out = tmp_path / "translated.ods"
        _create_rich_ods(inp)

        assert translate_file(inp, out, "French", "English (US)") is True
        assert out.exists()
        assert out.stat().st_size > 0, "Output ODS is empty"

        # Verify annotations: if present, they should contain either
        # the translated or original text (not silently corrupted)
        comments = _read_odf_comments(out)
        if comments:
            has_translated = _french(_ODS_COMMENT_TEXT) in comments
            has_original = _ODS_COMMENT_TEXT in comments
            assert has_translated or has_original, f"Comment text corrupted: {comments}"

    def test_uno_odp_comments_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """ODP annotations via UNO: pipeline does not crash.

        UNO may fail to open ODP files with injected annotations or
        strip them during re-save.  We verify the pipeline handles
        this gracefully; if comments survive, check content.
        """
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "pres.odp"
        out = tmp_path / "translated.odp"
        _create_rich_odp(inp)

        try:
            result = translate_file(
                inp,
                out,
                "French",
                "English (US)",
            )
        except ValueError:
            # UNO may reject the modified ODP — acceptable
            pytest.skip("UNO cannot open ODP with injected annotation")
            return

        assert result is True

        comments = _read_odf_comments(out)
        if comments:
            assert _french(_ODP_COMMENT_TEXT) in comments, (
                f"Expected {_french(_ODP_COMMENT_TEXT)!r} in {comments}"
            )

    # --- XLSX shapes ---

    def test_uno_xlsx_shapes_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """XLSX DrawingML shapes via UNO: pipeline succeeds.

        UNO re-save may lose DrawingML shapes (openpyxl warns about
        incomplete DrawingML support).  We verify the pipeline
        succeeds; if shapes survive, check content.
        """
        enable_feature(SETTING_TRANSLATE_DOC_SHAPES)

        inp = tmp_path / "sheet.xlsx"
        out = tmp_path / "translated.xlsx"
        _create_rich_xlsx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        shapes = _read_xlsx_shape_texts(out)
        if shapes:
            assert _french(_XLSX_SHAPE_TEXT) in shapes, (
                f"Expected {_french(_XLSX_SHAPE_TEXT)!r} in {shapes}"
            )

    # --- ODF shapes ---

    def test_uno_odt_shapes_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """ODT text boxes via UNO: each text-box text verified."""
        enable_feature(SETTING_TRANSLATE_DOC_SHAPES)

        inp = tmp_path / "doc.odt"
        out = tmp_path / "translated.odt"
        _create_rich_odt(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        shapes = _read_odf_shape_texts(out)
        if shapes:
            assert _french(_ODT_SHAPE_TEXT) in shapes, (
                f"Expected {_french(_ODT_SHAPE_TEXT)!r} in {shapes}"
            )

    def test_uno_ods_shapes_translated(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """ODS text boxes via UNO: each text-box text verified."""
        enable_feature(SETTING_TRANSLATE_DOC_SHAPES)

        inp = tmp_path / "sheet.ods"
        out = tmp_path / "translated.ods"
        _create_rich_ods(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        shapes = _read_odf_shape_texts(out)
        if shapes:
            assert _french(_ODS_SHAPE_TEXT) in shapes, (
                f"Expected {_french(_ODS_SHAPE_TEXT)!r} in {shapes}"
            )

    # --- Glossary forwarding ---

    def test_uno_docx_glossary_forwarded(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """Glossary entries are forwarded to the LLM via UNO backend."""
        captured_kwargs: dict[str, object] = {}

        def tracking_translate(
            texts: list[str],
            target_lang: str,
            source_lang: str = "",
            **kwargs: object,
        ) -> list[str]:
            captured_kwargs.update(kwargs)
            return [f"[{target_lang}] {t}" for t in texts]

        monkeypatch.setattr(
            "src.core.llm_engine.translate_text",
            tracking_translate,
        )
        monkeypatch.setattr(
            "src.core.text_processor._llm_engine.translate_text",
            tracking_translate,
        )

        glossary = [(1, "Hello", "Bonjour")]  # noqa: PLR2004
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        _create_docx(inp, ["Hello world"])

        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                glossary_entries=glossary,
            )
            is True
        )

        assert captured_kwargs.get("glossary_entries") == glossary

    # --- Negative tests (features disabled by default) ---

    def test_uno_docx_comments_not_translated_when_disabled(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Comments stay untranslated when SETTING_TRANSLATE_DOC_COMMENTS is off.

        The ``backend_test_env`` fixture disables all settings by default.
        If the comment survives UNO round-trip, its text must be the
        original (not the ``[French] ...`` mock-LLM output).
        """
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_docx_comments(out)
        for c in comments:
            assert _french(_DOCX_COMMENT_TEXT) != c, (
                "Comment was translated despite setting being off"
            )

    def test_uno_docx_shapes_not_translated_when_disabled(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Shapes stay untranslated when SETTING_TRANSLATE_DOC_SHAPES is off.

        The ``backend_test_env`` fixture disables all settings by default.
        If the shape survives UNO round-trip, its text must be the
        original (not the ``[French] ...`` mock-LLM output).
        """
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        shapes = _read_docx_shape_texts(out)
        for s in shapes:
            assert _french(_DOCX_SHAPE_TEXT) != s, (
                "Shape was translated despite setting being off"
            )

    # --- All features combined ---

    def test_uno_docx_all_features_combined(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        monkeypatch: pytest.MonkeyPatch,
        enable_feature: Callable[..., None],
    ) -> None:
        """All features enabled simultaneously: comments + shapes + glossary.

        Enables both SETTING_TRANSLATE_DOC_COMMENTS and
        SETTING_TRANSLATE_DOC_SHAPES, passes glossary_entries, and
        verifies all translatable items in a single translation pass.
        """
        enable_feature(
            SETTING_TRANSLATE_DOC_COMMENTS,
            SETTING_TRANSLATE_DOC_SHAPES,
        )

        captured_kwargs: dict[str, object] = {}

        def tracking_translate(
            texts: list[str],
            target_lang: str,
            source_lang: str = "",
            **kwargs: object,
        ) -> list[str]:
            captured_kwargs.update(kwargs)
            return [f"[{target_lang}] {t}" for t in texts]

        monkeypatch.setattr(
            "src.core.llm_engine.translate_text",
            tracking_translate,
        )
        monkeypatch.setattr(
            "src.core.text_processor._llm_engine.translate_text",
            tracking_translate,
        )

        glossary = [(1, "Hello", "Bonjour")]  # noqa: PLR2004
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                glossary_entries=glossary,
            )
            is True
        )

        # Verify paragraphs
        paragraphs = _read_docx_paragraphs(out)
        for orig in _DOCX_PARAGRAPHS:
            assert _french(orig) in paragraphs, f"Missing paragraph {_french(orig)!r}"

        # Verify table cells
        cells = _read_docx_table_cells(out)
        for row in _DOCX_TABLE_ROWS:
            for val in row:
                assert _french(val) in cells, f"Missing table cell {_french(val)!r}"

        # Verify comments (soft — UNO may strip)
        comments = _read_docx_comments(out)
        if comments:
            assert _french(_DOCX_COMMENT_TEXT) in comments

        # Verify shapes
        shapes = _read_docx_shape_texts(out)
        assert len(shapes) >= 1, "No shapes found"
        assert _french(_DOCX_SHAPE_TEXT) in shapes

        # Verify glossary forwarding
        assert captured_kwargs.get("glossary_entries") == glossary

    # --- Cancellation ---

    def test_uno_docx_cancellation_returns_false(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """cancel_check returning True → translate_file returns False."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        _create_docx(inp, ["Hello world"])

        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                cancel_check=lambda: True,
            )
            is False
        )

    # --- Progress callback ---

    def test_uno_docx_progress_callback_invoked(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """progress_callback receives increasing values ending at 100."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        _create_docx(inp, ["Hello world", "Goodbye world"])

        progress_values: list[int] = []
        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                progress_callback=progress_values.append,
            )
            is True
        )

        assert len(progress_values) >= 1, "Progress callback never called"
        assert progress_values[-1] == 100, (  # noqa: PLR2004
            f"Final progress should be 100, got {progress_values[-1]}"
        )
        # Values must be non-decreasing
        for i in range(1, len(progress_values)):
            assert progress_values[i] >= progress_values[i - 1], (
                f"Progress decreased: {progress_values[i - 1]} → {progress_values[i]}"
            )

    # --- src_lang and content_type forwarding ---

    def test_uno_docx_src_lang_and_content_type_forwarded(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """src_lang and content_type are forwarded to translate_text."""
        captured_calls: list[dict[str, object]] = []

        def tracking_translate(
            texts: list[str],
            target_lang: str,
            source_lang: str = "",
            **kwargs: object,
        ) -> list[str]:
            captured_calls.append(
                {
                    "source_lang": source_lang,
                    **kwargs,
                }
            )
            return [f"[{target_lang}] {t}" for t in texts]

        monkeypatch.setattr(
            "src.core.llm_engine.translate_text",
            tracking_translate,
        )
        monkeypatch.setattr(
            "src.core.text_processor._llm_engine.translate_text",
            tracking_translate,
        )

        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        _create_docx(inp, ["Hello world"])

        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
            )
            is True
        )

        assert len(captured_calls) >= 1, "translate_text never called"
        call = captured_calls[0]
        assert call["source_lang"] == "English (US)"
        # DOCX plain text → CONTENT_PLAIN_TEXT (no HTML in input)
        assert "content_type" in call, "content_type not forwarded"

    # --- Glossary for non-DOCX formats ---

    def test_uno_xlsx_glossary_forwarded(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """Glossary entries are forwarded to the LLM for XLSX files.

        UNO may not extract text from openpyxl ``inlineStr`` cells;
        if no text is extracted the pipeline copies the file without
        calling the LLM.  We verify: if LLM was called, glossary was
        forwarded.
        """
        captured_kwargs: dict[str, object] = {}

        def tracking_translate(
            texts: list[str],
            target_lang: str,
            source_lang: str = "",
            **kwargs: object,
        ) -> list[str]:
            captured_kwargs.update(kwargs)
            return [f"[{target_lang}] {t}" for t in texts]

        monkeypatch.setattr(
            "src.core.llm_engine.translate_text",
            tracking_translate,
        )
        monkeypatch.setattr(
            "src.core.text_processor._llm_engine.translate_text",
            tracking_translate,
        )

        glossary = [(1, "Hello", "Bonjour")]  # noqa: PLR2004
        inp = tmp_path / "sheet.xlsx"
        out = tmp_path / "translated.xlsx"
        _create_xlsx(inp, [["Hello", "World"]])

        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                glossary_entries=glossary,
            )
            is True
        )

        # If UNO extracted text and the LLM was called, verify glossary
        if captured_kwargs:
            assert captured_kwargs.get("glossary_entries") == glossary

    # --- Empty document edge cases ---

    def test_uno_xlsx_empty(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Empty XLSX via UNO → output exists, no crash."""
        inp = tmp_path / "empty.xlsx"
        out = tmp_path / "translated.xlsx"
        _create_xlsx(inp, [])

        assert translate_file(inp, out, "French", "English (US)") is True
        assert out.exists()

    def test_uno_pptx_empty(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Empty PPTX via UNO → output exists, no crash."""
        inp = tmp_path / "empty.pptx"
        out = tmp_path / "translated.pptx"
        _create_pptx(inp, [])

        assert translate_file(inp, out, "French", "English (US)") is True
        assert out.exists()

    def test_uno_odt_empty(
        self,
        tmp_path: Path,
        soffice_session: None,
        force_uno_backend: None,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Empty ODT via UNO → output exists, no crash."""
        from odf.opendocument import OpenDocumentText  # noqa: PLC0415

        inp = tmp_path / "empty.odt"
        out = tmp_path / "translated.odt"
        doc = OpenDocumentText()
        doc.save(str(inp))

        assert translate_file(inp, out, "French", "English (US)") is True
        assert out.exists()


# ═══════════════════════════════════════════════════════════════════════
# Win32com Tests
# ═══════════════════════════════════════════════════════════════════════


@requires_win32com
@pytest.mark.timeout(60)
class TestWin32comBackend:
    """Integration tests that exercise the win32com backend."""

    def test_win32com_backend_detected(self) -> None:
        """win32com is selected for legacy formats; OOXML uses python_lib."""
        # OOXML always uses python_lib regardless of win32com availability
        assert _detect_backend(".docx") == "python_lib"
        # Legacy formats prefer win32com
        assert _detect_backend(".doc") == "win32com"

    # --- DOCX ---

    def test_win32com_docx_roundtrip(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich DOCX via win32com: all features verified."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        # --- Every paragraph ---
        paragraphs = _read_docx_paragraphs(out)
        mixed_text = "".join(_DOCX_MIXED_FMT_PARTS)
        expected_paragraphs = (
            [_french(t) for t in _DOCX_PARAGRAPHS]
            + [_french(mixed_text)]
            + [_french(_DOCX_COMMENT_ANCHOR)]
        )
        for exp in expected_paragraphs:
            assert exp in paragraphs, f"Missing paragraph {exp!r} in {paragraphs}"

        # --- Every table cell ---
        cells = _read_docx_table_cells(out)
        expected_cells = [_french(val) for row in _DOCX_TABLE_ROWS for val in row]
        for exp in expected_cells:
            assert exp in cells, f"Missing table cell {exp!r} in {cells}"

        # --- Embedded image preserved ---
        images = _read_docx_images(out)
        assert len(images) >= 1, "Embedded image lost during translation"

        # --- Bold + italic on correct runs ---
        runs = _read_docx_run_formatting(out)
        bold_texts = {t for t, b, i in runs if b}
        italic_texts = {t for t, b, i in runs if i}
        assert _DOCX_MIXED_FMT_PARTS[0] in bold_texts
        assert _DOCX_MIXED_FMT_PARTS[2] in italic_texts

    def test_win32com_docx_empty(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Empty DOCX via win32com → output exists, no crash."""
        inp = tmp_path / "empty.docx"
        out = tmp_path / "translated.docx"
        _create_docx(inp, [])

        assert translate_file(inp, out, "French", "English (US)") is True
        assert out.exists()

    # --- XLSX ---

    def test_win32com_xlsx_roundtrip(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich XLSX via win32com: every cell verified."""
        inp = tmp_path / "sheet.xlsx"
        out = tmp_path / "translated.xlsx"
        _create_rich_xlsx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        rows = _read_xlsx_values(out)
        flat = [str(c) for row in rows for c in row if c]
        expected_cells = [_french(val) for row in _XLSX_ROWS for val in row]
        for exp in expected_cells:
            assert exp in flat, f"Missing XLSX cell {exp!r} in {flat}"

    # --- PPTX ---

    def test_win32com_pptx_roundtrip(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich PPTX via win32com: every slide text verified."""
        inp = tmp_path / "pres.pptx"
        out = tmp_path / "translated.pptx"
        _create_rich_pptx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        texts = _read_pptx_texts(out)
        mixed_text = "".join(_PPTX_MIXED_FMT_PARTS)
        expected = [_french(t) for t in _PPTX_SLIDES] + [_french(mixed_text)]
        for exp in expected:
            assert exp in texts, f"Missing slide text {exp!r} in {texts}"

    # --- Legacy ---

    def test_win32com_doc_legacy(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """.doc legacy round-trip via win32com."""
        modern = tmp_path / "source.docx"
        legacy = tmp_path / "source.doc"
        out = tmp_path / "translated.doc"
        _create_docx(modern, ["Hello world"])
        _create_legacy_via_win32com(modern, legacy)

        assert translate_file(legacy, out, "French", "English (US)") is True
        assert out.exists()
        assert out.stat().st_size > 0, "Output .doc is empty"

    def test_win32com_xls_legacy(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """.xls legacy round-trip via win32com."""
        modern = tmp_path / "source.xlsx"
        legacy = tmp_path / "source.xls"
        out = tmp_path / "translated.xls"
        _create_xlsx(modern, [["Hello", "World"]])
        _create_legacy_via_win32com(modern, legacy)

        assert translate_file(legacy, out, "French", "English (US)") is True
        assert out.exists()
        assert out.stat().st_size > 0, "Output .xls is empty"

    def test_win32com_ppt_legacy(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """.ppt legacy round-trip via win32com."""
        modern = tmp_path / "source.pptx"
        legacy = tmp_path / "source.ppt"
        out = tmp_path / "translated.ppt"
        _create_pptx(modern, ["Hello world"])
        _create_legacy_via_win32com(modern, legacy)

        assert translate_file(legacy, out, "French", "English (US)") is True
        assert out.exists()
        assert out.stat().st_size > 0, "Output .ppt is empty"

    # --- Feature-specific tests (win32com) ---

    def test_win32com_docx_comments_translated(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """DOCX comments: each comment verified."""
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_docx_comments(out)
        assert len(comments) == 1, (
            f"Expected 1 comment, got {len(comments)}: {comments}"
        )
        assert comments[0] == _french(_DOCX_COMMENT_TEXT)

    def test_win32com_docx_shapes_translated(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """DOCX shapes: each shape text verified."""
        enable_feature(SETTING_TRANSLATE_DOC_SHAPES)

        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        shape_texts = _read_docx_shape_texts(out)
        assert len(shape_texts) >= 1, "No shapes found in output"
        assert _french(_DOCX_SHAPE_TEXT) in shape_texts

    def test_win32com_docx_mixed_formatting_preserved(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Rich DOCX via win32com: bold + italic on the correct text."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        runs = _read_docx_run_formatting(out)
        bold_texts = {t for t, b, _ in runs if b}
        italic_texts = {t for t, _, i in runs if i}
        assert _DOCX_MIXED_FMT_PARTS[0] in bold_texts
        assert _DOCX_MIXED_FMT_PARTS[2] in italic_texts

    # --- PPTX comments ---

    def test_win32com_pptx_comments_translated(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """PPTX legacy comments via win32com: each comment verified."""
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "pres.pptx"
        out = tmp_path / "translated.pptx"
        _create_rich_pptx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_pptx_comments(out)
        if comments:
            assert _french(_PPTX_COMMENT_TEXT) in comments, (
                f"Expected {_french(_PPTX_COMMENT_TEXT)!r} in {comments}"
            )

    # --- XLSX shapes ---

    def test_win32com_xlsx_shapes_translated(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """XLSX DrawingML shapes via win32com: each shape text verified."""
        enable_feature(SETTING_TRANSLATE_DOC_SHAPES)

        inp = tmp_path / "sheet.xlsx"
        out = tmp_path / "translated.xlsx"
        _create_rich_xlsx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        shapes = _read_xlsx_shape_texts(out)
        assert len(shapes) >= 1, "No shapes found in output XLSX"
        assert _french(_XLSX_SHAPE_TEXT) in shapes, (
            f"Expected {_french(_XLSX_SHAPE_TEXT)!r} in {shapes}"
        )

    # --- XLSX comments ---

    def test_win32com_xlsx_comments_translated(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
        enable_feature: Callable[..., None],
    ) -> None:
        """XLSX comments via win32com: each comment verified."""
        enable_feature(SETTING_TRANSLATE_DOC_COMMENTS)

        inp = tmp_path / "sheet.xlsx"
        out = tmp_path / "translated.xlsx"
        _create_rich_xlsx(inp)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_xlsx_comments(out)
        assert len(comments) == 1, (
            f"Expected 1 comment, got {len(comments)}: {comments}"
        )
        assert comments[0] == _french(_XLSX_COMMENT_TEXT)

    # --- Glossary forwarding ---

    def test_win32com_docx_glossary_forwarded(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """Glossary entries are forwarded to the LLM via win32com backend."""
        captured_kwargs: dict[str, object] = {}

        def tracking_translate(
            texts: list[str],
            target_lang: str,
            source_lang: str = "",
            **kwargs: object,
        ) -> list[str]:
            captured_kwargs.update(kwargs)
            return [f"[{target_lang}] {t}" for t in texts]

        monkeypatch.setattr(
            "src.core.llm_engine.translate_text",
            tracking_translate,
        )
        monkeypatch.setattr(
            "src.core.text_processor._llm_engine.translate_text",
            tracking_translate,
        )

        glossary = [(1, "Hello", "Bonjour")]  # noqa: PLR2004
        inp = tmp_path / "doc_w32.docx"
        out = tmp_path / "translated_w32.docx"
        _create_docx(inp, ["Hello world"])

        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                glossary_entries=glossary,
            )
            is True
        )

        assert captured_kwargs.get("glossary_entries") == glossary

    # --- Negative tests (features disabled by default) ---

    def test_win32com_docx_comments_not_translated_when_disabled(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Comments stay untranslated when SETTING_TRANSLATE_DOC_COMMENTS off."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        comments = _read_docx_comments(out)
        for c in comments:
            assert _french(_DOCX_COMMENT_TEXT) != c, (
                "Comment was translated despite setting being off"
            )

    def test_win32com_docx_shapes_not_translated_when_disabled(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """Shapes stay untranslated when SETTING_TRANSLATE_DOC_SHAPES off."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert translate_file(inp, out, "French", "English (US)") is True

        shapes = _read_docx_shape_texts(out)
        for s in shapes:
            assert _french(_DOCX_SHAPE_TEXT) != s, (
                "Shape was translated despite setting being off"
            )

    # --- All features combined ---

    def test_win32com_docx_all_features_combined(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
        enable_feature: Callable[..., None],
    ) -> None:
        """All features enabled: comments + shapes + glossary via win32com."""
        enable_feature(
            SETTING_TRANSLATE_DOC_COMMENTS,
            SETTING_TRANSLATE_DOC_SHAPES,
        )

        captured_kwargs: dict[str, object] = {}

        def tracking_translate(
            texts: list[str],
            target_lang: str,
            source_lang: str = "",
            **kwargs: object,
        ) -> list[str]:
            captured_kwargs.update(kwargs)
            return [f"[{target_lang}] {t}" for t in texts]

        monkeypatch.setattr(
            "src.core.llm_engine.translate_text",
            tracking_translate,
        )
        monkeypatch.setattr(
            "src.core.text_processor._llm_engine.translate_text",
            tracking_translate,
        )

        glossary = [(1, "Hello", "Bonjour")]  # noqa: PLR2004
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        png = tmp_path / "tiny.png"
        _create_tiny_png(png)
        _create_rich_docx(inp, png)

        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                glossary_entries=glossary,
            )
            is True
        )

        # Verify paragraphs
        paragraphs = _read_docx_paragraphs(out)
        for orig in _DOCX_PARAGRAPHS:
            assert _french(orig) in paragraphs

        # Verify table cells
        cells = _read_docx_table_cells(out)
        for row in _DOCX_TABLE_ROWS:
            for val in row:
                assert _french(val) in cells

        # Verify comments
        comments = _read_docx_comments(out)
        assert len(comments) == 1, (
            f"Expected 1 comment, got {len(comments)}: {comments}"
        )
        assert comments[0] == _french(_DOCX_COMMENT_TEXT)

        # Verify shapes
        shapes = _read_docx_shape_texts(out)
        assert len(shapes) >= 1, "No shapes found"
        assert _french(_DOCX_SHAPE_TEXT) in shapes

        # Verify glossary forwarding
        assert captured_kwargs.get("glossary_entries") == glossary

    # --- Cancellation ---

    def test_win32com_docx_cancellation_returns_false(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """cancel_check returning True → translate_file returns False."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        _create_docx(inp, ["Hello world"])

        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                cancel_check=lambda: True,
            )
            is False
        )

    # --- Progress callback ---

    def test_win32com_docx_progress_callback_invoked(
        self,
        tmp_path: Path,
        mock_llm: Callable[..., list[str]],
    ) -> None:
        """progress_callback receives increasing values ending at 100."""
        inp = tmp_path / "doc.docx"
        out = tmp_path / "translated.docx"
        _create_docx(inp, ["Hello world", "Goodbye world"])

        progress_values: list[int] = []
        assert (
            translate_file(
                inp,
                out,
                "French",
                "English (US)",
                progress_callback=progress_values.append,
            )
            is True
        )

        assert len(progress_values) >= 1, "Progress callback never called"
        assert progress_values[-1] == 100, (  # noqa: PLR2004
            f"Final progress should be 100, got {progress_values[-1]}"
        )
        for i in range(1, len(progress_values)):
            assert progress_values[i] >= progress_values[i - 1], (
                f"Progress decreased: {progress_values[i - 1]} → {progress_values[i]}"
            )

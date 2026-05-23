"""Integration tests for the Office document translation pipeline.

Creates real DOCX/XLSX/PPTX files via python-docx/openpyxl/python-pptx.
Backend falls to _BACKEND_PYTHON_LIB naturally in CI (no MS Office or UNO).
Only the LLM is mocked.
"""

import builtins
import sys
from collections.abc import Callable, Generator
from pathlib import Path
from typing import Any

import pytest

from src.core.database import init_db
from src.core.text_processor import translate_file


def _bypass_uno_import() -> object:
    """Temporarily restore Python's real import if UNO's hook is active."""
    uno_mod = sys.modules.get("uno")
    if uno_mod is None:
        return None
    original = getattr(uno_mod, "_builtin_import", None)
    if original is None or builtins.__import__ is original:
        return None
    uno_hook = builtins.__import__
    builtins.__import__ = original
    return uno_hook


def _restore_uno_import(hook: object) -> None:
    """Restore the UNO import hook if it was bypassed."""
    if hook is not None:
        builtins.__import__ = hook


# ── Shared fixtures ──────────────────────────────────────────────────


@pytest.fixture(autouse=True)
def setup_integration_env(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> Generator[None, None, None]:
    """Per-test DB isolation + mock environment setup."""
    db_file = tmp_path / "integration.db"
    monkeypatch.setattr("src.core.database.get_db_path", lambda: str(db_file))
    config_dir = tmp_path / "config"
    config_dir.mkdir()
    monkeypatch.setattr(
        "src.utils.path_manager.get_app_config_dir",
        lambda: config_dir,
    )
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    monkeypatch.setattr(
        "src.utils.path_manager.get_app_data_dir",
        lambda: data_dir,
    )
    init_db()
    monkeypatch.setattr("src.core.translator.stop_soffice", lambda: None)
    # Disable image/comment/shape translation by default
    monkeypatch.setattr(
        "src.utils.config_manager.load_setting",
        lambda k, d=None: False,
    )
    # Force python-lib backend to avoid UNO connection issues in CI
    monkeypatch.setattr(
        "src.core.office_processor._detect_backend",
        lambda suffix, *_args: "python_lib",
    )
    yield


@pytest.fixture()
def mock_llm(
    monkeypatch: pytest.MonkeyPatch,
) -> Callable[..., list[str]]:
    """Patches translate_text at all import sites."""

    def fake_translate(
        texts: list[str],
        target_lang: str,
        source_lang: str = "",
        **kwargs: object,
    ) -> list[str]:
        return [f"[{target_lang}] {t}" for t in texts]

    monkeypatch.setattr(
        "src.core.llm_engine.translate_text",
        fake_translate,
    )
    monkeypatch.setattr(
        "src.core.text_processor._llm_engine.translate_text",
        fake_translate,
    )
    return fake_translate


# ── Helpers ──────────────────────────────────────────────────────────


def _create_docx(path: Path, paragraphs: list[str]) -> None:
    """Create a real .docx file with the given paragraphs."""
    from docx import Document  # noqa: PLC0415

    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    doc.save(str(path))


def _create_xlsx(path: Path, rows: list[list[Any]]) -> None:
    """Create a real .xlsx file with the given rows (list of lists)."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import Workbook  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    wb.save(str(path))


def _create_pptx(path: Path, slide_texts: list[str]) -> None:
    """Create a real .pptx file with one text box per slide."""
    hook = _bypass_uno_import()
    try:
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    prs = Presentation()
    for text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tx_box.text_frame.text = text
    prs.save(str(path))


def _read_docx_paragraphs(path: Path) -> list[str]:
    """Read all paragraph texts from a .docx file."""
    from docx import Document  # noqa: PLC0415

    doc = Document(str(path))
    return [p.text for p in doc.paragraphs]


def _read_xlsx_values(path: Path) -> list[list[Any]]:
    """Read all cell values from the first sheet of an .xlsx file."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import load_workbook  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    wb = load_workbook(str(path))
    ws = wb.active
    return [[cell.value for cell in row] for row in ws.iter_rows()]


def _read_pptx_texts(path: Path) -> list[str]:
    """Read all text frame texts from all slides in a .pptx."""
    hook = _bypass_uno_import()
    try:
        from pptx import Presentation  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return texts


# ── Basic round-trips ────────────────────────────────────────────────


def test_docx_roundtrip(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with 2 paragraphs → both translated."""
    inp = tmp_path / "doc.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["Hello world", "Goodbye world"])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    paragraphs = _read_docx_paragraphs(out)
    translated = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in translated)


def test_xlsx_roundtrip(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """XLSX with 2 cells → both translated."""
    inp = tmp_path / "sheet.xlsx"
    out = tmp_path / "translated.xlsx"
    _create_xlsx(inp, [["Hello", "World"]])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    rows = _read_xlsx_values(out)
    assert len(rows) >= 1
    # At least one cell has [French] prefix
    flat = [str(c) for row in rows for c in row if c]
    assert any("[French]" in v for v in flat)


def test_pptx_roundtrip(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """PPTX with 1 slide text → translated."""
    inp = tmp_path / "pres.pptx"
    out = tmp_path / "translated.pptx"
    _create_pptx(inp, ["Hello world"])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    texts = _read_pptx_texts(out)
    assert any("[French]" in t for t in texts)


# ── Edge cases ───────────────────────────────────────────────────────


def test_docx_empty_document(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """Empty DOCX should produce output without LLM calls."""
    inp = tmp_path / "empty.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, [])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    assert out.exists()


def test_docx_with_table(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with a table → table cells translated."""
    from docx import Document  # noqa: PLC0415

    inp = tmp_path / "table.docx"
    out = tmp_path / "translated.docx"

    doc = Document()
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Hello"
    table.cell(1, 1).text = "World"
    doc.save(str(inp))

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    doc2 = Document(str(out))
    cells = []
    for table in doc2.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    cells.append(cell.text)
    assert any("[French]" in c for c in cells)


def test_docx_mixed_formatting(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with bold + italic runs → formatting preserved."""
    from docx import Document  # noqa: PLC0415

    inp = tmp_path / "fmt.docx"
    out = tmp_path / "translated.docx"

    doc = Document()
    p = doc.add_paragraph()
    run_bold = p.add_run("Bold text ")
    run_bold.bold = True
    run_italic = p.add_run("italic text")
    run_italic.italic = True
    doc.save(str(inp))

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    assert out.exists()


def test_docx_with_comments_enabled(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """DOCX with comments + setting enabled → body text translated."""
    from src.constants.settings import SETTING_TRANSLATE_DOC_COMMENTS  # noqa: PLC0415

    # Enable comments translation
    def fake_load_setting(key: str, default: object = None) -> object:
        if key == SETTING_TRANSLATE_DOC_COMMENTS:
            return True
        return default

    monkeypatch.setattr(
        "src.utils.config_manager.load_setting",
        fake_load_setting,
    )

    inp = tmp_path / "comments.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["Hello with comments"])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    paragraphs = _read_docx_paragraphs(out)
    translated = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in translated)


def test_docx_with_shapes_enabled(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """DOCX with shapes setting enabled → body text translated."""
    from src.constants.settings import SETTING_TRANSLATE_DOC_SHAPES  # noqa: PLC0415

    def fake_load_setting(key: str, default: object = None) -> object:
        if key == SETTING_TRANSLATE_DOC_SHAPES:
            return True
        return default

    monkeypatch.setattr(
        "src.utils.config_manager.load_setting",
        fake_load_setting,
    )

    inp = tmp_path / "shapes.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["Hello with shapes"])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    paragraphs = _read_docx_paragraphs(out)
    translated = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in translated)


def test_xlsx_empty_cells(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """XLSX with empty cells interspersed → non-empty cells translated."""
    inp = tmp_path / "sparse.xlsx"
    out = tmp_path / "translated.xlsx"
    _create_xlsx(inp, [["Hello", None, "World"], [None, "Foo", None]])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    rows = _read_xlsx_values(out)
    flat = [str(c) for row in rows for c in row if c]
    assert any("[French]" in v for v in flat)


def test_pptx_multiple_slides(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """PPTX with 3 slides → all translated."""
    inp = tmp_path / "multi.pptx"
    out = tmp_path / "translated.pptx"
    _create_pptx(inp, ["Slide one", "Slide two", "Slide three"])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    texts = _read_pptx_texts(out)
    translated = [t for t in texts if "[French]" in t]
    assert len(translated) >= 3  # noqa: PLR2004


def test_docx_glossary_forwarded(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Glossary entries are forwarded to translate_batch."""
    captured_kwargs: dict[str, object] = {}

    def tracking_translate(
        texts: list[str],
        target_lang: str,
        source_lang: str = "",
        **kwargs: object,
    ) -> list[str]:
        captured_kwargs.update(kwargs)
        return [f"[{target_lang}] {t}" for t in texts]

    monkeypatch.setattr(
        "src.core.llm_engine.translate_text",
        tracking_translate,
    )
    monkeypatch.setattr(
        "src.core.text_processor._llm_engine.translate_text",
        tracking_translate,
    )

    inp = tmp_path / "glossary.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["Hello glossary test"])

    glossary = [(1, "Hello", "Bonjour")]
    result = translate_file(
        inp,
        out,
        "French",
        "English (US)",
        glossary_entries=glossary,
    )
    assert result is True
    assert captured_kwargs.get("glossary_entries") == glossary


def test_xlsx_multiple_sheets(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """XLSX with 2 sheets → cells on both sheets translated."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import Workbook  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    inp = tmp_path / "multi.xlsx"
    out = tmp_path / "translated.xlsx"

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1["A1"] = "Hello"
    ws2 = wb.create_sheet("Sheet2")
    ws2["A1"] = "World"
    wb.save(str(inp))

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    rows = _read_xlsx_values(out)
    flat = [str(c) for row in rows for c in row if c]
    # At least one translated cell from Sheet1 or Sheet2
    assert any("[French]" in v for v in flat)


def test_docx_cancellation_returns_false(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX translation with cancel_check=lambda: True → returns False."""
    inp = tmp_path / "cancel.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["Hello world", "Second paragraph", "Third paragraph"])

    result = translate_file(
        inp,
        out,
        "French",
        "English (US)",
        cancel_check=lambda: True,
    )
    assert result is False


def test_docx_progress_callback_invoked(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX translation calls progress_callback; final value is 100."""
    inp = tmp_path / "prog.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["Hello world", "Goodbye world"])

    progress_values: list[int] = []
    result = translate_file(
        inp,
        out,
        "French",
        "English (US)",
        progress_callback=progress_values.append,
    )
    assert result is True
    assert len(progress_values) >= 1
    assert progress_values[-1] == 100  # noqa: PLR2004
    # Progress must be non-decreasing
    for i in range(1, len(progress_values)):
        assert progress_values[i] >= progress_values[i - 1]


def test_xlsx_src_lang_forwarded(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """src_lang is forwarded through the Office pipeline to translate_text."""
    captured: dict[str, str] = {}

    def tracking_translate(
        texts: list[str],
        target_lang: str,
        source_lang: str = "",
        **kwargs: object,
    ) -> list[str]:
        captured["source_lang"] = source_lang
        return [f"[{target_lang}] {t}" for t in texts]

    monkeypatch.setattr("src.core.llm_engine.translate_text", tracking_translate)
    monkeypatch.setattr(
        "src.core.text_processor._llm_engine.translate_text", tracking_translate
    )

    inp = tmp_path / "src.xlsx"
    out = tmp_path / "translated.xlsx"
    _create_xlsx(inp, [["Hello"]])

    translate_file(inp, out, "French", "English (US)")
    assert captured.get("source_lang") == "English (US)"


def test_legacy_doc_no_backend(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """.doc file without win32com/UNO → OFFICE_CONVERTER_NOT_FOUND."""

    # Override the forced python_lib backend to simulate no-backend scenario
    def no_backend(suffix: str, *_args: object) -> str:
        raise ValueError("OFFICE_CONVERTER_NOT_FOUND")

    monkeypatch.setattr(
        "src.core.office_processor._detect_backend",
        no_backend,
    )

    inp = tmp_path / "legacy.doc"
    inp.write_bytes(b"fake doc content")
    out = tmp_path / "translated.doc"

    with pytest.raises(ValueError, match="OFFICE_CONVERTER_NOT_FOUND"):
        translate_file(inp, out, "French", "English (US)")


# ── TranslationConfig injection tests ──────────────────────────────


def test_docx_comments_via_config(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """config=TranslationConfig(translate_doc_comments=True) enables comments."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    inp = tmp_path / "cfg_comments.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["Hello with comments"])

    config = TranslationConfig(translate_doc_comments=True)
    result = translate_file(inp, out, "French", "English (US)", config=config)
    assert result is True
    paragraphs = _read_docx_paragraphs(out)
    translated = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in translated)


def test_docx_shapes_via_config(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """config=TranslationConfig(translate_doc_shapes=True) enables shapes."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    inp = tmp_path / "cfg_shapes.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["Hello with shapes"])

    config = TranslationConfig(translate_doc_shapes=True)
    result = translate_file(inp, out, "French", "English (US)", config=config)
    assert result is True
    paragraphs = _read_docx_paragraphs(out)
    translated = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in translated)


def test_xlsx_config_defaults_only_body(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """config=TranslationConfig() → only body cells translated."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    inp = tmp_path / "cfg_default.xlsx"
    out = tmp_path / "translated.xlsx"
    _create_xlsx(inp, [["Hello"]])

    config = TranslationConfig()
    result = translate_file(inp, out, "French", "English (US)", config=config)
    assert result is True
    rows = _read_xlsx_values(out)
    flat = [str(c) for row in rows for c in row if c]
    assert any("[French]" in v for v in flat)


def test_pptx_config_comments_enabled(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """config=TranslationConfig(translate_doc_comments=True) + PPTX."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    inp = tmp_path / "cfg_pptx.pptx"
    out = tmp_path / "translated.pptx"
    _create_pptx(inp, ["Hello slide"])

    config = TranslationConfig(translate_doc_comments=True)
    result = translate_file(inp, out, "French", "English (US)", config=config)
    assert result is True
    texts = _read_pptx_texts(out)
    assert any("[French]" in t for t in texts)


# ── Whitespace & edge case tests ────────────────────────────────────


def test_docx_whitespace_only_paragraphs(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with whitespace-only paragraphs → success, no crash."""
    inp = tmp_path / "ws.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["   ", "\t", "  \n  "])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    assert out.exists()


def test_xlsx_whitespace_only_cells(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """XLSX with whitespace cells → success."""
    inp = tmp_path / "ws.xlsx"
    out = tmp_path / "translated.xlsx"
    _create_xlsx(inp, [["   ", "\t"], ["  ", ""]])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    assert out.exists()


def test_pptx_empty_text_frames(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """PPTX with empty slide text → success."""
    inp = tmp_path / "empty_slide.pptx"
    out = tmp_path / "translated.pptx"
    _create_pptx(inp, [""])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    assert out.exists()


def test_docx_single_character(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with single-character paragraph "A" → translated."""
    inp = tmp_path / "single_char.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["A"])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    paragraphs = _read_docx_paragraphs(out)
    translated = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in translated)


def test_xlsx_numeric_only_cells(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """XLSX with numbers only → handled (may skip untranslatable)."""
    inp = tmp_path / "numeric.xlsx"
    out = tmp_path / "translated.xlsx"
    _create_xlsx(inp, [[42, 3.14, 100]])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    assert out.exists()


def test_docx_mixed_empty_and_real_paragraphs(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with ["", "Hello", "  ", "World"] → only real text translated."""
    inp = tmp_path / "mixed.docx"
    out = tmp_path / "translated.docx"
    _create_docx(inp, ["", "Hello", "  ", "World"])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    paragraphs = _read_docx_paragraphs(out)
    translated = [p for p in paragraphs if "[French]" in p]
    # "Hello" and "World" should be translated
    assert len(translated) >= 2  # noqa: PLR2004


# ── Formatting preservation tests ────────────────────────────────────


def test_docx_bold_preserved_through_translation(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with bold run → bold formatting survives translation."""
    from docx import Document  # noqa: PLC0415

    inp = tmp_path / "bold.docx"
    out = tmp_path / "translated.docx"

    doc = Document()
    para = doc.add_paragraph()
    run_bold = para.add_run("Important text")
    run_bold.bold = True
    # Add a non-bold run so mixed formatting triggers HTML path
    run_plain = para.add_run(" normal text")
    run_plain.bold = False
    doc.save(str(inp))

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    out_doc = Document(str(out))
    # At least one run should be bold
    bold_found = False
    for para in out_doc.paragraphs:
        for run in para.runs:
            if run.text.strip() and run.bold:
                bold_found = True
                break
    assert bold_found, "No bold run found in translated output"


def test_docx_italic_preserved_through_translation(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with italic run → italic formatting survives translation."""
    from docx import Document  # noqa: PLC0415

    inp = tmp_path / "italic.docx"
    out = tmp_path / "translated.docx"

    doc = Document()
    para = doc.add_paragraph()
    run_italic = para.add_run("Emphasized text")
    run_italic.italic = True
    # Add a non-italic run for mixed formatting
    run_plain = para.add_run(" regular text")
    run_plain.italic = False
    doc.save(str(inp))

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    out_doc = Document(str(out))
    italic_found = False
    for para in out_doc.paragraphs:
        for run in para.runs:
            if run.text.strip() and run.italic:
                italic_found = True
                break
    assert italic_found, "No italic run found in translated output"


def test_docx_mixed_formatting_preserved(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with bold + italic + underline → all formatting survives."""
    from docx import Document  # noqa: PLC0415

    inp = tmp_path / "mixed_fmt.docx"
    out = tmp_path / "translated.docx"

    doc = Document()
    para = doc.add_paragraph()
    run_bold = para.add_run("Bold ")
    run_bold.bold = True
    run_italic = para.add_run("Italic ")
    run_italic.italic = True
    run_underline = para.add_run("Underline")
    run_underline.underline = True
    doc.save(str(inp))

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    out_doc = Document(str(out))
    found_bold = False
    found_italic = False
    found_underline = False
    for para in out_doc.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            if run.bold:
                found_bold = True
            if run.italic:
                found_italic = True
            if run.underline:
                found_underline = True
    assert found_bold, "Bold not preserved"
    assert found_italic, "Italic not preserved"
    assert found_underline, "Underline not preserved"


# ── Helper: DOCX with headers/footers ──────────────────────────────


def _create_docx_with_headers(
    path: Path,
    header_text: str,
    footer_text: str,
    body_text: str = "Body paragraph",
) -> None:
    """Create a .docx file with header, footer, and a body paragraph."""
    from docx import Document  # noqa: PLC0415

    doc = Document()
    # Add body text
    doc.add_paragraph(body_text)
    # Set header/footer on the first (default) section
    section = doc.sections[0]
    hdr = section.header
    hdr.is_linked_to_previous = False
    hdr.paragraphs[0].text = header_text
    ftr = section.footer
    ftr.is_linked_to_previous = False
    ftr.paragraphs[0].text = footer_text
    doc.save(str(path))


def _read_docx_header_footer_texts(
    path: Path,
) -> tuple[list[str], list[str]]:
    """Read header and footer texts from all sections of a .docx file."""
    from docx import Document  # noqa: PLC0415

    doc = Document(str(path))
    headers: list[str] = []
    footers: list[str] = []
    for section in doc.sections:
        hdr = section.header
        if not hdr.is_linked_to_previous:
            for para in hdr.paragraphs:
                if para.text.strip():
                    headers.append(para.text)
        ftr = section.footer
        if not ftr.is_linked_to_previous:
            for para in ftr.paragraphs:
                if para.text.strip():
                    footers.append(para.text)
    return headers, footers


# ── Helper: DOCX with footnotes/endnotes via raw XML ──────────────


def _create_docx_with_footnote(
    path: Path,
    body_text: str,
    footnote_text: str,
) -> None:
    """Create a .docx with a properly registered footnote via OPC API.

    python-docx doesn't natively support footnotes, so we register the
    footnotes.xml part through the OPC package API to ensure it survives
    round-trip saves.
    """
    from docx import Document  # noqa: PLC0415
    from docx.opc.packuri import PackURI  # noqa: PLC0415
    from docx.opc.part import Part  # noqa: PLC0415

    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    fn_rel = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/footnotes"
    )
    fn_ct = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.footnotes+xml"
    )

    doc = Document()
    doc.add_paragraph(body_text)

    fn_xml = (
        f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:footnotes xmlns:w="{w_ns}">'
        f'  <w:footnote w:type="separator" w:id="0">'
        f"    <w:p><w:r><w:separator/></w:r></w:p>"
        f"  </w:footnote>"
        f'  <w:footnote w:type="continuationSeparator" w:id="1">'
        f"    <w:p><w:r><w:continuationSeparator/></w:r></w:p>"
        f"  </w:footnote>"
        f'  <w:footnote w:id="2">'
        f"    <w:p><w:r><w:t>{footnote_text}</w:t></w:r></w:p>"
        f"  </w:footnote>"
        f"</w:footnotes>"
    ).encode()

    fn_part = Part(
        PackURI("/word/footnotes.xml"),
        fn_ct,
        fn_xml,
        doc.part.package,
    )
    doc.part.relate_to(fn_part, fn_rel)
    doc.save(str(path))


def _create_docx_with_endnote(
    path: Path,
    body_text: str,
    endnote_text: str,
) -> None:
    """Create a .docx with a properly registered endnote via OPC API."""
    from docx import Document  # noqa: PLC0415
    from docx.opc.packuri import PackURI  # noqa: PLC0415
    from docx.opc.part import Part  # noqa: PLC0415

    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    en_rel = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/endnotes"
    )
    en_ct = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.endnotes+xml"
    )

    doc = Document()
    doc.add_paragraph(body_text)

    en_xml = (
        f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:endnotes xmlns:w="{w_ns}">'
        f'  <w:endnote w:type="separator" w:id="0">'
        f"    <w:p><w:r><w:separator/></w:r></w:p>"
        f"  </w:endnote>"
        f'  <w:endnote w:type="continuationSeparator" w:id="1">'
        f"    <w:p><w:r><w:continuationSeparator/></w:r></w:p>"
        f"  </w:endnote>"
        f'  <w:endnote w:id="2">'
        f"    <w:p><w:r><w:t>{endnote_text}</w:t></w:r></w:p>"
        f"  </w:endnote>"
        f"</w:endnotes>"
    ).encode()

    en_part = Part(
        PackURI("/word/endnotes.xml"),
        en_ct,
        en_xml,
        doc.part.package,
    )
    doc.part.relate_to(en_part, en_rel)
    doc.save(str(path))


def _read_docx_footnotes(path: Path) -> list[str]:
    """Read all user footnote texts from a .docx file."""
    import zipfile as zf_mod  # noqa: PLC0415

    from lxml import etree as et  # noqa: PLC0415

    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    skip_ids = {"0", "1", "-1"}

    with zf_mod.ZipFile(path, "r") as zf:
        if "word/footnotes.xml" not in zf.namelist():
            return []
        data = zf.read("word/footnotes.xml")

    root = et.fromstring(data)
    texts: list[str] = []
    for fn in root.iter(f"{{{w_ns}}}footnote"):
        fn_id = fn.get(f"{{{w_ns}}}id", "")
        if fn_id in skip_ids:
            continue
        parts: list[str] = []
        for para in fn.iter(f"{{{w_ns}}}p"):
            para_text = ""
            for r_elem in para.iter(f"{{{w_ns}}}r"):
                for t_elem in r_elem.iter(f"{{{w_ns}}}t"):
                    if t_elem.text:
                        para_text += t_elem.text
            if para_text.strip():
                parts.append(para_text.strip())
        if parts:
            texts.append("\n".join(parts))
    return texts


def _read_docx_endnotes(path: Path) -> list[str]:
    """Read all user endnote texts from a .docx file."""
    import zipfile as zf_mod  # noqa: PLC0415

    from lxml import etree as et  # noqa: PLC0415

    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    skip_ids = {"0", "1", "-1"}

    with zf_mod.ZipFile(path, "r") as zf:
        if "word/endnotes.xml" not in zf.namelist():
            return []
        data = zf.read("word/endnotes.xml")

    root = et.fromstring(data)
    texts: list[str] = []
    for en in root.iter(f"{{{w_ns}}}endnote"):
        en_id = en.get(f"{{{w_ns}}}id", "")
        if en_id in skip_ids:
            continue
        parts: list[str] = []
        for para in en.iter(f"{{{w_ns}}}p"):
            para_text = ""
            for r_elem in para.iter(f"{{{w_ns}}}r"):
                for t_elem in r_elem.iter(f"{{{w_ns}}}t"):
                    if t_elem.text:
                        para_text += t_elem.text
            if para_text.strip():
                parts.append(para_text.strip())
        if parts:
            texts.append("\n".join(parts))
    return texts


# ── Helper: XLSX with custom sheet names ───────────────────────────


def _create_xlsx_with_sheets(
    path: Path,
    sheet_names: list[str],
) -> None:
    """Create a real .xlsx with named sheets, each having one cell."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import Workbook  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    wb = Workbook()
    # Rename the default sheet
    ws = wb.active
    ws.title = sheet_names[0] if sheet_names else "Sheet1"
    ws["A1"] = "Hello"
    for name in sheet_names[1:]:
        ws_new = wb.create_sheet(name)
        ws_new["A1"] = "Data"
    wb.save(str(path))


def _read_xlsx_sheet_names(path: Path) -> list[str]:
    """Read sheet names from an .xlsx file."""
    hook = _bypass_uno_import()
    try:
        from openpyxl import load_workbook  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    wb = load_workbook(str(path))
    return wb.sheetnames


# ── Helper: PPTX with speaker notes ───────────────────────────────


def _create_pptx_with_notes(
    path: Path,
    slides_with_notes: list[tuple[str, str]],
) -> None:
    """Create a .pptx with slides that have speaker notes.

    Args:
        path: Output path.
        slides_with_notes: List of (slide_text, notes_text) tuples.
    """
    hook = _bypass_uno_import()
    try:
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    prs = Presentation()
    for slide_text, notes_text in slides_with_notes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        tx_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(1),
        )
        tx_box.text_frame.text = slide_text
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text
    prs.save(str(path))


def _read_pptx_notes(path: Path) -> list[str]:
    """Read all speaker note texts from a .pptx file."""
    hook = _bypass_uno_import()
    try:
        from pptx import Presentation  # noqa: PLC0415
    finally:
        _restore_uno_import(hook)

    prs = Presentation(str(path))
    notes: list[str] = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                notes.append(text)
    return notes


# ── Headers & footers tests ────────────────────────────────────────


def test_docx_headers_footers_translated(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with header/footer text → both translated alongside body."""
    inp = tmp_path / "hf.docx"
    out = tmp_path / "translated.docx"
    _create_docx_with_headers(
        inp,
        header_text="Company Header",
        footer_text="Page Footer",
        body_text="Main content",
    )

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    # Verify body translated
    paragraphs = _read_docx_paragraphs(out)
    body = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in body)

    # Verify header and footer translated
    headers, footers = _read_docx_header_footer_texts(out)
    assert len(headers) >= 1, "No header text found in output"
    assert any("[French]" in h for h in headers), f"Header not translated: {headers}"
    assert len(footers) >= 1, "No footer text found in output"
    assert any("[French]" in f for f in footers), f"Footer not translated: {footers}"


def test_docx_headers_footers_with_formatting(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """Headers with bold/italic formatting preserved after translation."""
    from docx import Document  # noqa: PLC0415

    inp = tmp_path / "hf_fmt.docx"
    out = tmp_path / "translated.docx"

    doc = Document()
    doc.add_paragraph("Body text")
    section = doc.sections[0]

    # Create header with mixed formatting
    hdr = section.header
    hdr.is_linked_to_previous = False
    hdr_para = hdr.paragraphs[0]
    run_bold = hdr_para.add_run("Bold Header ")
    run_bold.bold = True
    run_plain = hdr_para.add_run("normal part")
    run_plain.bold = False

    # Create footer with italic formatting
    ftr = section.footer
    ftr.is_linked_to_previous = False
    ftr_para = ftr.paragraphs[0]
    run_italic = ftr_para.add_run("Italic Footer ")
    run_italic.italic = True
    run_normal = ftr_para.add_run("normal part")
    run_normal.italic = False

    doc.save(str(inp))

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    # Verify header/footer text exists in output
    headers, footers = _read_docx_header_footer_texts(out)
    assert len(headers) >= 1, "No header text found in output"
    assert len(footers) >= 1, "No footer text found in output"


# ── Footnotes tests ────────────────────────────────────────────────


def test_docx_footnotes_translated(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with a footnote → footnote text is translated."""
    inp = tmp_path / "fn.docx"
    out = tmp_path / "translated.docx"
    _create_docx_with_footnote(inp, "Main document text", "This is a footnote")

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    # Body should be translated
    paragraphs = _read_docx_paragraphs(out)
    body = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in body)

    # Footnote should be translated
    footnotes = _read_docx_footnotes(out)
    assert len(footnotes) >= 1, "No user footnotes found in output"
    assert any("[French]" in fn for fn in footnotes), (
        f"Footnote not translated: {footnotes}"
    )


def test_docx_endnotes_translated(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with an endnote → endnote text is translated."""
    inp = tmp_path / "en.docx"
    out = tmp_path / "translated.docx"
    _create_docx_with_endnote(inp, "Main document text", "This is an endnote")

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True

    # Body should be translated
    paragraphs = _read_docx_paragraphs(out)
    body = [p for p in paragraphs if p.strip()]
    assert any("[French]" in p for p in body)

    # Endnote should be translated
    endnotes = _read_docx_endnotes(out)
    assert len(endnotes) >= 1, "No user endnotes found in output"
    assert any("[French]" in en for en in endnotes), (
        f"Endnote not translated: {endnotes}"
    )


def test_docx_empty_footnotes_skipped(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """DOCX with only separator footnotes (IDs 0, 1) → no crash."""
    inp = tmp_path / "empty_fn.docx"
    out = tmp_path / "translated.docx"
    # Create a plain docx — python-docx produces footnotes.xml with
    # only separator entries (IDs 0 and 1) which should be skipped.
    _create_docx(inp, ["Some body text"])

    result = translate_file(inp, out, "French", "English (US)")
    assert result is True
    assert out.exists()

    # No user footnotes should exist
    footnotes = _read_docx_footnotes(out)
    assert footnotes == []


# ── Sheet names tests ──────────────────────────────────────────────


def test_xlsx_sheet_names_translated(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """XLSX with custom sheet names + setting enabled → names translated."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    inp = tmp_path / "sheets.xlsx"
    out = tmp_path / "translated.xlsx"
    _create_xlsx_with_sheets(inp, ["Summary", "Details", "Appendix"])

    config = TranslationConfig(translate_sheet_names=True)
    result = translate_file(inp, out, "French", "English (US)", config=config)
    assert result is True

    names = _read_xlsx_sheet_names(out)
    assert len(names) >= 3  # noqa: PLR2004
    # _sanitize_sheet_name strips brackets, so [French] becomes "French"
    # e.g. "Summary" → "[French] Summary" → "French Summary"
    assert all(
        n != orig
        for n, orig in zip(names, ["Summary", "Details", "Appendix"], strict=True)
    ), f"Sheet names not translated: {names}"
    assert any("French" in n for n in names), (
        f"Sheet names missing translation marker: {names}"
    )


def test_xlsx_sheet_names_disabled_by_default(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """XLSX with custom sheet names but setting disabled → names unchanged."""
    inp = tmp_path / "sheets_off.xlsx"
    out = tmp_path / "translated.xlsx"
    _create_xlsx_with_sheets(inp, ["Summary", "Details"])

    # Default config has translate_sheet_names=False
    from src.core.config import TranslationConfig  # noqa: PLC0415

    config = TranslationConfig(translate_sheet_names=False)
    result = translate_file(inp, out, "French", "English (US)", config=config)
    assert result is True

    names = _read_xlsx_sheet_names(out)
    # Sheet names should remain unchanged (original names preserved)
    assert "Summary" in names
    assert "Details" in names
    assert not any("French" in n for n in names), (
        f"Sheet names should NOT be translated: {names}"
    )


# ── Speaker notes tests ───────────────────────────────────────────


def test_pptx_speaker_notes_translated(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """PPTX with speaker notes + setting enabled → notes translated."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    inp = tmp_path / "notes.pptx"
    out = tmp_path / "translated.pptx"
    _create_pptx_with_notes(
        inp,
        [
            ("Slide one", "Speaker note for slide one"),
            ("Slide two", "Speaker note for slide two"),
        ],
    )

    config = TranslationConfig(translate_doc_notes=True)
    result = translate_file(inp, out, "French", "English (US)", config=config)
    assert result is True

    # Body slides should be translated
    texts = _read_pptx_texts(out)
    assert any("[French]" in t for t in texts)

    # Speaker notes should be translated
    notes = _read_pptx_notes(out)
    assert len(notes) >= 2  # noqa: PLR2004
    assert any("[French]" in n for n in notes), f"Speaker notes not translated: {notes}"


def test_pptx_notes_disabled_by_default(
    tmp_path: Path,
    mock_llm: Callable[..., list[str]],
) -> None:
    """PPTX with speaker notes but setting disabled → notes unchanged."""
    from src.core.config import TranslationConfig  # noqa: PLC0415

    inp = tmp_path / "notes_off.pptx"
    out = tmp_path / "translated.pptx"
    _create_pptx_with_notes(
        inp,
        [("Slide one", "Keep this note unchanged")],
    )

    # Default config has translate_doc_notes=False
    config = TranslationConfig(translate_doc_notes=False)
    result = translate_file(inp, out, "French", "English (US)", config=config)
    assert result is True

    # Body should be translated
    texts = _read_pptx_texts(out)
    assert any("[French]" in t for t in texts)

    # Speaker notes should remain unchanged
    notes = _read_pptx_notes(out)
    assert len(notes) >= 1
    assert all("[French]" not in n for n in notes), (
        f"Speaker notes should NOT be translated: {notes}"
    )

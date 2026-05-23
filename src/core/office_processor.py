"""Office document processing for DOCX, XLSX, PPTX, ODT, ODS, ODP and legacy formats.

Uses a 3-tier backend system:
  1. win32com (Windows + MS Office)
  2. LibreOffice UNO API (cross-platform)
  3. python-docx / openpyxl / python-pptx / odfpy (modern + ODF formats)

Legacy formats (.doc, .xls, .ppt) require backend 1 or 2.
"""

from __future__ import annotations

import contextlib
import html
import logging
import os
import shutil
import sys
import tempfile
import zipfile
from collections.abc import Callable, Generator
from copy import copy as _shallow_copy
from pathlib import Path
from typing import TYPE_CHECKING, Any

from lxml import etree

if TYPE_CHECKING:
    from src.core.config import TranslationConfig

from src.constants.errors import base_error_tag
from src.constants.llm import CONTENT_DATA_VALUES, CONTENT_HTML, CONTENT_PLAIN_TEXT
from src.core.checkpoint import (
    hash_office_image,
    load_office_image_checkpoint,
    save_office_image_checkpoint,
)
from src.core.llm_engine import translate_batch
from src.core.office_formatter import (
    _DRAWINGML_FORMAT_ATTRS,
    _DRAWINGML_NS,
    _FORMATTING_HTML_RE,
    _HEX_TO_WD_HIGHLIGHT_INDEX,
    _HIGHLIGHT_COLORS,  # noqa: F401 (re-export for tests)
    _HYPERLINK_RELTYPE,
    _MC_ALTERNATE_CONTENT,  # noqa: F401 (re-export for tests)
    _ODF_NS,
    _PPTX_FORMAT_ATTRS,  # noqa: F401 (re-export for tests)
    _STRIP_FORMAT_TAGS_RE,
    _W_HYPERLINK_TAG,
    _WD_HIGHLIGHT_INDEX_TO_HEX,
    _WORDML_NS,
    _apply_drawingml_format_attrs,
    _apply_pptx_format_attrs,  # noqa: F401 (re-export for tests)
    _build_span_style,  # noqa: F401 (re-export for tests)
    _clear_run_text_only,  # noqa: F401 (re-export for tests)
    _color_hex_to_int,
    _color_hex_to_win32com,
    _drawingml_to_html,
    _FormattedSegment,
    _has_drawingml_hyperlinks,
    _has_drawingml_mixed_formatting,
    _has_mixed_formatting,
    _has_odf_text_box_mixed_formatting,
    _has_pptx_hyperlinks,
    _has_pptx_mixed_formatting,
    _inject_html_runs,
    _inject_pptx_html_runs,
    _int_to_color_hex,
    _odf_text_box_to_html,
    _para_has_hyperlinks,
    _parse_docx_char_styles,
    _parse_html_formatting,
    _parse_span_style,  # noqa: F401 (re-export for tests)
    _pptx_plain_fallback,  # noqa: F401 (re-export for tests)
    _pptx_runs_to_html,
    _read_docx_run_bg_hex,  # noqa: F401 (re-export for tests)
    _read_docx_run_color_hex,  # noqa: F401 (re-export for tests)
    _read_docx_run_size_pt,  # noqa: F401 (re-export for tests)
    _read_drawingml_rpr_formatting,  # noqa: F401 (re-export for tests)
    _read_odf_span_formatting,  # noqa: F401 (re-export for tests)
    _read_pptx_run_formatting,  # noqa: F401 (re-export for tests)
    _read_pptx_run_full_formatting,  # noqa: F401 (re-export for tests)
    _read_wml_rpr_formatting,
    _replace_paragraph_text,
    _run_has_visual_content,  # noqa: F401 (re-export for tests)
    _runs_to_html,
    _set_rpr_color,  # noqa: F401 (re-export for tests)
    _set_rpr_flag,  # noqa: F401 (re-export for tests)
    _set_rpr_font_size,  # noqa: F401 (re-export for tests)
    _set_rpr_underline,  # noqa: F401 (re-export for tests)
    _win32com_color_to_hex,
    _wrap_with_tags,
)
from src.core.office_lifecycle import (
    _APP_EXCEL,
    _APP_PPT,
    _APP_WORD,
    _ensure_soffice_running,  # noqa: F401 (re-export for tests)
    _find_soffice_binary,  # noqa: F401 (re-export for tests)
    _get_uno_desktop,
    _get_uno_search_paths,
    _win32com_close,
    _win32com_open,
    stop_soffice,  # noqa: F401 (re-export for tests)
)

logger = logging.getLogger("office_processor")

# OOXML formats handled natively by python-docx / openpyxl / python-pptx
_OOXML_EXTENSIONS = {".docx", ".xlsx", ".pptx"}

# ODF formats — best handled by UNO, fallback to odfpy (python_lib)
_ODF_EXTENSIONS = {".odt", ".ods", ".odp"}

# Extensions that can be handled by pure-Python libraries
_MODERN_EXTENSIONS = _OOXML_EXTENSIONS | _ODF_EXTENSIONS

# Legacy binary OLE2 formats that require win32com or UNO
_LEGACY_EXTENSIONS = {".doc", ".xls", ".ppt"}

# Mapping from legacy → modern format (used for both auto-conversion
# after translation and round-trip image conversion of legacy files)
LEGACY_CONVERT_MAP: dict[str, str] = {
    ".doc": ".docx",
    ".xls": ".xlsx",
    ".ppt": ".pptx",
}

# Mapping from ODF → modern format for auto-conversion after translation
ODF_CONVERT_MAP: dict[str, str] = {
    ".odt": ".docx",
    ".ods": ".xlsx",
    ".odp": ".pptx",
}

# UNO cell type constant for text cells
_UNO_CELL_TYPE_TEXT = 2

# DOM nodeType constants (used by odfpy element traversal)
_NODE_TYPE_ELEMENT = 1
_NODE_TYPE_TEXT = 3

# Backend identifiers
_BACKEND_WIN32COM = "win32com"
_BACKEND_UNO = "uno"
_BACKEND_PYTHON_LIB = "python_lib"

# ---------------------------------------------------------------------------
# Backend detection
# ---------------------------------------------------------------------------


def _detect_backend(suffix: str, libreoffice_path: str = "") -> str:
    """Detects the best available backend for the given file extension.

    Priority order depends on format family:
    - OOXML (.docx/.xlsx/.pptx): python_lib immediately (lightweight).
    - ODF (.odt/.ods/.odp): UNO → win32com → python_lib (odfpy).
    - Legacy (.doc/.xls/.ppt): win32com → UNO → error.

    Args:
        suffix: Lowercase file extension (e.g. ".docx", ".doc").
        libreoffice_path: User-configured LibreOffice path; forwarded
            to ``_get_uno_search_paths()``.

    Returns:
        str: One of the backend identifiers.

    Raises:
        ValueError: If no backend is available for the format.
    """
    # 1. OOXML — pure-Python libraries handle these natively
    if suffix in _OOXML_EXTENSIONS:
        return _BACKEND_PYTHON_LIB

    # 2. ODF — UNO → win32com → python_lib (odfpy).
    #    UNO is preferred because ODF is LibreOffice's native format (full
    #    spec compliance, per-run formatting preservation).  MS Office can
    #    open ODF but treats it as a foreign format (may mangle styles).
    #    odfpy is the lightweight last resort — preserves paragraph-level
    #    styles but loses per-run mixed formatting within paragraphs.
    if suffix in _ODF_EXTENSIONS:
        # Try UNO first (best ODF fidelity)
        for p in _get_uno_search_paths(libreoffice_path):
            if p not in sys.path and Path(p).is_dir():
                sys.path.append(p)
        try:
            import uno  # noqa: F401, PLC0415

            return _BACKEND_UNO
        except ImportError:
            pass

        # Try win32com
        try:
            import win32com.client  # noqa: F401, PLC0415

            return _BACKEND_WIN32COM
        except ImportError:
            pass

        # Fall back to odfpy via python_lib
        return _BACKEND_PYTHON_LIB

    # 3. Legacy — requires an external office suite
    try:
        import win32com.client  # noqa: F401, PLC0415

        return _BACKEND_WIN32COM
    except ImportError:
        pass

    for p in _get_uno_search_paths(libreoffice_path):
        if p not in sys.path and Path(p).is_dir():
            sys.path.append(p)
    try:
        import uno  # noqa: F401, PLC0415

        return _BACKEND_UNO
    except ImportError:
        pass

    raise ValueError("OFFICE_CONVERTER_NOT_FOUND")


# ---------------------------------------------------------------------------
# win32com backend — context manager & helpers
# ---------------------------------------------------------------------------

# Script detection and font substitution delegated to shared font utility.
from src.utils.font_utils import (  # noqa: E402, I001
    classify_generic_family as _classify_generic_family,
    detect_script as _detect_script_family,
    get_font_for_language as _get_font_for_language,
)


def _substitute_font(
    original_font: str,
    original_text: str,
    translated_text: str,
    target_lang: str = "",
) -> str | None:
    """Determines the font name to use after translation.

    When the original and translated texts share the same script family,
    the original font name is returned unchanged.  When scripts differ
    (e.g. Latin → CJK), a compatible font from the same generic family
    (serif / sans-serif / monospace) is selected for the target language.

    Args:
        original_font: The source document's font name.
        original_text: Text before translation.
        translated_text: Text after translation.
        target_lang: Target language name (used for font selection).

    Returns:
        The font name to apply, or ``None`` to clear the font (lets the
        application pick a default).
    """
    if not original_text or not translated_text:
        return original_font

    if _detect_script_family(original_text) == _detect_script_family(
        translated_text,
    ):
        return original_font

    # Scripts differ — select a compatible font for the target language
    if target_lang:
        generic = _classify_generic_family(font_name=original_font)
        return _get_font_for_language(target_lang, generic)

    # No target_lang available — clear font and let the app substitute
    return None


def _save_win32com_font(font_obj: object) -> dict[str, object]:
    """Saves font properties from a win32com Font object.

    Reads each property in WIN32COM_FONT_PROPERTIES and stores non-undefined
    values.  Properties that raise (e.g. on merged cells) are silently skipped.

    Args:
        font_obj: A win32com Range.Font COM object.

    Returns:
        dict: Mapping of property name to saved value.
    """
    from src.constants.office import (  # noqa: PLC0415
        WIN32COM_FONT_PROPERTIES,
        WIN32COM_UNDEFINED,
    )

    saved: dict[str, object] = {}
    for prop in WIN32COM_FONT_PROPERTIES:
        try:
            val = getattr(font_obj, prop)
            if val != WIN32COM_UNDEFINED:
                saved[prop] = val
        except Exception:  # noqa: BLE001
            continue
    return saved


def _restore_win32com_font(
    font_obj: object,
    saved: dict[str, object],
    *,
    original_text: str = "",
    translated_text: str = "",
    target_lang: str = "",
) -> None:
    """Restores previously saved font properties to a win32com Font object.

    Sets each property independently so a single failure does not prevent
    other properties from being restored.

    When *target_lang* is provided and ``"Name"`` is present in *saved*,
    the font name is substituted via :func:`_substitute_font` when the
    source and target scripts differ.

    Args:
        font_obj: A win32com Range.Font COM object.
        saved: Mapping of property name to value (from _save_win32com_font).
        original_text: The text before translation (for script detection).
        translated_text: The text after translation (for script detection).
        target_lang: Target language name for font substitution.
    """
    # Substitute font name when scripts are incompatible
    if "Name" in saved and original_text and translated_text:
        new_name = _substitute_font(
            str(saved["Name"]),
            original_text,
            translated_text,
            target_lang,
        )
        if new_name is None:
            saved = {k: v for k, v in saved.items() if k != "Name"}
        elif new_name != saved["Name"]:
            saved = {**saved, "Name": new_name}

    for prop, val in saved.items():
        try:
            setattr(font_obj, prop, val)
        except Exception:  # noqa: BLE001
            continue


# -- win32com text extract/inject functions ---------------------------------


def _read_win32com_char_formatting(
    char_range: object,
) -> tuple[bool, bool, bool, bool, bool, bool, float | None, str | None, str | None]:
    """Reads inline formatting from a single win32com Word character range.

    Args:
        char_range: A win32com ``Range`` object for a single character.

    Returns:
        Tuple of (bold, italic, underline, strike, superscript, subscript,
        font_size_pt, color_hex, bg_color_hex).
        Properties equal to ``WIN32COM_UNDEFINED`` are treated as False/None.
    """
    from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415

    font = char_range.Font
    # Bold / Italic — boolean-like; treat undefined as False
    raw_bold = font.Bold
    bold = bool(raw_bold) if raw_bold != WIN32COM_UNDEFINED else False
    raw_italic = font.Italic
    italic = bool(raw_italic) if raw_italic != WIN32COM_UNDEFINED else False
    # Underline — enum (0 = None, 1 = Single, …); undefined → False
    raw_ul = font.Underline
    underline = raw_ul not in (0, WIN32COM_UNDEFINED)
    # StrikeThrough — capital 'T'; boolean-like
    raw_strike = font.StrikeThrough
    strike = bool(raw_strike) if raw_strike != WIN32COM_UNDEFINED else False
    # Size
    raw_size = font.Size
    size: float | None = (
        float(raw_size) if raw_size != WIN32COM_UNDEFINED and raw_size > 0 else None
    )
    # Color — direct BGR Long
    raw_color = font.Color
    color: str | None = (
        _win32com_color_to_hex(int(raw_color))
        if raw_color != WIN32COM_UNDEFINED
        else None
    )
    # Background — try Shading first (arbitrary colour), then HighlightColorIndex
    bg_color: str | None = None
    try:
        shading_bgr = char_range.Shading.BackgroundPatternColor
        # wdColorAutomatic = -16777216; undefined = WIN32COM_UNDEFINED
        if shading_bgr not in (-16777216, WIN32COM_UNDEFINED) and shading_bgr >= 0:
            bg_color = _win32com_color_to_hex(int(shading_bgr))
    except Exception:  # noqa: BLE001
        pass
    if bg_color is None:
        try:
            hl_idx = char_range.HighlightColorIndex
            if hl_idx and hl_idx != WIN32COM_UNDEFINED:
                bg_color = _WD_HIGHLIGHT_INDEX_TO_HEX.get(int(hl_idx))
        except Exception:  # noqa: BLE001
            pass
    # Superscript / Subscript — boolean-like, treat undefined as False
    raw_sup = font.Superscript
    superscript = bool(raw_sup) if raw_sup != WIN32COM_UNDEFINED else False
    raw_sub = font.Subscript
    subscript = bool(raw_sub) if raw_sub != WIN32COM_UNDEFINED else False
    return (
        bold,
        italic,
        underline,
        strike,
        superscript,
        subscript,
        size,
        color,
        bg_color,
    )


def _has_win32com_range_mixed_formatting(rng: object) -> bool:
    """Checks whether a win32com Range has mixed per-character formatting.

    Uses a quick-exit via ``rng.Font.Bold == WIN32COM_UNDEFINED`` before
    falling back to full character-level iteration.  Returns ``False`` on
    any COM exception (conservative: assume uniform formatting).

    Args:
        rng: A win32com ``Range`` object.

    Returns:
        True if at least two characters have different formatting.
    """
    try:
        from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415

        # Quick check — COM already knows if bold is mixed across the range
        try:
            if rng.Font.Bold == WIN32COM_UNDEFINED:
                return True
        except Exception:  # noqa: BLE001
            pass

        count = rng.Characters.Count
        if count <= 1:
            return False

        first_sig = None
        for i in range(1, count + 1):
            ch = rng.Characters(i)
            if not ch.Text.strip():
                continue
            sig = _read_win32com_char_formatting(ch)
            if first_sig is None:
                first_sig = sig
            elif sig != first_sig:
                return True
        return False
    except Exception:  # noqa: BLE001
        return False  # COM error — conservatively treat as uniform


def _has_win32com_range_hyperlinks(rng: object) -> bool:
    """Checks whether a win32com Range contains hyperlinks.

    Args:
        rng: A win32com ``Range`` object.

    Returns:
        True if the range has at least one hyperlink.
    """
    try:
        return rng.Hyperlinks.Count > 0
    except Exception:  # noqa: BLE001
        return False


def _win32com_range_runs_to_html(rng: object) -> str:
    r"""Converts a win32com Range's characters to inline HTML.

    Groups consecutive characters with identical formatting and hyperlink
    URL into runs, skipping paragraph marks (``\r``), then emits HTML via
    ``_wrap_with_tags``.  Characters inside a hyperlink are tagged with
    ``<a href="...">`` so the LLM can preserve links during translation.

    Args:
        rng: A win32com ``Range`` object.

    Returns:
        HTML string representing the range's formatted text.
    """
    count = rng.Characters.Count

    # Build hyperlink position map (relative offset → URL)
    hyperlink_map: dict[int, str] = {}
    try:
        rng_start = rng.Start
        for hl_idx in range(1, rng.Hyperlinks.Count + 1):
            hl = rng.Hyperlinks(hl_idx)
            url = hl.Address
            if url:
                hl_start = hl.Range.Start - rng_start
                hl_end = hl.Range.End - rng_start
                for pos in range(hl_start, hl_end):
                    hyperlink_map[pos] = url
    except Exception:  # noqa: BLE001
        pass

    # Pass 1: group consecutive chars with same formatting + URL into runs
    # Each run is an 11-tuple:
    #   (text, bold, italic, underline, strike, sup, sub, size, color, bg, url)
    run_data: list[
        tuple[
            str,
            bool,
            bool,
            bool,
            bool,
            bool,
            bool,
            float | None,
            str | None,
            str | None,
            str | None,
        ]
    ] = []
    for i in range(1, count + 1):
        ch = rng.Characters(i)
        text = ch.Text
        # Skip paragraph marks
        if text == "\r":
            continue
        sig = _read_win32com_char_formatting(ch)
        url = hyperlink_map.get(i - 1)  # 0-based offset
        sig_with_url = (*sig, url)
        if run_data and run_data[-1][1:] == sig_with_url:
            # Same formatting + URL — extend current run
            prev = run_data[-1]
            run_data[-1] = (prev[0] + text, *sig_with_url)
        else:
            run_data.append((text, *sig_with_url))

    if not run_data:
        return ""

    # Detect variation — base is always None for safe roundtrip
    sizes = [d[7] for d in run_data]
    colors = [d[8] for d in run_data]
    bgs = [d[9] for d in run_data]
    has_size_variation = len(set(sizes)) > 1
    has_color_variation = len(set(colors)) > 1
    has_bg_variation = len(set(bgs)) > 1
    # base_size/color/bg are always None so every run with an explicit value
    # gets its own <span>.  Using most-common as base loses that value during
    # injection when the first run is not the most-common one.
    base_size = None
    base_color = None
    base_bg = None

    # Pass 2: emit HTML
    parts: list[str] = []
    for text, bold, italic, underline, strike, sup, sub, sz, clr, bg, url in run_data:
        parts.append(
            _wrap_with_tags(
                html.escape(text),
                bold,
                italic,
                underline,
                strike,
                sz if sz != base_size else None,
                clr if clr != base_color else None,
                has_size_variation=has_size_variation,
                has_color_variation=has_color_variation,
                bg_color_hex=bg if bg != base_bg else None,
                has_bg_variation=has_bg_variation,
                hyperlink_url=url,
                superscript=sup,
                subscript=sub,
            )
        )
    return "".join(parts)


def _has_win32com_word_mixed_formatting(para: object) -> bool:
    """Checks whether a win32com Word paragraph has mixed per-char formatting.

    Delegates to ``_has_win32com_range_mixed_formatting`` on the
    paragraph's ``Range``.

    Args:
        para: A win32com ``Paragraph`` object.

    Returns:
        True if at least two characters have different formatting.
    """
    return _has_win32com_range_mixed_formatting(para.Range)


def _has_win32com_word_hyperlinks(para: object) -> bool:
    """Checks whether a win32com Word paragraph contains hyperlinks.

    Delegates to ``_has_win32com_range_hyperlinks`` on the paragraph's
    ``Range``.

    Args:
        para: A win32com ``Paragraph`` object.

    Returns:
        True if the paragraph has at least one hyperlink.
    """
    return _has_win32com_range_hyperlinks(para.Range)


def _win32com_word_runs_to_html(para: object) -> str:
    """Converts a win32com Word paragraph's characters to inline HTML.

    Delegates to ``_win32com_range_runs_to_html`` on the paragraph's
    ``Range``.

    Args:
        para: A win32com ``Paragraph`` object.

    Returns:
        HTML string representing the paragraph's formatted text.
    """
    return _win32com_range_runs_to_html(para.Range)


def _extract_win32com_word(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from a Word document via win32com.

    For paragraphs with mixed per-run formatting, inline HTML is emitted
    via ``_win32com_word_runs_to_html`` so the LLM can preserve it.

    Args:
        file_path: Path to the .doc or .docx file.

    Returns:
        list: (location_key, text) pairs.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, file_path)
    try:
        texts: list[tuple[str, str]] = []

        # Paragraphs
        for i in range(1, doc.Paragraphs.Count + 1):
            para = doc.Paragraphs(i)
            if _has_win32com_word_mixed_formatting(
                para
            ) or _has_win32com_word_hyperlinks(para):
                text = _win32com_word_runs_to_html(para)
            else:
                text = para.Range.Text.rstrip("\r")
            if text.strip():
                texts.append((f"para:{i}", text))

        # Table cells
        for t_idx in range(1, doc.Tables.Count + 1):
            table = doc.Tables(t_idx)
            for r_idx in range(1, table.Rows.Count + 1):
                for c_idx in range(1, table.Columns.Count + 1):
                    try:
                        cell = table.Cell(r_idx, c_idx)
                        # Table cells don't have a Paragraphs collection
                        # easily accessible like the doc level, so we use
                        # plain-text extraction for cells.
                        text = cell.Range.Text.rstrip("\r\x07")
                        if text.strip():
                            texts.append(
                                (f"table:{t_idx}:{r_idx}:{c_idx}", text),
                            )
                    except Exception:  # noqa: BLE001
                        # Merged cells may raise errors
                        continue

        return texts
    finally:
        _win32com_close(app, doc, pycom, save_close=True)


def _inject_win32com_word_html_runs(  # noqa: PLR0912, PLR0913, PLR0915
    doc: object,
    rng: object,
    html_text: str,
    original_text: str = "",
    *,
    is_cell: bool = False,
    target_lang: str = "",
) -> None:
    r"""Replaces a win32com Word range's text with HTML-formatted segments.

    Parses ``html_text`` via ``_parse_html_formatting``, sets the full
    plain text on the range, then applies per-segment formatting by
    creating sub-ranges via ``doc.Range(start, end)``.

    The original font Name is preserved on the whole range (unless
    source and target script families differ).

    Args:
        doc: The win32com Word ``Document`` COM object.
        rng: The target ``Range`` (paragraph or cell range).
        html_text: Translated text with inline ``<b>/<i>/<u>/<s>`` tags.
        original_text: The text before translation (for script detection).
        is_cell: True when injecting into a table cell (no trailing ``\r``).
        target_lang: Target language name for font substitution.
    """
    segments = _parse_html_formatting(html_text)
    if not segments:
        # Strip residual tags so literal HTML doesn't appear in the document
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        rng.Text = plain + ("" if is_cell else "\r")
        return

    # Save range start and base font Name BEFORE text assignment
    range_start = rng.Start
    saved_name: str | None = None
    try:
        name_val = rng.Font.Name
        from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415

        if name_val != WIN32COM_UNDEFINED:
            saved_name = str(name_val)
    except Exception:  # noqa: BLE001
        pass

    # Build full plain text
    full_text = "".join(seg.text for seg in segments)
    # Set text — paragraph ranges need trailing \r, cells do not
    rng.Text = full_text + ("" if is_cell else "\r")

    # Restore base font Name on the whole range (script-aware substitution)
    if saved_name:
        font_name = _substitute_font(
            saved_name,
            original_text,
            full_text,
            target_lang,
        )
        if font_name is not None:
            try:
                whole = doc.Range(range_start, range_start + len(full_text))
                whole.Font.Name = font_name
            except Exception:  # noqa: BLE001
                pass

    # Apply per-segment formatting via sub-ranges
    offset = 0
    for seg in segments:
        seg_len = len(seg.text)
        if seg_len == 0:
            continue
        sub_rng = doc.Range(range_start + offset, range_start + offset + seg_len)
        sub_rng.Font.Bold = seg.bold
        sub_rng.Font.Italic = seg.italic
        # Underline: 1 = wdUnderlineSingle, 0 = wdUnderlineNone
        sub_rng.Font.Underline = 1 if seg.underline else 0
        sub_rng.Font.StrikeThrough = seg.strike
        sub_rng.Font.Superscript = seg.superscript
        sub_rng.Font.Subscript = seg.subscript
        if seg.font_size_pt is not None:
            sub_rng.Font.Size = seg.font_size_pt
        if seg.color_hex is not None:
            with contextlib.suppress(Exception):
                sub_rng.Font.Color = _color_hex_to_win32com(seg.color_hex)
        # Background colour — prefer Shading for arbitrary colours, fall
        # back to HighlightColorIndex for the 16 predefined colours.
        if seg.bg_color_hex is not None:
            with contextlib.suppress(Exception):
                hl_idx = _HEX_TO_WD_HIGHLIGHT_INDEX.get(seg.bg_color_hex)
                if hl_idx is not None:
                    sub_rng.HighlightColorIndex = hl_idx
                else:
                    sub_rng.Shading.BackgroundPatternColor = _color_hex_to_win32com(
                        seg.bg_color_hex
                    )
        offset += seg_len

    # Create hyperlinks for segments with URLs (second pass to avoid
    # interfering with formatting sub-range positions)
    offset = 0
    for seg in segments:
        seg_len = len(seg.text)
        if seg_len == 0:
            continue
        if seg.hyperlink_url:
            with contextlib.suppress(Exception):
                hl_rng = doc.Range(
                    range_start + offset,
                    range_start + offset + seg_len,
                )
                doc.Hyperlinks.Add(Anchor=hl_rng, Address=seg.hyperlink_url)
        offset += seg_len


def _inject_win32com_word(  # noqa: PLR0912, PLR0915
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into a Word document via win32com.

    For translations containing inline HTML formatting tags, uses
    ``_inject_win32com_word_html_runs`` to preserve per-run formatting.
    Otherwise falls back to uniform font save/restore.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, file_path)
    try:
        # Paragraphs
        for i in range(1, doc.Paragraphs.Count + 1):
            key = f"para:{i}"
            if key in translations:
                rng = doc.Paragraphs(i).Range
                orig_text = rng.Text.rstrip("\r")
                translated = translations[key]
                if _FORMATTING_HTML_RE.search(translated):
                    _inject_win32com_word_html_runs(
                        doc,
                        rng,
                        translated,
                        orig_text,
                        target_lang=target_lang,
                    )
                else:
                    font_saved = _save_win32com_font(rng.Font)
                    # Save highlight colour (Range property, not Font)
                    try:
                        highlight_saved = rng.HighlightColorIndex
                    except Exception:  # noqa: BLE001
                        highlight_saved = None
                    # Save shading background (arbitrary colour)
                    try:
                        shading_saved = rng.Shading.BackgroundPatternColor
                    except Exception:  # noqa: BLE001
                        shading_saved = None
                    rng.Text = translated + "\r"
                    # Re-acquire range after text change
                    rng = doc.Paragraphs(i).Range
                    _restore_win32com_font(
                        rng.Font,
                        font_saved,
                        original_text=orig_text,
                        translated_text=translated,
                        target_lang=target_lang,
                    )
                    # Restore highlight colour
                    if highlight_saved is not None:
                        with contextlib.suppress(Exception):
                            rng.HighlightColorIndex = highlight_saved
                    # Restore shading background colour
                    if shading_saved is not None:
                        with contextlib.suppress(Exception):
                            rng.Shading.BackgroundPatternColor = shading_saved

        # Table cells
        for t_idx in range(1, doc.Tables.Count + 1):
            table = doc.Tables(t_idx)
            for r_idx in range(1, table.Rows.Count + 1):
                for c_idx in range(1, table.Columns.Count + 1):
                    key = f"table:{t_idx}:{r_idx}:{c_idx}"
                    if key in translations:
                        try:
                            cell = table.Cell(r_idx, c_idx)
                            orig_text = cell.Range.Text.rstrip("\r\x07")
                            translated = translations[key]
                            if _FORMATTING_HTML_RE.search(translated):
                                _inject_win32com_word_html_runs(
                                    doc,
                                    cell.Range,
                                    translated,
                                    orig_text,
                                    is_cell=True,
                                    target_lang=target_lang,
                                )
                            else:
                                font_saved = _save_win32com_font(
                                    cell.Range.Font,
                                )
                                try:
                                    hl_saved = cell.Range.HighlightColorIndex
                                except Exception:  # noqa: BLE001
                                    hl_saved = None
                                try:
                                    shd_saved = (
                                        cell.Range.Shading.BackgroundPatternColor
                                    )
                                except Exception:  # noqa: BLE001
                                    shd_saved = None
                                cell.Range.Text = translated
                                # Re-acquire cell range after text change
                                cell = table.Cell(r_idx, c_idx)
                                _restore_win32com_font(
                                    cell.Range.Font,
                                    font_saved,
                                    original_text=orig_text,
                                    translated_text=translated,
                                    target_lang=target_lang,
                                )
                                if hl_saved is not None:
                                    with contextlib.suppress(Exception):
                                        cell.Range.HighlightColorIndex = hl_saved
                                if shd_saved is not None:
                                    with contextlib.suppress(Exception):
                                        shading = cell.Range.Shading
                                        shading.BackgroundPatternColor = shd_saved
                        except Exception:  # noqa: BLE001
                            continue

        doc.SaveAs(str(output_path.resolve()))
    finally:
        _win32com_close(app, doc, pycom, save_close=True)


def _extract_win32com_excel(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from an Excel workbook via win32com.

    Args:
        file_path: Path to the .xls or .xlsx file.

    Returns:
        list: (location_key, text) pairs.
    """
    app, wb, pycom = _win32com_open(_APP_EXCEL, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for ws in wb.Worksheets:
            used = ws.UsedRange
            if used is None:
                continue
            for row in range(1, used.Rows.Count + 1):
                for col in range(1, used.Columns.Count + 1):
                    cell = ws.Cells(row, col)
                    val = cell.Value
                    if isinstance(val, str) and val.strip():
                        texts.append(
                            (f"sheet:{ws.Name}:{row}:{col}", val),
                        )
        return texts
    finally:
        _win32com_close(app, wb, pycom, save_close=True)


def _inject_win32com_excel(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into an Excel workbook via win32com.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    app, wb, pycom = _win32com_open(_APP_EXCEL, file_path)
    try:
        for ws in wb.Worksheets:
            used = ws.UsedRange
            if used is None:
                continue
            for row in range(1, used.Rows.Count + 1):
                for col in range(1, used.Columns.Count + 1):
                    key = f"sheet:{ws.Name}:{row}:{col}"
                    if key in translations:
                        cell = ws.Cells(row, col)
                        orig_text = str(cell.Value) if cell.Value else ""
                        font_saved = _save_win32com_font(cell.Font)
                        cell.Value = translations[key]
                        _restore_win32com_font(
                            cell.Font,
                            font_saved,
                            original_text=orig_text,
                            translated_text=translations[key],
                            target_lang=target_lang,
                        )

        wb.SaveAs(str(output_path.resolve()))
    finally:
        _win32com_close(app, wb, pycom, save_close=True)


def _read_win32com_ppt_run_formatting(
    run_range: object,
) -> tuple[bool, bool, bool, bool, bool, bool, float | None, str | None, str | None]:
    """Reads inline formatting from a win32com PPT run TextRange.

    PPT ``Font.Color`` is a ``ColorFormat`` object — the BGR integer
    is accessed via ``.RGB``.  PPT ``Font.Strikethrough`` is lowercase 't'.

    Superscript/subscript is detected via ``Font.BaselineOffset``:
    positive values indicate superscript, negative values indicate subscript.

    Background colour is read via ``Font.Highlight.ForeColor.RGB``
    (Office 365 / 2019+).  Older versions silently return ``None``.

    Args:
        run_range: A win32com PPT ``TextRange`` for a single run.

    Returns:
        Tuple of (bold, italic, underline, strike, superscript, subscript,
        font_size_pt, color_hex, bg_color_hex).
    """
    from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415

    font = run_range.Font
    # Bold / Italic — tri-state (msoTrue=-1, msoFalse=0, msoTriStateMixed=-2)
    raw_bold = font.Bold
    bold = bool(raw_bold == -1) if raw_bold != WIN32COM_UNDEFINED else False
    raw_italic = font.Italic
    italic = bool(raw_italic == -1) if raw_italic != WIN32COM_UNDEFINED else False
    # Underline — tri-state
    raw_ul = font.Underline
    underline = bool(raw_ul == -1) if raw_ul != WIN32COM_UNDEFINED else False
    # Strikethrough — lowercase 't' in PPT
    raw_strike = font.Strikethrough
    strike = bool(raw_strike == -1) if raw_strike != WIN32COM_UNDEFINED else False
    # Size
    raw_size = font.Size
    size: float | None = (
        float(raw_size) if raw_size != WIN32COM_UNDEFINED and raw_size > 0 else None
    )
    # Color — ColorFormat object; BGR integer via .RGB
    try:
        raw_color = font.Color.RGB
        color: str | None = (
            _win32com_color_to_hex(int(raw_color))
            if raw_color != WIN32COM_UNDEFINED
            else None
        )
    except Exception:  # noqa: BLE001
        color = None
    # Background — Font.Highlight (Office 365 / 2019+); graceful fallback
    bg_color: str | None = None
    try:
        raw_hl = font.Highlight.ForeColor.RGB
        if isinstance(raw_hl, int) and raw_hl >= 0:
            bg_color = _win32com_color_to_hex(raw_hl)
    except Exception:  # noqa: BLE001
        pass
    # Superscript / Subscript — via BaselineOffset (positive = super, negative = sub)
    superscript = False
    subscript = False
    try:
        bl = font.BaselineOffset
        if isinstance(bl, (int, float)) and bl != WIN32COM_UNDEFINED:
            if bl > 0:
                superscript = True
            elif bl < 0:
                subscript = True
    except Exception:  # noqa: BLE001
        pass
    return (
        bold,
        italic,
        underline,
        strike,
        superscript,
        subscript,
        size,
        color,
        bg_color,
    )


def _has_win32com_ppt_mixed_formatting(para_range: object) -> bool:
    """Checks whether a win32com PPT paragraph has mixed per-run formatting.

    Iterates ``para_range.Runs(i)`` (1-based) and compares formatting tuples.

    Args:
        para_range: A win32com PPT ``TextRange`` for a paragraph.

    Returns:
        True if at least two runs have different formatting.
    """
    runs_collection = para_range.Runs()
    count = runs_collection.Count
    if count <= 1:
        return False

    first_sig = None
    for i in range(1, count + 1):
        run = para_range.Runs(i)
        if not run.Text.strip():
            continue
        sig = _read_win32com_ppt_run_formatting(run)
        if first_sig is None:
            first_sig = sig
        elif sig != first_sig:
            return True
    return False


def _has_win32com_ppt_hyperlinks(para_range: object) -> bool:
    """Checks whether a win32com PPT paragraph has hyperlinked runs.

    Iterates ``para_range.Runs(i)`` and checks each run's
    ``ActionSettings(ppMouseClick).Hyperlink.Address``.

    Args:
        para_range: A win32com PPT ``TextRange`` for a paragraph.

    Returns:
        True if at least one run has a non-empty hyperlink address.
    """
    try:
        runs_collection = para_range.Runs()
        count = runs_collection.Count
        for i in range(1, count + 1):
            run = para_range.Runs(i)
            if not run.Text.strip():
                continue
            try:
                url = run.ActionSettings(1).Hyperlink.Address
                if url:
                    return True
            except Exception:  # noqa: BLE001
                continue
    except Exception:  # noqa: BLE001
        pass
    return False


def _win32com_ppt_runs_to_html(para_range: object) -> str:
    """Converts a win32com PPT paragraph's runs to inline HTML.

    Two-pass: first collects run data, then emits HTML with ``<span>``
    only when size/colour actually vary.

    Args:
        para_range: A win32com PPT ``TextRange`` for a paragraph.

    Returns:
        HTML string representing the paragraph's formatted text.
    """
    runs_collection = para_range.Runs()
    count = runs_collection.Count

    # Pass 1: collect run data (11-tuple includes superscript/subscript + URL)
    run_data: list[
        tuple[
            str,
            bool,
            bool,
            bool,
            bool,
            bool,
            bool,
            float | None,
            str | None,
            str | None,
            str | None,
        ]
    ] = []
    for i in range(1, count + 1):
        run = para_range.Runs(i)
        text = run.Text
        if not text:
            continue
        b, it, u, s, sup, sub, sz, clr, bg = _read_win32com_ppt_run_formatting(run)
        # Read hyperlink URL via ActionSettings(ppMouseClick)
        url: str | None = None
        try:
            addr = run.ActionSettings(1).Hyperlink.Address
            if addr:
                url = addr
        except Exception:  # noqa: BLE001
            pass
        run_data.append((text, b, it, u, s, sup, sub, sz, clr, bg, url))

    if not run_data:
        return ""

    # Detect variation — base is always None for safe roundtrip
    sizes = [d[7] for d in run_data]
    colors = [d[8] for d in run_data]
    bgs = [d[9] for d in run_data]
    has_size_variation = len(set(sizes)) > 1
    has_color_variation = len(set(colors)) > 1
    has_bg_variation = len(set(bgs)) > 1
    # base_size/color/bg are always None so every run with an explicit value
    # gets its own <span>.  Using most-common as base loses that value during
    # injection when the first run is not the most-common one.
    base_size = None
    base_color = None
    base_bg = None

    # Pass 2: emit HTML
    parts: list[str] = []
    for text, bold, italic, underline, strike, sup, sub, sz, clr, bg, url in run_data:
        parts.append(
            _wrap_with_tags(
                html.escape(text),
                bold,
                italic,
                underline,
                strike,
                sz if sz != base_size else None,
                clr if clr != base_color else None,
                has_size_variation=has_size_variation,
                has_color_variation=has_color_variation,
                bg_color_hex=bg if bg != base_bg else None,
                has_bg_variation=has_bg_variation,
                hyperlink_url=url,
                superscript=sup,
                subscript=sub,
            )
        )
    return "".join(parts)


def _extract_win32com_ppt(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from a PowerPoint presentation via win32com.

    For paragraphs with mixed per-run formatting or hyperlinks, inline
    HTML is emitted via ``_win32com_ppt_runs_to_html`` so the LLM can
    preserve them.

    Args:
        file_path: Path to the .ppt or .pptx file.

    Returns:
        list: (location_key, text) pairs.
    """
    app, prs, pycom = _win32com_open(_APP_PPT, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for s_idx in range(1, prs.Slides.Count + 1):
            slide = prs.Slides(s_idx)
            for sh_idx in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(sh_idx)
                if not shape.HasTextFrame:
                    continue
                tf = shape.TextFrame
                for p_idx in range(1, tf.TextRange.Paragraphs().Count + 1):
                    para = tf.TextRange.Paragraphs(p_idx)
                    if _has_win32com_ppt_mixed_formatting(
                        para
                    ) or _has_win32com_ppt_hyperlinks(para):
                        text = _win32com_ppt_runs_to_html(para)
                    else:
                        text = para.Text
                    if text.strip():
                        texts.append(
                            (f"slide:{s_idx}:{sh_idx}:{p_idx}", text),
                        )
        return texts
    finally:
        _win32com_close(app, prs, pycom, save_close=True)


def _inject_win32com_ppt_html_runs(  # noqa: PLR0912, PLR0915
    tf: object,
    p_idx: int,
    html_text: str,
    original_text: str = "",
    *,
    target_lang: str = "",
) -> None:
    """Replaces a win32com PPT paragraph's text with HTML-formatted segments.

    Parses ``html_text`` via ``_parse_html_formatting``, sets the full
    plain text on the paragraph, then applies per-segment formatting
    using ``para_rng.Characters(offset + 1, length)`` (1-based).

    The original font Name is preserved on the whole paragraph (unless
    source and target script families differ).

    Args:
        tf: The win32com PPT ``TextFrame`` COM object.
        p_idx: 1-based paragraph index within the text frame.
        html_text: Translated text with inline ``<b>/<i>/<u>/<s>`` tags.
        original_text: The text before translation (for script detection).
        target_lang: Target language name for font substitution.
    """
    segments = _parse_html_formatting(html_text)
    if not segments:
        # Strip residual tags so literal HTML doesn't appear in the document
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        tf.TextRange.Paragraphs(p_idx).Text = plain
        return

    # Save base font Name BEFORE text assignment
    saved_name: str | None = None
    try:
        para_before = tf.TextRange.Paragraphs(p_idx)
        from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415

        name_val = para_before.Font.Name
        if name_val != WIN32COM_UNDEFINED:
            saved_name = str(name_val)
    except Exception:  # noqa: BLE001
        pass

    # Build full plain text
    full_text = "".join(seg.text for seg in segments)
    tf.TextRange.Paragraphs(p_idx).Text = full_text

    # Re-acquire paragraph range after text change
    para_rng = tf.TextRange.Paragraphs(p_idx)

    # Restore base font Name on the whole paragraph (script-aware)
    if saved_name:
        font_name = _substitute_font(
            saved_name,
            original_text,
            full_text,
            target_lang,
        )
        if font_name is not None:
            with contextlib.suppress(Exception):
                para_rng.Font.Name = font_name

    # Apply per-segment formatting via Characters(start, length) — 1-based
    offset = 0
    for seg in segments:
        seg_len = len(seg.text)
        if seg_len == 0:
            continue
        char_rng = para_rng.Characters(offset + 1, seg_len)
        char_rng.Font.Bold = -1 if seg.bold else 0
        char_rng.Font.Italic = -1 if seg.italic else 0
        char_rng.Font.Underline = -1 if seg.underline else 0
        # PPT uses lowercase 't' for Strikethrough
        char_rng.Font.Strikethrough = -1 if seg.strike else 0
        # Superscript / Subscript — via BaselineOffset (0.3 / -0.25 / 0.0)
        if seg.superscript:
            with contextlib.suppress(Exception):
                char_rng.Font.BaselineOffset = 0.3
        elif seg.subscript:
            with contextlib.suppress(Exception):
                char_rng.Font.BaselineOffset = -0.25
        else:
            with contextlib.suppress(Exception):
                char_rng.Font.BaselineOffset = 0.0
        if seg.font_size_pt is not None:
            char_rng.Font.Size = seg.font_size_pt
        if seg.color_hex is not None:
            with contextlib.suppress(Exception):
                char_rng.Font.Color.RGB = _color_hex_to_win32com(seg.color_hex)
        # Background — Font.Highlight (Office 365 / 2019+); silently skipped
        # on older versions that lack the Highlight property.
        if seg.bg_color_hex is not None:
            with contextlib.suppress(Exception):
                char_rng.Font.Highlight.ForeColor.RGB = _color_hex_to_win32com(
                    seg.bg_color_hex,
                )
        # Hyperlink — ActionSettings(ppMouseClick); silently skipped on
        # older Office versions that may not support programmatic hyperlinks.
        if seg.hyperlink_url:
            with contextlib.suppress(Exception):
                char_rng.ActionSettings(1).Hyperlink.Address = seg.hyperlink_url
        offset += seg_len


def _inject_win32com_ppt(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into a PowerPoint presentation via win32com.

    For translations containing inline HTML formatting tags, uses
    ``_inject_win32com_ppt_html_runs`` to preserve per-run formatting.
    Otherwise falls back to uniform font save/restore.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    app, prs, pycom = _win32com_open(_APP_PPT, file_path)
    try:
        for s_idx in range(1, prs.Slides.Count + 1):
            slide = prs.Slides(s_idx)
            for sh_idx in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(sh_idx)
                if not shape.HasTextFrame:
                    continue
                tf = shape.TextFrame
                for p_idx in range(1, tf.TextRange.Paragraphs().Count + 1):
                    key = f"slide:{s_idx}:{sh_idx}:{p_idx}"
                    if key in translations:
                        translated = translations[key]
                        para_rng = tf.TextRange.Paragraphs(p_idx)
                        orig_text = para_rng.Text
                        if _FORMATTING_HTML_RE.search(translated):
                            _inject_win32com_ppt_html_runs(
                                tf,
                                p_idx,
                                translated,
                                orig_text,
                                target_lang=target_lang,
                            )
                        else:
                            font_saved = _save_win32com_font(para_rng.Font)
                            # Save highlight colour (Font.Highlight is Office 365+)
                            highlight_saved = None
                            try:
                                raw_hl = para_rng.Font.Highlight.ForeColor.RGB
                                if isinstance(raw_hl, int) and raw_hl >= 0:
                                    highlight_saved = raw_hl
                            except Exception:  # noqa: BLE001
                                pass
                            para_rng.Text = translated
                            # Re-acquire paragraph range after text change
                            para_rng = tf.TextRange.Paragraphs(p_idx)
                            _restore_win32com_font(
                                para_rng.Font,
                                font_saved,
                                original_text=orig_text,
                                translated_text=translated,
                                target_lang=target_lang,
                            )
                            # Restore highlight colour
                            if highlight_saved is not None:
                                with contextlib.suppress(Exception):
                                    para_rng.Font.Highlight.ForeColor.RGB = (
                                        highlight_saved
                                    )

        prs.SaveAs(str(output_path.resolve()))
    finally:
        _win32com_close(app, prs, pycom, save_close=True)


# -- win32com comment extract/inject functions ------------------------------


def _extract_win32com_word_comments(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts comments from a Word document via win32com.

    Only top-level comments (where Ancestor is None) are extracted.

    Args:
        file_path: Path to the .doc file.

    Returns:
        list: (location_key, text) pairs with keys like 'comment:{index}'.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for i in range(1, doc.Comments.Count + 1):
            c = doc.Comments(i)
            # Skip replies (only top-level comments)
            try:
                if c.Ancestor is not None:
                    continue
            except Exception:  # noqa: BLE001
                pass  # Older Word versions may not have Ancestor
            text_rng = c.Range
            if _has_win32com_range_mixed_formatting(
                text_rng
            ) or _has_win32com_range_hyperlinks(text_rng):
                text = _win32com_range_runs_to_html(text_rng)
            else:
                text = text_rng.Text
            if text and text.strip():
                texts.append((f"comment:{c.Index}", text))
        return texts
    finally:
        _win32com_close(app, doc, pycom, save_close=True)


def _inject_win32com_word_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into a Word document via win32com.

    Args:
        output_path: Path to the .doc file to modify in place.
        translations: Mapping of 'comment:{index}' to translated text.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, output_path)
    try:
        for i in range(1, doc.Comments.Count + 1):
            c = doc.Comments(i)
            key = f"comment:{c.Index}"
            if key in translations:
                translated = translations[key]
                orig_text = c.Range.Text
                if _FORMATTING_HTML_RE.search(translated):
                    # HTML with formatting/hyperlinks — use rich injection
                    _inject_win32com_word_html_runs(
                        doc,
                        c.Range,
                        translated,
                        orig_text,
                        is_cell=True,
                    )
                else:
                    # Plain text — save/restore font properties
                    font_saved = _save_win32com_font(c.Range.Font)
                    c.Range.Text = translated
                    _restore_win32com_font(
                        c.Range.Font,
                        font_saved,
                        original_text=orig_text,
                        translated_text=translated,
                    )

        doc.Save()
    finally:
        _win32com_close(app, doc, pycom, save_close=True)


def _extract_win32com_excel_comments(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts cell comments from an Excel workbook via win32com.

    Args:
        file_path: Path to the .xls file.

    Returns:
        list: (location_key, text) pairs with keys like
              'comment:{sheet}:{row}:{col}'.
    """
    app, wb, pycom = _win32com_open(_APP_EXCEL, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for ws in wb.Worksheets:
            for comment in ws.Comments:
                text = comment.Text()
                if text and text.strip():
                    cell = comment.Parent
                    texts.append(
                        (
                            f"comment:{ws.Name}:{cell.Row}:{cell.Column}",
                            text,
                        )
                    )
        return texts
    finally:
        _win32com_close(app, wb, pycom, save_close=True)


def _inject_win32com_excel_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into an Excel workbook via win32com.

    Args:
        output_path: Path to the .xls file to modify in place.
        translations: Mapping of 'comment:{sheet}:{row}:{col}' to translated text.
    """
    app, wb, pycom = _win32com_open(_APP_EXCEL, output_path)
    try:
        for ws in wb.Worksheets:
            for comment in ws.Comments:
                cell = comment.Parent
                key = f"comment:{ws.Name}:{cell.Row}:{cell.Column}"
                if key in translations:
                    comment.Text(translations[key])

        wb.Save()
    finally:
        _win32com_close(app, wb, pycom, save_close=True)


def _extract_win32com_ppt_comments(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts comments from a PowerPoint presentation via win32com.

    Args:
        file_path: Path to the .ppt file.

    Returns:
        list: (location_key, text) pairs with keys like
              'comment:{slide_idx}:{comment_idx}'.
    """
    app, prs, pycom = _win32com_open(_APP_PPT, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for s_idx in range(1, prs.Slides.Count + 1):
            slide = prs.Slides(s_idx)
            for c_idx in range(1, slide.Comments.Count + 1):
                c = slide.Comments(c_idx)
                text = c.Text
                if text and text.strip():
                    # 0-based slide index for consistency with PPTX keys
                    texts.append(
                        (
                            f"comment:{s_idx - 1}:{c.Index}",
                            text,
                        )
                    )
        return texts
    finally:
        _win32com_close(app, prs, pycom, save_close=True)


def _inject_win32com_ppt_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into a PowerPoint presentation via win32com.

    Comment.Text in PowerPoint COM may be read-only. Falls back to
    deleting and re-adding with the same author and metadata.

    Args:
        output_path: Path to the .ppt file to modify in place.
        translations: Mapping of 'comment:{slide_idx}:{comment_idx}'
                      to translated text.
    """
    app, prs, pycom = _win32com_open(_APP_PPT, output_path)
    try:
        for s_idx in range(1, prs.Slides.Count + 1):
            slide = prs.Slides(s_idx)
            # Collect comments to modify (iterate in reverse for safe deletion)
            for c_idx in range(slide.Comments.Count, 0, -1):
                c = slide.Comments(c_idx)
                key = f"comment:{s_idx - 1}:{c.Index}"
                if key not in translations:
                    continue
                # Try direct text assignment first
                try:
                    c.Text = translations[key]
                except Exception:  # noqa: BLE001
                    # Fallback: delete and re-add with same metadata
                    author = c.Author
                    author_initials = c.AuthorInitials
                    dt = c.DateTime
                    left = c.Left
                    top = c.Top
                    c.Delete()
                    slide.Comments.Add(
                        left,
                        top,
                        author,
                        author_initials,
                        translations[key],
                        dt,
                    )

        prs.Save()
    finally:
        _win32com_close(app, prs, pycom, save_close=True)


# ---------------------------------------------------------------------------
# LibreOffice UNO backend
# ---------------------------------------------------------------------------


def _uno_file_url(path: Path) -> str:
    """Converts a file path to a file:/// URL for UNO.

    Args:
        path: File path to convert.

    Returns:
        str: The file URL.
    """
    return path.resolve().as_uri()


def _uno_open(file_path: Path) -> object:
    """Opens a document via LibreOffice UNO in hidden mode.

    Args:
        file_path: Path to the document.

    Returns:
        The UNO document object. Caller MUST call ``doc.close(True)``
        in a ``finally`` block.
    """
    from com.sun.star.beans import PropertyValue  # noqa: PLC0415

    desktop = _get_uno_desktop()
    props = (PropertyValue("Hidden", 0, True, 0),)
    doc = desktop.loadComponentFromURL(
        _uno_file_url(file_path),
        "_blank",
        0,
        props,
    )
    if doc is None:
        raise RuntimeError(
            f"UNO failed to open {file_path.name} (loadComponentFromURL returned None)"
        )
    return doc


def _uno_save(doc: object, output_path: Path) -> None:
    """Saves a UNO document preserving its original format.

    Reads the ``FilterName`` from the document's own ``MediaDescriptor``
    (set during import) and passes it to ``storeToURL`` so UNO writes in
    the same format as the source file rather than defaulting to ODF.
    Falls back to a hardcoded lookup if the descriptor is unavailable.

    Args:
        doc: The UNO document object.
        output_path: Destination file path.
    """
    from com.sun.star.beans import PropertyValue  # noqa: PLC0415

    # Try to get the filter from the document's own media descriptor
    filter_name = ""
    try:
        for prop in doc.getArgs():
            if prop.Name == "FilterName":
                filter_name = prop.Value
                break
    except Exception:  # noqa: BLE001
        pass

    # Fall back to extension-based lookup
    if not filter_name:
        filter_name = _UNO_FILTER_NAMES.get(output_path.suffix.lower(), "")

    props: list[object] = [PropertyValue("Overwrite", 0, True, 0)]
    if filter_name:
        props.append(PropertyValue("FilterName", 0, filter_name, 0))
    doc.storeToURL(_uno_file_url(output_path), tuple(props))


def _save_uno_char_props(text_obj: object) -> dict[str, object]:
    """Saves character formatting properties from a UNO text object.

    Reads each property in UNO_CHAR_PROPERTIES via getPropertyValue().
    Properties that raise are silently skipped.

    Args:
        text_obj: A UNO object supporting XPropertySet (paragraph, cell, etc.).

    Returns:
        dict: Mapping of property name to saved value.
    """
    from src.constants.office import UNO_CHAR_PROPERTIES  # noqa: PLC0415

    saved: dict[str, object] = {}
    for prop in UNO_CHAR_PROPERTIES:
        try:
            saved[prop] = text_obj.getPropertyValue(prop)
        except Exception:  # noqa: BLE001
            continue
    return saved


def _restore_uno_char_props(
    text_obj: object,
    saved: dict[str, object],
    *,
    original_text: str = "",
    translated_text: str = "",
    target_lang: str = "",
) -> None:
    """Restores previously saved character properties to a UNO text object.

    Sets each property independently so a single failure does not prevent
    other properties from being restored.

    When *target_lang* is provided and ``"CharFontName"`` is present in
    *saved*, the font name is substituted via :func:`_substitute_font`
    when the source and target scripts differ.

    Args:
        text_obj: A UNO object supporting XPropertySet.
        saved: Mapping of property name to value (from _save_uno_char_props).
        original_text: The text before translation (for script detection).
        translated_text: The text after translation (for script detection).
        target_lang: Target language name for font substitution.
    """
    # Substitute font name when scripts are incompatible
    if "CharFontName" in saved and original_text and translated_text:
        new_name = _substitute_font(
            str(saved["CharFontName"]),
            original_text,
            translated_text,
            target_lang,
        )
        if new_name is None:
            saved = {k: v for k, v in saved.items() if k != "CharFontName"}
        elif new_name != saved["CharFontName"]:
            saved = {**saved, "CharFontName": new_name}

    for prop, val in saved.items():
        try:
            text_obj.setPropertyValue(prop, val)
        except Exception:  # noqa: BLE001
            continue


# -- UNO per-run formatting constants ----------------------------------------

_UNO_WEIGHT_BOLD = 150.0
_UNO_WEIGHT_NORMAL = 100.0
_UNO_SLANT_ITALIC = 2
_UNO_SLANT_NONE = 0
_UNO_UNDERLINE_SINGLE = 1
_UNO_UNDERLINE_NONE = 0
_UNO_STRIKEOUT_SINGLE = 1
_UNO_STRIKEOUT_NONE = 0

# Properties set per-segment by _inject_uno_html_runs (excluded from base restore)
_UNO_FORMATTING_PROPS = {
    "CharWeight",
    "CharPosture",
    "CharUnderline",
    "CharStrikeout",
    "CharHighlight",
    "CharBackColor",
    "CharEscapement",
    "CharEscapementHeight",
}


def _read_uno_effective_formatting(
    obj: object,
) -> tuple[bool, bool, bool, bool, bool, bool]:
    """Reads the effective (resolved) formatting from a UNO text object.

    Returns the *effective* values, which include formatting inherited
    from paragraph/character styles.

    Note: UNO's ``CharPosture`` returns a ``uno.Enum`` (FontSlant) object,
    **not** a plain integer.  Comparing ``enum != 0`` always evaluates to
    ``True``, so we detect the enum via its ``.value`` string attribute
    (e.g. ``"NONE"``, ``"ITALIC"``).

    Args:
        obj: A UNO object supporting getPropertyValue (paragraph, portion,
            or text cursor).

    Returns:
        (bold, italic, underline, strike, superscript, subscript) booleans.
    """
    try:
        bold = obj.getPropertyValue("CharWeight") > _UNO_WEIGHT_NORMAL
    except Exception:  # noqa: BLE001
        bold = False
    try:
        posture = obj.getPropertyValue("CharPosture")
        # UNO returns a uno.Enum object for CharPosture (FontSlant);
        # .value gives the string name (e.g. "NONE", "ITALIC").
        if hasattr(posture, "value") and isinstance(posture.value, str):
            italic = posture.value != "NONE"
        else:
            italic = posture != _UNO_SLANT_NONE
    except Exception:  # noqa: BLE001
        italic = False
    try:
        underline = obj.getPropertyValue("CharUnderline") != _UNO_UNDERLINE_NONE
    except Exception:  # noqa: BLE001
        underline = False
    try:
        strike = obj.getPropertyValue("CharStrikeout") != _UNO_STRIKEOUT_NONE
    except Exception:  # noqa: BLE001
        strike = False
    # Superscript / Subscript via CharEscapement
    superscript = False
    subscript = False
    try:
        escapement = obj.getPropertyValue("CharEscapement")
        if isinstance(escapement, (int, float)):
            if escapement > 0:
                superscript = True
            elif escapement < 0:
                subscript = True
    except Exception:  # noqa: BLE001
        pass
    return (bold, italic, underline, strike, superscript, subscript)


def _read_uno_portion_formatting(
    portion: object,
) -> tuple[bool, bool, bool, bool, bool, bool]:
    """Reads effective inline formatting flags from a UNO text portion.

    Delegates to ``_read_uno_effective_formatting`` which handles the
    ``uno.Enum`` comparison for ``CharPosture``.

    Args:
        portion: A UNO TextPortion object (XPropertySet).

    Returns:
        (bold, italic, underline, strike, superscript, subscript) booleans.
    """
    return _read_uno_effective_formatting(portion)


def _read_uno_portion_bg_hex(portion: object) -> str | None:
    """Reads background/highlight colour from a UNO text portion.

    Checks ``CharHighlight`` first, then ``CharBackColor``.
    Both are integer RGB values; ``-1`` / ``0xFFFFFFFF`` means no colour.

    Args:
        portion: A UNO TextPortion object (XPropertySet).

    Returns:
        Lowercase hex colour string like ``"#ffff00"``, or None.
    """
    for prop in ("CharHighlight", "CharBackColor"):
        try:
            val = int(portion.getPropertyValue(prop))
            # -1 and 0xFFFFFFFF mean "no colour" / "transparent"
            if val < 0 or val > 0xFFFFFF:  # noqa: PLR2004
                continue
            # Skip black (0) for CharBackColor as it usually means "auto"
            if prop == "CharBackColor" and val == 0:
                continue
            return _int_to_color_hex(val)
        except Exception:  # noqa: BLE001
            continue
    return None


def _read_uno_portion_full_formatting(
    portion: object,
) -> tuple[
    bool,
    bool,
    bool,
    bool,
    bool,
    bool,
    float | None,
    str | None,
    str | None,
]:
    """Reads formatting flags plus font size, colour and bg from a UNO portion.

    Extends ``_read_uno_portion_formatting`` with ``CharHeight`` (float pt),
    ``CharColor`` (int → hex), and background colour via
    ``_read_uno_portion_bg_hex``.

    Args:
        portion: A UNO TextPortion object (XPropertySet).

    Returns:
        (bold, italic, underline, strike, superscript, subscript,
        font_size_pt, color_hex, bg_color_hex).
    """
    bold, italic, underline, strike, sup, sub = _read_uno_effective_formatting(portion)
    try:
        font_size_pt = float(portion.getPropertyValue("CharHeight"))
    except Exception:  # noqa: BLE001
        font_size_pt = None
    try:
        color_hex = _int_to_color_hex(int(portion.getPropertyValue("CharColor")))
    except Exception:  # noqa: BLE001
        color_hex = None
    bg_color_hex = _read_uno_portion_bg_hex(portion)
    return (
        bold,
        italic,
        underline,
        strike,
        sup,
        sub,
        font_size_pt,
        color_hex,
        bg_color_hex,
    )


def _has_uno_mixed_formatting(para: object) -> bool:
    """Checks whether a UNO paragraph has text portions with differing formatting.

    Compares each portion's full formatting (bold, italic, underline, strike,
    superscript, subscript, font size, colour, background colour).  Only
    considers portions with TextPortionType == "Text" and non-empty text.
    Returns False if 0 or 1 text portions remain.

    Args:
        para: A UNO paragraph supporting createEnumeration().

    Returns:
        True if at least two text portions have different formatting.
    """
    sigs: list[
        tuple[
            bool,
            bool,
            bool,
            bool,
            bool,
            bool,
            float | None,
            str | None,
            str | None,
        ]
    ] = []
    portion_enum = para.createEnumeration()
    while portion_enum.hasMoreElements():
        portion = portion_enum.nextElement()
        if portion.getPropertyValue("TextPortionType") != "Text":
            continue
        if not portion.getString():
            continue
        sigs.append(_read_uno_portion_full_formatting(portion))
    if len(sigs) <= 1:
        return False
    return len(set(sigs)) > 1


def _has_uno_hyperlinks(para: object) -> bool:
    """Checks whether a UNO paragraph has any portions with hyperlinks.

    Args:
        para: A UNO paragraph supporting createEnumeration().

    Returns:
        True if at least one text portion has a non-empty HyperLinkURL.
    """
    portion_enum = para.createEnumeration()
    while portion_enum.hasMoreElements():
        portion = portion_enum.nextElement()
        if portion.getPropertyValue("TextPortionType") != "Text":
            continue
        if not portion.getString():
            continue
        try:
            url = portion.getPropertyValue("HyperLinkURL")
            if url:
                return True
        except Exception:  # noqa: BLE001
            continue
    return False


def _uno_runs_to_html(para: object) -> str:
    """Converts a UNO paragraph's text portions to inline HTML.

    Two-pass approach: first collects all portion data to detect
    size/colour/bg variation, then emits HTML with ``<span>`` only when
    needed.  Portions with hyperlinks are wrapped in ``<a href="...">`` tags.

    Args:
        para: A UNO paragraph supporting createEnumeration().

    Returns:
        HTML string representing the paragraph's formatted text.
    """
    # Pass 1: collect portion data (text, bold, italic, underline, strike,
    # superscript, subscript, size, color, bg, url)
    portion_data: list[
        tuple[
            str,
            bool,
            bool,
            bool,
            bool,
            bool,
            bool,
            float | None,
            str | None,
            str | None,
            str | None,
        ]
    ] = []
    portion_enum = para.createEnumeration()
    while portion_enum.hasMoreElements():
        portion = portion_enum.nextElement()
        if portion.getPropertyValue("TextPortionType") != "Text":
            continue
        text = portion.getString()
        if not text:
            continue
        b, i, u, s, sup, sub, sz, clr, bg = _read_uno_portion_full_formatting(portion)
        # Read hyperlink URL from the portion
        url: str | None = None
        with contextlib.suppress(Exception):
            hlink_url = portion.getPropertyValue("HyperLinkURL")
            if hlink_url:
                url = hlink_url
        portion_data.append((text, b, i, u, s, sup, sub, sz, clr, bg, url))

    if not portion_data:
        return ""

    # Detect variation — base is always None for safe roundtrip
    sizes = [d[7] for d in portion_data]
    colors = [d[8] for d in portion_data]
    bgs = [d[9] for d in portion_data]
    has_size_variation = len(set(sizes)) > 1
    has_color_variation = len(set(colors)) > 1
    has_bg_variation = len(set(bgs)) > 1
    # base_size/color/bg are always None so every run with an explicit value
    # gets its own <span>.  Using most-common as base loses that value during
    # injection when the first run is not the most-common one.
    base_size = None
    base_color = None
    base_bg = None

    # Pass 2: emit HTML with <a> grouping for consecutive same-URL portions
    parts: list[str] = []
    current_url: str | None = None
    for (
        text,
        bold,
        italic,
        underline,
        strike,
        sup,
        sub,
        sz,
        clr,
        bg,
        url,
    ) in portion_data:
        # Close previous <a> if URL changed
        if url != current_url:
            if current_url is not None:
                parts.append("</a>")
            if url is not None:
                parts.append(f'<a href="{html.escape(url, quote=True)}">')
            current_url = url

        parts.append(
            _wrap_with_tags(
                html.escape(text),
                bold,
                italic,
                underline,
                strike,
                sz if sz != base_size else None,
                clr if clr != base_color else None,
                has_size_variation=has_size_variation,
                has_color_variation=has_color_variation,
                bg_color_hex=bg if bg != base_bg else None,
                has_bg_variation=has_bg_variation,
                superscript=sup,
                subscript=sub,
            )
        )
    # Close trailing <a> tag
    if current_url is not None:
        parts.append("</a>")
    return "".join(parts)


def _save_uno_first_portion_props(para: object) -> dict[str, object]:
    """Reads UNO_CHAR_PROPERTIES from the first text portion of a paragraph.

    This captures the actual font properties (name, size, colour) from the
    first run rather than from the paragraph level, which may differ.

    Args:
        para: A UNO paragraph supporting createEnumeration().

    Returns:
        dict mapping property names to values.  Empty if no text portion found.
    """
    from src.constants.office import UNO_CHAR_PROPERTIES  # noqa: PLC0415

    portion_enum = para.createEnumeration()
    while portion_enum.hasMoreElements():
        portion = portion_enum.nextElement()
        if portion.getPropertyValue("TextPortionType") != "Text":
            continue
        if not portion.getString():
            continue
        saved: dict[str, object] = {}
        for prop in UNO_CHAR_PROPERTIES:
            try:
                saved[prop] = portion.getPropertyValue(prop)
            except Exception:  # noqa: BLE001
                continue
        return saved
    return {}


def _inject_uno_html_runs(  # noqa: PLR0912, PLR0915
    para: object,
    html_text: str,
    base_props: dict[str, object],
    *,
    target_lang: str = "",
) -> None:
    """Replaces a UNO paragraph's text with HTML-formatted segments.

    Parses ``html_text`` via ``_parse_html_formatting``, sets the full
    plain text on the paragraph, then applies per-segment formatting via
    a text cursor.

    Base properties (font name, size, colour) from *base_props* are restored
    on the whole paragraph first, excluding the four formatting properties
    that are applied per-segment.  ``CharFontName`` is substituted with a
    compatible font when original and translated script families differ.

    Args:
        para: A UNO paragraph object.
        html_text: Translated text with inline <b>/<i>/<u>/<s> tags.
        base_props: Saved properties from ``_save_uno_first_portion_props``.
        target_lang: Target language name for font substitution.
    """
    segments = _parse_html_formatting(html_text)
    if not segments:
        # Strip residual tags so literal HTML doesn't appear in the document
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        para.setString(plain)
        return

    # Build full plain text from segments
    full_text = "".join(seg.text for seg in segments)
    para.setString(full_text)

    # Remove internal-only key before restoring properties to UNO
    original_text = str(base_props.pop("__original_text__", "") or "")

    # Restore base props (font name, size, colour) — skip formatting props
    para_text = para.getText()
    full_cursor = para_text.createTextCursorByRange(para.getStart())
    full_cursor.gotoStartOfParagraph(False)
    full_cursor.gotoEndOfParagraph(True)
    for prop, val in base_props.items():
        if prop in _UNO_FORMATTING_PROPS:
            continue
        if prop == "CharFontName":
            font_name = _substitute_font(
                str(val),
                original_text,
                full_text,
                target_lang,
            )
            if font_name is None:
                continue
            val = font_name  # noqa: PLW2901
        try:
            full_cursor.setPropertyValue(prop, val)
        except Exception:  # noqa: BLE001
            continue

    # Apply per-segment formatting
    offset = 0
    for seg in segments:
        seg_len = len(seg.text)
        if seg_len == 0:
            continue
        cursor = para_text.createTextCursorByRange(para.getStart())
        cursor.gotoStartOfParagraph(False)
        cursor.goRight(offset, False)
        cursor.goRight(seg_len, True)
        try:
            cursor.setPropertyValue(
                "CharWeight",
                _UNO_WEIGHT_BOLD if seg.bold else _UNO_WEIGHT_NORMAL,
            )
            cursor.setPropertyValue(
                "CharPosture",
                _UNO_SLANT_ITALIC if seg.italic else _UNO_SLANT_NONE,
            )
            cursor.setPropertyValue(
                "CharUnderline",
                _UNO_UNDERLINE_SINGLE if seg.underline else _UNO_UNDERLINE_NONE,
            )
            cursor.setPropertyValue(
                "CharStrikeout",
                _UNO_STRIKEOUT_SINGLE if seg.strike else _UNO_STRIKEOUT_NONE,
            )
            # Per-run superscript / subscript
            if seg.superscript:
                cursor.setPropertyValue("CharEscapement", 33)
                cursor.setPropertyValue("CharEscapementHeight", 58)
            elif seg.subscript:
                cursor.setPropertyValue("CharEscapement", -33)
                cursor.setPropertyValue("CharEscapementHeight", 58)
            else:
                cursor.setPropertyValue("CharEscapement", 0)
                cursor.setPropertyValue("CharEscapementHeight", 100)
            # Per-run font size override
            if seg.font_size_pt is not None:
                cursor.setPropertyValue("CharHeight", seg.font_size_pt)
            # Per-run text colour override
            if seg.color_hex is not None:
                cursor.setPropertyValue(
                    "CharColor",
                    _color_hex_to_int(seg.color_hex),
                )
            # Per-run background colour
            if seg.bg_color_hex is not None:
                bg_int = _color_hex_to_int(seg.bg_color_hex)
                cursor.setPropertyValue("CharHighlight", bg_int)
                cursor.setPropertyValue("CharBackColor", bg_int)
            else:
                # Clear highlight for segments without bg colour
                cursor.setPropertyValue("CharHighlight", -1)
                cursor.setPropertyValue("CharBackColor", -1)
            # Per-run hyperlink (empty string clears any existing link)
            cursor.setPropertyValue(
                "HyperLinkURL",
                seg.hyperlink_url or "",
            )
        except Exception:  # noqa: BLE001
            pass
        offset += seg_len


def _inject_uno_impress_html_runs(  # noqa: PLR0912, PLR0915
    para: object,
    html_text: str,
    base_props: dict[str, object],
    *,
    target_lang: str = "",
) -> None:
    """Impress-specific variant of ``_inject_uno_html_runs``.

    Impress text cursors do not implement ``XParagraphCursor``
    (no ``gotoStartOfParagraph``/``gotoEndOfParagraph``).  This
    function uses pure offset-based positioning via ``goRight``
    from the paragraph start range instead.

    Args:
        para: A UNO Impress paragraph object.
        html_text: Translated text with inline HTML tags.
        base_props: Saved properties from ``_save_uno_first_portion_props``.
        target_lang: Target language name for font substitution.
    """
    segments = _parse_html_formatting(html_text)
    if not segments:
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        para.setString(plain)
        return

    full_text = "".join(seg.text for seg in segments)
    para.setString(full_text)

    # Remove internal-only key before restoring properties to UNO
    original_text = str(base_props.pop("__original_text__", "") or "")

    # Restore base props on the full paragraph range (offset-based)
    para_text = para.getText()
    full_cursor = para_text.createTextCursorByRange(para.getStart())
    full_cursor.goRight(len(full_text), True)
    for prop, val in base_props.items():
        if prop in _UNO_FORMATTING_PROPS:
            continue
        if prop == "CharFontName":
            font_name = _substitute_font(
                str(val),
                original_text,
                full_text,
                target_lang,
            )
            if font_name is None:
                continue
            val = font_name  # noqa: PLW2901
        try:
            full_cursor.setPropertyValue(prop, val)
        except Exception:  # noqa: BLE001
            continue

    # Apply per-segment formatting (offset-based, no paragraph cursor)
    offset = 0
    for seg in segments:
        seg_len = len(seg.text)
        if seg_len == 0:
            continue
        cursor = para_text.createTextCursorByRange(para.getStart())
        cursor.goRight(offset, False)
        cursor.goRight(seg_len, True)
        try:
            cursor.setPropertyValue(
                "CharWeight",
                _UNO_WEIGHT_BOLD if seg.bold else _UNO_WEIGHT_NORMAL,
            )
            cursor.setPropertyValue(
                "CharPosture",
                _UNO_SLANT_ITALIC if seg.italic else _UNO_SLANT_NONE,
            )
            cursor.setPropertyValue(
                "CharUnderline",
                _UNO_UNDERLINE_SINGLE if seg.underline else _UNO_UNDERLINE_NONE,
            )
            cursor.setPropertyValue(
                "CharStrikeout",
                _UNO_STRIKEOUT_SINGLE if seg.strike else _UNO_STRIKEOUT_NONE,
            )
            # Per-run superscript / subscript
            if seg.superscript:
                cursor.setPropertyValue("CharEscapement", 33)
                cursor.setPropertyValue("CharEscapementHeight", 58)
            elif seg.subscript:
                cursor.setPropertyValue("CharEscapement", -33)
                cursor.setPropertyValue("CharEscapementHeight", 58)
            else:
                cursor.setPropertyValue("CharEscapement", 0)
                cursor.setPropertyValue("CharEscapementHeight", 100)
            if seg.font_size_pt is not None:
                cursor.setPropertyValue("CharHeight", seg.font_size_pt)
            if seg.color_hex is not None:
                cursor.setPropertyValue(
                    "CharColor",
                    _color_hex_to_int(seg.color_hex),
                )
            if seg.bg_color_hex is not None:
                bg_int = _color_hex_to_int(seg.bg_color_hex)
                cursor.setPropertyValue("CharHighlight", bg_int)
                cursor.setPropertyValue("CharBackColor", bg_int)
            else:
                cursor.setPropertyValue("CharHighlight", -1)
                cursor.setPropertyValue("CharBackColor", -1)
            # Per-run hyperlink (empty string clears any existing link)
            cursor.setPropertyValue(
                "HyperLinkURL",
                seg.hyperlink_url or "",
            )
        except Exception:  # noqa: BLE001
            pass
        offset += seg_len


def _inject_uno_impress_para_text(
    para: object,
    text: str,
    *,
    target_lang: str = "",
) -> None:
    """Injects translated text into a single UNO Impress paragraph.

    Uses ``_inject_uno_impress_html_runs`` for HTML-tagged text,
    plain ``setString`` with property save/restore otherwise.

    Args:
        para: A UNO Impress paragraph object.
        text: Translated text (plain or HTML-tagged).
        target_lang: Target language name for font substitution.
    """
    if _FORMATTING_HTML_RE.search(text):
        orig_text = para.getString()
        base_props = _save_uno_first_portion_props(para)
        base_props["__original_text__"] = orig_text
        _inject_uno_impress_html_runs(
            para,
            text,
            base_props,
            target_lang=target_lang,
        )
    else:
        orig_text = para.getString()
        char_saved = _save_uno_char_props(para)
        para.setString(text)
        _restore_uno_char_props(
            para,
            char_saved,
            original_text=orig_text,
            translated_text=text,
            target_lang=target_lang,
        )


# -- UNO text extract/inject functions --------------------------------------


def _extract_uno_writer(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from a Writer document via UNO.

    When a paragraph has mixed per-run formatting (e.g. bold + italic
    portions), the text is encoded as inline HTML so the LLM can
    preserve formatting tags.

    Args:
        file_path: Path to the .doc or .docx file.

    Returns:
        list: (location_key, text) pairs — plain text or inline HTML.
    """
    doc = _uno_open(file_path)
    texts: list[tuple[str, str]] = []
    try:
        text_content = doc.getText()
        enum = text_content.createEnumeration()
        p_idx = 0
        while enum.hasMoreElements():
            para = enum.nextElement()
            # Skip text tables (handled separately)
            if para.supportsService("com.sun.star.text.TextTable"):
                continue
            plain = para.getString()
            if plain.strip():
                text = (
                    _uno_runs_to_html(para)
                    if _has_uno_mixed_formatting(para) or _has_uno_hyperlinks(para)
                    else plain
                )
                texts.append((f"para:{p_idx}", text))
            p_idx += 1

        # Tables — for single-paragraph cells, apply mixed-formatting
        # or hyperlink check
        tables = doc.getTextTables()
        for t_idx in range(tables.getCount()):
            table = tables.getByIndex(t_idx)
            cell_names = table.getCellNames()
            for cell_name in cell_names:
                cell = table.getCellByName(cell_name)
                cell_text_obj = cell.getText()
                cell_enum = cell_text_obj.createEnumeration()
                paras: list[object] = []
                while cell_enum.hasMoreElements():
                    paras.append(cell_enum.nextElement())
                if len(paras) == 1 and paras[0].getString().strip():
                    cp = paras[0]
                    text = (
                        _uno_runs_to_html(cp)
                        if _has_uno_mixed_formatting(cp) or _has_uno_hyperlinks(cp)
                        else cp.getString()
                    )
                    texts.append(
                        (f"table:{t_idx}:{cell_name}", text),
                    )
                else:
                    text = cell.getString()
                    if text.strip():
                        texts.append(
                            (f"table:{t_idx}:{cell_name}", text),
                        )
    finally:
        doc.close(True)

    return texts


def _inject_uno_para_text(
    para: object,
    text: str,
    *,
    target_lang: str = "",
) -> None:
    """Injects translated text into a single UNO paragraph.

    Dispatches to ``_inject_uno_html_runs`` when *text* contains inline
    HTML formatting tags, otherwise uses plain ``setString`` with
    paragraph-level property save/restore.

    Args:
        para: A UNO paragraph object.
        text: Translated text (plain or HTML-tagged).
        target_lang: Target language name for font substitution.
    """
    if _FORMATTING_HTML_RE.search(text):
        orig_text = para.getString()
        base_props = _save_uno_first_portion_props(para)
        base_props["__original_text__"] = orig_text
        _inject_uno_html_runs(
            para,
            text,
            base_props,
            target_lang=target_lang,
        )
    else:
        orig_text = para.getString()
        char_saved = _save_uno_char_props(para)
        para.setString(text)
        _restore_uno_char_props(
            para,
            char_saved,
            original_text=orig_text,
            translated_text=text,
            target_lang=target_lang,
        )


def _inject_uno_cell_text(
    cell: object,
    text: str,
    *,
    target_lang: str = "",
) -> None:
    """Injects translated text into a UNO table cell.

    For single-paragraph cells with HTML tags, dispatches to
    ``_inject_uno_html_runs``.  Otherwise uses plain ``setString``
    with cell-level property save/restore.

    Args:
        cell: A UNO table cell object.
        text: Translated text (plain or HTML-tagged).
        target_lang: Target language name for font substitution.
    """
    if _FORMATTING_HTML_RE.search(text):
        cell_text_obj = cell.getText()
        cell_enum = cell_text_obj.createEnumeration()
        paras: list[object] = []
        while cell_enum.hasMoreElements():
            paras.append(cell_enum.nextElement())
        if len(paras) == 1:
            _inject_uno_para_text(paras[0], text, target_lang=target_lang)
            return
    # Plain text or multi-paragraph cell fallback
    orig_text = cell.getString()
    char_saved = _save_uno_char_props(cell)
    cell.setString(text)
    _restore_uno_char_props(
        cell,
        char_saved,
        original_text=orig_text,
        translated_text=text,
        target_lang=target_lang,
    )


def _inject_uno_writer(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into a Writer document via UNO.

    When the translated text contains inline HTML formatting tags
    (``<b>``, ``<i>``, ``<u>``, ``<s>``), per-segment formatting is
    applied via ``_inject_uno_html_runs``.  Otherwise, plain text is
    set with paragraph-level property restore.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    doc = _uno_open(file_path)
    try:
        text_content = doc.getText()
        enum = text_content.createEnumeration()
        p_idx = 0
        while enum.hasMoreElements():
            para = enum.nextElement()
            if para.supportsService("com.sun.star.text.TextTable"):
                continue
            key = f"para:{p_idx}"
            if key in translations:
                _inject_uno_para_text(
                    para,
                    translations[key],
                    target_lang=target_lang,
                )
            p_idx += 1

        # Tables
        tables = doc.getTextTables()
        for t_idx in range(tables.getCount()):
            table = tables.getByIndex(t_idx)
            for cell_name in table.getCellNames():
                key = f"table:{t_idx}:{cell_name}"
                if key in translations:
                    cell = table.getCellByName(cell_name)
                    _inject_uno_cell_text(
                        cell,
                        translations[key],
                        target_lang=target_lang,
                    )

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


def _extract_uno_calc(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from a Calc spreadsheet via UNO.

    Args:
        file_path: Path to the .xls or .xlsx file.

    Returns:
        list: (location_key, text) pairs.
    """
    doc = _uno_open(file_path)
    texts: list[tuple[str, str]] = []
    try:
        for s_idx in range(doc.getSheets().getCount()):
            sheet = doc.getSheets().getByIndex(s_idx)
            sheet_name = sheet.getName()
            cursor = sheet.createCursor()
            cursor.gotoStartOfUsedArea(False)
            cursor.gotoEndOfUsedArea(True)

            for row in range(
                cursor.getRangeAddress().StartRow, cursor.getRangeAddress().EndRow + 1
            ):
                for col in range(
                    cursor.getRangeAddress().StartColumn,
                    cursor.getRangeAddress().EndColumn + 1,
                ):
                    cell = sheet.getCellByPosition(col, row)
                    if cell.getType() == _UNO_CELL_TYPE_TEXT:
                        val = cell.getString()
                        if val.strip():
                            texts.append(
                                (f"sheet:{sheet_name}:{row}:{col}", val),
                            )
    finally:
        doc.close(True)

    return texts


def _inject_uno_calc(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into a Calc spreadsheet via UNO.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    doc = _uno_open(file_path)
    try:
        for s_idx in range(doc.getSheets().getCount()):
            sheet = doc.getSheets().getByIndex(s_idx)
            sheet_name = sheet.getName()
            cursor = sheet.createCursor()
            cursor.gotoStartOfUsedArea(False)
            cursor.gotoEndOfUsedArea(True)

            for row in range(
                cursor.getRangeAddress().StartRow, cursor.getRangeAddress().EndRow + 1
            ):
                for col in range(
                    cursor.getRangeAddress().StartColumn,
                    cursor.getRangeAddress().EndColumn + 1,
                ):
                    key = f"sheet:{sheet_name}:{row}:{col}"
                    if key in translations:
                        cell = sheet.getCellByPosition(col, row)
                        orig_text = cell.getString()
                        char_saved = _save_uno_char_props(cell)
                        cell.setString(translations[key])
                        _restore_uno_char_props(
                            cell,
                            char_saved,
                            original_text=orig_text,
                            translated_text=translations[key],
                            target_lang=target_lang,
                        )

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


def _extract_uno_impress(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from an Impress presentation via UNO.

    When any paragraph within a shape has mixed per-run formatting, the
    entire shape is extracted as inline HTML via ``_uno_runs_to_html``
    (paragraphs joined by newlines).  Otherwise, plain text is returned.

    Args:
        file_path: Path to the .ppt or .pptx file.

    Returns:
        list: (location_key, text) pairs.
    """
    doc = _uno_open(file_path)
    texts: list[tuple[str, str]] = []
    try:
        for s_idx in range(doc.getDrawPages().getCount()):
            page = doc.getDrawPages().getByIndex(s_idx)
            for sh_idx in range(page.getCount()):
                shape = page.getByIndex(sh_idx)
                if not shape.supportsService(
                    "com.sun.star.drawing.Text",
                ):
                    continue
                # Enumerate paragraphs and check for mixed formatting
                # or hyperlinks
                para_enum = shape.createEnumeration()
                paras: list[object] = []
                while para_enum.hasMoreElements():
                    paras.append(para_enum.nextElement())
                if any(
                    _has_uno_mixed_formatting(p) or _has_uno_hyperlinks(p)
                    for p in paras
                ):
                    text = "\n".join(_uno_runs_to_html(p) for p in paras)
                else:
                    text = shape.getString()
                if text.strip():
                    texts.append(
                        (f"slide:{s_idx}:{sh_idx}", text),
                    )
    finally:
        doc.close(True)

    return texts


def _inject_uno_impress(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into an Impress presentation via UNO.

    When the translated text contains inline HTML formatting tags,
    dispatches to ``_inject_uno_impress_para_text`` for per-run
    formatting on each paragraph (lines separated by newlines).
    Uses offset-based cursor positioning instead of XParagraphCursor
    methods.  Otherwise, uses plain ``setString`` with shape-level
    property save/restore.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    doc = _uno_open(file_path)
    try:
        for s_idx in range(doc.getDrawPages().getCount()):
            page = doc.getDrawPages().getByIndex(s_idx)
            for sh_idx in range(page.getCount()):
                shape = page.getByIndex(sh_idx)
                if not shape.supportsService(
                    "com.sun.star.drawing.Text",
                ):
                    continue
                key = f"slide:{s_idx}:{sh_idx}"
                if key not in translations:
                    continue
                translation = translations[key]
                handled = False
                if _FORMATTING_HTML_RE.search(translation):
                    para_enum = shape.createEnumeration()
                    paras: list[object] = []
                    while para_enum.hasMoreElements():
                        paras.append(para_enum.nextElement())
                    # Split HTML by newlines to match extraction's
                    # per-paragraph _uno_runs_to_html join.
                    lines = translation.split("\n")
                    if paras:
                        for p_idx, para in enumerate(paras):
                            line = lines[p_idx] if p_idx < len(lines) else ""
                            _inject_uno_impress_para_text(
                                para,
                                line,
                                target_lang=target_lang,
                            )
                        handled = True
                if not handled:
                    orig_text = shape.getString()
                    char_saved = _save_uno_char_props(shape)
                    shape.setString(translation)
                    _restore_uno_char_props(
                        shape,
                        char_saved,
                        original_text=orig_text,
                        translated_text=translation,
                        target_lang=target_lang,
                    )

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


# -- UNO comment extract/inject functions -----------------------------------


def _extract_uno_writer_comments(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts annotation comments from a Writer document via UNO.

    Enumerates text fields and filters by Annotation service.

    Args:
        file_path: Path to the .doc file.

    Returns:
        list: (location_key, text) pairs with keys like 'comment:{index}'.
    """
    doc = _uno_open(file_path)
    texts: list[tuple[str, str]] = []
    try:
        fields = doc.getTextFields()
        enum = fields.createEnumeration()
        idx = 0
        while enum.hasMoreElements():
            field = enum.nextElement()
            if field.supportsService(
                "com.sun.star.text.TextField.Annotation",
            ):
                text = field.getPropertyValue("Content")
                if text and text.strip():
                    texts.append((f"comment:{idx}", text))
                idx += 1
    finally:
        doc.close(True)

    return texts


def _inject_uno_writer_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into a Writer document via UNO.

    Args:
        output_path: Path to the .doc file to modify in place.
        translations: Mapping of 'comment:{index}' to translated text.
    """
    doc = _uno_open(output_path)
    try:
        fields = doc.getTextFields()
        enum = fields.createEnumeration()
        idx = 0
        while enum.hasMoreElements():
            field = enum.nextElement()
            if field.supportsService(
                "com.sun.star.text.TextField.Annotation",
            ):
                key = f"comment:{idx}"
                if key in translations:
                    field.setPropertyValue("Content", translations[key])
                idx += 1

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


def _extract_uno_calc_comments(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts cell annotations from a Calc spreadsheet via UNO.

    Args:
        file_path: Path to the .xls file.

    Returns:
        list: (location_key, text) pairs with keys like
              'comment:{sheet}:{row}:{col}' (1-based for XLSX compatibility).
    """
    doc = _uno_open(file_path)
    texts: list[tuple[str, str]] = []
    try:
        for s_idx in range(doc.getSheets().getCount()):
            sheet = doc.getSheets().getByIndex(s_idx)
            sheet_name = sheet.getName()
            annotations = sheet.getAnnotations()
            for a_idx in range(annotations.getCount()):
                annotation = annotations.getByIndex(a_idx)
                text = annotation.getString()
                if text and text.strip():
                    pos = annotation.getPosition()
                    # 1-based row/col for XLSX key compatibility
                    texts.append(
                        (
                            f"comment:{sheet_name}:{pos.Row + 1}:{pos.Column + 1}",
                            text,
                        )
                    )
    finally:
        doc.close(True)

    return texts


def _inject_uno_calc_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into a Calc spreadsheet via UNO.

    Args:
        output_path: Path to the .xls file to modify in place.
        translations: Mapping of 'comment:{sheet}:{row}:{col}'
                      to translated text.
    """
    doc = _uno_open(output_path)
    try:
        for s_idx in range(doc.getSheets().getCount()):
            sheet = doc.getSheets().getByIndex(s_idx)
            sheet_name = sheet.getName()
            annotations = sheet.getAnnotations()
            for a_idx in range(annotations.getCount()):
                annotation = annotations.getByIndex(a_idx)
                pos = annotation.getPosition()
                key = f"comment:{sheet_name}:{pos.Row + 1}:{pos.Column + 1}"
                if key in translations:
                    orig_text = annotation.getString()
                    char_saved = _save_uno_char_props(annotation)
                    annotation.setString(translations[key])
                    _restore_uno_char_props(
                        annotation,
                        char_saved,
                        original_text=orig_text,
                        translated_text=translations[key],
                    )

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


def _extract_uno_impress_comments(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts annotations from an Impress presentation via UNO.

    Args:
        file_path: Path to the .ppt file.

    Returns:
        list: (location_key, text) pairs with keys like
              'comment:{page_idx}:{anno_idx}'.
    """
    doc = _uno_open(file_path)
    texts: list[tuple[str, str]] = []
    try:
        for p_idx in range(doc.getDrawPages().getCount()):
            page = doc.getDrawPages().getByIndex(p_idx)
            try:
                annotations = page.getAnnotations()
            except Exception:  # noqa: BLE001
                continue  # Page may not support annotations
            enum = annotations.createEnumeration()
            a_idx = 0
            while enum.hasMoreElements():
                annotation = enum.nextElement()
                text = annotation.TextRange.getString()
                if text and text.strip():
                    texts.append(
                        (
                            f"comment:{p_idx}:{a_idx}",
                            text,
                        )
                    )
                a_idx += 1
    finally:
        doc.close(True)

    return texts


def _inject_uno_impress_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into an Impress presentation via UNO.

    Args:
        output_path: Path to the .ppt file to modify in place.
        translations: Mapping of 'comment:{page_idx}:{anno_idx}'
                      to translated text.
    """
    doc = _uno_open(output_path)
    try:
        for p_idx in range(doc.getDrawPages().getCount()):
            page = doc.getDrawPages().getByIndex(p_idx)
            try:
                annotations = page.getAnnotations()
            except Exception:  # noqa: BLE001
                continue
            enum = annotations.createEnumeration()
            a_idx = 0
            while enum.hasMoreElements():
                annotation = enum.nextElement()
                key = f"comment:{p_idx}:{a_idx}"
                if key in translations:
                    text_range = annotation.TextRange
                    orig_text = text_range.getString()
                    char_saved = _save_uno_char_props(text_range)
                    text_range.setString(translations[key])
                    _restore_uno_char_props(
                        text_range,
                        char_saved,
                        original_text=orig_text,
                        translated_text=translations[key],
                    )
                a_idx += 1

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


# ---------------------------------------------------------------------------
# Format conversion helpers (legacy ↔ modern)
# ---------------------------------------------------------------------------

# win32com SaveAs format codes by output extension
_WIN32COM_FORMAT_CODES: dict[str, tuple[str, int]] = {
    ".docx": ("Word.Application", 16),  # wdFormatXMLDocument
    ".doc": ("Word.Application", 0),  # wdFormatDocument
    ".xlsx": ("Excel.Application", 51),  # xlOpenXMLWorkbook
    ".xls": ("Excel.Application", 56),  # xlWorkbookNormal
    ".pptx": ("PowerPoint.Application", 24),  # ppSaveAsOpenXMLPresentation
    ".ppt": ("PowerPoint.Application", 1),  # ppSaveAsPPT
}

# UNO export filter names by output extension
_UNO_FILTER_NAMES: dict[str, str] = {
    ".docx": "MS Word 2007 XML",
    ".doc": "MS Word 97",
    ".xlsx": "Calc MS Excel 2007 XML",
    ".xls": "MS Excel 97",
    ".pptx": "Impress MS PowerPoint 2007 XML",
    ".ppt": "MS PowerPoint 97",
}


def _convert_with_win32com(input_path: Path, output_path: Path) -> None:
    """Converts an office file to another format using win32com SaveAs.

    Uses the output extension to determine the application and format code.

    Args:
        input_path: Path to the source file.
        output_path: Path for the converted file.
    """
    out_suffix = output_path.suffix.lower()
    app_name, fmt_code = _WIN32COM_FORMAT_CODES[out_suffix]

    app, doc_obj, pycom = _win32com_open(app_name, input_path)
    try:
        if app_name == _APP_PPT:
            doc_obj.SaveAs(str(output_path.resolve()), fmt_code)
        else:
            doc_obj.SaveAs(str(output_path.resolve()), FileFormat=fmt_code)
    finally:
        _win32com_close(app, doc_obj, pycom, save_close=True)


def _convert_with_uno(input_path: Path, output_path: Path) -> None:
    """Converts an office file to another format using LibreOffice UNO.

    Uses the output extension to select the export filter name.

    Args:
        input_path: Path to the source file.
        output_path: Path for the converted file.
    """
    from com.sun.star.beans import PropertyValue  # noqa: PLC0415

    out_suffix = output_path.suffix.lower()
    filter_name = _UNO_FILTER_NAMES[out_suffix]

    doc = _uno_open(input_path)
    try:
        store_props = (
            PropertyValue("FilterName", 0, filter_name, 0),
            PropertyValue("Overwrite", 0, True, 0),
        )
        doc.storeToURL(_uno_file_url(output_path), store_props)
    finally:
        doc.close(True)


def convert_to_modern_format(input_path: Path, output_path: Path) -> bool:
    """Converts a legacy/ODF office file to modern format (.docx/.xlsx/.pptx).

    Detects the available backend (win32com or UNO) and delegates to
    the appropriate conversion helper. Returns True on success, False
    on failure (logs a warning instead of raising).

    Args:
        input_path: Path to the translated file in legacy/ODF format.
        output_path: Path for the converted modern format file.

    Returns:
        bool: True if conversion succeeded, False otherwise.
    """
    # 1. Try win32com
    try:
        import win32com.client  # noqa: F401, PLC0415

        _convert_with_win32com(input_path, output_path)
        return True
    except ImportError:
        pass
    except Exception:  # noqa: BLE001
        logger.warning(
            "win32com conversion failed for %s", input_path.name, exc_info=True
        )
        return False

    # 2. Try UNO
    try:
        import uno  # noqa: F401, PLC0415

        _convert_with_uno(input_path, output_path)
        return True
    except ImportError:
        pass
    except Exception:  # noqa: BLE001
        logger.warning("UNO conversion failed for %s", input_path.name, exc_info=True)
        return False

    logger.warning(
        "No backend available to convert %s to modern format", input_path.name
    )
    return False


# ---------------------------------------------------------------------------
# python-docx / openpyxl / python-pptx / odfpy backend (modern + ODF formats)
# ---------------------------------------------------------------------------


# -- ODF helpers (used by odfpy extract/inject functions) -------------------

# Lazy-cached qnames for ODF Tab, LineBreak, Span, and A (hyperlink) elements.
# Computed once on first call to avoid creating throwaway ODF objects
# per invocation (matters for large ODS files with thousands of cells).
_ODF_TAB_QNAME: tuple[str, str] | None = None
_ODF_LB_QNAME: tuple[str, str] | None = None
_ODF_SPAN_QNAME: tuple[str, str] | None = None
_ODF_A_QNAME: tuple[str, str] | None = None


def _odf_qnames() -> tuple[
    tuple[str, str], tuple[str, str], tuple[str, str], tuple[str, str]
]:
    """Returns cached (tab_qname, linebreak_qname, span_qname, a_qname)."""
    global _ODF_TAB_QNAME, _ODF_LB_QNAME, _ODF_SPAN_QNAME, _ODF_A_QNAME  # noqa: PLW0603
    if _ODF_TAB_QNAME is None:
        from odf.text import A, LineBreak, Span, Tab  # noqa: PLC0415

        _ODF_TAB_QNAME = Tab().qname
        _ODF_LB_QNAME = LineBreak().qname
        _ODF_SPAN_QNAME = Span().qname
        _ODF_A_QNAME = A(href="").qname
    return _ODF_TAB_QNAME, _ODF_LB_QNAME, _ODF_SPAN_QNAME, _ODF_A_QNAME


def _odf_element_text(
    element: object,
    *,
    preserve_links: bool = False,
) -> str:
    """Recursively extracts all text content from an ODF element.

    Walks the element's childNodes tree. Text nodes (nodeType == 3) have
    their data collected. Element nodes (nodeType == 1) are recursed into.
    Tab elements produce a tab character; line-break elements produce a
    newline.

    When *preserve_links* is True, ``<text:a>`` hyperlinks are emitted as
    ``<a href="url">text</a>`` HTML tags instead of plain text.  This is
    used during extraction so the LLM sees (and preserves) hyperlink
    structure.

    Args:
        element: An odfpy element node.
        preserve_links: If True, emit ``<a>`` HTML for ``<text:a>`` elements.

    Returns:
        str: The concatenated text content (may contain ``<a>`` HTML when
             *preserve_links* is True).
    """
    tab_qn, lb_qn, _, a_qn = _odf_qnames()

    parts: list[str] = []
    for child in element.childNodes:
        if child.nodeType == _NODE_TYPE_TEXT:
            parts.append(child.data)
        elif child.nodeType == _NODE_TYPE_ELEMENT:
            if child.qname == tab_qn:
                parts.append("\t")
            elif child.qname == lb_qn:
                parts.append("\n")
            elif preserve_links and child.qname == a_qn:
                # ODF hyperlink — emit <a href="url">text</a>
                href = (
                    getattr(child, "attributes", {}).get(
                        ("http://www.w3.org/1999/xlink", "href"),
                    )
                    or ""
                )
                link_text = _odf_element_text(child)
                if href:
                    parts.append(
                        f'<a href="{html.escape(href, quote=True)}">'
                        f"{html.escape(link_text)}</a>",
                    )
                else:
                    parts.append(link_text)
            else:
                parts.append(
                    _odf_element_text(child, preserve_links=preserve_links),
                )
    return "".join(parts)


def _odf_replace_text(
    element: object,
    new_text: str,
    *,
    target_lang: str = "",
) -> None:
    """Replaces all text content in an ODF element with new text.

    Preserves the first ``<text:span>``'s ``stylename`` attribute so that
    character formatting (bold, italic, font, etc.) is retained.  If no
    span is found, falls back to plain ``addText()``.

    When *new_text* contains ``<a href="...">`` HTML tags (from hyperlink
    preservation during extraction), parses them via
    ``_parse_html_formatting`` and creates ``<text:a>`` elements with the
    correct ``xlink:href`` attribute.

    Note:
        odfpy's removeChild() cannot handle text nodes (nodeType == 3)
        because its internal cache assertion requires Element instances.
        We manually clear childNodes and only update caches for Elements.

    Args:
        element: An odfpy element node (typically a P or H element).
        new_text: The replacement text (may contain ``<a>`` HTML).
        target_lang: Target language name for font substitution.
    """
    from odf.text import A as OdfA  # noqa: PLC0415
    from odf.text import Span as OdfSpan  # noqa: PLC0415

    # Find the first <text:span> and save its stylename before clearing
    _, _, span_qname, _ = _odf_qnames()
    saved_stylename = None
    for child in element.childNodes:
        if child.nodeType == _NODE_TYPE_ELEMENT and child.qname == span_qname:
            saved_stylename = getattr(child, "attributes", {}).get(
                ("urn:oasis:names:tc:opendocument:xmlns:text:1.0", "style-name"),
            )
            break

    # Preserve the span style if one was found (font substitution is
    # handled at a higher level via _substitute_font)
    compatible_scripts = bool(saved_stylename)

    # Clear all children
    old_children = list(element.childNodes)
    element.childNodes = []
    for child in old_children:
        if child.nodeType == _NODE_TYPE_ELEMENT and element.ownerDocument:
            element.ownerDocument.remove_from_caches(child)
        child.parentNode = None

    # Check for hyperlink HTML and rebuild with <text:a> elements
    if _FORMATTING_HTML_RE.search(new_text):
        segments = _parse_html_formatting(new_text)
        if segments:
            for seg in segments:
                if seg.hyperlink_url:
                    link_el = OdfA(href=seg.hyperlink_url)
                    link_el.addText(seg.text)
                    element.addElement(link_el)
                elif compatible_scripts:
                    span = OdfSpan(stylename=saved_stylename)
                    span.addText(seg.text)
                    element.addElement(span)
                else:
                    element.addText(seg.text)
            return
        # No parseable segments — strip residual tags and fall through
        new_text = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", new_text))

    # Re-add text, wrapped in a Span if we had a stylename and scripts
    # are compatible (otherwise let the app pick a suitable font)
    if compatible_scripts:
        span = OdfSpan(stylename=saved_stylename)
        span.addText(new_text)
        element.addElement(span)
    else:
        element.addText(new_text)


_ODF_TABLE_CELL_QNAME = (
    "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
    "table-cell",
)


def _is_inside_table_cell(element: object) -> bool:
    """Checks if an ODF element is nested inside a table cell.

    Args:
        element: An odfpy element node.

    Returns:
        bool: True if a TableCell ancestor is found.
    """
    parent = element.parentNode
    while parent is not None:
        if getattr(parent, "qname", None) == _ODF_TABLE_CELL_QNAME:
            return True
        parent = getattr(parent, "parentNode", None)
    return False


def _resolve_para_hyperlink_rels(para: object) -> dict[str, str]:
    """Resolves hyperlink ``r:id`` values to URLs for a paragraph.

    Scans ``para._element`` for ``<w:hyperlink>`` children, looks up
    each ``r:id`` in the document's relationship collection, and
    returns a mapping of ``r:id`` → target URL.

    Args:
        para: A python-docx Paragraph object.

    Returns:
        dict mapping ``r:id`` strings to URL strings.  Empty if no
        external hyperlinks exist.
    """
    from docx.oxml.ns import qn  # noqa: PLC0415

    rels: dict[str, str] = {}
    r_id_attr = qn("r:id")
    for child in para._element:
        if child.tag != _W_HYPERLINK_TAG:
            continue
        r_id = child.get(r_id_attr)
        if not r_id:
            continue
        with contextlib.suppress(KeyError, AttributeError):
            rels[r_id] = para.part.rels[r_id].target_ref
    return rels


def _extract_para_with_links(
    para: object,
) -> str:
    """Extracts text from a paragraph, preserving hyperlinks as ``<a>`` tags.

    Uses the HTML path (``_runs_to_html``) when the paragraph has mixed
    formatting or ``<w:hyperlink>`` children.  Falls back to ``para.text``
    for simple uniform-formatting paragraphs without hyperlinks.

    Args:
        para: A python-docx Paragraph object.

    Returns:
        Plain text or inline HTML string.
    """
    has_links = _para_has_hyperlinks(para._element)
    if _has_mixed_formatting(para) or has_links:
        hyperlink_rels = _resolve_para_hyperlink_rels(para) if has_links else None
        return _runs_to_html(para, hyperlink_rels)
    return para.text


def _extract_python_docx(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from a DOCX file via python-docx.

    Extracts paragraph text and table cell text. Each paragraph or
    cell with non-empty text gets a unique location key. When a
    paragraph has mixed formatting or hyperlinks, the text is encoded
    as inline HTML so the LLM can preserve formatting and link tags.

    Args:
        file_path: Path to the .docx file.

    Returns:
        list: (location_key, text) pairs — plain text or inline HTML.
    """
    from docx import Document  # noqa: PLC0415

    doc = Document(str(file_path))
    texts: list[tuple[str, str]] = []

    # Paragraphs
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            text = _extract_para_with_links(para)
            texts.append((f"para:{i}", text))

    # Tables — deduplicate merged cells (python-docx returns the same
    # Cell object for each column a merged cell spans).  Extract each
    # paragraph individually so injection maps back 1-to-1.
    # Note: store actual _tc elements (not id()) to prevent lxml proxy
    # garbage-collection from causing id reuse across different cells.
    seen_tcs: set[object] = set()
    for t_idx, table in enumerate(doc.tables):
        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                tc = cell._tc
                if tc in seen_tcs:
                    continue
                seen_tcs.add(tc)
                for p_idx, para in enumerate(cell.paragraphs):
                    if para.text.strip():
                        text = _extract_para_with_links(para)
                        texts.append(
                            (
                                f"table:{t_idx}:{r_idx}:{c_idx}:{p_idx}",
                                text,
                            )
                        )

    return texts


_ODF_NS_STYLE = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
_ODF_NS_FO = "urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0"
_ODF_NS_OFFICE = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"


def _set_odf_default_rtl(file_path: Path) -> None:
    """Rewrites *file_path* (an ODF zip) so paragraphs default to RTL.

    Adds — or extends — the ``<style:default-style style:family="paragraph">``
    block in ``styles.xml`` to set ``style:writing-mode="rl-tb"`` and
    ``fo:text-align="end"``.  Idempotent: running on an already-RTL
    document is a no-op.
    """
    if not file_path.exists():
        return
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            try:
                styles_xml = zf.read("styles.xml")
            except KeyError:
                return
            other_files = {
                name: zf.read(name) for name in zf.namelist() if name != "styles.xml"
            }
    except (OSError, zipfile.BadZipFile):
        return

    try:
        root = etree.fromstring(styles_xml)
    except etree.XMLSyntaxError:
        return

    nsmap = {"style": _ODF_NS_STYLE, "fo": _ODF_NS_FO, "office": _ODF_NS_OFFICE}
    styles_block = root.find(f"{{{_ODF_NS_OFFICE}}}styles")
    if styles_block is None:
        return

    default_para = styles_block.find(
        "style:default-style[@style:family='paragraph']",
        namespaces=nsmap,
    )
    if default_para is None:
        default_para = etree.SubElement(
            styles_block,
            f"{{{_ODF_NS_STYLE}}}default-style",
            attrib={f"{{{_ODF_NS_STYLE}}}family": "paragraph"},
        )

    para_props = default_para.find("style:paragraph-properties", namespaces=nsmap)
    if para_props is None:
        para_props = etree.SubElement(
            default_para,
            f"{{{_ODF_NS_STYLE}}}paragraph-properties",
        )
    para_props.set(f"{{{_ODF_NS_STYLE}}}writing-mode", "rl-tb")
    para_props.set(f"{{{_ODF_NS_FO}}}text-align", "end")

    new_styles_xml = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )

    try:
        with zipfile.ZipFile(file_path, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("styles.xml", new_styles_xml)
            for name, data in other_files.items():
                zf.writestr(name, data)
    except OSError:
        return


def _set_docx_paragraph_rtl(para: object) -> None:
    """Adds ``<w:bidi/>`` to the paragraph and ``<w:rtl/>`` to every run.

    Word and LibreOffice Writer use these flags to flip paragraph
    direction and shape mirrored punctuation (parens, quotes) at run
    boundaries.  Without them an Arabic / Hebrew paragraph renders
    flush-left with broken punctuation.
    """
    from docx.oxml import OxmlElement  # noqa: PLC0415
    from docx.oxml.ns import qn  # noqa: PLC0415

    pPr = para._element.get_or_add_pPr()  # type: ignore[attr-defined]  # noqa: N806
    if pPr.find(qn("w:bidi")) is None:
        pPr.append(OxmlElement("w:bidi"))
    for run in para.runs:  # type: ignore[attr-defined]
        rPr = run._element.get_or_add_rPr()  # noqa: N806
        if rPr.find(qn("w:rtl")) is None:
            rPr.append(OxmlElement("w:rtl"))


def _inject_python_docx(  # noqa: PLR0912
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into a DOCX file via python-docx.

    When the translated text contains inline HTML formatting tags
    (``<b>``, ``<i>``, ``<u>``, ``<s>``, ``<a>``), ``_inject_html_runs``
    creates per-run formatting and hyperlink wrappers.  Otherwise falls
    back to ``_replace_paragraph_text``.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name; when RTL, every paragraph in
            the document is marked with ``<w:bidi/>`` and every run with
            ``<w:rtl/>``.
    """
    from docx import Document  # noqa: PLC0415

    from src.constants.languages import is_rtl_language  # noqa: PLC0415

    doc = Document(str(file_path))
    doc_part = doc.part
    is_rtl = is_rtl_language(target_lang)

    def _inject_para(para: object, text: str) -> None:
        """Injects translated text into a single paragraph."""
        if _FORMATTING_HTML_RE.search(text):
            _inject_html_runs(para, text, part=doc_part)
        else:
            _replace_paragraph_text(para, text)

    # Paragraphs
    for i, para in enumerate(doc.paragraphs):
        key = f"para:{i}"
        if key in translations:
            _inject_para(para, translations[key])

    # Tables — inject per-paragraph to match per-paragraph extraction
    for t_idx, table in enumerate(doc.tables):
        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                for p_idx, para in enumerate(cell.paragraphs):
                    key = f"table:{t_idx}:{r_idx}:{c_idx}:{p_idx}"
                    if key in translations:
                        _inject_para(para, translations[key])

    if is_rtl:
        # Mark every paragraph (including untranslated blanks) as RTL so
        # the whole document reads right-to-left.
        for para in doc.paragraphs:
            _set_docx_paragraph_rtl(para)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        _set_docx_paragraph_rtl(para)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))


def _extract_python_xlsx(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from an XLSX file via openpyxl.

    Iterates all sheets and collects cells with string values.

    Args:
        file_path: Path to the .xlsx file.

    Returns:
        list: (location_key, text) pairs.
    """
    from openpyxl import load_workbook  # noqa: PLC0415

    wb = load_workbook(str(file_path))
    texts: list[tuple[str, str]] = []

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value.strip():
                    texts.append(
                        (
                            f"sheet:{ws.title}:{cell.row}:{cell.column}",
                            cell.value,
                        ),
                    )

    wb.close()
    return texts


def _inject_python_xlsx(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into an XLSX file via openpyxl.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    from openpyxl import load_workbook  # noqa: PLC0415

    from src.constants.languages import is_rtl_language  # noqa: PLC0415

    wb = load_workbook(str(file_path))
    is_rtl = is_rtl_language(target_lang)

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                key = f"sheet:{ws.title}:{cell.row}:{cell.column}"
                if key in translations:
                    original = cell.value or ""
                    new_text = translations[key]
                    cell.value = new_text
                    # Substitute font name when scripts differ so the app
                    # uses a font that supports the target script.
                    if isinstance(original, str) and cell.font and cell.font.name:
                        new_font_name = _substitute_font(
                            cell.font.name,
                            original,
                            new_text,
                            target_lang,
                        )
                        if new_font_name is not None and new_font_name != (
                            cell.font.name or ""
                        ):
                            # ``cell.font`` is a StyleProxy whose .copy() is
                            # deprecated; ``copy(font)`` returns a real
                            # mutable Font we can edit then re-assign.
                            new_font = _shallow_copy(cell.font)
                            new_font.name = new_font_name
                            cell.font = new_font
        if is_rtl:
            # Flip the sheet view so columns read right→left.
            ws.sheet_view.rightToLeft = True

    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output_path))
    wb.close()

    # openpyxl drops embedded objects (xl/embeddings/) on save — restore them
    _restore_xlsx_embeddings(file_path, output_path)


def _walk_pptx_text_shapes(
    shapes: Any,  # noqa: ANN401 — python-pptx shape types are duck-typed (GroupShape vs leaf)
    parent_path: str = "",
) -> Generator[tuple[str, Any], None, None]:
    """Yield ``(shape_path, leaf_shape)`` for every text-bearing shape.

    Recurses into shape groups via duck-typing on ``.shapes``: a
    ``GroupShape`` exposes child shapes there, a regular text box
    doesn't.  The returned ``shape_path`` is a dotted index chain
    (``"0"``, ``"0.1"``, ``"2.0.3"``, …) so leaf positions stay stable
    across runs and survive the extract → inject round trip.
    """
    for i, shape in enumerate(shapes):
        path = f"{parent_path}.{i}" if parent_path else str(i)
        children = getattr(shape, "shapes", None)
        # GroupShape: recurse into children rather than skipping the
        # whole branch (groups themselves carry no text frame).
        if children is not None and not getattr(shape, "has_text_frame", False):
            yield from _walk_pptx_text_shapes(children, path)
            continue
        if getattr(shape, "has_text_frame", False):
            yield path, shape


def _extract_python_pptx(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from a PPTX file via python-pptx.

    Iterates slides and recurses through shape groups, then walks
    paragraphs and runs of every text frame.  Each non-empty paragraph
    gets a location key encoding the slide + dotted shape path + para
    index so grouped text round-trips through inject.  Paragraphs with
    mixed formatting or hyperlinks are encoded as inline HTML so the
    LLM can preserve formatting and ``<a>`` tags.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        list: (location_key, text) pairs — plain text or inline HTML.
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(file_path))
    texts: list[tuple[str, str]] = []

    for s_idx, slide in enumerate(prs.slides):
        for shape_path, shape in _walk_pptx_text_shapes(slide.shapes):
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                if not para.text.strip():
                    continue
                # Use HTML encoding when runs have mixed formatting
                # or hyperlinks
                if _has_pptx_mixed_formatting(para) or _has_pptx_hyperlinks(para):
                    text = _pptx_runs_to_html(para)
                else:
                    text = para.text
                texts.append(
                    (f"slide:{s_idx}:{shape_path}:{p_idx}", text),
                )

    return texts


def _set_pptx_paragraph_rtl(para: object) -> None:
    """Adds ``rtl="1"`` to a python-pptx paragraph's ``<a:pPr>``.

    PowerPoint and Keynote use this attribute to flip text-frame
    paragraph direction.  Idempotent.
    """
    pPr = para._pPr  # type: ignore[attr-defined]  # noqa: N806
    if pPr is None:
        pPr = para._p.get_or_add_pPr()  # type: ignore[attr-defined]  # noqa: N806
    pPr.set("rtl", "1")


def _inject_python_pptx(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into a PPTX file via python-pptx.

    For each translated paragraph: puts all text in the first run
    and clears other runs (preserves first run's formatting).
    When the translated text contains ``<a>`` tags, hyperlink
    relationships are created via the slide part.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language; when RTL, every paragraph in every
            text frame is marked with ``rtl="1"``.
    """
    from pptx import Presentation  # noqa: PLC0415

    from src.constants.languages import is_rtl_language  # noqa: PLC0415

    prs = Presentation(str(file_path))
    is_rtl = is_rtl_language(target_lang)

    for s_idx, slide in enumerate(prs.slides):
        slide_part = slide.part
        for shape_path, shape in _walk_pptx_text_shapes(slide.shapes):
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                key = f"slide:{s_idx}:{shape_path}:{p_idx}"
                if key in translations:
                    translated = translations[key]
                    if _FORMATTING_HTML_RE.search(translated):
                        _inject_pptx_html_runs(para, translated, part=slide_part)
                    elif para.runs:
                        para.runs[0].text = translated
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.text = translated
                if is_rtl:
                    _set_pptx_paragraph_rtl(para)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# -- odfpy extract/inject functions -----------------------------------------


def _extract_python_odt(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from an ODT file via odfpy.

    Extracts body paragraphs, headings, and table cell text. Paragraphs
    inside table cells are excluded from body paragraph counting (they
    are handled via the table iteration).

    Args:
        file_path: Path to the .odt file.

    Returns:
        list: (location_key, text) pairs.
    """
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.table import Table, TableCell, TableRow  # noqa: PLC0415
    from odf.text import H, P  # noqa: PLC0415

    doc = odf_load(str(file_path))
    texts: list[tuple[str, str]] = []
    body = doc.body

    # Paragraphs and headings (exclude those inside table cells)
    p_idx = 0
    for element in body.getElementsByType(P) + body.getElementsByType(H):
        if _is_inside_table_cell(element):
            continue
        text = _odf_element_text(element, preserve_links=True).strip()
        if text:
            texts.append((f"para:{p_idx}", text))
        p_idx += 1

    # Tables
    tables = body.getElementsByType(Table)
    for t_idx, table in enumerate(tables):
        for r_idx, row in enumerate(table.getElementsByType(TableRow)):
            c_idx = 0
            for cell in row.getElementsByType(TableCell):
                repeat = int(
                    cell.getAttribute("numbercolumnsrepeated") or "1",
                )
                cell_paras = cell.getElementsByType(P)
                cell_text = ""
                if cell_paras:
                    cell_text = "\n".join(
                        _odf_element_text(p, preserve_links=True) for p in cell_paras
                    ).strip()
                if cell_text:
                    texts.append(
                        (f"table:{t_idx}:{r_idx}:{c_idx}", cell_text),
                    )
                c_idx += repeat

    return texts


def _inject_python_odt(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into an ODT file via odfpy.

    For paragraphs and headings: replaces all child text with the
    translated text (inline formatting is not preserved, matching
    UNO backend behavior).
    For table cells: replaces text in the first paragraph element.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.table import Table, TableCell, TableRow  # noqa: PLC0415
    from odf.text import H, P  # noqa: PLC0415

    doc = odf_load(str(file_path))
    body = doc.body

    # Paragraphs and headings
    p_idx = 0
    for element in body.getElementsByType(P) + body.getElementsByType(H):
        if _is_inside_table_cell(element):
            continue
        key = f"para:{p_idx}"
        if key in translations:
            _odf_replace_text(
                element,
                translations[key],
                target_lang=target_lang,
            )
        p_idx += 1

    # Tables
    tables = body.getElementsByType(Table)
    for t_idx, table in enumerate(tables):
        for r_idx, row in enumerate(table.getElementsByType(TableRow)):
            c_idx = 0
            for cell in row.getElementsByType(TableCell):
                repeat = int(
                    cell.getAttribute("numbercolumnsrepeated") or "1",
                )
                key = f"table:{t_idx}:{r_idx}:{c_idx}"
                if key in translations:
                    cell_paras = cell.getElementsByType(P)
                    if cell_paras:
                        _odf_replace_text(
                            cell_paras[0],
                            translations[key],
                            target_lang=target_lang,
                        )
                c_idx += repeat

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))

    from src.constants.languages import is_rtl_language  # noqa: PLC0415

    if is_rtl_language(target_lang):
        _set_odf_default_rtl(output_path)


def _extract_python_ods(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from an ODS file via odfpy.

    Iterates all sheets and collects cells with string text content.
    Uses the same key format as _extract_python_xlsx:
    ``sheet:{name}:{row}:{col}`` with 1-based indices.

    Args:
        file_path: Path to the .ods file.

    Returns:
        list: (location_key, text) pairs.
    """
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.table import Table, TableCell, TableRow  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = odf_load(str(file_path))
    texts: list[tuple[str, str]] = []

    for table in doc.spreadsheet.getElementsByType(Table):
        sheet_name = table.getAttribute("name")
        for r_idx, row in enumerate(table.getElementsByType(TableRow)):
            # Skip large repeated empty rows (common in ODS files)
            row_repeat = int(
                row.getAttribute("numberrowsrepeated") or "1",
            )
            if row_repeat > 1:
                has_content = any(
                    _odf_element_text(p).strip()
                    for cell in row.getElementsByType(TableCell)
                    for p in cell.getElementsByType(P)
                )
                if not has_content:
                    continue

            c_idx = 0
            for cell in row.getElementsByType(TableCell):
                col_repeat = int(
                    cell.getAttribute("numbercolumnsrepeated") or "1",
                )
                value_type = cell.getAttribute("valuetype")
                paras = cell.getElementsByType(P)

                # Only extract string cells (match openpyxl behavior)
                if value_type == "string" and paras:
                    text = "\n".join(
                        _odf_element_text(p, preserve_links=True) for p in paras
                    ).strip()
                    if text:
                        texts.append(
                            (
                                f"sheet:{sheet_name}:{r_idx + 1}:{c_idx + 1}",
                                text,
                            ),
                        )

                c_idx += col_repeat

    return texts


def _inject_python_ods(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into an ODS file via odfpy.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.table import Table, TableCell, TableRow  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = odf_load(str(file_path))

    for table in doc.spreadsheet.getElementsByType(Table):
        sheet_name = table.getAttribute("name")
        for r_idx, row in enumerate(table.getElementsByType(TableRow)):
            row_repeat = int(
                row.getAttribute("numberrowsrepeated") or "1",
            )
            if row_repeat > 1:
                has_content = any(
                    _odf_element_text(p).strip()
                    for cell in row.getElementsByType(TableCell)
                    for p in cell.getElementsByType(P)
                )
                if not has_content:
                    continue

            c_idx = 0
            for cell in row.getElementsByType(TableCell):
                col_repeat = int(
                    cell.getAttribute("numbercolumnsrepeated") or "1",
                )
                key = f"sheet:{sheet_name}:{r_idx + 1}:{c_idx + 1}"
                if key in translations:
                    paras = cell.getElementsByType(P)
                    if paras:
                        _odf_replace_text(
                            paras[0],
                            translations[key],
                            target_lang=target_lang,
                        )
                c_idx += col_repeat

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))

    from src.constants.languages import is_rtl_language  # noqa: PLC0415

    if is_rtl_language(target_lang):
        _set_odf_default_rtl(output_path)


def _extract_python_odp(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from an ODP file via odfpy.

    Iterates presentation pages, draw frames, and paragraphs within.
    Each non-empty paragraph gets a location key using the same format
    as _extract_python_pptx: ``slide:{s}:{sh}:{p}``.

    Args:
        file_path: Path to the .odp file.

    Returns:
        list: (location_key, text) pairs.
    """
    from odf.draw import Frame, Page  # noqa: PLC0415
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = odf_load(str(file_path))
    texts: list[tuple[str, str]] = []

    for s_idx, page in enumerate(doc.getElementsByType(Page)):
        for sh_idx, frame in enumerate(page.getElementsByType(Frame)):
            for p_idx, para in enumerate(frame.getElementsByType(P)):
                text = _odf_element_text(
                    para,
                    preserve_links=True,
                ).strip()
                if text:
                    texts.append(
                        (f"slide:{s_idx}:{sh_idx}:{p_idx}", text),
                    )

    return texts


def _inject_python_odp(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Injects translations into an ODP file via odfpy.

    For each translated paragraph: replaces all text content,
    matching UNO backend behavior.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    from odf.draw import Frame, Page  # noqa: PLC0415
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = odf_load(str(file_path))

    for s_idx, page in enumerate(doc.getElementsByType(Page)):
        for sh_idx, frame in enumerate(page.getElementsByType(Frame)):
            for p_idx, para in enumerate(frame.getElementsByType(P)):
                key = f"slide:{s_idx}:{sh_idx}:{p_idx}"
                if key in translations:
                    _odf_replace_text(
                        para,
                        translations[key],
                        target_lang=target_lang,
                    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))

    from src.constants.languages import is_rtl_language  # noqa: PLC0415

    if is_rtl_language(target_lang):
        _set_odf_default_rtl(output_path)


# -- python_lib wrapper/router functions ------------------------------------


def _extract_python_word(file_path: Path) -> list[tuple[str, str]]:
    """Routes word-category extraction based on file extension.

    Args:
        file_path: Path to the document (.docx or .odt).

    Returns:
        list: (location_key, text) pairs.
    """
    if file_path.suffix.lower() == ".odt":
        return _extract_python_odt(file_path)
    return _extract_python_docx(file_path)


def _inject_python_word(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Routes word-category injection based on file extension.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    if file_path.suffix.lower() == ".odt":
        _inject_python_odt(file_path, output_path, translations, target_lang)
    else:
        _inject_python_docx(file_path, output_path, translations, target_lang)


def _extract_python_excel(file_path: Path) -> list[tuple[str, str]]:
    """Routes excel-category extraction based on file extension.

    Args:
        file_path: Path to the spreadsheet (.xlsx or .ods).

    Returns:
        list: (location_key, text) pairs.
    """
    if file_path.suffix.lower() == ".ods":
        return _extract_python_ods(file_path)
    return _extract_python_xlsx(file_path)


def _inject_python_excel(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Routes excel-category injection based on file extension.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    if file_path.suffix.lower() == ".ods":
        _inject_python_ods(file_path, output_path, translations, target_lang)
    else:
        _inject_python_xlsx(file_path, output_path, translations, target_lang)


def _extract_python_ppt(file_path: Path) -> list[tuple[str, str]]:
    """Routes ppt-category extraction based on file extension.

    Args:
        file_path: Path to the presentation (.pptx or .odp).

    Returns:
        list: (location_key, text) pairs.
    """
    if file_path.suffix.lower() == ".odp":
        return _extract_python_odp(file_path)
    return _extract_python_pptx(file_path)


def _inject_python_ppt(
    file_path: Path,
    output_path: Path,
    translations: dict[str, str],
    target_lang: str = "",
) -> None:
    """Routes ppt-category injection based on file extension.

    Args:
        file_path: Source file path.
        output_path: Output file path.
        translations: Mapping of location_key to translated text.
        target_lang: Target language name for font substitution.
    """
    if file_path.suffix.lower() == ".odp":
        _inject_python_odp(file_path, output_path, translations, target_lang)
    else:
        _inject_python_pptx(file_path, output_path, translations, target_lang)


# ---------------------------------------------------------------------------
# Backend dispatch tables
# ---------------------------------------------------------------------------

# Map file category → (extract_fn, inject_fn) per backend
_WORD_EXTENSIONS = {".doc", ".docx", ".odt"}
_EXCEL_EXTENSIONS = {".xls", ".xlsx", ".ods"}
_PPT_EXTENSIONS = {".ppt", ".pptx", ".odp"}

_EXTRACTORS: dict[str, dict[str, Callable]] = {
    _BACKEND_WIN32COM: {
        "word": _extract_win32com_word,
        "excel": _extract_win32com_excel,
        "ppt": _extract_win32com_ppt,
    },
    _BACKEND_UNO: {
        "word": _extract_uno_writer,
        "excel": _extract_uno_calc,
        "ppt": _extract_uno_impress,
    },
    _BACKEND_PYTHON_LIB: {
        "word": _extract_python_word,
        "excel": _extract_python_excel,
        "ppt": _extract_python_ppt,
    },
}

_INJECTORS: dict[str, dict[str, Callable]] = {
    _BACKEND_WIN32COM: {
        "word": _inject_win32com_word,
        "excel": _inject_win32com_excel,
        "ppt": _inject_win32com_ppt,
    },
    _BACKEND_UNO: {
        "word": _inject_uno_writer,
        "excel": _inject_uno_calc,
        "ppt": _inject_uno_impress,
    },
    _BACKEND_PYTHON_LIB: {
        "word": _inject_python_word,
        "excel": _inject_python_excel,
        "ppt": _inject_python_ppt,
    },
}


def _get_file_category(suffix: str) -> str:
    """Returns the file category for dispatch.

    Args:
        suffix: Lowercase file extension.

    Returns:
        str: "word", "excel", or "ppt".

    Raises:
        ValueError: If the extension is not an office format.
    """
    if suffix in _WORD_EXTENSIONS:
        return "word"
    if suffix in _EXCEL_EXTENSIONS:
        return "excel"
    if suffix in _PPT_EXTENSIONS:
        return "ppt"
    raise ValueError(f"Unsupported office extension: {suffix}")


# ---------------------------------------------------------------------------
# Image translation in documents
# ---------------------------------------------------------------------------

# MIME types of raster images supported for in-document translation
_SUPPORTED_IMAGE_TYPES = {
    "image/png",
    "image/jpeg",
    "image/bmp",
    "image/webp",
    "image/tiff",
}

# MIME type → file extension mapping
_IMAGE_TYPE_TO_EXT: dict[str, str] = {
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/bmp": ".bmp",
    "image/webp": ".webp",
    "image/tiff": ".tiff",
}

# File extension → MIME content type (reverse of _IMAGE_TYPE_TO_EXT)
_EXT_TO_IMAGE_CONTENT_TYPE: dict[str, str] = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".bmp": "image/bmp",
    ".webp": "image/webp",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
}

# Extensions that support embedded image translation
# Modern/ODF formats use zipfile directly; legacy formats use round-trip conversion
_IMAGE_EXTENSIONS = {
    ".docx",
    ".xlsx",
    ".pptx",
    ".odt",
    ".ods",
    ".odp",
    ".epub",
    ".doc",
    ".xls",
    ".ppt",
}

# Document suffix → media directory prefixes inside the ZIP archive.
# OOXML files normally store images under ``word/media/``, ``xl/media/``,
# or ``ppt/media/``.  However, some producers (e.g. certain Word versions)
# place them at the root ``media/`` directory and reference them via
# absolute paths in ``.rels`` (``Target="/media/image.png"``).  Both
# locations are valid OOXML, so we check multiple prefixes.
_SUFFIX_TO_MEDIA_PREFIXES: dict[str, tuple[str, ...]] = {
    ".docx": ("word/media/", "media/"),
    ".xlsx": ("xl/media/", "media/"),
    ".pptx": ("ppt/media/", "media/"),
    ".odt": ("Pictures/",),
    ".ods": ("Pictures/",),
    ".odp": ("Pictures/",),
    ".epub": ("",),  # images can live anywhere in EPUB archives
}

# LLM errors that will definitely repeat for every image — stop immediately.
# Matched against the BASE tag (before any ``:Service`` suffix the
# engine appends for AUTH_ERROR), so ``"AUTH_ERROR:Gemini"`` still
# qualifies as fatal — without this, the suffix-bearing variant
# would slip past the set-membership check and the loop would
# skip-with-warning through every image while the same auth error
# keeps repeating, finishing "Done" with zero translated images.
_FATAL_LLM_ERRORS = {"AUTH_ERROR", "QUOTA_ERROR", "VISION_NOT_SUPPORTED"}


def _is_fatal_llm_error(error_tag: str) -> bool:
    """Returns True when *error_tag* is in ``_FATAL_LLM_ERRORS``.

    Delegates to :func:`src.constants.errors.base_error_tag` to strip
    the optional ``:Service`` suffix the engine appends to AUTH_ERROR
    so ``"AUTH_ERROR:Gemini"`` matches as fatal alongside the bare
    ``"AUTH_ERROR"``.
    """
    return base_error_tag(error_tag) in _FATAL_LLM_ERRORS


def _should_translate_images(
    suffix: str,
    backend: str,
    config: TranslationConfig | None = None,
) -> bool:
    """Checks whether image translation should be attempted for this file.

    Returns True when the setting is enabled, OCR is configured, and the
    format supports embedded image translation.  Modern/ODF formats use
    ``zipfile`` directly; legacy formats (.doc, .xls, .ppt) use round-trip
    conversion to a modern format first.

    Args:
        suffix: Lowercase file extension (e.g. ".docx").
        backend: The detected backend identifier (unused, kept for API
                 consistency with ``_should_translate_comments``).
        config: Optional TranslationConfig snapshot; falls back to load_setting().

    Returns:
        bool: True if image translation should proceed.
    """
    if suffix not in _IMAGE_EXTENSIONS:
        return False

    if config is not None:
        return config.should_translate_images

    from src.constants.settings import SETTING_TRANSLATE_DOC_IMAGES  # noqa: PLC0415
    from src.utils.config_manager import (  # noqa: PLC0415
        check_ocr_setup,
        load_setting,
    )

    return bool(load_setting(SETTING_TRANSLATE_DOC_IMAGES, False)) and check_ocr_setup()


_COMMENT_EXTENSIONS = {
    ".docx",
    ".xlsx",
    ".pptx",
    ".odt",
    ".ods",
    ".odp",
    ".doc",
    ".xls",
    ".ppt",
}


def _should_translate_comments(
    suffix: str,
    backend: str,
    config: TranslationConfig | None = None,
) -> bool:
    """Checks whether comment translation should be attempted for this file.

    Returns True when the setting is enabled and the format supports
    comment extraction. Comment handling uses its own libraries
    (python-docx, openpyxl, python-pptx, zipfile+lxml) independently
    of the text-extraction backend, so no backend restriction is needed.

    Args:
        suffix: Lowercase file extension (e.g. ".docx").
        backend: The detected backend identifier (unused, kept for API
                 consistency with ``_should_translate_images``).
        config: Optional TranslationConfig snapshot; falls back to load_setting().

    Returns:
        bool: True if comment translation should proceed.
    """
    if suffix not in _COMMENT_EXTENSIONS:
        return False

    if config is not None:
        return config.translate_doc_comments

    from src.constants.settings import SETTING_TRANSLATE_DOC_COMMENTS  # noqa: PLC0415
    from src.utils.config_manager import load_setting  # noqa: PLC0415

    return bool(load_setting(SETTING_TRANSLATE_DOC_COMMENTS, False))


# Extensions that support shape / text-box translation.
# PPT formats are excluded — their main extractors already handle shapes.
_SHAPE_EXTENSIONS = {".docx", ".xlsx", ".odt", ".ods", ".doc", ".xls"}


def _should_translate_shapes(
    suffix: str,
    backend: str,
    config: TranslationConfig | None = None,
) -> bool:
    """Checks whether shape/text-box translation should be attempted.

    Returns True when the setting is enabled and the format supports
    shape extraction.  PPT formats are excluded because their primary
    extractors already handle shapes.

    Args:
        suffix: Lowercase file extension (e.g. ".docx").
        backend: The detected backend identifier (unused, kept for API
                 consistency with ``_should_translate_images``).
        config: Optional TranslationConfig snapshot; falls back to load_setting().

    Returns:
        bool: True if shape translation should proceed.
    """
    if suffix not in _SHAPE_EXTENSIONS:
        return False

    if config is not None:
        return config.translate_doc_shapes

    from src.constants.settings import SETTING_TRANSLATE_DOC_SHAPES  # noqa: PLC0415
    from src.utils.config_manager import load_setting  # noqa: PLC0415

    return bool(load_setting(SETTING_TRANSLATE_DOC_SHAPES, False))


def _extract_comments(  # noqa: PLR0911
    file_path: Path,
    suffix: str,
    backend: str,
) -> list[tuple[str, str]]:
    """Extracts comments from an office file.

    Args:
        file_path: Path to the office file.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.

    Returns:
        list: (location_key, text) pairs for comments.
    """
    # Modern formats (use python-docx / openpyxl / python-pptx / lxml)
    if suffix == ".docx":
        return _extract_docx_comments(file_path)
    if suffix == ".xlsx":
        return _extract_xlsx_comments(file_path)
    if suffix == ".pptx":
        return _extract_pptx_comments(file_path)
    if suffix in {".odt", ".ods", ".odp"}:
        return _extract_odf_comments(file_path)

    # Legacy formats (use win32com or UNO)
    if suffix == ".doc":
        if backend == _BACKEND_WIN32COM:
            return _extract_win32com_word_comments(file_path)
        return _extract_uno_writer_comments(file_path)
    if suffix == ".xls":
        if backend == _BACKEND_WIN32COM:
            return _extract_win32com_excel_comments(file_path)
        return _extract_uno_calc_comments(file_path)
    if suffix == ".ppt":
        if backend == _BACKEND_WIN32COM:
            return _extract_win32com_ppt_comments(file_path)
        return _extract_uno_impress_comments(file_path)

    return []


def _extract_docx_comments(file_path: Path) -> list[tuple[str, str]]:  # noqa: PLR0912, PLR0915
    """Extracts comments from a DOCX file via low-level XML access.

    Detects ``<w:hyperlink>`` elements within comment paragraphs and emits
    ``<a href="...">`` HTML tags so that hyperlinks are preserved through
    the LLM translation round-trip.  Hyperlink URLs are resolved from the
    comments part's ``.rels`` file.

    Args:
        file_path: Path to the .docx file.

    Returns:
        list: (location_key, text) pairs with keys like 'comment:{id}'.
              Text may contain ``<a>`` HTML when hyperlinks are present.
    """
    from docx import Document  # noqa: PLC0415
    from docx.oxml.ns import qn  # noqa: PLC0415

    texts: list[tuple[str, str]] = []
    doc = Document(str(file_path))

    try:
        comments_part = doc.part.package.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/comments"
        )
    except KeyError:
        return texts

    if comments_part is None:
        return texts

    # Parse comment hyperlink relationships from word/_rels/comments.xml.rels
    comment_hyperlink_rels: dict[str, str] = {}
    try:
        with zipfile.ZipFile(str(file_path), "r") as zf:
            rels_path = _get_rels_path("word/comments.xml")
            if rels_path in zf.namelist():
                comment_hyperlink_rels = _parse_hyperlink_rels(
                    zf.read(rels_path),
                )
    except Exception:  # noqa: BLE001
        pass

    w_hyperlink_tag = qn("w:hyperlink")
    w_r_tag = qn("w:r")
    w_t_tag = qn("w:t")
    r_id_attr = qn("r:id")

    root = etree.fromstring(comments_part.blob)
    for comment_el in root.findall(qn("w:comment")):
        c_id = comment_el.get(qn("w:id"))

        # Check if any paragraph contains <w:hyperlink> children
        p_elements = comment_el.findall(qn("w:p"))
        has_links = comment_hyperlink_rels and any(
            p_el.findall(w_hyperlink_tag) for p_el in p_elements
        )

        if has_links:
            # Build HTML preserving hyperlinks as <a> tags
            p_htmls: list[str] = []
            for p_el in p_elements:
                parts: list[str] = []
                for child in p_el:
                    if child.tag == w_hyperlink_tag:
                        r_id = child.get(r_id_attr, "")
                        url = comment_hyperlink_rels.get(r_id, "")
                        link_text = "".join(
                            t.text for t in child.iter(w_t_tag) if t.text
                        )
                        if url and link_text:
                            parts.append(
                                f'<a href="{html.escape(url, quote=True)}">'
                                f"{html.escape(link_text)}</a>",
                            )
                        elif link_text:
                            parts.append(html.escape(link_text))
                    elif child.tag == w_r_tag:
                        t_parts = [t.text for t in child.iter(w_t_tag) if t.text]
                        if t_parts:
                            parts.append(html.escape("".join(t_parts)))
                if parts:
                    p_htmls.append("".join(parts))
            text = "\n".join(p_htmls).strip()
        else:
            # Original plain-text extraction
            p_texts: list[str] = []
            for p_el in p_elements:
                t_parts = [t.text for t in p_el.iter(w_t_tag) if t.text]
                if t_parts:
                    p_texts.append("".join(t_parts))
            text = "\n".join(p_texts).strip()

        if text:
            texts.append((f"comment:{c_id}", text))

    return texts


def _extract_xlsx_comments(file_path: Path) -> list[tuple[str, str]]:
    """Extracts cell comments from an XLSX file via openpyxl.

    Args:
        file_path: Path to the .xlsx file.

    Returns:
        list: (location_key, text) pairs with keys like
              'comment:{sheet}:{row}:{col}'.
    """
    from openpyxl import load_workbook  # noqa: PLC0415

    texts: list[tuple[str, str]] = []
    wb = load_workbook(str(file_path))

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.comment and cell.comment.text and cell.comment.text.strip():
                    texts.append(
                        (
                            f"comment:{ws.title}:{cell.row}:{cell.column}",
                            cell.comment.text,
                        ),
                    )

    wb.close()
    return texts


def _inject_comments(  # noqa: PLR0912
    output_path: Path,
    translations: dict[str, str],
    suffix: str,
    backend: str,
) -> None:
    """Injects translated comments back into the output document.

    Args:
        output_path: Path to the output file (already written by inject_fn).
        translations: Mapping of location_key to translated text.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.
    """
    comment_keys = {k: v for k, v in translations.items() if k.startswith("comment:")}
    if not comment_keys:
        return

    # Modern formats
    if suffix == ".docx":
        _inject_docx_comments(output_path, comment_keys)
    elif suffix == ".xlsx":
        _inject_xlsx_comments(output_path, comment_keys)
    elif suffix == ".pptx":
        _inject_pptx_comments(output_path, comment_keys)
    elif suffix in {".odt", ".ods", ".odp"}:
        _inject_odf_comments(output_path, comment_keys)
    # Legacy formats
    elif suffix == ".doc":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_word_comments(output_path, comment_keys)
        else:
            _inject_uno_writer_comments(output_path, comment_keys)
    elif suffix == ".xls":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_excel_comments(output_path, comment_keys)
        else:
            _inject_uno_calc_comments(output_path, comment_keys)
    elif suffix == ".ppt":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_ppt_comments(output_path, comment_keys)
        else:
            _inject_uno_impress_comments(output_path, comment_keys)


def _inject_docx_comments(  # noqa: PLR0912, PLR0915
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into a DOCX file via low-level XML.

    When a translation contains ``<a href="...">`` tags, the comment's
    paragraphs are rebuilt with ``<w:hyperlink>`` elements and the
    corresponding relationships are added to
    ``word/_rels/comments.xml.rels``.  Plain-text translations use the
    simpler ``<w:t>`` replacement path.

    Args:
        output_path: Path to the .docx file to modify in place.
        translations: Mapping of 'comment:{id}' to translated text.
    """
    from docx import Document  # noqa: PLC0415
    from docx.oxml.ns import qn  # noqa: PLC0415

    doc = Document(str(output_path))

    try:
        comments_part = doc.part.package.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/comments"
        )
    except KeyError:
        return

    if comments_part is None:
        return

    # Check if any translation contains formatting HTML (hyperlinks, etc.)
    has_any_html = any(
        _FORMATTING_HTML_RE.search(v)
        for k, v in translations.items()
        if k.startswith("comment:")
    )

    root = etree.fromstring(comments_part.blob)
    for comment_el in root.findall(qn("w:comment")):
        c_id = comment_el.get(qn("w:id"))
        key = f"comment:{c_id}"
        if key not in translations:
            continue

        translated = translations[key]

        if _FORMATTING_HTML_RE.search(translated):
            # HTML with formatting/hyperlinks — rebuild comment paragraphs
            _inject_docx_comment_html(
                comment_el,
                translated,
                comments_part,
                qn,
            )
        else:
            # Plain-text injection — original approach
            t_elements = list(comment_el.iter(qn("w:t")))
            if t_elements:
                lines = translated.split("\n")
                first_t = t_elements[0]
                first_t.text = lines[0]
                parent_r = first_t.getparent()

                # Insert <w:br/> and new <w:t> for each remaining line.
                for line in lines[1:]:
                    br = etree.Element(qn("w:br"))
                    parent_r.append(br)
                    t_new = etree.Element(qn("w:t"))
                    t_new.text = line
                    parent_r.append(t_new)

                for t_el in t_elements[1:]:
                    t_el.text = ""

    # CommentsPart inherits from XmlPart, whose blob property serializes
    # from _element (not _blob). Update the element tree directly.
    comments_part._element = root

    # If hyperlinks were added, the rels file is updated via
    # comments_part.relate_to() which python-docx handles on save.
    doc.save(str(output_path))

    # When HTML translations were present, the rels file was updated via
    # comments_part.relate_to() above.  However, python-docx may not
    # persist comments.xml.rels for certain part types.  Verify and patch
    # the ZIP directly if needed.
    if has_any_html:
        _patch_docx_comment_rels(output_path, comments_part)


def _inject_docx_comment_html(  # noqa: PLR0912, PLR0915
    comment_el: object,
    html_text: str,
    comments_part: object,
    qn: object,
) -> None:
    """Rebuilds a single comment element's paragraphs from HTML.

    Parses ``html_text`` via ``_parse_html_formatting`` to obtain
    ``_FormattedSegment`` objects.  Segments with ``hyperlink_url`` are
    wrapped in ``<w:hyperlink>`` elements with relationship IDs created
    via ``comments_part.relate_to()``.

    Args:
        comment_el: The ``<w:comment>`` lxml element to modify in-place.
        html_text: Translated HTML string (may contain ``<a>`` tags).
        comments_part: The python-docx comments ``Part`` object, used to
            create OPC relationships for external hyperlinks.
        qn: The python-docx ``qn()`` namespace resolver.
    """
    import copy  # noqa: PLC0415

    segments = _parse_html_formatting(html_text)
    if not segments:
        # Strip residual tags — fall back to plain text
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        t_elements = list(comment_el.iter(qn("w:t")))
        if t_elements:
            t_elements[0].text = plain
            for t_el in t_elements[1:]:
                t_el.text = ""
        return

    w_p_tag = qn("w:p")
    w_r_tag = qn("w:r")
    w_rpr_tag = qn("w:rPr")
    w_t_tag = qn("w:t")
    w_hyperlink_tag = qn("w:hyperlink")
    r_id_attr = qn("r:id")
    w_anchor_attr = qn("w:anchor")
    xml_space_attr = "{http://www.w3.org/XML/1998/namespace}space"

    # Save the base rPr from the first run (font, theme, etc.)
    existing_paras = comment_el.findall(w_p_tag)
    base_rpr: object | None = None
    for p_el in existing_paras:
        for r_el in p_el.findall(w_r_tag):
            found = r_el.find(w_rpr_tag)
            if found is not None:
                base_rpr = copy.deepcopy(found)
                break
        if base_rpr is not None:
            break
    # Also check inside <w:hyperlink> for base rPr
    if base_rpr is None:
        for p_el in existing_paras:
            for hl in p_el.findall(w_hyperlink_tag):
                for r_el in hl.findall(w_r_tag):
                    found = r_el.find(w_rpr_tag)
                    if found is not None:
                        base_rpr = copy.deepcopy(found)
                        break
                if base_rpr is not None:
                    break
            if base_rpr is not None:
                break

    # Save pPr (paragraph properties) from the first paragraph
    base_ppr: object | None = None
    if existing_paras:
        ppr = existing_paras[0].find(qn("w:pPr"))
        if ppr is not None:
            base_ppr = copy.deepcopy(ppr)

    # Split segments into paragraph groups by '\n'
    para_groups: list[list[_FormattedSegment]] = [[]]
    for seg in segments:
        lines = seg.text.split("\n")
        for j, line_text in enumerate(lines):
            if j > 0:
                para_groups.append([])
            if line_text:
                para_groups[-1].append(seg._replace(text=line_text))
    # Remove trailing empty groups
    while para_groups and not para_groups[-1]:
        para_groups.pop()
    if not para_groups:
        return

    # Remove all existing <w:p> children from the comment
    for p_el in existing_paras:
        comment_el.remove(p_el)

    def _make_run(seg: _FormattedSegment) -> object:
        """Creates a ``<w:r>`` element from a formatted segment."""
        r_el = etree.Element(w_r_tag)
        if base_rpr is not None:
            r_el.append(copy.deepcopy(base_rpr))
        t_el = etree.Element(w_t_tag)
        t_el.text = seg.text
        if seg.text.startswith(" ") or seg.text.endswith(" "):
            t_el.set(xml_space_attr, "preserve")
        r_el.append(t_el)
        return r_el

    # Build new paragraphs
    for group in para_groups:
        p_el = etree.SubElement(comment_el, w_p_tag)
        if base_ppr is not None:
            p_el.append(copy.deepcopy(base_ppr))

        # Group segments by hyperlink URL and create runs
        current_url: str | None = None
        hyperlink_elem: object | None = None

        for seg in group:
            new_r = _make_run(seg)

            if seg.hyperlink_url:
                if seg.hyperlink_url != current_url:
                    # Start a new <w:hyperlink> group
                    hyperlink_elem = etree.Element(w_hyperlink_tag)
                    if seg.hyperlink_url.startswith("#"):
                        # Internal bookmark anchor
                        hyperlink_elem.set(
                            w_anchor_attr,
                            seg.hyperlink_url[1:],
                        )
                    else:
                        # External URL — create OPC relationship
                        try:
                            r_id = comments_part.relate_to(
                                seg.hyperlink_url,
                                _HYPERLINK_RELTYPE,
                                is_external=True,
                            )
                            hyperlink_elem.set(r_id_attr, r_id)
                        except Exception:  # noqa: BLE001
                            # Cannot create relationship — fall back to
                            # plain run
                            hyperlink_elem = None
                    if hyperlink_elem is not None:
                        p_el.append(hyperlink_elem)
                    current_url = seg.hyperlink_url
                if hyperlink_elem is not None:
                    hyperlink_elem.append(new_r)
                else:
                    # Fallback: attach as plain run
                    p_el.append(new_r)
            else:
                if current_url is not None:
                    current_url = None
                    hyperlink_elem = None
                p_el.append(new_r)


def _patch_docx_comment_rels(
    output_path: Path,
    comments_part: object,
) -> None:
    """Ensures ``word/_rels/comments.xml.rels`` is persisted in the DOCX ZIP.

    ``python-docx`` may not serialize ``.rels`` for the comments part
    when saved via ``doc.save()``.  This function verifies and patches
    the ZIP directly if the rels data is missing or stale.

    Args:
        output_path: Path to the saved .docx file.
        comments_part: The python-docx comments Part (with .rels data).
    """
    # Access the rels XML that python-docx should have written
    try:
        rels = comments_part.rels
        if not rels:
            return
    except Exception:  # noqa: BLE001
        return

    # Check for any hyperlink rels
    hyperlink_rels: list[tuple[str, str]] = []
    for rel in rels.values():
        try:
            if rel.reltype == _HYPERLINK_RELTYPE and rel.is_external:
                hyperlink_rels.append((rel.rId, rel.target_ref))
        except Exception:  # noqa: BLE001
            continue

    if not hyperlink_rels:
        return

    # Verify the ZIP actually contains the rels file
    rels_path = "word/_rels/comments.xml.rels"
    try:
        with zipfile.ZipFile(str(output_path), "r") as zf:
            if rels_path in zf.namelist():
                # Already present — python-docx handled it
                return
    except Exception:  # noqa: BLE001
        return

    # Not present — build rels XML and inject into the ZIP
    rels_xml: bytes | None = None
    for _r_id, url in hyperlink_rels:
        rels_xml, _ = _add_hyperlink_to_rels(rels_xml, url)

    if rels_xml is None:
        return

    try:
        with zipfile.ZipFile(str(output_path), "r") as zf:
            all_items = zf.infolist()
            file_data = {item.filename: zf.read(item.filename) for item in all_items}

        # Add the rels entry
        file_data[rels_path] = rels_xml
        # Create a ZipInfo for the new entry
        import zipfile as zf_mod  # noqa: PLC0415

        rels_info = zf_mod.ZipInfo(rels_path)
        all_items.append(rels_info)

        _rewrite_zip_content(output_path, file_data, all_items)
    except Exception:  # noqa: BLE001
        logger.warning("Failed to patch comments.xml.rels into %s", output_path)


def _inject_xlsx_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into an XLSX file via openpyxl.

    Args:
        output_path: Path to the .xlsx file to modify in place.
        translations: Mapping of 'comment:{sheet}:{row}:{col}' to translated text.
    """
    from openpyxl import load_workbook  # noqa: PLC0415

    wb = load_workbook(str(output_path))

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                key = f"comment:{ws.title}:{cell.row}:{cell.column}"
                if key in translations and cell.comment:
                    cell.comment.text = translations[key]

    wb.save(str(output_path))
    wb.close()


# Relationship type for legacy slide comments (PowerPoint 2007–2019)
_COMMENTS_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
)

# Relationship type for modern threaded comments (PowerPoint 365, 2021+)
_MODERN_COMMENTS_REL_TYPE = (
    "http://schemas.microsoft.com/office/2018/10/relationships/comments"
)

# PresentationML namespace (legacy comments)
_PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Modern comment namespace (PowerPoint 2018/8)
_PPTX_MODERN_NS = "http://schemas.microsoft.com/office/powerpoint/2018/8/main"

# WordprocessingShape namespace (DOCX text boxes inside <wps:txbx>)
_WPS_NS = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"

# SpreadsheetML namespace (for sheet/drawing relationships)
_SPREADSHEETML_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"

# SpreadsheetML relationship namespace
_RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _get_rels_path(part_path: str) -> str:
    """Returns the ``.rels`` path for a given XML part path inside a ZIP.

    E.g. ``'word/document.xml'`` → ``'word/_rels/document.xml.rels'``.

    Args:
        part_path: Path of the XML part inside the ZIP.

    Returns:
        Path of the corresponding ``.rels`` file.
    """
    parent, name = part_path.rsplit("/", 1) if "/" in part_path else ("", part_path)
    prefix = f"{parent}/_rels/" if parent else "_rels/"
    return f"{prefix}{name}.rels"


def _parse_hyperlink_rels(rels_xml: bytes) -> dict[str, str]:
    """Parses a ``.rels`` XML file into ``{r_id: url}`` for hyperlinks.

    Only external hyperlink relationships (``TargetMode="External"``) are
    included.

    Args:
        rels_xml: Raw bytes of the ``.rels`` file.

    Returns:
        dict mapping relationship IDs to target URLs.
    """
    root = etree.fromstring(rels_xml)
    rels: dict[str, str] = {}
    for rel in root:
        if (
            rel.get("Type") == _HYPERLINK_RELTYPE
            and rel.get("TargetMode") == "External"
        ):
            rels[rel.get("Id", "")] = rel.get("Target", "")
    return rels


def _add_hyperlink_to_rels(
    rels_xml: bytes | None,
    url: str,
) -> tuple[bytes, str]:
    """Adds a hyperlink relationship to a ``.rels`` XML file.

    If *rels_xml* is ``None``, creates a new ``Relationships`` document.

    Args:
        rels_xml: Existing ``.rels`` XML bytes, or ``None``.
        url: The target URL for the hyperlink.

    Returns:
        Tuple of ``(updated_rels_xml_bytes, new_r_id)``.
    """
    if rels_xml is None:
        root = etree.Element(
            f"{{{_RELS_NS}}}Relationships",
            nsmap={None: _RELS_NS},
        )
    else:
        root = etree.fromstring(rels_xml)

    # Find next available rId
    existing_ids = {rel.get("Id") for rel in root}
    counter = 1
    while f"rId{counter}" in existing_ids:
        counter += 1
    r_id = f"rId{counter}"

    # Add new relationship element
    rel_el = etree.SubElement(root, f"{{{_RELS_NS}}}Relationship")
    rel_el.set("Id", r_id)
    rel_el.set("Type", _HYPERLINK_RELTYPE)
    rel_el.set("Target", url)
    rel_el.set("TargetMode", "External")

    new_xml = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )
    return new_xml, r_id


def _extract_drawingml_text(tx_body_el: object) -> str:
    """Extracts plain text from a DrawingML ``<txBody>`` element.

    Iterates ``<a:p>`` paragraphs and joins ``<a:t>`` runs within each.
    Paragraphs are separated by newlines, and ``<a:br/>`` tags are preserved
    as newlines.

    Args:
        tx_body_el: An lxml element representing ``<txBody>``.

    Returns:
        str: The concatenated plain text.
    """
    a_p_tag = f"{{{_DRAWINGML_NS}}}p"
    a_t_tag = f"{{{_DRAWINGML_NS}}}t"
    a_br_tag = f"{{{_DRAWINGML_NS}}}br"

    paragraphs: list[str] = []
    for p_el in tx_body_el.findall(a_p_tag):
        parts: list[str] = []
        for child in p_el.iter():
            if child.tag == a_t_tag and child.text:
                parts.append(child.text)
            elif child.tag == a_br_tag:
                # Preserve explicit DrawingML line breaks (<a:br/>) as newlines.
                parts.append("\n")
        if parts:
            paragraphs.append("".join(parts))
    return "\n".join(paragraphs)


def _inject_drawingml_text(tx_body_el: object, new_text: str) -> None:
    """Replaces text in a DrawingML ``<txBody>`` element.

    Puts all translated text in the first ``<a:t>`` of the first
    ``<a:r>`` in the first ``<a:p>``, and clears remaining ``<a:t>``
    elements. Handles newlines by inserting ``<a:br/>`` and new ``<a:r>``
    elements.

    Args:
        tx_body_el: An lxml element representing ``<txBody>``.
        new_text: The replacement text.
    """
    a_t_tag = f"{{{_DRAWINGML_NS}}}t"
    a_br_tag = f"{{{_DRAWINGML_NS}}}br"

    t_elements = list(tx_body_el.iter(a_t_tag))
    if t_elements:
        lines = new_text.split("\n")
        first_t = t_elements[0]
        first_t.text = lines[0]
        parent_r = first_t.getparent()
        parent_p = parent_r.getparent()

        current_r = parent_r
        # Insert <a:br/> and new <a:r> for each remaining line.
        for line in lines[1:]:
            br = etree.Element(a_br_tag)
            idx = parent_p.index(current_r)
            parent_p.insert(idx + 1, br)

            new_r = etree.Element(parent_r.tag)
            new_t = etree.Element(a_t_tag)
            new_t.text = line
            new_r.append(new_t)
            parent_p.insert(idx + 2, new_r)

            current_r = new_r

        for t_el in t_elements[1:]:
            t_el.text = ""


def _inject_drawingml_html_runs(  # noqa: PLR0912, PLR0915
    tx_body_el: object,
    html_text: str,
    rels_adder: Callable[[str], str] | None = None,
) -> None:
    """Replaces DrawingML ``<a:txBody>`` runs with HTML-formatted segments.

    Parses ``html_text`` via ``_parse_html_formatting``, clears existing
    ``<a:r>`` elements, and rebuilds runs with per-segment ``<a:rPr>``
    formatting.  Falls back to ``_inject_drawingml_text`` if no HTML
    formatting tags are detected.

    When *rels_adder* is provided, segments with ``hyperlink_url`` get an
    ``<a:hlinkClick>`` element created inside ``<a:rPr>`` with a
    relationship ID returned by the callback.

    Args:
        tx_body_el: An lxml element representing ``<a:txBody>``.
        html_text: Translated text with inline ``<b>/<i>/<u>/<s>/<a>`` tags.
        rels_adder: Callback that accepts a URL string and returns a
            relationship ID (``r:id``) for the hyperlink.  ``None`` disables
            hyperlink injection.
    """
    import copy  # noqa: PLC0415

    # Fallback: no formatting tags → plain text
    if not _FORMATTING_HTML_RE.search(html_text):
        _inject_drawingml_text(tx_body_el, html_text)
        return

    segments = _parse_html_formatting(html_text)
    if not segments:
        # Strip residual tags so literal HTML doesn't appear in the document
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        _inject_drawingml_text(tx_body_el, plain)
        return

    a_p_tag = f"{{{_DRAWINGML_NS}}}p"
    a_r_tag = f"{{{_DRAWINGML_NS}}}r"
    a_rpr_tag = f"{{{_DRAWINGML_NS}}}rPr"
    a_t_tag = f"{{{_DRAWINGML_NS}}}t"
    a_hlink_tag = f"{{{_DRAWINGML_NS}}}hlinkClick"
    r_id_attr = (
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )

    p_elements = tx_body_el.findall(a_p_tag)
    if not p_elements:
        return

    # Save base rPr from the first run (deep copy, strip formatting attrs)
    base_rpr = None
    a_highlight_tag = f"{{{_DRAWINGML_NS}}}highlight"
    for r_el in p_elements[0].findall(a_r_tag):
        rpr = r_el.find(a_rpr_tag)
        if rpr is not None:
            base_rpr = copy.deepcopy(rpr)
            for attr in _DRAWINGML_FORMAT_ATTRS:
                base_rpr.attrib.pop(attr, None)
            # Strip baseline (superscript/subscript) from base
            base_rpr.attrib.pop("baseline", None)
            # Strip highlight from base so bg doesn't spread to all runs
            hl = base_rpr.find(a_highlight_tag)
            if hl is not None:
                base_rpr.remove(hl)
            # Strip hlinkClick from base so hyperlinks don't spread
            hlink = base_rpr.find(a_hlink_tag)
            if hlink is not None:
                base_rpr.remove(hlink)
            break

    # Split segments into paragraph groups (newlines → new <a:p>)
    para_groups: list[list[_FormattedSegment]] = [[]]
    for seg in segments:
        if "\n" in seg.text:
            parts = seg.text.split("\n")
            for idx, part in enumerate(parts):
                if idx > 0:
                    para_groups.append([])
                if part:
                    para_groups[-1].append(seg._replace(text=part))
        else:
            para_groups[-1].append(seg)

    # Save the first <a:p> as a template (for pPr / endParaRPr inheritance),
    # then remove all existing <a:p> elements from <a:txBody>.
    a_ppr_tag = f"{{{_DRAWINGML_NS}}}pPr"
    first_p = p_elements[0]
    saved_ppr = first_p.find(a_ppr_tag)
    if saved_ppr is not None:
        saved_ppr = copy.deepcopy(saved_ppr)

    for p_el in p_elements:
        tx_body_el.remove(p_el)

    # Rebuild <a:p> elements — one per paragraph group
    for group in para_groups:
        new_p = etree.SubElement(tx_body_el, a_p_tag)
        if saved_ppr is not None:
            new_p.insert(0, copy.deepcopy(saved_ppr))

        for seg in group:
            new_r = etree.SubElement(new_p, a_r_tag)
            # Build rPr: start from base copy, then apply formatting
            if base_rpr is not None:
                new_rpr = copy.deepcopy(base_rpr)
                new_r.insert(0, new_rpr)
            else:
                new_rpr = etree.SubElement(new_r, a_rpr_tag)

            _apply_drawingml_format_attrs(new_rpr, seg)

            # Add <a:hlinkClick> for hyperlinks
            if seg.hyperlink_url and rels_adder is not None:
                r_id = rels_adder(seg.hyperlink_url)
                hlink_elem = etree.SubElement(new_rpr, a_hlink_tag)
                hlink_elem.set(r_id_attr, r_id)

            # Create <a:t> with whitespace preservation
            new_t = etree.SubElement(new_r, a_t_tag)
            new_t.text = seg.text
            new_t.set(
                "{http://www.w3.org/XML/1998/namespace}space",
                "preserve",
            )


def _extract_pptx_legacy_comments(
    prs: object,
) -> list[tuple[str, str]]:
    """Extracts legacy comments from an already-opened Presentation.

    Legacy comments use ``<p:cm>`` elements with ``<p:text>`` children
    (PowerPoint 2007–2019).

    Args:
        prs: A python-pptx Presentation object.

    Returns:
        list: (location_key, text) pairs with keys like
              ``'comment:{slide_idx}:{comment_idx}'``.
    """
    texts: list[tuple[str, str]] = []

    for s_idx, slide in enumerate(prs.slides):
        for rel in slide.part.rels.values():
            if rel.reltype != _COMMENTS_REL_TYPE:
                continue

            root = etree.fromstring(rel.target_part.blob)
            cm_tag = f"{{{_PPTX_NS}}}cm"
            text_tag = f"{{{_PPTX_NS}}}text"

            for cm_el in root.iter(cm_tag):
                idx = cm_el.get("idx", "0")
                text_el = cm_el.find(text_tag)
                if text_el is not None and text_el.text and text_el.text.strip():
                    texts.append(
                        (f"comment:{s_idx}:{idx}", text_el.text),
                    )

    return texts


def _extract_pptx_modern_comments(
    prs: object,
) -> list[tuple[str, str]]:
    """Extracts modern threaded comments from an already-opened Presentation.

    Modern comments use ``<p188:cm>`` elements with ``<txBody>`` rich
    text and an optional ``<replyLst>`` (PowerPoint 365, 2021+).

    Args:
        prs: A python-pptx Presentation object.

    Returns:
        list: (location_key, text) pairs. Main comments use keys like
              ``'comment:{slide_idx}:{cm_id}'``; replies use
              ``'comment:{slide_idx}:{cm_id}:reply:{reply_id}'``.
    """
    texts: list[tuple[str, str]] = []
    cm_tag = f"{{{_PPTX_MODERN_NS}}}cm"
    tx_body_tag = f"{{{_PPTX_MODERN_NS}}}txBody"
    reply_lst_tag = f"{{{_PPTX_MODERN_NS}}}replyLst"
    reply_tag = f"{{{_PPTX_MODERN_NS}}}reply"

    for s_idx, slide in enumerate(prs.slides):
        for rel in slide.part.rels.values():
            if rel.reltype != _MODERN_COMMENTS_REL_TYPE:
                continue

            root = etree.fromstring(rel.target_part.blob)

            for cm_el in root.iter(cm_tag):
                cm_id = cm_el.get("id", "")

                # Extract main comment text from <txBody>
                tx_body = cm_el.find(tx_body_tag)
                if tx_body is not None:
                    text = _extract_drawingml_text(tx_body).strip()
                    if text:
                        texts.append(
                            (f"comment:{s_idx}:{cm_id}", text),
                        )

                # Extract reply text from <replyLst>/<reply>
                reply_lst = cm_el.find(reply_lst_tag)
                if reply_lst is not None:
                    for reply_el in reply_lst.findall(reply_tag):
                        reply_id = reply_el.get("id", "")
                        reply_body = reply_el.find(tx_body_tag)
                        if reply_body is not None:
                            reply_text = _extract_drawingml_text(
                                reply_body,
                            ).strip()
                            if reply_text:
                                texts.append(
                                    (
                                        f"comment:{s_idx}:{cm_id}:reply:{reply_id}",
                                        reply_text,
                                    )
                                )

    return texts


def _extract_pptx_comments(file_path: Path) -> list[tuple[str, str]]:
    """Extracts comments from a PPTX file via low-level XML on slide parts.

    Handles both legacy comments (``<p:cm>``) and modern threaded
    comments (``<p188:cm>``). A single file typically uses one format,
    but both are checked.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        list: (location_key, text) pairs.
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(file_path))
    texts: list[tuple[str, str]] = []
    texts.extend(_extract_pptx_legacy_comments(prs))
    texts.extend(_extract_pptx_modern_comments(prs))
    return texts


def _inject_pptx_legacy_comments(
    prs: object,
    translations: dict[str, str],
) -> bool:
    """Injects translated text into legacy PPTX comments.

    Args:
        prs: A python-pptx Presentation object.
        translations: Mapping of location keys to translated text.

    Returns:
        bool: True if any comment was modified.
    """
    modified = False

    for s_idx, slide in enumerate(prs.slides):
        for rel in slide.part.rels.values():
            if rel.reltype != _COMMENTS_REL_TYPE:
                continue

            root = etree.fromstring(rel.target_part.blob)
            cm_tag = f"{{{_PPTX_NS}}}cm"
            text_tag = f"{{{_PPTX_NS}}}text"
            part_modified = False

            for cm_el in root.iter(cm_tag):
                idx = cm_el.get("idx", "0")
                key = f"comment:{s_idx}:{idx}"
                if key in translations:
                    text_el = cm_el.find(text_tag)
                    if text_el is not None:
                        text_el.text = translations[key]
                        part_modified = True

            if part_modified:
                rel.target_part._blob = etree.tostring(
                    root,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone=True,
                )
                modified = True

    return modified


def _inject_pptx_modern_comments(
    prs: object,
    translations: dict[str, str],
) -> bool:
    """Injects translated text into modern threaded PPTX comments.

    Args:
        prs: A python-pptx Presentation object.
        translations: Mapping of location keys to translated text.

    Returns:
        bool: True if any comment was modified.
    """
    modified = False
    cm_tag = f"{{{_PPTX_MODERN_NS}}}cm"
    tx_body_tag = f"{{{_PPTX_MODERN_NS}}}txBody"
    reply_lst_tag = f"{{{_PPTX_MODERN_NS}}}replyLst"
    reply_tag = f"{{{_PPTX_MODERN_NS}}}reply"

    for s_idx, slide in enumerate(prs.slides):
        for rel in slide.part.rels.values():
            if rel.reltype != _MODERN_COMMENTS_REL_TYPE:
                continue

            root = etree.fromstring(rel.target_part.blob)
            part_modified = False

            for cm_el in root.iter(cm_tag):
                cm_id = cm_el.get("id", "")

                # Inject main comment text
                key = f"comment:{s_idx}:{cm_id}"
                if key in translations:
                    tx_body = cm_el.find(tx_body_tag)
                    if tx_body is not None:
                        _inject_drawingml_text(
                            tx_body,
                            translations[key],
                        )
                        part_modified = True

                # Inject reply text
                reply_lst = cm_el.find(reply_lst_tag)
                if reply_lst is not None:
                    for reply_el in reply_lst.findall(reply_tag):
                        reply_id = reply_el.get("id", "")
                        rkey = f"comment:{s_idx}:{cm_id}:reply:{reply_id}"
                        if rkey in translations:
                            reply_body = reply_el.find(tx_body_tag)
                            if reply_body is not None:
                                _inject_drawingml_text(
                                    reply_body,
                                    translations[rkey],
                                )
                                part_modified = True

            if part_modified:
                rel.target_part._blob = etree.tostring(
                    root,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone=True,
                )
                modified = True

    return modified


def _inject_pptx_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into a PPTX file via low-level XML.

    Handles both legacy and modern threaded comment formats.

    Args:
        output_path: Path to the .pptx file to modify in place.
        translations: Mapping of location keys to translated text.
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(output_path))
    modified = _inject_pptx_legacy_comments(prs, translations)
    modified |= _inject_pptx_modern_comments(prs, translations)
    if modified:
        prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Shared ZIP helpers
# ---------------------------------------------------------------------------


def _rewrite_zip_content(
    output_path: Path,
    file_data: dict[str, bytes],
    all_items: list[object],
) -> None:
    """Atomically rewrites a ZIP archive with modified file data.

    Writes to a temporary file then replaces the original.  Used by all
    zip-based inject functions (DOCX/XLSX/ODF shapes and comments).

    Args:
        output_path: Path to the ZIP file to overwrite.
        file_data: Mapping of archive entry names to their (possibly
                   modified) content bytes.
        all_items: Original ``ZipInfo`` list from the source archive,
                   used to preserve compression metadata and entry order.
    """
    tmp_path = output_path.with_suffix(output_path.suffix + ".tmp")
    try:
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zf_out:
            for item in all_items:
                zf_out.writestr(item, file_data[item.filename])
        shutil.move(str(tmp_path), str(output_path))
    except Exception:
        tmp_path.unlink(missing_ok=True)
        raise


# Relationship types for embedded OLE/package objects
_OLE_OBJECT_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject"
)
_PACKAGE_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package"
)
_EMBED_REL_TYPES = {_OLE_OBJECT_RELTYPE, _PACKAGE_RELTYPE}


def _patch_rels_for_embeddings(
    file_data: dict[str, bytes],
    src_rels: dict[str, bytes],
    new_items: list[zipfile.ZipInfo],
) -> None:
    """Restores embedding relationship entries into output rels files.

    Args:
        file_data: Mutable mapping of output ZIP entries (modified in place).
        src_rels: Source ``xl/_rels/*.rels`` files that may reference
                  embedding targets.
        new_items: Accumulator for new ZIP entries (appended if a rels
                   file is entirely missing from *file_data*).
    """
    for rels_name, src_rels_bytes in src_rels.items():
        src_rels_xml = etree.fromstring(src_rels_bytes)
        embed_rels = [
            rel for rel in src_rels_xml if rel.get("Type") in _EMBED_REL_TYPES
        ]
        if not embed_rels:
            continue

        if rels_name in file_data:
            out_rels_xml = etree.fromstring(file_data[rels_name])
        else:
            out_rels_xml = etree.Element(
                "Relationships",
                xmlns=_RELS_NS,
            )
            new_items.append(zipfile.ZipInfo(rels_name))

        existing_ids = {r.get("Id") for r in out_rels_xml}
        for rel in embed_rels:
            if rel.get("Id") not in existing_ids:
                out_rels_xml.append(rel)
        file_data[rels_name] = etree.tostring(
            out_rels_xml,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )


def _restore_xlsx_embeddings(
    source_path: Path,
    output_path: Path,
) -> None:
    """Restores embedded objects that openpyxl drops during save.

    openpyxl does not preserve OLE/package embedded objects stored
    under ``xl/embeddings/`` or their relationship and content-type
    entries.  This function reads those artefacts from *source_path*
    and patches them back into *output_path* after openpyxl's save.

    Args:
        source_path: Original XLSX before openpyxl processing.
        output_path: XLSX written by openpyxl (modified in place).
    """
    # 1. Identify embedding entries in the source ZIP
    with zipfile.ZipFile(source_path, "r") as src_zf:
        src_names = set(src_zf.namelist())
        embed_names = [n for n in src_names if n.startswith("xl/embeddings/")]
        if not embed_names:
            return
        embed_data = {n: src_zf.read(n) for n in embed_names}
        # Read source rels files that may reference embeddings
        src_rels: dict[str, bytes] = {}
        for n in src_names:
            if n.startswith("xl/_rels/") and n.endswith(".rels"):
                src_rels[n] = src_zf.read(n)
        src_ct = src_zf.read("[Content_Types].xml")

    # 2. Read the output ZIP
    with zipfile.ZipFile(output_path, "r") as out_zf:
        all_items = out_zf.infolist()
        file_data: dict[str, bytes] = {
            item.filename: out_zf.read(item.filename) for item in all_items
        }

    # 3. Restore embedding files
    new_items: list[zipfile.ZipInfo] = []
    for name, data in embed_data.items():
        if name not in file_data:
            file_data[name] = data
            new_items.append(zipfile.ZipInfo(name))

    if not new_items:
        return  # All embeddings already present

    # 4. Restore relationship entries pointing to embeddings
    _patch_rels_for_embeddings(file_data, src_rels, new_items)

    # 5. Restore content-type overrides for embeddings
    src_ct_xml = etree.fromstring(src_ct)
    if "[Content_Types].xml" in file_data:
        out_ct_xml = etree.fromstring(file_data["[Content_Types].xml"])
        existing_parts = {o.get("PartName") for o in out_ct_xml if o.get("PartName")}
        for override in src_ct_xml:
            part_name = override.get("PartName", "")
            if "/xl/embeddings/" in part_name and part_name not in existing_parts:
                out_ct_xml.append(override)
        file_data["[Content_Types].xml"] = etree.tostring(
            out_ct_xml,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )

    # 6. Atomic rewrite
    _rewrite_zip_content(output_path, file_data, all_items + new_items)


# ---------------------------------------------------------------------------
# Shared ODF paragraph helpers (used by both comments and shapes)
# ---------------------------------------------------------------------------


def _extract_odf_paragraph_text(  # noqa: PLR0912
    parent: object,
    text_p_tag: str,
) -> str:
    """Extracts concatenated paragraph text from an ODF element.

    Works with any element that contains ``<text:p>`` children, such as
    ``<draw:text-box>`` and ``<office:annotation>``.  Handles mixed
    content: direct text, child element text, and tail text.

    ODF hyperlinks (``<text:a xlink:href="URL">``) are emitted as
    ``<a href="URL">text</a>`` so that downstream HTML-aware injection
    can reconstruct them.

    Args:
        parent: An lxml element containing ``<text:p>`` children.
        text_p_tag: Fully-qualified ``{namespace}p`` tag string.

    Returns:
        str: Paragraphs joined by newlines, stripped.  May contain
             ``<a>`` HTML when hyperlinks are present.
    """
    text_ns = _ODF_NS["text"]
    text_a_tag = f"{{{text_ns}}}a"
    xlink_href = f"{{{_ODF_NS['xlink']}}}href"

    text_parts: list[str] = []
    for para in parent.findall(f".//{text_p_tag}"):
        parts: list[str] = []
        if para.text:
            parts.append(para.text)
        for child in para:
            if child.tag.endswith("line-break"):
                # WARNING: <text:line-break/> must be handled explicitly;
                # without this, adjacent lines are concatenated with no separator.
                parts.append("\n")
                if child.tail:
                    parts.append(child.tail)
            elif child.tag == text_a_tag:
                # ODF hyperlink — emit <a href="url">text</a>
                url = child.get(xlink_href, "")
                link_text = child.text or ""
                if url:
                    parts.append(
                        f'<a href="{html.escape(url, quote=True)}">'
                        f"{html.escape(link_text)}</a>",
                    )
                else:
                    parts.append(html.escape(link_text) if link_text else "")
                if child.tail:
                    parts.append(child.tail)
            else:
                if child.text:
                    parts.append(child.text)
                if child.tail:
                    parts.append(child.tail)
        if parts:
            text_parts.append("".join(parts))
    return "\n".join(text_parts).strip()


def _inject_odf_paragraph_text(
    parent: object,
    new_text: str,
    text_p_tag: str,
    *,
    target_lang: str = "",
) -> bool:
    """Replaces text in an ODF element that contains ``<text:p>`` children.

    Puts the translated text in the first ``<text:p>``, clears its
    children, and removes any extra ``<text:p>`` elements. Handles newlines
    by creating additional ``<text:p>`` elements. Works for
    both ``<draw:text-box>`` and ``<office:annotation>``.

    Preserves the first ``<text:span>``'s attributes so that character
    formatting (font name, size, bold, etc.) is retained.

    Args:
        parent: An lxml element containing ``<text:p>`` children.
        new_text: The replacement text.
        text_p_tag: Fully-qualified ``{namespace}p`` tag string.
        target_lang: Target language name for font substitution.

    Returns:
        bool: True if the element was modified.
    """
    paras = parent.findall(text_p_tag)
    if not paras:
        return False

    # Check for hyperlink HTML — route to the HTML-aware helper
    if _FORMATTING_HTML_RE.search(new_text):
        return _inject_odf_paragraph_text_html(
            parent,
            new_text,
            text_p_tag,
            paras,
            target_lang=target_lang,
        )

    # Derive the <text:span> tag from the <text:p> tag (same namespace)
    text_ns = text_p_tag.rsplit("}", 1)[0] + "}"
    text_span_tag = f"{text_ns}span"

    # Save the first <text:span>'s attributes before clearing
    first_span = paras[0].find(f".//{text_span_tag}")
    span_attribs: dict[str, str] | None = None
    if first_span is not None:
        span_attribs = dict(first_span.attrib)

    lines = new_text.split("\n")

    # Update first paragraph — wrap in span if formatting was found
    for child in list(paras[0]):
        paras[0].remove(child)
    if span_attribs:
        paras[0].text = None
        span_el = etree.SubElement(paras[0], text_span_tag, attrib=span_attribs)
        span_el.text = lines[0]
    else:
        paras[0].text = lines[0]

    # Remove remaining original paragraphs
    for extra_p in paras[1:]:
        parent.remove(extra_p)

    # WARNING: LibreOffice ignores literal '\n' inside a single <text:p>;
    # each line needs its own <text:p> element.
    for line in lines[1:]:
        new_p = etree.Element(text_p_tag)
        if span_attribs:
            span_el = etree.SubElement(new_p, text_span_tag, attrib=span_attribs)
            span_el.text = line
        else:
            new_p.text = line
        parent.append(new_p)

    return True


def _inject_odf_paragraph_text_html(  # noqa: PLR0912
    parent: object,
    new_text: str,
    text_p_tag: str,
    paras: list[object],
    *,
    target_lang: str = "",
) -> bool:
    """Injects HTML-formatted text (with hyperlinks) into ODF paragraphs.

    Parses *new_text* via ``_parse_html_formatting`` and reconstructs
    ``<text:p>`` children.  Segments with ``hyperlink_url`` become
    ``<text:a xlink:href="...">`` elements; plain segments become direct
    text or ``<text:span>`` elements (preserving character style).

    Falls back to plain-text injection when parsing yields no segments.

    Args:
        parent: An lxml element containing ``<text:p>`` children.
        new_text: Translated text with inline ``<a>`` HTML tags.
        text_p_tag: Fully-qualified ``{namespace}p`` tag string.
        paras: Pre-found ``<text:p>`` children from *parent*.
        target_lang: Target language name for font substitution.

    Returns:
        bool: True if the element was modified.
    """
    segments = _parse_html_formatting(new_text)
    if not segments:
        # Strip residual tags so literal HTML doesn't appear in the document
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", new_text))
        return _inject_odf_paragraph_text(
            parent,
            plain,
            text_p_tag,
            target_lang=target_lang,
        )

    text_ns = text_p_tag.rsplit("}", 1)[0] + "}"
    text_span_tag = f"{text_ns}span"
    text_a_tag = f"{text_ns}a"
    xlink_href = f"{{{_ODF_NS['xlink']}}}href"
    xlink_type = f"{{{_ODF_NS['xlink']}}}type"

    # Save the first <text:span>'s attributes for non-link text
    first_span = paras[0].find(f".//{text_span_tag}")
    span_attribs: dict[str, str] | None = None
    if first_span is not None:
        span_attribs = dict(first_span.attrib)

    # Split segments by newlines into paragraph groups
    para_groups: list[list[_FormattedSegment]] = [[]]
    for seg in segments:
        if "\n" in seg.text:
            sub_lines = seg.text.split("\n")
            for line_idx, line in enumerate(sub_lines):
                if line:
                    para_groups[-1].append(seg._replace(text=line))
                if line_idx < len(sub_lines) - 1:
                    para_groups.append([])
        else:
            para_groups[-1].append(seg)

    # Clear existing paragraphs
    for child in list(paras[0]):
        paras[0].remove(child)
    paras[0].text = None
    for extra_p in paras[1:]:
        parent.remove(extra_p)

    # Build paragraph content from segments
    def _fill_para(p_el: object, group: list[_FormattedSegment]) -> None:
        """Populates a ``<text:p>`` with segments.

        Creates ``<text:a>`` for hyperlinks and plain text or
        ``<text:span>`` for other content.
        """
        for seg in group:
            if seg.hyperlink_url:
                a_el = etree.SubElement(p_el, text_a_tag)
                a_el.set(xlink_href, seg.hyperlink_url)
                a_el.set(xlink_type, "simple")
                a_el.text = seg.text
            elif span_attribs:
                span_el = etree.SubElement(
                    p_el,
                    text_span_tag,
                    attrib=span_attribs,
                )
                span_el.text = seg.text
            else:
                # Append as direct text or tail of last child
                children = list(p_el)
                if children:
                    last = children[-1]
                    last.tail = (last.tail or "") + seg.text
                else:
                    p_el.text = (p_el.text or "") + seg.text

    # Fill first paragraph
    _fill_para(paras[0], para_groups[0] if para_groups else [])

    # Create additional paragraphs for remaining groups
    for group in para_groups[1:]:
        new_p = etree.Element(text_p_tag)
        _fill_para(new_p, group)
        parent.append(new_p)

    return True


# ---------------------------------------------------------------------------
# ODF (.odt / .ods / .odp) comment helpers — ZIP + lxml
# ---------------------------------------------------------------------------


def _extract_odf_comments(file_path: Path) -> list[tuple[str, str]]:
    """Extracts annotation text from an ODF file (.odt, .ods, .odp).

    Opens the ZIP archive, parses ``content.xml``, and collects text
    from all ``<office:annotation>`` elements.

    Args:
        file_path: Path to the ODF file.

    Returns:
        list: (location_key, text) pairs with keys like
              'comment:{annotation_name}' or 'comment:{index}'.
    """
    texts: list[tuple[str, str]] = []

    with zipfile.ZipFile(file_path, "r") as zf:
        content_xml = zf.read("content.xml")

    root = etree.fromstring(content_xml)
    office_name_attr = f"{{{_ODF_NS['office']}}}name"
    text_p_tag = f"{{{_ODF_NS['text']}}}p"

    for idx, annotation in enumerate(
        root.findall(".//office:annotation", _ODF_NS),
    ):
        ann_name = annotation.get(office_name_attr, str(idx))
        text = _extract_odf_paragraph_text(annotation, text_p_tag)
        if text:
            texts.append((f"comment:{ann_name}", text))

    return texts


def _inject_odf_comments(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated comments into an ODF file (.odt, .ods, .odp).

    Reads the ZIP archive, modifies ``<office:annotation>`` text in
    ``content.xml``, and writes the archive back.

    Args:
        output_path: Path to the ODF file to modify in place.
        translations: Mapping of ``'comment:{name}'`` to translated text.
    """
    with zipfile.ZipFile(output_path, "r") as zf:
        content_xml = zf.read("content.xml")
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item.filename) for item in all_items}

    root = etree.fromstring(content_xml)
    office_name_attr = f"{{{_ODF_NS['office']}}}name"
    text_p_tag = f"{{{_ODF_NS['text']}}}p"
    modified = False

    for idx, annotation in enumerate(
        root.findall(".//office:annotation", _ODF_NS),
    ):
        ann_name = annotation.get(office_name_attr, str(idx))
        key = f"comment:{ann_name}"
        if key not in translations:
            continue

        # Reuse the shared ODF text injection helper
        modified |= _inject_odf_paragraph_text(
            annotation,
            translations[key],
            text_p_tag,
        )

    if not modified:
        return

    file_data["content.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
    )
    _rewrite_zip_content(output_path, file_data, all_items)


# ---------------------------------------------------------------------------
# Shape / text-box extraction and injection
# ---------------------------------------------------------------------------


def _extract_shapes(  # noqa: PLR0911
    file_path: Path,
    suffix: str,
    backend: str,
) -> list[tuple[str, str]]:
    """Extracts text from shapes and text boxes in an office file.

    Args:
        file_path: Path to the office file.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.

    Returns:
        list: (location_key, text) pairs for shape text.
    """
    # Modern / ODF formats
    if suffix == ".docx":
        # Always use ZIP+lxml for DOCX shapes: it reads raw WordprocessingML
        # run properties (including character-style resolution) so mixed
        # formatting in text boxes is correctly detected and emitted as HTML.
        # UNO preserves <wps:txbx> elements when saving DOCX, so the same
        # ZIP+lxml injector works regardless of backend.
        return _extract_docx_shapes(file_path)
    if suffix == ".xlsx":
        return _extract_xlsx_shapes(file_path)
    if suffix == ".odt":
        return _extract_odt_shapes(file_path)
    if suffix == ".ods":
        return _extract_ods_shapes(file_path)

    # Legacy formats — use win32com or UNO
    if suffix == ".doc":
        if backend == _BACKEND_WIN32COM:
            return _extract_win32com_word_shapes(file_path)
        return _extract_uno_writer_shapes(file_path)
    if suffix == ".xls":
        if backend == _BACKEND_WIN32COM:
            return _extract_win32com_excel_shapes(file_path)
        return _extract_uno_calc_shapes(file_path)

    return []


def _inject_shapes(  # noqa: PLR0912
    output_path: Path,
    translations: dict[str, str],
    suffix: str,
    backend: str,
) -> None:
    """Injects translated shape text back into the output document.

    Args:
        output_path: Path to the output file (already written by inject_fn).
        translations: Mapping of location_key to translated text.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.
    """
    shape_keys = {k: v for k, v in translations.items() if k.startswith("shape:")}
    if not shape_keys:
        return

    # Modern / ODF formats
    if suffix == ".docx":
        # Always use ZIP+lxml for DOCX shapes to match the ZIP+lxml extractor
        # and correctly handle inline HTML formatting tags.
        _inject_docx_shapes(output_path, shape_keys)
    elif suffix == ".xlsx":
        _inject_xlsx_shapes(output_path, shape_keys)
    elif suffix == ".odt":
        _inject_odt_shapes(output_path, shape_keys)
    elif suffix == ".ods":
        _inject_ods_shapes(output_path, shape_keys)
    # Legacy formats
    elif suffix == ".doc":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_word_shapes(output_path, shape_keys)
        else:
            _inject_uno_writer_shapes(output_path, shape_keys)
    elif suffix == ".xls":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_excel_shapes(output_path, shape_keys)
        else:
            _inject_uno_calc_shapes(output_path, shape_keys)


# ---------------------------------------------------------------------------
# Sheet name translation
# ---------------------------------------------------------------------------

# Extensions that support sheet-name translation.
_SHEET_NAME_EXTENSIONS = {".xlsx", ".ods", ".xls"}

# Characters invalid in Excel sheet names.
_INVALID_SHEET_NAME_CHARS = set("\\/*?:[]")

# Maximum sheet name length in Excel/Calc.
_MAX_SHEET_NAME_LEN = 31


def _sanitize_sheet_name(name: str) -> str:
    """Sanitises a translated sheet name for Excel/Calc compatibility.

    Removes invalid characters and truncates to 31 characters.

    Args:
        name: Raw translated sheet name.

    Returns:
        str: Sanitised name, or "Sheet" if the result is empty.
    """
    sanitized = "".join(c for c in name if c not in _INVALID_SHEET_NAME_CHARS)
    sanitized = sanitized.strip()
    return sanitized[:_MAX_SHEET_NAME_LEN] if sanitized else "Sheet"


def _should_translate_sheet_names(
    suffix: str,
    backend: str,
    config: TranslationConfig | None = None,
) -> bool:
    """Checks whether sheet-name translation should be attempted.

    Args:
        suffix: Lowercase file extension.
        backend: The detected backend identifier (unused).
        config: Optional TranslationConfig; falls back to load_setting().

    Returns:
        bool: True if sheet-name translation should proceed.
    """
    if suffix not in _SHEET_NAME_EXTENSIONS:
        return False

    if config is not None:
        return config.translate_sheet_names

    from src.constants.settings import SETTING_TRANSLATE_SHEET_NAMES  # noqa: PLC0415
    from src.utils.config_manager import load_setting  # noqa: PLC0415

    return bool(load_setting(SETTING_TRANSLATE_SHEET_NAMES, False))


def _extract_sheet_names(
    file_path: Path,
    suffix: str,
    backend: str,
) -> list[tuple[str, str]]:
    """Extracts sheet names from a spreadsheet file.

    Args:
        file_path: Path to the spreadsheet.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.

    Returns:
        list: (location_key, text) pairs with keys like 'sheetname:{name}'.
    """
    if suffix == ".xlsx":
        return _extract_xlsx_sheet_names(file_path)
    if suffix == ".ods":
        return _extract_ods_sheet_names(file_path)
    # Legacy .xls
    if backend == _BACKEND_WIN32COM:
        return _extract_win32com_excel_sheet_names(file_path)
    return _extract_uno_calc_sheet_names(file_path)


def _inject_sheet_names(
    output_path: Path,
    translations: dict[str, str],
    suffix: str,
    backend: str,
) -> None:
    """Injects translated sheet names back into the output spreadsheet.

    Args:
        output_path: Path to the output file.
        translations: Mapping of location_key to translated text.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.
    """
    name_keys = {k: v for k, v in translations.items() if k.startswith("sheetname:")}
    if not name_keys:
        return

    if suffix == ".xlsx":
        _inject_xlsx_sheet_names(output_path, name_keys)
    elif suffix == ".ods":
        _inject_ods_sheet_names(output_path, name_keys)
    elif suffix == ".xls":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_excel_sheet_names(output_path, name_keys)
        else:
            _inject_uno_calc_sheet_names(output_path, name_keys)


def _extract_xlsx_sheet_names(file_path: Path) -> list[tuple[str, str]]:
    """Extracts sheet names from an XLSX file via ZIP+lxml.

    Reads only ``xl/workbook.xml`` (a few KB) instead of loading the
    full workbook through openpyxl, which would parse all cell data.

    Args:
        file_path: Path to the .xlsx file.

    Returns:
        list: (location_key, sheet_name) pairs.
    """
    with zipfile.ZipFile(file_path, "r") as zf:
        if "xl/workbook.xml" not in zf.namelist():
            return []
        data = zf.read("xl/workbook.xml")

    root = etree.fromstring(data)
    texts: list[tuple[str, str]] = []
    for sheet_el in root.iter(f"{{{_SPREADSHEETML_NS}}}sheet"):
        name = sheet_el.get("name", "")
        if name and name.strip():
            texts.append((f"sheetname:{name}", name))
    return texts


def _inject_xlsx_sheet_names(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated sheet names into XLSX via ZIP+lxml.

    Uses direct XML manipulation to avoid openpyxl's lossy round-trip
    (which would drop restored embedded objects).

    Args:
        output_path: Path to the .xlsx file to modify in place.
        translations: Mapping of 'sheetname:{name}' to translated name.
    """
    with zipfile.ZipFile(output_path, "r") as zf:
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item) for item in all_items}

    if "xl/workbook.xml" not in file_data:
        return

    root = etree.fromstring(file_data["xl/workbook.xml"])
    modified = False
    for sheet_el in root.iter(f"{{{_SPREADSHEETML_NS}}}sheet"):
        name = sheet_el.get("name", "")
        key = f"sheetname:{name}"
        if key in translations:
            sheet_el.set("name", _sanitize_sheet_name(translations[key]))
            modified = True

    if not modified:
        return

    file_data["xl/workbook.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )
    _rewrite_zip_content(output_path, file_data, all_items)


def _extract_ods_sheet_names(file_path: Path) -> list[tuple[str, str]]:
    """Extracts sheet names from an ODS file via odfpy.

    Args:
        file_path: Path to the .ods file.

    Returns:
        list: (location_key, sheet_name) pairs.
    """
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.table import Table  # noqa: PLC0415

    doc = odf_load(str(file_path))
    texts: list[tuple[str, str]] = []
    for table in doc.spreadsheet.getElementsByType(Table):
        name = table.getAttribute("name")
        if name and name.strip():
            texts.append((f"sheetname:{name}", name))
    return texts


def _inject_ods_sheet_names(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated sheet names into ODS via ZIP+lxml.

    Args:
        output_path: Path to the .ods file to modify in place.
        translations: Mapping of 'sheetname:{name}' to translated name.
    """
    table_ns = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"

    with zipfile.ZipFile(output_path, "r") as zf:
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item) for item in all_items}

    if "content.xml" not in file_data:
        return

    root = etree.fromstring(file_data["content.xml"])
    modified = False
    for table_el in root.iter(f"{{{table_ns}}}table"):
        name = table_el.get(f"{{{table_ns}}}name", "")
        key = f"sheetname:{name}"
        if key in translations:
            table_el.set(
                f"{{{table_ns}}}name",
                _sanitize_sheet_name(translations[key]),
            )
            modified = True

    if not modified:
        return

    file_data["content.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )
    _rewrite_zip_content(output_path, file_data, all_items)


def _extract_win32com_excel_sheet_names(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts sheet names from an XLS file via win32com.

    Args:
        file_path: Path to the .xls file.

    Returns:
        list: (location_key, sheet_name) pairs.
    """
    app, wb, pycom = _win32com_open(_APP_EXCEL, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for i in range(1, wb.Worksheets.Count + 1):
            name = wb.Worksheets(i).Name
            if name and name.strip():
                texts.append((f"sheetname:{name}", name))
        return texts
    finally:
        _win32com_close(app, wb, pycom)


def _inject_win32com_excel_sheet_names(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated sheet names into an XLS file via win32com.

    Args:
        output_path: Path to the .xls file to modify in place.
        translations: Mapping of 'sheetname:{name}' to translated name.
    """
    app, wb, pycom = _win32com_open(_APP_EXCEL, output_path)
    try:
        for i in range(1, wb.Worksheets.Count + 1):
            ws = wb.Worksheets(i)
            key = f"sheetname:{ws.Name}"
            if key in translations:
                ws.Name = _sanitize_sheet_name(translations[key])
        wb.Save()
    finally:
        _win32com_close(app, wb, pycom)


def _extract_uno_calc_sheet_names(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts sheet names from an XLS/ODS file via UNO.

    Args:
        file_path: Path to the spreadsheet.

    Returns:
        list: (location_key, sheet_name) pairs.
    """
    doc = _uno_open(file_path)
    try:
        sheets = doc.getSheets()
        texts: list[tuple[str, str]] = []
        for i in range(sheets.getCount()):
            name = sheets.getByIndex(i).getName()
            if name and name.strip():
                texts.append((f"sheetname:{name}", name))
        return texts
    finally:
        doc.close(True)


def _inject_uno_calc_sheet_names(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated sheet names into a spreadsheet via UNO.

    Args:
        output_path: Path to the spreadsheet to modify in place.
        translations: Mapping of 'sheetname:{name}' to translated name.
    """
    doc = _uno_open(output_path)
    try:
        sheets = doc.getSheets()
        for i in range(sheets.getCount()):
            sheet = sheets.getByIndex(i)
            key = f"sheetname:{sheet.getName()}"
            if key in translations:
                sheet.setName(_sanitize_sheet_name(translations[key]))
        _uno_save(doc, output_path)
    finally:
        doc.close(True)


# ---------------------------------------------------------------------------
# Speaker notes translation
# ---------------------------------------------------------------------------

# Extensions that support speaker-notes translation.
_NOTES_EXTENSIONS = {".pptx", ".odp", ".ppt"}


def _should_translate_notes(
    suffix: str,
    backend: str,
    config: TranslationConfig | None = None,
) -> bool:
    """Checks whether speaker-notes translation should be attempted.

    Args:
        suffix: Lowercase file extension.
        backend: The detected backend identifier (unused).
        config: Optional TranslationConfig; falls back to load_setting().

    Returns:
        bool: True if speaker-notes translation should proceed.
    """
    if suffix not in _NOTES_EXTENSIONS:
        return False

    if config is not None:
        return config.translate_doc_notes

    from src.constants.settings import SETTING_TRANSLATE_DOC_NOTES  # noqa: PLC0415
    from src.utils.config_manager import load_setting  # noqa: PLC0415

    return bool(load_setting(SETTING_TRANSLATE_DOC_NOTES, False))


def _extract_notes(
    file_path: Path,
    suffix: str,
    backend: str,
) -> list[tuple[str, str]]:
    """Extracts speaker notes from a presentation file.

    Args:
        file_path: Path to the presentation.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.

    Returns:
        list: (location_key, text) pairs with keys like 'note:{slide}:{para}'.
    """
    if suffix == ".pptx":
        return _extract_pptx_notes(file_path)
    if suffix == ".odp":
        return _extract_odp_notes(file_path)
    # Legacy .ppt
    if backend == _BACKEND_WIN32COM:
        return _extract_win32com_ppt_notes(file_path)
    return _extract_uno_impress_notes(file_path)


def _inject_notes(
    output_path: Path,
    translations: dict[str, str],
    suffix: str,
    backend: str,
) -> None:
    """Injects translated speaker notes back into the output presentation.

    Args:
        output_path: Path to the output file.
        translations: Mapping of location_key to translated text.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.
    """
    note_keys = {k: v for k, v in translations.items() if k.startswith("note:")}
    if not note_keys:
        return

    if suffix == ".pptx":
        _inject_pptx_notes(output_path, note_keys)
    elif suffix == ".odp":
        _inject_odp_notes(output_path, note_keys)
    elif suffix == ".ppt":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_ppt_notes(output_path, note_keys)
        else:
            _inject_uno_impress_notes(output_path, note_keys)


def _extract_pptx_notes(file_path: Path) -> list[tuple[str, str]]:
    """Extracts speaker notes from a PPTX file via python-pptx.

    Paragraphs with mixed formatting or hyperlinks are encoded as
    inline HTML so the LLM can preserve them.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        list: (location_key, text) pairs.
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(file_path))
    texts: list[tuple[str, str]] = []
    for s_idx, slide in enumerate(prs.slides):
        if not slide.has_notes_slide:
            continue
        notes_frame = slide.notes_slide.notes_text_frame
        for p_idx, para in enumerate(notes_frame.paragraphs):
            if not para.text.strip():
                continue
            if _has_pptx_mixed_formatting(para) or _has_pptx_hyperlinks(para):
                text = _pptx_runs_to_html(para)
            else:
                text = para.text
            texts.append((f"note:{s_idx}:{p_idx}", text))
    return texts


def _inject_pptx_notes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated speaker notes into a PPTX file via python-pptx.

    Args:
        output_path: Path to the .pptx file to modify in place.
        translations: Mapping of 'note:{slide}:{para}' to translated text.
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(output_path))
    for s_idx, slide in enumerate(prs.slides):
        if not slide.has_notes_slide:
            continue
        notes_frame = slide.notes_slide.notes_text_frame
        slide_part = slide.notes_slide.part
        for p_idx, para in enumerate(notes_frame.paragraphs):
            key = f"note:{s_idx}:{p_idx}"
            if key not in translations:
                continue
            translated = translations[key]
            if _FORMATTING_HTML_RE.search(translated):
                _inject_pptx_html_runs(para, translated, part=slide_part)
            elif para.runs:
                para.runs[0].text = translated
                for run in para.runs[1:]:
                    run.text = ""
            else:
                para.text = translated
    prs.save(str(output_path))


def _extract_odp_notes(file_path: Path) -> list[tuple[str, str]]:
    """Extracts speaker notes from an ODP file via odfpy.

    Args:
        file_path: Path to the .odp file.

    Returns:
        list: (location_key, text) pairs.
    """
    from odf.draw import Frame, Page  # noqa: PLC0415
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.presentation import Notes  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = odf_load(str(file_path))
    texts: list[tuple[str, str]] = []
    for s_idx, page in enumerate(doc.getElementsByType(Page)):
        for notes_elem in page.getElementsByType(Notes):
            for frame in notes_elem.getElementsByType(Frame):
                for p_idx, para in enumerate(frame.getElementsByType(P)):
                    text = _odf_element_text(
                        para,
                        preserve_links=True,
                    ).strip()
                    if text:
                        texts.append((f"note:{s_idx}:{p_idx}", text))
    return texts


def _inject_odp_notes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated speaker notes into an ODP file via odfpy.

    Args:
        output_path: Path to the .odp file to modify in place.
        translations: Mapping of 'note:{slide}:{para}' to translated text.
    """
    from odf.draw import Frame, Page  # noqa: PLC0415
    from odf.opendocument import load as odf_load  # noqa: PLC0415
    from odf.presentation import Notes  # noqa: PLC0415
    from odf.text import P  # noqa: PLC0415

    doc = odf_load(str(output_path))
    for s_idx, page in enumerate(doc.getElementsByType(Page)):
        for notes_elem in page.getElementsByType(Notes):
            for frame in notes_elem.getElementsByType(Frame):
                for p_idx, para in enumerate(frame.getElementsByType(P)):
                    key = f"note:{s_idx}:{p_idx}"
                    if key in translations:
                        _odf_replace_text(para, translations[key])
    doc.save(str(output_path))


def _extract_win32com_ppt_notes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts speaker notes from a PPT file via win32com.

    Iterates the notes page of each slide and extracts text from
    shapes that have text frames.

    Args:
        file_path: Path to the .ppt file.

    Returns:
        list: (location_key, text) pairs.
    """
    app, prs, pycom = _win32com_open(_APP_PPT, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for s_idx in range(1, prs.Slides.Count + 1):
            slide = prs.Slides(s_idx)
            try:
                notes_page = slide.NotesPage
            except Exception:  # noqa: BLE001
                continue
            for sh_idx in range(1, notes_page.Shapes.Count + 1):
                shape = notes_page.Shapes(sh_idx)
                try:
                    if not shape.HasTextFrame:
                        continue
                    # ppPlaceholderBody = 2 — the notes text body
                    with contextlib.suppress(Exception):
                        if shape.PlaceholderFormat.Type != 2:  # noqa: PLR2004
                            continue
                    text_range = shape.TextFrame.TextRange
                    for p_idx in range(1, text_range.Paragraphs().Count + 1):
                        para = text_range.Paragraphs(p_idx, 1)
                        text = para.Text.rstrip("\r\n")
                        if text and text.strip():
                            texts.append(
                                (f"note:{s_idx - 1}:{p_idx - 1}", text),
                            )
                except Exception:  # noqa: BLE001
                    continue
        return texts
    finally:
        _win32com_close(app, prs, pycom)


def _inject_win32com_ppt_notes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated speaker notes into a PPT file via win32com.

    Args:
        output_path: Path to the .ppt file to modify in place.
        translations: Mapping of 'note:{slide}:{para}' to translated text.
    """
    app, prs, pycom = _win32com_open(_APP_PPT, output_path)
    try:
        for s_idx in range(1, prs.Slides.Count + 1):
            slide = prs.Slides(s_idx)
            try:
                notes_page = slide.NotesPage
            except Exception:  # noqa: BLE001
                continue
            for sh_idx in range(1, notes_page.Shapes.Count + 1):
                shape = notes_page.Shapes(sh_idx)
                try:
                    if not shape.HasTextFrame:
                        continue
                    with contextlib.suppress(Exception):
                        if shape.PlaceholderFormat.Type != 2:  # noqa: PLR2004
                            continue
                    text_range = shape.TextFrame.TextRange
                    for p_idx in range(1, text_range.Paragraphs().Count + 1):
                        para = text_range.Paragraphs(p_idx, 1)
                        key = f"note:{s_idx - 1}:{p_idx - 1}"
                        if key in translations:
                            para.Text = translations[key]
                except Exception:  # noqa: BLE001
                    continue
        prs.SaveAs(str(output_path))
    finally:
        _win32com_close(app, prs, pycom)


def _extract_uno_impress_notes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts speaker notes from a PPT/ODP file via UNO.

    Args:
        file_path: Path to the presentation.

    Returns:
        list: (location_key, text) pairs.
    """
    doc = _uno_open(file_path)
    try:
        texts: list[tuple[str, str]] = []
        pages = doc.getDrawPages()
        for s_idx in range(pages.getCount()):
            page = pages.getByIndex(s_idx)
            try:
                notes_page = page.getNotesPage()
            except Exception:  # noqa: BLE001
                continue
            # Notes page shape 1 is the notes text body
            for sh_idx in range(notes_page.getCount()):
                shape = notes_page.getByIndex(sh_idx)
                if not shape.supportsService("com.sun.star.drawing.Text"):
                    continue
                text_obj = shape.getText()
                paragraph_text = text_obj.getString().strip()
                if not paragraph_text:
                    continue
                # Extract per-paragraph
                enum = text_obj.createEnumeration()
                p_idx = 0
                while enum.hasMoreElements():
                    para = enum.nextElement()
                    text = para.getString().strip()
                    if text:
                        texts.append((f"note:{s_idx}:{p_idx}", text))
                    p_idx += 1
        return texts
    finally:
        doc.close(True)


def _inject_uno_impress_notes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated speaker notes into a presentation via UNO.

    Args:
        output_path: Path to the presentation to modify in place.
        translations: Mapping of 'note:{slide}:{para}' to translated text.
    """
    doc = _uno_open(output_path)
    try:
        pages = doc.getDrawPages()
        for s_idx in range(pages.getCount()):
            page = pages.getByIndex(s_idx)
            try:
                notes_page = page.getNotesPage()
            except Exception:  # noqa: BLE001
                continue
            for sh_idx in range(notes_page.getCount()):
                shape = notes_page.getByIndex(sh_idx)
                if not shape.supportsService("com.sun.star.drawing.Text"):
                    continue
                text_obj = shape.getText()
                enum = text_obj.createEnumeration()
                p_idx = 0
                while enum.hasMoreElements():
                    para = enum.nextElement()
                    key = f"note:{s_idx}:{p_idx}"
                    if key in translations:
                        para.setString(translations[key])
                    p_idx += 1
        _uno_save(doc, output_path)
    finally:
        doc.close(True)


# ---------------------------------------------------------------------------
# Header and footer translation
# ---------------------------------------------------------------------------

# Extensions that support header/footer translation.
_HEADER_FOOTER_EXTENSIONS = {".docx", ".odt", ".doc"}

# Header/footer type identifiers for location keys.
_HF_DEFAULT = "default"
_HF_FIRST = "first"
_HF_EVEN = "even"

# wdHeaderFooterPrimary=1, wdHeaderFooterFirstPage=2, wdHeaderFooterEvenPages=3
_HF_TYPE_MAP = {1: _HF_DEFAULT, 2: _HF_FIRST, 3: _HF_EVEN}


def _extract_headers_footers(
    file_path: Path,
    suffix: str,
    backend: str,
) -> list[tuple[str, str]]:
    """Extracts headers and footers from a word-processing document.

    Args:
        file_path: Path to the document.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.

    Returns:
        list: (location_key, text) pairs with keys like
              'header:{section}:{type}:{para}' or 'footer:...'.
    """
    if suffix == ".docx":
        return _extract_docx_headers_footers(file_path)
    if suffix == ".odt":
        return _extract_odt_headers_footers(file_path)
    # Legacy .doc
    if backend == _BACKEND_WIN32COM:
        return _extract_win32com_word_headers_footers(file_path)
    return _extract_uno_writer_headers_footers(file_path)


def _inject_headers_footers(
    output_path: Path,
    translations: dict[str, str],
    suffix: str,
    backend: str,
) -> None:
    """Injects translated headers/footers back into the output document.

    Args:
        output_path: Path to the output file.
        translations: Mapping of location_key to translated text.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.
    """
    hf_keys = {
        k: v
        for k, v in translations.items()
        if k.startswith("header:") or k.startswith("footer:")
    }
    if not hf_keys:
        return

    if suffix == ".docx":
        _inject_docx_headers_footers(output_path, hf_keys)
    elif suffix == ".odt":
        _inject_odt_headers_footers(output_path, hf_keys)
    elif suffix == ".doc":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_word_headers_footers(output_path, hf_keys)
        else:
            _inject_uno_writer_headers_footers(output_path, hf_keys)


def _extract_docx_hf_part(
    paragraphs: list[object],
    section_idx: int,
    hf_type: str,
    prefix: str,
) -> list[tuple[str, str]]:
    """Extracts text from a DOCX header/footer part's paragraphs.

    Args:
        paragraphs: List of python-docx Paragraph objects.
        section_idx: Section index.
        hf_type: Type identifier ('default', 'first', 'even').
        prefix: 'header' or 'footer'.

    Returns:
        list: (location_key, text) pairs.
    """
    texts: list[tuple[str, str]] = []
    for p_idx, para in enumerate(paragraphs):
        if para.text.strip():
            text = _extract_para_with_links(para)
            texts.append((f"{prefix}:{section_idx}:{hf_type}:{p_idx}", text))
    return texts


def _extract_docx_headers_footers(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts headers and footers from a DOCX file via python-docx.

    Extracts default, first-page, and even-page headers/footers from
    each section.  Paragraphs with mixed formatting or hyperlinks are
    encoded as inline HTML.

    Args:
        file_path: Path to the .docx file.

    Returns:
        list: (location_key, text) pairs.
    """
    from docx import Document  # noqa: PLC0415

    doc = Document(str(file_path))
    texts: list[tuple[str, str]] = []

    for s_idx, section in enumerate(doc.sections):
        # Default header/footer
        hdr = section.header
        if not hdr.is_linked_to_previous:
            texts.extend(
                _extract_docx_hf_part(
                    hdr.paragraphs,
                    s_idx,
                    _HF_DEFAULT,
                    "header",
                ),
            )
        ftr = section.footer
        if not ftr.is_linked_to_previous:
            texts.extend(
                _extract_docx_hf_part(
                    ftr.paragraphs,
                    s_idx,
                    _HF_DEFAULT,
                    "footer",
                ),
            )

        # First-page header/footer
        if section.different_first_page_header_footer:
            first_hdr = section.first_page_header
            if not first_hdr.is_linked_to_previous:
                texts.extend(
                    _extract_docx_hf_part(
                        first_hdr.paragraphs,
                        s_idx,
                        _HF_FIRST,
                        "header",
                    ),
                )
            first_ftr = section.first_page_footer
            if not first_ftr.is_linked_to_previous:
                texts.extend(
                    _extract_docx_hf_part(
                        first_ftr.paragraphs,
                        s_idx,
                        _HF_FIRST,
                        "footer",
                    ),
                )

        # Even-page header/footer
        if doc.settings.odd_and_even_pages_header_footer:
            even_hdr = section.even_page_header
            if not even_hdr.is_linked_to_previous:
                texts.extend(
                    _extract_docx_hf_part(
                        even_hdr.paragraphs,
                        s_idx,
                        _HF_EVEN,
                        "header",
                    ),
                )
            even_ftr = section.even_page_footer
            if not even_ftr.is_linked_to_previous:
                texts.extend(
                    _extract_docx_hf_part(
                        even_ftr.paragraphs,
                        s_idx,
                        _HF_EVEN,
                        "footer",
                    ),
                )

    return texts


def _inject_docx_headers_footers(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated headers/footers into a DOCX file via python-docx.

    Args:
        output_path: Path to the .docx file to modify in place.
        translations: Mapping of 'header:...' / 'footer:...' to translated text.
    """
    from docx import Document  # noqa: PLC0415

    doc = Document(str(output_path))
    doc_part = doc.part

    def _inject_hf_part(
        paragraphs: list[object],
        section_idx: int,
        hf_type: str,
        prefix: str,
    ) -> None:
        for p_idx, para in enumerate(paragraphs):
            key = f"{prefix}:{section_idx}:{hf_type}:{p_idx}"
            if key in translations:
                text = translations[key]
                if _FORMATTING_HTML_RE.search(text):
                    _inject_html_runs(para, text, part=doc_part)
                else:
                    _replace_paragraph_text(para, text)

    for s_idx, section in enumerate(doc.sections):
        hdr = section.header
        if not hdr.is_linked_to_previous:
            _inject_hf_part(hdr.paragraphs, s_idx, _HF_DEFAULT, "header")
        ftr = section.footer
        if not ftr.is_linked_to_previous:
            _inject_hf_part(ftr.paragraphs, s_idx, _HF_DEFAULT, "footer")

        if section.different_first_page_header_footer:
            first_hdr = section.first_page_header
            if not first_hdr.is_linked_to_previous:
                _inject_hf_part(
                    first_hdr.paragraphs,
                    s_idx,
                    _HF_FIRST,
                    "header",
                )
            first_ftr = section.first_page_footer
            if not first_ftr.is_linked_to_previous:
                _inject_hf_part(
                    first_ftr.paragraphs,
                    s_idx,
                    _HF_FIRST,
                    "footer",
                )

        if doc.settings.odd_and_even_pages_header_footer:
            even_hdr = section.even_page_header
            if not even_hdr.is_linked_to_previous:
                _inject_hf_part(
                    even_hdr.paragraphs,
                    s_idx,
                    _HF_EVEN,
                    "header",
                )
            even_ftr = section.even_page_footer
            if not even_ftr.is_linked_to_previous:
                _inject_hf_part(
                    even_ftr.paragraphs,
                    s_idx,
                    _HF_EVEN,
                    "footer",
                )

    doc.save(str(output_path))


def _build_odf_hf_map(style_ns: str) -> dict[str, tuple[str, str]]:
    """Builds an ODF header/footer element-tag → (prefix, type) lookup.

    Used by both ``_extract_odt_headers_footers`` and
    ``_inject_odt_headers_footers`` to avoid duplicating the mapping.
    """
    return {
        f"{{{style_ns}}}header": ("header", _HF_DEFAULT),
        f"{{{style_ns}}}footer": ("footer", _HF_DEFAULT),
        f"{{{style_ns}}}header-first": ("header", _HF_FIRST),
        f"{{{style_ns}}}footer-first": ("footer", _HF_FIRST),
        f"{{{style_ns}}}header-left": ("header", _HF_EVEN),
        f"{{{style_ns}}}footer-left": ("footer", _HF_EVEN),
    }


def _extract_odt_headers_footers(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts headers and footers from an ODT file via ZIP+lxml.

    ODT stores headers/footers in ``styles.xml`` under
    ``<style:master-page>`` elements.

    Args:
        file_path: Path to the .odt file.

    Returns:
        list: (location_key, text) pairs.
    """
    style_ns = _ODF_NS["style"]
    text_ns = _ODF_NS["text"]

    with zipfile.ZipFile(file_path, "r") as zf:
        if "styles.xml" not in zf.namelist():
            return []
        data = zf.read("styles.xml")

    root = etree.fromstring(data)
    texts: list[tuple[str, str]] = []
    hf_map = _build_odf_hf_map(style_ns)

    text_p_tag = f"{{{text_ns}}}p"
    for s_idx, master in enumerate(
        root.iter(f"{{{style_ns}}}master-page"),
    ):
        for child in master:
            if child.tag not in hf_map:
                continue
            prefix, hf_type = hf_map[child.tag]
            text = _extract_odf_paragraph_text(child, text_p_tag)
            if text and text.strip():
                texts.append(
                    (f"{prefix}:{s_idx}:{hf_type}:0", text),
                )

    return texts


def _inject_odt_headers_footers(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated headers/footers into an ODT file via ZIP+lxml.

    Args:
        output_path: Path to the .odt file to modify in place.
        translations: Mapping of 'header:...' / 'footer:...' to translated text.
    """
    style_ns = _ODF_NS["style"]
    text_ns = _ODF_NS["text"]

    with zipfile.ZipFile(output_path, "r") as zf:
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item) for item in all_items}

    if "styles.xml" not in file_data:
        return

    root = etree.fromstring(file_data["styles.xml"])
    text_p_tag = f"{{{text_ns}}}p"
    hf_map = _build_odf_hf_map(style_ns)

    modified = False
    for s_idx, master in enumerate(
        root.iter(f"{{{style_ns}}}master-page"),
    ):
        for child in master:
            if child.tag not in hf_map:
                continue
            prefix, hf_type = hf_map[child.tag]
            # We extract full hf text at para index 0
            key = f"{prefix}:{s_idx}:{hf_type}:0"
            if key in translations:
                _inject_odf_paragraph_text(
                    child,
                    translations[key],
                    text_p_tag,
                )
                modified = True

    if not modified:
        return

    file_data["styles.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )
    _rewrite_zip_content(output_path, file_data, all_items)


def _extract_win32com_word_headers_footers(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts headers/footers from a DOC file via win32com.

    Args:
        file_path: Path to the .doc file.

    Returns:
        list: (location_key, text) pairs.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for s_idx in range(1, doc.Sections.Count + 1):
            section = doc.Sections(s_idx)
            for hf_type_id, hf_type_str in _HF_TYPE_MAP.items():
                # Headers
                try:
                    hdr = section.Headers(hf_type_id)
                    if hdr.Exists:
                        for p_idx in range(
                            1,
                            hdr.Range.Paragraphs.Count + 1,
                        ):
                            text = hdr.Range.Paragraphs(p_idx).Range.Text.rstrip("\r\n")
                            if text and text.strip():
                                texts.append(
                                    (
                                        f"header:{s_idx - 1}:{hf_type_str}:{p_idx - 1}",
                                        text,
                                    )
                                )
                except Exception:  # noqa: BLE001
                    pass
                # Footers
                try:
                    ftr = section.Footers(hf_type_id)
                    if ftr.Exists:
                        for p_idx in range(
                            1,
                            ftr.Range.Paragraphs.Count + 1,
                        ):
                            text = ftr.Range.Paragraphs(p_idx).Range.Text.rstrip("\r\n")
                            if text and text.strip():
                                texts.append(
                                    (
                                        f"footer:{s_idx - 1}:{hf_type_str}:{p_idx - 1}",
                                        text,
                                    )
                                )
                except Exception:  # noqa: BLE001
                    pass
        return texts
    finally:
        _win32com_close(app, doc, pycom)


def _inject_win32com_word_headers_footers(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated headers/footers into a DOC file via win32com.

    Args:
        output_path: Path to the .doc file to modify in place.
        translations: Mapping of 'header:...' / 'footer:...' to translated text.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, output_path)
    try:
        for s_idx in range(1, doc.Sections.Count + 1):
            section = doc.Sections(s_idx)
            for hf_type_id, hf_type_str in _HF_TYPE_MAP.items():
                try:
                    hdr = section.Headers(hf_type_id)
                    if hdr.Exists:
                        for p_idx in range(
                            1,
                            hdr.Range.Paragraphs.Count + 1,
                        ):
                            key = f"header:{s_idx - 1}:{hf_type_str}:{p_idx - 1}"
                            if key in translations:
                                hdr.Range.Paragraphs(
                                    p_idx,
                                ).Range.Text = translations[key]
                except Exception:  # noqa: BLE001
                    pass
                try:
                    ftr = section.Footers(hf_type_id)
                    if ftr.Exists:
                        for p_idx in range(
                            1,
                            ftr.Range.Paragraphs.Count + 1,
                        ):
                            key = f"footer:{s_idx - 1}:{hf_type_str}:{p_idx - 1}"
                            if key in translations:
                                ftr.Range.Paragraphs(
                                    p_idx,
                                ).Range.Text = translations[key]
                except Exception:  # noqa: BLE001
                    pass
        doc.SaveAs(str(output_path))
    finally:
        _win32com_close(app, doc, pycom)


def _extract_uno_writer_headers_footers(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts headers/footers from a DOC/ODT file via UNO.

    UNO stores headers/footers on page styles. Each unique page style
    is treated as a "section" for key purposes.

    Note:
        Only *default* headers/footers are extracted. UNO exposes
        first-page (``HeaderTextFirst``) and even-page
        (``HeaderTextLeft``) properties, but they require additional
        page-style flags (``HeaderIsShared`` / ``FirstIsShared``) that
        vary across LibreOffice versions. Default-only is sufficient
        for the vast majority of DOC files.

    Args:
        file_path: Path to the document.

    Returns:
        list: (location_key, text) pairs.
    """
    doc = _uno_open(file_path)
    try:
        texts: list[tuple[str, str]] = []
        styles = doc.getStyleFamilies().getByName("PageStyles")
        for s_idx, style_name in enumerate(styles.getElementNames()):
            style = styles.getByName(style_name)

            # Default header
            if style.HeaderIsOn:
                try:
                    hdr_text = style.HeaderText
                    enum = hdr_text.createEnumeration()
                    p_idx = 0
                    while enum.hasMoreElements():
                        para = enum.nextElement()
                        text = para.getString().strip()
                        if text:
                            texts.append(
                                (
                                    f"header:{s_idx}:{_HF_DEFAULT}:{p_idx}",
                                    text,
                                )
                            )
                        p_idx += 1
                except Exception:  # noqa: BLE001
                    pass

            # Default footer
            if style.FooterIsOn:
                try:
                    ftr_text = style.FooterText
                    enum = ftr_text.createEnumeration()
                    p_idx = 0
                    while enum.hasMoreElements():
                        para = enum.nextElement()
                        text = para.getString().strip()
                        if text:
                            texts.append(
                                (
                                    f"footer:{s_idx}:{_HF_DEFAULT}:{p_idx}",
                                    text,
                                )
                            )
                        p_idx += 1
                except Exception:  # noqa: BLE001
                    pass

        return texts
    finally:
        doc.close(True)


def _inject_uno_writer_headers_footers(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated headers/footers into a document via UNO.

    Only default headers/footers are handled. See
    :func:`_extract_uno_writer_headers_footers` note on first/even-page
    limitation.

    Args:
        output_path: Path to the document to modify in place.
        translations: Mapping of 'header:...' / 'footer:...' to translated text.
    """
    doc = _uno_open(output_path)
    try:
        styles = doc.getStyleFamilies().getByName("PageStyles")
        for s_idx, style_name in enumerate(styles.getElementNames()):
            style = styles.getByName(style_name)

            if style.HeaderIsOn:
                try:
                    hdr_text = style.HeaderText
                    enum = hdr_text.createEnumeration()
                    p_idx = 0
                    while enum.hasMoreElements():
                        para = enum.nextElement()
                        key = f"header:{s_idx}:{_HF_DEFAULT}:{p_idx}"
                        if key in translations:
                            para.setString(translations[key])
                        p_idx += 1
                except Exception:  # noqa: BLE001
                    pass

            if style.FooterIsOn:
                try:
                    ftr_text = style.FooterText
                    enum = ftr_text.createEnumeration()
                    p_idx = 0
                    while enum.hasMoreElements():
                        para = enum.nextElement()
                        key = f"footer:{s_idx}:{_HF_DEFAULT}:{p_idx}"
                        if key in translations:
                            para.setString(translations[key])
                        p_idx += 1
                except Exception:  # noqa: BLE001
                    pass

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


# ---------------------------------------------------------------------------
# Footnote and endnote translation
# ---------------------------------------------------------------------------

# Extensions that support footnote/endnote translation.
_FOOTNOTE_EXTENSIONS = {".docx", ".odt", ".doc"}

# Footnote/endnote IDs to skip in DOCX (separator and continuationSeparator).
_DOCX_FN_SKIP_IDS = {"0", "1", "-1"}


def _extract_footnotes(
    file_path: Path,
    suffix: str,
    backend: str,
) -> list[tuple[str, str]]:
    """Extracts footnotes and endnotes from a word-processing document.

    Args:
        file_path: Path to the document.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.

    Returns:
        list: (location_key, text) pairs with keys like 'footnote:{id}'
              or 'endnote:{id}'.
    """
    if suffix == ".docx":
        return _extract_docx_footnotes(file_path)
    if suffix == ".odt":
        return _extract_odt_footnotes(file_path)
    # Legacy .doc
    if backend == _BACKEND_WIN32COM:
        return _extract_win32com_word_footnotes(file_path)
    return _extract_uno_writer_footnotes(file_path)


def _inject_footnotes(
    output_path: Path,
    translations: dict[str, str],
    suffix: str,
    backend: str,
) -> None:
    """Injects translated footnotes/endnotes into the output document.

    Args:
        output_path: Path to the output file.
        translations: Mapping of location_key to translated text.
        suffix: Lowercase file extension.
        backend: Backend identifier for legacy format dispatch.
    """
    fn_keys = {
        k: v
        for k, v in translations.items()
        if k.startswith("footnote:") or k.startswith("endnote:")
    }
    if not fn_keys:
        return

    if suffix == ".docx":
        _inject_docx_footnotes(output_path, fn_keys)
    elif suffix == ".odt":
        _inject_odt_footnotes(output_path, fn_keys)
    elif suffix == ".doc":
        if backend == _BACKEND_WIN32COM:
            _inject_win32com_word_footnotes(output_path, fn_keys)
        else:
            _inject_uno_writer_footnotes(output_path, fn_keys)


def _extract_docx_fn_xml(
    xml_data: bytes,
    element_tag: str,
    key_prefix: str,
) -> list[tuple[str, str]]:
    """Extracts text from DOCX footnote or endnote XML.

    Args:
        xml_data: Raw XML bytes of footnotes.xml or endnotes.xml.
        element_tag: Fully-qualified tag (e.g. ``{ns}footnote``).
        key_prefix: Key prefix ('footnote' or 'endnote').

    Returns:
        list: (location_key, text) pairs.
    """
    w_ns = _WORDML_NS
    root = etree.fromstring(xml_data)
    texts: list[tuple[str, str]] = []
    for elem in root.iter(element_tag):
        elem_id = elem.get(f"{{{w_ns}}}id", "")
        if elem_id in _DOCX_FN_SKIP_IDS:
            continue
        parts: list[str] = []
        for para in elem.iter(f"{{{w_ns}}}p"):
            para_text = ""
            for r_elem in para.iter(f"{{{w_ns}}}r"):
                for t_elem in r_elem.iter(f"{{{w_ns}}}t"):
                    if t_elem.text:
                        para_text += t_elem.text
            if para_text.strip():
                parts.append(para_text.strip())
        full_text = "\n".join(parts)
        if full_text:
            texts.append((f"{key_prefix}:{elem_id}", full_text))
    return texts


def _extract_docx_footnotes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts footnotes and endnotes from a DOCX file via ZIP+lxml.

    Reads ``word/footnotes.xml`` and ``word/endnotes.xml``.  IDs 0, 1,
    and -1 are internal separators and are skipped.

    Args:
        file_path: Path to the .docx file.

    Returns:
        list: (location_key, text) pairs.
    """
    texts: list[tuple[str, str]] = []
    w_ns = _WORDML_NS

    with zipfile.ZipFile(file_path, "r") as zf:
        names = zf.namelist()
        if "word/footnotes.xml" in names:
            texts.extend(
                _extract_docx_fn_xml(
                    zf.read("word/footnotes.xml"),
                    f"{{{w_ns}}}footnote",
                    "footnote",
                ),
            )
        if "word/endnotes.xml" in names:
            texts.extend(
                _extract_docx_fn_xml(
                    zf.read("word/endnotes.xml"),
                    f"{{{w_ns}}}endnote",
                    "endnote",
                ),
            )

    return texts


def _inject_docx_footnotes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated footnotes/endnotes into a DOCX file via ZIP+lxml.

    Args:
        output_path: Path to the .docx file to modify in place.
        translations: Mapping of 'footnote:{id}' / 'endnote:{id}' to
                      translated text.
    """
    w_ns = _WORDML_NS

    with zipfile.ZipFile(output_path, "r") as zf:
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item) for item in all_items}

    modified = False

    # Footnotes
    if "word/footnotes.xml" in file_data:
        root = etree.fromstring(file_data["word/footnotes.xml"])
        for fn in root.iter(f"{{{w_ns}}}footnote"):
            fn_id = fn.get(f"{{{w_ns}}}id", "")
            key = f"footnote:{fn_id}"
            if key not in translations:
                continue
            # Replace text in the first paragraph's runs
            paras = list(fn.iter(f"{{{w_ns}}}p"))
            if not paras:
                continue
            _inject_docx_fn_text(paras, translations[key], w_ns)
            modified = True
        if modified:
            file_data["word/footnotes.xml"] = etree.tostring(
                root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )

    # Endnotes
    endnote_modified = False
    if "word/endnotes.xml" in file_data:
        root = etree.fromstring(file_data["word/endnotes.xml"])
        for en in root.iter(f"{{{w_ns}}}endnote"):
            en_id = en.get(f"{{{w_ns}}}id", "")
            key = f"endnote:{en_id}"
            if key not in translations:
                continue
            paras = list(en.iter(f"{{{w_ns}}}p"))
            if not paras:
                continue
            _inject_docx_fn_text(paras, translations[key], w_ns)
            endnote_modified = True
        if endnote_modified:
            file_data["word/endnotes.xml"] = etree.tostring(
                root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            modified = True

    if not modified:
        return

    _rewrite_zip_content(output_path, file_data, all_items)


def _inject_docx_fn_text(
    paras: list[object],
    new_text: str,
    w_ns: str,
) -> None:
    """Replaces text in DOCX footnote/endnote paragraphs.

    Preserves the footnote-reference run (``<w:footnoteRef/>``) in the
    first paragraph and replaces text in subsequent runs.

    Args:
        paras: List of ``<w:p>`` lxml elements.
        new_text: Translated text (paragraphs separated by newlines).
        w_ns: WordprocessingML namespace URI.
    """
    text_lines = new_text.split("\n")

    for p_idx, para in enumerate(paras):
        if p_idx >= len(text_lines):
            # Clear extra paragraphs
            for r_elem in list(para.iter(f"{{{w_ns}}}r")):
                # Keep footnoteRef / endnoteRef runs
                if r_elem.find(f"{{{w_ns}}}footnoteRef") is not None:
                    continue
                if r_elem.find(f"{{{w_ns}}}endnoteRef") is not None:
                    continue
                for t_elem in r_elem.iter(f"{{{w_ns}}}t"):
                    t_elem.text = ""
            continue

        line = text_lines[p_idx]
        first_text_set = False
        for r_elem in para.iter(f"{{{w_ns}}}r"):
            # Skip the reference marker run
            if r_elem.find(f"{{{w_ns}}}footnoteRef") is not None:
                continue
            if r_elem.find(f"{{{w_ns}}}endnoteRef") is not None:
                continue
            for t_elem in r_elem.iter(f"{{{w_ns}}}t"):
                if not first_text_set:
                    t_elem.text = line
                    # Preserve leading/trailing whitespace
                    t_elem.set(
                        "{http://www.w3.org/XML/1998/namespace}space",
                        "preserve",
                    )
                    first_text_set = True
                else:
                    t_elem.text = ""


def _extract_odt_footnotes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts footnotes and endnotes from an ODT file via ZIP+lxml.

    ODT stores footnotes as ``<text:note>`` elements inline in
    ``content.xml``.  The ``text:note-class`` attribute distinguishes
    footnotes from endnotes.

    Args:
        file_path: Path to the .odt file.

    Returns:
        list: (location_key, text) pairs.
    """
    text_ns = _ODF_NS["text"]

    with zipfile.ZipFile(file_path, "r") as zf:
        if "content.xml" not in zf.namelist():
            return []
        data = zf.read("content.xml")

    root = etree.fromstring(data)
    texts: list[tuple[str, str]] = []
    text_p_tag = f"{{{text_ns}}}p"

    for note in root.iter(f"{{{text_ns}}}note"):
        note_id = note.get(f"{{{text_ns}}}id", "")
        note_class = note.get(f"{{{text_ns}}}note-class", "footnote")
        prefix = "endnote" if note_class == "endnote" else "footnote"

        # Extract text from note-body
        for body in note.iter(f"{{{text_ns}}}note-body"):
            text = _extract_odf_paragraph_text(body, text_p_tag)
            if text and text.strip():
                texts.append((f"{prefix}:{note_id}", text))

    return texts


def _inject_odt_footnotes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated footnotes/endnotes into an ODT file via ZIP+lxml.

    Args:
        output_path: Path to the .odt file to modify in place.
        translations: Mapping of 'footnote:{id}' / 'endnote:{id}' to
                      translated text.
    """
    text_ns = _ODF_NS["text"]

    with zipfile.ZipFile(output_path, "r") as zf:
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item) for item in all_items}

    if "content.xml" not in file_data:
        return

    root = etree.fromstring(file_data["content.xml"])
    text_p_tag = f"{{{text_ns}}}p"
    modified = False

    for note in root.iter(f"{{{text_ns}}}note"):
        note_id = note.get(f"{{{text_ns}}}id", "")
        note_class = note.get(f"{{{text_ns}}}note-class", "footnote")
        prefix = "endnote" if note_class == "endnote" else "footnote"
        key = f"{prefix}:{note_id}"

        if key not in translations:
            continue

        for body in note.iter(f"{{{text_ns}}}note-body"):
            _inject_odf_paragraph_text(
                body,
                translations[key],
                text_p_tag,
            )
            modified = True

    if not modified:
        return

    file_data["content.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )
    _rewrite_zip_content(output_path, file_data, all_items)


def _extract_win32com_word_footnotes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts footnotes and endnotes from a DOC file via win32com.

    Args:
        file_path: Path to the .doc file.

    Returns:
        list: (location_key, text) pairs.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, file_path)
    try:
        texts: list[tuple[str, str]] = []
        # Footnotes
        for i in range(1, doc.Footnotes.Count + 1):
            fn = doc.Footnotes(i)
            text = fn.Range.Text.rstrip("\r\n")
            if text and text.strip():
                texts.append((f"footnote:{i}", text))
        # Endnotes
        for i in range(1, doc.Endnotes.Count + 1):
            en = doc.Endnotes(i)
            text = en.Range.Text.rstrip("\r\n")
            if text and text.strip():
                texts.append((f"endnote:{i}", text))
        return texts
    finally:
        _win32com_close(app, doc, pycom)


def _inject_win32com_word_footnotes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated footnotes/endnotes into a DOC file via win32com.

    Args:
        output_path: Path to the .doc file to modify in place.
        translations: Mapping of 'footnote:{id}' / 'endnote:{id}' to
                      translated text.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, output_path)
    try:
        for i in range(1, doc.Footnotes.Count + 1):
            key = f"footnote:{i}"
            if key in translations:
                doc.Footnotes(i).Range.Text = translations[key]
        for i in range(1, doc.Endnotes.Count + 1):
            key = f"endnote:{i}"
            if key in translations:
                doc.Endnotes(i).Range.Text = translations[key]
        doc.SaveAs(str(output_path))
    finally:
        _win32com_close(app, doc, pycom)


def _extract_uno_writer_footnotes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts footnotes and endnotes from a DOC/ODT file via UNO.

    Args:
        file_path: Path to the document.

    Returns:
        list: (location_key, text) pairs.
    """
    doc = _uno_open(file_path)
    try:
        texts: list[tuple[str, str]] = []
        # Footnotes
        footnotes = doc.getFootnotes()
        for i in range(footnotes.getCount()):
            fn = footnotes.getByIndex(i)
            text = fn.getString().strip()
            if text:
                texts.append((f"footnote:{i + 1}", text))
        # Endnotes
        endnotes = doc.getEndnotes()
        for i in range(endnotes.getCount()):
            en = endnotes.getByIndex(i)
            text = en.getString().strip()
            if text:
                texts.append((f"endnote:{i + 1}", text))
        return texts
    finally:
        doc.close(True)


def _inject_uno_writer_footnotes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated footnotes/endnotes into a document via UNO.

    Args:
        output_path: Path to the document to modify in place.
        translations: Mapping of 'footnote:{id}' / 'endnote:{id}' to
                      translated text.
    """
    doc = _uno_open(output_path)
    try:
        footnotes = doc.getFootnotes()
        for i in range(footnotes.getCount()):
            key = f"footnote:{i + 1}"
            if key in translations:
                footnotes.getByIndex(i).setString(translations[key])
        endnotes = doc.getEndnotes()
        for i in range(endnotes.getCount()):
            key = f"endnote:{i + 1}"
            if key in translations:
                endnotes.getByIndex(i).setString(translations[key])
        _uno_save(doc, output_path)
    finally:
        doc.close(True)


# ---------------------------------------------------------------------------
# Win32com shape helpers
# ---------------------------------------------------------------------------


def _extract_win32com_word_shapes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts text from shapes/text boxes in a Word document via win32com.

    When a shape's text range has mixed per-run formatting, inline HTML is
    emitted via ``_win32com_range_runs_to_html`` so the LLM can preserve it.

    Args:
        file_path: Path to the .doc file.

    Returns:
        list: (location_key, text) pairs with keys like 'shape:{index}'.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for i in range(1, doc.Shapes.Count + 1):
            shape = doc.Shapes(i)
            try:
                if shape.TextFrame.HasText:
                    text_rng = shape.TextFrame.TextRange
                    raw_text = text_rng.Text
                    if raw_text and raw_text.strip():
                        if _has_win32com_range_mixed_formatting(
                            text_rng
                        ) or _has_win32com_range_hyperlinks(text_rng):
                            text = _win32com_range_runs_to_html(text_rng)
                        else:
                            text = raw_text
                        texts.append((f"shape:{i - 1}", text))
            except Exception:  # noqa: BLE001
                pass  # Shape has no text frame
        return texts
    finally:
        _win32com_close(app, doc, pycom, save_close=True)


def _inject_win32com_word_shapes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated text into Word shapes via win32com.

    When the translated text contains inline HTML formatting tags,
    per-segment formatting is applied via ``_inject_win32com_word_html_runs``.
    Otherwise, plain text is set with uniform font save/restore.

    Args:
        output_path: Path to the .doc file to modify in place.
        translations: Mapping of 'shape:{index}' to translated text.
    """
    app, doc, pycom = _win32com_open(_APP_WORD, output_path)
    try:
        for i in range(1, doc.Shapes.Count + 1):
            key = f"shape:{i - 1}"
            if key in translations:
                with contextlib.suppress(Exception):
                    text_rng = doc.Shapes(i).TextFrame.TextRange
                    orig_text = text_rng.Text
                    translation = translations[key]
                    if _FORMATTING_HTML_RE.search(translation):
                        _inject_win32com_word_html_runs(
                            doc,
                            text_rng,
                            translation,
                            orig_text,
                            is_cell=True,
                        )
                    else:
                        font_saved = _save_win32com_font(text_rng.Font)
                        text_rng.Text = translation
                        # Re-acquire range after text change
                        text_rng = doc.Shapes(i).TextFrame.TextRange
                        _restore_win32com_font(
                            text_rng.Font,
                            font_saved,
                            original_text=orig_text,
                            translated_text=translation,
                        )

        doc.Save()
    finally:
        _win32com_close(app, doc, pycom, save_close=True)


def _extract_win32com_excel_shapes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts text from shapes in an Excel workbook via win32com.

    When a shape's text range has mixed per-run formatting, inline HTML is
    emitted via ``_win32com_range_runs_to_html`` so the LLM can preserve it.

    Args:
        file_path: Path to the .xls file.

    Returns:
        list: (location_key, text) pairs with keys like
              'shape:{sheet_name}:{index}'.
    """
    app, wb, pycom = _win32com_open(_APP_EXCEL, file_path)
    try:
        texts: list[tuple[str, str]] = []
        for ws in wb.Worksheets:
            for i in range(1, ws.Shapes.Count + 1):
                shape = ws.Shapes(i)
                try:
                    if shape.TextFrame2.HasText:
                        text_rng = shape.TextFrame2.TextRange
                        raw_text = text_rng.Text
                        if raw_text and raw_text.strip():
                            if _has_win32com_range_mixed_formatting(text_rng):
                                text = _win32com_range_runs_to_html(text_rng)
                            else:
                                text = raw_text
                            texts.append(
                                (f"shape:{ws.Name}:{i - 1}", text),
                            )
                except Exception:  # noqa: BLE001
                    pass  # Shape has no text frame
        return texts
    finally:
        _win32com_close(app, wb, pycom, save_close=True)


def _inject_win32com_excel_html_runs(  # noqa: PLR0912
    text_rng: object,
    html_text: str,
    original_text: str = "",
    *,
    target_lang: str = "",
) -> None:
    """Replaces an Excel shape's text with HTML-formatted segments via win32com.

    Parses ``html_text`` via ``_parse_html_formatting``, sets the full
    plain text on the range, then applies per-segment formatting using
    ``Characters(start, length)`` sub-ranges (1-based indexing).

    Args:
        text_rng: A win32com ``TextRange2`` COM object.
        html_text: Translated text with inline ``<b>/<i>/<u>/<s>`` tags.
        original_text: The text before translation (for script detection).
        target_lang: Target language name for font substitution.
    """
    segments = _parse_html_formatting(html_text)
    if not segments:
        # Strip residual tags so literal HTML doesn't appear in the document
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        text_rng.Text = plain
        return

    # Save base font Name BEFORE text assignment
    saved_name: str | None = None
    try:
        from src.constants.office import WIN32COM_UNDEFINED  # noqa: PLC0415

        name_val = text_rng.Font.Name
        if name_val != WIN32COM_UNDEFINED:
            saved_name = str(name_val)
    except Exception:  # noqa: BLE001
        pass

    # Build full plain text and set it
    full_text = "".join(seg.text for seg in segments)
    text_rng.Text = full_text

    # Restore base font Name on the whole range (script-aware)
    if saved_name:
        font_name = _substitute_font(
            saved_name,
            original_text,
            full_text,
            target_lang,
        )
        if font_name is not None:
            with contextlib.suppress(Exception):
                text_rng.Font.Name = font_name

    # Apply per-segment formatting via Characters (1-based indexing)
    offset = 0
    for seg in segments:
        seg_len = len(seg.text)
        if seg_len == 0:
            continue
        try:
            char_rng = text_rng.Characters(offset + 1, seg_len)
            char_rng.Font.Bold = seg.bold
            char_rng.Font.Italic = seg.italic
            char_rng.Font.Underline = seg.underline
            char_rng.Font.StrikeThrough = seg.strike
            # Superscript / Subscript — via BaselineOffset
            if seg.superscript:
                with contextlib.suppress(Exception):
                    char_rng.Font.BaselineOffset = 0.3
            elif seg.subscript:
                with contextlib.suppress(Exception):
                    char_rng.Font.BaselineOffset = -0.25
            else:
                with contextlib.suppress(Exception):
                    char_rng.Font.BaselineOffset = 0.0
            if seg.font_size_pt is not None:
                char_rng.Font.Size = seg.font_size_pt
            if seg.color_hex is not None:
                char_rng.Font.Color = _color_hex_to_win32com(
                    seg.color_hex,
                )
            # Background — Font.Highlight (Office 365 / 2019+);
            # silently skipped on older versions.
            if seg.bg_color_hex is not None:
                with contextlib.suppress(Exception):
                    char_rng.Font.Highlight.ForeColor.RGB = _color_hex_to_win32com(
                        seg.bg_color_hex
                    )
        except Exception:  # noqa: BLE001
            pass  # Defensive — shape may not support all Font properties
        offset += seg_len


def _inject_win32com_excel_shapes(
    output_path: Path,
    translations: dict[str, str],
    *,
    target_lang: str = "",
) -> None:
    """Injects translated text into Excel shapes via win32com.

    When the translated text contains inline HTML formatting tags,
    per-segment formatting is applied via
    ``_inject_win32com_excel_html_runs``.  Otherwise, plain text is set
    with uniform font save/restore.

    Args:
        output_path: Path to the .xls file to modify in place.
        translations: Mapping of 'shape:{sheet_name}:{index}' to translated text.
        target_lang: Target language name for font substitution.
    """
    app, wb, pycom = _win32com_open(_APP_EXCEL, output_path)
    try:
        for ws in wb.Worksheets:
            for i in range(1, ws.Shapes.Count + 1):
                key = f"shape:{ws.Name}:{i - 1}"
                if key in translations:
                    with contextlib.suppress(Exception):
                        text_rng = ws.Shapes(i).TextFrame2.TextRange
                        orig_text = text_rng.Text
                        translation = translations[key]
                        if _FORMATTING_HTML_RE.search(translation):
                            _inject_win32com_excel_html_runs(
                                text_rng,
                                translation,
                                orig_text,
                                target_lang=target_lang,
                            )
                        else:
                            font_saved = _save_win32com_font(text_rng.Font)
                            text_rng.Text = translation
                            # Re-acquire range after text change
                            text_rng = ws.Shapes(i).TextFrame2.TextRange
                            _restore_win32com_font(
                                text_rng.Font,
                                font_saved,
                                original_text=orig_text,
                                translated_text=translation,
                                target_lang=target_lang,
                            )

        wb.Save()
    finally:
        _win32com_close(app, wb, pycom, save_close=True)


# ---------------------------------------------------------------------------
# UNO shape helpers
# ---------------------------------------------------------------------------


def _extract_uno_writer_shapes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts text from shapes/text boxes in a Writer document via UNO.

    When any paragraph within a shape has mixed per-run formatting, the
    entire shape is extracted as inline HTML via ``_uno_runs_to_html``
    (paragraphs joined by newlines).  Otherwise, plain text is returned.

    Args:
        file_path: Path to the .doc or .docx file.

    Returns:
        list: (location_key, text) pairs with keys like 'shape:{index}'.
    """
    doc = _uno_open(file_path)
    texts: list[tuple[str, str]] = []
    try:
        draw_page = doc.getDrawPage()
        for i in range(draw_page.getCount()):
            shape = draw_page.getByIndex(i)
            if shape.supportsService("com.sun.star.drawing.Text"):
                para_enum = shape.createEnumeration()
                paras: list[object] = []
                while para_enum.hasMoreElements():
                    paras.append(para_enum.nextElement())
                # Check any paragraph for mixed formatting or hyperlinks
                if any(
                    _has_uno_mixed_formatting(p) or _has_uno_hyperlinks(p)
                    for p in paras
                ):
                    text = "\n".join(_uno_runs_to_html(p) for p in paras)
                else:
                    text = shape.getString()
                if text and text.strip():
                    texts.append((f"shape:{i}", text))
    finally:
        doc.close(True)

    return texts


def _inject_uno_writer_shapes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated text into Writer shapes via UNO.

    When the translated text contains inline HTML formatting tags,
    dispatches to ``_inject_uno_para_text`` for per-run formatting on
    each paragraph (lines separated by newlines).  Otherwise, uses plain
    ``setString`` with shape-level property save/restore.

    Args:
        output_path: Path to the .doc or .docx file to modify in place.
        translations: Mapping of 'shape:{index}' to translated text.
    """
    doc = _uno_open(output_path)
    try:
        draw_page = doc.getDrawPage()
        for i in range(draw_page.getCount()):
            key = f"shape:{i}"
            if key in translations:
                shape = draw_page.getByIndex(i)
                if shape.supportsService("com.sun.star.drawing.Text"):
                    translation = translations[key]
                    handled = False
                    if _FORMATTING_HTML_RE.search(translation):
                        para_enum = shape.createEnumeration()
                        paras: list[object] = []
                        while para_enum.hasMoreElements():
                            paras.append(para_enum.nextElement())
                        # Split HTML by newlines to match extraction's
                        # per-paragraph _uno_runs_to_html join.
                        lines = translation.split("\n")
                        if paras:
                            for p_idx, para in enumerate(paras):
                                line = lines[p_idx] if p_idx < len(lines) else ""
                                _inject_uno_para_text(para, line)
                            handled = True
                    if not handled:
                        orig_text = shape.getString()
                        char_saved = _save_uno_char_props(shape)
                        shape.setString(translation)
                        _restore_uno_char_props(
                            shape,
                            char_saved,
                            original_text=orig_text,
                            translated_text=translation,
                        )

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


def _extract_uno_calc_shapes(
    file_path: Path,
) -> list[tuple[str, str]]:
    """Extracts text from shapes in a Calc spreadsheet via UNO.

    When any paragraph within a shape has mixed per-run formatting, the
    entire shape is extracted as inline HTML via ``_uno_runs_to_html``
    (paragraphs joined by newlines).  Otherwise, plain text is returned.

    Args:
        file_path: Path to the .xls file.

    Returns:
        list: (location_key, text) pairs with keys like
              'shape:{sheet_name}:{index}'.
    """
    doc = _uno_open(file_path)
    texts: list[tuple[str, str]] = []
    try:
        sheets = doc.getSheets()
        for s_idx in range(sheets.getCount()):
            sheet = sheets.getByIndex(s_idx)
            sheet_name = sheet.getName()
            draw_page = sheet.getDrawPage()
            for i in range(draw_page.getCount()):
                shape = draw_page.getByIndex(i)
                if shape.supportsService("com.sun.star.drawing.Text"):
                    para_enum = shape.createEnumeration()
                    paras: list[object] = []
                    while para_enum.hasMoreElements():
                        paras.append(para_enum.nextElement())
                    # Check any paragraph for mixed formatting or hyperlinks
                    if any(
                        _has_uno_mixed_formatting(p) or _has_uno_hyperlinks(p)
                        for p in paras
                    ):
                        text = "\n".join(_uno_runs_to_html(p) for p in paras)
                    else:
                        text = shape.getString()
                    if text and text.strip():
                        texts.append((f"shape:{sheet_name}:{i}", text))
    finally:
        doc.close(True)

    return texts


def _inject_uno_calc_shapes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated text into Calc shapes via UNO.

    When the translated text contains inline HTML formatting tags,
    dispatches to ``_inject_uno_para_text`` for per-run formatting on
    each paragraph (lines separated by newlines).  Otherwise, uses plain
    ``setString`` with shape-level property save/restore.

    Args:
        output_path: Path to the .xls file to modify in place.
        translations: Mapping of 'shape:{sheet_name}:{index}' to translated text.
    """
    doc = _uno_open(output_path)
    try:
        sheets = doc.getSheets()
        for s_idx in range(sheets.getCount()):
            sheet = sheets.getByIndex(s_idx)
            sheet_name = sheet.getName()
            draw_page = sheet.getDrawPage()
            for i in range(draw_page.getCount()):
                key = f"shape:{sheet_name}:{i}"
                if key in translations:
                    shape = draw_page.getByIndex(i)
                    if shape.supportsService("com.sun.star.drawing.Text"):
                        translation = translations[key]
                        handled = False
                        if _FORMATTING_HTML_RE.search(translation):
                            para_enum = shape.createEnumeration()
                            paras: list[object] = []
                            while para_enum.hasMoreElements():
                                paras.append(para_enum.nextElement())
                            if paras:
                                lines = translation.split("\n")
                                for p_idx, para in enumerate(paras):
                                    line = lines[p_idx] if p_idx < len(lines) else ""
                                    _inject_uno_para_text(para, line)
                                handled = True
                        if not handled:
                            orig_text = shape.getString()
                            char_saved = _save_uno_char_props(shape)
                            shape.setString(translation)
                            _restore_uno_char_props(
                                shape,
                                char_saved,
                                original_text=orig_text,
                                translated_text=translation,
                            )

        _uno_save(doc, output_path)
    finally:
        doc.close(True)


# ---------------------------------------------------------------------------
# DOCX shape helpers — ZIP + lxml
# ---------------------------------------------------------------------------


def _read_txbx_data(
    txbx_el: object,
) -> tuple[str, list[object]]:
    r"""Reads plain text and ``<w:t>`` elements from a single ``<wps:txbx>``.

    Iterates paragraph-by-paragraph to preserve structural newlines between
    paragraphs.

    Args:
        txbx_el: An lxml element for a ``<wps:txbx>`` text box.

    Returns:
        Tuple of (plain_text, t_elements) where plain_text is the stripped
        concatenated text of all paragraphs joined by ``'\n'``, and
        t_elements is the flat list of all ``<w:t>`` elements found.
    """
    w_p_tag = f"{{{_WORDML_NS}}}p"
    w_t_tag = f"{{{_WORDML_NS}}}t"

    p_texts: list[str] = []
    all_t_els: list[object] = []
    for p_el in txbx_el.findall(f".//{w_p_tag}"):
        t_els = list(p_el.iter(w_t_tag))
        t_parts = [t.text for t in t_els if t.text]
        if t_parts:
            p_texts.append("".join(t_parts))
        all_t_els.extend(t_els)

    plain = "\n".join(p_texts).strip()
    return plain, all_t_els


def _wps_txbx_to_text_or_html(  # noqa: PLR0912, PLR0915
    txbx_el: object,
    char_styles: dict[
        str, tuple[bool, bool, bool, bool, float | None, str | None, str | None]
    ]
    | None = None,
    hyperlink_rels: dict[str, str] | None = None,
) -> str:
    r"""Extracts text from a ``<wps:txbx>`` element, using HTML when formatting varies.

    Iterates direct children of each ``<w:p>`` paragraph — both ``<w:r>``
    runs and ``<w:hyperlink>`` wrappers.  If run formatting varies or any
    hyperlinks are present, wraps the text in inline HTML tags via
    ``_wrap_with_tags`` and ``<a href="...">`` tags.  Otherwise returns
    plain text identical to ``_read_txbx_data``.  Paragraphs are joined
    with ``'\n'``.

    Character-style references (``<w:rStyle>``) are resolved when
    *char_styles* is provided: the style supplies base formatting and
    direct ``<w:rPr>`` attributes override.

    All ``<w:t>`` elements within a single run are concatenated so that
    split runs (e.g. from spell-checking) do not silently drop text.

    Args:
        txbx_el: An lxml element for a ``<wps:txbx>`` text box.
        char_styles: Mapping of style IDs to formatting tuples, as returned
            by ``_parse_docx_char_styles``.  ``None`` disables style resolution.
        hyperlink_rels: Mapping of relationship IDs to target URLs,
            parsed from the part's ``.rels`` file.

    Returns:
        Plain text or inline-HTML string representing the text box content.
    """
    w = _WORDML_NS
    w_p_tag = f"{{{w}}}p"
    w_r_tag = f"{{{w}}}r"
    w_rpr_tag = f"{{{w}}}rPr"
    w_t_tag = f"{{{w}}}t"
    w_rstyle_tag = f"{{{w}}}rStyle"
    w_hyperlink_tag = f"{{{w}}}hyperlink"
    r_id_attr = (
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    w_anchor_attr = f"{{{w}}}anchor"

    # Collect run data: (text, bold, italic, underline, strike, size, color, bg, url)
    all_run_data: list[
        tuple[
            str,
            bool,
            bool,
            bool,
            bool,
            float | None,
            str | None,
            str | None,
            str | None,
        ]
    ] = []
    # Track which run indices span each paragraph
    para_run_spans: list[tuple[int, int]] = []
    has_hyperlinks = False

    def _resolve_run_formatting(
        r_el: object,
        url: str | None,
    ) -> None:
        """Collects formatting data from a single ``<w:r>`` element."""
        nonlocal has_hyperlinks
        # Concatenate all <w:t> elements in this run
        t_els = r_el.findall(w_t_tag)
        run_text = "".join(t.text for t in t_els if t.text)
        if not run_text:
            return
        rpr_el = r_el.find(w_rpr_tag)
        b, i, u, s, sz, clr, bg = _read_wml_rpr_formatting(rpr_el)

        # Resolve character style if present
        if char_styles and rpr_el is not None:
            rstyle_el = rpr_el.find(w_rstyle_tag)
            if rstyle_el is not None:
                sid = rstyle_el.get(f"{{{w}}}val") or rstyle_el.get("val") or ""
                if sid in char_styles:
                    sb, si, su, ss, ssz, sclr, sbg = char_styles[sid]
                    # Style provides base; direct formatting overrides.
                    if rpr_el.find(f"{{{w}}}b") is None:
                        b = sb
                    if rpr_el.find(f"{{{w}}}i") is None:
                        i = si
                    if rpr_el.find(f"{{{w}}}u") is None:
                        u = su
                    if rpr_el.find(f"{{{w}}}strike") is None:
                        s = ss
                    if rpr_el.find(f"{{{w}}}sz") is None:
                        sz = ssz
                    if rpr_el.find(f"{{{w}}}color") is None:
                        clr = sclr
                    if (
                        rpr_el.find(f"{{{w}}}shd") is None
                        and rpr_el.find(f"{{{w}}}highlight") is None
                    ):
                        bg = sbg

        if url is not None:
            has_hyperlinks = True
        all_run_data.append((run_text, b, i, u, s, sz, clr, bg, url))

    for p_el in txbx_el.findall(f".//{w_p_tag}"):
        para_start = len(all_run_data)
        for child in p_el:
            if child.tag == w_hyperlink_tag:
                # Resolve hyperlink URL from r:id or w:anchor
                url: str | None = None
                if hyperlink_rels:
                    rid = child.get(r_id_attr)
                    if rid and rid in hyperlink_rels:
                        url = hyperlink_rels[rid]
                if url is None:
                    anchor = child.get(w_anchor_attr)
                    if anchor:
                        url = f"#{anchor}"
                for r_el in child.findall(w_r_tag):
                    _resolve_run_formatting(r_el, url)
            elif child.tag == w_r_tag:
                _resolve_run_formatting(child, None)
        para_run_spans.append((para_start, len(all_run_data)))

    if not all_run_data:
        return ""

    # Check whether formatting varies across all runs
    sigs = {(d[1], d[2], d[3], d[4], d[5], d[6], d[7]) for d in all_run_data}
    if len(sigs) <= 1 and not has_hyperlinks:
        # Uniform formatting, no hyperlinks — return plain text for consistency
        plain, _ = _read_txbx_data(txbx_el)
        return plain

    # Mixed formatting or hyperlinks — emit inline HTML
    sizes = [d[5] for d in all_run_data]
    colors = [d[6] for d in all_run_data]
    bgs = [d[7] for d in all_run_data]
    has_size_variation = len(set(sizes)) > 1
    has_color_variation = len(set(colors)) > 1
    has_bg_variation = len(set(bgs)) > 1
    # base_size/color/bg are always None so every run with an explicit value
    # gets its own <span>.  Using most-common as base loses that value during
    # injection when the first run is not the most-common one.
    base_size = None
    base_color = None
    base_bg = None

    para_htmls: list[str] = []
    for start, end in para_run_spans:
        parts: list[str] = []
        current_url: str | None = None
        for text, bold, italic, underline, strike, sz, clr, bg, url in all_run_data[
            start:end
        ]:
            # Manage <a> tag transitions
            if url != current_url:
                if current_url is not None:
                    parts.append("</a>")
                if url is not None:
                    parts.append(f'<a href="{html.escape(url, quote=True)}">')
                current_url = url
            parts.append(
                _wrap_with_tags(
                    html.escape(text),
                    bold,
                    italic,
                    underline,
                    strike,
                    sz if sz != base_size else None,
                    clr if clr != base_color else None,
                    has_size_variation=has_size_variation,
                    has_color_variation=has_color_variation,
                    bg_color_hex=bg if bg != base_bg else None,
                    has_bg_variation=has_bg_variation,
                )
            )
        # Close trailing <a> tag
        if current_url is not None:
            parts.append("</a>")
        if parts:
            para_htmls.append("".join(parts))
    return "\n".join(para_htmls)


def _inject_wps_txbx_plain(
    txbx_el: object,
    plain_text: str,
    t_elements: list[object],
) -> None:
    r"""Injects plain text into a ``<wps:txbx>`` element in-place.

    Sets the first ``<w:t>`` element's text to the first line and appends
    ``<w:br/>`` and new ``<w:t>`` elements for subsequent lines.  Remaining
    original ``<w:t>`` elements are cleared.

    Args:
        txbx_el: An lxml element for the ``<wps:txbx>`` text box.
        plain_text: The translated plain text (lines separated by ``'\n'``).
        t_elements: Flat list of all ``<w:t>`` elements from the text box.
    """
    lines = plain_text.split("\n")
    first_t = t_elements[0]
    first_t.text = lines[0]
    parent_r = first_t.getparent()

    w_br_tag = f"{{{_WORDML_NS}}}br"
    w_t_tag = f"{{{_WORDML_NS}}}t"
    for line in lines[1:]:
        br = etree.Element(w_br_tag)
        parent_r.append(br)
        t_new = etree.Element(w_t_tag)
        t_new.text = line
        parent_r.append(t_new)

    for t_el in t_elements[1:]:
        t_el.text = ""


def _inject_wps_txbx_html_runs(  # noqa: PLR0912, PLR0915
    txbx_el: object,
    html_text: str,
    rels_adder: Callable[[str], str] | None = None,
) -> None:
    r"""Injects HTML-formatted text into a ``<wps:txbx>`` element in-place.

    Parses ``html_text`` via ``_parse_html_formatting`` to obtain
    ``_FormattedSegment`` objects.  Segments containing ``'\n'`` are split
    across multiple ``<w:p>`` elements.  Existing run children are cleared
    and replaced with new ``<w:r>/<w:rPr>/<w:t>`` elements.  Excess
    paragraphs are removed; new ones are cloned from the last existing
    paragraph when more are needed.

    When *rels_adder* is provided, segments with ``hyperlink_url`` are
    wrapped in ``<w:hyperlink>`` elements with the relationship ID
    returned by the callback.

    Args:
        txbx_el: An lxml element for the ``<wps:txbx>`` text box.
        html_text: Translated HTML string with inline formatting tags.
        rels_adder: Callback that accepts a URL string and returns a
            relationship ID (``r:id``) for the hyperlink.  ``None`` disables
            hyperlink injection.
    """
    import copy  # noqa: PLC0415

    segments = _parse_html_formatting(html_text)
    if not segments:
        # Strip residual tags and fall back to plain-text injection
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        _, t_elements = _read_txbx_data(txbx_el)
        if t_elements:
            _inject_wps_txbx_plain(txbx_el, plain, t_elements)
        return

    w = _WORDML_NS
    w_txbx_content_tag = f"{{{w}}}txbxContent"
    w_p_tag = f"{{{w}}}p"
    w_r_tag = f"{{{w}}}r"
    w_rpr_tag = f"{{{w}}}rPr"
    w_t_tag = f"{{{w}}}t"
    w_b_tag = f"{{{w}}}b"
    w_i_tag = f"{{{w}}}i"
    w_u_tag = f"{{{w}}}u"
    w_strike_tag = f"{{{w}}}strike"
    w_sz_tag = f"{{{w}}}sz"
    w_color_tag = f"{{{w}}}color"
    w_val_attr = f"{{{w}}}val"
    w_hyperlink_tag = f"{{{w}}}hyperlink"
    r_id_attr = (
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    w_anchor_attr = f"{{{w}}}anchor"
    xml_space_attr = "{http://www.w3.org/XML/1998/namespace}space"

    # Find <w:txbxContent> — direct container of paragraphs inside <wps:txbx>
    txbx_content = txbx_el.find(w_txbx_content_tag)
    if txbx_content is None:
        return

    # Get existing <w:p> direct children
    existing_paras = [c for c in txbx_content if c.tag == w_p_tag]
    if not existing_paras:
        return

    # Save the first run's rPr as a base formatting template for new runs
    w_shd_tag = f"{{{w}}}shd"
    w_highlight_tag = f"{{{w}}}highlight"
    first_r = existing_paras[0].find(f".//{w_r_tag}")
    base_rpr: object | None = None
    if first_r is not None:
        existing_rpr = first_r.find(w_rpr_tag)
        if existing_rpr is not None:
            base_rpr = copy.deepcopy(existing_rpr)
            # Strip highlight/shading from base so bg doesn't spread
            for tag in (w_shd_tag, w_highlight_tag):
                el = base_rpr.find(tag)
                if el is not None:
                    base_rpr.remove(el)

    # Split segments by '\n' into paragraph groups (one group per <w:p>)
    # Each entry is a _FormattedSegment (preserves hyperlink_url)
    para_groups: list[list[_FormattedSegment]] = [[]]
    for seg in segments:
        lines = seg.text.split("\n")
        for j, line_text in enumerate(lines):
            if j > 0:
                para_groups.append([])
            if line_text:
                para_groups[-1].append(seg._replace(text=line_text))
    # Remove trailing empty paragraph groups
    while para_groups and not para_groups[-1]:
        para_groups.pop()
    if not para_groups:
        return

    n_groups = len(para_groups)
    n_existing = len(existing_paras)

    # Remove excess paragraphs from txbxContent
    for p_el in existing_paras[n_groups:]:
        txbx_content.remove(p_el)

    # Clone the last existing paragraph if more paragraphs are needed
    para_template = existing_paras[-1]
    for _ in range(max(0, n_groups - n_existing)):
        new_p = copy.deepcopy(para_template)
        for child in list(new_p):
            if child.tag in (w_r_tag, w_hyperlink_tag):
                new_p.remove(child)
        txbx_content.append(new_p)

    # Re-collect paragraphs after structural changes
    updated_paras = [c for c in txbx_content if c.tag == w_p_tag]

    # Bold/italic/underline/strike/shd/highlight are always controlled by
    # the HTML.  Size and color are only present when they *vary*.
    _always_skip = {w_b_tag, w_i_tag, w_u_tag, w_strike_tag, w_shd_tag, w_highlight_tag}

    def _make_run(seg: _FormattedSegment) -> object:  # noqa: PLR0912
        """Creates a ``<w:r>`` element from a formatted segment."""
        r_el = etree.Element(w_r_tag)
        rpr = etree.Element(w_rpr_tag)

        # Determine which tags to skip when copying base properties.
        run_skip = set(_always_skip)
        if seg.font_size_pt is not None:
            run_skip.add(w_sz_tag)
        if seg.color_hex is not None:
            run_skip.add(w_color_tag)

        # Copy base properties (font name, theme, etc.) minus overridden ones
        if base_rpr is not None:
            for child in base_rpr:
                if child.tag not in run_skip:
                    rpr.append(copy.deepcopy(child))
        if seg.bold:
            etree.SubElement(rpr, w_b_tag)
        if seg.italic:
            etree.SubElement(rpr, w_i_tag)
        if seg.underline:
            u_el = etree.SubElement(rpr, w_u_tag)
            u_el.set(w_val_attr, "single")
        if seg.strike:
            etree.SubElement(rpr, w_strike_tag)
        if seg.font_size_pt is not None:
            sz_el = etree.SubElement(rpr, w_sz_tag)
            sz_el.set(w_val_attr, str(int(seg.font_size_pt * 2)))
        if seg.color_hex is not None:
            color_el = etree.SubElement(rpr, w_color_tag)
            color_el.set(w_val_attr, seg.color_hex.lstrip("#").upper())
        if seg.bg_color_hex is not None:
            shd_el = etree.SubElement(rpr, w_shd_tag)
            shd_el.set(w_val_attr, "clear")
            shd_el.set(f"{{{w}}}color", "auto")
            shd_el.set(f"{{{w}}}fill", seg.bg_color_hex.lstrip("#").upper())
        if len(rpr):
            r_el.append(rpr)
        t_el = etree.Element(w_t_tag)
        t_el.text = seg.text
        if seg.text.startswith(" ") or seg.text.endswith(" "):
            t_el.set(xml_space_attr, "preserve")
        r_el.append(t_el)
        return r_el

    # Replace runs in each paragraph with the corresponding formatted segments
    for p_idx, group in enumerate(para_groups):
        p_el = updated_paras[p_idx]
        # Remove all run and hyperlink children
        for child in list(p_el):
            if child.tag in (w_r_tag, w_hyperlink_tag):
                p_el.remove(child)

        # Group segments by hyperlink and create runs
        current_url: str | None = None
        hyperlink_elem: object | None = None

        for seg in group:
            new_r = _make_run(seg)

            if seg.hyperlink_url:
                if seg.hyperlink_url != current_url:
                    # Start a new <w:hyperlink> group
                    hyperlink_elem = etree.Element(w_hyperlink_tag)
                    if seg.hyperlink_url.startswith("#"):
                        # Internal bookmark anchor
                        hyperlink_elem.set(
                            w_anchor_attr,
                            seg.hyperlink_url[1:],
                        )
                    elif rels_adder is not None:
                        # External URL — create relationship
                        rid = rels_adder(seg.hyperlink_url)
                        hyperlink_elem.set(r_id_attr, rid)
                    else:
                        # No rels_adder — cannot create relationship
                        hyperlink_elem = None
                    if hyperlink_elem is not None:
                        p_el.append(hyperlink_elem)
                    current_url = seg.hyperlink_url
                if hyperlink_elem is not None:
                    hyperlink_elem.append(new_r)
                else:
                    # Fallback: attach as plain run
                    p_el.append(new_r)
            else:
                if current_url is not None:
                    current_url = None
                    hyperlink_elem = None
                p_el.append(new_r)


def _collect_wps_texts(
    root: object,
) -> list[tuple[str, list[object]]]:
    """Finds all ``<wps:txbx>`` text boxes and their ``<w:t>`` elements.

    Delegates per-element data reading to ``_read_txbx_data``.

    Args:
        root: lxml root element of an XML part.

    Returns:
        list: Pairs of (concatenated_text, list_of_wt_elements).
    """
    wps_txbx_tag = f"{{{_WPS_NS}}}txbx"

    results: list[tuple[str, list[object]]] = []
    for txbx in root.iter(wps_txbx_tag):
        plain, t_els = _read_txbx_data(txbx)
        if plain and t_els:
            results.append((plain, t_els))
    return results


def _extract_docx_shapes(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from shapes/text boxes in a DOCX file via ZIP + lxml.

    Parses ``word/document.xml`` and ``word/header*.xml`` / ``word/footer*.xml``
    looking for ``<wps:txbx>`` elements that contain ``<w:t>`` runs.
    When run formatting varies or hyperlinks are present within a text box,
    inline HTML is emitted so the LLM can preserve it.

    Args:
        file_path: Path to the .docx file.

    Returns:
        list: (location_key, text) pairs with keys like 'shape:{index}'.
    """
    texts: list[tuple[str, str]] = []
    shape_idx = 0
    wps_txbx_tag = f"{{{_WPS_NS}}}txbx"

    with zipfile.ZipFile(file_path, "r") as zf:
        # Parse character styles once for style-based formatting resolution
        char_styles = _parse_docx_char_styles(zf)
        namelist = set(zf.namelist())

        # Collect XML parts that may contain shapes
        parts = ["word/document.xml"]
        for name in namelist:
            if (
                name.startswith("word/header") or name.startswith("word/footer")
            ) and name.endswith(".xml"):
                parts.append(name)

        for part_name in parts:
            if part_name not in namelist:
                continue
            root = etree.fromstring(zf.read(part_name))
            # Parse hyperlink relationships for this XML part
            rels_path = _get_rels_path(part_name)
            hyperlink_rels: dict[str, str] = {}
            if rels_path in namelist:
                hyperlink_rels = _parse_hyperlink_rels(zf.read(rels_path))

            for txbx in root.iter(wps_txbx_tag):
                plain, t_els = _read_txbx_data(txbx)
                if not plain or not t_els:
                    continue
                text = _wps_txbx_to_text_or_html(
                    txbx,
                    char_styles,
                    hyperlink_rels=hyperlink_rels,
                )
                texts.append((f"shape:{shape_idx}", text))
                shape_idx += 1

    return texts


def _inject_docx_shapes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated text into DOCX shapes/text boxes via ZIP + lxml.

    When the translated text contains inline HTML formatting tags,
    ``_inject_wps_txbx_html_runs`` is used to rebuild ``<w:r>`` elements
    with per-segment ``<w:rPr>`` formatting.  When ``<a href="...">``
    tags are present, hyperlink relationships are added to the part's
    ``.rels`` file.  Otherwise, plain text is injected via
    ``_inject_wps_txbx_plain``.

    Args:
        output_path: Path to the .docx file to modify in place.
        translations: Mapping of ``'shape:{index}'`` to translated text.
    """
    with zipfile.ZipFile(output_path, "r") as zf:
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item.filename) for item in all_items}

    # Collect XML parts that may contain shapes (body + headers/footers)
    parts = ["word/document.xml"]
    for name in file_data:
        if (
            name.startswith("word/header") or name.startswith("word/footer")
        ) and name.endswith(".xml"):
            parts.append(name)

    modified = False
    shape_idx = 0
    wps_txbx_tag = f"{{{_WPS_NS}}}txbx"

    for part_name in parts:
        if part_name not in file_data:
            continue
        root = etree.fromstring(file_data[part_name])
        part_modified = False

        # Parse existing rels for hyperlink injection
        rels_path = _get_rels_path(part_name)
        current_rels: bytes | None = file_data.get(rels_path)

        def _make_rels_adder(
            rp: str = rels_path,
        ) -> Callable[[str], str]:
            """Creates a closure that adds hyperlink rels for a part."""

            def adder(url: str) -> str:
                nonlocal current_rels
                new_xml, r_id = _add_hyperlink_to_rels(current_rels, url)
                current_rels = new_xml
                file_data[rp] = new_xml
                return r_id

            return adder

        rels_adder = _make_rels_adder()

        for txbx in root.iter(wps_txbx_tag):
            plain, t_elements = _read_txbx_data(txbx)
            if not plain or not t_elements:
                continue
            key = f"shape:{shape_idx}"
            shape_idx += 1
            if key not in translations:
                continue

            translation = translations[key]
            if _FORMATTING_HTML_RE.search(translation):
                _inject_wps_txbx_html_runs(
                    txbx,
                    translation,
                    rels_adder=rels_adder,
                )
            else:
                _inject_wps_txbx_plain(txbx, translation, t_elements)
            part_modified = True

        if part_modified:
            file_data[part_name] = etree.tostring(
                root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            modified = True

    if not modified:
        return
    _rewrite_zip_content(output_path, file_data, all_items)


# ---------------------------------------------------------------------------
# XLSX shape helpers — ZIP + lxml
# (reuses _extract_drawingml_text / _inject_drawingml_text)
# ---------------------------------------------------------------------------


def _resolve_xlsx_sheet_drawings(
    zf: object,
) -> list[tuple[str, str]]:
    """Resolves sheet-name → drawing-path mappings from an XLSX ZIP.

    Reads ``xl/workbook.xml`` to get sheet names and
    ``xl/worksheets/_rels/sheet{N}.xml.rels`` to find associated drawings.

    Args:
        zf: An open ``zipfile.ZipFile`` for the XLSX.

    Returns:
        list: Pairs of (sheet_name, drawing_xml_path) like
              ``("Sheet1", "xl/drawings/drawing1.xml")``.
    """
    results: list[tuple[str, str]] = []

    # Get sheet names from workbook.xml
    wb_xml = zf.read("xl/workbook.xml")
    wb_root = etree.fromstring(wb_xml)
    sheet_tag = f"{{{_SPREADSHEETML_NS}}}sheet"
    sheets_tag = f"{{{_SPREADSHEETML_NS}}}sheets"

    sheets_el = wb_root.find(sheets_tag)
    if sheets_el is None:
        return results

    sheet_names: list[str] = []
    for sheet_el in sheets_el.findall(sheet_tag):
        name = sheet_el.get("name", "")
        if name:
            sheet_names.append(name)

    # For each sheet, look up its relationship to a drawing
    drawing_rel_type = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing"
    )
    namelist = set(zf.namelist())

    for idx, sheet_name in enumerate(sheet_names, start=1):
        rels_path = f"xl/worksheets/_rels/sheet{idx}.xml.rels"
        if rels_path not in namelist:
            continue

        rels_root = etree.fromstring(zf.read(rels_path))
        rel_tag = f"{{{_RELS_NS}}}Relationship"

        for rel in rels_root.findall(rel_tag):
            if rel.get("Type") == drawing_rel_type:
                target = rel.get("Target", "")
                # Resolve relative path (e.g. "../drawings/drawing1.xml")
                if target.startswith("../"):
                    drawing_path = "xl/" + target.removeprefix("../")
                elif target.startswith("/"):
                    drawing_path = target.removeprefix("/")
                else:
                    drawing_path = f"xl/worksheets/{target}"

                if drawing_path in namelist:
                    results.append((sheet_name, drawing_path))
                break

    return results


def _extract_xlsx_shapes(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from shapes in an XLSX file via ZIP + lxml.

    Uses DrawingML ``<a:txBody>`` elements within each sheet's drawing XML.
    When run formatting varies or hyperlinks are present within a shape,
    inline HTML is emitted via ``_drawingml_to_html`` so the LLM can
    preserve it.

    Args:
        file_path: Path to the .xlsx file.

    Returns:
        list: (location_key, text) pairs with keys like
              'shape:{sheet_name}:{index}'.
    """
    texts: list[tuple[str, str]] = []

    with zipfile.ZipFile(file_path, "r") as zf:
        sheet_drawings = _resolve_xlsx_sheet_drawings(zf)
        namelist = set(zf.namelist())

        a_txbody_tag = f"{{{_DRAWINGML_NS}}}txBody"

        for sheet_name, drawing_path in sheet_drawings:
            root = etree.fromstring(zf.read(drawing_path))
            # Parse hyperlink relationships for this drawing
            rels_path = _get_rels_path(drawing_path)
            hyperlink_rels: dict[str, str] = {}
            if rels_path in namelist:
                hyperlink_rels = _parse_hyperlink_rels(zf.read(rels_path))

            for shape_idx, tx_body in enumerate(root.iter(a_txbody_tag)):
                has_links = bool(
                    hyperlink_rels and _has_drawingml_hyperlinks(tx_body),
                )
                if _has_drawingml_mixed_formatting(tx_body) or has_links:
                    text = _drawingml_to_html(
                        tx_body,
                        hyperlink_rels=hyperlink_rels,
                    ).strip()
                else:
                    text = _extract_drawingml_text(tx_body).strip()
                if text:
                    texts.append(
                        (f"shape:{sheet_name}:{shape_idx}", text),
                    )

    return texts


def _inject_xlsx_shapes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated text into XLSX shapes via ZIP + lxml.

    When translated text contains ``<a href="...">`` tags, hyperlink
    relationships are added to the drawing's ``.rels`` file.

    Args:
        output_path: Path to the .xlsx file to modify in place.
        translations: Mapping of ``'shape:{sheet_name}:{index}'``
                      to translated text.
    """
    with zipfile.ZipFile(output_path, "r") as zf:
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item.filename) for item in all_items}
        sheet_drawings = _resolve_xlsx_sheet_drawings(zf)

    a_txbody_tag = f"{{{_DRAWINGML_NS}}}txBody"
    modified = False

    for sheet_name, drawing_path in sheet_drawings:
        root = etree.fromstring(file_data[drawing_path])
        drawing_modified = False

        # Parse existing rels for hyperlink injection
        rels_path = _get_rels_path(drawing_path)
        current_rels: bytes | None = file_data.get(rels_path)

        def _make_rels_adder(
            rp: str = rels_path,
        ) -> Callable[[str], str]:
            """Creates a closure that adds hyperlink rels for a drawing."""

            def adder(url: str) -> str:
                nonlocal current_rels
                new_xml, r_id = _add_hyperlink_to_rels(current_rels, url)
                current_rels = new_xml
                file_data[rp] = new_xml
                return r_id

            return adder

        rels_adder = _make_rels_adder()

        for shape_idx, tx_body in enumerate(root.iter(a_txbody_tag)):
            key = f"shape:{sheet_name}:{shape_idx}"
            if key not in translations:
                continue
            translation = translations[key]
            if _FORMATTING_HTML_RE.search(translation):
                _inject_drawingml_html_runs(
                    tx_body,
                    translation,
                    rels_adder=rels_adder,
                )
            else:
                _inject_drawingml_text(tx_body, translation)
            drawing_modified = True

        if drawing_modified:
            file_data[drawing_path] = etree.tostring(
                root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            modified = True

    if not modified:
        return

    _rewrite_zip_content(output_path, file_data, all_items)


# ---------------------------------------------------------------------------
# ODF shape helpers — ZIP + lxml
# ---------------------------------------------------------------------------


def _build_odf_style_map(root: object) -> dict[str, object]:
    """Builds a mapping of style names to ``<style:style>`` elements.

    Scans ``<office:automatic-styles>`` for ``<style:style>`` entries
    with ``style:family="text"`` and returns a dict keyed by
    ``style:name``.

    Args:
        root: The lxml root element of an ODF ``content.xml``.

    Returns:
        dict: Mapping of style name to the ``<style:style>`` element.
    """
    style_map: dict[str, object] = {}
    auto_styles = root.find(
        f"{{{_ODF_NS['office']}}}automatic-styles",
    )
    if auto_styles is None:
        return style_map

    style_tag = f"{{{_ODF_NS['style']}}}style"
    style_name_attr = f"{{{_ODF_NS['style']}}}name"
    style_family_attr = f"{{{_ODF_NS['style']}}}family"

    for style_el in auto_styles.findall(style_tag):
        if style_el.get(style_family_attr) == "text":
            name = style_el.get(style_name_attr)
            if name:
                style_map[name] = style_el
    return style_map


def _inject_odf_text_box_html_runs(  # noqa: PLR0912, PLR0915
    text_box_el: object,
    html_text: str,
    text_p_tag: str,
    auto_styles_el: object,
    style_counter: list[int],
) -> bool:
    r"""Injects HTML-formatted text into an ODF ``<draw:text-box>`` element.

    Parses ``html_text`` via ``_parse_html_formatting``.  For each unique
    formatting signature, generates a ``<style:style>`` entry in
    ``auto_styles_el`` and wraps the text in ``<text:span>`` with the
    corresponding ``text:style-name``.  Handles ``'\n'`` by creating
    multiple ``<text:p>`` elements.

    Falls back to ``_inject_odf_paragraph_text`` when no HTML tags are
    detected.

    Args:
        text_box_el: An lxml element for ``<draw:text-box>``.
        html_text: Translated text with inline formatting tags.
        text_p_tag: The fully-qualified ``<text:p>`` tag name.
        auto_styles_el: The ``<office:automatic-styles>`` element.
        style_counter: Mutable ``[int]`` counter for unique style names.

    Returns:
        True if the element was modified.
    """
    # Fallback: no formatting tags → plain text
    if not _FORMATTING_HTML_RE.search(html_text):
        return _inject_odf_paragraph_text(text_box_el, html_text, text_p_tag)

    segments = _parse_html_formatting(html_text)
    if not segments:
        # Strip residual tags so literal HTML doesn't appear in the document
        plain = html.unescape(_STRIP_FORMAT_TAGS_RE.sub("", html_text))
        return _inject_odf_paragraph_text(text_box_el, plain, text_p_tag)

    text_span_tag = f"{{{_ODF_NS['text']}}}span"
    text_a_tag = f"{{{_ODF_NS['text']}}}a"
    xlink_href_attr = f"{{{_ODF_NS['xlink']}}}href"
    xlink_type_attr = f"{{{_ODF_NS['xlink']}}}type"
    text_style_attr = f"{{{_ODF_NS['text']}}}style-name"
    style_tag = f"{{{_ODF_NS['style']}}}style"
    style_name_attr = f"{{{_ODF_NS['style']}}}name"
    style_family_attr = f"{{{_ODF_NS['style']}}}family"
    text_props_tag = f"{{{_ODF_NS['style']}}}text-properties"
    fo_ns = _ODF_NS["fo"]
    style_ns = _ODF_NS["style"]

    # Build unique styles for each distinct formatting signature (incl. bg)
    sig_to_style: dict[
        tuple[
            bool,
            bool,
            bool,
            bool,
            bool,
            bool,
            float | None,
            str | None,
            str | None,
        ],
        str,
    ] = {}
    for seg in segments:
        sig = (
            seg.bold,
            seg.italic,
            seg.underline,
            seg.strike,
            seg.superscript,
            seg.subscript,
            seg.font_size_pt,
            seg.color_hex,
            seg.bg_color_hex,
        )
        if sig in sig_to_style:
            continue
        # Skip style creation for default formatting (all False/None)
        if not any(
            [
                seg.bold,
                seg.italic,
                seg.underline,
                seg.strike,
                seg.superscript,
                seg.subscript,
                seg.font_size_pt,
                seg.color_hex,
                seg.bg_color_hex,
            ]
        ):
            sig_to_style[sig] = ""  # empty = no span wrapping
            continue

        style_name = f"_ft_{style_counter[0]}"
        style_counter[0] += 1

        style_el = etree.SubElement(auto_styles_el, style_tag)
        style_el.set(style_name_attr, style_name)
        style_el.set(style_family_attr, "text")

        tp = etree.SubElement(style_el, text_props_tag)
        if seg.bold:
            tp.set(f"{{{fo_ns}}}font-weight", "bold")
        if seg.italic:
            tp.set(f"{{{fo_ns}}}font-style", "italic")
        if seg.underline:
            tp.set(f"{{{style_ns}}}text-underline-style", "solid")
            tp.set(f"{{{style_ns}}}text-underline-width", "auto")
        if seg.strike:
            tp.set(f"{{{style_ns}}}text-line-through-style", "solid")
        if seg.superscript:
            tp.set(f"{{{style_ns}}}text-position", "super 58%")
        elif seg.subscript:
            tp.set(f"{{{style_ns}}}text-position", "sub 58%")
        if seg.font_size_pt is not None:
            tp.set(f"{{{fo_ns}}}font-size", f"{seg.font_size_pt}pt")
        if seg.color_hex is not None:
            tp.set(f"{{{fo_ns}}}color", seg.color_hex.lower())
        if seg.bg_color_hex is not None:
            tp.set(f"{{{fo_ns}}}background-color", seg.bg_color_hex.lower())

        sig_to_style[sig] = style_name

    # Split segments by newlines into paragraph groups
    para_groups: list[list[_FormattedSegment]] = [[]]
    for seg in segments:
        if "\n" in seg.text:
            lines = seg.text.split("\n")
            for line_idx, line in enumerate(lines):
                if line:
                    para_groups[-1].append(seg._replace(text=line))
                if line_idx < len(lines) - 1:
                    para_groups.append([])
        else:
            para_groups[-1].append(seg)

    # Save paragraph style-names from existing <text:p> elements
    existing_p = text_box_el.findall(text_p_tag)
    saved_p_styles: list[str] = []
    for p_el in existing_p:
        saved_p_styles.append(p_el.get(text_style_attr, ""))

    for p_el in existing_p:
        text_box_el.remove(p_el)

    # Build new <text:p> elements, restoring paragraph style-names
    for g_idx, group in enumerate(para_groups):
        new_p = etree.SubElement(text_box_el, text_p_tag)
        # Re-apply paragraph style from corresponding original paragraph
        if g_idx < len(saved_p_styles) and saved_p_styles[g_idx]:
            new_p.set(text_style_attr, saved_p_styles[g_idx])
        elif saved_p_styles:
            # Fall back to last known paragraph style
            new_p.set(text_style_attr, saved_p_styles[-1])
        for seg in group:
            # Hyperlink segments get <text:a> regardless of formatting
            if seg.hyperlink_url:
                a_el = etree.SubElement(new_p, text_a_tag)
                a_el.set(xlink_href_attr, seg.hyperlink_url)
                a_el.set(xlink_type_attr, "simple")
                a_el.text = seg.text
                continue

            sig = (
                seg.bold,
                seg.italic,
                seg.underline,
                seg.strike,
                seg.superscript,
                seg.subscript,
                seg.font_size_pt,
                seg.color_hex,
                seg.bg_color_hex,
            )
            style_name = sig_to_style.get(sig, "")
            if style_name:
                span = etree.SubElement(new_p, text_span_tag)
                span.set(text_style_attr, style_name)
                span.text = seg.text
            else:
                # Default formatting — append as direct text
                children = list(new_p)
                if children:
                    # Append as tail of last child
                    last = children[-1]
                    last.tail = (last.tail or "") + seg.text
                else:
                    new_p.text = (new_p.text or "") + seg.text

    return True


def _extract_odt_shapes(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from ``<draw:text-box>`` elements in an ODT file.

    When span formatting varies within a text box, inline HTML is emitted
    via ``_odf_text_box_to_html`` so the LLM can preserve it.

    Args:
        file_path: Path to the .odt file.

    Returns:
        list: (location_key, text) pairs with keys like ``'shape:{index}'``.
    """
    texts: list[tuple[str, str]] = []

    with zipfile.ZipFile(file_path, "r") as zf:
        content_xml = zf.read("content.xml")

    root = etree.fromstring(content_xml)
    text_p_tag = f"{{{_ODF_NS['text']}}}p"
    style_map = _build_odf_style_map(root)

    for idx, text_box in enumerate(
        root.findall(".//draw:text-box", _ODF_NS),
    ):
        if _has_odf_text_box_mixed_formatting(text_box, style_map, text_p_tag):
            text = _odf_text_box_to_html(text_box, style_map, text_p_tag)
        else:
            text = _extract_odf_paragraph_text(text_box, text_p_tag)
        if text:
            texts.append((f"shape:{idx}", text))

    return texts


def _inject_odt_shapes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated text into ``<draw:text-box>`` elements in an ODT.

    When the translated text contains inline HTML formatting tags,
    ``_inject_odf_text_box_html_runs`` is used to create styled spans.
    Otherwise, falls back to ``_inject_odf_paragraph_text``.

    Args:
        output_path: Path to the .odt file to modify in place.
        translations: Mapping of ``'shape:{index}'`` to translated text.
    """
    with zipfile.ZipFile(output_path, "r") as zf:
        content_xml = zf.read("content.xml")
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item.filename) for item in all_items}

    root = etree.fromstring(content_xml)
    text_p_tag = f"{{{_ODF_NS['text']}}}p"
    modified = False

    # Find or create <office:automatic-styles>
    auto_styles = root.find(
        f"{{{_ODF_NS['office']}}}automatic-styles",
    )
    if auto_styles is None:
        auto_styles = etree.SubElement(
            root,
            f"{{{_ODF_NS['office']}}}automatic-styles",
        )
    style_counter: list[int] = [0]

    for idx, text_box in enumerate(
        root.findall(".//draw:text-box", _ODF_NS),
    ):
        key = f"shape:{idx}"
        if key in translations:
            translation = translations[key]
            if _FORMATTING_HTML_RE.search(translation):
                modified |= _inject_odf_text_box_html_runs(
                    text_box,
                    translation,
                    text_p_tag,
                    auto_styles,
                    style_counter,
                )
            else:
                modified |= _inject_odf_paragraph_text(
                    text_box,
                    translation,
                    text_p_tag,
                )

    if not modified:
        return

    file_data["content.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
    )
    _rewrite_zip_content(output_path, file_data, all_items)


def _extract_ods_shapes(file_path: Path) -> list[tuple[str, str]]:
    """Extracts text from ``<draw:text-box>`` elements in an ODS file.

    Iterates per ``<table:table>`` to produce sheet-qualified keys.
    When span formatting varies within a text box, inline HTML is emitted
    via ``_odf_text_box_to_html`` so the LLM can preserve it.

    Args:
        file_path: Path to the .ods file.

    Returns:
        list: (location_key, text) pairs with keys like
              ``'shape:{sheet_name}:{index}'``.
    """
    texts: list[tuple[str, str]] = []

    with zipfile.ZipFile(file_path, "r") as zf:
        content_xml = zf.read("content.xml")

    root = etree.fromstring(content_xml)
    table_tag = f"{{{_ODF_NS['table']}}}table"
    table_name_attr = f"{{{_ODF_NS['table']}}}name"
    text_p_tag = f"{{{_ODF_NS['text']}}}p"
    style_map = _build_odf_style_map(root)

    for table in root.iter(table_tag):
        sheet_name = table.get(table_name_attr, "Sheet")
        for shape_idx, text_box in enumerate(
            table.findall(".//draw:text-box", _ODF_NS),
        ):
            if _has_odf_text_box_mixed_formatting(
                text_box,
                style_map,
                text_p_tag,
            ):
                text = _odf_text_box_to_html(text_box, style_map, text_p_tag)
            else:
                text = _extract_odf_paragraph_text(text_box, text_p_tag)
            if text:
                texts.append(
                    (f"shape:{sheet_name}:{shape_idx}", text),
                )

    return texts


def _inject_ods_shapes(
    output_path: Path,
    translations: dict[str, str],
) -> None:
    """Injects translated text into ``<draw:text-box>`` elements in an ODS.

    When the translated text contains inline HTML formatting tags,
    ``_inject_odf_text_box_html_runs`` is used to create styled spans.
    Otherwise, falls back to ``_inject_odf_paragraph_text``.

    Args:
        output_path: Path to the .ods file to modify in place.
        translations: Mapping of ``'shape:{sheet_name}:{index}'``
                      to translated text.
    """
    with zipfile.ZipFile(output_path, "r") as zf:
        content_xml = zf.read("content.xml")
        all_items = zf.infolist()
        file_data = {item.filename: zf.read(item.filename) for item in all_items}

    root = etree.fromstring(content_xml)
    table_tag = f"{{{_ODF_NS['table']}}}table"
    table_name_attr = f"{{{_ODF_NS['table']}}}name"
    text_p_tag = f"{{{_ODF_NS['text']}}}p"
    modified = False

    # Find or create <office:automatic-styles>
    auto_styles = root.find(
        f"{{{_ODF_NS['office']}}}automatic-styles",
    )
    if auto_styles is None:
        auto_styles = etree.SubElement(
            root,
            f"{{{_ODF_NS['office']}}}automatic-styles",
        )
    style_counter: list[int] = [0]

    for table in root.iter(table_tag):
        sheet_name = table.get(table_name_attr, "Sheet")
        for shape_idx, text_box in enumerate(
            table.findall(".//draw:text-box", _ODF_NS),
        ):
            key = f"shape:{sheet_name}:{shape_idx}"
            if key in translations:
                translation = translations[key]
                if _FORMATTING_HTML_RE.search(translation):
                    modified |= _inject_odf_text_box_html_runs(
                        text_box,
                        translation,
                        text_p_tag,
                        auto_styles,
                        style_counter,
                    )
                else:
                    modified |= _inject_odf_paragraph_text(
                        text_box,
                        translation,
                        text_p_tag,
                    )

    if not modified:
        return

    file_data["content.xml"] = etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
    )
    _rewrite_zip_content(output_path, file_data, all_items)


def _translate_single_image(  # noqa: PLR0913
    image_bytes: bytes,
    content_type: str,
    target_lang: str,
    src_lang: str,
    glossary_entries: list[tuple[int, str, str]] | None,
    ocr_method: str,
    *,
    provider: str | None = None,
    model: str | None = None,
) -> bytes | None:
    """Translates a single image using the OCR → LLM → render pipeline.

    Writes the image to a temp file, processes it, and returns the
    translated image bytes. Returns None if the image has no translatable
    text or rendering fails. Does not catch ValueError so that fatal LLM
    errors can propagate to the caller.

    Args:
        image_bytes: Raw image data.
        content_type: MIME type (e.g. "image/png").
        target_lang: Target language name.
        src_lang: Source language name.
        glossary_entries: Optional glossary entries.
        ocr_method: OCR method name (e.g. "TesseractOCR").
        provider: Optional LLM provider override.
        model: Optional LLM model override.

    Returns:
        bytes | None: Translated image bytes, or None.
    """
    from src.core.image_processor import process_image_translation  # noqa: PLC0415
    from src.core.layout_analysis import merge_to_paragraphs  # noqa: PLC0415
    from src.core.llm_engine import translate_image_content  # noqa: PLC0415
    from src.core.ocr_engine import run_ocr  # noqa: PLC0415

    ext = _IMAGE_TYPE_TO_EXT.get(content_type)
    if ext is None:
        return None

    with tempfile.TemporaryDirectory(prefix="ftrans_img_") as tmp_dir:
        tmp_path = Path(tmp_dir)
        input_path = tmp_path / f"input{ext}"
        output_path = tmp_path / f"output{ext}"

        input_path.write_bytes(image_bytes)

        # 1. OCR
        ocr_results = run_ocr(str(input_path), method=ocr_method, src_lang=src_lang)
        if not ocr_results:
            return None  # No text detected

        raw_ocr_results = list(ocr_results)

        # 2. LLM translation (may raise ValueError for fatal errors)
        paragraph_data = translate_image_content(
            str(input_path),
            ocr_results,
            target_lang,
            src_lang,
            glossary_entries=glossary_entries,
            provider=provider,
            model=model,
        )

        # 3. Merge paragraphs
        merged_results, translations, raw_fragments = merge_to_paragraphs(
            paragraph_data,
            raw_ocr_results,
            ocr_method,
        )
        if not merged_results:
            return None  # No translatable text after merge

        # 4. Render translated image
        success = process_image_translation(
            str(input_path),
            str(output_path),
            merged_results,
            translations,
            target_lang=target_lang,
            raw_ocr_results=raw_fragments,
            ocr_method=ocr_method,
        )

        if success and output_path.exists():
            return output_path.read_bytes()
        return None


def _translate_zip_images(  # noqa: PLR0913, PLR0912
    output_path: Path,
    suffix: str,
    target_lang: str,
    src_lang: str,
    glossary_entries: list[tuple[int, str, str]] | None,
    ocr_method: str,
    progress_callback: Callable[[int], None] | None,
    cancel_check: Callable[[], bool] | None,
    *,
    provider: str | None = None,
    model: str | None = None,
    checkpoint_dir: Path | None = None,
) -> None:
    """Translates images embedded in an Office document using zipfile.

    Opens the document as a ZIP archive, identifies raster images in the
    known media directory, translates each via the OCR → LLM → render
    pipeline, replaces the originals in memory, and rewrites the archive
    atomically (write to ``.tmp``, then ``shutil.move``).

    Supports ``.docx``, ``.xlsx``, ``.pptx``, ``.odt``, ``.ods``, ``.odp``,
    and ``.epub``.

    **Skip-with-warning** policy for non-fatal per-image errors: a
    bad image (e.g. ``IMAGE_TOO_LARGE``, an unreadable JPEG header,
    a vision model returning empty text) leaves the original image in
    place and the loop continues.  The user gets a document with
    most images translated and the broken ones in their source form,
    rather than one stubborn image blocking the whole document.
    Fatal LLM errors (AUTH_ERROR, QUOTA_ERROR, VISION_NOT_SUPPORTED)
    still break out immediately — those indicate the entire pipeline
    can't continue, not "this one image won't translate".

    When ``checkpoint_dir`` is provided, each image's translated bytes
    are persisted under ``<checkpoint_dir>/office_images/<sha256>.bin``
    and consulted on re-runs.  This means an interrupted batch (50/100
    images done, then a quota error or cancellation) only retries the
    remaining 50 on resume instead of redoing the whole document.  The
    SHA256 of the source bytes is the cache key, so duplicate images
    (e.g. a company logo repeated on every page) deduplicate naturally.

    Args:
        output_path: Path to the saved translated document (modified in place).
        suffix: Lowercase file extension (e.g. ".docx").
        target_lang: Target language name.
        src_lang: Source language name.
        glossary_entries: Optional glossary entries.
        ocr_method: OCR method name (e.g. "TesseractOCR").
        progress_callback: Called with 0-100 for the image phase.
        cancel_check: Returns True if the task was cancelled.
        provider: Optional LLM provider override.
        model: Optional LLM model override.
        checkpoint_dir: Task storage directory for per-image cache.
            ``None`` disables caching (used by tests and the
            legacy-format conversion path's intermediate temp file).
    """
    media_prefixes = _SUFFIX_TO_MEDIA_PREFIXES.get(suffix)
    if media_prefixes is None:
        if progress_callback:
            progress_callback(100)
        return

    # Read the entire ZIP into memory (same pattern as _inject_odf_comments)
    with zipfile.ZipFile(output_path, "r") as zf:
        all_items = zf.infolist()
        file_data: dict[str, bytes] = {
            item.filename: zf.read(item.filename) for item in all_items
        }

    # Identify translatable images by path prefix + extension
    image_entries: list[str] = [
        fn
        for fn in file_data
        if any(fn.startswith(p) for p in media_prefixes)
        and Path(fn).suffix.lower() in _EXT_TO_IMAGE_CONTENT_TYPE
    ]

    if not image_entries:
        if progress_callback:
            progress_callback(100)
        return

    total = len(image_entries)
    modified = False

    for i, filename in enumerate(image_entries):
        if cancel_check and cancel_check():
            break

        image_bytes = file_data[filename]
        content_type = _EXT_TO_IMAGE_CONTENT_TYPE[Path(filename).suffix.lower()]

        # Consult the per-image cache before paying for OCR + LLM.
        # The hash is content-addressed, so identical images anywhere
        # in the document (or a sibling document in the same run)
        # collapse to one translation.  Cache misses fall through to
        # the live pipeline; cache hits skip straight to substitution.
        # ``image_hash`` is None when caching is disabled (tests; the
        # legacy round-trip's intermediate temp file when the caller
        # explicitly passes ``checkpoint_dir=None``).  Keeping the two
        # nullables paired lets the explicit ``is not None`` guards
        # at the save callsite (below) double as a type-narrowing
        # contract, preventing accidental writes to a None directory.
        image_hash: str | None = None
        translated: bytes | None = None
        if checkpoint_dir is not None:
            image_hash = hash_office_image(image_bytes)
            translated = load_office_image_checkpoint(checkpoint_dir, image_hash)

        if translated is None:
            try:
                translated = _translate_single_image(
                    image_bytes,
                    content_type,
                    target_lang,
                    src_lang,
                    glossary_entries,
                    ocr_method,
                    provider=provider,
                    model=model,
                )
            except ValueError as e:
                error_tag = str(e)
                # ``_is_fatal_llm_error`` strips the ``:Service``
                # suffix (e.g. ``"AUTH_ERROR:Gemini"``) so the
                # suffix-bearing variants the engine raises still
                # qualify as fatal.
                if _is_fatal_llm_error(error_tag):
                    logger.error(
                        "Fatal error translating image %s in %s: %s",
                        filename,
                        output_path.name,
                        error_tag,
                    )
                    raise
                # Non-fatal per-image failure (the LLM-level
                # ``@retry_api_call`` has already burnt its 3
                # attempts).  Examples: ``IMAGE_TOO_LARGE``,
                # OCR returned no text, vision model returned empty,
                # transient ``CONNECTION_ERROR`` that won't recover.
                # Skip the image — keep the original in place — and
                # let the rest of the document complete.  Contrast
                # with fatal LLM errors above (AUTH / QUOTA /
                # VISION_NOT_SUPPORTED) which raise because they
                # block every remaining image too.  WARNING level
                # surfaces the issue in ``app.log`` without marking
                # the whole document failed.
                logger.warning(
                    "Skipping image %s in %s: %s — keeping original",
                    filename,
                    output_path.name,
                    error_tag,
                )
                translated = None
            except Exception:
                # Same skip-with-warning policy for unexpected
                # exceptions — a single bad image shouldn't sink the
                # whole document.
                logger.warning(
                    "Unexpected error translating image %s in %s — keeping original",
                    filename,
                    output_path.name,
                    exc_info=True,
                )
                translated = None

            # Persist the freshly translated bytes so a resume after
            # cancellation / transient failure skips them next run.
            # Both guards are needed: ``image_hash`` is None when
            # caching is disabled, and ``translated`` is None when
            # the LLM call failed.
            if (
                translated is not None
                and checkpoint_dir is not None
                and image_hash is not None
            ):
                save_office_image_checkpoint(
                    checkpoint_dir,
                    image_hash,
                    translated,
                )

        if translated is not None:
            file_data[filename] = translated
            modified = True

        if progress_callback:
            progress_callback(int(((i + 1) / total) * 100))

    # Atomic rewrite if any images were modified.  Per the skip-with-
    # warning policy above, non-fatal per-image failures do not raise
    # — the document completes with originals in place for the
    # untranslatable ones.
    if modified:
        _rewrite_zip_content(output_path, file_data, all_items)


def _translate_legacy_images(  # noqa: PLR0913
    output_path: Path,
    suffix: str,
    backend: str,
    target_lang: str,
    src_lang: str,
    glossary_entries: list[tuple[int, str, str]] | None,
    ocr_method: str,
    progress_callback: Callable[[int], None] | None,
    cancel_check: Callable[[], bool] | None,
    *,
    provider: str | None = None,
    model: str | None = None,
    checkpoint_dir: Path | None = None,
) -> None:
    """Translates images in legacy office files via round-trip conversion.

    Converts the legacy file (.doc/.xls/.ppt) to its modern equivalent
    (.docx/.xlsx/.pptx), runs the existing ZIP-based image pipeline on
    the modern file, then converts back to the legacy format.

    Args:
        output_path: Path to the saved legacy document (modified in place).
        suffix: Lowercase legacy extension (e.g. ".doc").
        backend: Backend identifier ("win32com" or "uno").
        target_lang: Target language name.
        src_lang: Source language name.
        glossary_entries: Optional glossary entries.
        ocr_method: OCR method name (e.g. "TesseractOCR").
        progress_callback: Called with 0-100 for the image phase.
        cancel_check: Returns True if the task was cancelled.
        provider: Optional LLM provider override.
        model: Optional LLM model override.
        checkpoint_dir: Task storage directory for per-image cache.
            Forwarded to ``_translate_zip_images``; image hashes are
            content-keyed so caching still works across the legacy ↔
            modern round-trip as long as the conversion preserves the
            embedded raster bytes (which UNO does; win32com is best
            effort).
    """
    modern_suffix = LEGACY_CONVERT_MAP[suffix]
    convert_fn = (
        _convert_with_win32com if backend == _BACKEND_WIN32COM else _convert_with_uno
    )

    # Create temp file with the modern extension
    tmp_fd, tmp_str = tempfile.mkstemp(suffix=modern_suffix)
    tmp_path = Path(tmp_str)
    try:
        os.close(tmp_fd)

        # 1. Convert legacy → modern
        convert_fn(output_path, tmp_path)
        if not tmp_path.exists() or tmp_path.stat().st_size == 0:
            msg = f"Legacy-to-modern conversion produced no output: {tmp_path}"
            raise RuntimeError(msg)

        # 2. Translate images in the modern ZIP-based format.
        # ``checkpoint_dir`` is the *task* storage dir, not the modern
        # temp file — image hashes are content-keyed so caching still
        # works across the legacy ↔ modern round-trip as long as the
        # extracted image bytes are stable (which they are: legacy
        # ↔ modern conversion preserves the embedded raster blobs).
        _translate_zip_images(
            tmp_path,
            modern_suffix,
            target_lang,
            src_lang,
            glossary_entries,
            ocr_method,
            progress_callback,
            cancel_check,
            provider=provider,
            model=model,
            checkpoint_dir=checkpoint_dir,
        )

        # 3. Convert modern → legacy (overwrite output)
        convert_fn(tmp_path, output_path)
    finally:
        tmp_path.unlink(missing_ok=True)


def _translate_doc_images(  # noqa: PLR0913
    output_path: Path,
    suffix: str,
    backend: str,
    target_lang: str,
    src_lang: str,
    glossary_entries: list[tuple[int, str, str]] | None,
    progress_callback: Callable[[int], None] | None,
    cancel_check: Callable[[], bool] | None,
    config: TranslationConfig | None = None,
    *,
    provider: str | None = None,
    model: str | None = None,
    checkpoint_dir: Path | None = None,
) -> None:
    """Translates images embedded in an Office document.

    For modern/ODF formats: uses the ZIP-based image pipeline directly.
    For legacy formats (.doc/.xls/.ppt): converts to modern format first,
    runs the ZIP pipeline, then converts back.

    Args:
        output_path: Path to the saved translated document.
        suffix: Lowercase file extension (e.g. ".docx", ".doc").
        backend: Backend identifier for legacy format conversion.
        target_lang: Target language name.
        src_lang: Source language name.
        glossary_entries: Optional glossary entries.
        progress_callback: Called with 0-100 for the image phase.
        cancel_check: Returns True if the task was cancelled.
        config: Optional TranslationConfig snapshot; falls back to load_setting().
        provider: Optional LLM provider override.
        model: Optional LLM model override.
        checkpoint_dir: Task storage directory for per-image cache.
            Forwarded to the underlying ZIP pipeline; ``None`` disables
            caching.
    """
    if config is not None:
        ocr_method = config.ocr_method
    else:
        from src.constants.ocr import OCR_METHOD_TESSERACT  # noqa: PLC0415
        from src.constants.settings import SETTING_OCR_METHOD  # noqa: PLC0415
        from src.utils.config_manager import load_setting  # noqa: PLC0415

        ocr_method = load_setting(SETTING_OCR_METHOD, OCR_METHOD_TESSERACT)

    if suffix in _LEGACY_EXTENSIONS:
        _translate_legacy_images(
            output_path,
            suffix,
            backend,
            target_lang,
            src_lang,
            glossary_entries,
            ocr_method,
            progress_callback,
            cancel_check,
            provider=provider,
            model=model,
            checkpoint_dir=checkpoint_dir,
        )
    else:
        _translate_zip_images(
            output_path,
            suffix,
            target_lang,
            src_lang,
            glossary_entries,
            ocr_method,
            progress_callback,
            cancel_check,
            provider=provider,
            model=model,
            checkpoint_dir=checkpoint_dir,
        )


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def process_office_file(  # noqa: PLR0913, PLR0912, PLR0915
    file_path: Path,
    output_path: Path,
    target_lang: str,
    src_lang: str = "",
    progress_callback: Callable[[int], None] | None = None,
    glossary_entries: list[tuple[int, str, str]] | None = None,
    cancel_check: Callable[[], bool] | None = None,
    checkpoint_dir: Path | None = None,
    config: TranslationConfig | None = None,
    *,
    provider: str | None = None,
    model: str | None = None,
) -> bool:
    """Translates an Office document using the best available backend.

    Extracts translatable text, translates via LLM, and injects
    translations back into a copy of the document.

    Args:
        file_path: Path to the source office file.
        output_path: Path to write the translated file.
        target_lang: Target language name.
        src_lang: Source language name.
        progress_callback: Called with 0-100 progress percentage.
        glossary_entries: Optional glossary entries for translation.
        cancel_check: Returns True if the task was cancelled.
        checkpoint_dir: Directory for saving/loading checkpoints.
        config: Optional TranslationConfig for dependency injection.
        provider: Optional LLM provider override (Gemini / Custom).
        model: Optional LLM model override.

    Returns:
        bool: True on success, False if cancelled.

    Raises:
        ValueError: On backend or processing errors.
    """
    suffix = file_path.suffix.lower()
    lo_path = config.libreoffice_path if config is not None else ""
    backend = _detect_backend(suffix, lo_path)
    category = _get_file_category(suffix)
    do_images = _should_translate_images(suffix, backend, config)
    do_comments = _should_translate_comments(suffix, backend, config)
    do_shapes = _should_translate_shapes(suffix, backend, config)
    do_notes = _should_translate_notes(suffix, backend, config)
    do_sheet_names = _should_translate_sheet_names(suffix, backend, config)

    logger.debug(
        "Processing %s with backend=%s, category=%s"
        " (images=%s, comments=%s, shapes=%s, notes=%s, sheet_names=%s)",
        file_path.name,
        backend,
        category,
        do_images,
        do_comments,
        do_shapes,
        do_notes,
        do_sheet_names,
    )

    # Extract translatable text
    try:
        extract_fn = _EXTRACTORS[backend][category]
        texts = extract_fn(file_path)
    except ValueError:
        raise
    except RuntimeError as e:
        # UNO connection failure (soffice not installed / not startable)
        logger.error(
            "UNO connection failed for %s: %s",
            file_path.name,
            e,
        )
        raise ValueError("OFFICE_CONVERTER_NOT_FOUND") from e
    except Exception as e:
        logger.error(
            "Failed to extract text from %s: %s",
            file_path.name,
            e,
        )
        raise ValueError("TEXT_READ_ERROR") from e

    # Extract comments if enabled
    if do_comments:
        try:
            texts.extend(_extract_comments(file_path, suffix, backend))
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Comment extraction failed for %s: %s",
                file_path.name,
                e,
            )

    # Extract shapes / text boxes if enabled
    if do_shapes:
        try:
            texts.extend(_extract_shapes(file_path, suffix, backend))
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Shape extraction failed for %s: %s",
                file_path.name,
                e,
            )

    # Extract headers and footers (always-on for word-processor formats)
    if suffix in _HEADER_FOOTER_EXTENSIONS:
        try:
            texts.extend(
                _extract_headers_footers(file_path, suffix, backend),
            )
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Header/footer extraction failed for %s: %s",
                file_path.name,
                e,
            )

    # Extract footnotes and endnotes (always-on for word-processor formats)
    if suffix in _FOOTNOTE_EXTENSIONS:
        try:
            texts.extend(
                _extract_footnotes(file_path, suffix, backend),
            )
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Footnote extraction failed for %s: %s",
                file_path.name,
                e,
            )

    # Extract speaker notes if enabled
    if do_notes:
        try:
            texts.extend(_extract_notes(file_path, suffix, backend))
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Speaker notes extraction failed for %s: %s",
                file_path.name,
                e,
            )

    # Extract sheet names if enabled
    if do_sheet_names:
        try:
            texts.extend(
                _extract_sheet_names(file_path, suffix, backend),
            )
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Sheet name extraction failed for %s: %s",
                file_path.name,
                e,
            )

    if not texts:
        # No translatable text — copy file as-is
        output_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(file_path, output_path)

        # Still translate embedded images if enabled
        if do_images:
            try:
                _translate_doc_images(
                    output_path,
                    suffix,
                    backend,
                    target_lang,
                    src_lang,
                    glossary_entries,
                    progress_callback,
                    cancel_check,
                    config=config,
                    provider=provider,
                    model=model,
                    checkpoint_dir=checkpoint_dir,
                )
            except ValueError:
                # Fatal LLM error (auth, quota, service down) — propagate
                raise
            except Exception as e:  # noqa: BLE001
                logger.warning(
                    "Image translation failed for %s: %s",
                    file_path.name,
                    e,
                )
        if not do_images and progress_callback:
            progress_callback(100)
        return True

    # Check cancellation before translation
    if cancel_check and cancel_check():
        return False

    # Wrap progress callback: text gets 0-70%, images get 70-100%
    text_progress = progress_callback
    if do_images and progress_callback:

        def text_progress(pct: int) -> None:
            progress_callback(int(pct * 0.7))

    # Translate
    keys = [k for k, _ in texts]
    values = [v for _, v in texts]

    # Excel cells are short isolated values; word/ppt contain prose paragraphs.
    # Mixed-formatting paragraphs or text-box shapes may contain inline HTML
    # tags — use CONTENT_HTML when any value has formatting tags, even for
    # Excel (shapes in .xlsx/.xls/.ods can have bold+italic text boxes).
    has_html = any(_FORMATTING_HTML_RE.search(v) for v in values)
    if has_html:
        ct = CONTENT_HTML
    elif category == "excel":
        ct = CONTENT_DATA_VALUES
    else:
        ct = CONTENT_PLAIN_TEXT
    translated_values = translate_batch(
        values,
        target_lang,
        src_lang,
        text_progress,
        glossary_entries,
        cancel_check,
        checkpoint_dir=checkpoint_dir,
        content_type=ct,
        provider=provider,
        model=model,
    )
    if translated_values is None:
        return False  # Cancelled

    # Build translation map
    translations = dict(zip(keys, translated_values, strict=True))

    # Inject translations back
    try:
        inject_fn = _INJECTORS[backend][category]
        inject_fn(file_path, output_path, translations, target_lang)
    except ValueError:
        raise
    except RuntimeError as e:
        # UNO connection failure (soffice not installed / not startable)
        logger.error(
            "UNO connection failed for %s: %s",
            file_path.name,
            e,
        )
        raise ValueError("OFFICE_CONVERTER_NOT_FOUND") from e
    except Exception as e:
        logger.error(
            "Failed to write translated %s: %s",
            file_path.name,
            e,
        )
        raise ValueError("TEXT_WRITE_ERROR") from e

    # Inject translated comments if enabled
    if do_comments:
        try:
            _inject_comments(output_path, translations, suffix, backend)
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Comment injection failed for %s: %s",
                file_path.name,
                e,
            )

    # Inject translated shapes / text boxes if enabled
    if do_shapes:
        try:
            _inject_shapes(output_path, translations, suffix, backend)
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Shape injection failed for %s: %s",
                file_path.name,
                e,
            )

    # Inject translated headers/footers (always-on)
    if suffix in _HEADER_FOOTER_EXTENSIONS:
        try:
            _inject_headers_footers(
                output_path,
                translations,
                suffix,
                backend,
            )
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Header/footer injection failed for %s: %s",
                file_path.name,
                e,
            )

    # Inject translated footnotes/endnotes (always-on)
    if suffix in _FOOTNOTE_EXTENSIONS:
        try:
            _inject_footnotes(
                output_path,
                translations,
                suffix,
                backend,
            )
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Footnote injection failed for %s: %s",
                file_path.name,
                e,
            )

    # Inject translated speaker notes if enabled
    if do_notes:
        try:
            _inject_notes(output_path, translations, suffix, backend)
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Speaker notes injection failed for %s: %s",
                file_path.name,
                e,
            )

    # Inject translated sheet names if enabled (must be after main inject
    # since cell location keys reference original sheet names)
    if do_sheet_names:
        try:
            _inject_sheet_names(
                output_path,
                translations,
                suffix,
                backend,
            )
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Sheet name injection failed for %s: %s",
                file_path.name,
                e,
            )

    # Translate embedded images if enabled
    if do_images:

        def img_progress(pct: int) -> None:
            if progress_callback:
                progress_callback(70 + int(pct * 0.3))

        try:
            _translate_doc_images(
                output_path,
                suffix,
                backend,
                target_lang,
                src_lang,
                glossary_entries,
                img_progress,
                cancel_check,
                config=config,
                provider=provider,
                model=model,
                checkpoint_dir=checkpoint_dir,
            )
        except ValueError:
            # Fatal LLM error (auth, quota, service down) — propagate
            raise
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Image translation failed for %s: %s",
                file_path.name,
                e,
            )

    return True
